"""
Close Cockpit
-------------
Orchestrates a full quarter-close package from a single run: an updated
version of the prior deck, a variance analysis derived from the same
numbers that were just updated (no separate upload needed), a
benchmarking slide if peer data is provided, and an action item tracker
if meeting notes are provided. Everything is assembled into a client-ready
executive package alongside the updated deck.

This reuses the extraction, matching, and period-detection logic that was
built and hardened in Deck Refresh, and the benchmarking and meeting-notes
logic from Engagement Hub. Every generated PowerPoint text box goes through
a single helper that sets word_wrap and auto_size explicitly, the fix that
was needed after a real overflow bug surfaced earlier in this codebase.

Single self-contained Flask application. No external services.
"""

import os
import re
import io
import json
import time
import uuid
import shutil
import zipfile
import subprocess
import tempfile
import threading
from statistics import median
from collections import Counter

from flask import Flask, request, render_template, send_file, redirect, url_for, flash, abort
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from rapidfuzz import fuzz
import openpyxl
from openpyxl.styles import Font, PatternFill, Alignment, Border, Side

BASE_DIR = os.path.dirname(os.path.abspath(__file__))
SESSIONS_DIR = os.path.join(BASE_DIR, "sessions")
os.makedirs(SESSIONS_DIR, exist_ok=True)

app = Flask(__name__)
app.secret_key = "close-cockpit-local-secret"
app.config["MAX_CONTENT_LENGTH"] = 80 * 1024 * 1024

KPMG_BLUE = RGBColor(0x00, 0x33, 0x8D)
KPMG_COBALT = RGBColor(0x00, 0x91, 0xDA)
KPMG_LIGHT = RGBColor(0xF2, 0xF5, 0xFA)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x1A, 0x1A, 0x1A)
GREY = RGBColor(0x77, 0x77, 0x77)
GOOD = RGBColor(0x1E, 0x84, 0x49)
BAD = RGBColor(0xC0, 0x39, 0x2B)
GOOD_HEX, BAD_HEX = "1E8449", "C0392B"
PALETTE = [KPMG_BLUE, KPMG_COBALT, RGBColor(0x00, 0xA3, 0x93), RGBColor(0x5A, 0x28, 0x8C), RGBColor(0x8A, 0x8D, 0x8F)]


@app.after_request
def _no_cache(response):
    response.headers["Cache-Control"] = "no-store, no-cache, must-revalidate, max-age=0"
    return response


# ---------------------------------------------------------------------------
# SHARED PPTX HELPERS (every text box sets word_wrap + auto_size explicitly)
# ---------------------------------------------------------------------------

def add_textbox(slide, x, y, w, h, text, size=14, bold=False, italic=False,
                 color=DARK, align=PP_ALIGN.LEFT, font="Arial", wrap=True):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = wrap
    tf.auto_size = MSO_AUTO_SIZE.NONE
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font
    return box


def add_header_bar(slide, prs, title_text, subtitle_text=None, tag_text=None):
    bar = slide.shapes.add_shape(1, 0, 0, prs.slide_width, Inches(0.9))
    bar.fill.solid()
    bar.fill.fore_color.rgb = KPMG_BLUE
    bar.line.fill.background()
    tf = bar.text_frame
    tf.word_wrap = False
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.margin_left = Inches(0.3)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = "KPMG"
    r.font.size = Pt(22)
    r.font.bold = True
    r.font.color.rgb = WHITE
    r.font.name = "Arial"

    add_textbox(slide, 3.4, 0.12, 9.6, 0.65, title_text, size=16, bold=True,
                color=WHITE, align=PP_ALIGN.RIGHT)
    if subtitle_text:
        add_textbox(slide, 0.4, 1.0, 9.5, 0.35, subtitle_text, size=11, italic=True, color=GREY)
    if tag_text:
        tag_box = slide.shapes.add_shape(1, Inches(11.2), Inches(1.02), Inches(1.75), Inches(0.3))
        tag_box.fill.solid()
        tag_box.fill.fore_color.rgb = KPMG_COBALT
        tag_box.line.fill.background()
        tf2 = tag_box.text_frame
        tf2.word_wrap = False
        tf2.auto_size = MSO_AUTO_SIZE.NONE
        p2 = tf2.paragraphs[0]
        p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run()
        r2.text = tag_text
        r2.font.size = Pt(8.5)
        r2.font.bold = True
        r2.font.color.rgb = WHITE
        r2.font.name = "Arial"


def add_kpi_tile(slide, x, y, w, h, value, label, accent=KPMG_BLUE):
    tile = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    tile.fill.solid()
    tile.fill.fore_color.rgb = KPMG_LIGHT
    tile.line.color.rgb = accent
    tile.line.width = Pt(1)
    tf = tile.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.margin_left = Inches(0.15)
    tf.margin_right = Inches(0.15)
    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    r1 = p1.add_run()
    r1.text = value
    r1.font.size = Pt(22)
    r1.font.bold = True
    r1.font.color.rgb = KPMG_BLUE
    r1.font.name = "Arial"
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = label
    r2.font.size = Pt(10.5)
    r2.font.color.rgb = DARK
    r2.font.name = "Arial"
    return tile


def add_bullets(slide, x, y, w, h, bullets, size=12):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = "•  " + b
        run.font.size = Pt(size)
        run.font.color.rgb = DARK
        run.font.name = "Arial"
    return box


def add_table(slide, x, y, w, h, headers, rows, col_widths=None,
              row_colorizer=None, header_color=KPMG_BLUE, font_size=10.5):
    n_rows = 1 + len(rows)
    shape = slide.shapes.add_table(n_rows, len(headers), Inches(x), Inches(y), Inches(w), Inches(h))
    table = shape.table
    if col_widths:
        for i, cw in enumerate(col_widths):
            table.columns[i].width = Inches(cw)
    for c, htext in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = str(htext)
        run = cell.text_frame.paragraphs[0].runs[0]
        run.font.size = Pt(font_size + 1)
        run.font.bold = True
        run.font.color.rgb = WHITE
        run.font.name = "Arial"
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_color
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT
    for r_idx, row in enumerate(rows, start=1):
        for c_idx, val in enumerate(row):
            cell = table.cell(r_idx, c_idx)
            cell.text = str(val)
            run = cell.text_frame.paragraphs[0].runs[0]
            run.font.size = Pt(font_size)
            run.font.name = "Arial"
            run.font.color.rgb = DARK
            cell.fill.solid()
            cell.fill.fore_color.rgb = KPMG_LIGHT if (r_idx % 2 == 0) else WHITE
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER if c_idx > 0 else PP_ALIGN.LEFT
            if row_colorizer:
                override = row_colorizer(r_idx - 1, c_idx, val)
                if override:
                    run.font.color.rgb = override
                    run.font.bold = True
    return table


def new_slide(prs=None):
    if prs is None:
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    return prs, slide


def style_header_row(ws, row, ncols, fill="00338D"):
    for c in range(1, ncols + 1):
        cell = ws.cell(row=row, column=c)
        cell.font = Font(bold=True, color="FFFFFF")
        cell.fill = PatternFill("solid", fgColor=fill)
        cell.alignment = Alignment(horizontal="center" if c > 1 else "left")


def thin_border():
    thin = Side(style="thin", color="D7DEEA")
    return Border(left=thin, right=thin, top=thin, bottom=thin)


# ===========================================================================
# STAGE 1: DECK MATCHING ENGINE
# Adapted from Deck Refresh. Finds every number in the deck, matches it
# against new data with a keyword-conflict guard, and detects reporting
# period changes to rewrite headings.
# ===========================================================================

NUMBER_RE = re.compile(
    r"(?P<prefix>[$€£]?)\s*(?P<num>-?\d{1,3}(?:,\d{3})*(?:\.\d+)?|\-?\d+(?:\.\d+)?)"
    r"\s*(?P<mag>[KkMmBbTt])?\s*(?P<suffix>%?)"
)
PERIOD_WORDS = re.compile(
    r"\b(q1|q2|q3|q4|quarter\s*[1-4]|fy\s*\d{2,4}|20\d{2}|h1|h2|jan|feb|mar|apr|may|jun|jul|aug|sep|sept|oct|nov|dec|"
    r"january|february|march|april|june|july|august|september|october|november|december|ytd|mtd)\b",
    re.IGNORECASE,
)
QUARTER_TOKEN_RE = re.compile(r"\bQ[1-4]\b", re.IGNORECASE)
FY_TOKEN_RE = re.compile(r"\bFY\s?\d{2,4}\b", re.IGNORECASE)
YEAR_TOKEN_RE = re.compile(r"\b20\d{2}\b")
KEY_METRIC_WORDS = {
    "revenue", "cost", "costs", "expense", "expenses", "margin", "income",
    "profit", "budget", "actual", "variance", "ebitda", "ebit", "loss",
}


def normalize_label(text):
    t = (text or "").lower()
    t = re.sub(r"\([^)]*\)", " ", t)
    t = PERIOD_WORDS.sub(" ", t)
    t = re.sub(r"[^a-z0-9\s]", " ", t)
    t = re.sub(r"\s+", " ", t).strip()
    return t


def looks_numeric(text):
    text = (text or "").strip()
    if not text:
        return False
    m = NUMBER_RE.fullmatch(text)
    return m is not None and re.search(r"\d", text) is not None


def extract_number_value(text):
    if text is None:
        return None
    if isinstance(text, (int, float)):
        return float(text)
    m = NUMBER_RE.search(str(text))
    if not m:
        return None
    raw = m.group("num").replace(",", "")
    try:
        return float(raw)
    except ValueError:
        return None


def format_new_value(original_text, new_value):
    original_text = str(original_text)
    currency = ""
    for sym in ("$", "€", "£"):
        if sym in original_text:
            currency = sym
            break
    has_percent = "%" in original_text
    mag_match = re.search(r"\d\s*([KkMmBbTt])\b", original_text)
    magnitude = mag_match.group(1) if mag_match else ""
    has_comma = "," in original_text
    decimals = 0
    m = re.search(r"\.(\d+)", original_text)
    if m:
        decimals = len(m.group(1))
    is_negative = new_value < 0
    abs_val = abs(new_value)
    if decimals:
        num_str = f"{abs_val:,.{decimals}f}" if has_comma else f"{abs_val:.{decimals}f}"
    else:
        if abs_val == int(abs_val):
            num_str = f"{int(abs_val):,}" if has_comma else f"{int(abs_val)}"
        else:
            num_str = f"{abs_val:,.2f}" if has_comma else f"{abs_val:.2f}"
    sign = "-" if is_negative else ""
    result = f"{sign}{currency}{num_str}{magnitude}"
    if has_percent:
        result += "%"
    return result


def slide_title(slide):
    for shape in slide.shapes:
        if shape.has_text_frame and shape.is_placeholder:
            try:
                if shape.placeholder_format.type is not None and "TITLE" in str(shape.placeholder_format.type):
                    return shape.text_frame.text.strip()
            except Exception:
                pass
    for shape in slide.shapes:
        if shape.has_text_frame and shape.text_frame.text.strip():
            return shape.text_frame.text.strip().split("\n")[0]
    return ""


def extract_targets_pptx(pptx_path):
    prs = Presentation(pptx_path)
    targets = []
    for s_idx, slide in enumerate(prs.slides):
        title = slide_title(slide)
        for sh_idx, shape in enumerate(slide.shapes):
            if shape.has_text_frame:
                all_paras = shape.text_frame.paragraphs
                for p_idx, para in enumerate(all_paras):
                    for r_idx, run in enumerate(para.runs):
                        if looks_numeric(run.text):
                            other_text = "".join(
                                rr.text for k, rr in enumerate(para.runs) if k != r_idx
                            ).strip()
                            if not other_text:
                                sibling_text = " ".join(
                                    pp.text.strip() for pi, pp in enumerate(all_paras)
                                    if pi != p_idx and pp.text.strip()
                                ).strip()
                                other_text = sibling_text
                            label = other_text or title or shape.name
                            targets.append({
                                "id": f"s{s_idx}-sh{sh_idx}-p{p_idx}-r{r_idx}",
                                "kind": "run", "slide": s_idx + 1, "label": label,
                                "context": f"Slide {s_idx+1}: {title}" if title else f"Slide {s_idx+1}",
                                "original": run.text.strip(),
                                "location": (s_idx, sh_idx, "para", p_idx, r_idx),
                            })
            if shape.has_chart:
                chart = shape.chart
                chart_title = ""
                try:
                    if chart.has_title and chart.chart_title.text_frame.text.strip():
                        chart_title = chart.chart_title.text_frame.text.strip()
                except Exception:
                    pass
                try:
                    plot = chart.plots[0]
                    categories = [str(c) for c in plot.categories]
                    series_list = list(chart.series)
                    multi_series = len(series_list) > 1
                    for se_idx, series in enumerate(series_list):
                        values = list(series.values)
                        for c_idx, cat in enumerate(categories):
                            if c_idx >= len(values) or values[c_idx] is None:
                                continue
                            val = values[c_idx]
                            label = f"{series.name} {cat}".strip() if multi_series else cat
                            targets.append({
                                "id": f"s{s_idx}-sh{sh_idx}-chart-se{se_idx}-c{c_idx}",
                                "kind": "chart_point", "slide": s_idx + 1, "label": label,
                                "context": (f"Slide {s_idx+1}: {chart_title or title} (chart)"
                                            if (chart_title or title) else f"Slide {s_idx+1} (chart)"),
                                "original": (f"{val:g}" if isinstance(val, float) else str(val)),
                                "location": (s_idx, sh_idx, "chart", se_idx, c_idx),
                            })
                except Exception:
                    pass
            if shape.has_table:
                table = shape.table
                nrows, ncols = len(table.rows), len(table.columns)
                col_headers = [table.cell(0, c).text.strip() for c in range(ncols)]
                for r in range(1, nrows):
                    row_header = table.cell(r, 0).text.strip()
                    for c in range(1, ncols):
                        cell_text = table.cell(r, c).text.strip()
                        if looks_numeric(cell_text):
                            label = f"{row_header} {col_headers[c]}".strip()
                            targets.append({
                                "id": f"s{s_idx}-sh{sh_idx}-tbl-r{r}-c{c}",
                                "kind": "table_cell", "slide": s_idx + 1,
                                "label": label or f"{shape.name} row{r} col{c}",
                                "context": (f"Slide {s_idx+1}: {title} (table)" if title else f"Slide {s_idx+1} (table)"),
                                "original": cell_text,
                                "location": (s_idx, sh_idx, "table", r, c),
                            })
    return prs, targets


def extract_from_spreadsheet(path):
    pairs = []
    try:
        if path.lower().endswith(".csv"):
            frames = {"csv": pd.read_csv(path)}
        else:
            raw = pd.read_excel(path, sheet_name=None)
            frames = raw if isinstance(raw, dict) else {"Sheet1": raw}
    except Exception:
        return pairs
    for sheet_name, frame in frames.items():
        frame = frame.dropna(how="all").dropna(axis=1, how="all")
        if frame.empty:
            continue
        cols = [str(c) for c in frame.columns]
        first_col_is_label = (
            not pd.api.types.is_numeric_dtype(frame.iloc[:, 0]) if frame.shape[1] > 0 else False
        )
        if frame.shape[1] == 2 and first_col_is_label:
            for _, row in frame.iterrows():
                label = str(row.iloc[0]).strip()
                val = row.iloc[1]
                if pd.notna(val) and label and label.lower() != "nan":
                    try:
                        pairs.append((label, float(val)))
                    except (ValueError, TypeError):
                        pass
            continue
        for _, row in frame.iterrows():
            row_label = str(row.iloc[0]).strip() if first_col_is_label else ""
            start = 1 if first_col_is_label else 0
            data_col_count = len(cols) - start
            for c in range(start, len(cols)):
                val = row.iloc[c]
                if pd.isna(val):
                    continue
                try:
                    val = float(val)
                except (ValueError, TypeError):
                    continue
                col_label = cols[c]
                combined = f"{row_label} {col_label}".strip()
                if combined:
                    pairs.append((combined, val))
                if data_col_count == 1 and row_label:
                    pairs.append((row_label, val))
                if len(frame) == 1 and col_label and col_label.lower() not in ("unnamed: 0",):
                    pairs.append((col_label, val))
    return pairs


def match_targets_to_source(targets, source_pairs, threshold=60):
    norm_sources = [(normalize_label(lbl), lbl, val) for lbl, val in source_pairs]
    matches = []
    for t in targets:
        norm_t = normalize_label(t["label"])
        t_tokens = set(norm_t.split())
        t_key = t_tokens & KEY_METRIC_WORDS
        best, best_score, best_metric = None, 0, None
        top_seen_score = 0
        if norm_t:
            for norm_s, orig_s, val in norm_sources:
                if not norm_s:
                    continue
                s_tokens = set(norm_s.split())
                if not s_tokens:
                    continue
                score = fuzz.token_sort_ratio(norm_t, norm_s)
                top_seen_score = max(top_seen_score, score)
                s_key = s_tokens & KEY_METRIC_WORDS
                if t_key and not (t_key <= s_key):
                    continue
                sym_diff = len(t_tokens ^ s_tokens)
                actual_bonus = 0 if "actual" in s_tokens else 1
                metric = (sym_diff, actual_bonus, -score)
                if best_metric is None or metric < best_metric:
                    best_metric, best_score, best = metric, score, (orig_s, val)
        entry = dict(t)
        if best and best_score >= threshold:
            entry["matched_label"] = best[0]
            entry["new_value"] = best[1]
            entry["score"] = round(best_score, 1)
        else:
            entry["matched_label"] = None
            entry["new_value"] = None
            entry["score"] = round(top_seen_score, 1)
        if entry["new_value"] is not None:
            entry["new_text_preview"] = format_new_value(entry["original"], entry["new_value"])
        else:
            entry["new_text_preview"] = None
        matches.append(entry)
    return matches


def detect_period_tokens(texts):
    q_counter, fy_counter, year_counter = Counter(), Counter(), Counter()
    for t in texts:
        if not t:
            continue
        t = str(t)
        for m in FY_TOKEN_RE.finditer(t):
            fy_counter[re.sub(r"\s+", "", m.group(0)).upper()] += 1
        remainder = FY_TOKEN_RE.sub(" ", t)
        for m in QUARTER_TOKEN_RE.finditer(remainder):
            q_counter[m.group(0).upper()] += 1
        remainder2 = QUARTER_TOKEN_RE.sub(" ", remainder)
        for m in YEAR_TOKEN_RE.finditer(remainder2):
            year_counter[m.group(0)] += 1
    return q_counter, fy_counter, year_counter


def collect_all_pptx_text(prs):
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                texts.append(shape.text_frame.text)
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        texts.append(cell.text_frame.text)
            if shape.has_chart:
                try:
                    if shape.chart.has_title:
                        texts.append(shape.chart.chart_title.text_frame.text)
                except Exception:
                    pass
    return texts


def compute_period_replacements(deck_texts, source_pairs, pasted_text=""):
    deck_q, deck_fy, deck_year = detect_period_tokens(deck_texts)
    src_texts = [lbl for lbl, _ in source_pairs]
    if pasted_text:
        src_texts.append(pasted_text)
    src_q, src_fy, src_year = detect_period_tokens(src_texts)
    replacements = {}
    for deck_counter, src_counter in ((deck_q, src_q), (deck_fy, src_fy), (deck_year, src_year)):
        if not deck_counter or not src_counter:
            continue
        old = deck_counter.most_common(1)[0][0]
        new = src_counter.most_common(1)[0][0]
        if old.upper() != new.upper():
            replacements[old] = new
    return replacements


def _build_period_pattern(replacements):
    if not replacements:
        return None, None
    keys_sorted = sorted(replacements.keys(), key=len, reverse=True)
    pattern = re.compile(r"\b(" + "|".join(re.escape(k) for k in keys_sorted) + r")\b", re.IGNORECASE)
    lookup = {k.upper(): v for k, v in replacements.items()}
    return pattern, lookup


def _replace_text(text, pattern, lookup):
    if not text or not pattern:
        return text
    return pattern.sub(lambda m: lookup.get(m.group(0).upper(), m.group(0)), text)


def apply_edits_pptx(pptx_path, out_path, confirmed_matches, period_replacements=None):
    prs = Presentation(pptx_path)
    slides = list(prs.slides)
    chart_edits = {}
    pattern, lookup = _build_period_pattern(period_replacements or {})

    for m in confirmed_matches:
        loc = m["location"]
        s_idx, sh_idx = loc[0], loc[1]
        shape = list(slides[s_idx].shapes)[sh_idx]
        if loc[2] == "para":
            p_idx, r_idx = loc[3], loc[4]
            shape.text_frame.paragraphs[p_idx].runs[r_idx].text = m["new_text"]
        elif loc[2] == "table":
            r, c = loc[3], loc[4]
            cell = shape.table.cell(r, c)
            if cell.text_frame.paragraphs and cell.text_frame.paragraphs[0].runs:
                cell.text_frame.paragraphs[0].runs[0].text = m["new_text"]
                for extra in cell.text_frame.paragraphs[0].runs[1:]:
                    extra.text = ""
            else:
                cell.text = m["new_text"]
        elif loc[2] == "chart":
            se_idx, c_idx = loc[3], loc[4]
            numeric_val = extract_number_value(m["new_text"])
            if numeric_val is None:
                continue
            chart_edits.setdefault((s_idx, sh_idx), {})[(se_idx, c_idx)] = numeric_val

    charts_needing_period_update = set()
    if pattern:
        for slide_idx, slide in enumerate(slides):
            for sh_idx, shape in enumerate(slide.shapes):
                if not shape.has_chart:
                    continue
                chart = shape.chart
                try:
                    plot = chart.plots[0]
                    names_and_cats = [str(c) for c in plot.categories] + [s.name or "" for s in chart.series]
                    if any(pattern.search(t) for t in names_and_cats):
                        charts_needing_period_update.add((slide_idx, sh_idx))
                except Exception:
                    pass

    for (s_idx, sh_idx) in set(chart_edits.keys()) | charts_needing_period_update:
        shape = list(slides[s_idx].shapes)[sh_idx]
        chart = shape.chart
        plot = chart.plots[0]
        categories = [_replace_text(str(c), pattern, lookup) for c in plot.categories]
        new_chart_data = CategoryChartData()
        new_chart_data.categories = categories
        edits = chart_edits.get((s_idx, sh_idx), {})
        for se_idx, series in enumerate(chart.series):
            values = list(series.values)
            for c_idx in range(len(values)):
                if (se_idx, c_idx) in edits:
                    values[c_idx] = edits[(se_idx, c_idx)]
            new_name = _replace_text(series.name or "", pattern, lookup)
            new_chart_data.add_series(new_name, values)
        chart.replace_data(new_chart_data)

    if pattern:
        def process_text_frame(tf):
            for para in tf.paragraphs:
                for run in para.runs:
                    if run.text and pattern.search(run.text):
                        run.text = _replace_text(run.text, pattern, lookup)
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    process_text_frame(shape.text_frame)
                if shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            process_text_frame(cell.text_frame)
                if shape.has_chart:
                    try:
                        if shape.chart.has_title:
                            process_text_frame(shape.chart.chart_title.text_frame)
                    except Exception:
                        pass

    prs.save(out_path)


# ===========================================================================
# STAGE 2: VARIANCE ANALYSIS
# Derived directly from the deck-matching results (Stage 1) -- the same
# old/new number pairs that were just used to update the deck are reused
# here for variance commentary, so no separate upload is needed.
# ===========================================================================

COST_LIKE_WORDS = ["cost", "expense", "tax", "interest", "loss", "spend", "churn",
                    "attrition", "days sales", "turnover", "leverage", "debt"]
REVENUE_LIKE_WORDS = ["revenue", "income", "profit", "margin", "contribution", "gross",
                       "ebitda", "ebit", "earnings", "growth", "utilization", "rate", "retention"]


def classify_direction(label):
    l = label.lower()
    if any(w in l for w in REVENUE_LIKE_WORDS):
        return "higher_is_good"
    if any(w in l for w in COST_LIKE_WORDS):
        return "lower_is_good"
    return "neutral"


def build_variance_rows(matches, comparison_label, materiality_pct):
    rows = []
    for m in matches:
        if m["matched_label"] is None or m["new_value"] is None:
            continue
        old_val = extract_number_value(m["original"])
        if old_val is None:
            continue
        new_val = m["new_value"]
        variance = new_val - old_val
        variance_pct = (variance / old_val * 100) if old_val != 0 else None
        is_material = variance_pct is not None and abs(variance_pct) >= materiality_pct

        direction = classify_direction(m["label"])
        if direction == "higher_is_good":
            favorable = variance > 0
        elif direction == "lower_is_good":
            favorable = variance < 0
        else:
            favorable = None
        tag = "Favorable" if favorable is True else ("Unfavorable" if favorable is False else
                                                       ("Notable" if is_material else "Neutral"))

        rows.append({
            "label": m["label"], "context": m["context"], "old": old_val, "new": new_val,
            "old_display": m["original"], "new_display": m["new_text_preview"],
            "variance": variance, "variance_pct": variance_pct, "is_material": is_material,
            "tag": tag,
        })
    return dedupe_variance_rows(rows)


def dedupe_variance_rows(rows):
    """The same underlying metric can legitimately appear twice in a deck
    (once in a table, once in a chart), producing two rows with identical
    values but differently-ordered labels (e.g. "Digital Campaigns New
    Engagements" vs "New Engagements Digital Campaigns"). Collapse those
    into one row using the value pair as the signature, so every
    downstream output -- the workbook, the package table, the commentary --
    doesn't repeat the same insight twice under a reshuffled label."""
    seen = set()
    deduped = []
    for r in rows:
        signature = (round(r["old"], 2), round(r["new"], 2))
        if signature in seen:
            continue
        seen.add(signature)
        deduped.append(r)
    return deduped


def draft_variance_commentary(rows, comparison_label):
    material = sorted([r for r in rows if r["is_material"]],
                       key=lambda r: abs(r["variance_pct"] or 0), reverse=True)
    commentary = []
    for r in material[:8]:
        direction_word = "up" if r["variance"] > 0 else "down"
        pct_text = f"{abs(r['variance_pct']):.1f}%" if r["variance_pct"] is not None else "n/a"
        base = f"{r['label']} came in {direction_word} {pct_text} versus {comparison_label}"
        if r["tag"] == "Favorable":
            commentary.append(base + ", a favorable variance. [Add driver commentary]")
        elif r["tag"] == "Unfavorable":
            commentary.append(base + ", an unfavorable variance that warrants explanation. [Add driver commentary]")
        else:
            commentary.append(base + ". [Add context]")
    return commentary


# ===========================================================================
# STAGE 3: BENCHMARKING (optional, requires peer data upload)
# Ported from Engagement Hub.
# ===========================================================================

def classify_metric_direction(label):
    l = label.lower()
    if any(w in l for w in REVENUE_LIKE_WORDS):
        return "higher_is_good"
    if any(w in l for w in COST_LIKE_WORDS):
        return "lower_is_good"
    return "higher_is_good"


def percentile_rank(target, peers, higher_is_good=True):
    if not peers:
        return None
    if higher_is_good:
        beaten = sum(1 for p in peers if target >= p)
    else:
        beaten = sum(1 for p in peers if target <= p)
    return beaten / len(peers) * 100


def quartile_label(pct):
    if pct is None:
        return "n/a"
    if pct >= 75:
        return "Top Quartile"
    if pct >= 50:
        return "Above Median"
    if pct >= 25:
        return "Below Median"
    return "Bottom Quartile"


QUARTILE_CSS = {"Top Quartile": "top", "Above Median": "above",
                 "Below Median": "below", "Bottom Quartile": "bottom", "n/a": "neutral"}


def analyze_benchmark(df, metric_col, target_col, peer_cols):
    rows = []
    for _, row in df.iterrows():
        metric = str(row[metric_col]).strip()
        if not metric or pd.isna(row[target_col]):
            continue
        try:
            target_val = float(row[target_col])
        except (ValueError, TypeError):
            continue
        peers = []
        for pc in peer_cols:
            v = row.get(pc)
            if pd.notna(v):
                try:
                    peers.append(float(v))
                except (ValueError, TypeError):
                    pass
        if not peers:
            continue
        higher_is_good = classify_metric_direction(metric) == "higher_is_good"
        pct = percentile_rank(target_val, peers, higher_is_good)
        rows.append({
            "metric": metric, "target": target_val, "peers": peers,
            "peer_median": median(peers), "peer_min": min(peers), "peer_max": max(peers),
            "percentile": pct, "quartile": quartile_label(pct),
            "quartile_class": QUARTILE_CSS[quartile_label(pct)], "higher_is_good": higher_is_good,
        })
    return rows


# ===========================================================================
# STAGE 4: ACTION TRACKER (optional, requires meeting notes)
# Ported from Engagement Hub.
# ===========================================================================

BULLET_RE = re.compile(r'^\s*[-*•]\s+|^\s*\d+[\.\)]\s+')
ACTION_RE = re.compile(
    r'\b(will (?:send|share|review|confirm|finalize|circulate|draft|prepare|follow up|'
    r'schedule|update|provide|reach out|complete|deliver|escalate|coordinate)|'
    r'needs to|action item|next step|follow[- ]?up|to-?do)\b',
    re.IGNORECASE)
URGENT_RE = re.compile(r'\b(urgent|asap|critical|immediately|high priority)\b', re.IGNORECASE)
DATE_RE = re.compile(
    r'\bby\s+(?:next\s+)?(?:monday|tuesday|wednesday|thursday|friday|saturday|sunday|'
    r'eod|cob|end of day|end of week|\d{1,2}/\d{1,2}(?:/\d{2,4})?|[A-Z][a-z]+\s+\d{1,2}(?:st|nd|rd|th)?)\b',
    re.IGNORECASE)


def extract_owner(line, attendees):
    for a in attendees:
        if a and re.search(r'\b' + re.escape(a) + r'\b', line, re.IGNORECASE):
            return a
    return "Unassigned"


def extract_due(line):
    m = DATE_RE.search(line)
    return m.group(0) if m else "Not specified"


def parse_meeting_notes(text, attendees):
    lines = [l.strip() for l in text.splitlines() if l.strip()]
    items = []
    for line in lines:
        clean = BULLET_RE.sub("", line).strip()
        if len(clean) < 8:
            continue
        is_bullet = bool(BULLET_RE.match(line))
        actionable = bool(ACTION_RE.search(clean))
        if not is_bullet and not actionable:
            continue
        items.append({
            "text": clean, "owner": extract_owner(clean, attendees), "due": extract_due(clean),
            "priority": "High" if URGENT_RE.search(clean) else "Medium", "include": actionable,
        })
    return items


# ===========================================================================
# EXECUTIVE PACKAGE ASSEMBLY
# Builds one cohesive multi-slide deck from whatever stages ran: cover,
# executive summary (stats from every stage that ran), variance detail,
# benchmarking (if peer data was provided), and action items (if meeting
# notes were provided). Meant to be presented alongside the updated deck,
# not merged into it -- merging would mean abandoning Deck Refresh's core
# guarantee of exact formatting preservation on the original file.
# ===========================================================================

def build_executive_package(out_path, config, stats, variance_rows, variance_commentary,
                             benchmark_rows=None, action_items=None):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Cover ---
    prs, cover = new_slide(prs)
    bar = cover.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = KPMG_BLUE
    bar.line.fill.background()
    add_textbox(cover, 0.75, 0.55, 4, 0.5, "KPMG", size=22, bold=True, color=WHITE)
    add_textbox(cover, 0.75, 2.6, 10.5, 1.1, config["title"], size=34, bold=True, color=WHITE)
    if config.get("subtitle"):
        add_textbox(cover, 0.75, 3.6, 10.5, 0.6, config["subtitle"], size=15, color=RGBColor(0xCF, 0xE3, 0xF7))
    add_textbox(cover, 0.75, 6.9, 8, 0.35, "Quarter close package — generated automatically, reviewed before finalizing",
                size=10, italic=True, color=RGBColor(0xCF, 0xE3, 0xF7))

    # --- Slide 2: Executive Summary ---
    prs, summary_slide = new_slide(prs)
    add_header_bar(summary_slide, prs, "Executive Summary", config.get("subtitle"), tag_text="AUTOMATED PACKAGE")

    kpis = [
        (str(stats["matched_count"]), "Figures updated"),
        (str(stats["material_count"]), "Material variances"),
        (f"{stats['favorable_count']} / {stats['unfavorable_count']}", "Favorable / Unfavorable"),
    ]
    if benchmark_rows is not None:
        top_q = sum(1 for r in benchmark_rows if r["quartile"] == "Top Quartile")
        kpis.append((f"{top_q} / {len(benchmark_rows)}", "Metrics in top quartile"))
    if action_items is not None:
        kpis.append((str(len(action_items)), "Action items identified"))

    tile_w = min(3.9, (12.5 - 0.3 * (len(kpis) - 1)) / len(kpis))
    for i, (value, label) in enumerate(kpis):
        add_kpi_tile(summary_slide, 0.4 + i * (tile_w + 0.3), 1.5, tile_w, 1.4, value, label)

    add_textbox(summary_slide, 0.4, 3.2, 6, 0.35, "Top Highlights", size=14, bold=True, color=KPMG_BLUE)
    top_bullets = variance_commentary[:5]
    add_bullets(summary_slide, 0.4, 3.6, 12.5, max(1.5, 0.5 * len(top_bullets)), top_bullets, size=12)

    # --- Slide 3: Variance Detail ---
    prs, var_slide = new_slide(prs)
    add_header_bar(var_slide, prs, "Variance Detail", f"Material variances vs. {config['comparison_label']}",
                    tag_text="AUTOMATED PACKAGE")
    material_rows = sorted([r for r in variance_rows if r["is_material"]],
                            key=lambda r: abs(r["variance_pct"] or 0), reverse=True)[:12]
    table_rows = [[r["label"], r["old_display"], r["new_display"],
                   f"{r['variance_pct']:+.1f}%" if r["variance_pct"] is not None else "n/a", r["tag"]]
                  for r in material_rows]

    def colorize_var(row_idx, col_idx, val):
        if col_idx in (3, 4):
            if val == "Favorable" or (isinstance(val, str) and val.startswith("+") and material_rows[row_idx]["tag"] == "Favorable"):
                return GOOD
            if val == "Unfavorable" or (isinstance(val, str) and material_rows[row_idx]["tag"] == "Unfavorable"):
                return BAD
        return None

    table_h = min(5.5, 0.4 * (len(table_rows) + 1))
    add_table(var_slide, 0.4, 1.5, 12.5, table_h,
              ["Line Item", config["comparison_label"], "New", "Change", "Positioning"], table_rows,
              col_widths=[5.2, 2.3, 2.3, 1.6, 1.1], row_colorizer=colorize_var)

    # --- Slide 4: Benchmarking (optional) ---
    if benchmark_rows:
        prs, bench_slide = new_slide(prs)
        add_header_bar(bench_slide, prs, "Benchmarking", f"{config.get('target_name', 'Target')} vs. peer set",
                        tag_text="AUTOMATED PACKAGE")
        b_table_rows = []
        for r in benchmark_rows:
            b_table_rows.append([r["metric"], f"{r['target']:,.1f}", f"{r['peer_median']:,.1f}",
                                  f"{r['percentile']:.0f}th pct" if r["percentile"] is not None else "n/a",
                                  r["quartile"]])

        def colorize_bench(row_idx, col_idx, val):
            if col_idx == 4:
                if val == "Top Quartile":
                    return GOOD
                if val == "Bottom Quartile":
                    return BAD
            return None

        b_table_h = min(3.3, 0.4 * (len(b_table_rows) + 1))
        add_table(bench_slide, 0.4, 1.5, 12.5, b_table_h,
                  ["Metric", config.get("target_name", "Target"), "Peer Median", "Percentile", "Positioning"],
                  b_table_rows, col_widths=[3.7, 2.2, 2.4, 2.2, 2.0], row_colorizer=colorize_bench)

        chart_y = 1.5 + b_table_h + 0.3
        chart_data = CategoryChartData()
        chart_data.categories = [r["metric"] for r in benchmark_rows]
        chart_data.add_series("Percentile", tuple(r["percentile"] or 0 for r in benchmark_rows))
        gf = bench_slide.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, Inches(0.4), Inches(chart_y),
                                           Inches(12.5), Inches(prs.slide_height / 914400 - chart_y - 0.3), chart_data)
        chart = gf.chart
        chart.has_legend = False
        for i, pt in enumerate(chart.series[0].points):
            pt.format.fill.solid()
            pt.format.fill.fore_color.rgb = PALETTE[i % len(PALETTE)]
        chart.value_axis.minimum_scale = 0
        chart.value_axis.maximum_scale = 100
        plot = chart.plots[0]
        plot.has_data_labels = True
        plot.data_labels.number_format = '0"th"'
        plot.data_labels.number_format_is_linked = False
        plot.data_labels.font.size = Pt(9)

    # --- Slide 5: Action Items (optional) ---
    if action_items:
        prs, action_slide = new_slide(prs)
        add_header_bar(action_slide, prs, "Action Items & Next Steps", config.get("meeting_title", ""),
                        tag_text="AUTOMATED PACKAGE")
        by_owner = {}
        for item in action_items:
            by_owner.setdefault(item["owner"], []).append(item)
        a_rows = []
        for owner in sorted(by_owner.keys(), key=lambda o: (o == "Unassigned", o)):
            for item in by_owner[owner]:
                a_rows.append([item["text"], owner, item["due"], item["priority"]])

        def colorize_action(row_idx, col_idx, val):
            if col_idx == 3 and val == "High":
                return BAD
            return None

        a_table_h = min(5.5, 0.4 * (len(a_rows) + 1))
        add_table(action_slide, 0.4, 1.5, 12.5, a_table_h,
                  ["Action Item", "Owner", "Due Date", "Priority"], a_rows,
                  col_widths=[6.5, 2.3, 2.0, 1.7], row_colorizer=colorize_action)

    prs.save(out_path)


# ===========================================================================
# SUPPORTING WORKBOOK / EMAIL EXPORTS
# ===========================================================================

def build_variance_xlsx(variance_rows, commentary, title, comparison_label, out_path):
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Variance Analysis"
    ws.merge_cells("A1:E1")
    ws["A1"] = title
    ws["A1"].font = Font(bold=True, size=14, color="00338D")

    headers = ["Line Item", comparison_label, "New Value", "Variance %", "Positioning"]
    header_row = 3
    style_header_row(ws, header_row, len(headers))
    for c, h in enumerate(headers, start=1):
        ws.cell(row=header_row, column=c, value=h)

    border = thin_border()
    sorted_rows = sorted(variance_rows, key=lambda r: abs(r["variance_pct"] or 0), reverse=True)
    for i, r in enumerate(sorted_rows, start=1):
        row_n = header_row + i
        vals = [r["label"], r["old_display"], r["new_display"],
                f"{r['variance_pct']:+.1f}%" if r["variance_pct"] is not None else "n/a", r["tag"]]
        for c, v in enumerate(vals, start=1):
            cell = ws.cell(row=row_n, column=c, value=v)
            cell.border = border
            cell.alignment = Alignment(horizontal="center" if c > 1 else "left")
            if c == 5:
                if v == "Favorable":
                    cell.font = Font(color=GOOD_HEX, bold=True)
                elif v == "Unfavorable":
                    cell.font = Font(color=BAD_HEX, bold=True)
        if i % 2 == 0:
            for c in range(1, 6):
                ws.cell(row=row_n, column=c).fill = PatternFill("solid", fgColor="F2F5FA")
    ws.column_dimensions["A"].width = 40
    for col in "BCDE":
        ws.column_dimensions[col].width = 16

    ws2 = wb.create_sheet("Commentary")
    ws2["A1"] = "Key Variance Commentary"
    ws2["A1"].font = Font(bold=True, size=13, color="00338D")
    for i, c in enumerate(commentary):
        ws2.cell(row=3 + i, column=1, value=c).alignment = Alignment(wrap_text=True)
    ws2.column_dimensions["A"].width = 100
    wb.save(out_path)


def build_benchmark_xlsx(rows, title, target_name, out_path):
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Benchmarking"
    ws.merge_cells("A1:F1")
    ws["A1"] = title
    ws["A1"].font = Font(bold=True, size=14, color="00338D")
    headers = ["Metric", target_name, "Peer Median", "Peer Min", "Peer Max", "Percentile", "Positioning"]
    header_row = 3
    style_header_row(ws, header_row, len(headers))
    for c, h in enumerate(headers, start=1):
        ws.cell(row=header_row, column=c, value=h)
    border = thin_border()
    for i, r in enumerate(rows, start=1):
        row_n = header_row + i
        vals = [r["metric"], r["target"], r["peer_median"], r["peer_min"], r["peer_max"],
                f"{r['percentile']:.0f}" if r["percentile"] is not None else "n/a", r["quartile"]]
        for c, v in enumerate(vals, start=1):
            cell = ws.cell(row=row_n, column=c, value=v)
            cell.border = border
            cell.alignment = Alignment(horizontal="center" if c > 1 else "left")
            if c == 7:
                if v == "Top Quartile":
                    cell.font = Font(color=GOOD_HEX, bold=True)
                elif v == "Bottom Quartile":
                    cell.font = Font(color=BAD_HEX, bold=True)
        if i % 2 == 0:
            for c in range(1, 8):
                ws.cell(row=row_n, column=c).fill = PatternFill("solid", fgColor="F2F5FA")
    ws.column_dimensions["A"].width = 26
    for col in "BCDEFG":
        ws.column_dimensions[col].width = 15
    wb.save(out_path)


def build_action_tracker_xlsx(items, meeting_title, out_path):
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Action Tracker"
    ws.merge_cells("A1:E1")
    ws["A1"] = f"{meeting_title} — Action Item Tracker"
    ws["A1"].font = Font(bold=True, size=14, color="00338D")
    headers = ["Action Item", "Owner", "Due Date", "Priority", "Status"]
    header_row = 3
    style_header_row(ws, header_row, len(headers))
    for c, h in enumerate(headers, start=1):
        ws.cell(row=header_row, column=c, value=h)
    border = thin_border()
    for i, item in enumerate(items, start=1):
        row_n = header_row + i
        vals = [item["text"], item["owner"], item["due"], item["priority"], "Open"]
        for c, v in enumerate(vals, start=1):
            cell = ws.cell(row=row_n, column=c, value=v)
            cell.border = border
            cell.alignment = Alignment(wrap_text=(c == 1), vertical="top")
            if c == 4 and v == "High":
                cell.font = Font(color=BAD_HEX, bold=True)
        if i % 2 == 0:
            for c in range(1, 6):
                ws.cell(row=row_n, column=c).fill = PatternFill("solid", fgColor="F2F5FA")
    ws.column_dimensions["A"].width = 55
    ws.column_dimensions["B"].width = 18
    ws.column_dimensions["C"].width = 16
    ws.column_dimensions["D"].width = 12
    ws.column_dimensions["E"].width = 12
    wb.save(out_path)


def build_followup_email(items, meeting_title, sender_name, stats):
    by_owner = {}
    for item in items:
        by_owner.setdefault(item["owner"], []).append(item)
    lines = []
    lines.append(f"Subject: Quarter close package and action items — {meeting_title}")
    lines.append("")
    lines.append("Hi all,")
    lines.append("")
    lines.append(f"Attached is the quarter close package: the updated deck, a variance analysis "
                  f"({stats['material_count']} material variances flagged), and the action items below.")
    lines.append("")
    for owner in sorted(by_owner.keys(), key=lambda o: (o == "Unassigned", o)):
        lines.append(f"{owner}:")
        for item in by_owner[owner]:
            due_suffix = f" ({item['due']})" if item["due"] != "Not specified" else ""
            lines.append(f"  - {item['text']}{due_suffix}")
        lines.append("")
    lines.append("Please flag if anything above needs correction.")
    lines.append("")
    lines.append("Best,")
    lines.append(sender_name or "[Your name]")
    return "\n".join(lines)


# ===========================================================================
# REAL SLIDE PREVIEW (LibreOffice + PyMuPDF)
# ===========================================================================

SOFFICE_PATH = shutil.which("soffice") or shutil.which("libreoffice")
_soffice_lock = threading.Lock()
try:
    import fitz
    HAVE_FITZ = True
except ImportError:
    HAVE_FITZ = False


def render_pptx_to_images(pptx_path, out_dir, prefix):
    if not (SOFFICE_PATH and HAVE_FITZ):
        return None
    os.makedirs(out_dir, exist_ok=True)
    profile_dir = tempfile.mkdtemp(prefix="lo_profile_")
    try:
        with _soffice_lock:
            result = subprocess.run(
                [SOFFICE_PATH, "--headless", "--norestore",
                 f"-env:UserInstallation=file://{profile_dir}",
                 "--convert-to", "pdf", "--outdir", out_dir, pptx_path],
                capture_output=True, timeout=120,
            )
        pdf_path = os.path.join(out_dir, os.path.splitext(os.path.basename(pptx_path))[0] + ".pdf")
        if result.returncode != 0 or not os.path.exists(pdf_path):
            return None
        doc = fitz.open(pdf_path)
        paths = []
        for i, page in enumerate(doc):
            pix = page.get_pixmap(matrix=fitz.Matrix(1.8, 1.8))
            img_path = os.path.join(out_dir, f"{prefix}_{i+1}.png")
            pix.save(img_path)
            paths.append(img_path)
        doc.close()
        os.remove(pdf_path)
        return paths if paths else None
    except Exception:
        return None
    finally:
        shutil.rmtree(profile_dir, ignore_errors=True)


def _session_dir(sid):
    d = os.path.join(SESSIONS_DIR, sid)
    if not os.path.isdir(d):
        abort(404)
    return d


def _save_json(sess_dir, name, data):
    with open(os.path.join(sess_dir, name), "w") as f:
        json.dump(data, f)


def _load_json(sess_dir, name, default=None):
    path = os.path.join(sess_dir, name)
    if not os.path.exists(path):
        return default
    with open(path) as f:
        return json.load(f)


def _new_session():
    sid = uuid.uuid4().hex[:12]
    sess_dir = os.path.join(SESSIONS_DIR, sid)
    os.makedirs(sess_dir, exist_ok=True)
    return sid, sess_dir


# ===========================================================================
# ROUTES
# ===========================================================================

@app.route("/", methods=["GET"])
def index():
    return render_template("index.html")


@app.route("/process", methods=["POST"])
def process():
    deck_file = request.files.get("deck_file")
    data_file = request.files.get("data_file")
    if not deck_file or deck_file.filename == "" or not data_file or data_file.filename == "":
        flash("Please upload both the prior-quarter deck and the new data file.")
        return redirect(url_for("index"))

    title = request.form.get("title", "").strip() or "Quarter Close Package"
    subtitle = request.form.get("subtitle", "").strip()
    comparison_label = request.form.get("comparison_label", "").strip() or "Prior Period"
    target_name = request.form.get("target_name", "").strip() or "Client"
    meeting_title = request.form.get("meeting_title", "").strip() or "Quarter Close Call"
    sender_name = request.form.get("sender_name", "").strip()
    attendees_raw = request.form.get("attendees", "").strip()
    attendees = [a.strip() for a in attendees_raw.split(",") if a.strip()]
    try:
        materiality_pct = float(request.form.get("materiality_pct", "5") or 5)
    except ValueError:
        materiality_pct = 5.0

    sid, sess_dir = _new_session()
    deck_path = os.path.join(sess_dir, "original.pptx")
    deck_file.save(deck_path)
    data_path = os.path.join(sess_dir, "data_" + data_file.filename)
    data_file.save(data_path)

    # --- Stage 1: deck matching ---
    prs, targets = extract_targets_pptx(deck_path)
    pairs = extract_from_spreadsheet(data_path)
    if not pairs:
        flash("Couldn't read usable data from the new data file.")
        shutil.rmtree(sess_dir, ignore_errors=True)
        return redirect(url_for("index"))
    matches = match_targets_to_source(targets, pairs)
    deck_texts = collect_all_pptx_text(prs)
    period_replacements = compute_period_replacements(deck_texts, pairs)

    # --- Stage 2: variance analysis (derived from Stage 1) ---
    variance_rows = build_variance_rows(matches, comparison_label, materiality_pct)
    variance_commentary = draft_variance_commentary(variance_rows, comparison_label)

    # --- Stage 3: benchmarking (optional) ---
    benchmark_rows = None
    peer_file = request.files.get("peer_file")
    if peer_file and peer_file.filename:
        peer_path = os.path.join(sess_dir, "peers_" + peer_file.filename)
        peer_file.save(peer_path)
        try:
            bdf = pd.read_csv(peer_path) if peer_path.lower().endswith(".csv") else pd.read_excel(peer_path)
            bdf = bdf.dropna(how="all").dropna(axis=1, how="all")
            cols = list(bdf.columns)
            numeric_cols = [c for c in cols if pd.api.types.is_numeric_dtype(bdf[c])]
            non_numeric_cols = [c for c in cols if c not in numeric_cols]
            if len(numeric_cols) >= 2 and non_numeric_cols:
                benchmark_rows = analyze_benchmark(bdf, non_numeric_cols[0], numeric_cols[0], numeric_cols[1:])
        except Exception:
            benchmark_rows = None

    # --- Stage 4: action tracker (optional) ---
    action_items = None
    notes_text = request.form.get("notes_text", "").strip()
    notes_file = request.files.get("notes_file")
    if notes_file and notes_file.filename:
        try:
            notes_text = (notes_text + "\n" + notes_file.read().decode("utf-8", errors="ignore")).strip()
        except Exception:
            pass
    if notes_text:
        action_items = parse_meeting_notes(notes_text, attendees)

    _save_json(sess_dir, "matches.json", matches)
    _save_json(sess_dir, "period_replacements.json", period_replacements)
    _save_json(sess_dir, "variance_rows.json", variance_rows)
    _save_json(sess_dir, "variance_commentary.json", variance_commentary)
    if benchmark_rows is not None:
        _save_json(sess_dir, "benchmark_rows.json", benchmark_rows)
    if action_items is not None:
        _save_json(sess_dir, "action_items.json", action_items)
    _save_json(sess_dir, "config.json", {
        "title": title, "subtitle": subtitle, "comparison_label": comparison_label,
        "target_name": target_name, "meeting_title": meeting_title, "sender_name": sender_name,
        "materiality_pct": materiality_pct, "has_benchmark": benchmark_rows is not None,
        "has_actions": action_items is not None,
    })
    return redirect(url_for("review", sid=sid))


@app.route("/review/<sid>", methods=["GET"])
def review(sid):
    sess_dir = _session_dir(sid)
    matches = _load_json(sess_dir, "matches.json")
    config = _load_json(sess_dir, "config.json")
    if matches is None:
        flash("That session has expired. Please start again.")
        return redirect(url_for("index"))

    matched = [m for m in matches if m["matched_label"]]
    unmatched = [m for m in matches if not m["matched_label"]]
    period_replacements = _load_json(sess_dir, "period_replacements.json", {})
    variance_rows = _load_json(sess_dir, "variance_rows.json", [])
    variance_commentary = _load_json(sess_dir, "variance_commentary.json", [])
    material_count = sum(1 for r in variance_rows if r["is_material"])
    favorable_count = sum(1 for r in variance_rows if r["tag"] == "Favorable")
    unfavorable_count = sum(1 for r in variance_rows if r["tag"] == "Unfavorable")
    benchmark_rows = _load_json(sess_dir, "benchmark_rows.json", None)
    action_items = _load_json(sess_dir, "action_items.json", None)

    return render_template(
        "review.html", sid=sid, config=config,
        matched=matched, unmatched=unmatched, period_replacements=period_replacements,
        variance_commentary=variance_commentary, material_count=material_count,
        favorable_count=favorable_count, unfavorable_count=unfavorable_count,
        benchmark_rows=benchmark_rows, action_items=action_items,
    )


@app.route("/generate/<sid>", methods=["POST"])
def generate(sid):
    sess_dir = _session_dir(sid)
    matches = _load_json(sess_dir, "matches.json")
    config = _load_json(sess_dir, "config.json")
    period_replacements = _load_json(sess_dir, "period_replacements.json", {})
    variance_rows = _load_json(sess_dir, "variance_rows.json", [])
    benchmark_rows = _load_json(sess_dir, "benchmark_rows.json", None)
    action_items = _load_json(sess_dir, "action_items.json", None)
    if matches is None:
        flash("That session has expired. Please start again.")
        return redirect(url_for("index"))

    # Confirmed commentary (editable/excludable on the review page)
    saved_commentary = _load_json(sess_dir, "variance_commentary.json", [])
    confirmed_commentary = []
    for i in range(len(saved_commentary)):
        if request.form.get(f"include_commentary_{i}") == "on":
            confirmed_commentary.append(request.form.get(f"commentary_{i}", saved_commentary[i]).strip())

    if benchmark_rows is not None:
        benchmark_rows = [r for i, r in enumerate(benchmark_rows)
                           if request.form.get(f"include_bench_{i}") == "on"] or benchmark_rows

    if action_items is not None:
        confirmed_actions = []
        for i, item in enumerate(action_items):
            if request.form.get(f"include_action_{i}") == "on":
                entry = dict(item)
                entry["text"] = request.form.get(f"action_text_{i}", item["text"]).strip() or item["text"]
                entry["owner"] = request.form.get(f"action_owner_{i}", item["owner"]).strip() or "Unassigned"
                entry["due"] = request.form.get(f"action_due_{i}", item["due"]).strip() or "Not specified"
                confirmed_actions.append(entry)
        action_items = confirmed_actions

    matched = [dict(m, new_text=m["new_text_preview"]) for m in matches if m["matched_label"]]
    material_count = sum(1 for r in variance_rows if r["is_material"])
    favorable_count = sum(1 for r in variance_rows if r["tag"] == "Favorable")
    unfavorable_count = sum(1 for r in variance_rows if r["tag"] == "Unfavorable")
    stats = {"matched_count": len(matched), "material_count": material_count,
             "favorable_count": favorable_count, "unfavorable_count": unfavorable_count}

    deck_path = os.path.join(sess_dir, "original.pptx")
    updated_deck_path = os.path.join(sess_dir, "deck_updated.pptx")
    apply_edits_pptx(deck_path, updated_deck_path, matched, period_replacements)

    package_path = os.path.join(sess_dir, "executive_package.pptx")
    build_executive_package(package_path, config, stats, variance_rows, confirmed_commentary,
                             benchmark_rows, action_items)

    build_variance_xlsx(variance_rows, confirmed_commentary, config["title"],
                         config["comparison_label"], os.path.join(sess_dir, "variance_analysis.xlsx"))
    if benchmark_rows:
        build_benchmark_xlsx(benchmark_rows, config["title"], config["target_name"],
                              os.path.join(sess_dir, "benchmark_analysis.xlsx"))
    if action_items:
        build_action_tracker_xlsx(action_items, config["meeting_title"],
                                   os.path.join(sess_dir, "action_tracker.xlsx"))
        email_text = build_followup_email(action_items, config["meeting_title"], config["sender_name"], stats)
        with open(os.path.join(sess_dir, "followup_email.txt"), "w") as f:
            f.write(email_text)

    render_dir = os.path.join(sess_dir, "render")
    package_images = render_pptx_to_images(package_path, render_dir, "package")
    rendering_ok = package_images is not None

    zip_path = os.path.join(sess_dir, "quarter_close_package.zip")
    with zipfile.ZipFile(zip_path, "w", zipfile.ZIP_DEFLATED) as zf:
        for fname in ["deck_updated.pptx", "executive_package.pptx", "variance_analysis.xlsx",
                      "benchmark_analysis.xlsx", "action_tracker.xlsx", "followup_email.txt"]:
            fpath = os.path.join(sess_dir, fname)
            if os.path.exists(fpath):
                zf.write(fpath, fname)

    _save_json(sess_dir, "meta.json", {
        "rendering_ok": rendering_ok, "slide_count": len(package_images) if package_images else 0,
        "matched_count": len(matched), "material_count": material_count,
        "favorable_count": favorable_count, "unfavorable_count": unfavorable_count,
        "has_benchmark": bool(benchmark_rows), "has_actions": bool(action_items),
        "benchmark_top_count": sum(1 for r in (benchmark_rows or []) if r["quartile"] == "Top Quartile"),
        "action_count": len(action_items) if action_items else 0,
    })
    return redirect(url_for("result", sid=sid))


@app.route("/result/<sid>", methods=["GET"])
def result(sid):
    sess_dir = _session_dir(sid)
    meta = _load_json(sess_dir, "meta.json")
    if meta is None:
        flash("That session has expired. Please start again.")
        return redirect(url_for("index"))
    return render_template("result.html", sid=sid, meta=meta,
                            soffice_missing=not (SOFFICE_PATH and HAVE_FITZ))


@app.route("/slide_image/<sid>/<int:n>", methods=["GET"])
def slide_image(sid, n):
    sess_dir = _session_dir(sid)
    path = os.path.join(sess_dir, "render", f"package_{n}.png")
    if not os.path.exists(path):
        abort(404)
    return send_file(path, mimetype="image/png")


@app.route("/download/<sid>/<which>", methods=["GET"])
def download(sid, which):
    sess_dir = _session_dir(sid)
    files = {
        "deck": ("deck_updated.pptx", "application/vnd.openxmlformats-officedocument.presentationml.presentation"),
        "package": ("executive_package.pptx", "application/vnd.openxmlformats-officedocument.presentationml.presentation"),
        "variance": ("variance_analysis.xlsx", "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"),
        "benchmark": ("benchmark_analysis.xlsx", "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"),
        "actions": ("action_tracker.xlsx", "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"),
        "email": ("followup_email.txt", "text/plain"),
        "zip": ("quarter_close_package.zip", "application/zip"),
    }
    if which not in files:
        abort(404)
    fname, mimetype = files[which]
    path = os.path.join(sess_dir, fname)
    if not os.path.exists(path):
        abort(404)
    return send_file(path, as_attachment=True, download_name=fname, mimetype=mimetype)


def _open_browser():
    import webbrowser
    import sys as _sys
    time.sleep(1.2)
    url = "http://127.0.0.1:5080"
    chrome_candidates = []
    if _sys.platform == "darwin":
        chrome_candidates.append("open -a 'Google Chrome' %s")
    elif _sys.platform.startswith("win"):
        for p in (r"C:\Program Files\Google\Chrome\Application\chrome.exe",
                  r"C:\Program Files (x86)\Google\Chrome\Application\chrome.exe"):
            if os.path.exists(p):
                chrome_candidates.append(p.replace("\\", "\\\\") + " %s")
    else:
        chrome_candidates += ["google-chrome %s", "chromium-browser %s", "chromium %s"]
    for template in chrome_candidates:
        try:
            webbrowser.get(template).open(url)
            return
        except webbrowser.Error:
            continue
    webbrowser.open(url)


if __name__ == "__main__":
    import sys
    if os.environ.get("WERKZEUG_RUN_MAIN") != "true":
        threading.Thread(target=_open_browser, daemon=True).start()
    print("Starting Close Cockpit at http://127.0.0.1:5080")
    app.run(host="127.0.0.1", port=5080, debug=("--debug" in sys.argv))
