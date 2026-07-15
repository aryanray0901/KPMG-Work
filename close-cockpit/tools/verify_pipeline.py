"""
Verifies the full Close Cockpit pipeline against sample data: deck
matching, variance derivation and deduplication, benchmarking, action
item parsing, and the generated executive package's rendering (checking
for off-slide or overlapping text, the same check that caught a real
layout bug during development).

Usage: python tools/verify_pipeline.py
"""
import os
import sys
import shutil
import subprocess
import tempfile

sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

import pandas as pd
from pptx import Presentation
from app import (
    extract_targets_pptx, extract_from_spreadsheet, match_targets_to_source,
    collect_all_pptx_text, compute_period_replacements, apply_edits_pptx,
    build_variance_rows, draft_variance_commentary, analyze_benchmark,
    parse_meeting_notes, build_executive_package, build_variance_xlsx,
    build_benchmark_xlsx, build_action_tracker_xlsx, build_followup_email,
    SOFFICE_PATH, HAVE_FITZ,
)

BASE = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
SAMPLES = os.path.join(BASE, "sample_files")
OUT = tempfile.mkdtemp(prefix="cc_verify_")
failures = []


def check(name, condition, detail=""):
    status = "PASS" if condition else "FAIL"
    print(f"[{status}] {name}" + (f" — {detail}" if detail and not condition else ""))
    if not condition:
        failures.append(name)


def render_and_check(pptx_path, label):
    if not (SOFFICE_PATH and HAVE_FITZ):
        print(f"[SKIP] {label} render check — LibreOffice/PyMuPDF not available")
        return
    import fitz
    profile_dir = tempfile.mkdtemp(prefix="lo_verify_")
    try:
        result = subprocess.run(
            [SOFFICE_PATH, "--headless", "--norestore",
             f"-env:UserInstallation=file://{profile_dir}",
             "--convert-to", "pdf", "--outdir", OUT, pptx_path],
            capture_output=True, timeout=120,
        )
        pdf_path = os.path.join(OUT, os.path.splitext(os.path.basename(pptx_path))[0] + ".pdf")
        check(f"{label}: converts to PDF", result.returncode == 0 and os.path.exists(pdf_path))
        if not os.path.exists(pdf_path):
            return
        doc = fitz.open(pdf_path)
        issues = 0
        for page in doc:
            spans = []
            for b in page.get_text("dict")["blocks"]:
                for l in b.get("lines", []):
                    for s in l.get("spans", []):
                        if s["text"].strip():
                            spans.append((s["bbox"], s["text"]))
            pw, ph = page.rect.width, page.rect.height
            for bbox, _ in spans:
                x0, y0, x1, y1 = bbox
                if x0 < -1 or y0 < -1 or x1 > pw + 1 or y1 > ph + 1:
                    issues += 1
            for i in range(len(spans)):
                for j in range(i + 1, len(spans)):
                    b1, _ = spans[i]
                    b2, _ = spans[j]
                    ox = max(0, min(b1[2], b2[2]) - max(b1[0], b2[0]))
                    oy = max(0, min(b1[3], b2[3]) - max(b1[1], b2[1]))
                    a1 = (b1[2] - b1[0]) * (b1[3] - b1[1])
                    a2 = (b2[2] - b2[0]) * (b2[3] - b2[1])
                    if ox * oy > 0.3 * min(a1, a2) and ox * oy > 20:
                        issues += 1
        check(f"{label}: no off-slide or overlapping text across {len(doc)} slides", issues == 0, f"{issues} issue(s)")
        doc.close()
    finally:
        shutil.rmtree(profile_dir, ignore_errors=True)


print("=== Stage 1: Deck matching ===")
deck_path = os.path.join(SAMPLES, "kpmg_advisory_q3_original.pptx")
data_path = os.path.join(SAMPLES, "kpmg_advisory_q4_data.xlsx")
prs, targets = extract_targets_pptx(deck_path)
pairs = extract_from_spreadsheet(data_path)
matches = match_targets_to_source(targets, pairs)
matched = [m for m in matches if m["matched_label"]]
check("matches all 120 known targets", len(matched) == 120, f"got {len(matched)}")
deck_texts = collect_all_pptx_text(Presentation(deck_path))
period_replacements = compute_period_replacements(deck_texts, pairs)
check("detects Q3 -> Q4 period change", period_replacements == {"Q3": "Q4"}, f"got {period_replacements}")

updated_path = os.path.join(OUT, "deck_updated.pptx")
confirmed = [dict(m, new_text=m["new_text_preview"]) for m in matched]
apply_edits_pptx(deck_path, updated_path, confirmed, period_replacements)
check("updated deck file created", os.path.exists(updated_path))
render_and_check(updated_path, "Updated deck")

print("\n=== Stage 2: Variance analysis ===")
variance_rows = build_variance_rows(matches, "Budget", 5.0)
old_new_pairs = [(round(r["old"], 2), round(r["new"], 2)) for r in variance_rows]
check("variance rows contain no duplicate value pairs",
      len(old_new_pairs) == len(set(old_new_pairs)), f"{len(old_new_pairs)} rows, {len(set(old_new_pairs))} unique")
commentary = draft_variance_commentary(variance_rows, "Budget")
check("drafts commentary for material variances", len(commentary) > 0, f"got {len(commentary)}")
material_count = sum(1 for r in variance_rows if r["is_material"])
favorable_count = sum(1 for r in variance_rows if r["tag"] == "Favorable")
unfavorable_count = sum(1 for r in variance_rows if r["tag"] == "Unfavorable")

print("\n=== Stage 3: Benchmarking ===")
bdf = pd.read_excel(os.path.join(SAMPLES, "sample_peer_benchmark.xlsx"))
benchmark_rows = analyze_benchmark(bdf, "Metric", "Target", ["Peer A", "Peer B", "Peer C"])
check("benchmarks all metric rows", len(benchmark_rows) == 4, f"got {len(benchmark_rows)}")
dso_row = next(r for r in benchmark_rows if "Days Sales" in r["metric"])
check("classifies Days Sales Outstanding as lower-is-better", dso_row["higher_is_good"] is False)

print("\n=== Stage 4: Action tracker ===")
with open(os.path.join(SAMPLES, "sample_meeting_notes.txt")) as f:
    notes_text = f.read()
attendees = ["Sarah Chen", "James Cole", "Priya Nair", "Marcus Lee"]
action_items = parse_meeting_notes(notes_text, attendees)
check("parses expected number of action items", len(action_items) == 4, f"got {len(action_items)}")
check("assigns owners only on explicit attendee match",
      all(item["owner"] in attendees for item in action_items),
      f"owners: {[i['owner'] for i in action_items]}")

print("\n=== Assembly: Executive package ===")
config = {"title": "Q4 FY26 Quarter Close Package", "subtitle": "Advisory Practice - Meridian Engagement",
          "comparison_label": "Budget", "target_name": "Meridian Retail Group", "meeting_title": "Q4 Close Call"}
stats = {"matched_count": len(matched), "material_count": material_count,
         "favorable_count": favorable_count, "unfavorable_count": unfavorable_count}
package_path = os.path.join(OUT, "executive_package.pptx")
build_executive_package(package_path, config, stats, variance_rows, commentary, benchmark_rows, action_items)
p = Presentation(package_path)
check("package has 5 slides (cover, summary, variance, benchmark, actions)", len(p.slides) == 5, f"got {len(p.slides)}")
render_and_check(package_path, "Executive package")

build_variance_xlsx(variance_rows, commentary, config["title"], config["comparison_label"],
                     os.path.join(OUT, "variance.xlsx"))
build_benchmark_xlsx(benchmark_rows, config["title"], config["target_name"], os.path.join(OUT, "benchmark.xlsx"))
build_action_tracker_xlsx(action_items, config["meeting_title"], os.path.join(OUT, "actions.xlsx"))
build_followup_email(action_items, config["meeting_title"], "Alex Rivera", stats)
check("all supporting workbooks created",
      all(os.path.exists(os.path.join(OUT, f)) for f in ["variance.xlsx", "benchmark.xlsx", "actions.xlsx"]))

shutil.rmtree(OUT, ignore_errors=True)

print("\n" + "=" * 40)
if failures:
    print(f"FAIL: {len(failures)} check(s) failed: {failures}")
    sys.exit(1)
else:
    print("ALL CHECKS PASSED")
