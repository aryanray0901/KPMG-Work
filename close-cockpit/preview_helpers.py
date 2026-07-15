"""Reusable local preview helpers for generated PowerPoint and Excel files."""

import html
import math
import os
import shutil
import subprocess
import tempfile
import threading
from datetime import date, datetime, time as dt_time

from openpyxl import load_workbook
from openpyxl.styles.numbers import is_date_format
from openpyxl.utils import get_column_letter
from pptx import Presentation

SOFFICE_PATH = shutil.which("soffice") or shutil.which("libreoffice")
POWERSHELL_PATH = shutil.which("powershell") or shutil.which("pwsh")
_soffice_lock = threading.Lock()
_powerpoint_lock = threading.Lock()

try:
    import fitz
    HAVE_FITZ = True
except ImportError:
    fitz = None
    HAVE_FITZ = False


def _render_pptx_with_powerpoint(pptx_path, out_dir, prefix):
    if os.name != "nt" or not POWERSHELL_PATH:
        return None
    os.makedirs(out_dir, exist_ok=True)
    script_path = os.path.join(out_dir, f"render_{prefix}.ps1")
    script = r'''
param(
    [Parameter(Mandatory=$true)][string]$InputPath,
    [Parameter(Mandatory=$true)][string]$OutputDir,
    [Parameter(Mandatory=$true)][string]$Prefix
)
$ErrorActionPreference = "Stop"
$powerpoint = $null
$presentation = $null
try {
    $powerpoint = New-Object -ComObject PowerPoint.Application
    $presentation = $powerpoint.Presentations.Open($InputPath, -1, 0, 0)
    $slideWidth = [double]$presentation.PageSetup.SlideWidth
    $slideHeight = [double]$presentation.PageSetup.SlideHeight
    $exportWidth = 1600
    $exportHeight = [int][Math]::Round($exportWidth * ($slideHeight / $slideWidth))
    for ($i = 1; $i -le $presentation.Slides.Count; $i++) {
        $outputPath = Join-Path $OutputDir ("{0}_{1}.png" -f $Prefix, $i)
        $presentation.Slides.Item($i).Export($outputPath, "PNG", $exportWidth, $exportHeight)
    }
}
finally {
    if ($presentation -ne $null) {
        $presentation.Close()
        [void][System.Runtime.InteropServices.Marshal]::ReleaseComObject($presentation)
    }
    if ($powerpoint -ne $null) {
        $powerpoint.Quit()
        [void][System.Runtime.InteropServices.Marshal]::ReleaseComObject($powerpoint)
    }
    [GC]::Collect()
    [GC]::WaitForPendingFinalizers()
}
'''
    try:
        with open(script_path, "w", encoding="utf-8") as f:
            f.write(script)
        with _powerpoint_lock:
            result = subprocess.run(
                [POWERSHELL_PATH, "-NoProfile", "-ExecutionPolicy", "Bypass",
                 "-File", script_path, "-InputPath", os.path.abspath(pptx_path),
                 "-OutputDir", os.path.abspath(out_dir), "-Prefix", prefix],
                capture_output=True, timeout=180,
            )
        if result.returncode != 0:
            return None
        paths = []
        index = 1
        while True:
            path = os.path.join(out_dir, f"{prefix}_{index}.png")
            if not os.path.exists(path):
                break
            paths.append(path)
            index += 1
        return paths or None
    except Exception:
        return None
    finally:
        try:
            os.remove(script_path)
        except OSError:
            pass


def _convert_pptx_to_pdf(pptx_path, out_dir):
    if not SOFFICE_PATH:
        return None
    profile_dir = tempfile.mkdtemp(prefix="lo_profile_")
    try:
        with _soffice_lock:
            result = subprocess.run(
                [SOFFICE_PATH, "--headless", "--norestore",
                 f"-env:UserInstallation=file://{profile_dir}",
                 "--convert-to", "pdf", "--outdir", out_dir, pptx_path],
                capture_output=True, timeout=120,
            )
        pdf_path = os.path.join(out_dir, os.path.splitext(os.path.basename(pptx_path))[0] + ".pdf")
        return pdf_path if result.returncode == 0 and os.path.exists(pdf_path) else None
    except Exception:
        return None
    finally:
        shutil.rmtree(profile_dir, ignore_errors=True)


def _render_pptx_with_libreoffice(pptx_path, out_dir, prefix):
    if not (SOFFICE_PATH and HAVE_FITZ):
        return None
    os.makedirs(out_dir, exist_ok=True)
    pdf_path = _convert_pptx_to_pdf(pptx_path, out_dir)
    if not pdf_path:
        return None
    try:
        doc = fitz.open(pdf_path)
        paths = []
        for index, page in enumerate(doc, start=1):
            pix = page.get_pixmap(matrix=fitz.Matrix(2, 2), alpha=False)
            img_path = os.path.join(out_dir, f"{prefix}_{index}.png")
            pix.save(img_path)
            paths.append(img_path)
        doc.close()
        return paths or None
    except Exception:
        return None
    finally:
        try:
            os.remove(pdf_path)
        except OSError:
            pass


def render_pptx_to_images(pptx_path, out_dir, prefix):
    """Render every slide. PowerPoint is preferred on Windows, then LibreOffice."""
    paths = _render_pptx_with_powerpoint(pptx_path, out_dir, prefix)
    if paths:
        return paths, "Microsoft PowerPoint"
    paths = _render_pptx_with_libreoffice(pptx_path, out_dir, prefix)
    if paths:
        return paths, "LibreOffice"
    return None, None


def presentation_text_manifest(pptx_path):
    """Return a complete text/table/chart fallback when native slide rendering is unavailable."""
    prs = Presentation(pptx_path)
    slides = []
    for slide_number, slide in enumerate(prs.slides, start=1):
        blocks = []
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                text = "\n".join(
                    p.text.strip() for p in shape.text_frame.paragraphs if p.text and p.text.strip()
                ).strip()
                if text:
                    blocks.append({"type": "text", "content": text})
            if getattr(shape, "has_table", False):
                table_rows = []
                for row in shape.table.rows:
                    table_rows.append([cell.text.strip() for cell in row.cells])
                if table_rows:
                    blocks.append({"type": "table", "rows": table_rows})
            if getattr(shape, "has_chart", False):
                chart = shape.chart
                categories = []
                try:
                    categories = [str(c.label) for c in chart.plots[0].categories]
                except Exception:
                    pass
                series_rows = []
                for series in chart.series:
                    try:
                        values = [None if v is None else v for v in series.values]
                    except Exception:
                        values = []
                    series_rows.append({"name": str(series.name), "values": values})
                if categories or series_rows:
                    blocks.append({"type": "chart", "categories": categories, "series": series_rows})
        slides.append({"number": slide_number, "blocks": blocks})
    return slides


def workbook_manifest(path):
    wb = load_workbook(path, data_only=False, read_only=False)
    sheets = []
    for index, ws in enumerate(wb.worksheets):
        sheets.append({
            "index": index,
            "name": ws.title,
            "rows": max(ws.max_row, 1),
            "columns": max(ws.max_column, 1),
        })
    wb.close()
    return {"filename": os.path.basename(path), "sheets": sheets}


def _argb_to_hex(color):
    if color is None or getattr(color, "type", None) != "rgb":
        return None
    rgb = getattr(color, "rgb", None)
    if not rgb:
        return None
    rgb = str(rgb)
    if len(rgb) == 8:
        rgb = rgb[2:]
    if len(rgb) != 6 or rgb.upper() in {"000000", "FFFFFF"}:
        return rgb.upper() if len(rgb) == 6 else None
    return rgb.upper()


def _format_value(cell):
    value = cell.value
    if value is None:
        return ""
    if isinstance(value, (datetime, date, dt_time)):
        if isinstance(value, datetime):
            return value.strftime("%Y-%m-%d %H:%M")
        if isinstance(value, date):
            return value.strftime("%Y-%m-%d")
        return value.strftime("%H:%M")
    if isinstance(value, bool):
        return "TRUE" if value else "FALSE"
    if isinstance(value, str):
        return value
    if isinstance(value, (int, float)):
        fmt = cell.number_format or "General"
        if is_date_format(fmt):
            return str(value)
        if "%" in fmt:
            decimals = 0
            if "." in fmt:
                decimals = len(fmt.split(".", 1)[1].split("%", 1)[0].replace("#", "0"))
            return f"{value * 100:.{decimals}f}%"
        decimals = 0
        if "." in fmt:
            decimal_part = fmt.split(".", 1)[1]
            decimals = sum(1 for ch in decimal_part if ch in "0#")
            decimals = min(decimals, 6)
        use_commas = "," in fmt or abs(value) >= 1000
        if isinstance(value, int) or (isinstance(value, float) and value.is_integer() and decimals == 0):
            text = f"{int(value):,}" if use_commas else str(int(value))
        else:
            text = f"{value:,.{decimals or 2}f}" if use_commas else f"{value:.{decimals or 2}f}"
        if "$" in fmt:
            text = "$" + text
        elif "€" in fmt:
            text = "€" + text
        elif "£" in fmt:
            text = "£" + text
        return text
    return str(value)


def _cell_style(cell):
    styles = []
    fill = _argb_to_hex(cell.fill.fgColor) if cell.fill and cell.fill.fill_type else None
    font_color = _argb_to_hex(cell.font.color) if cell.font and cell.font.color else None
    if fill:
        styles.append(f"background-color:#{fill}")
    if font_color:
        styles.append(f"color:#{font_color}")
    if cell.font and cell.font.bold:
        styles.append("font-weight:700")
    if cell.font and cell.font.italic:
        styles.append("font-style:italic")
    if cell.alignment:
        if cell.alignment.horizontal in {"left", "center", "right"}:
            styles.append(f"text-align:{cell.alignment.horizontal}")
        if cell.alignment.wrap_text:
            styles.append("white-space:pre-wrap")
    return ";".join(styles)


def workbook_page(path, sheet_index=0, page=1, rows_per_page=50, column_page=1, columns_per_page=24):
    wb = load_workbook(path, data_only=False, read_only=False)
    if sheet_index < 0 or sheet_index >= len(wb.worksheets):
        wb.close()
        raise IndexError("Invalid sheet")
    ws = wb.worksheets[sheet_index]
    sheet_name = ws.title
    max_row = max(ws.max_row, 1)
    total_columns = max(ws.max_column, 1)
    pages = max(1, math.ceil(max_row / rows_per_page))
    column_pages = max(1, math.ceil(total_columns / columns_per_page))
    page = max(1, min(int(page), pages))
    column_page = max(1, min(int(column_page), column_pages))
    start_row = (page - 1) * rows_per_page + 1
    end_row = min(max_row, start_row + rows_per_page - 1)
    start_col = (column_page - 1) * columns_per_page + 1
    end_col = min(total_columns, start_col + columns_per_page - 1)

    merged_anchors = {}
    merged_children = set()
    for merged in ws.merged_cells.ranges:
        min_col, min_row, max_col_m, max_row_m = merged.bounds
        anchor = (min_row, min_col)
        merged_anchors[anchor] = (max_row_m - min_row + 1, max_col_m - min_col + 1)
        for row in range(min_row, max_row_m + 1):
            for col in range(min_col, max_col_m + 1):
                if (row, col) != anchor:
                    merged_children.add((row, col))

    parts = ['<div class="sheet-scroll"><table class="sheet-table"><thead><tr><th class="sheet-corner"></th>']
    for col in range(start_col, end_col + 1):
        parts.append(f'<th class="sheet-col-head">{get_column_letter(col)}</th>')
    parts.append('</tr></thead><tbody>')
    for row in range(start_row, end_row + 1):
        parts.append(f'<tr><th class="sheet-row-head">{row}</th>')
        for col in range(start_col, end_col + 1):
            if (row, col) in merged_children:
                continue
            cell = ws.cell(row=row, column=col)
            attrs = []
            if (row, col) in merged_anchors:
                rowspan, colspan = merged_anchors[(row, col)]
                visible_colspan = min(colspan, end_col - col + 1)
                visible_rowspan = min(rowspan, end_row - row + 1)
                if visible_rowspan > 1:
                    attrs.append(f'rowspan="{visible_rowspan}"')
                if visible_colspan > 1:
                    attrs.append(f'colspan="{visible_colspan}"')
            style = _cell_style(cell)
            if style:
                attrs.append(f'style="{html.escape(style, quote=True)}"')
            if cell.comment and cell.comment.text:
                attrs.append(f'title="{html.escape(cell.comment.text, quote=True)}"')
                attrs.append('class="has-comment"')
            text = html.escape(_format_value(cell)).replace("\n", "<br>")
            parts.append(f'<td {" ".join(attrs)}>{text}</td>')
        parts.append('</tr>')
    parts.append('</tbody></table></div>')
    wb.close()
    return {
        "html": "".join(parts),
        "sheet_name": sheet_name,
        "sheet_index": sheet_index,
        "page": page,
        "pages": pages,
        "start_row": start_row,
        "end_row": end_row,
        "total_rows": max_row,
        "column_page": column_page,
        "column_pages": column_pages,
        "start_column": start_col,
        "end_column": end_col,
        "start_column_letter": get_column_letter(start_col),
        "end_column_letter": get_column_letter(end_col),
        "total_columns": total_columns,
    }

