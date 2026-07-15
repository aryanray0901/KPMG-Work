"""
Engagement Hub
--------------
A consulting workbench combining three modules used routinely on client
engagements:

  1. Action Tracker    - turns meeting notes into a reviewed action item
                          list, an Excel tracker, a follow-up email draft,
                          and a one-page summary slide.
  2. Executive Summary  - turns a period-over-period data table into a
                          KPI-and-narrative executive summary slide.
  3. Benchmarking       - compares a target company's metrics against a
                          peer set and produces a quartile-ranked
                          benchmarking slide and workbook.

Single self-contained Flask application, no external services. All
PowerPoint generation goes through a shared text-box helper that fixes a
real bug found in an earlier version of this codebase: text boxes created
without explicit word-wrap and auto-size settings get expanded by
PowerPoint/LibreOffice to fit their text on one line, which can overflow
into neighboring shapes. Every text box here sets both explicitly.
"""

import os
import re
import io
import json
import time
import uuid
import shutil
import subprocess
import tempfile
import threading
from datetime import datetime
from statistics import median

from flask import Flask, request, render_template, send_file, redirect, url_for, flash, abort, jsonify
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
import openpyxl
from openpyxl.styles import Font, PatternFill, Alignment, Border, Side

BASE_DIR = os.path.dirname(os.path.abspath(__file__))
SESSIONS_DIR = os.path.join(BASE_DIR, "sessions")
os.makedirs(SESSIONS_DIR, exist_ok=True)

app = Flask(__name__)
app.secret_key = "engagement-hub-local-secret"
app.config["MAX_CONTENT_LENGTH"] = 50 * 1024 * 1024  # 50MB

KPMG_BLUE = RGBColor(0x00, 0x33, 0x8D)
KPMG_COBALT = RGBColor(0x00, 0x91, 0xDA)
KPMG_LIGHT = RGBColor(0xF2, 0xF5, 0xFA)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x1A, 0x1A, 0x1A)
GREY = RGBColor(0x77, 0x77, 0x77)
GOOD = RGBColor(0x1E, 0x84, 0x49)
BAD = RGBColor(0xC0, 0x39, 0x2B)
GOOD_HEX, BAD_HEX = "1E8449", "C0392B"
PALETTE = [KPMG_BLUE, KPMG_COBALT, RGBColor(0x00, 0xA3, 0x93), RGBColor(0x5A, 0x28, 0x8C), RGBColor(0x8A, 0x8D, 0x8F)]


@app.after_request
def _no_cache(response):
    response.headers["Cache-Control"] = "no-store, no-cache, must-revalidate, max-age=0"
    return response


# ---------------------------------------------------------------------------
# SHARED PPTX HELPERS
# Every text box created here explicitly sets word_wrap=True and
# auto_size=NONE, so PowerPoint/LibreOffice always respects the box
# dimensions instead of expanding to fit text on one line.
# ---------------------------------------------------------------------------

def add_textbox(slide, x, y, w, h, text, size=14, bold=False, italic=False,
                 color=DARK, align=PP_ALIGN.LEFT, font="Arial", wrap=True):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = wrap
    tf.auto_size = MSO_AUTO_SIZE.NONE
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font
    return box


def add_header_bar(slide, prs, title_text, subtitle_text=None):
    bar = slide.shapes.add_shape(1, 0, 0, prs.slide_width, Inches(0.9))
    bar.fill.solid()
    bar.fill.fore_color.rgb = KPMG_BLUE
    bar.line.fill.background()
    tf = bar.text_frame
    tf.word_wrap = False
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.margin_left = Inches(0.3)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = "KPMG"
    r.font.size = Pt(22)
    r.font.bold = True
    r.font.color.rgb = WHITE
    r.font.name = "Arial"

    add_textbox(slide, 3.4, 0.12, 9.6, 0.65, title_text, size=16, bold=True,
                color=WHITE, align=PP_ALIGN.RIGHT)

    if subtitle_text:
        add_textbox(slide, 0.4, 1.0, 12.5, 0.35, subtitle_text, size=11,
                    italic=True, color=GREY)


def add_kpi_tile(slide, x, y, w, h, value, label, accent=KPMG_BLUE):
    tile = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    tile.fill.solid()
    tile.fill.fore_color.rgb = KPMG_LIGHT
    tile.line.color.rgb = accent
    tile.line.width = Pt(1)
    tf = tile.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.margin_left = Inches(0.15)
    tf.margin_right = Inches(0.15)
    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    r1 = p1.add_run()
    r1.text = value
    r1.font.size = Pt(24)
    r1.font.bold = True
    r1.font.color.rgb = KPMG_BLUE
    r1.font.name = "Arial"
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = label
    r2.font.size = Pt(11)
    r2.font.color.rgb = DARK
    r2.font.name = "Arial"
    return tile


def add_bullets(slide, x, y, w, h, bullets, size=12):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = "•  " + b
        run.font.size = Pt(size)
        run.font.color.rgb = DARK
        run.font.name = "Arial"
    return box


def add_table(slide, x, y, w, h, headers, rows, col_widths=None,
              row_colorizer=None, header_color=KPMG_BLUE):
    n_rows = 1 + len(rows)
    shape = slide.shapes.add_table(n_rows, len(headers), Inches(x), Inches(y), Inches(w), Inches(h))
    table = shape.table
    if col_widths:
        for i, cw in enumerate(col_widths):
            table.columns[i].width = Inches(cw)
    for c, htext in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = str(htext)
        run = cell.text_frame.paragraphs[0].runs[0]
        run.font.size = Pt(11.5)
        run.font.bold = True
        run.font.color.rgb = WHITE
        run.font.name = "Arial"
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_color
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT
    for r_idx, row in enumerate(rows, start=1):
        for c_idx, val in enumerate(row):
            cell = table.cell(r_idx, c_idx)
            cell.text = str(val)
            run = cell.text_frame.paragraphs[0].runs[0]
            run.font.size = Pt(10.5)
            run.font.name = "Arial"
            run.font.color.rgb = DARK
            cell.fill.solid()
            cell.fill.fore_color.rgb = KPMG_LIGHT if (r_idx % 2 == 0) else WHITE
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER if c_idx > 0 else PP_ALIGN.LEFT
            if row_colorizer:
                override = row_colorizer(r_idx - 1, c_idx, val)
                if override:
                    run.font.color.rgb = override
                    run.font.bold = True
    return table


def new_slide(prs=None):
    if prs is None:
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    return prs, slide


# ---------------------------------------------------------------------------
# SHARED XLSX HELPERS
# ---------------------------------------------------------------------------

def style_header_row(ws, row, ncols, fill="00338D"):
    for c in range(1, ncols + 1):
        cell = ws.cell(row=row, column=c)
        cell.font = Font(bold=True, color="FFFFFF")
        cell.fill = PatternFill("solid", fgColor=fill)
        cell.alignment = Alignment(horizontal="center" if c > 1 else "left")


def thin_border():
    thin = Side(style="thin", color="D7DEEA")
    return Border(left=thin, right=thin, top=thin, bottom=thin)


# ---------------------------------------------------------------------------
# GENERATED FILE PREVIEWS
# ---------------------------------------------------------------------------

from preview_helpers import (
    SOFFICE_PATH, POWERSHELL_PATH, HAVE_FITZ, render_pptx_to_images,
    presentation_text_manifest, workbook_manifest, workbook_page,
)


def _session_dir(sid):
    d = os.path.join(SESSIONS_DIR, sid)
    if not os.path.isdir(d):
        abort(404)
    return d


# ===========================================================================
# MODULE 1: ACTION TRACKER
# Turns meeting notes into a reviewed action item list, an Excel tracker,
# a follow-up email draft, and a one-page summary slide.
# ===========================================================================

BULLET_RE = re.compile(r'^\s*[-*•]\s+|^\s*\d+[\.\)]\s+')
ACTION_RE = re.compile(
    r'\b(will (?:send|share|review|confirm|finalize|circulate|draft|prepare|follow up|'
    r'schedule|update|provide|reach out|complete|deliver|escalate|coordinate)|'
    r'needs to|action item|next step|follow[- ]?up|to-?do)\b',
    re.IGNORECASE)
URGENT_RE = re.compile(r'\b(urgent|asap|critical|immediately|high priority)\b', re.IGNORECASE)
DATE_RE = re.compile(
    r'\bby\s+(?:next\s+)?(?:monday|tuesday|wednesday|thursday|friday|saturday|sunday|'
    r'eod|cob|end of day|end of week|\d{1,2}/\d{1,2}(?:/\d{2,4})?|[A-Z][a-z]+\s+\d{1,2}(?:st|nd|rd|th)?)\b',
    re.IGNORECASE)


def extract_owner(line, attendees):
    """Only assigns an owner if an attendee's name is actually present in the
    line. Deliberately does not guess from capitalization patterns -- a
    wrong guess is worse than 'Unassigned' for something a reviewer relies
    on to route follow-ups."""
    for a in attendees:
        if a and re.search(r'\b' + re.escape(a) + r'\b', line, re.IGNORECASE):
            return a
    return "Unassigned"


def extract_due(line):
    m = DATE_RE.search(line)
    return m.group(0) if m else "Not specified"


def parse_meeting_notes(text, attendees):
    lines = [l.strip() for l in text.splitlines() if l.strip()]
    items = []
    for line in lines:
        clean = BULLET_RE.sub("", line).strip()
        if len(clean) < 8:
            continue
        is_bullet = bool(BULLET_RE.match(line))
        actionable = bool(ACTION_RE.search(clean))
        if not is_bullet and not actionable:
            continue
        items.append({
            "text": clean,
            "owner": extract_owner(clean, attendees),
            "due": extract_due(clean),
            "priority": "High" if URGENT_RE.search(clean) else "Medium",
            "include": actionable,
        })
    return items


def build_action_tracker_xlsx(items, meeting_title, meeting_date, out_path):
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Action Tracker"
    ws.merge_cells("A1:E1")
    ws["A1"] = f"{meeting_title} — Action Item Tracker"
    ws["A1"].font = Font(bold=True, size=14, color="00338D")
    ws["A2"] = f"Meeting date: {meeting_date}" if meeting_date else ""
    ws["A2"].font = Font(italic=True, size=10, color="777777")

    headers = ["Action Item", "Owner", "Due Date", "Priority", "Status"]
    header_row = 4
    style_header_row(ws, header_row, len(headers))
    for c, h in enumerate(headers, start=1):
        ws.cell(row=header_row, column=c, value=h)

    border = thin_border()
    for i, item in enumerate(items, start=1):
        r = header_row + i
        vals = [item["text"], item["owner"], item["due"], item["priority"], "Open"]
        for c, v in enumerate(vals, start=1):
            cell = ws.cell(row=r, column=c, value=v)
            cell.border = border
            cell.alignment = Alignment(wrap_text=(c == 1), vertical="top")
            if c == 4 and v == "High":
                cell.font = Font(color=BAD_HEX, bold=True)
        if i % 2 == 0:
            for c in range(1, 6):
                ws.cell(row=r, column=c).fill = PatternFill("solid", fgColor="F2F5FA")

    ws.column_dimensions["A"].width = 55
    ws.column_dimensions["B"].width = 18
    ws.column_dimensions["C"].width = 16
    ws.column_dimensions["D"].width = 12
    ws.column_dimensions["E"].width = 12
    wb.save(out_path)


def build_followup_email(items, meeting_title, meeting_date, sender_name):
    by_owner = {}
    for item in items:
        by_owner.setdefault(item["owner"], []).append(item)

    lines = []
    lines.append(f"Subject: Follow-up and action items — {meeting_title}")
    lines.append("")
    lines.append("Hi all,")
    lines.append("")
    lines.append(f"Thanks for the time {'on ' + meeting_date if meeting_date else 'today'}. "
                  f"Summary of action items below, organized by owner.")
    lines.append("")
    for owner in sorted(by_owner.keys(), key=lambda o: (o == "Unassigned", o)):
        lines.append(f"{owner}:")
        for item in by_owner[owner]:
            due_suffix = f" ({item['due']})" if item["due"] != "Not specified" else ""
            lines.append(f"  - {item['text']}{due_suffix}")
        lines.append("")
    lines.append("Please flag if anything above is missing or assigned incorrectly.")
    lines.append("")
    lines.append("Best,")
    lines.append(sender_name or "[Your name]")
    return "\n".join(lines)


def build_action_tracker_pptx(items, meeting_title, meeting_date, out_path):
    prs, slide = new_slide()
    add_header_bar(slide, prs, "Action Item Summary",
                    f"{meeting_title}" + (f" — {meeting_date}" if meeting_date else ""))

    by_owner = {}
    for item in items:
        by_owner.setdefault(item["owner"], []).append(item)

    rows = []
    for owner in sorted(by_owner.keys(), key=lambda o: (o == "Unassigned", o)):
        for item in by_owner[owner]:
            rows.append([item["text"], owner, item["due"], item["priority"]])

    def colorize(row_idx, col_idx, val):
        if col_idx == 3 and val == "High":
            return BAD
        return None

    add_table(slide, 0.4, 1.5, 12.5, min(5.5, 0.4 + 0.35 * len(rows)),
              ["Action Item", "Owner", "Due Date", "Priority"], rows,
              col_widths=[6.5, 2.3, 2.0, 1.7], row_colorizer=colorize)

    prs.save(out_path)


# ===========================================================================
# MODULE 2: EXECUTIVE SUMMARY
# Turns a period-over-period data table (label column + 2+ numeric columns
# in chronological order) into a KPI-and-narrative executive summary slide.
# ===========================================================================

def analyze_trend(df, label_col, period_cols, unit_suffix, metric_name="value"):
    first_col, last_col = period_cols[0], period_cols[-1]
    rows = []
    for _, row in df.iterrows():
        label = str(row[label_col]).strip()
        if not label or pd.isna(row[first_col]) or pd.isna(row[last_col]):
            continue
        first_val = float(row[first_col])
        last_val = float(row[last_col])
        change = last_val - first_val
        pct_change = (change / first_val * 100) if first_val != 0 else None
        rows.append({"label": label, "first": first_val, "last": last_val,
                      "change": change, "pct_change": pct_change})

    total_first = sum(r["first"] for r in rows)
    total_last = sum(r["last"] for r in rows)
    total_change_pct = ((total_last - total_first) / total_first * 100) if total_first != 0 else None

    for r in rows:
        r["share_of_total"] = (r["last"] / total_last * 100) if total_last != 0 else 0

    ranked_by_pct = sorted([r for r in rows if r["pct_change"] is not None],
                            key=lambda r: r["pct_change"], reverse=True)
    top_mover = ranked_by_pct[0] if ranked_by_pct else None
    bottom_mover = ranked_by_pct[-1] if ranked_by_pct else None
    top_contributor = max(rows, key=lambda r: r["share_of_total"]) if rows else None

    bullets = []
    if total_change_pct is not None:
        direction = "grew" if total_change_pct >= 0 else "declined"
        bullets.append(
            f"Total {metric_name} {direction} from ${total_first:,.1f}{unit_suffix} to "
            f"${total_last:,.1f}{unit_suffix} across the period, a change of {total_change_pct:+.1f}%."
        )
    if top_contributor:
        bullets.append(
            f"{top_contributor['label']} was the largest contributor in the most recent period, "
            f"representing {top_contributor['share_of_total']:.1f}% of total {metric_name}."
        )
    if top_mover and top_mover["pct_change"] is not None and top_mover["pct_change"] > 0:
        bullets.append(
            f"{top_mover['label']} posted the largest increase, up {top_mover['pct_change']:.1f}% "
            f"over the period."
        )
    if bottom_mover and bottom_mover["pct_change"] is not None and bottom_mover["pct_change"] < 0:
        bullets.append(
            f"{bottom_mover['label']} declined {abs(bottom_mover['pct_change']):.1f}% over the period "
            f"and may warrant follow-up. [Add context]"
        )

    return {
        "rows": rows, "total_first": total_first, "total_last": total_last,
        "total_change_pct": total_change_pct, "top_mover": top_mover,
        "bottom_mover": bottom_mover, "top_contributor": top_contributor,
        "bullets": bullets, "first_col": str(first_col), "last_col": str(last_col),
        "metric_name": metric_name,
    }


def build_summary_pptx(analysis, title, subtitle, label_col, unit_suffix, out_path):
    prs, slide = new_slide()
    add_header_bar(slide, prs, title, subtitle)

    metric_name = analysis.get("metric_name", "value")
    kpis = [
        (f"${analysis['total_last']:,.1f}{unit_suffix}", f"Total {metric_name} ({analysis['last_col']})"),
        (f"{analysis['total_change_pct']:+.1f}%" if analysis['total_change_pct'] is not None else "n/a",
         f"Change vs. {analysis['first_col']}"),
        (analysis['top_contributor']['label'] if analysis['top_contributor'] else "n/a",
         f"Largest contributor ({analysis['top_contributor']['share_of_total']:.1f}% of total)"
         if analysis['top_contributor'] else ""),
    ]
    tile_w, gap = 3.9, 0.35
    for i, (value, label) in enumerate(kpis):
        add_kpi_tile(slide, 0.4 + i * (tile_w + gap), 1.5, tile_w, 1.5, value, label)

    add_textbox(slide, 0.4, 3.35, 6, 0.35, "Key Highlights", size=14, bold=True, color=KPMG_BLUE)
    bullets_h = max(1.2, 0.5 * len(analysis["bullets"]))
    add_bullets(slide, 0.4, 3.75, 12.5, bullets_h, analysis["bullets"], size=12.5)

    table_y = 3.75 + bullets_h + 0.25
    rows_data = sorted(analysis["rows"], key=lambda r: r["last"], reverse=True)
    table_rows = [[r["label"], f"${r['first']:,.1f}{unit_suffix}", f"${r['last']:,.1f}{unit_suffix}",
                   f"{r['pct_change']:+.1f}%" if r["pct_change"] is not None else "n/a"] for r in rows_data]

    def colorize(row_idx, col_idx, val):
        if col_idx == 3 and isinstance(val, str) and val.endswith("%") and val != "n/a":
            return GOOD if val.startswith("+") else BAD
        return None

    table_h = min(prs.slide_height / 914400 - table_y - 0.3, 0.35 * (len(table_rows) + 1))
    add_table(slide, 0.4, table_y, 12.5, table_h,
              [label_col, analysis["first_col"], analysis["last_col"], "Change"], table_rows,
              col_widths=[5.5, 2.3, 2.3, 2.4], row_colorizer=colorize)

    prs.save(out_path)


def build_summary_xlsx(analysis, title, label_col, unit_suffix, out_path):
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Executive Summary"
    ws.merge_cells("A1:D1")
    ws["A1"] = title
    ws["A1"].font = Font(bold=True, size=14, color="00338D")

    headers = [label_col, analysis["first_col"], analysis["last_col"], "Change %"]
    header_row = 3
    style_header_row(ws, header_row, len(headers))
    for c, h in enumerate(headers, start=1):
        ws.cell(row=header_row, column=c, value=h)

    border = thin_border()
    rows_data = sorted(analysis["rows"], key=lambda r: r["last"], reverse=True)
    for i, r in enumerate(rows_data, start=1):
        row_n = header_row + i
        vals = [r["label"], f"${r['first']:,.1f}{unit_suffix}", f"${r['last']:,.1f}{unit_suffix}",
                f"{r['pct_change']:+.1f}%" if r["pct_change"] is not None else "n/a"]
        for c, v in enumerate(vals, start=1):
            cell = ws.cell(row=row_n, column=c, value=v)
            cell.border = border
            cell.alignment = Alignment(horizontal="center" if c > 1 else "left")
            if c == 4 and isinstance(v, str) and v.endswith("%") and v != "n/a":
                cell.font = Font(color=GOOD_HEX if v.startswith("+") else BAD_HEX, bold=True)
        if i % 2 == 0:
            for c in range(1, 5):
                ws.cell(row=row_n, column=c).fill = PatternFill("solid", fgColor="F2F5FA")

    ws.column_dimensions["A"].width = 30
    for col in "BCD":
        ws.column_dimensions[col].width = 16

    ws2 = wb.create_sheet("Highlights")
    ws2["A1"] = "Key Highlights"
    ws2["A1"].font = Font(bold=True, size=13, color="00338D")
    for i, b in enumerate(analysis["bullets"]):
        ws2.cell(row=3 + i, column=1, value=b).alignment = Alignment(wrap_text=True)
    ws2.column_dimensions["A"].width = 100

    wb.save(out_path)


# ===========================================================================
# MODULE 3: BENCHMARKING
# Compares a target company's metrics against a peer set and produces a
# quartile-ranked benchmarking slide and workbook.
# ===========================================================================

COST_LIKE_WORDS_B = ["cost", "expense", "tax", "interest", "loss", "spend", "churn",
                      "attrition", "days sales", "turnover", "leverage", "debt"]
REVENUE_LIKE_WORDS_B = ["revenue", "income", "profit", "margin", "contribution", "gross",
                         "ebitda", "ebit", "earnings", "growth", "utilization", "rate", "retention"]


def classify_metric_direction(label):
    l = label.lower()
    if any(w in l for w in REVENUE_LIKE_WORDS_B):
        return "higher_is_good"
    if any(w in l for w in COST_LIKE_WORDS_B):
        return "lower_is_good"
    return "higher_is_good"


def percentile_rank(target, peers, higher_is_good=True):
    if not peers:
        return None
    if higher_is_good:
        beaten = sum(1 for p in peers if target >= p)
    else:
        beaten = sum(1 for p in peers if target <= p)
    return beaten / len(peers) * 100


def quartile_label(pct):
    if pct is None:
        return "n/a"
    if pct >= 75:
        return "Top Quartile"
    if pct >= 50:
        return "Above Median"
    if pct >= 25:
        return "Below Median"
    return "Bottom Quartile"


QUARTILE_CSS = {
    "Top Quartile": "top", "Above Median": "above",
    "Below Median": "below", "Bottom Quartile": "bottom", "n/a": "neutral",
}


def analyze_benchmark(df, metric_col, target_col, peer_cols):
    rows = []
    for _, row in df.iterrows():
        metric = str(row[metric_col]).strip()
        if not metric or pd.isna(row[target_col]):
            continue
        try:
            target_val = float(row[target_col])
        except (ValueError, TypeError):
            continue
        peers = []
        for pc in peer_cols:
            v = row.get(pc)
            if pd.notna(v):
                try:
                    peers.append(float(v))
                except (ValueError, TypeError):
                    pass
        if not peers:
            continue
        higher_is_good = classify_metric_direction(metric) == "higher_is_good"
        pct = percentile_rank(target_val, peers, higher_is_good)
        rows.append({
            "metric": metric, "target": target_val, "peers": peers,
            "peer_median": median(peers), "peer_min": min(peers), "peer_max": max(peers),
            "percentile": pct, "quartile": quartile_label(pct),
            "quartile_class": QUARTILE_CSS[quartile_label(pct)], "higher_is_good": higher_is_good,
        })
    return rows


def build_benchmark_pptx(rows, title, subtitle, target_name, out_path):
    prs, slide = new_slide()
    add_header_bar(slide, prs, title, subtitle)

    table_rows = []
    for r in rows:
        table_rows.append([
            r["metric"], f"{r['target']:,.1f}", f"{r['peer_median']:,.1f}",
            f"{r['peer_min']:,.1f} - {r['peer_max']:,.1f}",
            f"{r['percentile']:.0f}th pct" if r["percentile"] is not None else "n/a",
            r["quartile"],
        ])

    def colorize(row_idx, col_idx, val):
        if col_idx == 5:
            if val == "Top Quartile":
                return GOOD
            if val == "Bottom Quartile":
                return BAD
        return None

    table_h = min(prs.slide_height / 914400 - 1.6 - 2.6, 0.4 * (len(table_rows) + 1))
    add_table(slide, 0.4, 1.5, 12.5, table_h,
              ["Metric", target_name, "Peer Median", "Peer Range", "Percentile", "Positioning"],
              table_rows, col_widths=[3.3, 1.8, 1.9, 2.3, 1.7, 1.5], row_colorizer=colorize)

    chart_y = 1.5 + table_h + 0.3
    chart_data = CategoryChartData()
    chart_data.categories = [r["metric"] for r in rows]
    chart_data.add_series("Percentile", tuple(r["percentile"] or 0 for r in rows))
    gf = slide.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, Inches(0.4), Inches(chart_y),
                                 Inches(12.5), Inches(prs.slide_height / 914400 - chart_y - 0.3), chart_data)
    chart = gf.chart
    chart.has_legend = False
    chart.has_title = True
    chart.chart_title.text_frame.text = f"{target_name} Percentile Rank vs. Peers"
    chart.chart_title.text_frame.paragraphs[0].font.size = Pt(12)
    chart.chart_title.text_frame.paragraphs[0].font.bold = True
    chart.chart_title.text_frame.paragraphs[0].font.color.rgb = KPMG_BLUE
    plot = chart.plots[0]
    plot.has_data_labels = True
    plot.data_labels.number_format = '0"th"'
    plot.data_labels.number_format_is_linked = False
    plot.data_labels.font.size = Pt(9)
    for i, pt in enumerate(chart.series[0].points):
        pt.format.fill.solid()
        pt.format.fill.fore_color.rgb = PALETTE[i % len(PALETTE)]
    chart.value_axis.minimum_scale = 0
    chart.value_axis.maximum_scale = 100

    prs.save(out_path)


def build_benchmark_xlsx(rows, title, target_name, out_path):
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Benchmarking"
    ws.merge_cells("A1:F1")
    ws["A1"] = title
    ws["A1"].font = Font(bold=True, size=14, color="00338D")

    headers = ["Metric", target_name, "Peer Median", "Peer Min", "Peer Max", "Percentile", "Positioning"]
    header_row = 3
    style_header_row(ws, header_row, len(headers))
    for c, h in enumerate(headers, start=1):
        ws.cell(row=header_row, column=c, value=h)

    border = thin_border()
    for i, r in enumerate(rows, start=1):
        row_n = header_row + i
        vals = [r["metric"], r["target"], r["peer_median"], r["peer_min"], r["peer_max"],
                f"{r['percentile']:.0f}" if r["percentile"] is not None else "n/a", r["quartile"]]
        for c, v in enumerate(vals, start=1):
            cell = ws.cell(row=row_n, column=c, value=v)
            cell.border = border
            cell.alignment = Alignment(horizontal="center" if c > 1 else "left")
            if c == 7:
                if v == "Top Quartile":
                    cell.font = Font(color=GOOD_HEX, bold=True)
                elif v == "Bottom Quartile":
                    cell.font = Font(color=BAD_HEX, bold=True)
        if i % 2 == 0:
            for c in range(1, 8):
                ws.cell(row=row_n, column=c).fill = PatternFill("solid", fgColor="F2F5FA")

    ws.column_dimensions["A"].width = 26
    for col in "BCDEFG":
        ws.column_dimensions[col].width = 15

    wb.save(out_path)


# ===========================================================================
# ROUTES: DASHBOARD
# ===========================================================================

@app.route("/", methods=["GET"])
def index():
    return render_template("index.html")


def _save_json(sess_dir, name, data):
    with open(os.path.join(sess_dir, name), "w") as f:
        json.dump(data, f)


def _load_json(sess_dir, name, default=None):
    path = os.path.join(sess_dir, name)
    if not os.path.exists(path):
        return default
    with open(path) as f:
        return json.load(f)


def _new_session():
    sid = uuid.uuid4().hex[:12]
    sess_dir = os.path.join(SESSIONS_DIR, sid)
    os.makedirs(sess_dir, exist_ok=True)
    return sid, sess_dir


# ===========================================================================
# ROUTES: MODULE 1 - ACTION TRACKER
# ===========================================================================

@app.route("/actions", methods=["GET"])
def actions_index():
    return render_template("actions_index.html")


@app.route("/actions/process", methods=["POST"])
def actions_process():
    notes_text = request.form.get("notes_text", "").strip()
    notes_file = request.files.get("notes_file")
    if notes_file and notes_file.filename:
        try:
            notes_text = (notes_text + "\n" + notes_file.read().decode("utf-8", errors="ignore")).strip()
        except Exception:
            pass
    if not notes_text:
        flash("Paste your meeting notes or upload a text file to continue.")
        return redirect(url_for("actions_index"))

    meeting_title = request.form.get("meeting_title", "").strip() or "Meeting"
    meeting_date = request.form.get("meeting_date", "").strip()
    sender_name = request.form.get("sender_name", "").strip()
    attendees_raw = request.form.get("attendees", "").strip()
    attendees = [a.strip() for a in attendees_raw.split(",") if a.strip()]

    items = parse_meeting_notes(notes_text, attendees)
    if not items:
        flash("No candidate action items found. Try formatting notes as bullet points, one per line.")
        return redirect(url_for("actions_index"))

    sid, sess_dir = _new_session()
    _save_json(sess_dir, "items.json", items)
    _save_json(sess_dir, "config.json", {
        "meeting_title": meeting_title, "meeting_date": meeting_date, "sender_name": sender_name,
    })
    return redirect(url_for("actions_review", sid=sid))


@app.route("/actions/review/<sid>", methods=["GET"])
def actions_review(sid):
    sess_dir = _session_dir(sid)
    items = _load_json(sess_dir, "items.json")
    config = _load_json(sess_dir, "config.json")
    if items is None:
        flash("That session has expired. Please start again.")
        return redirect(url_for("actions_index"))
    return render_template("actions_review.html", sid=sid, items=items, config=config)


@app.route("/actions/export/<sid>", methods=["POST"])
def actions_export(sid):
    sess_dir = _session_dir(sid)
    items = _load_json(sess_dir, "items.json")
    config = _load_json(sess_dir, "config.json")
    if items is None:
        flash("That session has expired. Please start again.")
        return redirect(url_for("actions_index"))

    confirmed = []
    for i, item in enumerate(items):
        if request.form.get(f"include_{i}") == "on":
            entry = dict(item)
            entry["text"] = request.form.get(f"text_{i}", item["text"]).strip() or item["text"]
            entry["owner"] = request.form.get(f"owner_{i}", item["owner"]).strip() or "Unassigned"
            entry["due"] = request.form.get(f"due_{i}", item["due"]).strip() or "Not specified"
            entry["priority"] = request.form.get(f"priority_{i}", item["priority"])
            confirmed.append(entry)

    if not confirmed:
        flash("Select at least one action item to export.")
        return redirect(url_for("actions_review", sid=sid))

    xlsx_path = os.path.join(sess_dir, "action_tracker.xlsx")
    pptx_path = os.path.join(sess_dir, "action_summary.pptx")
    build_action_tracker_xlsx(confirmed, config["meeting_title"], config["meeting_date"], xlsx_path)
    build_action_tracker_pptx(confirmed, config["meeting_title"], config["meeting_date"], pptx_path)
    email_text = build_followup_email(confirmed, config["meeting_title"], config["meeting_date"], config["sender_name"])
    with open(os.path.join(sess_dir, "followup_email.txt"), "w") as f:
        f.write(email_text)

    render_dir = os.path.join(sess_dir, "render")
    shutil.rmtree(render_dir, ignore_errors=True)
    images, render_engine = render_pptx_to_images(pptx_path, render_dir, "actions")
    _save_json(sess_dir, "meta.json", {
        "rendering_ok": images is not None, "slide_count": len(images) if images else 0,
        "render_engine": render_engine, "count": len(confirmed),
    })
    return redirect(url_for("actions_result", sid=sid))


@app.route("/actions/result/<sid>", methods=["GET"])
def actions_result(sid):
    sess_dir = _session_dir(sid)
    if not os.path.exists(os.path.join(sess_dir, "action_tracker.xlsx")):
        flash("That session has expired. Please start again.")
        return redirect(url_for("actions_index"))
    meta = _load_json(sess_dir, "meta.json", {})
    email_text = ""
    email_path = os.path.join(sess_dir, "followup_email.txt")
    if os.path.exists(email_path):
        with open(email_path) as f:
            email_text = f.read()
    preview_presentations = [{
        "id": "actions", "title": "Action summary slide", "filename": "action_summary.pptx",
        "download_url": url_for("actions_download", sid=sid, which="pptx"),
        "rendering_ok": meta.get("rendering_ok", False),
        "slide_count": meta.get("slide_count", 0), "render_engine": meta.get("render_engine"),
        "image_url": f"/presentation_image/{sid}/actions/__SLIDE__",
        "fallback_slides": presentation_text_manifest(os.path.join(sess_dir, "action_summary.pptx")),
    }]
    workbook_path = os.path.join(sess_dir, "action_tracker.xlsx")
    preview_workbooks = [{
        "id": "actions", "title": "Action tracker", "filename": "action_tracker.xlsx",
        "download_url": url_for("actions_download", sid=sid, which="xlsx"),
        "preview_url": f"/workbook_preview/{sid}/actions",
        "manifest": workbook_manifest(workbook_path),
    }]
    preview_texts = [{
        "title": "Follow-up email draft", "filename": "followup_email.txt",
        "download_url": url_for("actions_download", sid=sid, which="email"), "content": email_text,
    }]
    return render_template(
        "actions_result.html", sid=sid, count=meta.get("count", 0),
        preview_presentations=preview_presentations, preview_workbooks=preview_workbooks,
        preview_texts=preview_texts,
    )


@app.route("/actions/download/<sid>/<which>", methods=["GET"])
def actions_download(sid, which):
    sess_dir = _session_dir(sid)
    files = {
        "xlsx": ("action_tracker.xlsx", "action_tracker.xlsx",
                 "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"),
        "pptx": ("action_summary.pptx", "action_summary.pptx",
                 "application/vnd.openxmlformats-officedocument.presentationml.presentation"),
        "email": ("followup_email.txt", "followup_email.txt", "text/plain"),
    }
    if which not in files:
        abort(404)
    fname, dl_name, mimetype = files[which]
    path = os.path.join(sess_dir, fname)
    if not os.path.exists(path):
        abort(404)
    return send_file(path, as_attachment=True, download_name=dl_name, mimetype=mimetype)


# ===========================================================================
# ROUTES: MODULE 2 - EXECUTIVE SUMMARY
# ===========================================================================

@app.route("/summary", methods=["GET"])
def summary_index():
    return render_template("summary_index.html")


@app.route("/summary/process", methods=["POST"])
def summary_process():
    data_file = request.files.get("data_file")
    if not data_file or data_file.filename == "":
        flash("Please upload an Excel or CSV file with your period-over-period data.")
        return redirect(url_for("summary_index"))

    title = request.form.get("title", "").strip() or "Executive Summary"
    subtitle = request.form.get("subtitle", "").strip()
    metric_name = request.form.get("metric_name", "").strip() or "value"
    currency_suffix = request.form.get("currency_suffix", "").strip() or "M"

    sid, sess_dir = _new_session()
    data_path = os.path.join(sess_dir, "input_" + data_file.filename)
    data_file.save(data_path)

    try:
        df = pd.read_csv(data_path) if data_path.lower().endswith(".csv") else pd.read_excel(data_path)
    except Exception:
        flash("Couldn't read that file. Please upload a valid .xlsx or .csv.")
        shutil.rmtree(sess_dir, ignore_errors=True)
        return redirect(url_for("summary_index"))

    df = df.dropna(how="all").dropna(axis=1, how="all")
    cols = list(df.columns)
    numeric_cols = [c for c in cols if pd.api.types.is_numeric_dtype(df[c])]
    non_numeric_cols = [c for c in cols if c not in numeric_cols]
    if len(numeric_cols) < 2 or not non_numeric_cols:
        flash("Couldn't find a label column plus at least two numeric (period) columns.")
        shutil.rmtree(sess_dir, ignore_errors=True)
        return redirect(url_for("summary_index"))
    label_col = non_numeric_cols[0]

    analysis = analyze_trend(df, label_col, numeric_cols, currency_suffix, metric_name)
    if not analysis["rows"]:
        flash("No usable rows found in that file.")
        shutil.rmtree(sess_dir, ignore_errors=True)
        return redirect(url_for("summary_index"))

    _save_json(sess_dir, "analysis.json", analysis)
    _save_json(sess_dir, "config.json", {
        "title": title, "subtitle": subtitle, "label_col": str(label_col),
        "metric_name": metric_name, "currency_suffix": currency_suffix,
    })
    return redirect(url_for("summary_review", sid=sid))


@app.route("/summary/review/<sid>", methods=["GET"])
def summary_review(sid):
    sess_dir = _session_dir(sid)
    analysis = _load_json(sess_dir, "analysis.json")
    config = _load_json(sess_dir, "config.json")
    if analysis is None:
        flash("That session has expired. Please start again.")
        return redirect(url_for("summary_index"))
    return render_template("summary_review.html", sid=sid, analysis=analysis, config=config)


@app.route("/summary/export/<sid>", methods=["POST"])
def summary_export(sid):
    sess_dir = _session_dir(sid)
    analysis = _load_json(sess_dir, "analysis.json")
    config = _load_json(sess_dir, "config.json")
    if analysis is None:
        flash("That session has expired. Please start again.")
        return redirect(url_for("summary_index"))

    edited_bullets = []
    for i in range(len(analysis["bullets"])):
        if request.form.get(f"include_bullet_{i}") == "on":
            edited_bullets.append(request.form.get(f"bullet_{i}", analysis["bullets"][i]).strip())
    analysis["bullets"] = edited_bullets or analysis["bullets"]

    pptx_path = os.path.join(sess_dir, "executive_summary.pptx")
    xlsx_path = os.path.join(sess_dir, "executive_summary.xlsx")
    build_summary_pptx(analysis, config["title"], config["subtitle"], config["label_col"],
                        config["currency_suffix"], pptx_path)
    build_summary_xlsx(analysis, config["title"], config["label_col"], config["currency_suffix"], xlsx_path)

    render_dir = os.path.join(sess_dir, "render")
    shutil.rmtree(render_dir, ignore_errors=True)
    images, render_engine = render_pptx_to_images(pptx_path, render_dir, "summary")
    _save_json(sess_dir, "meta.json", {
        "rendering_ok": images is not None, "slide_count": len(images) if images else 0,
        "render_engine": render_engine,
    })
    return redirect(url_for("summary_result", sid=sid))


@app.route("/summary/result/<sid>", methods=["GET"])
def summary_result(sid):
    sess_dir = _session_dir(sid)
    if not os.path.exists(os.path.join(sess_dir, "executive_summary.pptx")):
        flash("That session has expired. Please start again.")
        return redirect(url_for("summary_index"))
    meta = _load_json(sess_dir, "meta.json", {})
    preview_presentations = [{
        "id": "summary", "title": "Executive summary slide", "filename": "executive_summary.pptx",
        "download_url": url_for("summary_download", sid=sid, which="pptx"),
        "rendering_ok": meta.get("rendering_ok", False),
        "slide_count": meta.get("slide_count", 0), "render_engine": meta.get("render_engine"),
        "image_url": f"/presentation_image/{sid}/summary/__SLIDE__",
        "fallback_slides": presentation_text_manifest(os.path.join(sess_dir, "executive_summary.pptx")),
    }]
    workbook_path = os.path.join(sess_dir, "executive_summary.xlsx")
    preview_workbooks = [{
        "id": "summary", "title": "Executive summary workbook", "filename": "executive_summary.xlsx",
        "download_url": url_for("summary_download", sid=sid, which="xlsx"),
        "preview_url": f"/workbook_preview/{sid}/summary",
        "manifest": workbook_manifest(workbook_path),
    }]
    return render_template(
        "summary_result.html", sid=sid, preview_presentations=preview_presentations,
        preview_workbooks=preview_workbooks, preview_texts=[],
    )


@app.route("/summary/download/<sid>/<which>", methods=["GET"])
def summary_download(sid, which):
    sess_dir = _session_dir(sid)
    files = {
        "pptx": ("executive_summary.pptx", "application/vnd.openxmlformats-officedocument.presentationml.presentation"),
        "xlsx": ("executive_summary.xlsx", "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"),
    }
    if which not in files:
        abort(404)
    fname, mimetype = files[which]
    path = os.path.join(sess_dir, fname)
    if not os.path.exists(path):
        abort(404)
    return send_file(path, as_attachment=True, download_name=fname, mimetype=mimetype)


# ===========================================================================
# ROUTES: MODULE 3 - BENCHMARKING
# ===========================================================================

@app.route("/benchmark", methods=["GET"])
def benchmark_index():
    return render_template("benchmark_index.html")


@app.route("/benchmark/process", methods=["POST"])
def benchmark_process():
    data_file = request.files.get("data_file")
    if not data_file or data_file.filename == "":
        flash("Please upload an Excel or CSV file with your benchmarking data.")
        return redirect(url_for("benchmark_index"))

    title = request.form.get("title", "").strip() or "Benchmarking Analysis"
    subtitle = request.form.get("subtitle", "").strip()
    target_name = request.form.get("target_name", "").strip() or "Target"

    sid, sess_dir = _new_session()
    data_path = os.path.join(sess_dir, "input_" + data_file.filename)
    data_file.save(data_path)

    try:
        df = pd.read_csv(data_path) if data_path.lower().endswith(".csv") else pd.read_excel(data_path)
    except Exception:
        flash("Couldn't read that file. Please upload a valid .xlsx or .csv.")
        shutil.rmtree(sess_dir, ignore_errors=True)
        return redirect(url_for("benchmark_index"))

    df = df.dropna(how="all").dropna(axis=1, how="all")
    cols = list(df.columns)
    numeric_cols = [c for c in cols if pd.api.types.is_numeric_dtype(df[c])]
    non_numeric_cols = [c for c in cols if c not in numeric_cols]
    if len(numeric_cols) < 2 or not non_numeric_cols:
        flash("Couldn't find a metric column, a target column, and at least one peer column.")
        shutil.rmtree(sess_dir, ignore_errors=True)
        return redirect(url_for("benchmark_index"))
    metric_col = non_numeric_cols[0]
    target_col = numeric_cols[0]
    peer_cols = numeric_cols[1:]

    rows = analyze_benchmark(df, metric_col, target_col, peer_cols)
    if not rows:
        flash("No usable rows found in that file.")
        shutil.rmtree(sess_dir, ignore_errors=True)
        return redirect(url_for("benchmark_index"))

    _save_json(sess_dir, "rows.json", rows)
    _save_json(sess_dir, "config.json", {
        "title": title, "subtitle": subtitle, "target_name": target_name,
        "metric_col": str(metric_col), "target_col": str(target_col),
        "peer_cols": [str(c) for c in peer_cols],
    })
    return redirect(url_for("benchmark_review", sid=sid))


@app.route("/benchmark/review/<sid>", methods=["GET"])
def benchmark_review(sid):
    sess_dir = _session_dir(sid)
    rows = _load_json(sess_dir, "rows.json")
    config = _load_json(sess_dir, "config.json")
    if rows is None:
        flash("That session has expired. Please start again.")
        return redirect(url_for("benchmark_index"))
    top_count = sum(1 for r in rows if r["quartile"] == "Top Quartile")
    bottom_count = sum(1 for r in rows if r["quartile"] == "Bottom Quartile")
    return render_template("benchmark_review.html", sid=sid, rows=rows, config=config,
                            top_count=top_count, bottom_count=bottom_count)


@app.route("/benchmark/export/<sid>", methods=["POST"])
def benchmark_export(sid):
    sess_dir = _session_dir(sid)
    rows = _load_json(sess_dir, "rows.json")
    config = _load_json(sess_dir, "config.json")
    if rows is None:
        flash("That session has expired. Please start again.")
        return redirect(url_for("benchmark_index"))

    included = [r for i, r in enumerate(rows) if request.form.get(f"include_{i}") == "on"]
    if not included:
        included = rows

    pptx_path = os.path.join(sess_dir, "benchmark_analysis.pptx")
    xlsx_path = os.path.join(sess_dir, "benchmark_analysis.xlsx")
    build_benchmark_pptx(included, config["title"], config["subtitle"], config["target_name"], pptx_path)
    build_benchmark_xlsx(included, config["title"], config["target_name"], xlsx_path)

    render_dir = os.path.join(sess_dir, "render")
    shutil.rmtree(render_dir, ignore_errors=True)
    images, render_engine = render_pptx_to_images(pptx_path, render_dir, "benchmark")
    _save_json(sess_dir, "meta.json", {
        "rendering_ok": images is not None, "slide_count": len(images) if images else 0,
        "render_engine": render_engine,
    })
    return redirect(url_for("benchmark_result", sid=sid))


@app.route("/benchmark/result/<sid>", methods=["GET"])
def benchmark_result(sid):
    sess_dir = _session_dir(sid)
    if not os.path.exists(os.path.join(sess_dir, "benchmark_analysis.pptx")):
        flash("That session has expired. Please start again.")
        return redirect(url_for("benchmark_index"))
    meta = _load_json(sess_dir, "meta.json", {})
    preview_presentations = [{
        "id": "benchmark", "title": "Benchmarking slide", "filename": "benchmark_analysis.pptx",
        "download_url": url_for("benchmark_download", sid=sid, which="pptx"),
        "rendering_ok": meta.get("rendering_ok", False),
        "slide_count": meta.get("slide_count", 0), "render_engine": meta.get("render_engine"),
        "image_url": f"/presentation_image/{sid}/benchmark/__SLIDE__",
        "fallback_slides": presentation_text_manifest(os.path.join(sess_dir, "benchmark_analysis.pptx")),
    }]
    workbook_path = os.path.join(sess_dir, "benchmark_analysis.xlsx")
    preview_workbooks = [{
        "id": "benchmark", "title": "Benchmarking workbook", "filename": "benchmark_analysis.xlsx",
        "download_url": url_for("benchmark_download", sid=sid, which="xlsx"),
        "preview_url": f"/workbook_preview/{sid}/benchmark",
        "manifest": workbook_manifest(workbook_path),
    }]
    return render_template(
        "benchmark_result.html", sid=sid, preview_presentations=preview_presentations,
        preview_workbooks=preview_workbooks, preview_texts=[],
    )


@app.route("/benchmark/download/<sid>/<which>", methods=["GET"])
def benchmark_download(sid, which):
    sess_dir = _session_dir(sid)
    files = {
        "pptx": ("benchmark_analysis.pptx", "application/vnd.openxmlformats-officedocument.presentationml.presentation"),
        "xlsx": ("benchmark_analysis.xlsx", "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"),
    }
    if which not in files:
        abort(404)
    fname, mimetype = files[which]
    path = os.path.join(sess_dir, fname)
    if not os.path.exists(path):
        abort(404)
    return send_file(path, as_attachment=True, download_name=fname, mimetype=mimetype)


@app.route("/presentation_image/<sid>/<which>/<int:n>", methods=["GET"])
def presentation_image(sid, which, n):
    if which not in {"actions", "summary", "benchmark"}:
        abort(404)
    path = os.path.join(_session_dir(sid), "render", f"{which}_{n}.png")
    if not os.path.exists(path):
        abort(404)
    return send_file(path, mimetype="image/png")


@app.route("/workbook_preview/<sid>/<which>", methods=["GET"])
def workbook_preview(sid, which):
    files = {
        "actions": "action_tracker.xlsx",
        "summary": "executive_summary.xlsx",
        "benchmark": "benchmark_analysis.xlsx",
    }
    if which not in files:
        abort(404)
    path = os.path.join(_session_dir(sid), files[which])
    if not os.path.exists(path):
        abort(404)
    try:
        sheet = int(request.args.get("sheet", 0))
        page = int(request.args.get("page", 1))
        column_page = int(request.args.get("column_page", 1))
        return jsonify(workbook_page(path, sheet_index=sheet, page=page, column_page=column_page))
    except (ValueError, IndexError):
        abort(404)


def _open_browser():
    import webbrowser
    import sys as _sys
    time.sleep(1.2)
    url = "http://127.0.0.1:5070"
    chrome_candidates = []
    if _sys.platform == "darwin":
        chrome_candidates.append("open -a 'Google Chrome' %s")
    elif _sys.platform.startswith("win"):
        for p in (r"C:\Program Files\Google\Chrome\Application\chrome.exe",
                  r"C:\Program Files (x86)\Google\Chrome\Application\chrome.exe"):
            if os.path.exists(p):
                chrome_candidates.append(p.replace("\\", "\\\\") + " %s")
    else:
        chrome_candidates += ["google-chrome %s", "chromium-browser %s", "chromium %s"]
    for template in chrome_candidates:
        try:
            webbrowser.get(template).open(url)
            return
        except webbrowser.Error:
            continue
    webbrowser.open(url)


if __name__ == "__main__":
    import sys
    if os.environ.get("WERKZEUG_RUN_MAIN") != "true":
        threading.Thread(target=_open_browser, daemon=True).start()
    print("Starting Engagement Hub at http://127.0.0.1:5070")
    app.run(host="127.0.0.1", port=5070, debug=("--debug" in sys.argv))
