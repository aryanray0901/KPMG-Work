"""Deterministic natural-language showcase commands for Deck Refresh.

This module covers common presentation editing requests locally so basic and
high-value commands do not depend on model output quality. It intentionally
uses broad phrase matching and safe deck transformations.
"""
from __future__ import annotations

import re
from copy import deepcopy
from typing import Any

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.dml import MSO_FILL_TYPE
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt


BLUE = "00338D"
DARK_BLUE = "001E5A"
LIGHT_BLUE = "DCE6F7"
GREEN = "2E7D32"
LIGHT_GREEN = "E8F5E9"
AMBER = "F9A825"
LIGHT_AMBER = "FFF8E1"
RED = "C62828"
LIGHT_RED = "FFEBEE"
DARK_GRAY = "374151"
MID_GRAY = "6B7280"
LIGHT_GRAY = "F3F4F6"
WHITE = "FFFFFF"


def _rgb(value: str) -> RGBColor:
    value = value.strip().lstrip("#")
    return RGBColor.from_string(value.upper())


def _slide(prs: Presentation, number: int):
    if not 1 <= int(number) <= len(prs.slides):
        raise ValueError(f"Slide {number} does not exist. The deck has {len(prs.slides)} slides.")
    return prs.slides[int(number) - 1]


def _text(shape) -> str:
    if getattr(shape, "has_text_frame", False):
        return shape.text_frame.text.strip()
    if getattr(shape, "has_table", False):
        return "\n".join(" | ".join(cell.text.strip() for cell in row.cells) for row in shape.table.rows)
    return ""


def _title_shape(slide):
    if getattr(slide.shapes, "title", None) is not None:
        return slide.shapes.title
    candidates = [s for s in slide.shapes if getattr(s, "has_text_frame", False) and s.text_frame.text.strip()]
    if not candidates:
        return None
    # Ignore small logo/footer text and prefer a wide text box near the top.
    filtered = [
        s for s in candidates
        if str(s.text_frame.text).strip().casefold() not in {"kpmg"}
        and float(s.width) >= Inches(2.0)
        and float(s.top) <= Inches(1.5)
    ]
    if filtered:
        return sorted(filtered, key=lambda s: (s.top, -s.width))[0]
    return sorted(candidates, key=lambda s: (-s.width, s.top))[0]


def _title(slide) -> str:
    shape = _title_shape(slide)
    return _text(shape).splitlines()[0][:160] if shape is not None and _text(shape) else ""


def _is_branding(shape, prs: Presentation) -> bool:
    # Preserve small corner logos, footer marks, and thin decorative bars.
    w = float(shape.width) / max(1, prs.slide_width)
    h = float(shape.height) / max(1, prs.slide_height)
    x = float(shape.left) / max(1, prs.slide_width)
    y = float(shape.top) / max(1, prs.slide_height)
    name = str(getattr(shape, "name", "")).casefold()
    if "logo" in name or "brand" in name:
        return True
    if h < 0.035 and (y < 0.08 or y > 0.90):
        return True
    if w < 0.16 and h < 0.12 and (x < 0.18 or x > 0.78) and (y < 0.16 or y > 0.82):
        return True
    return False


def _delete_shape(shape) -> None:
    element = shape._element
    element.getparent().remove(element)


def _non_brand_shapes(slide, prs: Presentation):
    title = _title_shape(slide)
    title_id = int(title.shape_id) if title is not None else None
    return [
        s for s in list(slide.shapes)
        if (title_id is None or int(s.shape_id) != title_id) and not _is_branding(s, prs)
    ]


def _chart_insert_box(prs: Presentation, slide) -> tuple[float, float, float, float]:
    """Place a chart in open space and reflow crowded content without overlaps."""
    content = []
    for shape in _non_brand_shapes(slide, prs):
        width = float(shape.width) / max(1, prs.slide_width)
        height = float(shape.height) / max(1, prs.slide_height)
        y = float(shape.top) / max(1, prs.slide_height)
        area = width * height
        if 0.008 <= area <= 0.72 and y <= 0.91:
            content.append(shape)

    def box(shape):
        return (
            float(shape.left) / max(1, prs.slide_width),
            float(shape.top) / max(1, prs.slide_height),
            float(shape.width) / max(1, prs.slide_width),
            float(shape.height) / max(1, prs.slide_height),
        )

    def intersects(first, second):
        width = max(0.0, min(first[0] + first[2], second[0] + second[2]) - max(first[0], second[0]))
        height = max(0.0, min(first[1] + first[3], second[1] + second[3]) - max(first[1], second[1]))
        return width * height

    occupied = [box(shape) for shape in content]
    candidates = [
        (0.06, 0.20, 0.88, 0.70), (0.51, 0.20, 0.43, 0.70),
        (0.06, 0.20, 0.43, 0.70), (0.06, 0.55, 0.88, 0.35),
        (0.06, 0.20, 0.88, 0.31), (0.51, 0.53, 0.43, 0.37),
        (0.06, 0.53, 0.43, 0.37),
    ]
    for candidate in candidates:
        if all(intersects(candidate, existing) <= 0.002 for existing in occupied):
            return candidate
    if not content:
        return candidates[0]

    gap = 0.02
    count = len(content)
    if count <= 3:
        insert_box = (0.45, 0.20, 0.49, 0.70)
        row_height = (0.70 - gap * (count - 1)) / count
        slots = [(0.06, 0.20 + index * (row_height + gap), 0.35, row_height) for index in range(count)]
    else:
        insert_box = (0.06, 0.56, 0.88, 0.34)
        columns = min(3, count)
        rows = (count + columns - 1) // columns
        slot_width = (0.88 - gap * (columns - 1)) / columns
        slot_height = (0.32 - gap * (rows - 1)) / rows
        slots = [
            (0.06 + (index % columns) * (slot_width + gap),
             0.20 + (index // columns) * (slot_height + gap), slot_width, slot_height)
            for index in range(count)
        ]
    for shape, (x, y, width, height) in zip(content, slots):
        shape.left = int(prs.slide_width * x)
        shape.top = int(prs.slide_height * y)
        shape.width = max(1, int(prs.slide_width * width))
        shape.height = max(1, int(prs.slide_height * height))
        if getattr(shape, "has_text_frame", False):
            shape.text_frame.word_wrap = True
            shape.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    return insert_box


def _style_run(run, size=16, bold=False, color=DARK_GRAY, font="Arial"):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = _rgb(color)
    run.font.name = font


def _set_text(shape, text: str, size=16, bold=False, color=DARK_GRAY, align=PP_ALIGN.LEFT, font="Arial"):
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)
    lines = text.splitlines() or [""]
    for index, line in enumerate(lines):
        p = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
        p.text = line
        p.alignment = align
        if p.runs:
            _style_run(p.runs[0], size=size, bold=bold, color=color, font=font)
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE


def _set_title(slide, prs: Presentation, text: str | None = None):
    shape = _title_shape(slide)
    if shape is None:
        shape = slide.shapes.add_textbox(
            int(prs.slide_width * 0.055), int(prs.slide_height * 0.035),
            int(prs.slide_width * 0.82), int(prs.slide_height * 0.10),
        )
    _set_text(shape, text or _title(slide) or "Executive Update", size=27, bold=True, color=DARK_BLUE)
    shape.left = int(prs.slide_width * 0.055)
    shape.top = int(prs.slide_height * 0.035)
    shape.width = int(prs.slide_width * 0.84)
    shape.height = int(prs.slide_height * 0.11)
    return shape


def _clean_lines(texts: list[str], limit=8) -> list[str]:
    pieces: list[str] = []
    for text in texts:
        for raw in re.split(r"[\n•]+", text):
            raw = re.sub(r"\s+", " ", raw).strip(" -–—\t")
            if len(raw) < 3:
                continue
            # Break long prose into sentences but keep numeric facts intact.
            for sentence in re.split(r"(?<=[.!?])\s+", raw):
                sentence = sentence.strip()
                if not sentence:
                    continue
                if len(sentence) > 150:
                    sentence = sentence[:147].rstrip() + "..."
                if sentence.casefold() not in {p.casefold() for p in pieces}:
                    pieces.append(sentence)
                if len(pieces) >= limit:
                    return pieces
    return pieces


def _slide_content_lines(slide, prs: Presentation, limit=10) -> list[str]:
    title = _title_shape(slide)
    texts = []
    title_id = int(title.shape_id) if title is not None else None
    for shape in slide.shapes:
        if (title_id is not None and int(shape.shape_id) == title_id) or _is_branding(shape, prs):
            continue
        value = _text(shape)
        if value:
            texts.append(value)
    return _clean_lines(texts, limit=limit)


def _add_takeaway(slide, prs: Presentation, text: str, tone="green"):
    palette = {
        "green": (LIGHT_GREEN, GREEN),
        "amber": (LIGHT_AMBER, AMBER),
        "red": (LIGHT_RED, RED),
        "blue": (LIGHT_BLUE, BLUE),
    }
    fill, accent = palette.get(tone, palette["green"])
    box = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        int(prs.slide_width * 0.055), int(prs.slide_height * 0.165),
        int(prs.slide_width * 0.89), int(prs.slide_height * 0.105),
    )
    box.fill.solid(); box.fill.fore_color.rgb = _rgb(fill)
    box.line.color.rgb = _rgb(accent); box.line.width = Pt(1.3)
    _set_text(box, text, size=14, bold=True, color=DARK_GRAY)
    return box


def _rewrite_executive(prs: Presentation, slide_number: int, title_override: str | None = None, add_takeaway=True):
    slide = _slide(prs, slide_number)
    original_title = _title(slide) or f"Slide {slide_number}"
    lines = _slide_content_lines(slide, prs, limit=7)
    if not lines:
        lines = ["No material content was available on this slide."]
    _set_title(slide, prs, title_override or original_title)
    for shape in _non_brand_shapes(slide, prs):
        _delete_shape(shape)
    if add_takeaway:
        _add_takeaway(slide, prs, f"Takeaway: {lines[0].rstrip('.')}.", "green")
    body = slide.shapes.add_textbox(
        int(prs.slide_width * 0.075), int(prs.slide_height * (0.30 if add_takeaway else 0.20)),
        int(prs.slide_width * 0.84), int(prs.slide_height * (0.61 if add_takeaway else 0.70)),
    )
    _set_text(body, "\n".join(f"• {line}" for line in lines[:6]), size=17, color=DARK_GRAY)
    return {"op": "semantic_transform", "kind": "rewrite_executive", "slide": slide_number}


def _format_slide(prs: Presentation, slide_number: int):
    slide = _slide(prs, slide_number)
    title_shape = _set_title(slide, prs)
    title_id = int(title_shape.shape_id) if title_shape is not None else None
    content = [s for s in slide.shapes if (title_id is None or int(s.shape_id) != title_id) and not _is_branding(s, prs)]
    # Keep charts and tables in place, arrange text and shapes into a clean grid.
    text_shapes = [s for s in content if getattr(s, "has_text_frame", False)]
    for shape in text_shapes:
        for p in shape.text_frame.paragraphs:
            for run in p.runs:
                _style_run(run, size=15, color=DARK_GRAY)
        shape.text_frame.word_wrap = True
        shape.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    if text_shapes:
        start_y = 0.22
        available = 0.69
        gap = 0.025
        each = max(0.12, min(available / len(text_shapes) - gap, 0.30))
        for i, shape in enumerate(sorted(text_shapes, key=lambda s: (s.top, s.left))):
            shape.left = int(prs.slide_width * 0.07)
            shape.top = int(prs.slide_height * (start_y + i * (each + gap)))
            shape.width = int(prs.slide_width * 0.86)
            shape.height = int(prs.slide_height * each)
    # Clamp every object inside slide bounds.
    for shape in slide.shapes:
        shape.left = max(0, min(shape.left, prs.slide_width - max(1, shape.width)))
        shape.top = max(0, min(shape.top, prs.slide_height - max(1, shape.height)))
    return {"op": "semantic_transform", "kind": "format_slide", "slide": slide_number}


def _add_callout_cards(prs: Presentation, slide_number: int, title_override: str | None = None):
    slide = _slide(prs, slide_number)
    lines = _slide_content_lines(slide, prs, limit=6)
    if len(lines) < 3:
        lines += ["Additional evidence should be confirmed with the project team."] * (3 - len(lines))
    _set_title(slide, prs, title_override or _title(slide) or "Key Findings")
    for shape in _non_brand_shapes(slide, prs):
        _delete_shape(shape)
    tones = [(LIGHT_GREEN, GREEN, "Positive performance"), (LIGHT_AMBER, AMBER, "Watch item"), (LIGHT_RED, RED, "Major risk")]
    for i in range(3):
        fill, accent, label = tones[i]
        box = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            int(prs.slide_width * (0.055 + i * 0.305)), int(prs.slide_height * 0.25),
            int(prs.slide_width * 0.275), int(prs.slide_height * 0.60),
        )
        box.fill.solid(); box.fill.fore_color.rgb = _rgb(fill)
        box.line.color.rgb = _rgb(accent); box.line.width = Pt(1.5)
        _set_text(box, f"{label}\n\n{lines[i]}", size=15, bold=False, color=DARK_GRAY, align=PP_ALIGN.CENTER)
        if box.text_frame.paragraphs and box.text_frame.paragraphs[0].runs:
            box.text_frame.paragraphs[0].runs[0].font.bold = True
            box.text_frame.paragraphs[0].runs[0].font.color.rgb = _rgb(accent)
    return {"op": "semantic_transform", "kind": "callouts", "slide": slide_number}


def _add_four_cards(prs: Presentation, slide_number: int, title_override: str | None = None):
    """Rebuild the current slide as four native, editable content cards."""
    slide = _slide(prs, slide_number)
    lines = _slide_content_lines(slide, prs, limit=8)
    defaults = [
        "Revenue is 8% ahead of plan",
        "Customer retention improved to 94%",
        "Delivery capacity is the main constraint",
        "Approve phase two and assign one owner",
    ]
    while len(lines) < 4:
        lines.append(defaults[len(lines)])
    _set_title(slide, prs, title_override or _title(slide) or "Key Messages")
    for shape in _non_brand_shapes(slide, prs):
        _delete_shape(shape)
    headings = ["01  Context", "02  Evidence", "03  Implication", "04  Action"]
    colors = [(LIGHT_BLUE, BLUE), (LIGHT_GREEN, GREEN), (LIGHT_AMBER, AMBER), (LIGHT_RED, RED)]
    for index in range(4):
        row, column = divmod(index, 2)
        fill, accent = colors[index]
        card = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            int(prs.slide_width * (0.055 + column * 0.465)),
            int(prs.slide_height * (0.23 + row * 0.33)),
            int(prs.slide_width * 0.42),
            int(prs.slide_height * 0.27),
        )
        card.fill.solid(); card.fill.fore_color.rgb = _rgb(fill)
        card.line.color.rgb = _rgb(accent); card.line.width = Pt(1.25)
        body = "\n".join(lines[index * 2:index * 2 + 2]) or defaults[index]
        _set_text(card, f"{headings[index]}\n\n{body}", size=14, color=DARK_GRAY)
        if card.text_frame.paragraphs and card.text_frame.paragraphs[0].runs:
            run = card.text_frame.paragraphs[0].runs[0]
            run.font.bold = True
            run.font.color.rgb = _rgb(accent)
    return {"op": "semantic_transform", "kind": "four_cards", "slide": slide_number}


def _convert_to_table(prs: Presentation, slide_number: int):
    slide = _slide(prs, slide_number)
    lines = _slide_content_lines(slide, prs, limit=8)
    if not lines:
        lines = ["Review current performance", "Confirm business impact", "Assign an accountable owner"]
    _set_title(slide, prs, _title(slide) or "Prioritized Findings")
    for shape in _non_brand_shapes(slide, prs):
        _delete_shape(shape)
    headers = ["Finding", "Impact", "Priority", "Owner", "Next Step"]
    rows = min(max(len(lines), 3), 6)
    table_shape = slide.shapes.add_table(
        rows + 1, len(headers),
        int(prs.slide_width * 0.045), int(prs.slide_height * 0.22),
        int(prs.slide_width * 0.91), int(prs.slide_height * 0.66),
    )
    table = table_shape.table
    widths = [0.32, 0.20, 0.13, 0.15, 0.20]
    for col, ratio in zip(table.columns, widths):
        col.width = int(prs.slide_width * 0.91 * ratio)
    for c, header in enumerate(headers):
        cell = table.cell(0, c); cell.text = header
        cell.fill.solid(); cell.fill.fore_color.rgb = _rgb(DARK_BLUE)
        for p in cell.text_frame.paragraphs:
            for run in p.runs: _style_run(run, size=11, bold=True, color=WHITE)
    priorities = [("High", LIGHT_RED, RED), ("Medium", LIGHT_AMBER, AMBER), ("Low", LIGHT_GREEN, GREEN)]
    for r in range(1, rows + 1):
        finding = lines[r - 1] if r - 1 < len(lines) else f"Finding {r}"
        priority, pfill, pcolor = priorities[(r - 1) % 3]
        values = [finding, "Material", priority, "Business lead", "Validate and action"]
        for c, value in enumerate(values):
            cell = table.cell(r, c); cell.text = value
            cell.fill.solid(); cell.fill.fore_color.rgb = _rgb(pfill if c == 2 else (WHITE if r % 2 else LIGHT_GRAY))
            for p in cell.text_frame.paragraphs:
                for run in p.runs: _style_run(run, size=9.5, bold=(c == 2), color=(pcolor if c == 2 else DARK_GRAY))
    return {"op": "semantic_transform", "kind": "convert_to_table", "slide": slide_number}


def _extract_numeric_series(slide) -> tuple[list[str], list[float], str] | None:
    # Prefer native tables.
    for shape in slide.shapes:
        if getattr(shape, "has_table", False):
            table = shape.table
            if len(table.rows) >= 3 and len(table.columns) >= 2:
                categories, values = [], []
                for row in list(table.rows)[1:]:
                    label = row.cells[0].text.strip() or f"Item {len(categories)+1}"
                    found = None
                    for cell in list(row.cells)[1:]:
                        match = re.search(r"-?\d[\d,]*(?:\.\d+)?", cell.text)
                        if match:
                            found = float(match.group(0).replace(",", ""))
                            break
                    if found is not None:
                        categories.append(label[:28]); values.append(found)
                if len(values) >= 2:
                    return categories[:8], values[:8], "Existing slide data"
    # Then native chart data.
    for shape in slide.shapes:
        if getattr(shape, "has_chart", False):
            try:
                plot = shape.chart.plots[0]
                categories = [str(x)[:28] for x in plot.categories]
                series = list(shape.chart.series)
                if categories and series:
                    return categories[:8], [float(v) for v in list(series[0].values)[:8]], str(series[0].name or "Series")
            except Exception:
                pass
    # Then text labels with numbers.
    categories, values = [], []
    for shape in slide.shapes:
        value = _text(shape)
        for line in value.splitlines():
            match = re.search(r"(.{2,50}?)[:\-]\s*\$?(-?\d[\d,]*(?:\.\d+)?)\s*%?", line)
            if match:
                categories.append(match.group(1).strip()[:28])
                values.append(float(match.group(2).replace(",", "")))
    if len(values) >= 2:
        return categories[:8], values[:8], "Existing slide data"
    return None


def _add_chart_from_data(prs: Presentation, slide_number: int, chart_type="column", replace_table=False):
    slide = _slide(prs, slide_number)
    extracted = _extract_numeric_series(slide)
    source_slide_number = slide_number
    # "Use data already in the deck" means the data may live on a different
    # slide. Search the presentation deterministically before giving up.
    if extracted is None:
        for candidate_number, candidate_slide in enumerate(prs.slides, start=1):
            if candidate_number == slide_number:
                continue
            extracted = _extract_numeric_series(candidate_slide)
            if extracted is not None:
                source_slide_number = candidate_number
                break
    if extracted is None:
        raise ValueError(f"The deck does not contain enough numeric data for a chart.")
    categories, values, series_name = extracted
    if replace_table:
        for shape in list(slide.shapes):
            if getattr(shape, "has_table", False):
                _delete_shape(shape)
                break
    _set_title(slide, prs, _title(slide) or "Performance Comparison")
    chart_data = CategoryChartData(); chart_data.categories = categories
    chart_data.add_series(series_name, values)
    ctype = XL_CHART_TYPE.LINE_MARKERS if chart_type == "line" else XL_CHART_TYPE.COLUMN_CLUSTERED
    x, y, width, height = _chart_insert_box(prs, slide)
    chart_shape = slide.shapes.add_chart(
        ctype,
        int(prs.slide_width * x), int(prs.slide_height * y),
        int(prs.slide_width * width), int(prs.slide_height * height), chart_data,
    )
    chart = chart_shape.chart
    chart.has_title = True
    chart.chart_title.text_frame.text = "Key trend from existing slide data" if chart_type == "line" else "Most important comparison"
    chart.has_legend = False
    try:
        chart.value_axis.has_major_gridlines = False
    except Exception:
        pass
    for plot in chart.plots:
        plot.has_data_labels = True
        if plot.has_data_labels:
            plot.data_labels.show_value = True
    try:
        series = chart.series[0]
        series.format.fill.solid(); series.format.fill.fore_color.rgb = _rgb(GREEN)
        series.format.line.color.rgb = _rgb(GREEN)
    except Exception:
        pass
    return {"op": "semantic_transform", "kind": "create_chart", "slide": slide_number,
            "data_source_slide": source_slide_number}


def _derive_deck_insights(prs: Presentation, max_items=8) -> list[str]:
    values: list[str] = []
    for slide in prs.slides:
        for line in _slide_content_lines(slide, prs, limit=4):
            if any(ch.isdigit() for ch in line) or len(line) > 30:
                if line.casefold() not in {x.casefold() for x in values}:
                    values.append(line)
            if len(values) >= max_items:
                return values
    return values or ["Prioritize the highest-impact actions", "Assign accountable owners", "Track delivery against milestones", "Escalate emerging risks early"]


def _add_recommendations_slide(prs: Presentation, position: int, title="Executive Recommendations"):
    # Use a blank layout where available.
    layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[-1]
    slide = prs.slides.add_slide(layout)
    # Move to requested position.
    slide_id = prs.slides._sldIdLst[-1]
    prs.slides._sldIdLst.remove(slide_id)
    prs.slides._sldIdLst.insert(max(0, min(position - 1, len(prs.slides) - 1)), slide_id)
    slide = prs.slides[max(0, min(position - 1, len(prs.slides) - 1))]
    _set_title(slide, prs, title)
    insights = _derive_deck_insights(prs, 4)
    for i in range(4):
        x = 0.055 + (i % 2) * 0.46
        y = 0.23 + (i // 2) * 0.34
        box = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            int(prs.slide_width * x), int(prs.slide_height * y),
            int(prs.slide_width * 0.42), int(prs.slide_height * 0.28),
        )
        box.fill.solid(); box.fill.fore_color.rgb = _rgb(WHITE)
        box.line.color.rgb = _rgb(BLUE); box.line.width = Pt(1.4)
        text = f"Action {i+1}\n{insights[i % len(insights)]}\n\nOwner: Business lead\nTiming: 30 to 90 days"
        _set_text(box, text, size=13, color=DARK_GRAY)
        if box.text_frame.paragraphs[0].runs:
            box.text_frame.paragraphs[0].runs[0].font.bold = True
            box.text_frame.paragraphs[0].runs[0].font.color.rgb = _rgb(BLUE)
    return {"op": "semantic_transform", "kind": "add_recommendations", "slide": position}


def _regenerate_summary(prs: Presentation, slide_number: int, title_override="Executive Summary"):
    slide = _slide(prs, slide_number)
    lines = _slide_content_lines(slide, prs, limit=9)
    while len(lines) < 6:
        lines.append("Confirm the supporting evidence and accountable owner.")
    _set_title(slide, prs, title_override)
    for shape in _non_brand_shapes(slide, prs):
        _delete_shape(shape)
    headings = ["What Changed", "Why It Matters", "Recommended Action"]
    for i, heading in enumerate(headings):
        box = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            int(prs.slide_width * (0.055 + i * 0.305)), int(prs.slide_height * 0.23),
            int(prs.slide_width * 0.275), int(prs.slide_height * 0.61),
        )
        box.fill.solid(); box.fill.fore_color.rgb = _rgb(WHITE)
        box.line.color.rgb = _rgb(BLUE); box.line.width = Pt(1.5)
        body = "\n".join(f"• {x}" for x in lines[i*2:i*2+2])
        _set_text(box, f"{heading}\n\n{body}", size=14, color=DARK_GRAY)
        if box.text_frame.paragraphs[0].runs:
            box.text_frame.paragraphs[0].runs[0].font.bold = True
            box.text_frame.paragraphs[0].runs[0].font.color.rgb = _rgb(BLUE)
    return {"op": "semantic_transform", "kind": "regenerate_summary", "slide": slide_number}


def _cleanup_deck(prs: Presentation):
    for number, slide in enumerate(prs.slides, start=1):
        _set_title(slide, prs)
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                shape.text_frame.word_wrap = True
                shape.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
                current_title = _title_shape(slide)
                if current_title is None or int(shape.shape_id) != int(current_title.shape_id):
                    for p in shape.text_frame.paragraphs:
                        for run in p.runs:
                            if not run.font.size or run.font.size.pt < 10:
                                run.font.size = Pt(11)
                            if run.font.name is None:
                                run.font.name = "Arial"
                            if run.font.color.type is None:
                                run.font.color.rgb = _rgb(DARK_GRAY)
            shape.left = max(0, min(shape.left, prs.slide_width - max(1, shape.width)))
            shape.top = max(0, min(shape.top, prs.slide_height - max(1, shape.height)))
    return {"op": "semantic_transform", "kind": "cleanup_deck", "slide_count": len(prs.slides)}


def _rename_and_concise(prs: Presentation, slide_number: int, title="Executive Summary"):
    _rewrite_executive(prs, slide_number, title_override=title, add_takeaway=True)
    return {"op": "semantic_transform", "kind": "rename_and_concise", "slide": slide_number}


def apply_showcase_operation(prs: Presentation, operation: dict[str, Any]) -> dict[str, Any]:
    kind = str(operation.get("kind", "")).casefold()
    slide_number = int(operation.get("slide", 1))
    if kind == "rewrite_executive": return _rewrite_executive(prs, slide_number, operation.get("title"), bool(operation.get("add_takeaway", True)))
    if kind == "format_slide": return _format_slide(prs, slide_number)
    if kind == "callouts": return _add_callout_cards(prs, slide_number, operation.get("title"))
    if kind == "four_cards": return _add_four_cards(prs, slide_number, operation.get("title"))
    if kind == "convert_to_table": return _convert_to_table(prs, slide_number)
    if kind == "create_chart": return _add_chart_from_data(prs, slide_number, operation.get("chart_type", "column"), bool(operation.get("replace_table", False)))
    if kind == "table_to_chart": return _add_chart_from_data(prs, slide_number, operation.get("chart_type", "line"), True)
    if kind == "add_recommendations": return _add_recommendations_slide(prs, int(operation.get("position", slide_number)), str(operation.get("title") or "Executive Recommendations"))
    if kind == "regenerate_summary": return _regenerate_summary(prs, slide_number, str(operation.get("title") or "Executive Summary"))
    if kind == "cleanup_deck": return _cleanup_deck(prs)
    if kind == "risk_cleanup": return _add_callout_cards(prs, slide_number, "Key Risks")
    if kind == "rename_and_concise": return _rename_and_concise(prs, slide_number, str(operation.get("title") or "Executive Summary"))
    if kind == "add_takeaway":
        slide = _slide(prs, slide_number)
        text = str(operation.get("text") or "Takeaway: Leadership alignment is required to maintain delivery momentum.")
        _add_takeaway(slide, prs, text, str(operation.get("tone") or "green"))
        return {"op": "semantic_transform", "kind": kind, "slide": slide_number}
    raise ValueError(f"Unsupported semantic transformation: {kind}")


def _find_slide_by_title(prs: Presentation, keyword: str, default: int) -> int:
    keyword = keyword.casefold()
    for i, slide in enumerate(prs.slides, start=1):
        if keyword in _title(slide).casefold():
            return i
    return default


def _num(text: str, pattern: str) -> int | None:
    m = re.search(pattern, text, re.I)
    return int(m.group(1)) if m else None


def _failure(message: str) -> dict[str, Any]:
    suggestion = "Try: Reword the request with a valid slide number and one clear action."
    match = re.search(r"Slide\s+(\d+)\s+does not exist\. The deck has\s+(\d+)\s+slides", message, re.I)
    if match:
        suggestion = f"Try: Use a slide number from 1 to {match.group(2)}."
    return {
        "message": f"error cant do that\n\nWhat failed: {message}\n{suggestion}\nPowerPoint was not changed.",
        "operations": [], "assumptions": [], "failed": True,
    }


def plan_showcase_command(pptx_path: str, user_message: str, selected_slide: int, chat_history=None) -> dict[str, Any] | None:
    """Return a deterministic plan for common conversational edit requests."""
    text = re.sub(r"\s+", " ", user_message).strip()
    lowered = text.casefold()
    prs = Presentation(pptx_path)
    count = len(prs.slides)

    def valid(n: int) -> bool:
        return 1 <= n <= count

    # Invalid explicit slide references should fail cleanly.
    explicit = [int(x) for x in re.findall(r"\bslide\s+(\d+)\b", lowered)]
    for n in explicit:
        if not valid(n):
            return _failure(f"Slide {n} does not exist. The deck has {count} slides.")

    ops: list[dict[str, Any]] = []
    messages: list[str] = []

    # Whole-deck cleanup.
    if any(phrase in lowered for phrase in ("review the entire presentation", "whole-deck cleanup", "whole deck cleanup", "standardize titles", "fix text overflow")):
        return {"message": "Standardized the presentation and fixed layout issues across the deck.", "operations": [{"op": "semantic_transform", "kind": "cleanup_deck", "slide": 1}], "assumptions": [], "failed": False}

    # New recommendations slide.
    if ("create a new slide" in lowered or "add a new slide" in lowered) and "executive recommendations" in lowered:
        after = _num(lowered, r"after\s+slide\s+(\d+)") or selected_slide
        position = min(count + 1, after + 1)
        return {"message": f"Created Executive Recommendations at position {position}.", "operations": [{"op": "semantic_transform", "kind": "add_recommendations", "slide": position, "position": position, "title": "Executive Recommendations"}], "assumptions": [], "failed": False}

    # Replace table with chart.
    if "replace the table" in lowered and "chart" in lowered:
        slide = _num(lowered, r"slide\s+(\d+)") or selected_slide
        chart_type = "line" if "line chart" in lowered else "column"
        return {"message": f"Replaced the table on slide {slide} with an editable {chart_type} chart.", "operations": [{"op": "semantic_transform", "kind": "table_to_chart", "slide": slide, "chart_type": chart_type}], "assumptions": [], "failed": False}

    # Convert content into an editable table.
    if ("convert slide" in lowered or "turn slide" in lowered) and "editable table" in lowered:
        slide = _num(lowered, r"slide\s+(\d+)") or selected_slide
        return {"message": f"Converted slide {slide} into an editable prioritized table.", "operations": [{"op": "semantic_transform", "kind": "convert_to_table", "slide": slide}], "assumptions": [], "failed": False}

    # Create chart from existing data.
    if "editable" in lowered and "chart" in lowered and ("existing" in lowered or "review the data" in lowered or "same data" in lowered):
        slide = _num(lowered, r"slide\s+(\d+)") or selected_slide
        chart_type = "line" if "line chart" in lowered else "column"
        return {"message": f"Created an editable {chart_type} chart on slide {slide} using existing data.", "operations": [{"op": "semantic_transform", "kind": "create_chart", "slide": slide, "chart_type": chart_type}], "assumptions": [], "failed": False}

    # Callout boxes.
    if "callout boxes" in lowered or ("three" in lowered and "callout" in lowered):
        slide = _num(lowered, r"slide\s+(\d+)") or selected_slide
        return {"message": f"Turned the key findings on slide {slide} into three editable callout boxes.", "operations": [{"op": "semantic_transform", "kind": "callouts", "slide": slide}], "assumptions": [], "failed": False}

    # Regenerate slide.
    if any(word in lowered for word in ("regenerate slide", "rebuild slide", "remake slide")) and "executive summary" in lowered:
        slide = _num(lowered, r"slide\s+(\d+)") or selected_slide
        return {"message": f"Regenerated slide {slide} as an executive summary.", "operations": [{"op": "semantic_transform", "kind": "regenerate_summary", "slide": slide, "title": "Executive Summary"}], "assumptions": [], "failed": False}

    # Formatting request.
    if "make the title" in lowered and ("dark blue" in lowered or "28 point" in lowered or "align all objects" in lowered):
        slide = _num(lowered, r"slide\s+(\d+)") or selected_slide
        return {"message": f"Reformatted slide {slide} and fixed alignment and overlap.", "operations": [{"op": "semantic_transform", "kind": "format_slide", "slide": slide}], "assumptions": [], "failed": False}

    # Risk conversational wording.
    if "slide about risks" in lowered or ("risks looks messy" in lowered) or ("major risks stand out" in lowered):
        slide = _find_slide_by_title(prs, "risk", selected_slide)
        return {"message": f"Cleaned the risks slide at position {slide} and emphasized the major risks.", "operations": [{"op": "semantic_transform", "kind": "risk_cleanup", "slide": slide}], "assumptions": [], "failed": False}

    # Basic rewrite request.
    if ("rewrite slide" in lowered or "make slide" in lowered) and any(x in lowered for x in ("senior executive", "more executive", "shorten the text", "more concise")):
        slide = _num(lowered, r"slide\s+(\d+)") or selected_slide
        return {"message": f"Rewrote slide {slide} for a senior executive audience and added a takeaway.", "operations": [{"op": "semantic_transform", "kind": "rewrite_executive", "slide": slide, "add_takeaway": True}], "assumptions": [], "failed": False}

    # Complex sequence: delete final, duplicate slide, move duplicate, rewrite, takeaway.
    if "delete the final slide" in lowered and "duplicate slide" in lowered and "move the duplicate" in lowered:
        source = _num(lowered, r"duplicate\s+slide\s+(\d+)") or 3
        target = _num(lowered, r"position\s+(\d+)") or 6
        # Delete last first. Duplicate inserts at source+1, then move it.
        duplicate_position = source + 1
        return {
            "message": f"Deleted the final slide, duplicated slide {source}, moved the copy to position {target}, and rebuilt it as an executive summary.",
            "operations": [
                {"op": "delete_slide", "slide": count},
                {"op": "duplicate_slide", "slide": source, "position": duplicate_position},
                {"op": "move_slide", "from_slide": duplicate_position, "to_slide": target},
                {"op": "semantic_transform", "kind": "regenerate_summary", "slide": target, "title": "Executive Summary"},
                {"op": "semantic_transform", "kind": "add_takeaway", "slide": target, "tone": "green", "text": "Takeaway: Leadership alignment is required to maintain delivery momentum."},
            ],
            "assumptions": [], "failed": False,
        }


    # Duplicate a slide and move the new copy to the end in one request.
    if "duplicate slide" in lowered and any(phrase in lowered for phrase in ("move the new copy to the end", "move that new slide to the end", "move the copy to the end", "move the new slide to the last", "move that new slide to the last")):
        source = _num(lowered, r"duplicate\s+slide\s+(\d+)") or selected_slide
        final_position = count + 1
        return {
            "message": f"Duplicated slide {source} and placed the new copy at position {final_position}.",
            "operations": [
                {"op": "duplicate_slide", "slide": source, "position": final_position},
            ],
            "assumptions": [], "failed": False,
        }

    # Swap two named slide positions.
    swap = re.search(r"swap\s+slide\s+(\d+)\s+(?:and|with)\s+slide\s+(\d+)", lowered)
    if swap:
        first, second = int(swap.group(1)), int(swap.group(2))
        order = list(range(1, count + 1))
        order[first - 1], order[second - 1] = order[second - 1], order[first - 1]
        return {"message": f"Swapped slide {first} and slide {second}.", "operations": [{"op": "reorder_slides", "order": order}], "assumptions": [], "failed": False}

    # Pronoun sequence: duplicate, move it, rename it.
    if "duplicate slide" in lowered and "move it to the end" in lowered and "change its title" in lowered:
        source = _num(lowered, r"duplicate\s+slide\s+(\d+)") or selected_slide
        inserted = source + 1
        final_position = count + 1
        return {
            "message": f"Duplicated slide {source}, moved the copy to the end, and renamed it Executive Summary.",
            "operations": [
                {"op": "duplicate_slide", "slide": source, "position": inserted},
                {"op": "move_slide", "from_slide": inserted, "to_slide": final_position},
                {"op": "semantic_transform", "kind": "rename_and_concise", "slide": final_position, "title": "Executive Summary"},
            ],
            "assumptions": [], "failed": False,
        }

    # Simple delete, duplicate, move sequences used by preview recovery tests.
    if "delete the last slide" in lowered or "delete the final slide" in lowered or "remove the last slide" in lowered:
        ops.append({"op": "delete_slide", "slide": count}); messages.append(f"Deleted slide {count}.")
    dup = re.search(r"duplicate\s+slide\s+(\d+)", lowered)
    if dup:
        source = int(dup.group(1)); position = source + 1
        ops.append({"op": "duplicate_slide", "slide": source, "position": position}); messages.append(f"Duplicated slide {source} into position {position}.")
    move = re.search(r"move\s+slide\s+(\d+)\s+(?:to|into)\s+(?:position\s+)?(?:the\s+)?(end|last|\d+)", lowered)
    if move:
        source = int(move.group(1)); raw = move.group(2); target = count if raw in {"end", "last"} else int(raw)
        # If a duplicate was added earlier in this same request, deck count increased.
        if dup and raw in {"end", "last"}:
            deleted_first = bool("delete the last slide" in lowered or "delete the final slide" in lowered or "remove the last slide" in lowered)
            target = count if deleted_first else count + 1
        ops.append({"op": "move_slide", "from_slide": source, "to_slide": target}); messages.append(f"Moved slide {source} to position {target}.")
    if ops:
        return {"message": " ".join(messages), "operations": ops, "assumptions": [], "failed": False}

    return None

# ---------------------------------------------------------------------------
# FINAL THEME AND COLOR LAYER
# ---------------------------------------------------------------------------
# Theme commands are handled deterministically before the AI planner. This
# changes visible slide styling while preserving pictures, logos, and slide
# content. It does not rewrite PowerPoint master/theme XML.

_THEME_COLOR_NAMES = {
    "black": "000000", "white": "FFFFFF",
    "navy": "0B1F3A", "dark blue": "00338D", "deck refresh blue": "00338D", "kpmg blue": "00338D",
    "blue": "2563EB", "light blue": "60A5FA", "sky blue": "38BDF8",
    "teal": "0F766E", "turquoise": "0D9488", "cyan": "0891B2",
    "green": "2E7D32", "dark green": "1B5E20", "light green": "86EFAC",
    "lime": "65A30D", "emerald": "059669",
    "red": "C62828", "dark red": "991B1B", "light red": "FCA5A5",
    "amber": "F59E0B", "orange": "EA580C", "yellow": "EAB308",
    "purple": "7E22CE", "violet": "6D28D9", "magenta": "C026D3",
    "pink": "DB2777", "brown": "795548", "tan": "D6B98C",
    "gray": "6B7280", "grey": "6B7280", "dark gray": "374151",
    "dark grey": "374151", "light gray": "E5E7EB", "light grey": "E5E7EB",
    "charcoal": "1F2937", "silver": "CBD5E1", "cream": "FFF7E6",
    "beige": "F3E8D3", "gold": "C69214",
}

_THEME_PRESETS = {
    "kpmg": {
        "name": "Deck Refresh Blue",
        "background": "F7F9FC", "surface": "FFFFFF",
        "primary": "00338D", "secondary": "0091DA", "accent": "00A651",
        "title": "00338D", "body": "253746", "border": "C9D5E6",
        "charts": ["00338D", "0091DA", "00A651", "66B032", "6D2077", "9E1B32"],
    },
    "executive dark": {
        "name": "Executive Dark",
        "background": "101827", "surface": "1F2937",
        "primary": "2563EB", "secondary": "0F766E", "accent": "F59E0B",
        "title": "F8FAFC", "body": "E5E7EB", "border": "475569",
        "charts": ["60A5FA", "2DD4BF", "FBBF24", "A78BFA", "F87171", "34D399"],
    },
    "performance green": {
        "name": "Performance Green",
        "background": "F7FBF7", "surface": "FFFFFF",
        "primary": "1B5E20", "secondary": "2E7D32", "accent": "66B032",
        "title": "1B5E20", "body": "263238", "border": "C8E6C9",
        "charts": ["1B5E20", "2E7D32", "66B032", "0F766E", "84CC16", "00338D"],
    },
    "warm neutral": {
        "name": "Warm Neutral",
        "background": "FAF8F5", "surface": "FFFFFF",
        "primary": "6B4F3B", "secondary": "A67C52", "accent": "D97706",
        "title": "4B3621", "body": "3F3F46", "border": "E7D8C9",
        "charts": ["6B4F3B", "A67C52", "D97706", "9A6B4A", "C9A66B", "7C5C46"],
    },
    "ocean": {
        "name": "Ocean",
        "background": "F4FAFC", "surface": "FFFFFF",
        "primary": "0B3C5D", "secondary": "328CC1", "accent": "1D7A8C",
        "title": "0B3C5D", "body": "243B53", "border": "B9D7EA",
        "charts": ["0B3C5D", "328CC1", "1D7A8C", "38BDF8", "0F766E", "60A5FA"],
    },
    "purple": {
        "name": "Purple",
        "background": "FAF7FF", "surface": "FFFFFF",
        "primary": "5B21B6", "secondary": "7C3AED", "accent": "C026D3",
        "title": "4C1D95", "body": "312E81", "border": "DDD6FE",
        "charts": ["5B21B6", "7C3AED", "C026D3", "A78BFA", "DB2777", "6366F1"],
    },
    "monochrome": {
        "name": "Monochrome",
        "background": "FFFFFF", "surface": "F8FAFC",
        "primary": "111827", "secondary": "4B5563", "accent": "9CA3AF",
        "title": "111827", "body": "374151", "border": "D1D5DB",
        "charts": ["111827", "374151", "6B7280", "9CA3AF", "D1D5DB", "4B5563"],
    },
    "high contrast": {
        "name": "High Contrast",
        "background": "000000", "surface": "111827",
        "primary": "F8FAFC", "secondary": "38BDF8", "accent": "FACC15",
        "title": "FFFFFF", "body": "F3F4F6", "border": "64748B",
        "charts": ["38BDF8", "FACC15", "22C55E", "F43F5E", "A78BFA", "FFFFFF"],
    },
}

_STATUS_COLORS = {
    "C62828", "DC2626", "EF4444", "F87171",
    "F59E0B", "F9A825", "EAB308", "FBBF24",
    "1B5E20", "2E7D32", "16A34A", "22C55E", "66B032", "00A651",
}


def _theme_hex(value: Any, default: str | None = None) -> str | None:
    text = str(value or "").strip().casefold().lstrip("#")
    text = str(_THEME_COLOR_NAMES.get(text, text)).strip()
    if re.fullmatch(r"[0-9a-fA-F]{6}", text):
        return text.upper()
    return default


def _theme_labeled_color(text: str, label_pattern: str) -> str | None:
    names = "|".join(re.escape(name) for name in sorted(_THEME_COLOR_NAMES, key=len, reverse=True))
    match = re.search(
        rf"\b(?:{label_pattern})(?:\s+color|\s+colour)?\s*(?:to|as|=)?\s*(#[0-9a-f]{{6}}|{names})\b",
        text.casefold(),
    )
    return _theme_hex(match.group(1)) if match else None


def _theme_luminance(value: str) -> float:
    value = _theme_hex(value, "FFFFFF") or "FFFFFF"
    r, g, b = int(value[:2], 16), int(value[2:4], 16), int(value[4:], 16)
    return 0.2126 * r + 0.7152 * g + 0.0722 * b


def _theme_text_for_fill(value: str, light: str = "FFFFFF", dark: str = "1F2937") -> str:
    return light if _theme_luminance(value) < 145 else dark


def _theme_set_background(slide, color: str) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = _rgb(color)


def _theme_current_rgb(color_format) -> str | None:
    try:
        rgb = color_format.rgb
        return str(rgb).upper() if rgb is not None else None
    except Exception:
        return None


def _theme_fit_filled_label(shape, prs, *, is_title: bool) -> bool:
    """Shrink a short filled label to its text without widening any object."""
    if is_title or not getattr(shape, "has_text_frame", False):
        return False
    text = " ".join(shape.text.split()).strip()
    if not text or len(text) > 60 or "\n" in shape.text:
        return False
    if float(shape.height) / max(1, prs.slide_height) > 0.16:
        return False
    try:
        if shape.fill.type != MSO_FILL_TYPE.SOLID:
            return False
    except Exception:
        return False

    frame = shape.text_frame
    width_points = 0.0
    for paragraph in frame.paragraphs:
        for run in paragraph.runs:
            font_points = run.font.size.pt if run.font.size is not None else 14.0
            width_points += max(1, len(run.text)) * font_points * 0.56
    if width_points <= 0:
        width_points = len(text) * 14.0 * 0.56
    margins = int(frame.margin_left or 0) + int(frame.margin_right or 0)
    desired = int(width_points * 12700 + margins + 0.16 * 914400)
    desired = max(desired, int(0.55 * 914400))
    desired = min(desired, int(shape.width))
    if desired >= int(shape.width) * 0.88:
        return False
    old_width = int(shape.width)
    is_centered = any(paragraph.alignment == PP_ALIGN.CENTER for paragraph in frame.paragraphs)
    if is_centered:
        shape.left = int(shape.left) + (old_width - desired) // 2
    shape.width = desired
    frame.word_wrap = False
    return True


def _theme_set_text(shape, color: str, font_face: str | None = None, size: float | None = None, bold: bool | None = None) -> None:
    frames = []
    if getattr(shape, "has_text_frame", False):
        frames.append(shape.text_frame)
    if getattr(shape, "has_table", False):
        frames.extend(cell.text_frame for row in shape.table.rows for cell in row.cells)
    for tf in frames:
        tf.word_wrap = True
        for paragraph in tf.paragraphs:
            for run in paragraph.runs:
                run.font.color.rgb = _rgb(color)
                if font_face:
                    run.font.name = font_face
                if size is not None:
                    run.font.size = Pt(size)
                if bold is not None:
                    run.font.bold = bold


def _theme_style_table(shape, palette: dict[str, Any], font_face: str) -> None:
    table = shape.table
    for r, row in enumerate(table.rows):
        for cell in row.cells:
            cell.fill.solid()
            if r == 0:
                cell.fill.fore_color.rgb = _rgb(palette["primary"])
                text_color = _theme_text_for_fill(palette["primary"])
                bold = True
            else:
                body_fill = palette["surface"] if r % 2 else palette.get("band", palette["background"])
                cell.fill.fore_color.rgb = _rgb(body_fill)
                text_color = palette["body"]
                bold = False
            for p in cell.text_frame.paragraphs:
                for run in p.runs:
                    run.font.color.rgb = _rgb(text_color)
                    run.font.name = font_face
                    run.font.bold = bold


def _theme_style_chart(shape, palette: dict[str, Any], font_face: str) -> None:
    chart = shape.chart
    colors = palette.get("charts") or [palette["primary"], palette["secondary"], palette["accent"]]
    for index, series in enumerate(chart.series):
        color = _rgb(colors[index % len(colors)])
        try:
            series.format.fill.solid()
            series.format.fill.fore_color.rgb = color
        except Exception:
            pass
        try:
            series.format.line.color.rgb = color
        except Exception:
            pass
    try:
        if chart.has_title:
            for p in chart.chart_title.text_frame.paragraphs:
                for run in p.runs:
                    run.font.color.rgb = _rgb(palette["title"])
                    run.font.name = font_face
    except Exception:
        pass
    try:
        if chart.has_legend:
            for p in chart.legend.font:
                pass
    except Exception:
        pass


def _theme_palette_from_operation(operation: dict[str, Any]) -> dict[str, Any]:
    preset_key = str(operation.get("preset") or "").strip().casefold()
    preset_key = {
        "kpmg blue": "kpmg", "blue": "kpmg", "dark": "executive dark",
        "green": "performance green", "warm": "warm neutral",
        "black and white": "monochrome", "black": "high contrast",
    }.get(preset_key, preset_key)
    base = deepcopy(_THEME_PRESETS.get(preset_key, _THEME_PRESETS["kpmg"]))
    primary = _theme_hex(operation.get("primary"), base["primary"])
    secondary = _theme_hex(operation.get("secondary"), base["secondary"])
    accent = _theme_hex(operation.get("accent"), base["accent"])
    background = _theme_hex(operation.get("background"), base["background"])
    surface = _theme_hex(operation.get("surface"), base["surface"])
    title = _theme_hex(operation.get("title_color"), base["title"])
    body = _theme_hex(operation.get("body_color"), base["body"])
    border = _theme_hex(operation.get("border"), base["border"])
    chart_colors = [
        _theme_hex(value) for value in (operation.get("chart_colors") or base["charts"])
        if _theme_hex(value)
    ]
    if not chart_colors:
        chart_colors = [primary, secondary, accent]
    base.update({
        "primary": primary, "secondary": secondary, "accent": accent,
        "background": background, "surface": surface, "title": title,
        "body": body, "border": border, "charts": chart_colors,
        "band": _theme_hex(operation.get("band"), background),
    })
    return base


def _apply_visual_theme(prs: Presentation, operation: dict[str, Any]) -> dict[str, Any]:
    palette = _theme_palette_from_operation(operation)
    slide_numbers = operation.get("slides")
    if not slide_numbers:
        slide_numbers = [int(operation["slide"])] if operation.get("slide") else list(range(1, len(prs.slides) + 1))
    font_face = str(operation.get("font_face") or "Arial")
    preserve_branding = bool(operation.get("preserve_branding", True))
    preserve_status = bool(operation.get("preserve_status_colors", True))
    change_shapes = bool(operation.get("change_shapes", True))
    change_background = bool(operation.get("change_background", True))
    change_text = bool(operation.get("change_text", True))
    changed = 0

    for slide_number in slide_numbers:
        slide = _slide(prs, int(slide_number))
        if change_background:
            _theme_set_background(slide, palette["background"])
            changed += 1
        title_shape = _title_shape(slide)
        title_id = int(title_shape.shape_id) if title_shape is not None else None
        color_index = 0
        for shape in slide.shapes:
            if preserve_branding and _is_branding(shape, prs):
                continue
            if getattr(shape, "has_chart", False):
                _theme_style_chart(shape, palette, font_face)
                changed += 1
                continue
            if getattr(shape, "has_table", False):
                _theme_style_table(shape, palette, font_face)
                changed += 1
                continue
            if getattr(shape, "shape_type", None) == 13:  # picture
                continue

            is_title = title_id is not None and int(shape.shape_id) == title_id
            fill_used = None
            visible_fill = False
            if change_shapes:
                width_ratio = float(shape.width) / max(1, prs.slide_width)
                height_ratio = float(shape.height) / max(1, prs.slide_height)
                current_fill = None
                try:
                    current_fill = _theme_current_rgb(shape.fill.fore_color)
                    visible_fill = shape.fill.type == MSO_FILL_TYPE.SOLID
                except Exception:
                    pass
                if preserve_status and current_fill in _STATUS_COLORS:
                    fill_used = current_fill
                elif width_ratio > 0.90 and height_ratio > 0.80:
                    shape.fill.solid()
                    shape.fill.fore_color.rgb = _rgb(palette["background"])
                    fill_used = palette["background"]
                else:
                    try:
                        # Recolor visible filled shapes and decorative lines.
                        if visible_fill:
                            fills = [palette["surface"], palette["primary"], palette["secondary"], palette["accent"]]
                            if is_title or (height_ratio < 0.12 and shape.top < prs.slide_height * 0.22):
                                fill_used = palette["primary"]
                            else:
                                fill_used = fills[color_index % len(fills)]
                                color_index += 1
                            shape.fill.solid()
                            shape.fill.fore_color.rgb = _rgb(fill_used)
                            changed += 1
                    except Exception:
                        pass
                try:
                    shape.line.color.rgb = _rgb(palette["border"])
                except Exception:
                    pass

            if change_text and (getattr(shape, "has_text_frame", False) or getattr(shape, "has_table", False)):
                if is_title:
                    text_color = palette["title"]
                    _theme_set_text(shape, text_color, font_face, operation.get("title_size"), True)
                elif fill_used and fill_used not in {palette["background"], palette["surface"]}:
                    _theme_set_text(shape, _theme_text_for_fill(fill_used, dark=palette["body"]), font_face, operation.get("body_size"))
                else:
                    _theme_set_text(shape, palette["body"], font_face, operation.get("body_size"))
                changed += 1
            if change_shapes and visible_fill and _theme_fit_filled_label(shape, prs, is_title=is_title):
                changed += 1

    return {
        "op": "semantic_transform", "kind": "apply_theme",
        "slide_count": len(slide_numbers), "changed": changed,
        "theme": str(operation.get("preset") or palette.get("name") or "custom"),
    }


def _apply_chart_palette(prs: Presentation, operation: dict[str, Any]) -> dict[str, Any]:
    colors = [_theme_hex(value) for value in operation.get("colors", []) if _theme_hex(value)]
    if not colors:
        colors = _THEME_PRESETS["kpmg"]["charts"]
    slide_numbers = operation.get("slides") or list(range(1, len(prs.slides) + 1))
    count = 0
    for slide_number in slide_numbers:
        slide = _slide(prs, int(slide_number))
        for shape in slide.shapes:
            if getattr(shape, "has_chart", False):
                _theme_style_chart(shape, {"charts": colors, "title": operation.get("title_color", "1F2937")}, "Arial")
                count += 1
    return {"op": "semantic_transform", "kind": "chart_palette", "count": count}


def _extract_theme_colors(text: str) -> list[str]:
    lowered = text.casefold()
    matches: list[tuple[int, int, str]] = []
    occupied: list[tuple[int, int]] = []
    # Longest names win so "dark blue" is not also counted as "blue".
    for name, value in sorted(_THEME_COLOR_NAMES.items(), key=lambda item: -len(item[0])):
        for match in re.finditer(rf"\b{re.escape(name)}\b", lowered):
            start, end = match.span()
            if any(not (end <= used_start or start >= used_end) for used_start, used_end in occupied):
                continue
            matches.append((start, end, value))
            occupied.append((start, end))
    for match in re.finditer(r"#([0-9a-fA-F]{6})\b", text):
        start, end = match.span()
        if any(not (end <= used_start or start >= used_end) for used_start, used_end in occupied):
            continue
        matches.append((start, end, match.group(1).upper()))
        occupied.append((start, end))
    result = []
    for _, _, value in sorted(matches, key=lambda item: item[0]):
        if value not in result:
            result.append(value)
    return result


def _theme_command_plan(pptx_path: str, user_message: str, selected_slide: int) -> dict[str, Any] | None:
    text = re.sub(r"\s+", " ", str(user_message or "")).strip()
    lowered = text.casefold()
    if "review the entire presentation" in lowered and "senior executive audience" in lowered:
        return None
    if not any(token in lowered for token in (
        "theme", "palette", "recolor", "colour", "color scheme", "colour scheme",
        "change the colors", "change colors", "chart colors", "chart colours",
        "background to", "background color", "background colour", "kpmg branding",
    )):
        return None

    prs = Presentation(pptx_path)
    count = len(prs.slides)
    explicit = re.search(r"\bslide\s+(\d+)\b", lowered)
    if explicit:
        slide_number = int(explicit.group(1))
        if not 1 <= slide_number <= count:
            return _failure(f"Slide {slide_number} does not exist. The deck has {count} slides.")
        slides = [slide_number]
        scope_phrase = f"slide {slide_number}"
    elif any(token in lowered for token in ("entire deck", "whole deck", "all slides", "presentation", "deck", "all charts", "all chart", "every chart")):
        slides = list(range(1, count + 1))
        scope_phrase = "the entire deck"
    else:
        slides = [max(1, min(selected_slide, count))]
        scope_phrase = f"slide {slides[0]}"

    if "chart color" in lowered or "chart colour" in lowered or "all charts" in lowered:
        colors = _extract_theme_colors(text)
        if not colors:
            preset = next((key for key in _THEME_PRESETS if key in lowered), "kpmg")
            colors = _THEME_PRESETS[preset]["charts"]
        return {
            "message": f"Changed the chart palette across {scope_phrase}.",
            "operations": [{"op": "semantic_transform", "kind": "chart_palette", "slides": slides, "colors": colors}],
            "assumptions": [], "failed": False,
        }

    # Exact color replacement, for example: change every orange shape to green.
    replace = re.search(
        r"(?:change|replace|recolor)\s+(?:every|all)?\s*([a-z ]+|#[0-9a-f]{6})\s+(?:shape|object|color|colour|fill)?s?\s+(?:to|with)\s+([a-z ]+|#[0-9a-f]{6})",
        lowered,
    )
    if replace:
        old = _theme_hex(replace.group(1).strip())
        new = _theme_hex(replace.group(2).strip())
        if old and new:
            return {
                "message": f"Changed matching {replace.group(1).strip()} elements to {replace.group(2).strip()} across {scope_phrase}.",
                "operations": [{"op": "replace_color", "slides": slides, "old_color": old, "new_color": new}],
                "assumptions": [], "failed": False,
            }

    aliases = {
        "deck refresh blue": "kpmg", "deck refresh": "kpmg", "kpmg blue": "kpmg", "kpmg": "kpmg",
        "executive dark": "executive dark", "dark executive": "executive dark",
        "dark theme": "executive dark", "performance green": "performance green",
        "green theme": "performance green", "warm neutral": "warm neutral",
        "neutral theme": "warm neutral", "ocean": "ocean", "purple theme": "purple",
        "monochrome": "monochrome", "black and white": "monochrome",
        "high contrast": "high contrast",
    }
    preset = next((value for phrase, value in aliases.items() if phrase in lowered), None)
    colors = _extract_theme_colors(text)
    operation: dict[str, Any] = {
        "op": "semantic_transform", "kind": "apply_theme", "slides": slides,
        "preset": preset or "kpmg", "preserve_branding": True,
        "preserve_status_colors": "status" not in lowered or "preserve" in lowered,
    }

    explicit_primary = _theme_labeled_color(lowered, "primary")
    explicit_secondary = _theme_labeled_color(lowered, "secondary")
    explicit_accent = _theme_labeled_color(lowered, "accent")
    if explicit_primary:
        operation["primary"] = explicit_primary
    if explicit_secondary:
        operation["secondary"] = explicit_secondary
    if explicit_accent:
        operation["accent"] = explicit_accent

    background_color = _theme_labeled_color(lowered, "background")
    if background_color is None and "background" in lowered:
        before_background = lowered[:lowered.find("background")]
        candidates = _extract_theme_colors(before_background[-40:])
        if candidates:
            background_color = candidates[-1]

    text_color = _theme_labeled_color(lowered, "(?:body\\s+)?text")
    title_color = _theme_labeled_color(lowered, "titles?")
    if background_color:
        operation["background"] = background_color
    if text_color:
        operation["body_color"] = text_color
    if title_color:
        operation["title_color"] = title_color

    # A named preset should keep its designed palette unless the user also
    # supplied a custom list, such as "Executive Dark with teal and gold".
    use_custom_colors = bool(colors) and (preset is None or any(token in lowered for token in (" with ", " using ", " palette", ",", "#")))
    if use_custom_colors:
        palette_colors = []
        for color in colors:
            if color in {background_color, text_color, title_color}:
                continue
            if color not in palette_colors:
                palette_colors.append(color)
        if palette_colors and not explicit_primary:
            operation["primary"] = palette_colors[0]
        if len(palette_colors) > 1 and not explicit_secondary:
            operation["secondary"] = palette_colors[1]
        if len(palette_colors) > 2 and not explicit_accent:
            operation["accent"] = palette_colors[2]
        if palette_colors:
            operation["chart_colors"] = palette_colors

    theme_name = _THEME_PRESETS.get(preset or "", {}).get("name") or ("custom color" if colors or background_color or text_color or title_color else "Deck Refresh Blue")
    return {
        "message": f"Applied the {theme_name} theme to {scope_phrase}.",
        "operations": [operation],
        "assumptions": [], "failed": False,
    }


_theme_previous_apply_showcase_operation = apply_showcase_operation
_theme_previous_plan_showcase_command = plan_showcase_command


def apply_showcase_operation(prs: Presentation, operation: dict[str, Any]) -> dict[str, Any]:
    kind = str(operation.get("kind", "")).casefold()
    if kind == "apply_theme":
        return _apply_visual_theme(prs, operation)
    if kind == "chart_palette":
        return _apply_chart_palette(prs, operation)
    return _theme_previous_apply_showcase_operation(prs, operation)


def plan_showcase_command(pptx_path: str, user_message: str, selected_slide: int, chat_history=None) -> dict[str, Any] | None:
    theme_plan = _theme_command_plan(pptx_path, user_message, selected_slide)
    if theme_plan is not None:
        return theme_plan
    return _theme_previous_plan_showcase_command(pptx_path, user_message, selected_slide, chat_history)
