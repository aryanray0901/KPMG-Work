"""Built-in PowerPoint preview fallback.

This renderer is intentionally conservative. It renders common PowerPoint
objects directly from the PPTX package when PowerPoint, Keynote, or
LibreOffice is unavailable or temporarily fails. Native renderers remain the
first choice because they provide exact Office fidelity.
"""

from __future__ import annotations

import io
import math
import os
from functools import lru_cache
from pathlib import Path
from typing import Iterable

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from chart_contrast import contrast_palette

EMU_PER_INCH = 914400


def _hex_to_rgb(value, default=(255, 255, 255)):
    if value is None:
        return default
    try:
        text = str(value).replace("#", "").strip()
        if len(text) == 6:
            return tuple(int(text[i:i + 2], 16) for i in (0, 2, 4))
    except Exception:
        pass
    return default


def _rgb_from_color_format(color_format, default=None):
    try:
        rgb = color_format.rgb
        if rgb is not None:
            return _hex_to_rgb(rgb, default or (0, 0, 0))
    except Exception:
        pass
    return default


def _fill_rgb(fill, default=None):
    try:
        if fill is not None and fill.type is not None:
            rgb = _rgb_from_color_format(fill.fore_color, None)
            if rgb is not None:
                return rgb
    except Exception:
        pass
    return default


def _line_rgb(shape, default=(185, 193, 205)):
    try:
        rgb = _rgb_from_color_format(shape.line.color, None)
        if rgb is not None:
            return rgb
    except Exception:
        pass
    return default


def _luminance(rgb):
    if not rgb:
        return 255
    return 0.2126 * rgb[0] + 0.7152 * rgb[1] + 0.0722 * rgb[2]


_FONTS_DIR = os.path.join(os.path.dirname(os.path.abspath(__file__)), "fonts")


@lru_cache(maxsize=128)
def _font_path(bold: bool, italic: bool):
    candidates = []
    # Bundled fonts first: these ship with the app itself, so rendering
    # never depends on whether the host OS happens to have fonts
    # installed. Vercel's serverless containers generally don't -- without
    # this, every text element silently fell back to PIL's built-in
    # default font, a tiny bitmap font that ignores the requested size
    # entirely, which is why text looked uniformly small everywhere.
    if bold and italic:
        candidates.append(os.path.join(_FONTS_DIR, "DejaVuSans-BoldOblique.ttf"))
    elif bold:
        candidates.append(os.path.join(_FONTS_DIR, "DejaVuSans-Bold.ttf"))
    elif italic:
        candidates.append(os.path.join(_FONTS_DIR, "DejaVuSans-Oblique.ttf"))
    else:
        candidates.append(os.path.join(_FONTS_DIR, "DejaVuSans.ttf"))
    if os.name == "nt":
        windir = os.environ.get("WINDIR", r"C:\Windows")
        fonts = os.path.join(windir, "Fonts")
        if bold and italic:
            candidates += [os.path.join(fonts, "arialbi.ttf"), os.path.join(fonts, "calibriz.ttf")]
        elif bold:
            candidates += [os.path.join(fonts, "arialbd.ttf"), os.path.join(fonts, "calibrib.ttf")]
        elif italic:
            candidates += [os.path.join(fonts, "ariali.ttf"), os.path.join(fonts, "calibrii.ttf")]
        else:
            candidates += [os.path.join(fonts, "arial.ttf"), os.path.join(fonts, "calibri.ttf")]
    elif os.sys.platform == "darwin":
        if bold:
            candidates += ["/System/Library/Fonts/Supplemental/Arial Bold.ttf"]
        else:
            candidates += ["/System/Library/Fonts/Supplemental/Arial.ttf"]
    if bold:
        candidates += ["/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf"]
    else:
        candidates += ["/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf"]
    for path in candidates:
        if os.path.isfile(path):
            return path
    return None


@lru_cache(maxsize=256)
def _font(size: int, bold: bool = False, italic: bool = False):
    size = max(7, min(int(size), 120))
    path = _font_path(bool(bold), bool(italic))
    if path:
        try:
            return ImageFont.truetype(path, size=size)
        except Exception:
            pass
    return ImageFont.load_default()


def _box(shape, prs, width, height):
    x = int(round(float(shape.left) / float(prs.slide_width) * width))
    y = int(round(float(shape.top) / float(prs.slide_height) * height))
    w = max(1, int(round(float(shape.width) / float(prs.slide_width) * width)))
    h = max(1, int(round(float(shape.height) / float(prs.slide_height) * height)))
    return x, y, x + w, y + h


def _pt_to_px(pt, prs, width):
    try:
        points = float(pt.pt)
    except Exception:
        try:
            points = float(pt)
        except Exception:
            points = 16.0
    slide_inches = float(prs.slide_width) / EMU_PER_INCH
    return max(7, int(round(points * (width / slide_inches) / 72.0)))


def _wrap_line(draw, text, font, max_width):
    text = str(text or "")
    if not text:
        return [""]
    words = text.split()
    if not words:
        return [""]
    lines = []
    current = words[0]
    for word in words[1:]:
        candidate = current + " " + word
        if draw.textbbox((0, 0), candidate, font=font)[2] <= max_width:
            current = candidate
        else:
            lines.append(current)
            current = word
    lines.append(current)
    return lines


def _paragraph_style(paragraph, prs, width, fallback_color):
    run = paragraph.runs[0] if paragraph.runs else None
    font_obj = run.font if run is not None else paragraph.font
    size = _pt_to_px(font_obj.size if font_obj and font_obj.size is not None else 16, prs, width)
    bold = bool(font_obj.bold) if font_obj and font_obj.bold is not None else False
    italic = bool(font_obj.italic) if font_obj and font_obj.italic is not None else False
    color = fallback_color
    if font_obj is not None:
        try:
            value = _rgb_from_color_format(font_obj.color, None)
            if value is not None:
                color = value
        except Exception:
            pass
    return _font(size, bold, italic), color


def _draw_text_frame(draw, shape, prs, width, height, fallback_fill=None):
    if not getattr(shape, "has_text_frame", False):
        return
    left, top, right, bottom = _box(shape, prs, width, height)
    pad_x = max(4, int((right - left) * 0.035))
    pad_y = max(3, int((bottom - top) * 0.035))
    available = max(8, right - left - 2 * pad_x)
    fallback_color = (255, 255, 255) if _luminance(fallback_fill) < 110 else (31, 41, 55)
    y = top + pad_y

    paragraphs = list(shape.text_frame.paragraphs)
    for paragraph in paragraphs:
        text = paragraph.text or ""
        if not text.strip():
            y += 4
            continue
        font, color = _paragraph_style(paragraph, prs, width, fallback_color)
        prefix = ""
        try:
            if paragraph.level > 0 or paragraph._p.pPr is not None and paragraph._p.pPr.find("{http://schemas.openxmlformats.org/drawingml/2006/main}buChar") is not None:
                prefix = "• "
        except Exception:
            pass
        lines = []
        for raw_line in text.splitlines() or [text]:
            lines.extend(_wrap_line(draw, prefix + raw_line, font, available))
            prefix = ""
        line_height = max(9, int(draw.textbbox((0, 0), "Ag", font=font)[3] * 1.18))
        block_height = line_height * len(lines)
        if y + block_height > bottom:
            # Fit text to the available height instead of drawing outside the shape.
            remaining = max(8, bottom - y)
            scale = max(0.55, min(1.0, remaining / max(1, block_height)))
            font = _font(max(7, int(getattr(font, "size", 12) * scale)))
            line_height = max(8, int(draw.textbbox((0, 0), "Ag", font=font)[3] * 1.15))
        alignment = str(getattr(paragraph, "alignment", "") or "").lower()
        for line in lines:
            if y + line_height > bottom + 2:
                break
            text_width = draw.textbbox((0, 0), line, font=font)[2]
            if "center" in alignment:
                x = left + max(pad_x, (right - left - text_width) // 2)
            elif "right" in alignment:
                x = right - pad_x - text_width
            else:
                x = left + pad_x
            draw.text((x, y), line, font=font, fill=color)
            y += line_height
        y += max(1, line_height // 5)
        if y >= bottom:
            break


def _draw_picture(canvas, shape, prs, width, height):
    try:
        source = Image.open(io.BytesIO(shape.image.blob)).convert("RGBA")
    except Exception:
        return False
    left, top, right, bottom = _box(shape, prs, width, height)
    target_w, target_h = max(1, right - left), max(1, bottom - top)
    source.thumbnail((target_w, target_h), Image.Resampling.LANCZOS)
    x = left + (target_w - source.width) // 2
    y = top + (target_h - source.height) // 2
    canvas.alpha_composite(source, (x, y))
    return True


def _draw_table(draw, shape, prs, width, height):
    left, top, right, bottom = _box(shape, prs, width, height)
    table = shape.table
    rows = len(table.rows)
    cols = len(table.columns)
    if rows <= 0 or cols <= 0:
        return
    col_widths = [max(1, int(float(col.width) / sum(float(c.width) for c in table.columns) * (right - left))) for col in table.columns]
    row_heights = [max(1, int(float(row.height) / sum(float(r.height) for r in table.rows) * (bottom - top))) for row in table.rows]
    y = top
    for r, row in enumerate(table.rows):
        x = left
        for c, cell in enumerate(row.cells):
            cw = col_widths[c] if c < len(col_widths) else (right - left) // cols
            rh = row_heights[r] if r < len(row_heights) else (bottom - top) // rows
            fill = _fill_rgb(cell.fill, (0, 51, 141) if r == 0 else ((247, 249, 252) if r % 2 == 0 else (255, 255, 255)))
            draw.rectangle((x, y, x + cw, y + rh), fill=fill, outline=(196, 203, 214), width=1)
            text = cell.text.strip()
            font = _font(max(7, min(16, int(rh * 0.24))), bold=(r == 0))
            color = (255, 255, 255) if _luminance(fill) < 120 else (31, 41, 55)
            lines = _wrap_line(draw, text, font, max(8, cw - 10))
            ty = y + 4
            for line in lines[: max(1, rh // max(8, getattr(font, "size", 10) + 2))]:
                draw.text((x + 5, ty), line, font=font, fill=color)
                ty += max(8, int(getattr(font, "size", 10) * 1.15))
            x += cw
        y += row_heights[r] if r < len(row_heights) else (bottom - top) // rows


def _chart_values(chart):
    try:
        plot = chart.plots[0]
        categories = [str(value) for value in plot.categories]
        series = []
        for item in chart.series:
            values = []
            for value in item.values:
                try:
                    values.append(float(value))
                except Exception:
                    values.append(0.0)
            series.append((str(item.name or "Series"), values))
        return categories, series
    except Exception:
        return [], []


def _chart_kind(chart, series):
    chart_type = str(getattr(chart, "chart_type", "")).casefold()
    names = [name.casefold() for name, _values in series]
    if names[:3] == ["base", "increase", "decrease"]:
        return "waterfall"
    if "pie" in chart_type or "doughnut" in chart_type:
        return "pie"
    if "scatter" in chart_type or "xy_" in chart_type:
        return "scatter"
    if "line" in chart_type:
        return "line"
    if "area" in chart_type:
        return "area"
    if "bar_" in chart_type:
        return "bar"
    return "column"


def _series_rgb(chart, index, fallback):
    try:
        series = chart.series[index]
        return _fill_rgb(series.format.fill, _rgb_from_color_format(series.format.line.color, fallback))
    except Exception:
        return fallback


def _point_rgb(chart, index, fallback):
    try:
        return _fill_rgb(chart.series[0].points[index].format.fill, fallback)
    except Exception:
        return fallback


def _scatter_values(chart):
    output = []
    for item in chart.series:
        try:
            x_nodes = item._ser.xpath("./c:xVal/c:numRef/c:numCache/c:pt/c:v")
            x_values = [float(node.text) for node in x_nodes]
            y_values = [float(value) for value in item.values]
            output.append((str(item.name or "Series"), x_values, y_values))
        except Exception:
            continue
    return output


def _draw_chart_legend(draw, entries, x, y, max_width, font, text_color=(75, 85, 99)):
    cursor_x, cursor_y = x, y
    row_height = max(12, getattr(font, "size", 9) + 5)
    for label, color in entries:
        label = str(label)[:24]
        label_width = draw.textbbox((0, 0), label, font=font)[2]
        item_width = 12 + label_width + 14
        if cursor_x > x and cursor_x + item_width > x + max_width:
            cursor_x = x
            cursor_y += row_height
        draw.rectangle((cursor_x, cursor_y + 2, cursor_x + 8, cursor_y + 10), fill=color)
        draw.text((cursor_x + 12, cursor_y), label, font=font, fill=text_color)
        cursor_x += item_width
    return cursor_y + row_height


def _draw_chart(draw, shape, prs, width, height, slide_background=(255, 255, 255)):
    left, top, right, bottom = _box(shape, prs, width, height)
    chart = shape.chart
    chart_background = _fill_rgb(getattr(getattr(chart, "format", None), "fill", None), slide_background)
    readable = contrast_palette(chart_background)
    text_color = readable["text"]
    secondary_text = readable["secondary"]
    axis_color = readable["axis"]
    draw.rectangle((left, top, right, bottom), fill=chart_background, outline=readable["gridline"], width=1)
    categories, series = _chart_values(chart)
    palette = [
        (0, 94, 184), (0, 145, 218), (0, 163, 161), (72, 54, 152),
        (188, 32, 75), (234, 170, 0), (0, 166, 81), (109, 32, 119),
        (30, 73, 226), (94, 201, 230), (155, 218, 243), (124, 135, 142),
    ]
    kind = _chart_kind(chart, series)
    title = ""
    try:
        if chart.has_title:
            title = chart.chart_title.text_frame.text.strip()
    except Exception:
        pass
    if title:
        font = _font(max(10, int((bottom - top) * 0.055)), bold=True)
        draw.text((left + 10, top + 6), title, font=font, fill=text_color)
    chart_top = top + (34 if title else 14)
    chart_left = left + 48
    chart_right = right - 14
    chart_bottom = bottom - 30
    all_values = [value for _, values in series for value in values]
    if not series or not all_values or (not categories and kind != "scatter"):
        draw.text((left + 12, chart_top + 10), "Editable chart", font=_font(12), fill=secondary_text)
        return

    if kind == "pie":
        values = [max(0.0, value) for value in series[0][1]]
        total = sum(values)
        if total <= 0:
            draw.text((left + 12, chart_top + 10), "Pie chart needs positive values", font=_font(12), fill=secondary_text)
            return
        legend_width = max(125, int((right - left) * 0.23))
        pie_right = right - legend_width
        diameter = max(40, min(pie_right - left - 28, bottom - chart_top - 18))
        pie_left = left + 16 + max(0, (pie_right - left - 24 - diameter) // 2)
        pie_top = chart_top + max(0, (bottom - chart_top - diameter) // 2)
        angle = -90.0
        for index, value in enumerate(values):
            sweep = 360.0 * value / total
            color = _point_rgb(chart, index, palette[index % len(palette)])
            draw.pieslice((pie_left, pie_top, pie_left + diameter, pie_top + diameter), angle, angle + sweep,
                          fill=color, outline=(255, 255, 255), width=2)
            if sweep >= 16:
                middle = math.radians(angle + sweep / 2)
                radius = diameter * 0.34
                label = f"{value / total:.0%}"
                label_font = _font(max(8, min(12, diameter // 32)), bold=True)
                tx = pie_left + diameter / 2 + math.cos(middle) * radius
                ty = pie_top + diameter / 2 + math.sin(middle) * radius
                box = draw.textbbox((0, 0), label, font=label_font)
                draw.text((int(tx - (box[2] - box[0]) / 2), int(ty - (box[3] - box[1]) / 2)),
                          label, font=label_font, fill=(255, 255, 255))
            angle += sweep
        legend_font = _font(max(7, min(10, (bottom - chart_top) // max(16, len(categories) * 3))))
        y = chart_top + 4
        for index, category in enumerate(categories):
            color = _point_rgb(chart, index, palette[index % len(palette)])
            draw.rectangle((pie_right + 8, y + 2, pie_right + 18, y + 12), fill=color)
            draw.text((pie_right + 23, y), str(category)[:18], font=legend_font, fill=secondary_text)
            y += max(13, getattr(legend_font, "size", 8) + 5)
        return

    if kind == "scatter":
        scatter = _scatter_values(chart)
        if not scatter:
            draw.text((left + 12, chart_top + 10), "Editable scatter plot", font=_font(12), fill=secondary_text)
            return
        all_x = [value for _name, x_values, _y_values in scatter for value in x_values]
        all_y = [value for _name, _x_values, y_values in scatter for value in y_values]
        x_min, x_max = min(all_x), max(all_x)
        y_min, y_max = min(all_y), max(all_y)
        x_span, y_span = x_max - x_min or 1.0, y_max - y_min or 1.0
        legend_font = _font(8)
        entries = [(name, _series_rgb(chart, index, palette[index % len(palette)])) for index, (name, _x, _y) in enumerate(scatter)]
        chart_top = _draw_chart_legend(draw, entries, chart_left, chart_top, chart_right - chart_left, legend_font, secondary_text) + 2
        draw.line((chart_left, chart_bottom, chart_right, chart_bottom), fill=axis_color, width=1)
        draw.line((chart_left, chart_top, chart_left, chart_bottom), fill=axis_color, width=1)
        for index, (_name, x_values, y_values) in enumerate(scatter):
            color = _series_rgb(chart, index, palette[index % len(palette)])
            for x_value, y_value in zip(x_values, y_values):
                x = chart_left + int((x_value - x_min) / x_span * max(1, chart_right - chart_left))
                y = chart_bottom - int((y_value - y_min) / y_span * max(1, chart_bottom - chart_top))
                draw.ellipse((x - 4, y - 4, x + 4, y + 4), fill=color, outline=(255, 255, 255), width=1)
        return

    legend_font = _font(8)
    if len(series) > 1 and kind != "waterfall":
        entries = [(name, _series_rgb(chart, index, palette[index % len(palette)])) for index, (name, _values) in enumerate(series)]
        chart_top = _draw_chart_legend(draw, entries, chart_left, chart_top, chart_right - chart_left, legend_font, secondary_text) + 2

    minimum = min(0.0, min(all_values))
    maximum = max(1.0, max(all_values))
    if kind == "waterfall" and len(series) >= 3:
        base, increase, decrease = series[0][1], series[1][1], series[2][1]
        waterfall_bounds = [
            value
            for index in range(len(categories))
            for value in (
                base[index] if index < len(base) else 0.0,
                (base[index] if index < len(base) else 0.0)
                + (increase[index] if index < len(increase) else 0.0)
                + (decrease[index] if index < len(decrease) else 0.0),
            )
        ]
        minimum = min(0.0, min(waterfall_bounds))
        maximum = max(1.0, max(waterfall_bounds))
    span = maximum - minimum or 1.0

    if kind == "bar":
        chart_left = left + max(70, int((right - left) * 0.14))
        draw.line((chart_left, chart_bottom, chart_right, chart_bottom), fill=axis_color, width=1)
        group_height = max(1, (chart_bottom - chart_top) / max(1, len(categories)))
        bar_height = max(2, int(group_height * 0.72 / max(1, len(series))))
        label_font = _font(max(7, min(10, int(group_height * 0.20))))
        zero_x = chart_left - int(minimum / span * max(1, chart_right - chart_left))
        for category_index, category in enumerate(categories):
            center = chart_top + group_height * (category_index + 0.5)
            label = str(category)[:14]
            label_width = draw.textbbox((0, 0), label, font=label_font)[2]
            draw.text((chart_left - label_width - 7, int(center - getattr(label_font, "size", 8) / 2)),
                      label, font=label_font, fill=secondary_text)
            for series_index, (_name, values) in enumerate(series):
                value = values[category_index] if category_index < len(values) else 0.0
                x_value = chart_left + int((value - minimum) / span * max(1, chart_right - chart_left))
                y0 = int(center - group_height * 0.36 + series_index * bar_height)
                color = _series_rgb(chart, series_index, palette[series_index % len(palette)])
                draw.rectangle((min(zero_x, x_value), y0, max(zero_x, x_value), y0 + bar_height - 1), fill=color)
        return

    draw.line((chart_left, chart_bottom, chart_right, chart_bottom), fill=axis_color, width=1)
    draw.line((chart_left, chart_top, chart_left, chart_bottom), fill=axis_color, width=1)
    group_width = max(1, (chart_right - chart_left) / max(1, len(categories)))
    series_count = max(1, len(series))
    bar_width = max(2, int(group_width * 0.72 / series_count))
    zero_y = chart_bottom - int((0.0 - minimum) / span * max(1, chart_bottom - chart_top))
    points_by_series = []
    for s_index, (_, values) in enumerate(series):
        color = _series_rgb(chart, s_index, palette[s_index % len(palette)])
        points = []
        for c_index, _category in enumerate(categories):
            value = values[c_index] if c_index < len(values) else 0.0
            x_center = chart_left + group_width * (c_index + 0.5)
            y_value = chart_bottom - int((value - minimum) / span * max(1, chart_bottom - chart_top))
            points.append((int(x_center), y_value))
            if kind == "column":
                x0 = int(x_center - group_width * 0.36 + s_index * bar_width)
                x1 = x0 + bar_width - 1
                draw.rectangle((x0, min(zero_y, y_value), x1, max(zero_y, y_value)), fill=color)
        points_by_series.append((color, points))
    if kind == "line":
        for color, points in points_by_series:
            if len(points) > 1:
                draw.line(points, fill=color, width=3)
            for x, y in points:
                draw.ellipse((x - 3, y - 3, x + 3, y + 3), fill=color)
    elif kind == "area":
        for color, points in points_by_series:
            polygon = [(points[0][0], chart_bottom), *points, (points[-1][0], chart_bottom)] if points else []
            light = tuple(int(channel + (255 - channel) * 0.68) for channel in color)
            if polygon:
                draw.polygon(polygon, fill=light)
        for color, points in points_by_series:
            if len(points) > 1:
                draw.line(points, fill=color, width=3)
    elif kind == "waterfall" and len(series) >= 3:
        base, increase, decrease = series[0][1], series[1][1], series[2][1]
        wf_width = max(3, int(group_width * 0.55))
        for index in range(len(categories)):
            base_value = base[index] if index < len(base) else 0.0
            up = increase[index] if index < len(increase) else 0.0
            down = decrease[index] if index < len(decrease) else 0.0
            start = base_value
            end = base_value + up if up else base_value + down
            y0 = chart_bottom - int((start - minimum) / span * max(1, chart_bottom - chart_top))
            y1 = chart_bottom - int((end - minimum) / span * max(1, chart_bottom - chart_top))
            x = int(chart_left + group_width * (index + 0.5))
            color = (0, 166, 81) if up else (188, 32, 75)
            draw.rectangle((x - wf_width // 2, min(y0, y1), x + wf_width // 2, max(y0, y1)), fill=color)
    label_font = _font(max(7, min(11, int(group_width * 0.12))))
    for index, category in enumerate(categories):
        if len(categories) > 8 and index % 2 == 1:
            continue
        label = category[:16]
        x_center = chart_left + group_width * (index + 0.5)
        text_width = draw.textbbox((0, 0), label, font=label_font)[2]
        draw.text((int(x_center - text_width / 2), chart_bottom + 5), label, font=label_font, fill=secondary_text)


def _draw_shape(canvas, draw, shape, prs, width, height, slide_background=(255, 255, 255)):
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        _draw_picture(canvas, shape, prs, width, height)
        return
    if getattr(shape, "has_table", False):
        _draw_table(draw, shape, prs, width, height)
        return
    if getattr(shape, "has_chart", False):
        _draw_chart(draw, shape, prs, width, height, slide_background)
        return
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        try:
            for child in shape.shapes:
                _draw_shape(canvas, draw, child, prs, width, height, slide_background)
        except Exception:
            pass
        return

    left, top, right, bottom = _box(shape, prs, width, height)
    fill = _fill_rgb(getattr(shape, "fill", None), None)
    outline = _line_rgb(shape, None)
    if fill is not None or outline is not None:
        draw.rectangle((left, top, right, bottom), fill=fill, outline=outline, width=1)
    _draw_text_frame(draw, shape, prs, width, height, fallback_fill=fill)


def _slide_background(slide):
    try:
        return _fill_rgb(slide.background.fill, (255, 255, 255))
    except Exception:
        return (255, 255, 255)


def render_pptx(pptx_path: str, out_dir: str, prefix: str, canvas_width: int = 1600):
    """Render common PowerPoint content to PNG files.

    Returns a list of generated image paths, or ``None`` if the file itself is
    unreadable. The function does not modify the PowerPoint.
    """
    try:
        prs = Presentation(pptx_path)
    except Exception:
        return None
    if not prs.slides:
        return None
    os.makedirs(out_dir, exist_ok=True)
    aspect = float(prs.slide_height) / float(prs.slide_width)
    canvas_height = max(600, int(round(canvas_width * aspect)))
    paths = []
    for index, slide in enumerate(prs.slides, start=1):
        background = _slide_background(slide)
        canvas = Image.new("RGBA", (canvas_width, canvas_height), background + (255,))
        draw = ImageDraw.Draw(canvas)
        for shape in slide.shapes:
            try:
                _draw_shape(canvas, draw, shape, prs, canvas_width, canvas_height, background)
            except Exception:
                # One unsupported object must not prevent the rest of the slide
                # or deck from being previewed.
                continue
        path = os.path.join(out_dir, f"{prefix}_{index}.png")
        canvas.convert("RGB").save(path, format="PNG", optimize=True)
        paths.append(path)
    return paths or None
