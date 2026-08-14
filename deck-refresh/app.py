"""
Deck Refresh
------------
Upload a PowerPoint OR Excel file, upload/paste new data, and this app finds
matching numbers by label, fuzzy-matches across period changes (Q3 -> Q4,
FY25 -> FY26, etc.), rewrites headings to match, and writes new values
straight into the existing cells, runs, and chart data so the original
formatting and file structure remain in place. The app runs locally and does
not send files to an external service.
"""

import os
import re
import io
import json
import time
import random
import uuid
import shutil
import subprocess
import tempfile
import threading
import sys
from pathlib import Path
from collections import Counter
from numbers import Number

from flask import Flask, request, render_template, send_file, redirect, url_for, flash, abort, jsonify
from werkzeug.utils import secure_filename
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from openpyxl import load_workbook
import pandas as pd
from rapidfuzz import fuzz
from dotenv import load_dotenv
from builtin_preview import render_pptx as render_pptx_builtin
from builder_ops import LAYOUTS, create_blank_deck, layout_operations, wizard
from chart_contrast import ensure_chart_contrast
import storage as blob_storage
from replacement_engine import (
    ReplacementError,
    apply_selected_deck_replacements,
    compare_deck_replacements,
    inspect_deck_structure,
    replace_deck_1to1,
    validate_replacement_structure,
)

load_dotenv(Path(__file__).with_name(".env"))

from pptx_editor import (
    EditorError,
    apply_operations as apply_editor_operations,
    plan_edit as plan_editor_edit,
    guaranteed_local_plan,
    diagnose_failure_message as diagnose_editor_failure,
)

BASE_DIR = os.path.dirname(os.path.abspath(__file__))
IS_VERCEL = bool(os.environ.get("VERCEL"))
SESSIONS_DIR = os.environ.get("DECK_REFRESH_SESSIONS_DIR") or (
    os.path.join(tempfile.gettempdir(), "deck_refresh_sessions")
    if IS_VERCEL else os.path.join(BASE_DIR, "sessions")
)
os.makedirs(SESSIONS_DIR, exist_ok=True)

app = Flask(__name__)
_secret = os.environ.get("FLASK_SECRET_KEY")
if not _secret:
    # Never fall back to a hardcoded secret when running on a public deployment.
    # Locally this just means sessions reset on restart, which is harmless.
    _secret = "deck-refresh-local-secret" if not IS_VERCEL else os.urandom(32).hex()
app.secret_key = _secret
app.config["MAX_CONTENT_LENGTH"] = 100 * 1024 * 1024  # 100MB

# ---------------------------------------------------------------------------
# REAL SLIDE RENDERING
# Uses Microsoft PowerPoint on Windows and macOS, then LibreOffice, with
# Apple Keynote as a final macOS fallback. Each engine renders the real PPTX.
# ---------------------------------------------------------------------------

def _first_existing(*candidates):
    for candidate in candidates:
        if not candidate:
            continue
        expanded = os.path.expanduser(candidate)
        if os.path.isfile(expanded):
            return expanded
    return None


SOFFICE_PATH = (
    os.environ.get("LIBREOFFICE_PATH")
    or shutil.which("soffice")
    or shutil.which("libreoffice")
    or _first_existing(
        "/Applications/LibreOffice.app/Contents/MacOS/soffice",
        "~/Applications/LibreOffice.app/Contents/MacOS/soffice",
        r"C:\Program Files\LibreOffice\program\soffice.exe",
        r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
    )
)
POWERSHELL_PATH = shutil.which("powershell") or shutil.which("pwsh")
OSASCRIPT_PATH = shutil.which("osascript")
_soffice_lock = threading.Lock()
_powerpoint_lock = threading.Lock()
_mac_office_lock = threading.Lock()
_keynote_lock = threading.Lock()

try:
    import fitz  # PyMuPDF
    HAVE_FITZ = True
except ImportError:
    HAVE_FITZ = False


def _mac_app_exists(app_name):
    if sys.platform != "darwin":
        return False
    return any(
        os.path.isdir(os.path.expanduser(path))
        for path in (
            f"/Applications/{app_name}.app",
            f"~/Applications/{app_name}.app",
        )
    )


def _pdf_to_slide_images(pdf_path, out_dir, prefix):
    if not HAVE_FITZ or not pdf_path or not os.path.exists(pdf_path):
        return None
    try:
        doc = fitz.open(pdf_path)
        paths = []
        for i, page in enumerate(doc):
            pix = page.get_pixmap(matrix=fitz.Matrix(2, 2), alpha=False)
            img_path = os.path.join(out_dir, f"{prefix}_{i + 1}.png")
            pix.save(img_path)
            paths.append(img_path)
        doc.close()
        return paths or None
    except Exception:
        return None


def _render_pptx_with_powerpoint(pptx_path, out_dir, prefix):
    """Render with the installed Microsoft PowerPoint desktop app on Windows.
    This uses PowerPoint's own export engine, so the browser preview matches
    what the user sees when opening the PPTX in PowerPoint.
    """
    if os.name != "nt" or not POWERSHELL_PATH:
        return None

    os.makedirs(out_dir, exist_ok=True)
    script_path = os.path.join(out_dir, f"render_{prefix}.ps1")
    script = r'''
param(
    [Parameter(Mandatory=$true)][string]$InputPath,
    [Parameter(Mandatory=$true)][string]$OutputDir,
    [Parameter(Mandatory=$true)][string]$Prefix
)
$ErrorActionPreference = "Stop"
$powerpoint = $null
$presentation = $null
try {
    $powerpoint = New-Object -ComObject PowerPoint.Application
    $presentation = $powerpoint.Presentations.Open($InputPath, -1, 0, 0)
    $slideWidth = [double]$presentation.PageSetup.SlideWidth
    $slideHeight = [double]$presentation.PageSetup.SlideHeight
    $exportWidth = 1600
    $exportHeight = [int][Math]::Round($exportWidth * ($slideHeight / $slideWidth))
    for ($i = 1; $i -le $presentation.Slides.Count; $i++) {
        $outputPath = Join-Path $OutputDir ("{0}_{1}.png" -f $Prefix, $i)
        $presentation.Slides.Item($i).Export($outputPath, "PNG", $exportWidth, $exportHeight)
    }
}
finally {
    if ($presentation -ne $null) {
        $presentation.Close()
        [void][System.Runtime.InteropServices.Marshal]::ReleaseComObject($presentation)
    }
    if ($powerpoint -ne $null) {
        $powerpoint.Quit()
        [void][System.Runtime.InteropServices.Marshal]::ReleaseComObject($powerpoint)
    }
    [GC]::Collect()
    [GC]::WaitForPendingFinalizers()
}
'''
    try:
        with open(script_path, "w", encoding="utf-8") as f:
            f.write(script)
        with _powerpoint_lock:
            result = subprocess.run(
                [POWERSHELL_PATH, "-NoProfile", "-ExecutionPolicy", "Bypass",
                 "-File", script_path, "-InputPath", os.path.abspath(pptx_path),
                 "-OutputDir", os.path.abspath(out_dir), "-Prefix", prefix],
                capture_output=True, timeout=180,
            )
        if result.returncode != 0:
            return None
        paths = []
        i = 1
        while True:
            path = os.path.join(out_dir, f"{prefix}_{i}.png")
            if not os.path.exists(path):
                break
            paths.append(path)
            i += 1
        return paths or None
    except Exception:
        return None
    finally:
        try:
            os.remove(script_path)
        except OSError:
            pass


def _mac_render_job_root(app_kind):
    """Use an app-owned sandbox directory so Office and Keynote can write
    previews without asking for access to an arbitrary project folder.
    """
    home = os.path.expanduser("~")
    if app_kind == "powerpoint":
        candidates = [
            os.path.join(home, "Library", "Containers", "com.microsoft.Powerpoint", "Data", "Documents", "DeckRefreshPreview"),
            os.path.join(home, "Library", "Group Containers", "UBF8T346G9.Office", "TemporaryItems", "DeckRefreshPreview"),
        ]
    else:
        candidates = [
            os.path.join(home, "Library", "Containers", "com.apple.iWork.Keynote", "Data", "Documents", "DeckRefreshPreview"),
            os.path.join(home, "Library", "Caches", "DeckRefreshPreview"),
        ]
    for candidate in candidates:
        try:
            os.makedirs(candidate, exist_ok=True)
            return candidate
        except OSError:
            continue
    return tempfile.gettempdir()


def _render_pptx_with_powerpoint_mac(pptx_path, out_dir, prefix):
    """Render a PPTX on macOS through the installed PowerPoint application.

    PowerPoint exports the presentation to PDF through AppleScript. PyMuPDF
    then turns each PDF page into a PNG. The first run may trigger macOS's
    Automation permission prompt for Terminal or Python.
    """
    if sys.platform != "darwin" or not OSASCRIPT_PATH or not HAVE_FITZ:
        return None
    if not _mac_app_exists("Microsoft PowerPoint"):
        return None

    os.makedirs(out_dir, exist_ok=True)
    job_root = _mac_render_job_root("powerpoint")
    job_dir = tempfile.mkdtemp(prefix="deck_refresh_", dir=job_root)
    input_copy = os.path.join(job_dir, "presentation.pptx")
    pdf_path = os.path.join(job_dir, "presentation.pdf")
    script_path = os.path.join(job_dir, "render.applescript")
    script = r'''
on run argv
    set inputPath to item 1 of argv
    set outputPath to item 2 of argv
    set inputFile to POSIX file inputPath
    set outputFile to POSIX file outputPath
    tell application "Microsoft PowerPoint"
        launch
        open inputFile
        delay 0.4
        set openedPresentation to active presentation
        save openedPresentation in outputFile as save as PDF
        close openedPresentation saving no
    end tell
    return outputPath
end run
'''
    try:
        shutil.copy2(pptx_path, input_copy)
        Path(pdf_path).touch()
        with open(script_path, "w", encoding="utf-8") as f:
            f.write(script)
        with _mac_office_lock:
            result = subprocess.run(
                [OSASCRIPT_PATH, script_path, input_copy, pdf_path],
                capture_output=True,
                text=True,
                timeout=180,
            )
        if result.returncode != 0 or not os.path.exists(pdf_path) or os.path.getsize(pdf_path) < 100:
            return None
        return _pdf_to_slide_images(pdf_path, out_dir, prefix)
    except Exception:
        return None
    finally:
        shutil.rmtree(job_dir, ignore_errors=True)


def _render_pptx_with_keynote(pptx_path, out_dir, prefix):
    """Render with Apple Keynote when PowerPoint and LibreOffice are absent."""
    if sys.platform != "darwin" or not OSASCRIPT_PATH or not HAVE_FITZ:
        return None
    if not _mac_app_exists("Keynote"):
        return None

    os.makedirs(out_dir, exist_ok=True)
    job_root = _mac_render_job_root("keynote")
    job_dir = tempfile.mkdtemp(prefix="deck_refresh_", dir=job_root)
    input_copy = os.path.join(job_dir, "presentation.pptx")
    pdf_path = os.path.join(job_dir, "presentation.pdf")
    script_path = os.path.join(job_dir, "render.applescript")
    script = r'''
on run argv
    set inputPath to item 1 of argv
    set outputPath to item 2 of argv
    set inputFile to POSIX file inputPath
    set outputFile to POSIX file outputPath
    tell application "Keynote"
        launch
        set openedDocument to open inputFile
        delay 0.4
        export openedDocument to outputFile as PDF
        close openedDocument saving no
    end tell
    return outputPath
end run
'''
    try:
        shutil.copy2(pptx_path, input_copy)
        with open(script_path, "w", encoding="utf-8") as f:
            f.write(script)
        with _keynote_lock:
            result = subprocess.run(
                [OSASCRIPT_PATH, script_path, input_copy, pdf_path],
                capture_output=True,
                text=True,
                timeout=180,
            )
        if result.returncode != 0 or not os.path.exists(pdf_path) or os.path.getsize(pdf_path) < 100:
            return None
        return _pdf_to_slide_images(pdf_path, out_dir, prefix)
    except Exception:
        return None
    finally:
        shutil.rmtree(job_dir, ignore_errors=True)


def _convert_pptx_to_pdf(pptx_path, out_dir):
    if not SOFFICE_PATH:
        return None
    profile_dir = tempfile.mkdtemp(prefix="lo_profile_")
    try:
        with _soffice_lock:
            result = subprocess.run(
                [SOFFICE_PATH, "--headless", "--norestore",
                 f"-env:UserInstallation=file://{profile_dir}",
                 "--convert-to", "pdf", "--outdir", out_dir, pptx_path],
                capture_output=True, timeout=120,
            )
        pdf_path = os.path.join(out_dir, os.path.splitext(os.path.basename(pptx_path))[0] + ".pdf")
        return pdf_path if (result.returncode == 0 and os.path.exists(pdf_path)) else None
    except Exception:
        return None
    finally:
        shutil.rmtree(profile_dir, ignore_errors=True)


def _render_pptx_with_libreoffice(pptx_path, out_dir, prefix):
    if not (SOFFICE_PATH and HAVE_FITZ):
        return None
    os.makedirs(out_dir, exist_ok=True)
    pdf_path = _convert_pptx_to_pdf(pptx_path, out_dir)
    if not pdf_path:
        return None
    try:
        return _pdf_to_slide_images(pdf_path, out_dir, prefix)
    finally:
        try:
            os.remove(pdf_path)
        except OSError:
            pass


def _renderer_help_text():
    if sys.platform == "darwin":
        if _mac_app_exists("Microsoft PowerPoint"):
            return (
                "Click Retry preview. If macOS asks, allow Terminal or Python to control "
                "Microsoft PowerPoint. You can also enable it in System Settings, Privacy & Security, Automation."
            )
        if _mac_app_exists("Keynote"):
            return "Click Retry preview and allow Terminal or Python to control Keynote if macOS asks."
        return "Install Microsoft PowerPoint, Apple Keynote, or LibreOffice, restart Deck Refresh, then retry the preview."
    if os.name == "nt":
        return "Install Microsoft PowerPoint or LibreOffice, restart Deck Refresh, then retry the preview."
    return "Install LibreOffice, restart Deck Refresh, then retry the preview."


def _any_renderer_available():
    if os.name == "nt" and POWERSHELL_PATH:
        return True
    if sys.platform == "darwin" and (_mac_app_exists("Microsoft PowerPoint") or _mac_app_exists("Keynote")):
        return True
    return bool(SOFFICE_PATH and HAVE_FITZ)


def _clear_render_prefix(out_dir, prefix):
    try:
        for name in os.listdir(out_dir):
            if name.startswith(f"{prefix}_") and name.lower().endswith(".png"):
                try:
                    os.remove(os.path.join(out_dir, name))
                except OSError:
                    pass
    except OSError:
        pass


def render_pptx_to_images(pptx_path, out_dir, prefix):
    """Return slide image paths and the rendering engine used.

    Native Office renderers are attempted first. Each native renderer gets a
    second attempt because PowerPoint and Keynote occasionally retain a COM or
    Automation lock for a moment after a prior export. A built-in renderer is
    the final fallback, so a successful edit still receives a fresh preview
    even when the desktop renderer is missing or temporarily unavailable.
    """
    engines = [
        (_render_pptx_with_powerpoint, "Microsoft PowerPoint"),
        (_render_pptx_with_powerpoint_mac, "Microsoft PowerPoint for Mac"),
        (_render_pptx_with_libreoffice, "LibreOffice"),
        (_render_pptx_with_keynote, "Apple Keynote"),
    ]
    for renderer, label in engines:
        for attempt in range(2):
            _clear_render_prefix(out_dir, prefix)
            paths = renderer(pptx_path, out_dir, prefix)
            if paths:
                return paths, label
            if attempt == 0:
                time.sleep(0.35)

    _clear_render_prefix(out_dir, prefix)
    builtin_paths = render_pptx_builtin(pptx_path, out_dir, prefix)
    if builtin_paths:
        return builtin_paths, "Built-in PowerPoint preview"

    return None, None
@app.after_request
def _no_cache(response):
    response.headers["Cache-Control"] = "no-store, no-cache, must-revalidate, max-age=0"
    return response


# ---------------------------------------------------------------------------
# SHARED HELPERS
# ---------------------------------------------------------------------------

NUMBER_RE = re.compile(
    r"(?P<prefix>[$€£]?)\s*(?P<num>-?\d{1,3}(?:,\d{3})*(?:\.\d+)?|\-?\d+(?:\.\d+)?)"
    r"\s*(?P<mag>[KkMmBbTt])?\s*(?P<suffix>%?)"
)
MAGNITUDE = {"K": 1_000, "M": 1_000_000, "B": 1_000_000_000, "T": 1_000_000_000_000}

PERIOD_WORDS = re.compile(
    r"\b(q1|q2|q3|q4|quarter\s*[1-4]|fy\s*\d{2,4}|20\d{2}|h1|h2|jan|feb|mar|apr|may|jun|jul|aug|sep|sept|oct|nov|dec|"
    r"january|february|march|april|june|july|august|september|october|november|december|ytd|mtd)\b",
    re.IGNORECASE,
)

QUARTER_TOKEN_RE = re.compile(r"\bQ[1-4]\b", re.IGNORECASE)
FY_TOKEN_RE = re.compile(r"\bFY\s?\d{2,4}\b", re.IGNORECASE)
YEAR_TOKEN_RE = re.compile(r"\b20\d{2}\b")


def normalize_label(text):
    t = (text or "").lower()
    t = re.sub(r"\([^)]*\)", " ", t)  # drop unit annotations like "($M)", "(%)"
    t = PERIOD_WORDS.sub(" ", t)
    t = re.sub(r"[^a-z0-9\s]", " ", t)
    t = re.sub(r"\s+", " ", t).strip()
    return t


def looks_numeric(text):
    text = (text or "").strip()
    if not text:
        return False
    m = NUMBER_RE.fullmatch(text)
    return m is not None and re.search(r"\d", text) is not None


def extract_number_value(text):
    if text is None:
        return None
    if isinstance(text, (int, float)):
        return float(text)
    m = NUMBER_RE.search(str(text))
    if not m:
        return None
    raw = m.group("num").replace(",", "")
    try:
        return float(raw)
    except ValueError:
        return None


def format_new_value(original_text, new_value):
    """Render new_value using the same style ($, %, commas, decimals, K/M/B/T) as original_text."""
    original_text = str(original_text)
    currency = ""
    for sym in ("$", "€", "£"):
        if sym in original_text:
            currency = sym
            break
    has_percent = "%" in original_text
    if has_percent and abs(new_value) <= 1:
        new_value *= 100
    mag_match = re.search(r"\d\s*([KkMmBbTt])\b", original_text)
    magnitude = mag_match.group(1) if mag_match else ""
    has_comma = "," in original_text
    decimals = 0
    m = re.search(r"\.(\d+)", original_text)
    if m:
        decimals = len(m.group(1))
    is_negative = new_value < 0
    abs_val = abs(new_value)
    if decimals:
        num_str = f"{abs_val:,.{decimals}f}" if has_comma else f"{abs_val:.{decimals}f}"
    else:
        if abs_val == int(abs_val):
            num_str = f"{int(abs_val):,}" if has_comma else f"{int(abs_val)}"
        else:
            num_str = f"{abs_val:,.2f}" if has_comma else f"{abs_val:.2f}"
    sign = "-" if is_negative else ""
    result = f"{sign}{currency}{num_str}{magnitude}"
    if has_percent:
        result += "%"
    return result


def detect_period_tokens(texts):
    q_counter, fy_counter, year_counter = Counter(), Counter(), Counter()
    for t in texts:
        if not t:
            continue
        t = str(t)
        for m in FY_TOKEN_RE.finditer(t):
            fy_counter[re.sub(r"\s+", "", m.group(0)).upper()] += 1
        remainder = FY_TOKEN_RE.sub(" ", t)
        for m in QUARTER_TOKEN_RE.finditer(remainder):
            q_counter[m.group(0).upper()] += 1
        remainder2 = QUARTER_TOKEN_RE.sub(" ", remainder)
        for m in YEAR_TOKEN_RE.finditer(remainder2):
            year_counter[m.group(0)] += 1
    return q_counter, fy_counter, year_counter


def compute_period_replacements(deck_texts, source_pairs, pasted_text):
    deck_q, deck_fy, deck_year = detect_period_tokens(deck_texts)
    src_texts = [lbl for lbl, _ in source_pairs]
    if pasted_text:
        src_texts.append(pasted_text)
    src_q, src_fy, src_year = detect_period_tokens(src_texts)

    replacements = {}
    for deck_counter, src_counter in ((deck_q, src_q), (deck_fy, src_fy), (deck_year, src_year)):
        if not deck_counter or not src_counter:
            continue
        old = deck_counter.most_common(1)[0][0]
        new = src_counter.most_common(1)[0][0]
        if old.upper() != new.upper():
            replacements[old] = new
    return replacements


KEY_METRIC_WORDS = {
    "revenue", "cost", "costs", "expense", "expenses", "margin", "income",
    "profit", "budget", "actual", "variance", "ebitda", "ebit", "loss",
}


def match_targets_to_source(targets, source_pairs, threshold=60):
    norm_sources = [(normalize_label(lbl), lbl, val) for lbl, val in source_pairs]
    matches = []
    for t in targets:
        norm_t = normalize_label(t["label"])
        t_tokens = set(norm_t.split())
        t_key = t_tokens & KEY_METRIC_WORDS
        best, best_score, best_metric = None, 0, None
        top_seen_score = 0
        if norm_t:
            for norm_s, orig_s, val in norm_sources:
                if not norm_s:
                    continue
                s_tokens = set(norm_s.split())
                if not s_tokens:
                    continue
                score = fuzz.token_sort_ratio(norm_t, norm_s)
                top_seen_score = max(top_seen_score, score)
                # A long shared descriptive phrase (e.g. "Technology
                # Consulting") can make "...Costs" and "...Revenue" look
                # highly similar by raw character overlap even though
                # they're different metrics. If the target names a specific
                # metric keyword, the candidate must contain that same
                # keyword too -- a bare label with no keyword at all, or one
                # naming a *different* keyword, must not satisfy it.
                s_key = s_tokens & KEY_METRIC_WORDS
                if t_key and not (t_key <= s_key):
                    continue
                # Among everything that survives the keyword gate, prefer
                # the candidate whose token set differs least from the
                # target's -- exact-word overlap is a far more reliable
                # financial-label signal than raw character similarity
                # (which can rank "Advisory Budget" above the correct
                # "Deal Advisory Actual" by a hair, or "Total Revenue"
                # above "Tax Revenue" for a target literally saying "Tax").
                # Ties broken by fuzzy score, then by preferring an
                # explicit "Actual" figure -- the sensible default when the
                # target doesn't say which period-type it wants.
                sym_diff = len(t_tokens ^ s_tokens)
                actual_bonus = 0 if "actual" in s_tokens else 1
                metric = (sym_diff, actual_bonus, -score)
                if best_metric is None or metric < best_metric:
                    best_metric, best_score, best = metric, score, (orig_s, val)
        entry = dict(t)
        if best and best_score >= threshold:
            entry["matched_label"] = best[0]
            entry["new_value"] = best[1]
            entry["score"] = round(best_score, 1)
        else:
            entry["matched_label"] = None
            entry["new_value"] = None
            entry["score"] = round(top_seen_score, 1)
        if entry["new_value"] is not None:
            entry["new_text_preview"] = format_new_value(entry["original"], entry["new_value"])
        else:
            entry["new_text_preview"] = None
        matches.append(entry)
    return matches


# ---------------------------------------------------------------------------
# SOURCE DATA EXTRACTION (label -> new value), shared by pptx & xlsx targets
# ---------------------------------------------------------------------------

def extract_from_spreadsheet(path):
    pairs = []
    try:
        if path.lower().endswith(".csv"):
            frames = {"csv": pd.read_csv(path)}
        else:
            raw = pd.read_excel(path, sheet_name=None)
            frames = raw if isinstance(raw, dict) else {"Sheet1": raw}
    except Exception:
        return pairs

    for sheet_name, frame in frames.items():
        frame = frame.dropna(how="all").dropna(axis=1, how="all")
        if frame.empty:
            continue
        cols = [str(c) for c in frame.columns]
        first_col_is_label = (
            not pd.api.types.is_numeric_dtype(frame.iloc[:, 0]) if frame.shape[1] > 0 else False
        )

        if frame.shape[1] == 2 and first_col_is_label:
            for _, row in frame.iterrows():
                label = str(row.iloc[0]).strip()
                val = row.iloc[1]
                if pd.notna(val) and label and label.lower() != "nan":
                    try:
                        pairs.append((label, float(val)))
                    except (ValueError, TypeError):
                        pass
            continue

        for _, row in frame.iterrows():
            row_label = str(row.iloc[0]).strip() if first_col_is_label else ""
            start = 1 if first_col_is_label else 0
            data_col_count = len(cols) - start
            for c in range(start, len(cols)):
                val = row.iloc[c]
                if pd.isna(val):
                    continue
                try:
                    val = float(val)
                except (ValueError, TypeError):
                    continue
                col_label = cols[c]
                combined = f"{row_label} {col_label}".strip()
                if combined:
                    pairs.append((combined, val))
                # Bare row/column labels are only unambiguous when there's a
                # single data column (or single data row) -- otherwise every
                # row sharing that column header (e.g. every line item's
                # "Q4 Actual") would collide under the same bare label with
                # different values, and whichever happened to be inserted
                # first would silently win.
                if data_col_count == 1 and row_label:
                    pairs.append((row_label, val))
                if len(frame) == 1 and col_label and col_label.lower() not in ("unnamed: 0",):
                    pairs.append((col_label, val))
    return pairs


RELATIVE_SYNONYMS = {
    "sales": "revenue", "turnover": "revenue", "income": "profit",
    "earnings": "profit", "expenses": "cost", "expense": "cost",
    "costs": "cost", "employees": "headcount", "staff": "headcount",
    "workers": "headcount", "clients": "customers", "customer": "customers",
    "actuals": "actual", "forecast": "budget", "plan": "budget",
    "planned": "budget", "target": "budget", "percentage": "percent",
    "pct": "percent", "rate": "percent", "geography": "region",
    "geographic": "region", "countries": "country", "companies": "company",
}
RELATIVE_GENERIC_TOKENS = {
    "data", "value", "values", "metric", "metrics", "figure", "figures",
    "sheet", "table", "chart", "slide", "series", "point", "row", "column",
}


def _relative_tokens(value):
    normalized = normalize_label(value)
    tokens = []
    for token in normalized.split():
        mapped = RELATIVE_SYNONYMS.get(token, token)
        if mapped not in RELATIVE_GENERIC_TOKENS:
            tokens.append(mapped)
    return tokens


def _raw_spreadsheet_rows(path):
    suffix = os.path.splitext(str(path))[1].lower()
    if suffix == ".csv":
        frame = pd.read_csv(path, header=None)
        return {"CSV": frame.where(pd.notna(frame), None).values.tolist()}
    if suffix in {".xlsx", ".xlsm"}:
        workbook = load_workbook(path, data_only=True, read_only=True)
        try:
            return {
                sheet.title: [list(row) for row in sheet.iter_rows(values_only=True)]
                for sheet in workbook.worksheets
            }
        finally:
            workbook.close()
    if suffix == ".xls":
        frames = pd.read_excel(path, sheet_name=None, header=None)
        return {
            name: frame.where(pd.notna(frame), None).values.tolist()
            for name, frame in frames.items()
        }
    return {}


def _relative_numeric_value(value):
    if isinstance(value, Number) and not isinstance(value, bool):
        try:
            return float(value)
        except (TypeError, ValueError):
            return None
    text = str(value or "").strip()
    if looks_numeric(text):
        return extract_number_value(text)
    return None


def _nearest_text_left(rows, row_index, column_index, limit=6):
    row = rows[row_index]
    for offset in range(1, limit + 1):
        index = column_index - offset
        if index < 0:
            break
        value = row[index] if index < len(row) else None
        text = str(value or "").strip()
        if text and _relative_numeric_value(value) is None:
            return text
    return ""


def _nearest_text_above(rows, row_index, column_index, limit=12):
    for offset in range(1, limit + 1):
        index = row_index - offset
        if index < 0:
            break
        row = rows[index]
        value = row[column_index] if column_index < len(row) else None
        text = str(value or "").strip()
        if text and _relative_numeric_value(value) is None:
            return text
    return ""


def extract_relative_spreadsheet_fields(path):
    """Read numeric fields from loose, multi-row, and inconsistently headed workbooks."""
    candidates = []
    for sheet_name, rows in _raw_spreadsheet_rows(path).items():
        for row_index, row in enumerate(rows):
            for column_index, raw_value in enumerate(row):
                value = _relative_numeric_value(raw_value)
                if value is None:
                    continue
                row_label = _nearest_text_left(rows, row_index, column_index)
                column_label = _nearest_text_above(rows, row_index, column_index)
                aliases = []
                for alias in (
                    f"{row_label} {column_label}",
                    f"{sheet_name} {row_label} {column_label}",
                    f"{sheet_name} {row_label}",
                    f"{sheet_name} {column_label}",
                    row_label,
                    column_label,
                    sheet_name,
                ):
                    clean = re.sub(r"\s+", " ", alias).strip()
                    if clean and clean.casefold() not in {item.casefold() for item in aliases}:
                        aliases.append(clean)
                label = aliases[0] if aliases else f"{sheet_name} row {row_index + 1} column {column_index + 1}"
                candidates.append({
                    "id": f"{sheet_name}!R{row_index + 1}C{column_index + 1}",
                    "sheet": sheet_name,
                    "row": row_index + 1,
                    "column": column_index + 1,
                    "row_label": row_label,
                    "column_label": column_label,
                    "label": label,
                    "aliases": aliases or [label],
                    "value": value,
                })
    return candidates


def _relative_match_score(target, candidate):
    target_label = str(target.get("label") or "")
    target_context = str(target.get("context") or "")
    target_tokens = set(_relative_tokens(target_label))
    context_tokens = set(_relative_tokens(target_context))
    best = 0.0
    for alias in candidate["aliases"]:
        alias_tokens = set(_relative_tokens(alias))
        token_score = fuzz.token_set_ratio(" ".join(target_tokens), " ".join(alias_tokens)) if target_tokens and alias_tokens else 0
        text_score = fuzz.WRatio(normalize_label(target_label), normalize_label(alias)) if target_label else 0
        context_score = fuzz.token_set_ratio(" ".join(context_tokens), " ".join(alias_tokens)) if context_tokens and alias_tokens else 0
        overlap = len(target_tokens & alias_tokens)
        score = max(token_score, text_score) * 0.78 + context_score * 0.12 + min(10, overlap * 5)
        best = max(best, score)

    original = str(target.get("original") or "")
    value = float(candidate["value"])
    if "%" in original and (abs(value) <= 1.5 or abs(value) <= 100):
        best += 4
    if any(symbol in original for symbol in ("$", "€", "£")):
        best += 3
    return min(100.0, best)


def _target_group_key(target):
    location = target.get("location") or ()
    return tuple(location[:3]) if len(location) >= 3 else (target.get("slide"), target.get("kind"))


def _choose_relative_sheet(group, by_sheet):
    group_text = " ".join(str(target.get("label") or "") for target in group)
    best_sheet = None
    best_score = float("-inf")
    for sheet_name, candidates in by_sheet.items():
        sample_labels = " ".join(candidate["label"] for candidate in candidates[:40])
        lexical = fuzz.token_set_ratio(normalize_label(group_text), normalize_label(f"{sheet_name} {sample_labels}"))
        count_ratio = min(len(group), len(candidates)) / max(len(group), len(candidates), 1)
        score = lexical * 0.65 + count_ratio * 35
        if score > best_score:
            best_score = score
            best_sheet = sheet_name
    return best_sheet


def _different_relative_candidate(target, ranked):
    original = str(target.get("original") or "")
    for score, candidate in ranked:
        if format_new_value(original, candidate["value"]) != original:
            return score, candidate
    return ranked[0] if ranked else (0.0, None)


def match_targets_relative(targets, candidates):
    """Return a replacement for each comparable target using labels, context, then position."""
    if not targets or not candidates:
        return []
    by_sheet = {}
    for candidate in candidates:
        by_sheet.setdefault(candidate["sheet"], []).append(candidate)

    groups = {}
    for target in targets:
        groups.setdefault(_target_group_key(target), []).append(target)

    results = []
    for group in groups.values():
        selected_sheet = _choose_relative_sheet(group, by_sheet)
        sheet_candidates = by_sheet.get(selected_sheet, candidates)
        for index, target in enumerate(group):
            ranked = sorted(
                ((_relative_match_score(target, candidate), candidate) for candidate in candidates),
                key=lambda item: (-item[0], item[1]["sheet"], item[1]["row"], item[1]["column"]),
            )
            score, candidate = _different_relative_candidate(target, ranked[:25])
            method = "Similar variable"
            if score < 42:
                relative_index = 0 if len(group) == 1 else round(index * (len(sheet_candidates) - 1) / (len(group) - 1))
                positional = sheet_candidates[max(0, min(relative_index, len(sheet_candidates) - 1))]
                positional_ranked = [(score, positional)] + ranked
                score, candidate = _different_relative_candidate(target, positional_ranked)
                method = "Relative position"
            if candidate is None:
                continue
            new_text = format_new_value(target.get("original", ""), candidate["value"])
            if new_text == str(target.get("original") or ""):
                continue
            confidence = "High" if score >= 72 else ("Medium" if score >= 48 else "Review")
            result = dict(target)
            result.update({
                "matched_label": candidate["label"],
                "source_field": candidate["id"],
                "source_sheet": candidate["sheet"],
                "new_value": candidate["value"],
                "new_text_preview": new_text,
                "new_text": new_text,
                "score": round(score, 1),
                "match_method": method,
                "confidence": confidence,
            })
            results.append(result)
    return results


def extract_from_text(text):
    pairs = []
    for line in text.splitlines():
        line = line.strip()
        if not line:
            continue
        parts = re.split(r"[:=\t]|(?:\s-\s)", line, maxsplit=1)
        if len(parts) == 2:
            label, rest = parts
            val = extract_number_value(rest)
            if val is not None and label.strip():
                pairs.append((label.strip(), val))
    return pairs


def gather_source_pairs(data_path, pasted_text):
    source_pairs = []
    if data_path:
        fname = data_path.lower()
        if fname.endswith((".xlsx", ".xls", ".csv")):
            source_pairs += extract_from_spreadsheet(data_path)
        elif fname.endswith(".pptx"):
            _, targets = extract_targets_pptx(data_path)
            for t in targets:
                val = extract_number_value(t["original"])
                if val is not None and t["label"]:
                    source_pairs.append((t["label"], val))
    if pasted_text:
        source_pairs += extract_from_text(pasted_text)
    return source_pairs


# ---------------------------------------------------------------------------
# PPTX: target extraction / editing / snapshot
# ---------------------------------------------------------------------------

def slide_title(slide):
    for shape in slide.shapes:
        if shape.has_text_frame and shape.is_placeholder:
            try:
                if shape.placeholder_format.type is not None and "TITLE" in str(shape.placeholder_format.type):
                    return shape.text_frame.text.strip()
            except Exception:
                pass
    for shape in slide.shapes:
        if shape.has_text_frame and shape.text_frame.text.strip():
            return shape.text_frame.text.strip().split("\n")[0]
    return ""


def extract_targets_pptx(pptx_path):
    prs = Presentation(pptx_path)
    targets = []
    for s_idx, slide in enumerate(prs.slides):
        title = slide_title(slide)
        for sh_idx, shape in enumerate(slide.shapes):
            if shape.has_text_frame:
                all_paras = shape.text_frame.paragraphs
                for p_idx, para in enumerate(all_paras):
                    for r_idx, run in enumerate(para.runs):
                        if looks_numeric(run.text):
                            other_text = "".join(
                                rr.text for k, rr in enumerate(para.runs) if k != r_idx
                            ).strip()
                            if not other_text:
                                # check sibling paragraphs in the same text box
                                # (e.g. a KPI tile: value on one line, label below it)
                                sibling_text = " ".join(
                                    pp.text.strip() for pi, pp in enumerate(all_paras)
                                    if pi != p_idx and pp.text.strip()
                                ).strip()
                                other_text = sibling_text
                            label = other_text or title or shape.name
                            targets.append({
                                "id": f"s{s_idx}-sh{sh_idx}-p{p_idx}-r{r_idx}",
                                "kind": "run",
                                "slide": s_idx + 1,
                                "label": label,
                                "context": f"Slide {s_idx+1}: {title}" if title else f"Slide {s_idx+1}",
                                "original": run.text.strip(),
                                "location": (s_idx, sh_idx, "para", p_idx, r_idx),
                            })

            if shape.has_chart:
                chart = shape.chart
                chart_title = ""
                try:
                    if chart.has_title and chart.chart_title.text_frame.text.strip():
                        chart_title = chart.chart_title.text_frame.text.strip()
                except Exception:
                    pass
                try:
                    plot = chart.plots[0]
                    categories = [str(c) for c in plot.categories]
                    series_list = list(chart.series)
                    multi_series = len(series_list) > 1
                    for se_idx, series in enumerate(series_list):
                        values = list(series.values)
                        for c_idx, cat in enumerate(categories):
                            if c_idx >= len(values) or values[c_idx] is None:
                                continue
                            val = values[c_idx]
                            label = f"{series.name} {cat}".strip() if multi_series else cat
                            targets.append({
                                "id": f"s{s_idx}-sh{sh_idx}-chart-se{se_idx}-c{c_idx}",
                                "kind": "chart_point",
                                "slide": s_idx + 1,
                                "label": label,
                                "context": (f"Slide {s_idx+1}: {chart_title or title} (chart)"
                                            if (chart_title or title) else f"Slide {s_idx+1} (chart)"),
                                "original": (f"{val:g}" if isinstance(val, float) else str(val)),
                                "location": (s_idx, sh_idx, "chart", se_idx, c_idx),
                            })
                except Exception:
                    pass

            if shape.has_table:
                table = shape.table
                nrows, ncols = len(table.rows), len(table.columns)
                col_headers = [table.cell(0, c).text.strip() for c in range(ncols)]
                for r in range(1, nrows):
                    row_header = table.cell(r, 0).text.strip()
                    for c in range(1, ncols):
                        cell_text = table.cell(r, c).text.strip()
                        if looks_numeric(cell_text):
                            label = f"{row_header} {col_headers[c]}".strip()
                            targets.append({
                                "id": f"s{s_idx}-sh{sh_idx}-tbl-r{r}-c{c}",
                                "kind": "table_cell",
                                "slide": s_idx + 1,
                                "label": label or f"{shape.name} row{r} col{c}",
                                "context": (f"Slide {s_idx+1}: {title} (table)" if title else f"Slide {s_idx+1} (table)"),
                                "original": cell_text,
                                "location": (s_idx, sh_idx, "table", r, c),
                            })
    return prs, targets


def collect_all_pptx_text(prs):
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                texts.append(shape.text_frame.text)
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        texts.append(cell.text_frame.text)
            if shape.has_chart:
                try:
                    if shape.chart.has_title:
                        texts.append(shape.chart.chart_title.text_frame.text)
                except Exception:
                    pass
    return texts


def _build_period_pattern(replacements):
    if not replacements:
        return None, None
    keys_sorted = sorted(replacements.keys(), key=len, reverse=True)
    pattern = re.compile(r"\b(" + "|".join(re.escape(k) for k in keys_sorted) + r")\b", re.IGNORECASE)
    lookup = {k.upper(): v for k, v in replacements.items()}
    return pattern, lookup


def _replace_text(text, pattern, lookup):
    if not text or not pattern:
        return text
    return pattern.sub(lambda m: lookup.get(m.group(0).upper(), m.group(0)), text)


def apply_period_replacements_pptx(prs, replacements):
    pattern, lookup = _build_period_pattern(replacements)
    if not pattern:
        return

    def process_text_frame(tf):
        for para in tf.paragraphs:
            for run in para.runs:
                if run.text and pattern.search(run.text):
                    run.text = _replace_text(run.text, pattern, lookup)

    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                process_text_frame(shape.text_frame)
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        process_text_frame(cell.text_frame)
            if shape.has_chart:
                try:
                    if shape.chart.has_title:
                        process_text_frame(shape.chart.chart_title.text_frame)
                except Exception:
                    pass


def apply_edits_pptx(pptx_path, out_path, confirmed_matches, period_replacements=None):
    prs = Presentation(pptx_path)
    slides = list(prs.slides)
    chart_edits = {}
    pattern, lookup = _build_period_pattern(period_replacements or {})

    for m in confirmed_matches:
        loc = m["location"]
        s_idx, sh_idx = loc[0], loc[1]
        shape = list(slides[s_idx].shapes)[sh_idx]
        if loc[2] == "para":
            p_idx, r_idx = loc[3], loc[4]
            shape.text_frame.paragraphs[p_idx].runs[r_idx].text = m["new_text"]
        elif loc[2] == "table":
            r, c = loc[3], loc[4]
            cell = shape.table.cell(r, c)
            if cell.text_frame.paragraphs and cell.text_frame.paragraphs[0].runs:
                cell.text_frame.paragraphs[0].runs[0].text = m["new_text"]
                for extra in cell.text_frame.paragraphs[0].runs[1:]:
                    extra.text = ""
            else:
                cell.text = m["new_text"]
        elif loc[2] == "chart":
            se_idx, c_idx = loc[3], loc[4]
            numeric_val = extract_number_value(m["new_text"])
            if numeric_val is None:
                continue
            chart_edits.setdefault((s_idx, sh_idx), {})[(se_idx, c_idx)] = numeric_val

    # Rebuild data for every chart that had at least one confirmed numeric
    # edit OR needs its series/category names updated to the new period.
    charts_needing_period_update = set()
    if pattern:
        for slide_idx, slide in enumerate(slides):
            for sh_idx, shape in enumerate(slide.shapes):
                if not shape.has_chart:
                    continue
                chart = shape.chart
                try:
                    plot = chart.plots[0]
                    names_and_cats = [str(c) for c in plot.categories] + [s.name or "" for s in chart.series]
                    if any(pattern.search(t) for t in names_and_cats):
                        charts_needing_period_update.add((slide_idx, sh_idx))
                except Exception:
                    pass

    for (s_idx, sh_idx) in set(chart_edits.keys()) | charts_needing_period_update:
        shape = list(slides[s_idx].shapes)[sh_idx]
        chart = shape.chart
        plot = chart.plots[0]
        categories = [_replace_text(str(c), pattern, lookup) for c in plot.categories]
        new_chart_data = CategoryChartData()
        new_chart_data.categories = categories
        edits = chart_edits.get((s_idx, sh_idx), {})
        for se_idx, series in enumerate(chart.series):
            values = list(series.values)
            for c_idx in range(len(values)):
                if (se_idx, c_idx) in edits:
                    values[c_idx] = edits[(se_idx, c_idx)]
            new_name = _replace_text(series.name or "", pattern, lookup)
            new_chart_data.add_series(new_name, values)
        chart.replace_data(new_chart_data)

    apply_period_replacements_pptx(prs, period_replacements or {})
    ensure_chart_contrast(prs)
    prs.save(out_path)


def snapshot_pptx(path):
    """Everything needed to render a before/after comparison in the browser:
    per slide, every table (as a 2D grid) and every chart (categories/series/
    colors) as plain data."""
    prs = Presentation(path)
    slides_out = []
    for slide in prs.slides:
        title = slide_title(slide)
        tables, charts = [], []
        for shape in slide.shapes:
            if shape.has_table:
                t = shape.table
                rows = [[t.cell(r, c).text for c in range(len(t.columns))] for r in range(len(t.rows))]
                tables.append(rows)
            if shape.has_chart:
                chart = shape.chart
                try:
                    plot = chart.plots[0]
                    cats = [str(c) for c in plot.categories]
                    series_out = []
                    for series in chart.series:
                        vals = [float(v) if v is not None else 0.0 for v in series.values]
                        colors = []
                        for pt in series.points:
                            try:
                                colors.append("#" + str(pt.format.fill.fore_color.rgb))
                            except Exception:
                                colors.append(None)
                        series_out.append({"name": series.name or "", "values": vals, "colors": colors})
                    chart_type = "pie" if "PIE" in str(chart.chart_type) else (
                        "line" if "LINE" in str(chart.chart_type) else "bar"
                    )
                    charts.append({"type": chart_type, "categories": cats, "series": series_out})
                except Exception:
                    pass
        slides_out.append({"title": title or f"Slide", "tables": tables, "charts": charts})
    return slides_out


# ---------------------------------------------------------------------------
# XLSX: target extraction / editing / snapshot
# Excel charts are normally bound to live cell ranges, so simply updating
# cell values (and never touching chart XML) makes the embedded chart
# reflect the new numbers automatically, with all original formatting,
# when the file is opened in Excel.
# ---------------------------------------------------------------------------

def extract_targets_xlsx(path):
    wb = load_workbook(path, data_only=False)
    targets = []
    for ws in wb.worksheets:
        sheet_id = re.sub(r"[^A-Za-z0-9]+", "_", ws.title).strip("_") or "sheet"
        max_row, max_col = ws.max_row, ws.max_column
        if max_row < 2 or max_col < 2:
            continue
        headers = [ws.cell(row=1, column=c).value for c in range(1, max_col + 1)]
        for r in range(2, max_row + 1):
            row_label = ws.cell(row=r, column=1).value
            row_label = str(row_label).strip() if row_label is not None else ""
            for c in range(2, max_col + 1):
                cell = ws.cell(row=r, column=c)
                val = cell.value
                if isinstance(val, (int, float)) and not isinstance(val, bool):
                    col_header = headers[c - 1]
                    col_header = str(col_header).strip() if col_header is not None else ""
                    label = f"{row_label} {col_header}".strip()
                    original = f"{val:,.2f}" if (val != int(val)) else f"{int(val):,}"
                    targets.append({
                        "id": f"{sheet_id}-r{r}-c{c}",
                        "kind": "xlsx_cell",
                        "sheet": ws.title,
                        "label": label or col_header or row_label,
                        "context": ws.title,
                        "original": original,
                        "location": (ws.title, r, c),
                    })
    wb.close()
    return targets


def collect_all_xlsx_text(path):
    wb = load_workbook(path, data_only=False)
    texts = []
    for ws in wb.worksheets:
        texts.append(ws.title)
        for row in ws.iter_rows():
            for cell in row:
                if isinstance(cell.value, str):
                    texts.append(cell.value)
    wb.close()
    return texts


def apply_period_replacements_xlsx(wb, replacements):
    if not replacements:
        return
    keys_sorted = sorted(replacements.keys(), key=len, reverse=True)
    pattern = re.compile(r"\b(" + "|".join(re.escape(k) for k in keys_sorted) + r")\b", re.IGNORECASE)
    lookup = {k.upper(): v for k, v in replacements.items()}

    def sub_func(m):
        return lookup.get(m.group(0).upper(), m.group(0))

    for ws in wb.worksheets:
        for row in ws.iter_rows():
            for cell in row:
                if isinstance(cell.value, str) and pattern.search(cell.value):
                    cell.value = pattern.sub(sub_func, cell.value)


def apply_edits_xlsx(xlsx_path, out_path, confirmed_matches, period_replacements=None):
    wb = load_workbook(xlsx_path)
    for m in confirmed_matches:
        sheet, r, c = m["location"]
        ws = wb[sheet]
        ws.cell(row=r, column=c).value = m["new_value"]
    apply_period_replacements_xlsx(wb, period_replacements or {})
    wb.save(out_path)
    wb.close()


def snapshot_xlsx(path):
    """Per-sheet grid of values, plus a simple reconstructed chart (first
    text column as categories, first numeric column as values) so the
    browser can show a comparable visual. This is a rendering aid only —
    the real embedded Excel chart (with its original formatting) is left
    completely untouched in the actual file and updates itself in Excel."""
    wb = load_workbook(path, data_only=True)
    sheets_out = []
    for ws in wb.worksheets:
        max_row, max_col = min(ws.max_row, 60), min(ws.max_column, 12)
        rows = []
        for r in range(1, max_row + 1):
            rows.append([ws.cell(row=r, column=c).value for c in range(1, max_col + 1)])

        chart = None
        if len(rows) >= 2 and len(rows[0]) >= 2:
            cats, vals = [], []
            for r in rows[1:]:
                label = r[0]
                num = None
                for cell_val in r[1:]:
                    if isinstance(cell_val, (int, float)) and not isinstance(cell_val, bool):
                        num = float(cell_val)
                        break
                if label is not None and num is not None:
                    cats.append(str(label))
                    vals.append(num)
            if cats and vals:
                chart = {"type": "bar", "categories": cats, "series": [{"name": "Value", "values": vals, "colors": []}]}

        sheets_out.append({"title": ws.title, "tables": [rows], "charts": [chart] if chart else []})
    wb.close()
    return sheets_out


# ---------------------------------------------------------------------------
# ROUTES
# ---------------------------------------------------------------------------

def _session_dir(sid):
    d = os.path.join(SESSIONS_DIR, sid)
    if not os.path.isdir(d):
        abort(404)
    return d


def _save_session_meta(sess_dir, **kwargs):
    path = os.path.join(sess_dir, "meta.json")
    data = {}
    if os.path.exists(path):
        with open(path) as f:
            data = json.load(f)
    data.update(kwargs)
    with open(path, "w") as f:
        json.dump(data, f)
    return data


def _load_session_meta(sess_dir):
    path = os.path.join(sess_dir, "meta.json")
    if not os.path.exists(path):
        return {}
    with open(path) as f:
        return json.load(f)


# ---------------------------------------------------------------------------
# BLOB-BACKED SESSION PERSISTENCE (Vercel only)
#
# Vercel Functions don't share a persistent /tmp across invocations, so a
# session directory created in one request may not exist when the next
# request for the same session arrives. These helpers mirror a session's
# local working directory to Vercel Blob storage so it survives regardless
# of which container ends up serving each request. Locally, or when Blob
# storage isn't configured, these are no-ops and behavior is unchanged.
# ---------------------------------------------------------------------------

def _blob_prefix(sid):
    return f"sessions/{sid}/"


def _hydrate_session(sid):
    """Pull a session's files from Blob storage into local disk if missing."""
    if not (IS_VERCEL and blob_storage.ENABLED):
        return
    local_dir = os.path.join(SESSIONS_DIR, sid)
    marker = os.path.join(local_dir, ".hydrated")
    if os.path.exists(marker):
        return  # this container already has this session's files
    # A page can trigger a burst of ~20-40 simultaneous requests (e.g. every
    # slide thumbnail loading at once), each potentially landing on a fresh
    # container that all try to list the same Blob prefix in the same
    # instant. A little jitter spreads that burst out instead of every
    # container hammering the Blob API at the exact same millisecond.
    time.sleep(random.uniform(0, 0.2))
    if os.path.exists(marker):
        return  # another request in this container finished while we waited
    try:
        blobs = blob_storage.list_prefix(_blob_prefix(sid))
    except Exception:
        return  # listing failed even after storage.py's own retries; give up
    if not blobs:
        return  # unknown session id; normal 404 handling takes over
    prefix = _blob_prefix(sid)
    for item in blobs:
        rel = item.get("pathname", "")[len(prefix):]
        if not rel:
            continue
        local_path = os.path.join(local_dir, rel)
        if not os.path.exists(local_path):
            try:
                blob_storage.download_to(item["url"], local_path)
            except Exception:
                # One file failing to download must not block the rest of
                # the session from hydrating.
                continue
    os.makedirs(local_dir, exist_ok=True)
    open(marker, "w").close()


def _persist_session(sid, _is_retry=False):
    """Mirror every file in a session's local directory up to Blob storage."""
    if not (IS_VERCEL and blob_storage.ENABLED):
        return
    local_dir = os.path.join(SESSIONS_DIR, sid)
    if not os.path.isdir(local_dir):
        return
    prefix = _blob_prefix(sid)
    attempted = 0
    succeeded = 0
    for root, _dirs, files in os.walk(local_dir):
        for name in files:
            if name == ".hydrated":
                continue
            local_path = os.path.join(root, name)
            rel = os.path.relpath(local_path, local_dir).replace(os.sep, "/")
            attempted += 1
            try:
                blob_storage.put_file(prefix + rel, local_path)
                succeeded += 1
            except Exception:
                # One file failing to upload (even after storage.py's own
                # retries) must not stop the rest of the session's files
                # from persisting. Losing one slide image is recoverable;
                # aborting the whole batch silently drops many.
                continue
    # If literally every file failed, this almost certainly isn't a
    # one-off blip on a single file -- it's sustained load on Blob's API
    # (e.g. this session's persist landing on top of another burst of
    # requests already in flight). A short backoff and one full retry of
    # the whole batch clears this in practice; without it, a session can
    # end up with nothing persisted at all, and every later request for
    # it 404s no matter what it asks for.
    if attempted > 0 and succeeded == 0 and not _is_retry:
        time.sleep(2.0)
        _persist_session(sid, _is_retry=True)


def _hydrate_single_file(sid, rel_path):
    """Download exactly one file from a session, without pulling the rest.

    Slide-image routes only ever need one specific file per request. If a
    page fires 20 parallel image requests, using the bulk _hydrate_session
    for each one means 20 containers each downloading up to 20 files --
    a burst of ~400 Blob operations in the same instant, which is very
    likely what pushes the API into transient failures. Fetching only the
    one needed file per request keeps that burst to ~20 operations instead.
    """
    if not (IS_VERCEL and blob_storage.ENABLED):
        return
    local_dir = os.path.join(SESSIONS_DIR, sid)
    local_path = os.path.join(local_dir, rel_path)
    if os.path.exists(local_path):
        return
    # Same reasoning as _hydrate_session: a page load can fire ~20 of these
    # in the same instant, each landing on a different container. Jitter
    # spreads the resulting Blob calls out instead of all firing at once.
    time.sleep(random.uniform(0, 0.4))
    if os.path.exists(local_path):
        return  # another request in this container fetched it while we waited
    exact_pathname = _blob_prefix(sid) + rel_path
    try:
        blobs = blob_storage.list_prefix(exact_pathname)
    except Exception:
        return
    for item in blobs:
        if item.get("pathname") == exact_pathname:
            try:
                blob_storage.download_to(item["url"], local_path)
            except Exception:
                pass
            return


@app.before_request
def _hydrate_session_from_blob():
    # slide_image and editor_slide_image hydrate only their one needed file
    # (see the explicit calls in those routes) rather than the whole session.
    if request.endpoint in ("slide_image", "editor_slide_image"):
        return
    sid = (request.view_args or {}).get("sid")
    if sid:
        _hydrate_session(sid)


@app.after_request
def _persist_session_to_blob(response):
    sid = (request.view_args or {}).get("sid")
    if sid and request.method == "POST" and response.status_code < 400:
        _persist_session(sid)
    return response


@app.route("/", methods=["GET"])
def index():
    return render_template("index.html")


@app.route("/health", methods=["GET"])
def health():
    return jsonify({"status": "ok", "service": "Deck Refresh"})


def _replacement_sheet_count(data_path):
    suffix = os.path.splitext(data_path)[1].lower()
    if suffix == ".csv":
        return 1
    if suffix in {".xlsx", ".xlsm"}:
        workbook = load_workbook(data_path, read_only=True, data_only=True)
        try:
            return len(workbook.sheetnames)
        finally:
            workbook.close()
    if suffix == ".xls":
        return len(pd.ExcelFile(data_path).sheet_names)
    return 1


def _apply_relative_matches_safely(original_path, pending_path, matches):
    """Apply compatible native edits and skip only an unsupported object group."""
    if not matches:
        shutil.copy2(original_path, pending_path)
        return []
    non_chart = [match for match in matches if match.get("kind") != "chart_point"]
    chart_groups = {}
    for match in matches:
        if match.get("kind") != "chart_point":
            continue
        location = match.get("location") or ()
        chart_groups.setdefault(tuple(location[:2]), []).append(match)
    batches = ([non_chart] if non_chart else []) + list(chart_groups.values())

    applied = []
    with tempfile.TemporaryDirectory(prefix="deck_refresh_relative_") as temp_dir:
        current = str(original_path)
        version = 0
        for batch in batches:
            version += 1
            candidate_path = os.path.join(temp_dir, f"relative_{version}.pptx")
            try:
                apply_edits_pptx(current, candidate_path, batch, {})
            except Exception:
                if any(match.get("kind") == "chart_point" for match in batch):
                    continue
                for match in batch:
                    version += 1
                    single_path = os.path.join(temp_dir, f"relative_{version}.pptx")
                    try:
                        apply_edits_pptx(current, single_path, [match], {})
                    except Exception:
                        continue
                    current = single_path
                    applied.append(match)
                continue
            current = candidate_path
            applied.extend(batch)
        shutil.copy2(current, pending_path)
    return applied


def _generic_replace_deck_1to1(original_path, data_path, pending_path, data_filename):
    """Match arbitrary native PPT data objects to workbook labels without moving objects."""
    original_path = str(original_path)
    data_path = str(data_path)
    pending_path = str(pending_path)
    prs, targets = extract_targets_pptx(original_path)
    candidates = extract_relative_spreadsheet_fields(data_path)
    confirmed = match_targets_relative(targets, candidates)

    confirmed = _apply_relative_matches_safely(original_path, pending_path, confirmed)
    validation = validate_replacement_structure(original_path, pending_path)
    if not validation.get("passed"):
        raise ReplacementError("The generic replacement changed the presentation structure, so Deck Refresh rejected it.")

    structure = inspect_deck_structure(original_path)
    kind_counts = Counter(item.get("kind") for item in confirmed)
    return {
        "profile": "Generic label-matched deck",
        "source_sheets": _replacement_sheet_count(data_path),
        "source_rows": len(candidates),
        "text_updates": kind_counts.get("run", 0),
        "chart_points": kind_counts.get("chart_point", 0),
        "table_cells": kind_counts.get("table_cell", 0),
        "total_updates": len(confirmed),
        "slide_count": structure["slide_count"],
        "charts": structure["charts"],
        "tables": structure["tables"],
        "geometry_preserved": True,
        "theme_preserved": True,
        "match_mode": "label_and_structure",
        "data_filename": data_filename,
        "relative_matches": [
            {
                "slide": match.get("slide"),
                "kind": match.get("kind"),
                "target_label": match.get("label"),
                "original": match.get("original"),
                "replacement": match.get("new_text"),
                "source_label": match.get("matched_label"),
                "source_field": match.get("source_field"),
                "source_sheet": match.get("source_sheet"),
                "method": match.get("match_method"),
                "confidence": match.get("confidence"),
                "score": match.get("score"),
            }
            for match in confirmed
        ],
        "analysis_message": (
            f"Deck Refresh proposed {len(confirmed)} relative replacements across {len({candidate['sheet'] for candidate in candidates})} worksheets."
            if confirmed else
            "Deck Refresh completed structural validation. No comparable numeric fields were present, so it proposed no value changes."
        ),
    }


def _attach_relative_sources(mapping_summary, relative_matches):
    by_type = {
        "chart_point": "Chart",
        "table_cell": "Table",
        "run": "Text",
    }
    unused = set(range(len(mapping_summary.get("entries", []))))
    for match in relative_matches or []:
        wanted_type = by_type.get(match.get("kind"))
        slide = int(match.get("slide") or 0)
        replacement = str(match.get("replacement") or "")
        candidates = []
        for index in unused:
            entry = mapping_summary["entries"][index]
            entry_type = entry.get("object_type")
            type_match = entry_type == wanted_type or (wanted_type == "Text" and entry_type in {"Text", "Figure"})
            if entry.get("slide") != slide or not type_match:
                continue
            value_match = bool(replacement and replacement in str(entry.get("new_value") or ""))
            candidates.append((0 if value_match else 1, index))
        if not candidates:
            continue
        _, selected = min(candidates)
        unused.remove(selected)
        entry = mapping_summary["entries"][selected]
        entry.update({
            "source_label": match.get("source_label"),
            "source_field": match.get("source_field"),
            "source_sheet": match.get("source_sheet"),
            "match_method": match.get("method"),
            "match_confidence": match.get("confidence"),
            "match_score": match.get("score"),
        })
    mapping_summary["relative_match_count"] = len(relative_matches or [])
    mapping_summary["review_confidence_count"] = sum(
        1 for match in relative_matches or [] if match.get("confidence") == "Review"
    )
    return mapping_summary


@app.route("/replace1to1", methods=["POST"])
def replace_one_to_one():
    presentation_file = request.files.get("primary_file")
    data_file = request.files.get("data_file")
    if not presentation_file or not presentation_file.filename:
        flash("Upload the PowerPoint presentation to refresh.")
        return redirect(url_for("index"))
    if not presentation_file.filename.lower().endswith(".pptx"):
        flash("The 1:1 refresh requires a .pptx presentation.")
        return redirect(url_for("index"))
    if not data_file or not data_file.filename:
        flash("Upload the matching Excel or CSV replacement file.")
        return redirect(url_for("index"))
    data_suffix = os.path.splitext(data_file.filename)[1].lower()
    if data_suffix not in {".xlsx", ".xlsm", ".xls", ".csv"}:
        flash("Replacement data must be an .xlsx, .xlsm, .xls, or .csv file.")
        return redirect(url_for("index"))

    sid = uuid.uuid4().hex[:12]
    sess_dir = os.path.join(SESSIONS_DIR, sid)
    os.makedirs(sess_dir, exist_ok=True)
    original_path = os.path.join(sess_dir, "original.pptx")
    pending_path = os.path.join(sess_dir, "pending_updated.pptx")
    safe_data_name = secure_filename(data_file.filename) or f"replacement{data_suffix}"
    data_path = os.path.join(sess_dir, safe_data_name)
    presentation_file.save(original_path)
    data_file.save(data_path)

    try:
        try:
            report = replace_deck_1to1(
                original_path,
                data_path,
                pending_path,
                data_filename=data_file.filename,
            )
        except ReplacementError:
            report = _generic_replace_deck_1to1(
                original_path,
                data_path,
                pending_path,
                data_file.filename,
            )
        mapping_summary = compare_deck_replacements(original_path, pending_path)
        if report.get("relative_matches") is not None:
            mapping_summary = _attach_relative_sources(mapping_summary, report.get("relative_matches"))
        if not mapping_summary.get("entries") and report.get("profile") != "Generic label-matched deck":
            raise ReplacementError("The replacement data did not change any figures, chart points, table cells, or text.")
        deck_structure = inspect_deck_structure(original_path, mapping_summary.get("entries", []))
        structure_validation = validate_replacement_structure(original_path, pending_path)
        if not structure_validation.get("passed"):
            raise ReplacementError("The proposed replacement failed presentation structure validation.")
    except ReplacementError as exc:
        shutil.rmtree(sess_dir, ignore_errors=True)
        flash(str(exc))
        return redirect(url_for("index"))
    except Exception:
        shutil.rmtree(sess_dir, ignore_errors=True)
        flash("Deck Refresh could not complete a safe 1:1 replacement. The presentation was not changed.")
        return redirect(url_for("index"))

    with open(os.path.join(sess_dir, "replacement_mappings.json"), "w", encoding="utf-8") as handle:
        json.dump(mapping_summary, handle, indent=2)
    with open(os.path.join(sess_dir, "replacement_structure.json"), "w", encoding="utf-8") as handle:
        json.dump(
            {"deck": deck_structure, "validation": structure_validation},
            handle,
            indent=2,
        )

    _save_session_meta(
        sess_dir,
        file_type="pptx",
        original_filename=presentation_file.filename,
        data_filename=data_file.filename,
        applied_count=0,
        rendering_ok=False,
        slide_count=report.get("slide_count", 0),
        render_engine=None,
        replacement_mode=True,
        replacement_report=report,
        replacement_mapping_summary={key: value for key, value in mapping_summary.items() if key != "entries"},
        replacement_reviewed=False,
    )
    _persist_session(sid)
    return redirect(url_for("replacement_review", sid=sid))


@app.route("/replace1to1/review/<sid>", methods=["GET"])
def replacement_review(sid):
    sess_dir = _session_dir(sid)
    mapping_path = os.path.join(sess_dir, "replacement_mappings.json")
    structure_path = os.path.join(sess_dir, "replacement_structure.json")
    pending_path = os.path.join(sess_dir, "pending_updated.pptx")
    if not all(os.path.exists(path) for path in (mapping_path, structure_path, pending_path)):
        flash("That replacement review has expired. Upload the files again.")
        return redirect(url_for("index"))
    with open(mapping_path, encoding="utf-8") as handle:
        mapping_summary = json.load(handle)
    with open(structure_path, encoding="utf-8") as handle:
        structure_review = json.load(handle)
    entries = mapping_summary.get("entries", [])
    groups = []
    current = None
    for entry in entries:
        if current is None or current["slide"] != entry["slide"]:
            current = {
                "slide": entry["slide"],
                "slide_title": entry["slide_title"],
                "entries": [],
            }
            groups.append(current)
        current["entries"].append(entry)
    meta = _load_session_meta(sess_dir)
    return render_template(
        "replacement_review.html",
        sid=sid,
        groups=groups,
        summary=mapping_summary,
        deck_structure=structure_review.get("deck", {}),
        structure_validation=structure_review.get("validation", {}),
        report=meta.get("replacement_report", {}),
        original_filename=meta.get("original_filename", "presentation.pptx"),
        data_filename=meta.get("data_filename", "replacement data"),
    )


@app.route("/replace1to1/apply/<sid>", methods=["POST"])
def replacement_apply(sid):
    sess_dir = _session_dir(sid)
    original_path = os.path.join(sess_dir, "original.pptx")
    pending_path = os.path.join(sess_dir, "pending_updated.pptx")
    mapping_path = os.path.join(sess_dir, "replacement_mappings.json")
    if not all(os.path.exists(path) for path in (original_path, pending_path, mapping_path)):
        flash("That replacement review has expired. Upload the files again.")
        return redirect(url_for("index"))
    with open(mapping_path, encoding="utf-8") as handle:
        mapping_summary = json.load(handle)
    valid_ids = {int(entry["id"]) for entry in mapping_summary.get("entries", [])}
    approval = request.form.get("approval")
    if approval == "approve_all":
        selected_ids = valid_ids
    elif approval == "approve_selected":
        selected_ids = {
            int(value) for value in request.form.getlist("mapping_id")
            if str(value).isdigit() and int(value) in valid_ids
        }
        if not selected_ids:
            flash("Select at least one proposed change to apply.")
            return redirect(url_for("replacement_review", sid=sid))
    else:
        flash("Review the presentation structure, then approve the selected changes.")
        return redirect(url_for("replacement_review", sid=sid))

    updated_path = os.path.join(sess_dir, "updated.pptx")
    try:
        applied_count = apply_selected_deck_replacements(
            original_path,
            pending_path,
            mapping_summary.get("entries", []),
            selected_ids,
            updated_path,
        )
    except ReplacementError as exc:
        flash(str(exc))
        return redirect(url_for("replacement_review", sid=sid))

    meta = _load_session_meta(sess_dir)
    report = meta.get("replacement_report", {})
    rendering_ok = False
    render_engine = None
    slide_count = report.get("slide_count", 0)
    render_dir = os.path.join(sess_dir, "render")
    shutil.rmtree(render_dir, ignore_errors=True)
    original_images, original_engine = render_pptx_to_images(original_path, render_dir, "original")
    updated_images, updated_engine = render_pptx_to_images(updated_path, render_dir, "updated")
    if original_images and updated_images and len(original_images) == len(updated_images):
        rendering_ok = True
        slide_count = len(original_images)
        render_engine = original_engine if original_engine == updated_engine else f"{original_engine} / {updated_engine}"

    _save_session_meta(
        sess_dir,
        applied_count=applied_count,
        rendering_ok=rendering_ok,
        slide_count=slide_count,
        render_engine=render_engine,
        replacement_reviewed=True,
    )
    return redirect(url_for("result", sid=sid))


@app.route("/process", methods=["POST"])
def process():
    primary_file = request.files.get("primary_file")
    if not primary_file or primary_file.filename == "":
        flash("Please upload a PowerPoint (.pptx) or Excel (.xlsx) file.")
        return redirect(url_for("index"))

    fname = primary_file.filename.lower()
    if fname.endswith(".pptx"):
        file_type = "pptx"
    elif fname.endswith((".xlsx", ".xls")):
        file_type = "xlsx"
    else:
        flash("That file type isn't supported. Upload a .pptx or .xlsx file.")
        return redirect(url_for("index"))

    sid = uuid.uuid4().hex[:12]
    sess_dir = os.path.join(SESSIONS_DIR, sid)
    os.makedirs(sess_dir, exist_ok=True)

    primary_path = os.path.join(sess_dir, "original." + ("pptx" if file_type == "pptx" else "xlsx"))
    primary_file.save(primary_path)

    data_path = None
    data_file = request.files.get("data_file")
    if data_file and data_file.filename:
        data_path = os.path.join(sess_dir, "data_" + data_file.filename)
        data_file.save(data_path)

    pasted_text = request.form.get("pasted_text", "").strip()
    source_pairs = gather_source_pairs(data_path, pasted_text)

    if not source_pairs:
        flash("No usable data found in what you provided (Excel/CSV/text/PPTX). Please check the format.")
        shutil.rmtree(sess_dir, ignore_errors=True)
        return redirect(url_for("index"))

    if file_type == "pptx":
        prs, targets = extract_targets_pptx(primary_path)
        deck_texts = collect_all_pptx_text(prs)
    else:
        targets = extract_targets_xlsx(primary_path)
        deck_texts = collect_all_xlsx_text(primary_path)

    matches = match_targets_to_source(targets, source_pairs)
    period_replacements = compute_period_replacements(deck_texts, source_pairs, pasted_text)

    with open(os.path.join(sess_dir, "matches.json"), "w") as f:
        json.dump(matches, f)
    with open(os.path.join(sess_dir, "period_replacements.json"), "w") as f:
        json.dump(period_replacements, f)
    _save_session_meta(sess_dir, file_type=file_type, original_filename=primary_file.filename,
                        source_count=len(source_pairs))
    _persist_session(sid)

    return redirect(url_for("review", sid=sid))


@app.route("/review/<sid>", methods=["GET"])
def review(sid):
    sess_dir = _session_dir(sid)
    matches_path = os.path.join(sess_dir, "matches.json")
    if not os.path.exists(matches_path):
        flash("That session has expired. Please upload again.")
        return redirect(url_for("index"))

    with open(matches_path) as f:
        matches = json.load(f)
    period_path = os.path.join(sess_dir, "period_replacements.json")
    period_replacements = {}
    if os.path.exists(period_path):
        with open(period_path) as f:
            period_replacements = json.load(f)

    meta = _load_session_meta(sess_dir)
    matched = [m for m in matches if m["matched_label"]]
    unmatched = [m for m in matches if not m["matched_label"]]

    return render_template(
        "preview.html",
        sid=sid,
        matched=matched,
        unmatched=unmatched,
        source_count=meta.get("source_count", 0),
        period_replacements=period_replacements,
        file_type=meta.get("file_type", "pptx"),
    )


@app.route("/apply/<sid>", methods=["POST"])
def apply(sid):
    sess_dir = _session_dir(sid)
    matches_path = os.path.join(sess_dir, "matches.json")
    if not os.path.exists(matches_path):
        flash("That session has expired. Please upload again.")
        return redirect(url_for("index"))

    meta = _load_session_meta(sess_dir)
    file_type = meta.get("file_type", "pptx")
    primary_path = os.path.join(sess_dir, "original." + ("pptx" if file_type == "pptx" else "xlsx"))

    with open(matches_path) as f:
        matches = json.load(f)
    period_path = os.path.join(sess_dir, "period_replacements.json")
    period_replacements = {}
    if os.path.exists(period_path):
        with open(period_path) as f:
            period_replacements = json.load(f)

    confirmed = []
    for m in matches:
        if request.form.get(f"apply_{m['id']}") == "on" and m.get("new_text_preview"):
            override = request.form.get(f"value_{m['id']}", "").strip()
            new_text = override if override else m["new_text_preview"]
            entry = dict(m)
            entry["new_text"] = new_text
            if entry.get("new_value") is None:
                entry["new_value"] = extract_number_value(new_text)
            confirmed.append(entry)

    out_path = os.path.join(sess_dir, "updated." + ("pptx" if file_type == "pptx" else "xlsx"))
    if file_type == "pptx":
        apply_edits_pptx(primary_path, out_path, confirmed, period_replacements)
    else:
        apply_edits_xlsx(primary_path, out_path, confirmed, period_replacements)

    rendering_ok = False
    slide_count = 0
    if file_type == "pptx":
        render_dir = os.path.join(sess_dir, "render")
        shutil.rmtree(render_dir, ignore_errors=True)
        orig_images, orig_engine = render_pptx_to_images(primary_path, render_dir, "original")
        new_images, new_engine = render_pptx_to_images(out_path, render_dir, "updated")
        if orig_images and new_images and len(orig_images) == len(new_images):
            rendering_ok = True
            slide_count = len(orig_images)
            render_engine = orig_engine if orig_engine == new_engine else f"{orig_engine} / {new_engine}"
        else:
            render_engine = None
    else:
        render_engine = None

    _save_session_meta(
        sess_dir,
        applied_count=len(confirmed),
        rendering_ok=rendering_ok,
        slide_count=slide_count,
        render_engine=render_engine,
    )
    return redirect(url_for("result", sid=sid))


@app.route("/result/<sid>", methods=["GET"])
def result(sid):
    sess_dir = _session_dir(sid)
    meta = _load_session_meta(sess_dir)
    file_type = meta.get("file_type", "pptx")
    primary_path = os.path.join(sess_dir, "original." + ("pptx" if file_type == "pptx" else "xlsx"))
    updated_path = os.path.join(sess_dir, "updated." + ("pptx" if file_type == "pptx" else "xlsx"))
    if not os.path.exists(updated_path):
        flash("That session has expired. Please upload again.")
        return redirect(url_for("index"))

    rendering_ok = meta.get("rendering_ok", False)

    if rendering_ok:
        return render_template(
            "result.html",
            sid=sid,
            file_type=file_type,
            rendering_ok=True,
            slide_count=meta.get("slide_count", 0),
            applied_count=meta.get("applied_count", 0),
            render_engine=meta.get("render_engine", "PowerPoint renderer"),
            original_filename=meta.get("original_filename", "original.pptx"),
            updated_filename="updated_" + meta.get("original_filename", "presentation.pptx"),
            replacement_mode=meta.get("replacement_mode", False),
            replacement_report=meta.get("replacement_report", {}),
            data_filename=meta.get("data_filename", "replacement data"),
        )

    # Fallback: no LibreOffice/PyMuPDF available (or xlsx), show the
    # reconstructed data/chart comparison instead of real slide images.
    if file_type == "pptx":
        old_snap = snapshot_pptx(primary_path)
        new_snap = snapshot_pptx(updated_path)
    else:
        old_snap = snapshot_xlsx(primary_path)
        new_snap = snapshot_xlsx(updated_path)

    tabs = []
    for i in range(max(len(old_snap), len(new_snap))):
        old_s = old_snap[i] if i < len(old_snap) else {"title": "", "tables": [], "charts": []}
        new_s = new_snap[i] if i < len(new_snap) else {"title": "", "tables": [], "charts": []}
        tabs.append({
            "label": new_s.get("title") or old_s.get("title") or f"{'Slide' if file_type=='pptx' else 'Sheet'} {i+1}",
            "old": old_s,
            "new": new_s,
        })

    return render_template(
        "result.html",
        sid=sid,
        file_type=file_type,
        rendering_ok=False,
        tabs_json=json.dumps(tabs),
        applied_count=meta.get("applied_count", 0),
        soffice_missing=not _any_renderer_available(),
        replacement_mode=meta.get("replacement_mode", False),
        replacement_report=meta.get("replacement_report", {}),
        data_filename=meta.get("data_filename", "replacement data"),
    )


@app.route("/slide_image/<sid>/<which>/<int:n>", methods=["GET"])
def slide_image(sid, which, n):
    prefix = "original" if which == "original" else "updated"
    _hydrate_single_file(sid, os.path.join("render", f"{prefix}_{n}.png"))
    sess_dir = _session_dir(sid)
    path = os.path.join(sess_dir, "render", f"{prefix}_{n}.png")
    if not os.path.exists(path):
        abort(404)
    return send_file(path, mimetype="image/png")


@app.route("/download/<sid>/<which>", methods=["GET"])
def download(sid, which):
    sess_dir = _session_dir(sid)
    meta = _load_session_meta(sess_dir)
    file_type = meta.get("file_type", "pptx")
    ext = "pptx" if file_type == "pptx" else "xlsx"
    mimetype = (
        "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        if file_type == "pptx"
        else "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
    )
    if which == "original":
        path = os.path.join(sess_dir, f"original.{ext}")
        name = "original." + ext
    elif which == "updated":
        path = os.path.join(sess_dir, f"updated.{ext}")
        name = "updated." + ext
    else:
        abort(404)
    if not os.path.exists(path):
        abort(404)
    return send_file(path, as_attachment=True, download_name=name, mimetype=mimetype)


# ---------------------------------------------------------------------------
# CHAT-DRIVEN POWERPOINT EDITOR
# ---------------------------------------------------------------------------

_editor_locks = {}


def _editor_lock(sid):
    lock = _editor_locks.get(sid)
    if lock is None:
        lock = threading.Lock()
        _editor_locks[sid] = lock
    return lock


def _editor_version_path(sess_dir, version):
    return os.path.join(sess_dir, "editor_versions", f"version_{int(version):04d}.pptx")


def _editor_render_dir(sess_dir, version):
    return os.path.join(sess_dir, "editor_render", f"version_{int(version):04d}")


def _render_editor_version(sess_dir, version):
    """Render one editor version without destroying a usable preview first.

    Rendering happens in a temporary directory. The new images replace the
    version's current render directory only after every expected slide image
    exists. This prevents one failed export from leaving the browser in a
    permanently broken state.
    """
    pptx_path = _editor_version_path(sess_dir, version)
    final_dir = _editor_render_dir(sess_dir, version)
    slide_count = len(Presentation(pptx_path).slides)
    temp_dir = final_dir + f".tmp_{uuid.uuid4().hex[:8]}"
    shutil.rmtree(temp_dir, ignore_errors=True)
    os.makedirs(temp_dir, exist_ok=True)
    try:
        paths, engine = render_pptx_to_images(pptx_path, temp_dir, "slide")
        rendering_ok = bool(paths and len(paths) == slide_count and all(os.path.exists(path) for path in paths))
        if rendering_ok:
            backup_dir = final_dir + ".old"
            shutil.rmtree(backup_dir, ignore_errors=True)
            if os.path.isdir(final_dir):
                try:
                    os.replace(final_dir, backup_dir)
                except OSError:
                    shutil.rmtree(final_dir, ignore_errors=True)
            os.replace(temp_dir, final_dir)
            shutil.rmtree(backup_dir, ignore_errors=True)
            return {
                "slide_count": slide_count,
                "rendering_ok": True,
                "render_engine": engine,
                "render_help": "",
                "preview_revision": time.time_ns(),
            }
    finally:
        shutil.rmtree(temp_dir, ignore_errors=True)

    # Preserve an already-complete preview for this version if the refresh
    # attempt failed. This is mainly used by manual Retry preview requests.
    existing = [
        os.path.join(final_dir, f"slide_{number}.png")
        for number in range(1, slide_count + 1)
    ]
    has_existing = bool(existing and all(os.path.exists(path) for path in existing))
    return {
        "slide_count": slide_count,
        "rendering_ok": has_existing,
        "render_engine": "Previous preview" if has_existing else None,
        "render_help": "" if has_existing else _renderer_help_text(),
        "preview_revision": time.time_ns(),
    }


def _create_editor_session(source_path, original_filename):
    sid = uuid.uuid4().hex[:12]
    sess_dir = os.path.join(SESSIONS_DIR, sid)
    versions_dir = os.path.join(sess_dir, "editor_versions")
    os.makedirs(versions_dir, exist_ok=True)
    first_version = _editor_version_path(sess_dir, 0)
    shutil.copy2(source_path, first_version)
    Presentation(first_version)
    render_state = _render_editor_version(sess_dir, 0)
    _save_session_meta(
        sess_dir,
        editor=True,
        original_filename=os.path.basename(original_filename or "presentation.pptx"),
        editor_version=0,
        editor_max_version=0,
        **render_state,
    )
    with open(os.path.join(sess_dir, "editor_chat.json"), "w", encoding="utf-8") as f:
        json.dump([], f)
    _persist_session(sid)
    return sid


def _load_editor_chat(sess_dir):
    path = os.path.join(sess_dir, "editor_chat.json")
    if not os.path.exists(path):
        return []
    try:
        with open(path, encoding="utf-8") as f:
            data = json.load(f)
        return data if isinstance(data, list) else []
    except (OSError, json.JSONDecodeError):
        return []


def _save_editor_chat(sess_dir, history):
    with open(os.path.join(sess_dir, "editor_chat.json"), "w", encoding="utf-8") as f:
        json.dump(history[-60:], f, ensure_ascii=False)


def _editor_state(sess_dir, selected_slide=None):
    meta = _load_session_meta(sess_dir)
    version = int(meta.get("editor_version", 0))
    max_version = int(meta.get("editor_max_version", version))
    slide_count = int(meta.get("slide_count", 0))
    selected = int(selected_slide or 1)
    selected = max(1, min(selected, max(1, slide_count)))
    rendering_ok = bool(meta.get("rendering_ok", False))
    return {
        "version": version,
        "max_version": max_version,
        "slide_count": slide_count,
        "selected_slide": selected,
        "can_undo": version > 0,
        "can_redo": version < max_version,
        "rendering_ok": rendering_ok,
        "render_engine": meta.get("render_engine"),
        "render_help": "" if rendering_ok else (meta.get("render_help") or _renderer_help_text()),
        "preview_revision": int(meta.get("preview_revision") or 0),
        "openai_ready": bool(
            os.environ.get("OPENAI_API_KEY")
            or os.environ.get("AI_GATEWAY_API_KEY")
            or os.environ.get("VERCEL_OIDC_TOKEN")
        ),
    }


def _commit_editor_operations(sess_dir, operations):
    meta = _load_session_meta(sess_dir)
    current = int(meta.get("editor_version", 0))
    max_version = int(meta.get("editor_max_version", current))
    next_version = current + 1

    source = _editor_version_path(sess_dir, current)
    target = _editor_version_path(sess_dir, next_version)
    temp_target = target + ".tmp.pptx"
    try:
        result = apply_editor_operations(source, temp_target, operations)
        effective_count = len([op for op in operations if op.get("op") != "noop"])
        applied_count = len([item for item in result.get("applied", []) if item.get("op") != "noop"])
        if result.get("skipped") or applied_count != effective_count:
            result = {
                "applied": [],
                "skipped": result.get("skipped", []) or [{"reason": "The full edit transaction did not apply."}],
                "slide_count": int(meta.get("slide_count", 1)),
                "unchanged": True,
                "failed": True,
            }
        if result.get("unchanged"):
            # Do not create a fake version when every requested operation was
            # rejected or skipped. The current deck and preview stay untouched.
            result.update({
                "version": current,
                "rendering_ok": bool(meta.get("rendering_ok", False)),
                "render_engine": meta.get("render_engine"),
                "render_help": meta.get("render_help") or _renderer_help_text(),
            })
            return result

        for version in range(next_version, max_version + 1):
            try:
                os.remove(_editor_version_path(sess_dir, version))
            except OSError:
                pass
            shutil.rmtree(_editor_render_dir(sess_dir, version), ignore_errors=True)

        os.replace(temp_target, target)
    finally:
        try:
            os.remove(temp_target)
        except OSError:
            pass

    render_state = _render_editor_version(sess_dir, next_version)
    _save_session_meta(
        sess_dir,
        editor_version=next_version,
        editor_max_version=next_version,
        **render_state,
    )
    result.update(render_state)
    result["version"] = next_version
    return result


def _editor_selected_slide_after_result(selected_slide, result):
    selected = int(selected_slide)
    for item in result.get("applied", []) if result else []:
        op = str(item.get("op", ""))
        if op == "delete_slide":
            deleted = int(item.get("slide", selected))
            if deleted < selected:
                selected -= 1
        elif op == "move_slide":
            source = int(item.get("from_slide", selected))
            destination = int(item.get("to_slide", selected))
            if selected == source:
                selected = destination
            elif source < selected <= destination:
                selected -= 1
            elif destination <= selected < source:
                selected += 1
        elif op in {"add_slide", "duplicate_slide"} and item.get("slide") is not None:
            selected = int(item["slide"])
    return max(1, selected)


def _editor_result_message(plan, result, failure_message=None):
    planned_message = str((plan or {}).get("message") or "").strip()
    if (plan or {}).get("failed"):
        return planned_message or failure_message or "error cant do that"
    if result is None:
        return planned_message or failure_message or "error cant do that"

    applied = [item for item in result.get("applied", []) if item.get("op") != "noop"]
    skipped = result.get("skipped", []) or []
    if result.get("failed") or skipped or not applied:
        return failure_message or "error cant do that"
    return planned_message or f"Applied {len(applied)} edit{'s' if len(applied) != 1 else ''}."


@app.route("/editor/new", methods=["POST"])
def editor_new():
    temp_dir = tempfile.mkdtemp(prefix="deck_refresh_new_")
    source = os.path.join(temp_dir, "Untitled presentation.pptx")
    try:
        create_blank_deck(source)
        sid = _create_editor_session(source, "Untitled presentation.pptx")
    finally:
        shutil.rmtree(temp_dir, ignore_errors=True)
    return redirect(url_for("editor", sid=sid, builder=1))


@app.route("/editor/upload", methods=["POST"])
def editor_upload():
    uploaded = request.files.get("editor_file")
    if not uploaded or not uploaded.filename:
        flash("Choose a PowerPoint file to open in the editor.")
        return redirect(url_for("index"))
    if not uploaded.filename.lower().endswith(".pptx"):
        flash("The editor supports .pptx files.")
        return redirect(url_for("index"))

    temp_dir = tempfile.mkdtemp(prefix="deck_editor_upload_")
    temp_path = os.path.join(temp_dir, "upload.pptx")
    try:
        uploaded.save(temp_path)
        sid = _create_editor_session(temp_path, uploaded.filename)
    except Exception as exc:
        flash(f"The PowerPoint could not be opened: {exc}")
        return redirect(url_for("index"))
    finally:
        shutil.rmtree(temp_dir, ignore_errors=True)
    return redirect(url_for("editor", sid=sid))


@app.route("/editor/from-result/<source_sid>/<which>", methods=["GET"])
def editor_from_result(source_sid, which):
    source_dir = _session_dir(source_sid)
    meta = _load_session_meta(source_dir)
    if meta.get("file_type") != "pptx":
        abort(404)
    if which == "updated":
        source_path = os.path.join(source_dir, "updated.pptx")
        filename = "updated_" + meta.get("original_filename", "presentation.pptx")
    elif which == "original":
        source_path = os.path.join(source_dir, "original.pptx")
        filename = meta.get("original_filename", "presentation.pptx")
    else:
        abort(404)
    if not os.path.exists(source_path):
        abort(404)
    sid = _create_editor_session(source_path, filename)
    return redirect(url_for("editor", sid=sid))


@app.route("/editor/<sid>", methods=["GET"])
def editor(sid):
    sess_dir = _session_dir(sid)
    meta = _load_session_meta(sess_dir)
    if not meta.get("editor"):
        abort(404)
    state = _editor_state(sess_dir, request.args.get("slide", 1))
    return render_template(
        "editor.html",
        sid=sid,
        state=state,
        filename=meta.get("original_filename", "presentation.pptx"),
        history=_load_editor_chat(sess_dir),
        layouts=LAYOUTS,
        open_builder=request.args.get("builder") == "1",
    )


@app.route("/editor/wizard/<sid>", methods=["POST"])
def editor_wizard(sid):
    sess_dir = _session_dir(sid)
    if not _load_session_meta(sess_dir).get("editor"):
        abort(404)
    payload = request.get_json(silent=True) or {}
    message = str(payload.get("message", "")).strip()
    if not message:
        return jsonify({"ok": False, "error": "Describe what you need to communicate."}), 400
    return jsonify({"ok": True, "recommendation": wizard(message)})


@app.route("/editor/build/<sid>", methods=["POST"])
def editor_build(sid):
    sess_dir = _session_dir(sid)
    meta = _load_session_meta(sess_dir)
    if not meta.get("editor"):
        abort(404)
    layout = str(request.form.get("layout", "blank")).strip().casefold()
    chart_type = str(request.form.get("chart_type", "")).strip().casefold()
    data_mode = str(request.form.get("data_mode", "blank")).strip().casefold()
    title = str(request.form.get("title", "")).strip()
    content = str(request.form.get("content", "")).strip()
    selected_slide = int(request.form.get("selected_slide") or 1)
    placement = str(request.form.get("placement", "new")).strip().casefold()
    if placement not in {"new", "current"}:
        return jsonify({"ok": False, "error": "Choose New slide or Current slide."}), 400
    smart = request.form.get("smart") == "true" or layout == "smart_insert"
    upload = request.files.get("data_file")
    temp_dir = tempfile.mkdtemp(prefix="deck_refresh_builder_")
    data_path = None
    try:
        if upload and upload.filename:
            suffix = Path(upload.filename).suffix.casefold()
            if suffix not in {".xlsx", ".xls", ".xlsm", ".csv", ".tsv"}:
                return jsonify({"ok": False, "error": "Upload an Excel or CSV file."}), 400
            data_path = os.path.join(temp_dir, "source" + suffix)
            upload.save(data_path)
        if smart and not data_path:
            return jsonify({"ok": False, "error": "Smart Insert needs an Excel or CSV file."}), 400
        if data_mode == "upload" and not data_path:
            return jsonify({"ok": False, "error": "Choose an Excel or CSV file, or select Blank chart."}), 400
        state = _editor_state(sess_dir, selected_slide)
        selected_slide = max(1, min(selected_slide, state["slide_count"]))
        empty_starter = state["slide_count"] == 1 and int(state["version"]) == 0 and meta.get("original_filename") == "Untitled presentation.pptx"
        position = (1 if empty_starter else min(selected_slide + 1, state["slide_count"] + 1)) if placement == "new" else selected_slide
        operations, data_profile, chosen_layout = layout_operations(
            layout, position, title, content, data_path, smart,
            chart_type_override=chart_type or None,
        )
        if placement == "current":
            current_operations = []
            for operation in operations:
                if operation.get("op") == "add_slide":
                    continue
                operation = dict(operation)
                if operation.get("slide") == position:
                    operation["slide"] = selected_slide
                if operation.get("op") == "add_chart":
                    operation["auto_fit"] = True
                    operation.pop("x", None)
                    operation.pop("y", None)
                    operation.pop("width", None)
                    operation.pop("height", None)
                    if title and not str(operation.get("title") or "").strip():
                        operation["title"] = title
                current_operations.append(operation)
            if not any(operation.get("op") == "add_chart" for operation in current_operations):
                return jsonify({"ok": False, "error": "Current slide placement is available for chart options."}), 400
            operations = current_operations
        elif empty_starter:
            operations.append({"op": "delete_slide", "slide": 2})
        with _editor_lock(sid):
            result = _commit_editor_operations(sess_dir, operations)
        if result.get("failed") or result.get("skipped"):
            reasons = "; ".join(str(item.get("reason")) for item in result.get("skipped", []))
            return jsonify({"ok": False, "error": reasons or "The slide could not be created.",
                            "state": _editor_state(sess_dir, selected_slide)}), 400
        new_state = _editor_state(sess_dir, position)
        detected = ""
        if data_profile:
            series = ", ".join(str(value) for value in data_profile.get("numeric_columns", [])[:4])
            detected = f" Detected {len(data_profile['headers'])} columns. X axis: {data_profile['category_column']}."
            if series:
                detected += f" Y series: {series}."
        created_name = f"{chart_type} chart" if chart_type else chosen_layout.replace("_", " ")
        destination = f"on slide {selected_slide}" if placement == "current" else "on a new slide"
        blank = "blank " if data_mode != "upload" and not data_path else ""
        return jsonify({"ok": True, "message": f"Added a {blank}native editable {created_name} {destination}.{detected}",
                        "profile": data_profile, "state": new_state})
    except (ValueError, EditorError) as exc:
        return jsonify({"ok": False, "error": str(exc), "state": _editor_state(sess_dir, selected_slide)}), 400
    finally:
        shutil.rmtree(temp_dir, ignore_errors=True)


@app.route("/editor/chat/<sid>", methods=["POST"])
def editor_chat(sid):
    sess_dir = _session_dir(sid)
    meta = _load_session_meta(sess_dir)
    if not meta.get("editor"):
        abort(404)
    if request.mimetype and request.mimetype.startswith("multipart/form-data"):
        message = str(request.form.get("message", "")).strip()
        selected_slide = int(request.form.get("selected_slide") or 1)
        uploaded_files = request.files.getlist("attachments")
    else:
        payload = request.get_json(silent=True) or {}
        message = str(payload.get("message", "")).strip()
        selected_slide = int(payload.get("selected_slide") or 1)
        uploaded_files = []
    if not message:
        return jsonify({"ok": True, "message": "Type an edit request first.", "state": _editor_state(sess_dir, selected_slide)})

    attachment_paths = []
    if uploaded_files:
        attachment_dir = os.path.join(sess_dir, "editor_attachments")
        os.makedirs(attachment_dir, exist_ok=True)
        for uploaded in uploaded_files[:10]:
            if not uploaded or not uploaded.filename:
                continue
            safe_name = secure_filename(uploaded.filename) or f"attachment_{uuid.uuid4().hex}"
            target = os.path.join(attachment_dir, f"{uuid.uuid4().hex[:8]}_{safe_name}")
            uploaded.save(target)
            attachment_paths.append(target)

    with _editor_lock(sid):
        meta = _load_session_meta(sess_dir)
        version = int(meta.get("editor_version", 0))
        slide_count = int(meta.get("slide_count", 1))
        selected_slide = max(1, min(selected_slide, max(1, slide_count)))
        current_path = _editor_version_path(sess_dir, version)
        render_dir = _editor_render_dir(sess_dir, version)
        image_path = os.path.join(render_dir, f"slide_{selected_slide}.png")
        deck_image_paths = [
            os.path.join(render_dir, f"slide_{slide_number}.png")
            for slide_number in range(1, slide_count + 1)
            if os.path.exists(os.path.join(render_dir, f"slide_{slide_number}.png"))
        ]
        history = _load_editor_chat(sess_dir)
        history.append({"role": "user", "content": message})
        try:
            plan = plan_editor_edit(
                current_path,
                message,
                selected_slide,
                image_path if os.path.exists(image_path) else None,
                history,
                deck_image_paths,
                attachment_paths,
            )
            operations = plan.get("operations", [])
            effective_operations = [operation for operation in operations if operation.get("op") != "noop"]
            result = None
            failure_message = None
            if effective_operations:
                result = _commit_editor_operations(sess_dir, effective_operations)
                new_selected = _editor_selected_slide_after_result(selected_slide, result)
                state = _editor_state(sess_dir, min(new_selected, result.get("slide_count", new_selected)))
                applied = [item for item in result.get("applied", []) if item.get("op") != "noop"]
                skipped = result.get("skipped", []) or []
                if result.get("failed") or skipped or len(applied) != len(effective_operations):
                    reasons = [str(item.get("reason") or "An operation was skipped.") for item in skipped]
                    if len(applied) != len(effective_operations):
                        reasons.append(f"Only {len(applied)} of {len(effective_operations)} operations applied.")
                    failure_message = diagnose_editor_failure(
                        current_path,
                        message,
                        selected_slide,
                        history,
                        "; ".join(reasons) or "The final edit transaction did not apply completely.",
                        "Saving the validated edit plan to the working PowerPoint",
                    )
            else:
                state = _editor_state(sess_dir, selected_slide)
            assistant_message = _editor_result_message(plan, result, failure_message)
            history.append({"role": "assistant", "content": assistant_message})
            _save_editor_chat(sess_dir, history)
            return jsonify({
                "ok": True,
                "message": assistant_message,
                "operations": operations,
                "state": state,
            })
        except EditorError as exc:
            assistant_message = diagnose_editor_failure(
                current_path,
                message,
                selected_slide,
                history,
                str(exc),
                "Applying the requested PowerPoint change",
            )
            history.append({"role": "assistant", "content": assistant_message})
            _save_editor_chat(sess_dir, history)
            return jsonify({"ok": True, "message": assistant_message, "operations": [], "state": _editor_state(sess_dir, selected_slide)})
        except Exception as exc:
            assistant_message = diagnose_editor_failure(
                current_path,
                message,
                selected_slide,
                history,
                str(exc),
                "Completing the PowerPoint edit transaction",
            )
            history.append({"role": "assistant", "content": assistant_message})
            _save_editor_chat(sess_dir, history)
            return jsonify({"ok": True, "message": assistant_message, "operations": [], "state": _editor_state(sess_dir, selected_slide)})


@app.route("/editor/quick/<sid>", methods=["POST"])
def editor_quick(sid):
    """Execute inspector controls through the deterministic local compiler."""
    sess_dir = _session_dir(sid)
    meta = _load_session_meta(sess_dir)
    if not meta.get("editor"):
        abort(404)
    payload = request.get_json(silent=True) or {}
    message = str(payload.get("command", "")).strip()
    selected_slide = int(payload.get("selected_slide") or 1)
    if not message:
        return jsonify({"ok": False, "error": "This control has no edit command."}), 400

    with _editor_lock(sid):
        state = _editor_state(sess_dir, selected_slide)
        selected_slide = state["selected_slide"]
        current_path = _editor_version_path(sess_dir, state["version"])
        plan = guaranteed_local_plan(current_path, message, selected_slide)
        operations = [operation for operation in (plan or {}).get("operations", []) if operation.get("op") != "noop"]
        if not plan or plan.get("failed") or not operations:
            return jsonify({"ok": False, "error": (plan or {}).get("message") or "This edit is not available.",
                            "state": state}), 400
        result = _commit_editor_operations(sess_dir, operations)
        if result.get("failed") or result.get("skipped"):
            reasons = "; ".join(str(item.get("reason") or "Operation skipped") for item in result.get("skipped", []))
            return jsonify({"ok": False, "error": reasons or "The edit could not be applied.",
                            "state": _editor_state(sess_dir, selected_slide)}), 400
        new_selected = _editor_selected_slide_after_result(selected_slide, result)
        new_state = _editor_state(sess_dir, min(new_selected, result.get("slide_count", new_selected)))
        history = _load_editor_chat(sess_dir)
        history.extend([
            {"role": "user", "content": message},
            {"role": "assistant", "content": plan["message"]},
        ])
        _save_editor_chat(sess_dir, history)
        return jsonify({"ok": True, "message": plan["message"], "operations": operations, "state": new_state})


@app.route("/editor/action/<sid>", methods=["POST"])
def editor_action(sid):
    sess_dir = _session_dir(sid)
    meta = _load_session_meta(sess_dir)
    if not meta.get("editor"):
        abort(404)
    payload = request.get_json(silent=True) or {}
    action = str(payload.get("action", "")).strip().lower()
    selected_slide = int(payload.get("selected_slide") or 1)

    with _editor_lock(sid):
        state = _editor_state(sess_dir, selected_slide)
        selected_slide = state["selected_slide"]
        version = state["version"]
        max_version = state["max_version"]

        if action == "undo":
            if version > 0:
                target_version = version - 1
                _save_session_meta(
                    sess_dir,
                    editor_version=target_version,
                    **_render_editor_version(sess_dir, target_version),
                )
                message = "Undid the last edit."
            else:
                message = "The deck is already at its first saved version."
            return jsonify({"ok": True, "message": message, "state": _editor_state(sess_dir, selected_slide)})

        if action == "redo":
            if version < max_version:
                target_version = version + 1
                _save_session_meta(
                    sess_dir,
                    editor_version=target_version,
                    **_render_editor_version(sess_dir, target_version),
                )
                message = "Redid the edit."
            else:
                message = "There is no later saved version to restore."
            return jsonify({"ok": True, "message": message, "state": _editor_state(sess_dir, selected_slide)})

        if action == "retry_preview":
            render_state = _render_editor_version(sess_dir, version)
            _save_session_meta(sess_dir, **render_state)
            message = "Preview is ready." if render_state.get("rendering_ok") else render_state.get("render_help")
            return jsonify({"ok": True, "message": message, "state": _editor_state(sess_dir, selected_slide)})

        operations = None
        new_selected = selected_slide
        message = "Slide updated."
        if action == "add_slide":
            new_selected = min(selected_slide + 1, state["slide_count"] + 1)
            operations, _, _ = layout_operations("text", new_selected, "New slide")
            message = f"Created slide {new_selected} with editable content."
        elif action == "duplicate_slide":
            new_selected = min(selected_slide + 1, state["slide_count"] + 1)
            operations = [{"op": "duplicate_slide", "slide": selected_slide, "position": new_selected}]
            message = f"Duplicated slide {selected_slide} as slide {new_selected}."
        elif action == "delete_slide":
            if state["slide_count"] == 1:
                operations = [{"op": "clear_slide", "slide": 1, "preserve_branding": False, "preserve_title": False}]
                new_selected = 1
                message = "Cleared the only slide and kept a valid blank presentation."
            else:
                operations = [{"op": "delete_slide", "slide": selected_slide}]
                new_selected = max(1, min(selected_slide, state["slide_count"] - 1))
                message = f"Deleted slide {selected_slide}."
        elif action == "move_left" and selected_slide > 1:
            operations = [{"op": "move_slide", "from_slide": selected_slide, "to_slide": selected_slide - 1}]
            new_selected = selected_slide - 1
            message = f"Moved the slide to position {new_selected}."
        elif action == "move_right" and selected_slide < state["slide_count"]:
            operations = [{"op": "move_slide", "from_slide": selected_slide, "to_slide": selected_slide + 1}]
            new_selected = selected_slide + 1
            message = f"Moved the slide to position {new_selected}."
        elif action == "move_left":
            return jsonify({"ok": True, "message": "This slide is already first.", "state": state})
        elif action == "move_right":
            return jsonify({"ok": True, "message": "This slide is already last.", "state": state})
        else:
            return jsonify({"ok": False, "error": "That action is not available."}), 400

        try:
            result = _commit_editor_operations(sess_dir, operations)
            if result.get("failed") or result.get("skipped") or result.get("unchanged"):
                reasons = "; ".join(str(item.get("reason") or "Operation skipped") for item in result.get("skipped", []))
                return jsonify({"ok": False, "error": reasons or "The slide action did not change the deck.",
                                "state": _editor_state(sess_dir, selected_slide)}), 400
            return jsonify({"ok": True, "message": message, "state": _editor_state(sess_dir, new_selected)})
        except EditorError as exc:
            return jsonify({"ok": False, "error": str(exc), "state": _editor_state(sess_dir, selected_slide)}), 400


@app.route("/editor_slide_image/<sid>/<int:version>/<int:n>", methods=["GET"])
def editor_slide_image(sid, version, n):
    rel = os.path.join("editor_render", f"version_{int(version):04d}", f"slide_{n}.png")
    _hydrate_single_file(sid, rel)
    sess_dir = _session_dir(sid)
    path = os.path.join(_editor_render_dir(sess_dir, version), f"slide_{n}.png")
    if not os.path.exists(path):
        abort(404)
    return send_file(path, mimetype="image/png")


@app.route("/editor/download/<sid>", methods=["GET"])
def editor_download(sid):
    sess_dir = _session_dir(sid)
    meta = _load_session_meta(sess_dir)
    if not meta.get("editor"):
        abort(404)
    version = int(meta.get("editor_version", 0))
    path = _editor_version_path(sess_dir, version)
    filename = meta.get("original_filename", "presentation.pptx")
    stem, _ = os.path.splitext(filename)
    return send_file(
        path,
        as_attachment=True,
        download_name=f"{stem}_edited.pptx",
        mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation",
    )

def _open_browser():
    import webbrowser
    import sys as _sys
    time.sleep(1.2)
    url = "http://127.0.0.1:5050"
    chrome_candidates = []
    if _sys.platform == "darwin":
        chrome_candidates.append("open -a 'Google Chrome' %s")
    elif _sys.platform.startswith("win"):
        for p in (r"C:\Program Files\Google\Chrome\Application\chrome.exe",
                  r"C:\Program Files (x86)\Google\Chrome\Application\chrome.exe"):
            if os.path.exists(p):
                chrome_candidates.append(p.replace("\\", "\\\\") + " %s")
    else:
        chrome_candidates += ["google-chrome %s", "chromium-browser %s", "chromium %s"]
    for template in chrome_candidates:
        try:
            webbrowser.get(template).open(url)
            return
        except webbrowser.Error:
            continue
    webbrowser.open(url)


if __name__ == "__main__":
    import sys
    if os.environ.get("WERKZEUG_RUN_MAIN") != "true":
        threading.Thread(target=_open_browser, daemon=True).start()
    print("Starting Deck Refresh at http://127.0.0.1:5050")
    app.run(host="127.0.0.1", port=5050, debug=("--debug" in sys.argv))
