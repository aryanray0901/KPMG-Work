"""Layout library and spreadsheet inference for Deck Refresh's slide builder."""

from __future__ import annotations

import re
from pathlib import Path
from typing import Any

import pandas as pd
from pptx import Presentation
from pptx.util import Inches


LAYOUTS = [
    "blank", "title_slide", "section_divider", "executive_summary", "agenda", "text",
    "two_column", "comparison", "timeline", "process_flow", "org_chart", "kpi_dashboard",
    "table", "bar_chart", "line_chart", "pie_chart", "area_chart", "waterfall_chart",
    "scatter_plot", "heatmap", "gantt_chart", "map", "swot", "2x2_matrix", "pyramid",
    "funnel", "risk_matrix", "financial_statement", "roadmap", "image_gallery", "quote",
    "closing_thank_you",
]

CHART_LAYOUTS = {"bar_chart", "line_chart", "pie_chart", "area_chart", "waterfall_chart", "scatter_plot"}
CHART_TYPE_BY_LAYOUT = {
    "bar_chart": "bar", "line_chart": "line", "pie_chart": "pie",
    "area_chart": "area", "waterfall_chart": "waterfall", "scatter_plot": "scatter",
}
CHART_LAYOUT_BY_TYPE = {
    "column": "bar_chart", "bar": "bar_chart", "line": "line_chart", "pie": "pie_chart",
    "area": "area_chart", "waterfall": "waterfall_chart", "scatter": "scatter_plot",
}


def create_blank_deck(path: str) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    while prs.slides:
        node = prs.slides._sldIdLst[-1]
        prs.part.drop_rel(node.rId)
        del prs.slides._sldIdLst[-1]
    prs.slides.add_slide(prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[-1])
    prs.save(path)


def read_frame(path: str) -> pd.DataFrame:
    suffix = Path(path).suffix.casefold()
    if suffix in {".csv", ".tsv"}:
        frame = pd.read_csv(path, sep="\t" if suffix == ".tsv" else ",")
    elif suffix in {".xlsx", ".xls", ".xlsm"}:
        sheets = pd.read_excel(path, sheet_name=None)
        frames = [item.dropna(axis=0, how="all").dropna(axis=1, how="all") for item in sheets.values()]
        frames = [item for item in frames if not item.empty]
        if not frames:
            raise ValueError("The workbook has no usable rows.")
        frame = max(frames, key=lambda item: item.shape[0] * item.shape[1])
    else:
        raise ValueError("Upload an Excel or CSV file.")
    frame = frame.dropna(axis=0, how="all").dropna(axis=1, how="all").reset_index(drop=True)
    frame.columns = [str(value).strip() or f"Column {index + 1}" for index, value in enumerate(frame.columns)]
    if frame.empty:
        raise ValueError("The data file has no usable rows.")
    return frame


def numeric(series: pd.Series) -> pd.Series:
    if pd.api.types.is_numeric_dtype(series):
        return pd.to_numeric(series, errors="coerce")
    raw = series.astype(str)
    values = pd.to_numeric(raw.str.replace(r"[$£€,%()]", "", regex=True).str.replace(",", "", regex=False), errors="coerce")
    values.loc[raw.str.match(r"^\(.*\)$") & values.notna()] *= -1
    return values


def profile(frame: pd.DataFrame) -> dict[str, Any]:
    headers = list(frame.columns)
    numeric_columns, date_columns, percentages, currencies = [], [], [], []
    for column in headers:
        values = frame[column]
        raw = values.dropna().astype(str)
        header = str(column).strip().casefold()
        if numeric(values).notna().mean() >= 0.65:
            numeric_columns.append(column)
        if ((not raw.empty and raw.str.contains("%", regex=False).mean() >= 0.25)
                or any(term in header for term in ("%", "percent", "margin", "rate", "share"))):
            percentages.append(column)
        if ((not raw.empty and raw.str.contains(r"[$£€]", regex=True).mean() >= 0.25)
                or any(term in header for term in ("revenue", "sales", "cost", "profit", "budget", "expense", "spend"))):
            currencies.append(column)
        if not pd.api.types.is_numeric_dtype(values):
            parsed = pd.to_datetime(values, errors="coerce", format="mixed")
            if parsed.notna().mean() >= 0.65:
                date_columns.append(column)
    categories = [column for column in headers if column not in numeric_columns]
    category = date_columns[0] if date_columns else (categories[0] if categories else headers[0])
    totals = [int(index) for index, value in frame[category].items()
              if str(value).strip().casefold() in {"total", "grand total", "subtotal", "net total"}]
    names = " ".join(headers).casefold()
    header_words = set(re.findall(r"[a-z0-9]+", names))
    if date_columns and header_words.intersection({"task", "start", "end", "duration"}):
        suggested, data_kind = "gantt_chart", "project"
    elif header_words.intersection({"survey", "response", "question", "rating"}):
        suggested, data_kind = "kpi_dashboard", "survey"
    elif date_columns and numeric_columns:
        suggested, data_kind = "kpi_dashboard", "time_series"
    elif len(numeric_columns) >= 2:
        suggested, data_kind = "kpi_dashboard", "multi_metric"
    elif numeric_columns:
        suggested, data_kind = "bar_chart", "categorical"
    else:
        suggested, data_kind = "table", "records"
    insight = "The data is ready for an editable PowerPoint visual."
    if numeric_columns:
        values = numeric(frame.drop(index=totals, errors="ignore")[numeric_columns[0]]).dropna()
        if len(values) > 1 and values.iloc[0] != 0:
            change = (values.iloc[-1] / values.iloc[0] - 1) * 100
            insight = f"{numeric_columns[0]} {'increased' if change >= 0 else 'decreased'} {abs(change):.1f}% from the first to latest observation."
    return {"headers": headers, "category_column": category, "numeric_columns": numeric_columns,
            "date_columns": date_columns, "percentage_columns": percentages,
            "currency_columns": currencies, "total_rows": totals,
            "suggested_layout": suggested, "data_kind": data_kind, "insight": insight}


def matrix(frame: pd.DataFrame, rows: int = 18, columns: int = 9) -> list[list[str]]:
    view = frame.iloc[:rows, :columns].fillna("")
    return [[str(value) for value in view.columns]] + [[str(value) for value in row] for row in view.values.tolist()]


def chart_fields(frame: pd.DataFrame, info: dict[str, Any], chart_type: str | None = None) -> dict[str, Any]:
    view = frame.drop(index=info["total_rows"], errors="ignore").head(30)
    category_values = view[info["category_column"]]
    if info["category_column"] in info["date_columns"]:
        parsed_categories = pd.to_datetime(category_values, errors="coerce", format="mixed")
        categories = [
            value.strftime("%b %Y") if not pd.isna(value) else str(original)[:60]
            for original, value in zip(category_values, parsed_categories)
        ]
    else:
        categories = [str(value)[:60] for value in category_values.fillna("")]
    numeric_columns = list(info["numeric_columns"])
    if info["category_column"] in numeric_columns and len(numeric_columns) > 1:
        numeric_columns.remove(info["category_column"])
    if chart_type == "scatter":
        x_column = info["numeric_columns"][0] if info["numeric_columns"] else None
        y_columns = [column for column in info["numeric_columns"] if column != x_column]
        if x_column is None or not y_columns:
            raise ValueError("A scatter plot needs at least two numeric columns.")
        x_values = numeric(view[x_column]).fillna(0).astype(float).tolist()
        series = [
            {"name": str(column), "values": numeric(view[column]).fillna(0).astype(float).tolist()}
            for column in y_columns[:6]
        ]
        return {"categories": [str(value) for value in x_values], "x_values": x_values, "series": series}
    series = []
    for column in numeric_columns[:6]:
        series.append({"name": str(column), "values": numeric(view[column]).fillna(0).astype(float).tolist()})
    if not series:
        raise ValueError("A chart needs at least one numeric column.")
    if chart_type in {"pie", "waterfall"}:
        series = series[:1]
    result = {"categories": categories, "series": series}
    return result


def chart_number_format(info: dict[str, Any], fields: dict[str, Any], chart_type: str) -> str:
    """Choose an axis format from the plotted value column, not display strings."""
    if chart_type == "scatter":
        source = info["numeric_columns"][0] if info["numeric_columns"] else ""
    else:
        source = fields["series"][0]["name"] if fields.get("series") else ""
    if source in info.get("currency_columns", []):
        return '"$"#,##0'
    if source in info.get("percentage_columns", []):
        return "0.0%"
    return "#,##0"


def wizard(message: str) -> dict[str, str]:
    text = message.casefold()
    rules = [
        (("compare", "versus", " vs ", "difference"), "comparison"),
        (("timeline", "milestone", "roadmap"), "roadmap"),
        (("process", "workflow", "steps"), "process_flow"),
        (("risk", "impact", "likelihood"), "risk_matrix"),
        (("summary", "executive", "recommendation"), "executive_summary"),
        (("kpi", "dashboard", "performance", "sales"), "kpi_dashboard"),
        (("organization", "org chart", "reporting"), "org_chart"),
        (("quote", "testimonial"), "quote"),
    ]
    layout = next((name for words, name in rules if any(word in text for word in words)), "text")
    return {"layout": layout, "title": message.strip().rstrip(".")[:100] or "New slide",
            "reason": f"{layout.replace('_', ' ').title()} best supports this message."}


def _shape(slide: int, x: float, y: float, w: float, h: float, text: str,
           fill="FFFFFF", line="D8E2EE", color="172B4D", font=15, rounded=True):
    return {"op": "add_shape", "slide": slide, "shape_type": "rounded_rectangle" if rounded else "rectangle",
            "x": x, "y": y, "width": w, "height": h, "text": text, "fill_color": fill,
            "line_color": line, "font_color": color, "font_size": font, "alignment": "center"}


def layout_operations(layout: str, position: int, title: str = "", content: str = "",
                      data_path: str | None = None, smart: bool = False,
                      chart_type_override: str | None = None) -> tuple[list[dict[str, Any]], dict[str, Any] | None, str]:
    if layout not in LAYOUTS and layout != "smart_insert":
        raise ValueError("Unknown slide layout.")
    frame = read_frame(data_path) if data_path else None
    info = profile(frame) if frame is not None else None
    if smart and info:
        layout = info["suggested_layout"]
    requested_chart_type = str(chart_type_override or "").strip().casefold()
    if requested_chart_type:
        if requested_chart_type not in CHART_LAYOUT_BY_TYPE:
            raise ValueError("Choose a supported chart type.")
        if layout not in CHART_LAYOUTS:
            raise ValueError("Chart type selection only works with chart slides.")
    pretty = layout.replace("_", " ").title()
    supplied_title = title.strip()
    title = supplied_title if layout in CHART_LAYOUTS else (supplied_title or pretty)
    if layout in {"title_slide", "section_divider", "closing_thank_you"}:
        defaults = {
            "title_slide": ("Leadership briefing", "Prepared for executive discussion"),
            "section_divider": ("Strategic priorities", "Key decisions and supporting analysis"),
            "closing_thank_you": ("Thank you", "Questions, decisions, and next steps"),
        }
        default_title, default_subtitle = defaults[layout]
        shown = supplied_title or default_title
        return ([{"op": "add_slide", "position": position, "title": shown,
                  "subtitle": content or default_subtitle, "body": [], "background_color": "00338D",
                  "title_color": "FFFFFF", "subtitle_color": "FFFFFF", "title_size": 36}], info, layout)
    ops: list[dict[str, Any]] = [{"op": "add_slide", "position": position, "title": title,
                                  "subtitle": "" if layout in CHART_LAYOUTS else (info["insight"] if info else ""), "body": []}]
    slide = position
    items = [line.strip(" •-") for line in content.splitlines() if line.strip()]
    if layout == "blank":
        ops.append({"op": "clear_slide", "slide": slide, "preserve_branding": False, "preserve_title": False})
    elif frame is None and layout in CHART_LAYOUTS:
        chart_type = requested_chart_type or CHART_TYPE_BY_LAYOUT[layout]
        categories = ["", "", ""]
        fields: dict[str, Any] = {
            "categories": categories,
            "series": [{"name": "Series 1", "values": [0, 0, 0]}],
        }
        if chart_type == "scatter":
            fields["x_values"] = [1, 2, 3]
        chart_y = .22 if supplied_title else .10
        chart_height = .68 if supplied_title else .82
        ops.append({"op": "add_chart", "slide": slide, "chart_type": chart_type, "title": "",
                    **fields, "x": .06, "y": chart_y, "width": .88, "height": chart_height,
                    "series_colors": ["FFFFFF", "00A651", "BC204B"] if chart_type == "waterfall" else ["005EB8", "0091DA", "00A3A1", "483698"],
                    "show_legend": False})
    elif frame is not None and layout in CHART_LAYOUTS:
        chart_type = requested_chart_type or CHART_TYPE_BY_LAYOUT[layout]
        fields = chart_fields(frame, info, chart_type)
        number_format = chart_number_format(info, fields, chart_type)
        if supplied_title:
            chart_title = ""
        elif chart_type == "scatter":
            chart_title = f"{fields['series'][0]['name']} vs {info['numeric_columns'][0]}"
        else:
            chart_title = str(fields["series"][0]["name"])
        chart_y = .22 if supplied_title else .10
        chart_height = .68 if supplied_title else .82
        ops.append({"op": "add_chart", "slide": slide, "chart_type": chart_type, "title": chart_title,
                    **fields, "x": .06, "y": chart_y, "width": .88, "height": chart_height,
                    "series_colors": ["FFFFFF", "00A651", "BC204B"] if chart_type == "waterfall" else ["005EB8", "0091DA", "00A3A1", "483698", "BC204B"],
                    "show_legend": chart_type == "pie" or len(fields["series"]) > 1,
                    "number_format": number_format, "x_number_format": number_format})
    elif frame is not None and layout == "gantt_chart":
        task_col = info["category_column"]
        for index, task in enumerate(frame[task_col].fillna("").astype(str).head(5)):
            ops.append({"op": "add_textbox", "slide": slide, "text": task, "x": .05, "y": .22 + index * .075,
                        "width": .20, "height": .05, "font_size": 11, "font_color": "172B4D", "no_fill": True, "no_line": True})
            ops.append(_shape(slide, .27 + index * .045, .22 + index * .075, .30, .045, "",
                              ["00338D", "005EB8", "0091DA", "00A3A1", "483698"][index], None, "FFFFFF", 10))
        ops.append({"op": "add_table", "slide": slide, "data": matrix(frame, rows=7), "x": .05, "y": .64,
                    "width": .90, "height": .24, "header_fill": "00338D", "header_font_color": "FFFFFF",
                    "band_fill": "F3F6FA", "banded_rows": True, "font_size": 8})
    elif frame is not None and layout in {"table", "financial_statement"}:
        ops.append({"op": "add_table", "slide": slide, "data": matrix(frame), "x": .05, "y": .25,
                    "width": .90, "height": .62, "header_fill": "00338D", "header_font_color": "FFFFFF",
                    "band_fill": "F3F6FA", "banded_rows": True, "font_size": 10})
    elif frame is not None and layout == "kpi_dashboard":
        fields = chart_fields(frame, info) if info["numeric_columns"] else None
        if info.get("data_kind") == "survey" and fields:
            pie_fields = {**fields, "series": fields["series"][:1]}
            ops.append({"op": "add_chart", "slide": slide, "chart_type": "pie", "title": "Response mix",
                        **pie_fields, "x": .05, "y": .25, "width": .38, "height": .48,
                        "series_colors": ["005EB8", "0091DA", "00A3A1", "483698"]})
            ops.append({"op": "add_chart", "slide": slide, "chart_type": "bar", "title": "Response comparison",
                        **fields, "x": .45, "y": .25, "width": .35, "height": .48,
                        "series_colors": ["005EB8", "0091DA", "00A3A1", "483698"]})
            ops.append(_shape(slide, .82, .25, .14, .48, "Insights\n\n" + info["insight"], "F3F6FA", "D8E2EE", font=13))
        else:
            labels = info["numeric_columns"][:4] or info["headers"][:4]
            for index, label in enumerate(labels):
                values = numeric(frame[label]).dropna() if label in info["numeric_columns"] else pd.Series(dtype=float)
                value = f"{values.iloc[-1]:,.1f}" if not values.empty else str(frame[label].iloc[-1])
                ops.append(_shape(slide, .05 + index * .225, .22, .205, .16, f"{label}\n{value}", "FFFFFF", "D8E2EE", "00338D", 16))
            if fields:
                ops.append({"op": "add_chart", "slide": slide, "chart_type": "line" if info["date_columns"] else "column",
                            "title": "Performance", **fields, "x": .05, "y": .43, "width": .62, "height": .43,
                            "series_colors": ["005EB8", "0091DA", "00A3A1", "483698"]})
            ops.append(_shape(slide, .70, .43, .25, .43, "Executive takeaway\n\n" + info["insight"], "F3F6FA", "D8E2EE", font=14))
    elif layout == "executive_summary":
        cards = [
            "Performance\nRevenue is 8% ahead of plan",
            "Customer\nRetention improved to 94%",
            "Delivery\nThree milestones close this quarter",
            "Decision\nApprove phase two funding",
        ]
        for index, value in enumerate(cards):
            ops.append(_shape(slide, .05 + index * .225, .23, .205, .18, value, "FFFFFF", "D8E2EE", "00338D", 15))
        ops.append(_shape(slide, .05, .46, .60, .37,
                          "Executive summary\n\nMomentum remains positive. Delivery capacity is the main constraint to the next phase.",
                          "FFFFFF", "D8E2EE", "172B4D", 17))
        ops.append(_shape(slide, .68, .46, .27, .37,
                          "Recommended action\n\nApprove funding and assign one accountable delivery owner.",
                          "EAF6F4", "B9DDD7", "0F766E", 14))
    elif layout == "kpi_dashboard":
        kpis = [("Revenue", "$12.4M"), ("Growth", "+8.6%"), ("Margin", "31.2%"), ("Open risks", "3")]
        for index, (label, value) in enumerate(kpis):
            ops.append(_shape(slide, .05 + index * .225, .22, .205, .17, f"{label}\n{value}", "FFFFFF", "D8E2EE", "00338D", 17))
        ops.append({"op": "add_chart", "slide": slide, "chart_type": "line", "title": "Quarterly revenue trend",
                    "categories": ["Q1", "Q2", "Q3", "Q4"],
                    "series": [{"name": "Revenue ($M)", "values": [8.2, 9.5, 10.7, 12.4]}],
                    "x": .05, "y": .44, "width": .62, "height": .42,
                    "series_colors": ["005EB8"], "show_legend": False, "show_data_labels": True})
        ops.append(_shape(slide, .70, .44, .25, .42,
                          "Executive takeaway\n\nRevenue rose each quarter. Protect margin while scaling delivery.",
                          "F3F6FA", "D8E2EE", font=14))
    elif layout == "table":
        sample = [["Risk", "Owner", "Impact", "Status"], ["Delivery capacity", "Program lead", "High", "Mitigate"],
                  ["User adoption", "Change lead", "Medium", "Monitor"], ["Budget variance", "Finance", "Low", "On track"]]
        ops.append({"op": "add_table", "slide": slide, "data": sample, "x": .06, "y": .26, "width": .88,
                    "height": .55, "header_fill": "00338D", "header_font_color": "FFFFFF",
                    "banded_rows": True, "band_fill": "F3F6FA", "font_size": 11})
    elif layout == "financial_statement":
        statement = [["Income statement", "Q3", "Q4", "FY total"], ["Revenue", "$10.7M", "$12.4M", "$41.8M"],
                     ["Cost of sales", "($6.9M)", "($7.8M)", "($26.5M)"], ["Gross profit", "$3.8M", "$4.6M", "$15.3M"],
                     ["Operating expense", "($2.1M)", "($2.3M)", "($8.6M)"], ["EBITDA", "$1.7M", "$2.3M", "$6.7M"]]
        ops.append({"op": "add_table", "slide": slide, "data": statement, "x": .06, "y": .24, "width": .88,
                    "height": .60, "header_fill": "00338D", "header_font_color": "FFFFFF",
                    "banded_rows": True, "band_fill": "F3F6FA", "font_size": 11})
    elif layout == "gantt_chart":
        tasks = [("Mobilize team", .27, .18), ("Complete discovery", .34, .25), ("Design solution", .44, .24),
                 ("Pilot release", .57, .18), ("Scale rollout", .68, .22)]
        colors = ["00338D", "005EB8", "0091DA", "00A3A1", "483698"]
        for index, (task, start, width) in enumerate(tasks):
            y = .24 + index * .095
            ops.append({"op": "add_textbox", "slide": slide, "text": task, "x": .05, "y": y,
                        "width": .20, "height": .055, "font_size": 11, "font_color": "172B4D", "no_fill": True, "no_line": True})
            ops.append(_shape(slide, start, y, width, .052, f"W{index + 1} to W{index + 3}", colors[index], None, "FFFFFF", 9))
        ops.append(_shape(slide, .74, .24, .20, .43,
                          "Status\n\nDiscovery complete\nDesign on track\nPilot starts next",
                          "F3F6FA", "D8E2EE", "172B4D", 12))
    elif layout == "comparison":
        if frame is not None and info["numeric_columns"]:
            fields = chart_fields(frame, info)
            ops.append({"op": "add_chart", "slide": slide, "chart_type": "column", "title": "Comparison",
                        **fields, "x": .05, "y": .25, "width": .58, "height": .50,
                        "series_colors": ["005EB8", "0091DA", "00A3A1"]})
        else:
            ops.append({"op": "add_chart", "slide": slide, "chart_type": "column", "title": "Q2 versus Q3 sales",
                        "categories": ["North", "South", "West"],
                        "series": [{"name": "Q2", "values": [4.2, 3.6, 3.1]}, {"name": "Q3", "values": [4.8, 4.0, 3.7]}],
                        "x": .05, "y": .25, "width": .58, "height": .50,
                        "series_colors": ["005EB8", "00A3A1"], "show_legend": True, "show_data_labels": True})
        ops.append(_shape(slide, .67, .25, .28, .22,
                          "Summary\nQ3 sales increased across all regions.", "FFFFFF", "D8E2EE", "00338D", 14))
        ops.append(_shape(slide, .67, .52, .28, .23,
                          "Recommendation\nPrioritize the North region while preserving West momentum.",
                          "EAF6F4", "B9DDD7", "0F766E", 14))
    elif layout == "two_column":
        columns = [
            "Option A\n\nFaster rollout\nHigher near-term cost\nLower delivery risk",
            "Option B\n\nPhased rollout\nLower near-term cost\nLonger time to value",
        ]
        for index, value in enumerate(columns):
            ops.append(_shape(slide, .05 + index * .46, .24, .42, .56, value,
                              "FFFFFF", "D8E2EE", "00338D", 16))
    elif layout == "swot":
        values = [
            "Strengths\nStrong client retention\nExperienced delivery team",
            "Weaknesses\nLimited capacity\nManual reporting effort",
            "Opportunities\nNew market demand\nAutomation at scale",
            "Threats\nCompetitive pricing\nExecution delays",
        ]
        for index, value in enumerate(values):
            ops.append(_shape(slide, .05 + (index % 2) * .46, .24 + (index // 2) * .31,
                              .42, .27, value, "FFFFFF", "D8E2EE", "00338D", 14))
    elif layout in {"timeline", "roadmap", "process_flow"}:
        defaults = {
            "timeline": ["Q1\nPlan", "Q2\nBuild", "Q3\nLaunch", "Q4\nScale"],
            "roadmap": ["Now\nMobilize", "Next\nPilot", "Later\nExpand", "Scale\nOptimize"],
            "process_flow": ["Intake", "Analyze", "Design", "Deliver", "Measure"],
        }
        steps = (items or defaults[layout])[:5]
        for index, item in enumerate(steps):
            width = .20 if len(steps) == 4 else .16
            gap = .225 if len(steps) == 4 else .19
            ops.append(_shape(slide, .04 + index * gap, .38, width, .18, item,
                              ["00338D", "005EB8", "0091DA", "00A3A1", "483698"][index], None, "FFFFFF", 13))
    elif layout == "2x2_matrix":
        quadrants = ["Strategic bets\nHigh impact, high effort", "Quick wins\nHigh impact, low effort",
                     "Deprioritize\nLow impact, high effort", "Fill-ins\nLow impact, low effort"]
        for index, value in enumerate(quadrants):
            row, col = divmod(index, 2)
            ops.append(_shape(slide, .10 + col * .41, .25 + row * .29, .39, .27, value,
                              ["DCE6F7", "D8EFF4", "F6DDE5", "EAF6F4"][index], "FFFFFF", "172B4D", 13, rounded=False))
        ops.append({"op": "add_textbox", "slide": slide, "text": "IMPACT", "x": .02, "y": .43,
                    "width": .07, "height": .08, "font_size": 10, "font_color": "5D6B7A", "bold": True, "no_fill": True, "no_line": True})
        ops.append({"op": "add_textbox", "slide": slide, "text": "EFFORT", "x": .47, "y": .84,
                    "width": .12, "height": .05, "font_size": 10, "font_color": "5D6B7A", "bold": True, "no_fill": True, "no_line": True})
    elif layout == "risk_matrix":
        cells = ["Monitor\nVendor dependency", "Mitigate\nAdoption delay", "Accept\nMinor scope change", "Escalate\nDelivery capacity"]
        fills = ["FFF4D6", "FCE8E8", "EAF6F4", "F6DDE5"]
        for index, value in enumerate(cells):
            row, col = divmod(index, 2)
            ops.append(_shape(slide, .10 + col * .41, .25 + row * .29, .39, .27, value,
                              fills[index], "FFFFFF", "172B4D", 13, rounded=False))
        ops.append({"op": "add_textbox", "slide": slide, "text": "IMPACT", "x": .02, "y": .43,
                    "width": .07, "height": .08, "font_size": 10, "font_color": "5D6B7A", "bold": True, "no_fill": True, "no_line": True})
        ops.append({"op": "add_textbox", "slide": slide, "text": "LIKELIHOOD", "x": .44, "y": .84,
                    "width": .18, "height": .05, "font_size": 10, "font_color": "5D6B7A", "bold": True, "no_fill": True, "no_line": True})
    elif layout == "heatmap":
        values = [[42, 55, 67, 73, 81], [38, 49, 61, 70, 76], [31, 44, 58, 64, 72], [27, 39, 51, 60, 68]]
        for row in range(4):
            for col in range(5):
                value = values[row][col]
                fill = "DCE6F7" if value < 45 else "B8D6EC" if value < 60 else "79B9DB" if value < 72 else "005EB8"
                ops.append(_shape(slide, .08 + col * (.82 / 5), .25 + row * (.58 / 4), .78 / 5, .52 / 4,
                                  f"{value}%", fill, "FFFFFF", "FFFFFF" if value >= 60 else "172B4D", 12, rounded=False))
    elif layout in {"pyramid", "funnel"}:
        widths = [.28, .42, .56, .70]
        if layout == "funnel": widths.reverse()
        labels = (["Vision", "Priorities", "Initiatives", "Capabilities"] if layout == "pyramid"
                  else ["Awareness 100%", "Qualified 65%", "Proposal 38%", "Converted 22%"])
        for index, width in enumerate(widths):
            ops.append(_shape(slide, .5 - width / 2, .23 + index * .16, width, .12, labels[index],
                              ["00338D", "005EB8", "0091DA", "00A3A1"][index], None, "FFFFFF", 14))
    elif layout == "org_chart":
        ops.append(_shape(slide, .38, .23, .24, .13, "Program sponsor", "00338D", None, "FFFFFF", 17))
        for index, label in enumerate(["Strategy lead", "Operations lead", "Technology lead"]):
            ops.append(_shape(slide, .09 + index * .30, .51, .25, .14, label, "FFFFFF", "D8E2EE", "00338D", 15))
    elif layout == "quote":
        ops.append({"op": "add_textbox", "slide": slide, "text": content or "Clear priorities turn analysis into measurable action.",
                    "x": .12, "y": .32, "width": .76, "height": .34, "font_size": 27, "font_color": "00338D",
                    "bold": True, "alignment": "center", "no_fill": True, "no_line": True})
    elif layout == "image_gallery":
        gallery = [("01", "Market overview", "DCE6F7"), ("02", "Customer experience", "D8EFF4"), ("03", "Delivery team", "EAF6F4")]
        for index, (number, caption, fill) in enumerate(gallery):
            ops.append(_shape(slide, .05 + index * .31, .26, .28, .44, number, fill, "D8E2EE", "00338D", 28))
            ops.append(_shape(slide, .05 + index * .31, .72, .28, .10, caption, "FFFFFF", "D8E2EE", "172B4D", 13))
    elif layout in {"agenda", "text"}:
        defaults = (["Situation and objective", "Key findings", "Recommendations and next steps"] if layout == "agenda"
                    else ["The business is growing ahead of plan", "Delivery capacity remains the main constraint", "Leadership approval is required for phase two"])
        for index, item in enumerate((items or defaults)[:8]):
            ops.append(_shape(slide, .08, .24 + index * .085, .06, .055, str(index + 1), "0091DA", None, "FFFFFF", 11))
            ops.append({"op": "add_textbox", "slide": slide, "text": item, "x": .17, "y": .24 + index * .085,
                        "width": .72, "height": .06, "font_size": 17, "font_color": "172B4D", "no_fill": True, "no_line": True})
    elif layout == "map":
        regions = [("West\n$3.7M", .11, .43, "DCE6F7"), ("Central\n$4.0M", .29, .36, "79B9DB"),
                   ("North\n$4.8M", .44, .28, "005EB8"), ("South\n$3.2M", .43, .55, "D8EFF4")]
        for label, x, y, fill in regions:
            ops.append(_shape(slide, x, y, .17, .14, label, fill, "FFFFFF", "FFFFFF" if fill == "005EB8" else "172B4D", 13))
        ops.append(_shape(slide, .70, .25, .24, .56,
                          "Regional insights\n\nNorth leads revenue\nCentral grew 11%\nWest margin is strongest",
                          "FFFFFF", "D8E2EE", "00338D", 14))
    return ops, info, layout
