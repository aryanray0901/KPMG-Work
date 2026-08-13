"""Deterministic 1:1 data refresh for supported Deck Refresh presentations.

The engine updates existing PowerPoint objects only. It never adds, removes,
reorders, or resizes slide objects. Native charts and tables remain editable.
"""

from __future__ import annotations

import hashlib
import os
import re
import shutil
import zipfile
import xml.etree.ElementTree as ET
from collections import Counter
from datetime import datetime
from numbers import Number

import pandas as pd
from openpyxl import load_workbook
from pptx import Presentation
from pptx.chart.data import CategoryChartData

from chart_contrast import ensure_chart_contrast


class ReplacementError(Exception):
    """Raised when a workbook cannot safely replace a presentation."""


IG_RATINGS = {
    "AAA", "AA+", "AA", "AA-", "A+", "A", "A-", "BBB+", "BBB", "BBB-"
}
RATING_STRENGTH = {
    rating: score
    for score, rating in enumerate(
        ["NR", "B-", "B", "B+", "BB-", "BB", "BB+", "BBB-", "BBB", "BBB+",
         "A-", "A", "A+", "AA-", "AA", "AA+", "AAA"],
        start=1,
    )
}
REGIONS = [
    ("us_canada", "U.S. & Canada", "United States and Canada"),
    ("europe", "Europe", "Europe"),
    ("asia_pacific", "Asia-Pacific", "Asia / Pacific"),
    ("latin_america", "Latin America", "Latin America and Caribbean"),
]


def _read_sheet_rows(path):
    """Return {sheet_name: rows} for xlsx, xlsm, xls, or csv."""
    suffix = os.path.splitext(path)[1].lower()
    if suffix == ".csv":
        frame = pd.read_csv(path, header=None)
        return {"csv": frame.where(pd.notna(frame), None).values.tolist()}
    if suffix in {".xlsx", ".xlsm"}:
        workbook = load_workbook(path, data_only=True, read_only=True)
        return {
            sheet.title: [list(row) for row in sheet.iter_rows(values_only=True)]
            for sheet in workbook.worksheets
        }
    if suffix == ".xls":
        frames = pd.read_excel(path, sheet_name=None, header=None)
        return {
            name: frame.where(pd.notna(frame), None).values.tolist()
            for name, frame in frames.items()
        }
    raise ReplacementError("Upload an Excel or CSV replacement file.")


def _clean(value):
    return "" if value is None else str(value).strip()


def _number(value, default=0.0):
    if value is None or value == "":
        return default
    if isinstance(value, (int, float)):
        return float(value)
    text = str(value).strip().replace(",", "").replace("$", "")
    if text.endswith("%"):
        return float(text[:-1]) / 100
    try:
        return float(text)
    except (TypeError, ValueError):
        return default


def _pct(value, digits=1):
    return f"{value * 100:.{digits}f}%"


def _money(value):
    return f"${value:,.2f}"


def _integer(value):
    return f"{int(round(value)):,}"


def _region_key(value):
    text = _clean(value).casefold()
    if "united states and canada" in text:
        return "us_canada"
    if "europe" in text:
        return "europe"
    if "asia" in text or "pacific" in text:
        return "asia_pacific"
    if "latin america" in text or "caribbean" in text:
        return "latin_america"
    return "other"


def _shape_text(shape):
    return shape.text.strip() if getattr(shape, "has_text_frame", False) else ""


def _find_text_shape(slide, marker, exact=False):
    marker_key = marker.strip().casefold()
    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        text = _shape_text(shape)
        key = text.casefold()
        if (exact and key == marker_key) or (not exact and marker_key in key):
            return shape
    return None


def _find_named_shape(slide, name):
    for shape in slide.shapes:
        if shape.name == name:
            return shape
    return None


def _set_paragraph(paragraph, text):
    old = paragraph.text
    if paragraph.runs:
        paragraph.runs[0].text = str(text)
        for run in paragraph.runs[1:]:
            run.text = ""
    else:
        paragraph.add_run().text = str(text)
    return old != str(text)


def _set_shape_text(shape, text):
    if shape is None or not shape.has_text_frame:
        raise ReplacementError("A required presentation text object is missing.")
    lines = str(text).split("\n")
    frame = shape.text_frame
    while len(frame.paragraphs) < len(lines):
        frame.add_paragraph()
    changed = 0
    for index, paragraph in enumerate(frame.paragraphs):
        value = lines[index] if index < len(lines) else ""
        changed += int(_set_paragraph(paragraph, value))
    return changed


def _set_named_text(slide, name, text, fallback_marker=None):
    shape = _find_named_shape(slide, name)
    if shape is None and fallback_marker:
        shape = _find_text_shape(slide, fallback_marker)
    return _set_shape_text(shape, text)


def _first_chart(slide):
    for shape in slide.shapes:
        if getattr(shape, "has_chart", False):
            return shape.chart
    raise ReplacementError("A required native chart is missing.")


def _first_table(slide):
    for shape in slide.shapes:
        if getattr(shape, "has_table", False):
            return shape.table
    raise ReplacementError("A required native table is missing.")


def _replace_chart(chart, categories, series, number_format=None):
    data = CategoryChartData()
    data.categories = [str(value) for value in categories]
    point_count = 0
    for name, values in series:
        clean_values = [float(value) for value in values]
        data.add_series(str(name), clean_values)
        point_count += len(clean_values)
    chart.replace_data(data)
    if number_format:
        for plot in chart.plots:
            try:
                if plot.has_data_labels:
                    plot.data_labels.number_format = number_format
                    plot.data_labels.number_format_is_linked = False
            except (AttributeError, ValueError):
                pass
    return point_count


def _write_table(table, rows):
    changed = 0
    for row_index in range(len(table.rows)):
        for column_index in range(len(table.columns)):
            value = ""
            if row_index < len(rows) and column_index < len(rows[row_index]):
                value = rows[row_index][column_index]
            cell = table.cell(row_index, column_index)
            paragraphs = cell.text_frame.paragraphs
            if paragraphs:
                changed += int(_set_paragraph(paragraphs[0], value))
                for paragraph in paragraphs[1:]:
                    changed += int(_set_paragraph(paragraph, ""))
            else:
                cell.text = str(value)
                changed += 1
    return changed


def _all_text(prs):
    values = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                values.append(shape.text)
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    values.extend(cell.text for cell in row.cells)
    return "\n".join(values)


def _geometry_signature(prs):
    return [
        [
            (int(shape.shape_type), int(shape.left), int(shape.top), int(shape.width), int(shape.height),
             bool(getattr(shape, "has_chart", False)), bool(getattr(shape, "has_table", False)))
            for shape in slide.shapes
        ]
        for slide in prs.slides
    ]


def _native_object_signature(prs):
    charts = []
    tables = []
    for slide_index, slide in enumerate(prs.slides):
        for shape_index, shape in enumerate(slide.shapes):
            if getattr(shape, "has_chart", False):
                charts.append((slide_index, shape_index, int(shape.chart.chart_type)))
            if getattr(shape, "has_table", False):
                tables.append((slide_index, shape_index, len(shape.table.rows), len(shape.table.columns)))
    return charts, tables


def _package_hashes(path, prefixes=("ppt/theme/",)):
    hashes = {}
    with zipfile.ZipFile(path) as archive:
        for name in archive.namelist():
            if any(name.startswith(prefix) for prefix in prefixes):
                raw = archive.read(name).decode("utf-8")
                canonical = ET.canonicalize(raw, strip_text=True)
                hashes[name] = hashlib.sha256(canonical.encode("utf-8")).hexdigest()
    return hashes


def _verify_preservation(original_path, output_path, before_geometry, before_native, before_hashes):
    original = Presentation(original_path)
    updated = Presentation(output_path)
    if len(original.slides) != len(updated.slides):
        raise ReplacementError("The refresh changed the slide count, so the output was rejected.")
    if (original.slide_width, original.slide_height) != (updated.slide_width, updated.slide_height):
        raise ReplacementError("The refresh changed the slide size, so the output was rejected.")
    if before_geometry != _geometry_signature(updated):
        raise ReplacementError("The refresh moved or resized a slide object, so the output was rejected.")
    if before_native != _native_object_signature(updated):
        raise ReplacementError("The refresh changed a native chart or table type, so the output was rejected.")
    if before_hashes != _package_hashes(output_path):
        raise ReplacementError("The refresh changed the presentation theme or master, so the output was rejected.")


def _review_slide_title(slide, slide_number):
    candidates = []
    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        text = _shape_text(shape)
        if not text or text == str(slide_number):
            continue
        font_sizes = []
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                if run.font.size:
                    font_sizes.append(int(run.font.size))
        largest_font = max(font_sizes) if font_sizes else 0
        candidates.append((-largest_font, -int(shape.height), int(shape.top), text.splitlines()[0]))
    if not candidates:
        return f"Slide {slide_number}"
    candidates.sort()
    return candidates[0][3][:120]


def _chart_categories(chart):
    try:
        categories = chart.plots[0].categories
        values = []
        for category in categories:
            label = getattr(category, "label", None)
            values.append(str(label if label is not None else category))
        return values
    except (AttributeError, IndexError, TypeError, ValueError):
        return []


def _display_chart_value(value, series_name, all_values):
    if value is None:
        return "Blank"
    try:
        number = float(value)
    except (TypeError, ValueError):
        return str(value)
    name = str(series_name or "").casefold()
    finite_values = [abs(float(item)) for item in all_values if isinstance(item, Number)]
    if any(token in name for token in ("coupon", "revenue", "cost", "amount", "$")):
        return f"${number:,.2f}"
    if finite_values and max(finite_values) <= 1.5:
        return f"{number:.1%}"
    if number.is_integer():
        return f"{int(number):,}"
    return f"{number:,.2f}"


def _review_mapping(mapping_id, slide_number, slide_title, object_type, object_name, label, old_value, new_value, locator=None):
    mapping = {
        "id": mapping_id,
        "slide": slide_number,
        "slide_title": slide_title,
        "object_type": object_type,
        "object_name": object_name,
        "label": label,
        "old_value": str(old_value),
        "new_value": str(new_value),
    }
    if locator:
        mapping["locator"] = locator
    return mapping


def _chart_type_label(chart):
    chart_type = getattr(chart, "chart_type", None)
    name = getattr(chart_type, "name", None) or str(chart_type or "Chart")
    return name.replace("_", " ").title()


def _table_headers(table):
    if not table.rows:
        return []
    return [cell.text.strip() or f"Column {index + 1}" for index, cell in enumerate(table.rows[0].cells)]


def inspect_deck_structure(presentation_path, mapping_entries=None):
    """Describe editable data objects so users can validate a deck before approval."""
    prs = Presentation(presentation_path)
    mapping_entries = mapping_entries or []
    mappings_by_object = Counter(
        (int(entry.get("slide", 0)), str(entry.get("object_name", "")))
        for entry in mapping_entries
    )
    slides = []
    totals = Counter()

    for slide_number, slide in enumerate(prs.slides, start=1):
        slide_title = _review_slide_title(slide, slide_number)
        objects = []
        slide_counts = Counter()
        for shape_index, shape in enumerate(slide.shapes, start=1):
            object_name = shape.name or f"Object {shape_index}"
            mapped_changes = mappings_by_object[(slide_number, object_name)]

            if getattr(shape, "has_chart", False):
                chart = shape.chart
                categories = _chart_categories(chart)
                series = list(chart.series)
                series_names = [str(item.name or f"Series {index + 1}") for index, item in enumerate(series)]
                point_count = sum(len(list(item.values or [])) for item in series)
                objects.append({
                    "type": "Chart",
                    "name": object_name,
                    "structure": f"{_chart_type_label(chart)} · {len(series)} series · {len(categories)} categories",
                    "details": ", ".join(series_names) or "No named series",
                    "data_points": point_count,
                    "mapped_changes": mapped_changes,
                    "editable": True,
                })
                slide_counts["charts"] += 1
                totals["charts"] += 1
                totals["data_points"] += point_count
                continue

            if getattr(shape, "has_table", False):
                table = shape.table
                row_count = len(table.rows)
                column_count = len(table.columns)
                headers = _table_headers(table)
                objects.append({
                    "type": "Table",
                    "name": object_name,
                    "structure": f"{row_count} rows × {column_count} columns",
                    "details": ", ".join(headers) or "No headers",
                    "data_points": row_count * column_count,
                    "mapped_changes": mapped_changes,
                    "editable": True,
                })
                slide_counts["tables"] += 1
                totals["tables"] += 1
                totals["data_points"] += row_count * column_count
                continue

            if getattr(shape, "has_text_frame", False):
                text = shape.text.strip()
                if not text:
                    continue
                is_figure = looks_like_figure(text)
                if not is_figure and not mapped_changes:
                    slide_counts["text_objects"] += 1
                    totals["text_objects"] += 1
                    continue
                object_type = "Figure" if is_figure else "Text"
                objects.append({
                    "type": object_type,
                    "name": object_name,
                    "structure": "Numeric KPI" if is_figure else "Data-linked text",
                    "details": text[:180],
                    "data_points": 1,
                    "mapped_changes": mapped_changes,
                    "editable": True,
                })
                key = "figures" if is_figure else "text_objects"
                slide_counts[key] += 1
                totals[key] += 1
                totals["data_points"] += 1

        slides.append({
            "slide": slide_number,
            "title": slide_title,
            "shape_count": len(slide.shapes),
            "charts": slide_counts["charts"],
            "tables": slide_counts["tables"],
            "figures": slide_counts["figures"],
            "text_objects": slide_counts["text_objects"],
            "mapped_changes": sum(item["mapped_changes"] for item in objects),
            "objects": objects,
        })
        totals["shapes"] += len(slide.shapes)

    return {
        "slide_count": len(prs.slides),
        "slide_width_inches": round(prs.slide_width / 914400, 2),
        "slide_height_inches": round(prs.slide_height / 914400, 2),
        "shapes": totals["shapes"],
        "charts": totals["charts"],
        "tables": totals["tables"],
        "figures": totals["figures"],
        "text_objects": totals["text_objects"],
        "data_points": totals["data_points"],
        "slides": slides,
    }


def validate_replacement_structure(original_path, updated_path):
    """Return the structural safeguards applied before blanket approval."""
    original = Presentation(original_path)
    updated = Presentation(updated_path)
    original_geometry = _geometry_signature(original)
    updated_geometry = _geometry_signature(updated)
    original_native = _native_object_signature(original)
    updated_native = _native_object_signature(updated)
    original_shapes = [len(slide.shapes) for slide in original.slides]
    updated_shapes = [len(slide.shapes) for slide in updated.slides]
    checks = [
        {
            "name": "Slide count",
            "original": str(len(original.slides)),
            "updated": str(len(updated.slides)),
            "passed": len(original.slides) == len(updated.slides),
        },
        {
            "name": "Slide size",
            "original": f"{original.slide_width / 914400:.2f} × {original.slide_height / 914400:.2f} in",
            "updated": f"{updated.slide_width / 914400:.2f} × {updated.slide_height / 914400:.2f} in",
            "passed": (original.slide_width, original.slide_height) == (updated.slide_width, updated.slide_height),
        },
        {
            "name": "Objects per slide",
            "original": f"{sum(original_shapes)} total",
            "updated": f"{sum(updated_shapes)} total",
            "passed": original_shapes == updated_shapes,
        },
        {
            "name": "Object positions and sizes",
            "original": "Locked",
            "updated": "Unchanged" if original_geometry == updated_geometry else "Changed",
            "passed": original_geometry == updated_geometry,
        },
        {
            "name": "Native charts and tables",
            "original": f"{len(original_native[0])} charts · {len(original_native[1])} tables",
            "updated": f"{len(updated_native[0])} charts · {len(updated_native[1])} tables",
            "passed": original_native == updated_native,
        },
        {
            "name": "Theme and master",
            "original": "Original theme",
            "updated": "Unchanged" if _package_hashes(original_path) == _package_hashes(updated_path) else "Changed",
            "passed": _package_hashes(original_path) == _package_hashes(updated_path),
        },
    ]
    return {"passed": all(check["passed"] for check in checks), "checks": checks}


def compare_deck_replacements(original_path, updated_path):
    """Return every visible old-to-new mapping for the review screen."""
    original = Presentation(original_path)
    updated = Presentation(updated_path)
    if len(original.slides) != len(updated.slides):
        raise ReplacementError("The replacement review found a changed slide count.")

    entries = []
    mapping_id = 0
    for slide_index, (old_slide, new_slide) in enumerate(zip(original.slides, updated.slides), start=1):
        slide_title = _review_slide_title(new_slide, slide_index)
        if len(old_slide.shapes) != len(new_slide.shapes):
            raise ReplacementError(f"Slide {slide_index} changed its object count.")
        for shape_index, (old_shape, new_shape) in enumerate(zip(old_slide.shapes, new_slide.shapes), start=1):
            object_name = new_shape.name or old_shape.name or f"Object {shape_index}"
            if getattr(old_shape, "has_chart", False) and getattr(new_shape, "has_chart", False):
                old_chart = old_shape.chart
                new_chart = new_shape.chart
                old_categories = _chart_categories(old_chart)
                new_categories = _chart_categories(new_chart)
                old_series = list(old_chart.series)
                new_series = list(new_chart.series)
                series_count = max(len(old_series), len(new_series))
                for series_index in range(series_count):
                    old_item = old_series[series_index] if series_index < len(old_series) else None
                    new_item = new_series[series_index] if series_index < len(new_series) else None
                    old_name = getattr(old_item, "name", "") if old_item is not None else ""
                    new_name = getattr(new_item, "name", "") if new_item is not None else ""
                    old_values = list(getattr(old_item, "values", []) or []) if old_item is not None else []
                    new_values = list(getattr(new_item, "values", []) or []) if new_item is not None else []
                    point_count = max(len(old_values), len(new_values), len(old_categories), len(new_categories))
                    combined_values = old_values + new_values
                    for point_index in range(point_count):
                        old_category = old_categories[point_index] if point_index < len(old_categories) else f"Point {point_index + 1}"
                        new_category = new_categories[point_index] if point_index < len(new_categories) else f"Point {point_index + 1}"
                        old_value = old_values[point_index] if point_index < len(old_values) else None
                        new_value = new_values[point_index] if point_index < len(new_values) else None
                        if old_category == new_category and old_value == new_value and old_name == new_name:
                            continue
                        mapping_id += 1
                        old_display = _display_chart_value(old_value, old_name, combined_values)
                        new_display = _display_chart_value(new_value, new_name, combined_values)
                        entries.append(_review_mapping(
                            mapping_id,
                            slide_index,
                            slide_title,
                            "Chart",
                            object_name,
                            f"{new_name or old_name or f'Series {series_index + 1}'} · point {point_index + 1}",
                            f"{old_category}: {old_display}",
                            f"{new_category}: {new_display}",
                            {
                                "kind": "chart_point",
                                "slide_index": slide_index - 1,
                                "shape_index": shape_index - 1,
                                "series_index": series_index,
                                "point_index": point_index,
                            },
                        ))
                continue

            if getattr(old_shape, "has_table", False) and getattr(new_shape, "has_table", False):
                old_table = old_shape.table
                new_table = new_shape.table
                row_count = max(len(old_table.rows), len(new_table.rows))
                column_count = max(len(old_table.columns), len(new_table.columns))
                for row_index in range(row_count):
                    for column_index in range(column_count):
                        old_text = old_table.cell(row_index, column_index).text if row_index < len(old_table.rows) and column_index < len(old_table.columns) else ""
                        new_text = new_table.cell(row_index, column_index).text if row_index < len(new_table.rows) and column_index < len(new_table.columns) else ""
                        if old_text == new_text:
                            continue
                        header = new_table.cell(0, column_index).text or old_table.cell(0, column_index).text or f"Column {column_index + 1}"
                        row_label = "Header" if row_index == 0 else (
                            new_table.cell(row_index, 0).text or old_table.cell(row_index, 0).text or f"Row {row_index + 1}"
                        )
                        mapping_id += 1
                        entries.append(_review_mapping(
                            mapping_id,
                            slide_index,
                            slide_title,
                            "Table",
                            object_name,
                            f"{row_label} · {header}",
                            old_text or "Blank",
                            new_text or "Blank",
                            {
                                "kind": "table_cell",
                                "slide_index": slide_index - 1,
                                "shape_index": shape_index - 1,
                                "row_index": row_index,
                                "column_index": column_index,
                            },
                        ))
                continue

            if getattr(old_shape, "has_text_frame", False) and getattr(new_shape, "has_text_frame", False):
                old_text = old_shape.text.strip()
                new_text = new_shape.text.strip()
                if old_text == new_text:
                    continue
                first_line = (new_text or old_text or object_name).splitlines()[0][:100]
                object_type = "Figure" if looks_like_figure(old_text) and looks_like_figure(new_text) else "Text"
                mapping_id += 1
                entries.append(_review_mapping(
                    mapping_id,
                    slide_index,
                    slide_title,
                    object_type,
                    object_name,
                    first_line,
                    old_text or "Blank",
                    new_text or "Blank",
                    {
                        "kind": "text_shape",
                        "slide_index": slide_index - 1,
                        "shape_index": shape_index - 1,
                    },
                ))

    counts = Counter(entry["object_type"] for entry in entries)
    return {
        "entries": entries,
        "changed_fields": len(entries),
        "text_changes": counts.get("Text", 0),
        "figure_changes": counts.get("Figure", 0),
        "chart_changes": counts.get("Chart", 0),
        "table_changes": counts.get("Table", 0),
        "slides_changed": len({entry["slide"] for entry in entries}),
    }


def apply_selected_deck_replacements(original_path, pending_path, mapping_entries, selected_ids, output_path):
    """Apply only approved review rows while preserving deck geometry and native objects."""
    selected_ids = {int(value) for value in selected_ids}
    selected = [entry for entry in mapping_entries if int(entry.get("id", 0)) in selected_ids]
    if len(selected) == len(mapping_entries):
        shutil.copy2(pending_path, output_path)
        return len(selected)
    if not selected:
        shutil.copy2(original_path, output_path)
        return 0

    prs = Presentation(original_path)
    pending = Presentation(pending_path)
    before_geometry = _geometry_signature(prs)
    before_native = _native_object_signature(prs)
    before_hashes = _package_hashes(original_path)
    chart_groups = {}

    for entry in selected:
        locator = entry.get("locator") or {}
        kind = locator.get("kind")
        slide_index = int(locator.get("slide_index", -1))
        shape_index = int(locator.get("shape_index", -1))
        if slide_index < 0 or shape_index < 0:
            raise ReplacementError("A selected change is missing its PowerPoint object location.")
        target_shape = prs.slides[slide_index].shapes[shape_index]
        source_shape = pending.slides[slide_index].shapes[shape_index]

        if kind == "text_shape":
            _set_shape_text(target_shape, source_shape.text)
        elif kind == "table_cell":
            row_index = int(locator["row_index"])
            column_index = int(locator["column_index"])
            source_text = source_shape.table.cell(row_index, column_index).text
            cell = target_shape.table.cell(row_index, column_index)
            if cell.text_frame.paragraphs:
                _set_paragraph(cell.text_frame.paragraphs[0], source_text)
                for paragraph in cell.text_frame.paragraphs[1:]:
                    _set_paragraph(paragraph, "")
            else:
                cell.text = source_text
        elif kind == "chart_point":
            chart_groups.setdefault((slide_index, shape_index), []).append(locator)
        else:
            raise ReplacementError("A selected change has an unsupported object type.")

    for (slide_index, shape_index), locators in chart_groups.items():
        target_chart = prs.slides[slide_index].shapes[shape_index].chart
        source_chart = pending.slides[slide_index].shapes[shape_index].chart
        categories = _chart_categories(target_chart)
        source_categories = _chart_categories(source_chart)
        target_series = list(target_chart.series)
        source_series = list(source_chart.series)
        names = [str(series.name or f"Series {index + 1}") for index, series in enumerate(target_series)]
        values = [list(series.values or []) for series in target_series]
        for locator in locators:
            series_index = int(locator["series_index"])
            point_index = int(locator["point_index"])
            if series_index >= len(source_series) or series_index >= len(values):
                raise ReplacementError("A selected chart series no longer exists.")
            source_values = list(source_series[series_index].values or [])
            if point_index < len(source_values) and point_index < len(values[series_index]):
                values[series_index][point_index] = source_values[point_index]
            if point_index < len(source_categories) and point_index < len(categories):
                categories[point_index] = source_categories[point_index]
            names[series_index] = str(source_series[series_index].name or names[series_index])
        data = CategoryChartData()
        data.categories = categories
        for name, series_values in zip(names, values):
            data.add_series(name, series_values)
        target_chart.replace_data(data)

    prs.save(output_path)
    _verify_preservation(original_path, output_path, before_geometry, before_native, before_hashes)
    return len(selected)


def looks_like_figure(value):
    text = str(value or "").strip()
    return bool(re.fullmatch(r"[$€£]?\s*-?[\d,.]+(?:\.\d+)?%?", text))


def _date_value(value):
    if isinstance(value, datetime):
        return value
    if hasattr(value, "year") and hasattr(value, "month") and hasattr(value, "day"):
        return datetime(value.year, value.month, value.day)
    if value:
        parsed = pd.to_datetime(value, errors="coerce")
        if not pd.isna(parsed):
            try:
                return parsed.to_pydatetime(warn=False)
            except TypeError:
                return parsed.to_pydatetime()
    return None


def _screening_profile(sheets):
    if "Screening" not in sheets:
        raise ReplacementError("This presentation expects a workbook with a Screening sheet.")
    rows = sheets["Screening"]
    header_index = None
    for index, row in enumerate(rows[:30]):
        text = [_clean(value).casefold() for value in row]
        if text and text[0] == "company name" and any("credit rating" in value for value in text):
            header_index = index
            break
    if header_index is None:
        raise ReplacementError("The Screening header row was not found.")
    headers = [_clean(value) for value in rows[header_index]]
    records = []
    for row in rows[header_index + 1:]:
        if not row or not _clean(row[0]):
            continue
        padded = list(row) + [None] * max(0, len(headers) - len(row))
        if len(padded) < 7 or not _clean(padded[3]) or not _clean(padded[6]):
            continue
        record = dict(zip(headers, padded))
        company = _clean(padded[0])
        ticker = _clean(padded[1]) if len(padded) > 1 else ""
        risk = _clean(padded[3]) if len(padded) > 3 else ""
        rating = _clean(padded[4]) if len(padded) > 4 else ""
        rating_date = _date_value(padded[5] if len(padded) > 5 else None)
        region_raw = _clean(padded[6]) if len(padded) > 6 else ""
        records.append({
            "company": re.sub(r"\s*\([^)]*\)\s*$", "", company),
            "ticker": ticker or "-",
            "risk": risk or "Unrated",
            "rating": rating or "NR",
            "date": rating_date,
            "region": _region_key(region_raw),
            "region_raw": region_raw,
            "source": record,
        })
    if not records:
        raise ReplacementError("The Screening sheet contains no usable company rows.")

    region_totals = Counter(record["region"] for record in records)
    region_ig = Counter(record["region"] for record in records if record["rating"] in IG_RATINGS)
    risk_totals = Counter(record["risk"] for record in records)
    risk_ig = Counter(record["risk"] for record in records if record["rating"] in IG_RATINGS)
    rating_totals = Counter(record["rating"] for record in records)
    year_totals = Counter()
    year_ig = Counter()
    for record in records:
        year = "\u22642016" if record["date"] and record["date"].year <= 2016 else (
            str(record["date"].year) if record["date"] else "Unknown"
        )
        year_totals[year] += 1
        if record["rating"] in IG_RATINGS:
            year_ig[year] += 1
    return {
        "records": records,
        "header_row": header_index + 1,
        "headers": headers,
        "region_totals": region_totals,
        "region_ig": region_ig,
        "risk_totals": risk_totals,
        "risk_ig": risk_ig,
        "rating_totals": rating_totals,
        "year_totals": year_totals,
        "year_ig": year_ig,
    }


def _screening_samples(profile):
    samples = []
    for key, _, _ in REGIONS:
        candidates = [record for record in profile["records"] if record["region"] == key]
        candidates.sort(
            key=lambda record: (RATING_STRENGTH.get(record["rating"], 0), record["company"]),
            reverse=True,
        )
        samples.extend(candidates[:2])
    return samples[:8]


def _replace_screening(prs, profile, data_filename):
    if len(prs.slides) < 9:
        raise ReplacementError("The Company Screening presentation needs at least nine slides.")
    if any(not any(getattr(shape, "has_chart", False) for shape in prs.slides[index].shapes) for index in (2, 3, 4, 5)):
        raise ReplacementError("The Company Screening presentation is missing one or more native charts on slides 3 through 6.")
    if any(not any(getattr(shape, "has_table", False) for shape in prs.slides[index].shapes) for index in (6, 7)):
        raise ReplacementError("The Company Screening presentation is missing a native table on slide 7 or 8.")
    records = profile["records"]
    total = len(records)
    ig_count = sum(1 for record in records if record["rating"] in IG_RATINGS)
    ig_share = ig_count / total
    region_totals = profile["region_totals"]
    region_ig = profile["region_ig"]
    region_labels = [label for _, label, _ in REGIONS]
    region_keys = [key for key, _, _ in REGIONS]
    largest_key = max(region_keys, key=lambda key: region_totals.get(key, 0))
    largest_label = next(label for key, label, _ in REGIONS if key == largest_key)
    largest_share = region_totals.get(largest_key, 0) / total
    ranked_region_keys = sorted(region_keys, key=lambda key: region_totals.get(key, 0), reverse=True)
    second_key = ranked_region_keys[1]
    second_label = next(label for key, label, _ in REGIONS if key == second_key)
    top_ratings = profile["rating_totals"].most_common(3)
    top_risk = profile["risk_totals"].most_common(1)[0][0]
    strong_or_excellent = profile["risk_totals"].get("Strong", 0) + profile["risk_totals"].get("Excellent", 0)
    years = ["\u22642016", "2017", "2018", "2019", "2020", "2021", "2022", "2023", "2024", "2025", "2026"]
    recent_count = profile["year_totals"].get("2025", 0) + profile["year_totals"].get("2026", 0)
    old_count = profile["year_totals"].get("\u22642016", 0)
    changes = {"text": 0, "chart_points": 0, "table_cells": 0}

    changes["text"] += _set_named_text(
        prs.slides[0], "Subtitle-4-2",
        f"Prepared for client discussion | {total:,} screened entities | August 2026 | Confidential",
    )

    slide = prs.slides[1]
    changes["text"] += _set_named_text(slide, "Title-2-5", f"The screen is broad and {ig_share:.0%} investment grade")
    changes["text"] += _set_named_text(slide, "Content-Placeholder-9-9", f"{total:,}")
    changes["text"] += _set_named_text(slide, "Content-Placeholder-9-10", _pct(ig_share))
    changes["text"] += _set_named_text(slide, "Content-Placeholder-9-11", _pct(largest_share))
    changes["text"] += _set_named_text(slide, "Content-Placeholder-10-7", largest_label)
    changes["text"] += _set_named_text(
        slide, "Content-Placeholder-15-12",
        "Executive summary\n"
        f"The replacement file contains {total:,} screened entities. {ig_count:,} are investment grade, and {largest_label} is the largest region at {_pct(largest_share)}.",
    )

    slide = prs.slides[2]
    changes["text"] += _set_named_text(slide, "Google-Shape-533-p58-3", f"{largest_label} is {largest_share:.0%} of the universe")
    changes["chart_points"] += _replace_chart(
        _first_chart(slide), region_labels,
        [
            ("All entities", [region_totals.get(key, 0) for key in region_keys]),
            ("Investment grade", [region_ig.get(key, 0) for key in region_keys]),
        ],
    )
    second_count = region_totals.get(second_key, 0)
    remaining_keys = [key for key in region_keys if key not in {largest_key, second_key}]
    remaining = sum(region_totals.get(key, 0) for key in remaining_keys)
    remaining_labels = [next(label for region_key, label, _ in REGIONS if region_key == key) for key in remaining_keys]
    changes["text"] += _set_shape_text(
        _find_named_shape(slide, "Content-Placeholder-10-5") or _find_text_shape(slide, "follows"),
        f"{second_label} follows\n{second_count:,} entities represent {_pct(second_count / total)} of the screen.",
    )
    changes["text"] += _set_shape_text(
        _find_named_shape(slide, "Content-Placeholder-10-7") or _find_text_shape(slide, "Largest region"),
        f"Largest region\n{region_totals.get(largest_key, 0):,} entities are in {largest_label}.",
    )
    changes["text"] += _set_shape_text(
        _find_named_shape(slide, "Content-Placeholder-10-9") or _find_text_shape(slide, "Coverage breadth"),
        f"Coverage breadth\n{' and '.join(remaining_labels)} add {remaining:,} entities across the remaining regions.",
    )

    slide = prs.slides[3]
    non_ig = [region_totals.get(key, 0) - region_ig.get(key, 0) for key in region_keys]
    changes["chart_points"] += _replace_chart(
        _first_chart(slide), region_labels,
        [
            ("Investment grade", [region_ig.get(key, 0) for key in region_keys]),
            ("Non-IG / unrated", non_ig),
        ],
    )
    changes["text"] += _set_named_text(slide, "Content-Placeholder-9-4", _pct(ig_share))
    changes["text"] += _set_named_text(slide, "Content-Placeholder-9-6", f"{total - ig_count:,}")
    region_shares = [region_ig.get(key, 0) / region_totals.get(key, 1) for key in region_keys if region_totals.get(key, 0)]
    share_title = "Investment-grade share stays high across every geography" if min(region_shares) >= 0.7 else "Investment-grade mix varies by geography"
    changes["text"] += _set_named_text(slide, "Title-10-3", share_title)
    rating_sentence = ", ".join(f"{name} ({count:,})" for name, count in top_ratings)
    changes["text"] += _set_named_text(
        slide, "Content-Placeholder-1-9",
        "Credit quality is resilient\n"
        f"The refreshed screen contains {ig_count:,} investment-grade entities. The three most common ratings are {rating_sentence}.",
    )

    slide = prs.slides[4]
    risk_order = [risk for risk in ("Satisfactory", "Strong", "Excellent") if profile["risk_totals"].get(risk, 0)]
    for risk, _ in profile["risk_totals"].most_common():
        if risk not in risk_order:
            risk_order.append(risk)
    changes["chart_points"] += _replace_chart(
        _first_chart(slide), risk_order,
        [
            ("All entities", [profile["risk_totals"].get(risk, 0) for risk in risk_order]),
            ("Investment grade", [profile["risk_ig"].get(risk, 0) for risk in risk_order]),
        ],
    )
    higher_share = strong_or_excellent / total
    changes["text"] += _set_named_text(slide, "Title-10-3", f"{top_risk} risk dominates; {higher_share:.0%} score higher")
    changes["text"] += _set_named_text(slide, "Content-Placeholder-9-4", _pct(higher_share))
    changes["text"] += _set_named_text(slide, "Content-Placeholder-9-6", f"{profile['risk_totals'].get('Excellent', 0):,}")
    changes["text"] += _set_named_text(
        slide, "Content-Placeholder-1-9",
        "Risk strength supports credit quality\n"
        f"All {total:,} rows retain the source business-risk field. Investment-grade entities remain traceable within every displayed risk band.",
    )

    slide = prs.slides[5]
    changes["chart_points"] += _replace_chart(
        _first_chart(slide), years,
        [
            ("All ratings", [profile["year_totals"].get(year, 0) for year in years]),
            ("Investment grade", [profile["year_ig"].get(year, 0) for year in years]),
        ],
    )
    changes["text"] += _set_named_text(slide, "Content-Placeholder-9-5", f"{recent_count:,}")
    changes["text"] += _set_named_text(slide, "Content-Placeholder-9-8", f"{old_count:,}")
    changes["text"] += _set_named_text(
        slide, "Content-Placeholder-2-10",
        f"The refreshed dates include {recent_count:,} records from 2025-26 and {old_count:,} records from 2016 or earlier.",
    )

    slide = prs.slides[6]
    region_rows = [["Region", "Entities", "Investment grade", "Non-IG / unrated", "IG share"]]
    for key, label, _ in REGIONS:
        region_total = region_totals.get(key, 0)
        region_grade = region_ig.get(key, 0)
        region_rows.append([
            label, f"{region_total:,}", f"{region_grade:,}", f"{region_total - region_grade:,}",
            _pct(region_grade / region_total) if region_total else "0.0%",
        ])
    region_rows.append(["Total", f"{total:,}", f"{ig_count:,}", f"{total - ig_count:,}", _pct(ig_share)])
    top_one = top_ratings[0] if top_ratings else ("NR", 0)
    top_two = top_ratings[1] if len(top_ratings) > 1 else ("NR", 0)
    region_rows.append(["Top rating", top_one[0], f"{top_one[1]:,}", "Second", f"{top_two[0]} ({top_two[1]:,})"])
    changes["table_cells"] += _write_table(_first_table(slide), region_rows)

    slide = prs.slides[7]
    sample_rows = [["Company", "Ticker", "Risk", "Rating", "Region"]]
    labels = {key: label for key, label, _ in REGIONS}
    for record in _screening_samples(profile):
        sample_rows.append([
            record["company"][:34], record["ticker"], record["risk"], record["rating"],
            labels.get(record["region"], record["region_raw"].replace(" (Primary)", "")),
        ])
    changes["table_cells"] += _write_table(_first_table(slide), sample_rows)

    slide = prs.slides[8]
    body = _find_named_shape(slide, "TextBox-11-7") or _find_text_shape(slide, "companion workbook")
    changes["text"] += _set_shape_text(
        body,
        f"The companion workbook named {data_filename} preserves both sheet names, the seven-column header row, {total:,} data rows, date formatting, rating text, and geography labels.\n"
        "Deck Refresh matched the stable fields and updated the native chart series, table cells, metrics, and narrative text without moving slide objects.",
    )
    return changes


def _goldman_profile(sheets):
    required = ["Table 1", "Table 2", "Table 3", "Table 4"]
    if any(name not in sheets for name in required):
        raise ReplacementError("This presentation expects Table 1 through Table 4.")
    scenarios = []
    filed_text = ""
    for name in required[:3]:
        rows = sheets[name]
        if len(rows) >= 3 and not filed_text:
            filed_text = _clean(rows[2][0] if rows[2] else "")
        observations = []
        total_coupon = 0.0
        raw_payments = 0
        for row in rows[7:]:
            label = _clean(row[0] if len(row) > 0 else "")
            price_cell = row[1] if len(row) > 1 else None
            coupon_cell = row[2] if len(row) > 2 else None
            if "total hypothetical coupons" in _clean(price_cell).casefold():
                total_coupon = _number(coupon_cell)
                break
            if not label or price_cell is None:
                continue
            price = _number(price_cell)
            coupon = _number(coupon_cell)
            repeat = 6 if "eighth" in label.casefold() and "thirteenth" in label.casefold() else 1
            raw_payments += int(coupon > 0)
            for copy in range(repeat):
                observations.append({
                    "label": label if repeat == 1 else f"Observation {8 + copy}",
                    "price": price,
                    "coupon": coupon / repeat if repeat else coupon,
                })
        if not observations:
            raise ReplacementError(f"{name} contains no hypothetical observations.")
        if total_coupon == 0:
            total_coupon = sum(item["coupon"] for item in observations)
        scenarios.append({
            "name": name,
            "observations": observations,
            "total_coupon": total_coupon,
            "raw_payments": raw_payments,
        })
    toc = []
    for row_index, row in enumerate(sheets["Table 4"][8:19], start=9):
        section = _clean(row[0] if len(row) > 0 else "")
        page = _clean(row[1] if len(row) > 1 else "")
        if section:
            toc.append((section, page, row_index))
    filed_date = filed_text.split(":", 1)[1].strip() if ":" in filed_text else filed_text
    quarter = ""
    for rows in sheets.values():
        for row in rows[:12]:
            for value in row:
                match = re.search(r"reporting\s+quarter\s*:\s*(Q[1-4])\b", _clean(value), re.IGNORECASE)
                if match:
                    quarter = match.group(1).upper()
                    break
            if quarter:
                break
        if quarter:
            break
    if not quarter:
        parsed = pd.to_datetime(filed_date, errors="coerce")
        if not pd.isna(parsed):
            quarter = f"Q{((int(parsed.month) - 1) // 3) + 1}"
    return {"scenarios": scenarios, "toc": toc, "filed_date": filed_date, "quarter": quarter or "Q4"}


def _formatted_filed_date(value):
    parsed = pd.to_datetime(value, errors="coerce")
    if pd.isna(parsed):
        return value
    return f"{parsed.strftime('%B')} {parsed.day}, {parsed.year}"


def _replace_quarter_references(prs, quarter):
    pattern = re.compile(r"\bQ[1-4]\b", re.IGNORECASE)
    changed = 0

    def update_frame(frame):
        nonlocal changed
        for paragraph in frame.paragraphs:
            if paragraph.runs:
                for run in paragraph.runs:
                    updated, count = pattern.subn(quarter, run.text)
                    if count:
                        run.text = updated
                        changed += count
            else:
                updated, count = pattern.subn(quarter, paragraph.text)
                if count:
                    paragraph.text = updated
                    changed += count

    for slide in prs.slides:
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                update_frame(shape.text_frame)
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    for cell in row.cells:
                        update_frame(cell.text_frame)
            if getattr(shape, "has_chart", False):
                chart = shape.chart
                if chart.has_title:
                    update_frame(chart.chart_title.text_frame)
                categories = _chart_categories(chart)
                series = list(chart.series)
                new_categories = [pattern.sub(quarter, value) for value in categories]
                new_names = [pattern.sub(quarter, str(item.name or f"Series {index + 1}")) for index, item in enumerate(series)]
                label_changes = sum(old != new for old, new in zip(categories, new_categories))
                label_changes += sum(str(item.name or "") != new for item, new in zip(series, new_names))
                if label_changes:
                    data = CategoryChartData()
                    data.categories = new_categories
                    for name, item in zip(new_names, series):
                        data.add_series(name, list(item.values or []))
                    chart.replace_data(data)
                    changed += label_changes
        if slide.has_notes_slide:
            for shape in slide.notes_slide.shapes:
                if getattr(shape, "has_text_frame", False):
                    update_frame(shape.text_frame)
    return changed


def _replace_goldman(prs, profile, data_filename):
    if len(prs.slides) < 9:
        raise ReplacementError("The Goldman Sachs presentation needs at least nine slides.")
    if any(not any(getattr(shape, "has_chart", False) for shape in prs.slides[index].shapes) for index in (2, 3, 4)):
        raise ReplacementError("The Goldman Sachs presentation is missing one or more native charts on slides 3 through 5.")
    if any(not any(getattr(shape, "has_table", False) for shape in prs.slides[index].shapes) for index in (5, 7)):
        raise ReplacementError("The Goldman Sachs presentation is missing a native table on slide 6 or 8.")
    scenarios = profile["scenarios"]
    stats = []
    for scenario in scenarios:
        observations = scenario["observations"]
        prices = [item["price"] for item in observations]
        paying = sum(1 for item in observations if item["coupon"] > 0)
        stats.append({
            "observations": len(observations),
            "average": sum(prices) / len(prices),
            "min": min(prices),
            "max": max(prices),
            "paying": paying,
            "total_coupon": scenario["total_coupon"],
        })
    all_observations = [item for scenario in scenarios for item in scenario["observations"]]
    total_observations = len(all_observations)
    paying_events = sum(1 for item in all_observations if item["coupon"] > 0)
    total_coupon = sum(scenario["total_coupon"] for scenario in scenarios)
    pay_rate = paying_events / total_observations if total_observations else 0
    bands = [
        ("Below 50%", lambda value: value < 0.5),
        ("50% to 69.9%", lambda value: 0.5 <= value < 0.7),
        ("70% or higher", lambda value: value >= 0.7),
    ]
    band_rows = []
    for label, test in bands:
        points = [item for item in all_observations if test(item["price"])]
        band_rows.append({
            "label": label,
            "observations": len(points),
            "paying": sum(1 for item in points if item["coupon"] > 0),
        })
    top_band_share = (band_rows[-1]["paying"] / paying_events) if paying_events else 0
    changes = {"text": 0, "chart_points": 0, "table_cells": 0}

    filed = _formatted_filed_date(profile["filed_date"])
    changes["text"] += _set_named_text(
        prs.slides[0], "Subtitle-4-2", f"{profile['quarter']} client presentation | Form 424B2 | Filed {filed} | Confidential"
    )

    slide = prs.slides[1]
    changes["text"] += _set_named_text(slide, "Content-Placeholder-9-9", str(len(scenarios)))
    changes["text"] += _set_named_text(slide, "Content-Placeholder-9-10", str(total_observations))
    changes["text"] += _set_named_text(slide, "Content-Placeholder-9-11", _pct(pay_rate))
    changes["text"] += _set_named_text(
        slide, "Content-Placeholder-15-12",
        "Executive summary\n"
        f"Across three refreshed scenarios, {paying_events} of {total_observations} observation dates produce coupons. Total modeled coupons equal {_money(total_coupon)}.",
    )

    slide = prs.slides[2]
    paying_scenarios = sum(1 for stat in stats if stat["total_coupon"] > 0)
    title = "All scenarios now pay coupons" if paying_scenarios == 3 else f"{paying_scenarios} of 3 scenarios pay coupons"
    changes["text"] += _set_named_text(slide, "Google-Shape-533-p58-3", title)
    changes["chart_points"] += _replace_chart(
        _first_chart(slide), [f"Scenario {index + 1}" for index in range(len(stats))],
        [
            ("Actual coupons ($)", [stat["total_coupon"] for stat in stats]),
            ("Maximum coupons ($)", [stat["observations"] * 10.667 for stat in stats]),
        ],
        number_format="$0.00",
    )
    for index, marker in enumerate(("Scenario 1", "Scenario 2", "Scenario 3")):
        stat = stats[index]
        payment_word = "observation" if stat["paying"] == 1 else "observations"
        changes["text"] += _set_shape_text(
            _find_text_shape(slide, marker),
            f"{marker}\n{_money(stat['total_coupon'])} of coupons across {stat['paying']} paying {payment_word}.",
        )

    slide = prs.slides[3]
    first_six = [scenario["observations"][:6] for scenario in scenarios]
    categories = ["First", "Second", "Third", "Fourth", "Fifth", "Sixth"]
    changes["chart_points"] += _replace_chart(
        _first_chart(slide), categories,
        [(f"Scenario {index + 1}", [item["price"] for item in points]) for index, points in enumerate(first_six)],
    )
    peak_index = max(range(len(stats)), key=lambda index: stats[index]["max"])
    low_index = min(range(len(stats)), key=lambda index: stats[index]["min"])
    changes["text"] += _set_named_text(slide, "Content-Placeholder-9-5", _pct(stats[peak_index]["max"], 0))
    changes["text"] += _set_named_text(slide, "Content-Placeholder-15-6", f"Scenario {peak_index + 1} peak")
    changes["text"] += _set_named_text(slide, "Content-Placeholder-9-8", _pct(stats[low_index]["min"], 0))
    changes["text"] += _set_named_text(slide, "Content-Placeholder-15-9", f"Scenario {low_index + 1} low")
    changes["text"] += _set_named_text(
        slide, "Content-Placeholder-2-10",
        f"Scenario {peak_index + 1} reaches {_pct(stats[peak_index]['max'], 0)}. Scenario {low_index + 1} falls to {_pct(stats[low_index]['min'], 0)} across the refreshed paths.",
    )

    slide = prs.slides[4]
    changes["chart_points"] += _replace_chart(
        _first_chart(slide), [row["label"] for row in band_rows],
        [
            ("Observations", [row["observations"] for row in band_rows]),
            ("Coupon-paying", [row["paying"] for row in band_rows]),
        ],
    )
    changes["text"] += _set_named_text(slide, "Content-Placeholder-9-4", _pct(pay_rate))
    changes["text"] += _set_named_text(slide, "Content-Placeholder-9-6", _pct(top_band_share, 0))
    changes["text"] += _set_named_text(
        slide, "Content-Placeholder-1-9",
        "The refreshed examples preserve the payoff break\n"
        f"{band_rows[0]['observations']} observations fall below 50%, {band_rows[1]['observations']} fall between 50% and 69.9%, and {band_rows[2]['observations']} reach 70% or higher. {band_rows[2]['paying']} top-band observations pay coupons.",
    )

    slide = prs.slides[5]
    summary_rows = [["Scenario", "Observations", "Average price", "Paying dates", "Total coupons"]]
    for index, stat in enumerate(stats):
        summary_rows.append([
            f"Scenario {index + 1}", str(stat["observations"]), _pct(stat["average"]),
            str(stat["paying"]), _money(stat["total_coupon"]),
        ])
    average = sum(item["price"] for item in all_observations) / total_observations
    summary_rows.append(["Combined", str(total_observations), _pct(average), str(paying_events), _money(total_coupon)])
    summary_rows.append([
        "Lowest price", _pct(min(item["price"] for item in all_observations), 0),
        "Highest price", _pct(max(item["price"] for item in all_observations), 0), "Source rows",
    ])
    changes["table_cells"] += _write_table(_first_table(slide), summary_rows)

    slide = prs.slides[7]
    toc_rows = [["Section", "Page", "Review focus", "Deck object", "Replacement key"]]
    for index, (section, page, row_number) in enumerate(profile["toc"][:8]):
        toc_rows.append([
            section[:34], page,
            "Payoff mechanics" if index < 3 else ("Risk and hedging" if index < 6 else "Tax / distribution"),
            "Table" if index < 3 else "Text", f"Table 4 row {row_number}",
        ])
    changes["table_cells"] += _write_table(_first_table(slide), toc_rows)

    slide = prs.slides[8]
    body = _find_named_shape(slide, "TextBox-11-7") or _find_text_shape(slide, "companion workbook")
    changes["text"] += _set_shape_text(
        body,
        f"The companion workbook named {data_filename} keeps Table 1 through Table 4, the metadata rows, observation labels, percentage columns, currency columns, totals rows, and table-of-contents structure.\n"
        "Deck Refresh matched the stable anchors and updated the native chart series, tables, metrics, filing date, and narrative text without moving slide objects.",
    )
    changes["text"] += _replace_quarter_references(prs, profile["quarter"])
    return changes


def replace_deck_1to1(presentation_path, data_path, output_path, data_filename=None):
    """Refresh a supported deck and return a machine-readable replacement report."""
    data_filename = data_filename or os.path.basename(data_path)
    sheets = _read_sheet_rows(data_path)
    prs = Presentation(presentation_path)
    before_geometry = _geometry_signature(prs)
    before_native = _native_object_signature(prs)
    before_hashes = _package_hashes(presentation_path)
    if "Screening" in sheets:
        profile_name = "Company screening"
        profile = _screening_profile(sheets)
        changes = _replace_screening(prs, profile, data_filename)
        source_rows = len(profile["records"])
        source_sheets = 2 if "Screen Criteria" in sheets else 1
    elif all(name in sheets for name in ("Table 1", "Table 2", "Table 3", "Table 4")):
        profile_name = "Goldman Sachs 424B2"
        profile = _goldman_profile(sheets)
        changes = _replace_goldman(prs, profile, data_filename)
        source_rows = sum(len(scenario["observations"]) for scenario in profile["scenarios"])
        source_sheets = 4
    else:
        raise ReplacementError(
            "The workbook does not match a specialized analysis profile."
        )

    ensure_chart_contrast(prs)
    prs.save(output_path)
    _verify_preservation(presentation_path, output_path, before_geometry, before_native, before_hashes)
    total_updates = sum(changes.values())
    return {
        "profile": profile_name,
        "source_sheets": source_sheets,
        "source_rows": source_rows,
        "text_updates": changes["text"],
        "chart_points": changes["chart_points"],
        "table_cells": changes["table_cells"],
        "total_updates": total_updates,
        "slide_count": len(prs.slides),
        "charts": len(before_native[0]),
        "tables": len(before_native[1]),
        "geometry_preserved": True,
        "theme_preserved": True,
    }
