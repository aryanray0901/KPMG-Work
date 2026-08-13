"""Automatic readable text styling for native PowerPoint charts."""

from __future__ import annotations

from typing import Any

from pptx.dml.color import RGBColor
from pptx.util import Pt


_WHITE = (255, 255, 255)
_INK = (17, 24, 39)


def _rgb_tuple(value: Any) -> tuple[int, int, int] | None:
    if value is None:
        return None
    try:
        text = str(value).strip().lstrip("#")
        if len(text) == 6:
            return tuple(int(text[index:index + 2], 16) for index in (0, 2, 4))
    except Exception:
        return None
    return None


def _fill_rgb(fill: Any) -> tuple[int, int, int] | None:
    try:
        return _rgb_tuple(fill.fore_color.rgb)
    except Exception:
        return None


def _luminance(color: tuple[int, int, int]) -> float:
    channels = []
    for item in color:
        value = item / 255.0
        channels.append(value / 12.92 if value <= 0.04045 else ((value + 0.055) / 1.055) ** 2.4)
    return 0.2126 * channels[0] + 0.7152 * channels[1] + 0.0722 * channels[2]


def _contrast_ratio(first: tuple[int, int, int], second: tuple[int, int, int]) -> float:
    high, low = sorted((_luminance(first), _luminance(second)), reverse=True)
    return (high + 0.05) / (low + 0.05)


def _blend(background: tuple[int, int, int], foreground: tuple[int, int, int], amount: float) -> tuple[int, int, int]:
    return tuple(round(bg + (fg - bg) * amount) for bg, fg in zip(background, foreground))


def contrast_palette(background: tuple[int, int, int]) -> dict[str, tuple[int, int, int]]:
    """Return WCAG-oriented chart colors for a visible background."""
    text = _WHITE if _contrast_ratio(background, _WHITE) >= _contrast_ratio(background, _INK) else _INK
    dark = text == _WHITE
    return {
        "background": background,
        "text": text,
        "secondary": _blend(background, text, 0.84 if dark else 0.72),
        "gridline": _blend(background, text, 0.30 if dark else 0.17),
        "axis": _blend(background, text, 0.48 if dark else 0.32),
    }


def _covers_point(shape: Any, x: float, y: float) -> bool:
    try:
        return shape.left <= x <= shape.left + shape.width and shape.top <= y <= shape.top + shape.height
    except Exception:
        return False


def visible_chart_background(slide: Any, chart_shape: Any) -> tuple[int, int, int]:
    """Resolve the color visually behind the center of a transparent chart."""
    background = _fill_rgb(getattr(getattr(slide, "background", None), "fill", None)) or _WHITE
    try:
        center_x = chart_shape.left + chart_shape.width / 2
        center_y = chart_shape.top + chart_shape.height / 2
        for shape in slide.shapes:
            if shape is chart_shape or getattr(shape, "shape_id", None) == getattr(chart_shape, "shape_id", None):
                break
            if _covers_point(shape, center_x, center_y):
                color = _fill_rgb(getattr(shape, "fill", None))
                if color is not None:
                    background = color
    except Exception:
        pass
    try:
        chart_fill = _fill_rgb(chart_shape.chart.format.fill)
        if chart_fill is not None:
            background = chart_fill
    except Exception:
        pass
    return background


def _set_font(font: Any, color: tuple[int, int, int], size: float | None = None, bold: bool | None = None) -> None:
    try:
        font.color.rgb = RGBColor(*color)
        if size is not None:
            font.size = Pt(size)
        if bold is not None:
            font.bold = bold
    except Exception:
        pass


def _set_text_frame(text_frame: Any, color: tuple[int, int, int], size: float | None = None, bold: bool | None = None) -> None:
    try:
        for paragraph in text_frame.paragraphs:
            _set_font(paragraph.font, color, size, bold)
            for run in paragraph.runs:
                _set_font(run.font, color, size, bold)
    except Exception:
        pass


def _set_line(line: Any, color: tuple[int, int, int], width: float | None = None) -> None:
    try:
        line.color.rgb = RGBColor(*color)
        if width is not None:
            line.width = Pt(width)
    except Exception:
        pass


def apply_chart_text_contrast(slide: Any, shape: Any, compact: bool | None = None) -> dict[str, Any]:
    """Make every native chart label readable without changing its data or series colors."""
    if not getattr(shape, "has_chart", False):
        return {"changed": False}
    try:
        chart = shape.chart
    except Exception:
        # Some duplicated third-party slides contain a stale chart relationship.
        # Contrast styling must never block the requested deck edit.
        return {"changed": False}
    if compact is None:
        try:
            compact = shape.width < 2_200_000 or shape.height < 1_450_000
        except Exception:
            compact = False
    palette = contrast_palette(visible_chart_background(slide, shape))
    title_size = 11 if compact else 22
    axis_size = 6 if compact else 10
    label_size = 6 if compact else 10

    try:
        if chart.has_title:
            _set_text_frame(chart.chart_title.text_frame, palette["text"], title_size, True)
    except Exception:
        pass
    try:
        if chart.has_legend:
            _set_font(chart.legend.font, palette["secondary"], axis_size, False)
    except Exception:
        pass

    for axis_name in ("value_axis", "category_axis"):
        try:
            axis = getattr(chart, axis_name)
            _set_font(axis.tick_labels.font, palette["secondary"], axis_size, False)
            _set_line(axis.format.line, palette["axis"], 0.75)
            if getattr(axis, "has_title", False):
                _set_text_frame(axis.axis_title.text_frame, palette["text"], axis_size, False)
            if getattr(axis, "has_major_gridlines", False):
                _set_line(axis.major_gridlines.format.line, palette["gridline"], 0.75)
        except Exception:
            pass

    try:
        for plot in chart.plots:
            if getattr(plot, "has_data_labels", False):
                _set_font(plot.data_labels.font, palette["text"], label_size, False)
    except Exception:
        pass
    return {
        "changed": True,
        "background": "#" + "".join(f"{item:02X}" for item in palette["background"]),
        "text": "#" + "".join(f"{item:02X}" for item in palette["text"]),
    }


def ensure_chart_contrast(presentation: Any) -> int:
    """Apply background-aware chart readability to a whole presentation."""
    count = 0
    for slide in presentation.slides:
        for shape in slide.shapes:
            if getattr(shape, "has_chart", False):
                apply_chart_text_contrast(slide, shape)
                count += 1
    return count
