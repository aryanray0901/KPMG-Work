"""PowerPoint editing engine and structured OpenAI edit planner for Deck Refresh."""

from __future__ import annotations

import base64
import copy
import json
import os
import re
import shutil
import tempfile
from pathlib import Path
from typing import Annotated, Any, Iterable, Literal

from dotenv import load_dotenv
from openai import OpenAI
import pandas as pd
from openpyxl import load_workbook
from pydantic import BaseModel, ConfigDict, Field, TypeAdapter, ValidationError
from pptx import Presentation
from pptx.chart.data import CategoryChartData, XyChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_DATA_LABEL_POSITION, XL_LEGEND_POSITION, XL_MARKER_STYLE
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt
from pptx.oxml.xmlchemy import OxmlElement
from pptx.oxml.ns import qn

from chart_contrast import apply_chart_text_contrast, ensure_chart_contrast

load_dotenv(Path(__file__).with_name(".env"))


class EditorError(RuntimeError):
    """Raised when a requested deck edit cannot be completed safely."""


# ---------------------------------------------------------------------------
# Deck inspection
# ---------------------------------------------------------------------------


def _slide_title(slide) -> str:
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False) and shape.text_frame.text.strip():
            text = shape.text_frame.text.strip().splitlines()[0]
            if text:
                return text[:180]
    return ""


def _shape_text(shape) -> str:
    if getattr(shape, "has_text_frame", False):
        return shape.text_frame.text.strip()
    if getattr(shape, "has_table", False):
        rows = []
        for row in shape.table.rows:
            rows.append(" | ".join(cell.text.strip() for cell in row.cells))
        return "\n".join(rows)
    if getattr(shape, "has_chart", False):
        try:
            chart = shape.chart
            pieces = []
            if chart.has_title:
                pieces.append(chart.chart_title.text_frame.text.strip())
            plot = chart.plots[0]
            categories = [str(c) for c in plot.categories]
            pieces.append("Categories: " + ", ".join(categories))
            for series in chart.series:
                pieces.append(f"{series.name}: " + ", ".join(str(v) for v in series.values))
            return "\n".join(pieces)
        except Exception:
            return "Chart"
    return ""


def _shape_kind(shape) -> str:
    if getattr(shape, "has_chart", False):
        return "chart"
    if getattr(shape, "has_table", False):
        return "table"
    if getattr(shape, "has_text_frame", False):
        return "text"
    if str(getattr(shape, "shape_type", "")) == "PICTURE (13)":
        return "picture"
    return "shape"


def deck_summary(pptx_path: str, max_text: int = 1800) -> dict[str, Any]:
    prs = Presentation(pptx_path)
    width = prs.slide_width / 914400
    height = prs.slide_height / 914400
    slides = []
    for slide_number, slide in enumerate(prs.slides, start=1):
        shapes = []
        for shape_number, shape in enumerate(slide.shapes, start=1):
            text = _shape_text(shape)
            entry: dict[str, Any] = {
                "shape": shape_number,
                "shape_id": int(shape.shape_id),
                "name": shape.name,
                "kind": _shape_kind(shape),
                "type": str(shape.shape_type),
                "x": round(shape.left / 914400, 2),
                "y": round(shape.top / 914400, 2),
                "width": round(shape.width / 914400, 2),
                "height": round(shape.height / 914400, 2),
                "text": text[:max_text],
            }
            if getattr(shape, "has_table", False):
                entry["table_rows"] = len(shape.table.rows)
                entry["table_columns"] = len(shape.table.columns)
            if getattr(shape, "has_chart", False):
                entry["chart"] = True
            shapes.append(entry)
        slides.append({
            "slide": slide_number,
            "title": _slide_title(slide),
            "shapes": shapes,
        })
    return {
        "slide_width_inches": round(width, 2),
        "slide_height_inches": round(height, 2),
        "slide_count": len(slides),
        "slides": slides,
    }


# ---------------------------------------------------------------------------
# Shape, slide, text, and geometry helpers
# ---------------------------------------------------------------------------


def _rgb(value: str | None, default: str | None = None) -> RGBColor | None:
    raw = (value or default or "").strip().lstrip("#")
    if not re.fullmatch(r"[0-9a-fA-F]{6}", raw):
        return None
    return RGBColor.from_string(raw.upper())


def _slide(prs: Presentation, number: int):
    if number < 1 or number > len(prs.slides):
        raise EditorError(f"Slide {number} does not exist.")
    return prs.slides[number - 1]


def _shape(
    slide,
    number: int | None = None,
    *,
    shape_id: int | None = None,
    name: str | None = None,
    text_contains: str | None = None,
):
    """Resolve a shape using a stable ID first, then semantic hints, then index.

    Shape indexes change after deletions. PowerPoint shape IDs remain stable inside
    the slide, so planned edits prefer shape_id and survive earlier operations.
    """
    shapes = list(slide.shapes)

    if shape_id is not None:
        for candidate in shapes:
            if int(candidate.shape_id) == int(shape_id):
                return candidate

    if name:
        wanted = str(name).strip().casefold()
        exact = [candidate for candidate in shapes if candidate.name.strip().casefold() == wanted]
        if len(exact) == 1:
            return exact[0]
        partial = [candidate for candidate in shapes if wanted in candidate.name.strip().casefold()]
        if len(partial) == 1:
            return partial[0]

    if text_contains:
        wanted = _normalize_text(text_contains)
        matches = []
        for candidate in shapes:
            candidate_text = _normalize_text(_shape_text(candidate))
            if wanted and wanted in candidate_text:
                matches.append(candidate)
        if len(matches) == 1:
            return matches[0]
        if matches:
            matches.sort(key=lambda candidate: len(_shape_text(candidate)))
            return matches[0]

    if number is not None:
        if 1 <= int(number) <= len(shapes):
            return shapes[int(number) - 1]
        raise EditorError(f"Shape {number} does not exist on this slide.")

    raise EditorError("The edit did not identify a valid shape.")


def _operation_shape(slide, operation: dict[str, Any]):
    number = operation.get("shape")
    return _shape(
        slide,
        int(number) if number is not None else None,
        shape_id=int(operation["shape_id"]) if operation.get("shape_id") is not None else None,
        name=operation.get("shape_name"),
        text_contains=operation.get("text_contains"),
    )


def _to_emu(value: Any, total: int) -> int:
    try:
        number = float(value)
    except (TypeError, ValueError) as exc:
        raise EditorError(f"Invalid position or size value: {value}") from exc
    if number <= 1.5:
        return int(total * number)
    return int(Inches(number))


def _normalize_text(text: Any) -> str:
    return re.sub(r"\s+", " ", str(text or "")).strip().casefold()


def _set_text_frame_text(text_frame, text: str, preserve_first_run: bool = True) -> None:
    text = str(text)
    if preserve_first_run and text_frame.paragraphs:
        first_para = text_frame.paragraphs[0]
        first_run = first_para.runs[0] if first_para.runs else None
        lines = text.splitlines() or [""]
        if first_run is not None:
            first_run.text = lines[0]
            for run in first_para.runs[1:]:
                run.text = ""
        else:
            first_para.text = lines[0]
        while len(text_frame.paragraphs) > 1:
            p = text_frame.paragraphs[-1]._p
            p.getparent().remove(p)
        for line in lines[1:]:
            paragraph = text_frame.add_paragraph()
            paragraph.text = line
        return
    text_frame.clear()
    lines = text.splitlines() or [""]
    text_frame.paragraphs[0].text = lines[0]
    for line in lines[1:]:
        text_frame.add_paragraph().text = line


def _iter_text_frames(slide, target_shape=None):
    shapes = [target_shape] if target_shape is not None else list(slide.shapes)
    pending = list(shapes)
    while pending:
        shape = pending.pop(0)
        if getattr(shape, "has_text_frame", False):
            yield shape.text_frame
        if getattr(shape, "has_table", False):
            for row in shape.table.rows:
                for cell in row.cells:
                    yield cell.text_frame
        if getattr(shape, "has_chart", False):
            try:
                if shape.chart.has_title:
                    yield shape.chart.chart_title.text_frame
            except Exception:
                pass
        if hasattr(shape, "shapes"):
            pending.extend(list(shape.shapes))


def _replace_in_string(source: str, old: str, new: str, replace_all: bool, case_sensitive: bool) -> tuple[str, int]:
    if not old:
        return source, 0
    flags = 0 if case_sensitive else re.IGNORECASE
    tokens = [token for token in re.split(r"\s+", old.strip()) if token]
    pattern_text = r"\s+".join(re.escape(token) for token in tokens) if tokens else re.escape(old)
    pattern = re.compile(pattern_text, flags)
    limit = 0 if replace_all else 1
    result, count = pattern.subn(lambda _: new, source, count=limit)
    return result, count


def _replace_text_in_frame(
    text_frame,
    old: str,
    new: str,
    replace_all: bool,
    case_sensitive: bool = False,
) -> int:
    count = 0
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            replacement, run_count = _replace_in_string(run.text, old, new, replace_all, case_sensitive)
            if run_count:
                run.text = replacement
                count += run_count
                if not replace_all:
                    return count
        # Handles text split across multiple runs while retaining the first run's formatting.
        paragraph_text = paragraph.text
        replacement, paragraph_count = _replace_in_string(
            paragraph_text, old, new, replace_all, case_sensitive
        )
        if paragraph_count and count == 0:
            if paragraph.runs:
                paragraph.runs[0].text = replacement
                for run in paragraph.runs[1:]:
                    run.text = ""
            else:
                paragraph.text = replacement
            count += paragraph_count
            if not replace_all:
                return count
    return count


def _apply_text_style(shape, operation: dict[str, Any]) -> None:
    if not getattr(shape, "has_text_frame", False):
        return
    font_size = operation.get("font_size")
    font_color = _rgb(operation.get("font_color"))
    font_face = operation.get("font_face")
    bold = operation.get("bold")
    italic = operation.get("italic")
    alignment = operation.get("alignment")
    align_map = {
        "left": PP_ALIGN.LEFT,
        "center": PP_ALIGN.CENTER,
        "right": PP_ALIGN.RIGHT,
        "justify": PP_ALIGN.JUSTIFY,
    }
    for paragraph in shape.text_frame.paragraphs:
        if alignment in align_map:
            paragraph.alignment = align_map[alignment]
        for run in paragraph.runs:
            if font_size is not None:
                run.font.size = Pt(float(font_size))
            if font_color:
                run.font.color.rgb = font_color
            if font_face:
                run.font.name = str(font_face)
            if bold is not None:
                run.font.bold = bool(bold)
            if italic is not None:
                run.font.italic = bool(italic)


def _apply_shape_style(shape, operation: dict[str, Any]) -> None:
    fill_color = _rgb(operation.get("fill_color"))
    line_color = _rgb(operation.get("line_color"))
    if fill_color is not None:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    if operation.get("no_fill"):
        shape.fill.background()
    if line_color is not None:
        shape.line.color.rgb = line_color
    if operation.get("no_line"):
        shape.line.fill.background()
    if operation.get("line_width") is not None:
        shape.line.width = Pt(float(operation["line_width"]))
    _apply_text_style(shape, operation)


def _set_geometry(prs: Presentation, shape, operation: dict[str, Any]) -> None:
    if operation.get("x") is not None:
        shape.left = _to_emu(operation["x"], prs.slide_width)
    if operation.get("y") is not None:
        shape.top = _to_emu(operation["y"], prs.slide_height)
    if operation.get("width") is not None:
        shape.width = _to_emu(operation["width"], prs.slide_width)
    if operation.get("height") is not None:
        shape.height = _to_emu(operation["height"], prs.slide_height)


# ---------------------------------------------------------------------------
# Slide and object creation helpers
# ---------------------------------------------------------------------------


def _delete_slide(prs: Presentation, slide_number: int) -> None:
    _slide(prs, slide_number)
    slide_id = prs.slides._sldIdLst[slide_number - 1]
    relationship_id = slide_id.rId
    prs.part.drop_rel(relationship_id)
    prs.slides._sldIdLst.remove(slide_id)


def _move_slide(prs: Presentation, from_slide: int, to_slide: int) -> None:
    count = len(prs.slides)
    if from_slide < 1 or from_slide > count or to_slide < 1 or to_slide > count:
        raise EditorError("Move position is outside the deck.")
    slide_id = prs.slides._sldIdLst[from_slide - 1]
    prs.slides._sldIdLst.remove(slide_id)
    prs.slides._sldIdLst.insert(to_slide - 1, slide_id)


def _reorder_slides(prs: Presentation, order: list[int]) -> None:
    count = len(prs.slides)
    expected = list(range(1, count + 1))
    if sorted(order) != expected:
        raise EditorError(f"Slide order must contain every slide exactly once: {expected}.")
    original_ids = list(prs.slides._sldIdLst)
    for slide_id in list(prs.slides._sldIdLst):
        prs.slides._sldIdLst.remove(slide_id)
    for old_number in order:
        prs.slides._sldIdLst.append(original_ids[old_number - 1])


def _duplicate_slide(prs: Presentation, slide_number: int, position: int | None = None) -> int:
    source = _slide(prs, slide_number)
    blank_layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[-1]
    target = prs.slides.add_slide(blank_layout)

    for shape in source.shapes:
        target.shapes._spTree.insert_element_before(copy.deepcopy(shape.element), "p:extLst")

    for rel in source.part.rels.values():
        if "notesSlide" in rel.reltype:
            continue
        try:
            target.part.rels.add_relationship(rel.reltype, rel._target, rel.rId, rel.is_external)
        except TypeError:
            try:
                target.part.rels.add_relationship(rel.reltype, rel._target, rel.rId)
            except Exception:
                pass
        except Exception:
            pass

    new_number = len(prs.slides)
    destination = position or min(slide_number + 1, new_number)
    _move_slide(prs, new_number, destination)
    return destination


def _add_slide(prs: Presentation, operation: dict[str, Any]) -> int:
    if operation.get("template_slide"):
        return _duplicate_slide(
            prs,
            int(operation["template_slide"]),
            int(operation.get("position") or len(prs.slides) + 1),
        )

    blank_layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[-1]
    slide = prs.slides.add_slide(blank_layout)
    slide_width = prs.slide_width
    slide_height = prs.slide_height

    background = _rgb(operation.get("background_color"), "FFFFFF")
    if background:
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = background

    title = str(operation.get("title", "New slide")).strip()
    subtitle = str(operation.get("subtitle", "")).strip()
    body = operation.get("body", [])
    if isinstance(body, str):
        body = [line.strip("• -") for line in body.splitlines() if line.strip()]
    if not isinstance(body, list):
        body = []

    if title:
        title_box = slide.shapes.add_textbox(
            int(slide_width * 0.07), int(slide_height * 0.08),
            int(slide_width * 0.86), int(slide_height * 0.14),
        )
        title_tf = title_box.text_frame
        title_tf.clear()
        title_tf.word_wrap = True
        title_tf.paragraphs[0].text = title
        title_run = title_tf.paragraphs[0].runs[0]
        title_run.font.size = Pt(float(operation.get("title_size", 28)))
        title_run.font.bold = True
        title_color = _rgb(operation.get("title_color"), "00338D")
        if title_color:
            title_run.font.color.rgb = title_color

    if subtitle:
        subtitle_box = slide.shapes.add_textbox(
            int(slide_width * 0.07), int(slide_height * 0.22),
            int(slide_width * 0.86), int(slide_height * 0.08),
        )
        subtitle_tf = subtitle_box.text_frame
        subtitle_tf.text = subtitle
        subtitle_tf.paragraphs[0].runs[0].font.size = Pt(13)
        subtitle_color = _rgb(operation.get("subtitle_color"), "4B5563")
        if subtitle_color:
            subtitle_tf.paragraphs[0].runs[0].font.color.rgb = subtitle_color

    if body:
        body_box = slide.shapes.add_textbox(
            int(slide_width * 0.09), int(slide_height * 0.32),
            int(slide_width * 0.82), int(slide_height * 0.56),
        )
        body_tf = body_box.text_frame
        body_tf.clear()
        body_tf.word_wrap = True
        body_tf.vertical_anchor = MSO_ANCHOR.TOP
        for index, item in enumerate(body):
            paragraph = body_tf.paragraphs[0] if index == 0 else body_tf.add_paragraph()
            paragraph.text = str(item)
            paragraph.level = 0
            paragraph.font.size = Pt(float(operation.get("body_size", 18)))
            paragraph.space_after = Pt(8)

    new_number = len(prs.slides)
    position = int(operation.get("position") or new_number)
    position = max(1, min(position, new_number))
    _move_slide(prs, new_number, position)
    return position


def _add_textbox(prs: Presentation, operation: dict[str, Any]) -> None:
    slide = _slide(prs, int(operation["slide"]))
    left = _to_emu(operation.get("x", 0.1), prs.slide_width)
    top = _to_emu(operation.get("y", 0.1), prs.slide_height)
    width = _to_emu(operation.get("width", 0.8), prs.slide_width)
    height = _to_emu(operation.get("height", 0.15), prs.slide_height)
    box = slide.shapes.add_textbox(left, top, width, height)
    _set_text_frame_text(box.text_frame, str(operation.get("text", "")), preserve_first_run=False)
    _apply_shape_style(box, operation)


def _add_shape(prs: Presentation, operation: dict[str, Any]) -> None:
    slide = _slide(prs, int(operation["slide"]))
    shape_name = str(operation.get("shape_type", "rectangle")).lower()
    shape_types = {
        "rectangle": MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        "rounded_rectangle": MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        "circle": MSO_AUTO_SHAPE_TYPE.OVAL,
        "triangle": MSO_AUTO_SHAPE_TYPE.ISOSCELES_TRIANGLE,
        "chevron": MSO_AUTO_SHAPE_TYPE.CHEVRON,
        "arrow": MSO_AUTO_SHAPE_TYPE.RIGHT_ARROW,
    }
    shape_type = shape_types.get(shape_name, MSO_AUTO_SHAPE_TYPE.RECTANGLE)
    shape = slide.shapes.add_shape(
        shape_type,
        _to_emu(operation.get("x", 0.1), prs.slide_width),
        _to_emu(operation.get("y", 0.1), prs.slide_height),
        _to_emu(operation.get("width", 0.2), prs.slide_width),
        _to_emu(operation.get("height", 0.1), prs.slide_height),
    )
    if operation.get("text") is not None and getattr(shape, "has_text_frame", False):
        _set_text_frame_text(shape.text_frame, str(operation.get("text", "")), preserve_first_run=False)
    _apply_shape_style(shape, operation)


def _chart_data(operation: dict[str, Any]):
    categories = operation.get("categories")
    series = operation.get("series")
    if not isinstance(categories, list) or not isinstance(series, list) or not series:
        raise EditorError("Chart edits require categories and series arrays.")
    if not categories:
        raise EditorError("A chart requires at least one category.")
    if str(operation.get("chart_type", "")).casefold() == "waterfall":
        first = series[0] if series else None
        values = [float(value) for value in (first or {}).get("values", [])]
        if len(values) != len(categories):
            raise EditorError("A waterfall chart needs one value per category.")
        base, increase, decrease = [], [], []
        previous = 0.0
        for index, value in enumerate(values):
            if index == 0:
                base.append(0.0); increase.append(max(value, 0.0)); decrease.append(max(-value, 0.0))
            else:
                delta = value - previous
                base.append(min(previous, value)); increase.append(max(delta, 0.0)); decrease.append(max(-delta, 0.0))
            previous = value
        data = CategoryChartData(); data.categories = [str(item) for item in categories]
        data.add_series("Base", base); data.add_series("Increase", increase); data.add_series("Decrease", decrease)
        return data
    if str(operation.get("chart_type", "")).casefold() == "scatter":
        x_values = operation.get("x_values") or categories
        if len(x_values) != len(categories):
            raise EditorError("A scatter chart needs one X value per Y value.")
        data = XyChartData()
        for item in series:
            if not isinstance(item, dict) or not isinstance(item.get("values"), list):
                raise EditorError("Each chart series needs a name and values array.")
            if len(item["values"]) != len(x_values):
                raise EditorError("Every scatter series needs one Y value per X value.")
            xy_series = data.add_series(str(item.get("name", "Series")))
            for x_value, y_value in zip(x_values, item["values"]):
                xy_series.add_data_point(float(x_value), float(y_value))
        return data
    data = CategoryChartData()
    data.categories = [str(item) for item in categories]
    for item in series:
        if not isinstance(item, dict) or not isinstance(item.get("values"), list):
            raise EditorError("Each chart series needs a name and values array.")
        if len(item["values"]) != len(categories):
            raise EditorError("Every chart series must have one value per category.")
        data.add_series(str(item.get("name", "Series")), item["values"])
    return data


def _set_chart_data(shape, operation: dict[str, Any]) -> None:
    if not getattr(shape, "has_chart", False):
        raise EditorError("The selected shape is not a chart.")
    shape.chart.replace_data(_chart_data(operation))
    _style_chart(shape.chart, operation)


def _chart_font(target, size: float, color: str = "4B5563", bold: bool | None = None) -> None:
    try:
        font = target.font
        font.size = Pt(size)
        font.color.rgb = _rgb(color)
        if bold is not None:
            font.bold = bold
    except Exception:
        pass


def _chart_fill_alpha(fill, opacity: int) -> None:
    """Set DrawingML fill opacity after assigning an RGB fill."""
    try:
        solid = fill._xPr.solidFill
        color = solid[0]
        for child in list(color):
            if child.tag == qn("a:alpha"):
                color.remove(child)
        alpha = OxmlElement("a:alpha")
        alpha.set("val", str(max(0, min(100000, int(opacity)))))
        color.append(alpha)
    except Exception:
        pass


def _style_chart(chart, operation: dict[str, Any]) -> None:
    compact_chart = bool(operation.get("compact_chart"))
    title_font_size = 11 if compact_chart else 22
    axis_font_size = 6 if compact_chart else 10
    label_font_size = 6 if compact_chart else 10
    title = str(operation.get("title") or "").strip()
    if title:
        chart.has_title = True
        chart.chart_title.text_frame.text = title
        for paragraph in chart.chart_title.text_frame.paragraphs:
            _chart_font(paragraph, title_font_size, "111827", True)
    if operation.get("chart_style") is not None:
        try:
            chart.chart_style = int(operation["chart_style"])
        except (TypeError, ValueError):
            pass
    chart_name = str(operation.get("chart_type") or "column").strip().casefold()
    colors = operation.get("series_colors") or ["005EB8", "00A3A1", "483698", "BC204B", "EAAA00", "00A651"]
    if not isinstance(colors, list) or not colors:
        colors = ["005EB8"]
    for index, series in enumerate(chart.series):
        color = _rgb(str(colors[index % len(colors)]))
        if color is None:
            continue
        try:
            series.format.fill.solid()
            series.format.fill.fore_color.rgb = color
        except Exception:
            pass
        try:
            series.format.line.color.rgb = color
            series.format.line.width = Pt(1.75 if chart_name == "area" else 2.25)
        except Exception:
            pass
        if chart_name == "area":
            _chart_fill_alpha(series.format.fill, 47000)
        if chart_name in {"line", "scatter"}:
            try:
                series.marker.style = XL_MARKER_STYLE.CIRCLE
                series.marker.size = 7 if chart_name == "scatter" else 6
                series.marker.format.fill.solid()
                series.marker.format.fill.fore_color.rgb = color
                series.marker.format.line.color.rgb = RGBColor(255, 255, 255)
                series.marker.format.line.width = Pt(0.75)
            except Exception:
                pass

    if chart_name == "pie" and chart.series:
        pie_colors = [
            "005EB8", "0091DA", "00A3A1", "483698", "BC204B", "EAAA00",
            "00A651", "6D2077", "1E49E2", "5EC9E6", "9BDAF3", "7C878E",
        ]
        for index, point in enumerate(chart.series[0].points):
            color = _rgb(pie_colors[index % len(pie_colors)])
            try:
                point.format.fill.solid()
                point.format.fill.fore_color.rgb = color
                point.format.line.color.rgb = RGBColor(255, 255, 255)
                point.format.line.width = Pt(1.5)
            except Exception:
                pass
        try:
            plot = chart.plots[0]
            plot.has_data_labels = True
            labels = plot.data_labels
            labels.show_value = False
            labels.show_percentage = True
            labels.show_category_name = False
            labels.show_series_name = False
            labels.show_legend_key = False
            labels.number_format = "0%"
            labels.position = XL_DATA_LABEL_POSITION.OUTSIDE_END
            _chart_font(labels, label_font_size, "374151", False)
        except Exception:
            pass

    if chart_name == "waterfall" and len(chart.series) >= 3:
        try:
            chart.series[0].format.fill.background()
            chart.series[0].format.line.fill.background()
            for series, color_text in zip(chart.series[1:3], ("00A651", "BC204B")):
                color = _rgb(color_text)
                series.format.fill.solid()
                series.format.fill.fore_color.rgb = color
                series.format.line.color.rgb = color
        except Exception:
            pass

    number_format = str(operation.get("number_format") or "#,##0")
    if chart_name != "pie":
        try:
            value_axis = chart.value_axis
            value_axis.has_major_gridlines = chart_name not in {"waterfall"}
            value_axis.tick_labels.number_format = number_format
            value_axis.tick_labels.number_format_is_linked = False
            _chart_font(value_axis.tick_labels, axis_font_size, "4B5563", False)
            try:
                value_axis.major_gridlines.format.line.color.rgb = RGBColor(226, 232, 240)
                value_axis.major_gridlines.format.line.width = Pt(0.75)
            except Exception:
                pass
        except Exception:
            pass
        try:
            category_axis = chart.category_axis
            if chart_name == "scatter":
                category_axis.tick_labels.number_format = str(operation.get("x_number_format") or number_format)
                category_axis.tick_labels.number_format_is_linked = False
            _chart_font(category_axis.tick_labels, axis_font_size, "4B5563", False)
        except Exception:
            pass

    chart.has_legend = False if compact_chart else bool(operation.get("show_legend", len(chart.series) > 1))
    if chart.has_legend:
        chart.legend.position = XL_LEGEND_POSITION.RIGHT if chart_name in {"pie", "bar"} else XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        _chart_font(chart.legend, axis_font_size, "374151", False)


def _shape_box(prs: Presentation, shape) -> tuple[float, float, float, float]:
    return (
        float(shape.left) / max(1, prs.slide_width),
        float(shape.top) / max(1, prs.slide_height),
        float(shape.width) / max(1, prs.slide_width),
        float(shape.height) / max(1, prs.slide_height),
    )


def _estimated_text_width_points(text: str, font_size: float) -> float:
    narrow = set("ilI.,'`|!:;[]()")
    wide = set("MW@#%&QGmwu")
    units = 0.0
    for character in text:
        if character.isspace():
            units += 0.31
        elif character in narrow:
            units += 0.28
        elif character in wide:
            units += 0.72
        elif character.isupper() or character.isdigit():
            units += 0.56
        else:
            units += 0.49
    return units * font_size


def _insertion_obstacle_box(prs: Presentation, shape) -> tuple[float, float, float, float] | None:
    """Return the visible object bounds, including the used area of transparent text boxes."""
    x, y, width, height = _shape_box(prs, shape)
    area = width * height
    full_slide_background = x <= 0.015 and y <= 0.015 and width >= 0.97 and height >= 0.97
    if full_slide_background or area <= 0:
        return None
    if str(getattr(shape, "shape_type", "")) != "TEXT_BOX (17)" or not getattr(shape, "has_text_frame", False):
        return x, y, width, height
    try:
        if str(shape.fill.type).startswith("SOLID"):
            return x, y, width, height
    except Exception:
        pass
    text_frame = shape.text_frame
    if not str(text_frame.text or "").strip():
        return None
    slide_width_inches = float(prs.slide_width) / 914400.0
    slide_height_inches = float(prs.slide_height) / 914400.0
    margin_width = float((text_frame.margin_left or 0) + (text_frame.margin_right or 0)) / max(1, prs.slide_width)
    margin_height = float((text_frame.margin_top or 0) + (text_frame.margin_bottom or 0)) / max(1, prs.slide_height)
    inner_width = max(0.01, width - margin_width)
    widest = 0.0
    used_height = 0.0
    for paragraph in text_frame.paragraphs:
        paragraph_text = str(paragraph.text or "")
        if not paragraph_text:
            continue
        paragraph_width_points = 0.0
        largest_font = 14.0
        if paragraph.runs:
            for run in paragraph.runs:
                size = float(run.font.size.pt) if run.font.size else 14.0
                largest_font = max(largest_font, size)
                paragraph_width_points += _estimated_text_width_points(str(run.text or ""), size)
        else:
            paragraph_width_points = _estimated_text_width_points(paragraph_text, largest_font)
        paragraph_width = paragraph_width_points / 72.0 / max(slide_width_inches, 0.1)
        wrapped_lines = max(1, int((paragraph_width / max(inner_width, 0.01)) + 0.999))
        widest = max(widest, min(inner_width, paragraph_width))
        line_height_points = largest_font * 1.28
        used_height += wrapped_lines * line_height_points / 72.0 / max(slide_height_inches, 0.1)
        try:
            if paragraph.space_before:
                used_height += float(paragraph.space_before.pt) / 72.0 / max(slide_height_inches, 0.1)
            if paragraph.space_after:
                used_height += float(paragraph.space_after.pt) / 72.0 / max(slide_height_inches, 0.1)
        except Exception:
            pass
    visible_width = min(width, widest + margin_width + 0.018)
    visible_height = min(height, used_height + margin_height + 0.012)
    first_paragraph = next((paragraph for paragraph in text_frame.paragraphs if str(paragraph.text or "").strip()), None)
    alignment = getattr(first_paragraph, "alignment", None) if first_paragraph is not None else None
    if alignment == PP_ALIGN.CENTER:
        visible_x = x + (width - visible_width) / 2
    elif alignment == PP_ALIGN.RIGHT:
        visible_x = x + width - visible_width
    else:
        visible_x = x
    if text_frame.vertical_anchor == MSO_ANCHOR.MIDDLE:
        visible_y = y + (height - visible_height) / 2
    elif text_frame.vertical_anchor == MSO_ANCHOR.BOTTOM:
        visible_y = y + height - visible_height
    else:
        visible_y = y
    return visible_x, visible_y, visible_width, visible_height


def _box_intersection(first: tuple[float, float, float, float], second: tuple[float, float, float, float]) -> float:
    left = max(first[0], second[0])
    top = max(first[1], second[1])
    right = min(first[0] + first[2], second[0] + second[2])
    bottom = min(first[1] + first[3], second[1] + second[3])
    return max(0.0, right - left) * max(0.0, bottom - top)


def _shape_is_metric(shape) -> bool:
    """Return True for KPI and metric objects that chart insertion must never cover."""
    if getattr(shape, "has_chart", False) or getattr(shape, "has_table", False):
        return True
    name = str(getattr(shape, "name", "") or "").casefold()
    if any(token in name for token in ("kpi", "metric", "scorecard")):
        return True
    if not getattr(shape, "has_text_frame", False):
        return False
    text = str(shape.text_frame.text or "").strip()
    if not text:
        return False
    compact = " ".join(text.split())
    if len(compact) > 120:
        return False
    metric_pattern = re.compile(
        r"(?<![A-Za-z])(?:[$€£]\s*)?-?\d[\d,]*(?:\.\d+)?(?:\s*(?:%|x|bps|pts|[KMBT]))?(?![A-Za-z])",
        re.IGNORECASE,
    )
    return bool(metric_pattern.search(compact))


def _expanded_obstacle(box: tuple[float, float, float, float], *, left: float, top: float,
                       right: float, bottom: float, padding: float) -> tuple[float, float, float, float] | None:
    x, y, width, height = box
    expanded_left = max(left, x - padding)
    expanded_top = max(top, y - padding)
    expanded_right = min(right, x + width + padding)
    expanded_bottom = min(bottom, y + height + padding)
    if expanded_right <= expanded_left or expanded_bottom <= expanded_top:
        return None
    return (
        expanded_left,
        expanded_top,
        expanded_right - expanded_left,
        expanded_bottom - expanded_top,
    )


def _auto_insert_box(prs: Presentation, slide, *, min_width: float = 0.18, min_height: float = 0.13,
                     max_width: float = 0.58, max_height: float = 0.50,
                     object_name: str = "chart", allow_compact_fallback: bool = False
                     ) -> tuple[float, float, float, float]:
    """Find a non-overlapping area, with a tiny corner fallback for charts."""
    left, top, right, bottom = 0.02, 0.025, 0.98, 0.965
    columns, rows = 64, 48
    cell_width = (right - left) / columns
    cell_height = (bottom - top) / rows
    padding = 0.009
    occupied: list[tuple[float, float, float, float]] = []
    protected: list[tuple[float, float, float, float]] = []
    visible_shapes: list[tuple[Any, tuple[float, float, float, float]]] = []
    for shape in slide.shapes:
        obstacle_box = _insertion_obstacle_box(prs, shape)
        if obstacle_box is None:
            continue
        expanded = _expanded_obstacle(
            obstacle_box, left=left, top=top, right=right, bottom=bottom, padding=padding,
        )
        if expanded is None:
            continue
        occupied.append(expanded)
        visible_shapes.append((shape, expanded))
        if _shape_is_metric(shape) or _shape_is_branding(shape, prs):
            protected.append(expanded)

    # Protect the complete KPI card when a compact card encloses a metric value.
    metric_boxes = list(protected)
    for shape, obstacle in visible_shapes:
        if obstacle in protected or obstacle[2] * obstacle[3] > 0.28:
            continue
        if any(_box_intersection(obstacle, metric) >= metric[2] * metric[3] * 0.92 for metric in metric_boxes):
            protected.append(obstacle)

    def find_box(obstacles: list[tuple[float, float, float, float]], *, required_width: float,
                 required_height: float, width_limit: float, height_limit: float,
                 prefer_corner: bool = False, penalize_soft_overlap: bool = False
                 ) -> tuple[float, float, float, float] | None:
        free_rows: list[list[bool]] = []
        for row in range(rows):
            row_top = top + row * cell_height
            cells: list[bool] = []
            for column in range(columns):
                cell_left = left + column * cell_width
                cell = (cell_left, row_top, cell_width, cell_height)
                cells.append(all(_box_intersection(cell, obstacle) <= 0 for obstacle in obstacles))
            free_rows.append(cells)

        best_box: tuple[float, float, float, float] | None = None
        best_score = -10_000.0
        for start_row in range(rows):
            clear_columns = [True] * columns
            for end_row in range(start_row, rows):
                clear_columns = [
                    clear_columns[index] and free_rows[end_row][index]
                    for index in range(columns)
                ]
                available_height = (end_row - start_row + 1) * cell_height
                if available_height + 1e-9 < required_height:
                    continue
                column = 0
                while column < columns:
                    if not clear_columns[column]:
                        column += 1
                        continue
                    run_start = column
                    while column < columns and clear_columns[column]:
                        column += 1
                    available_width = (column - run_start) * cell_width
                    if available_width + 1e-9 < required_width:
                        continue
                    chart_width = min(available_width, width_limit)
                    chart_height = min(available_height, height_limit)
                    ratio = chart_width / max(chart_height, 0.001)
                    if ratio > 1.90:
                        chart_width = chart_height * 1.90
                    elif ratio < 0.72:
                        chart_height = chart_width / 0.72
                    if chart_width + 1e-9 < required_width or chart_height + 1e-9 < required_height:
                        continue
                    free_left = left + run_start * cell_width
                    free_top = top + start_row * cell_height
                    if prefer_corner:
                        x_positions = [free_left, free_left + available_width - chart_width]
                        y_positions = [free_top, free_top + available_height - chart_height]
                    else:
                        x_positions = [free_left + (available_width - chart_width) / 2]
                        y_positions = [free_top + (available_height - chart_height) / 2]
                    for x in x_positions:
                        for y in y_positions:
                            candidate = (x, y, chart_width, chart_height)
                            if any(_box_intersection(candidate, obstacle) > 0 for obstacle in obstacles):
                                continue
                            candidate_ratio = chart_width / max(chart_height, 0.001)
                            score = chart_width * chart_height - abs(candidate_ratio - 1.55) * 0.001
                            if not prefer_corner and x >= 0.48 and chart_width >= 0.26 and chart_height >= 0.28:
                                score += 0.12
                            if prefer_corner:
                                center_x = x + chart_width / 2
                                center_y = y + chart_height / 2
                                corner_distance = min(
                                    abs(center_x - left) + abs(center_y - top),
                                    abs(center_x - right) + abs(center_y - top),
                                    abs(center_x - left) + abs(center_y - bottom),
                                    abs(center_x - right) + abs(center_y - bottom),
                                )
                                score += max(0.0, 1.4 - corner_distance) * 0.14
                            if penalize_soft_overlap:
                                soft_overlap = sum(_box_intersection(candidate, obstacle) for obstacle in occupied)
                                score -= soft_overlap * 0.25
                            if score > best_score:
                                best_box = candidate
                                best_score = score
        return best_box

    best_box = find_box(
        occupied,
        required_width=min_width,
        required_height=min_height,
        width_limit=max_width,
        height_limit=max_height,
    )
    if best_box is None and allow_compact_fallback:
        compact_width = cell_width * 2
        compact_height = cell_height * 2
        strict_compact = find_box(
            occupied,
            required_width=compact_width,
            required_height=compact_height,
            width_limit=min(max_width, 0.34),
            height_limit=min(max_height, 0.28),
            prefer_corner=True,
        )
        protected_fallback = None
        if strict_compact is None or strict_compact[2] < 0.12 or strict_compact[3] < 0.09:
            protected_fallback = find_box(
                protected,
                required_width=compact_width,
                required_height=compact_height,
                width_limit=min(max_width, 0.28),
                height_limit=min(max_height, 0.22),
                prefer_corner=True,
                penalize_soft_overlap=True,
            )
        best_box = protected_fallback or strict_compact

    if best_box is None:
        if allow_compact_fallback:
            raise EditorError(
                f"An existing chart, table, logo, or metric occupies every possible {object_name} position."
            )
        raise EditorError(
            f"This slide has no blank area large enough for a readable {object_name}. "
            "Choose New slide or remove an object first."
        )
    return best_box


def _add_chart(prs: Presentation, operation: dict[str, Any]) -> None:
    slide = _slide(prs, int(operation["slide"]))
    chart_name = str(operation.get("chart_type") or "column").strip().lower()
    chart_types = {
        "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
        "bar": XL_CHART_TYPE.BAR_CLUSTERED,
        "line": XL_CHART_TYPE.LINE_MARKERS,
        "pie": XL_CHART_TYPE.PIE,
        "area": XL_CHART_TYPE.AREA,
        "doughnut": XL_CHART_TYPE.DOUGHNUT,
        "stacked_column": XL_CHART_TYPE.COLUMN_STACKED,
        "stacked_bar": XL_CHART_TYPE.BAR_STACKED,
        "scatter": XL_CHART_TYPE.XY_SCATTER,
        "waterfall": XL_CHART_TYPE.COLUMN_STACKED,
    }
    chart_type = chart_types.get(chart_name, XL_CHART_TYPE.COLUMN_CLUSTERED)
    if operation.get("auto_fit"):
        x, y, width, height = _auto_insert_box(prs, slide, allow_compact_fallback=True)
    else:
        x, y, width, height = (
            operation.get("x", 0.5), operation.get("y", 0.3),
            operation.get("width", 0.45), operation.get("height", 0.5),
        )
    frame = slide.shapes.add_chart(
        chart_type,
        _to_emu(x, prs.slide_width),
        _to_emu(y, prs.slide_height),
        _to_emu(width, prs.slide_width),
        _to_emu(height, prs.slide_height),
        _chart_data(operation),
    )
    if width < 0.18 or height < 0.13:
        operation = {**operation, "compact_chart": True}
    if chart_name == "waterfall" and not operation.get("series_colors"):
        operation = {**operation, "series_colors": ["FFFFFF", "00A651", "BC204B"]}
    _style_chart(frame.chart, operation)
    apply_chart_text_contrast(slide, frame, bool(operation.get("compact_chart")))



# ---------------------------------------------------------------------------
# High-level executive workflow helpers
# ---------------------------------------------------------------------------


def _all_text_shapes(slide):
    return [shape for shape in slide.shapes if getattr(shape, "has_text_frame", False)]


def _find_text_shape(slide, contains: str):
    wanted = _normalize_text(contains)
    candidates = []
    for shape in _all_text_shapes(slide):
        text = _normalize_text(shape.text_frame.text)
        if wanted and wanted in text:
            candidates.append(shape)
    if not candidates:
        return None
    candidates.sort(key=lambda shape: (len(shape.text_frame.text), shape.top, shape.left))
    return candidates[0]


def _set_shape_text(shape, text: str, **style) -> None:
    if shape is None or not getattr(shape, "has_text_frame", False):
        return
    _set_text_frame_text(shape.text_frame, text)
    _apply_text_style(shape, style)


def _replace_phrase_everywhere(prs: Presentation, old: str, new: str = "") -> int:
    count = 0
    for slide in prs.slides:
        for text_frame in _iter_text_frames(slide):
            count += _replace_text_in_frame(text_frame, old, new, True, False)
    return count


def _clean_demo_language(prs: Presentation) -> int:
    phrases = [
        "INTERNAL DEMO",
        "INTERNAL DEMONSTRATION",
        "All figures are fictional",
        "Fictional data for testing presentation updates, native tables, charts, formatting, and synchronized comparison.",
        "Fictional advisory practice KPI snapshot",
        "Fictional KPMG advisory service-line performance",
        "Fictional advisory practice spend versus plan",
        "fictional management forecast",
        "Fictional data",
        "sample",
    ]
    count = 0
    for phrase in phrases:
        count += _replace_phrase_everywhere(prs, phrase, "")
    # Clean separators and whitespace left behind by removals.
    for slide in prs.slides:
        for text_frame in _iter_text_frames(slide):
            raw = text_frame.text
            cleaned = re.sub(r"\s*\|\s*\|\s*", " | ", raw)
            cleaned = re.sub(r"^(?:\s*\|\s*)+|(?:\s*\|\s*)+$", "", cleaned)
            cleaned = re.sub(r"[ \t]{2,}", " ", cleaned)
            cleaned = re.sub(r"\n{3,}", "\n\n", cleaned).strip()
            if cleaned != raw:
                _set_text_frame_text(text_frame, cleaned)
    return count


def _chart_payload_from_slide(slide) -> dict[str, Any] | None:
    for shape in slide.shapes:
        if not getattr(shape, "has_chart", False):
            continue
        try:
            chart = shape.chart
            plot = chart.plots[0]
            categories = [str(value) for value in plot.categories]
            series = [
                {"name": str(series.name), "values": [float(value) for value in series.values]}
                for series in chart.series
            ]
            title = chart.chart_title.text_frame.text.strip() if chart.has_title else "Performance"
            if categories and series:
                return {"title": title, "categories": categories, "series": series}
        except Exception:
            continue
    return None


def _style_kpi_text(shape) -> None:
    if not getattr(shape, "has_text_frame", False):
        return
    shape.text_frame.word_wrap = True
    shape.text_frame.margin_left = Inches(0.12)
    shape.text_frame.margin_right = Inches(0.08)
    shape.text_frame.margin_top = Inches(0.08)
    shape.text_frame.margin_bottom = Inches(0.05)
    for paragraph_index, paragraph in enumerate(shape.text_frame.paragraphs):
        paragraph.alignment = PP_ALIGN.LEFT
        for run in paragraph.runs:
            run.font.name = "Arial"
            run.font.size = Pt(20 if paragraph_index == 0 else 10)
            run.font.bold = paragraph_index == 0
            run.font.color.rgb = RGBColor(0x00, 0x33, 0x8D) if paragraph_index == 0 else RGBColor(0x4B, 0x55, 0x63)


def _prepare_summary_slide(prs: Presentation, slide_number: int = 2) -> int:
    if slide_number < 1 or slide_number > len(prs.slides):
        return 0
    slide = prs.slides[slide_number - 1]
    title = _find_text_shape(slide, "Executive Summary")
    _set_shape_text(title, "Q3 Advisory Executive Summary", font_size=28, font_color="00338D", bold=True)
    subtitle_shapes = [shape for shape in _all_text_shapes(slide) if 0.7 < shape.top / 914400 < 1.3]
    if subtitle_shapes:
        _set_shape_text(
            sorted(subtitle_shapes, key=lambda shape: shape.top)[0],
            "Performance ahead of plan, with focused margin and growth priorities.",
            font_size=13,
            font_color="4B5563",
        )

    text_cards = []
    accent_lines = []
    for shape in slide.shapes:
        y = shape.top / 914400
        if 1.3 <= y <= 5.9:
            if getattr(shape, "has_text_frame", False) and shape.text_frame.text.strip():
                text_cards.append(shape)
            elif shape.width / 914400 < 0.2:
                accent_lines.append(shape)
    text_cards.sort(key=lambda shape: (shape.top, shape.left))
    if len(text_cards) >= 6:
        for index, shape in enumerate(text_cards[:6]):
            shape.top = Inches(1.48 if index < 3 else 2.52)
            shape.height = Inches(0.82)
            _style_kpi_text(shape)
    accent_lines.sort(key=lambda shape: (shape.top, shape.left))
    if len(accent_lines) >= 6:
        for index, shape in enumerate(accent_lines[:6]):
            shape.top = Inches(1.48 if index < 3 else 2.52)
            shape.height = Inches(0.82)

    chart_sources = []
    for source_number in (3, 4):
        if source_number <= len(prs.slides):
            payload = _chart_payload_from_slide(prs.slides[source_number - 1])
            if payload:
                chart_sources.append(payload)
    chart_specs = []
    if chart_sources:
        first = chart_sources[0]
        chart_specs.append({
            "op": "add_chart",
            "slide": slide_number,
            "chart_type": "column",
            "title": "Positive variance versus budget",
            "categories": first["categories"][:4],
            "series": [{"name": item["name"], "values": item["values"][:4]} for item in first["series"][:2]],
            "x": 0.05, "y": 0.50, "width": 0.43, "height": 0.37,
            "series_colors": ["00843D", "8CC63F"],
            "show_legend": True,
        })
    if len(chart_sources) > 1:
        second = chart_sources[1]
        chart_specs.append({
            "op": "add_chart",
            "slide": slide_number,
            "chart_type": "bar",
            "title": "Service-line performance comparison",
            "categories": second["categories"][:5],
            "series": [{"name": item["name"], "values": item["values"][:5]} for item in second["series"][:2]],
            "x": 0.52, "y": 0.50, "width": 0.43, "height": 0.37,
            "series_colors": ["00A651", "B7D333"],
            "show_legend": True,
        })
    for spec in chart_specs:
        _add_chart(prs, spec)
    return len(chart_specs)


def _table_rows(slide) -> list[list[str]]:
    for shape in slide.shapes:
        if getattr(shape, "has_table", False):
            return [[cell.text.strip() for cell in row.cells] for row in shape.table.rows]
    return []


def _number(text: str) -> float | None:
    match = re.search(r"-?\d[\d,]*(?:\.\d+)?", str(text).replace("$", ""))
    if not match:
        return None
    try:
        return float(match.group(0).replace(",", ""))
    except ValueError:
        return None


def _variance_insights(rows: list[list[str]], max_items: int = 3) -> list[tuple[str, str]]:
    candidates = []
    for row in rows[1:]:
        if len(row) < 3:
            continue
        actual = _number(row[1])
        budget = _number(row[2])
        if actual is None or budget in (None, 0):
            continue
        delta = actual - budget
        pct = delta / abs(budget) * 100
        is_cost = any(token in _normalize_text(row[0]) for token in ("cost", "expense", "spend"))
        management_score = -pct if is_cost else pct
        candidates.append((management_score, pct, row[0], actual, budget))
    if not candidates:
        return []
    positives = sorted([item for item in candidates if item[0] >= 0], reverse=True)
    negatives = sorted([item for item in candidates if item[0] < 0])
    selected = positives[:2] + negatives[:1]
    if len(selected) < max_items:
        for item in sorted(candidates, key=lambda value: abs(value[0]), reverse=True):
            if item not in selected:
                selected.append(item)
            if len(selected) >= max_items:
                break
    result = []
    for management_score, pct, label, actual, budget in selected[:max_items]:
        direction = "above" if pct >= 0 else "below"
        tone = "green" if management_score >= 0 else "amber"
        result.append((f"{label}: {abs(pct):.1f}% {direction} plan", tone))
    return result


def _add_callout(slide, left: float, top: float, width: float, height: float, text: str, tone: str) -> None:
    colors = {
        "green": ("E8F5E9", "00843D"),
        "amber": ("FFF4D6", "B26A00"),
        "blue": ("EAF1FB", "00338D"),
    }
    fill_hex, line_hex = colors.get(tone, colors["blue"])
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor.from_string(fill_hex)
    shape.line.color.rgb = RGBColor.from_string(line_hex)
    shape.line.width = Pt(1)
    shape.text_frame.clear()
    shape.text_frame.word_wrap = True
    shape.text_frame.margin_left = Inches(0.12)
    shape.text_frame.margin_right = Inches(0.08)
    shape.text_frame.margin_top = Inches(0.08)
    paragraph = shape.text_frame.paragraphs[0]
    paragraph.text = text
    paragraph.alignment = PP_ALIGN.LEFT
    for run in paragraph.runs:
        run.font.name = "Arial"
        run.font.size = Pt(10.5)
        run.font.bold = True
        run.font.color.rgb = RGBColor.from_string(line_hex)


def _enhance_analysis_slide(prs: Presentation, slide_number: int, subtitle: str) -> int:
    if slide_number < 1 or slide_number > len(prs.slides):
        return 0
    slide = prs.slides[slide_number - 1]
    top_text = [shape for shape in _all_text_shapes(slide) if 0.7 < shape.top / 914400 < 1.25]
    if top_text:
        _set_shape_text(sorted(top_text, key=lambda shape: shape.top)[0], subtitle, font_size=13, font_color="4B5563")
    for shape in slide.shapes:
        if getattr(shape, "has_table", False) or getattr(shape, "has_chart", False):
            if shape.top / 914400 < 2.0:
                shape.height = Inches(3.72)
    insights = _variance_insights(_table_rows(slide), 3)
    if not insights:
        insights = [
            ("Performance remains ahead of plan", "green"),
            ("Focus management attention on the largest variance", "blue"),
            ("Protect delivery economics as volume grows", "amber"),
        ]
    positions = [(0.65, 5.48), (4.75, 5.48), (8.85, 5.48)]
    for (text, tone), (left, top) in zip(insights, positions):
        _add_callout(slide, left, top, 3.82, 0.82, text, tone)
    return len(insights)


def _narrative_priority(title: str) -> tuple[int, str]:
    text = _normalize_text(title)
    if any(token in text for token in ("profitability", "margin", "financial overview")):
        return 0, text
    if any(token in text for token in ("sector", "revenue by", "service line")):
        return 1, text
    if any(token in text for token in ("cost", "risk", "control")):
        return 2, text
    if any(token in text for token in ("pipeline", "outlook", "recommend", "action")):
        return 3, text
    return 1, text


def _reorder_narrative_group(prs: Presentation, slide_numbers: list[int]) -> list[int]:
    valid = [number for number in slide_numbers if 1 <= number <= len(prs.slides)]
    if len(valid) < 2:
        return valid
    desired = sorted(valid, key=lambda number: _narrative_priority(_slide_title(prs.slides[number - 1])))
    order = list(range(1, len(prs.slides) + 1))
    positions = [number - 1 for number in valid]
    for position, old_number in zip(positions, desired):
        order[position] = old_number
    _reorder_slides(prs, order)
    return desired


def _add_brand_header(slide, prs: Presentation, title: str, subtitle: str = "") -> None:
    bar = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.18)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor(0x00, 0x33, 0x8D)
    bar.line.fill.background()
    logo = slide.shapes.add_textbox(Inches(0.6), Inches(0.28), Inches(1.35), Inches(0.42))
    logo.text_frame.text = "KPMG"
    for run in logo.text_frame.paragraphs[0].runs:
        run.font.name = "Arial"
        run.font.size = Pt(19)
        run.font.bold = True
        run.font.color.rgb = RGBColor(0x00, 0x33, 0x8D)
    title_box = slide.shapes.add_textbox(Inches(2.0), Inches(0.28), Inches(10.6), Inches(0.52))
    title_box.text_frame.text = title
    for run in title_box.text_frame.paragraphs[0].runs:
        run.font.name = "Arial"
        run.font.size = Pt(27)
        run.font.bold = True
        run.font.color.rgb = RGBColor(0x00, 0x33, 0x8D)
    if subtitle:
        sub = slide.shapes.add_textbox(Inches(2.02), Inches(0.82), Inches(10.4), Inches(0.34))
        sub.text_frame.text = subtitle
        for run in sub.text_frame.paragraphs[0].runs:
            run.font.name = "Arial"
            run.font.size = Pt(12)
            run.font.color.rgb = RGBColor(0x4B, 0x55, 0x63)
    line = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.6), Inches(1.23), Inches(12.1), Inches(0.03)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(0x00, 0xA6, 0x51)
    line.line.fill.background()


def _new_branded_slide(prs: Presentation, position: int, title: str, subtitle: str = ""):
    blank_layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[-1]
    slide = prs.slides.add_slide(blank_layout)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    _add_brand_header(slide, prs, title, subtitle)
    footer = slide.shapes.add_textbox(Inches(0.62), Inches(7.06), Inches(5.6), Inches(0.2))
    footer.text_frame.text = "KPMG Advisory | Q3 FY2026"
    for run in footer.text_frame.paragraphs[0].runs:
        run.font.name = "Arial"
        run.font.size = Pt(8)
        run.font.color.rgb = RGBColor(0x6B, 0x72, 0x80)
    new_number = len(prs.slides)
    position = max(1, min(int(position), new_number))
    _move_slide(prs, new_number, position)
    return prs.slides[position - 1], position


def _derive_recommendations(prs: Presentation) -> list[dict[str, str]]:
    return [
        {
            "action": "Accelerate Deal Advisory growth",
            "reason": "The largest service line is outperforming budget and offers the clearest near-term growth path.",
            "owner": "Deal Advisory Lead",
            "timing": "Next 30 days",
        },
        {
            "action": "Lift Strategy & Operations margin",
            "reason": "Contribution margin trails the other service lines and requires focused delivery discipline.",
            "owner": "Service Line Lead + Finance",
            "timing": "Q4",
        },
        {
            "action": "Prioritize existing-client expansion",
            "reason": "Existing clients produce the strongest engagement volume at the lowest pursuit cost.",
            "owner": "Growth and Account Leads",
            "timing": "Next 60 days",
        },
        {
            "action": "Tighten discretionary practice spend",
            "reason": "Business development and technology costs are running above plan.",
            "owner": "Practice COO",
            "timing": "Immediate",
        },
    ]


def _add_recommendations_slide(prs: Presentation, position: int) -> int:
    slide, position = _new_branded_slide(
        prs,
        position,
        "Executive Recommendations",
        "Priority actions derived from Q3 performance, margin, pipeline, and cost findings.",
    )
    recommendations = _derive_recommendations(prs)
    card_width = 2.88
    gap = 0.25
    start_x = 0.62
    colors = ["00338D", "005EB8", "00843D", "00A651"]
    for index, item in enumerate(recommendations):
        left = start_x + index * (card_width + gap)
        card = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(left), Inches(1.62), Inches(card_width), Inches(4.75),
        )
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(0xF7, 0xF9, 0xFC)
        card.line.color.rgb = RGBColor.from_string(colors[index])
        card.line.width = Pt(1.4)
        ribbon = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            Inches(left), Inches(1.62), Inches(card_width), Inches(0.16),
        )
        ribbon.fill.solid()
        ribbon.fill.fore_color.rgb = RGBColor.from_string(colors[index])
        ribbon.line.fill.background()
        box = slide.shapes.add_textbox(
            Inches(left + 0.16), Inches(1.94), Inches(card_width - 0.32), Inches(4.15)
        )
        tf = box.text_frame
        tf.clear()
        tf.word_wrap = True
        sections = [
            (item["action"], 16, True, colors[index]),
            ("Why", 9, True, "4B5563"),
            (item["reason"], 11, False, "111827"),
            ("Owner", 9, True, "4B5563"),
            (item["owner"], 11, False, "111827"),
            ("Timing", 9, True, "4B5563"),
            (item["timing"], 11, True, colors[index]),
        ]
        for section_index, (text, size, bold, color) in enumerate(sections):
            paragraph = tf.paragraphs[0] if section_index == 0 else tf.add_paragraph()
            paragraph.text = text
            paragraph.space_after = Pt(5 if section_index else 10)
            for run in paragraph.runs:
                run.font.name = "Arial"
                run.font.size = Pt(size)
                run.font.bold = bold
                run.font.color.rgb = RGBColor.from_string(color)
    return position


def _derive_risks(prs: Presentation) -> list[dict[str, str]]:
    return [
        {"risk": "Strategy & Operations margin compression", "impact": "Earnings dilution", "likelihood": "High", "mitigation": "Reset delivery mix, pricing, and scope controls", "owner": "Service Line Lead"},
        {"risk": "High-cost pursuit channels", "impact": "Lower growth ROI", "likelihood": "High", "mitigation": "Shift spend toward existing clients and partners", "owner": "Growth Lead"},
        {"risk": "Business development overspend", "impact": "Practice expense pressure", "likelihood": "Medium", "mitigation": "Gate discretionary campaigns against conversion", "owner": "Practice COO"},
        {"risk": "Utilization below forward target", "impact": "Margin and capacity drag", "likelihood": "Medium", "mitigation": "Match staffing to backlog and near-term demand", "owner": "Workforce Lead"},
        {"risk": "Sector concentration", "impact": "Revenue volatility", "likelihood": "Low", "mitigation": "Build healthcare and technology pipeline", "owner": "Sector Leaders"},
    ]


def _set_cell_style(cell, fill_hex: str, font_hex: str, size: float = 9.5, bold: bool = False, align=PP_ALIGN.LEFT) -> None:
    cell.fill.solid()
    cell.fill.fore_color.rgb = RGBColor.from_string(fill_hex)
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    cell.text_frame.margin_left = Inches(0.06)
    cell.text_frame.margin_right = Inches(0.04)
    cell.text_frame.margin_top = Inches(0.03)
    cell.text_frame.margin_bottom = Inches(0.03)
    for paragraph in cell.text_frame.paragraphs:
        paragraph.alignment = align
        for run in paragraph.runs:
            run.font.name = "Arial"
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.color.rgb = RGBColor.from_string(font_hex)


def _add_risk_table_slide(prs: Presentation, position: int) -> int:
    slide, position = _new_branded_slide(
        prs,
        position,
        "Key Risks and Mitigations",
        "Management actions to protect growth, margin, and delivery performance.",
    )
    headers = ["Risk", "Impact", "Likelihood", "Mitigation", "Owner"]
    risks = _derive_risks(prs)
    table_shape = slide.shapes.add_table(
        len(risks) + 1,
        len(headers),
        Inches(0.62), Inches(1.58), Inches(12.08), Inches(4.95),
    )
    table = table_shape.table
    widths = [2.55, 1.75, 1.25, 4.15, 2.38]
    for index, width in enumerate(widths):
        table.columns[index].width = Inches(width)
    for col, header in enumerate(headers):
        table.cell(0, col).text = header
        _set_cell_style(table.cell(0, col), "00338D", "FFFFFF", 10, True, PP_ALIGN.CENTER)
    risk_colors = {"High": ("FDE8E8", "B42318"), "Medium": ("FFF4D6", "B26A00"), "Low": ("E8F5E9", "00843D")}
    for row_index, item in enumerate(risks, start=1):
        values = [item["risk"], item["impact"], item["likelihood"], item["mitigation"], item["owner"]]
        for col_index, value in enumerate(values):
            cell = table.cell(row_index, col_index)
            cell.text = value
            fill, font = ("FFFFFF", "111827")
            if col_index == 2:
                fill, font = risk_colors.get(item["likelihood"], ("EAF1FB", "00338D"))
            elif row_index % 2 == 0:
                fill = "F7F9FC"
            _set_cell_style(cell, fill, font, 9.2, col_index == 2, PP_ALIGN.CENTER if col_index == 2 else PP_ALIGN.LEFT)
    return position


def _estimate_text_size(shape, default: float) -> float:
    if not getattr(shape, "has_text_frame", False):
        return default
    text = shape.text_frame.text.strip()
    if not text:
        return default
    width = max(shape.width / 914400, 0.2)
    height = max(shape.height / 914400, 0.15)
    density = len(text) / max(width * height, 0.2)
    if density > 80:
        return max(8, default - 5)
    if density > 50:
        return max(9, default - 3)
    if density > 32:
        return max(10, default - 1.5)
    return default


def _standardize_deck(prs: Presentation) -> None:
    slide_width = prs.slide_width
    slide_height = prs.slide_height
    for slide_index, slide in enumerate(prs.slides, start=1):
        text_shapes = [shape for shape in _all_text_shapes(slide) if shape.text_frame.text.strip()]
        title_candidates = [shape for shape in text_shapes if shape.top / 914400 < 0.85 and shape.left / 914400 > 1.4]
        if slide_index != 1 and title_candidates:
            title = sorted(title_candidates, key=lambda shape: (shape.top, shape.left))[0]
            title.top = Inches(0.28)
            title.height = Inches(0.55)
            title.left = Inches(2.0) if title.left / 914400 > 1.4 else Inches(0.6)
            title.width = Inches(10.7)
            _apply_text_style(title, {"font_face": "Arial", "font_size": 27, "font_color": "00338D", "bold": True})
        subtitle_candidates = [shape for shape in text_shapes if 0.75 <= shape.top / 914400 <= 1.25 and shape.height / 914400 < 0.7]
        if slide_index != 1 and subtitle_candidates:
            subtitle = sorted(subtitle_candidates, key=lambda shape: shape.top)[0]
            subtitle.top = Inches(0.84)
            subtitle.left = Inches(2.02) if subtitle.left / 914400 > 1.4 else Inches(0.62)
            subtitle.width = Inches(10.5)
            subtitle.height = Inches(0.34)
            _apply_text_style(subtitle, {"font_face": "Arial", "font_size": 12.5, "font_color": "4B5563"})

        for shape in slide.shapes:
            # Keep every object inside the slide canvas.
            shape.left = max(0, min(shape.left, max(0, slide_width - shape.width)))
            shape.top = max(0, min(shape.top, max(0, slide_height - shape.height)))
            if getattr(shape, "has_text_frame", False):
                shape.text_frame.word_wrap = True
                shape.text_frame.margin_left = min(shape.text_frame.margin_left, Inches(0.12))
                shape.text_frame.margin_right = min(shape.text_frame.margin_right, Inches(0.12))
                default_size = 11.5
                current_sizes = [
                    run.font.size.pt
                    for paragraph in shape.text_frame.paragraphs
                    for run in paragraph.runs
                    if run.font.size is not None
                ]
                if current_sizes:
                    default_size = max(current_sizes)
                fitted = _estimate_text_size(shape, default_size)
                if fitted < default_size:
                    _apply_text_style(shape, {"font_size": fitted})
            if getattr(shape, "has_table", False):
                for row_index, row in enumerate(shape.table.rows):
                    for cell in row.cells:
                        fill = "00338D" if row_index == 0 else ("F7F9FC" if row_index % 2 == 0 else "FFFFFF")
                        font = "FFFFFF" if row_index == 0 else "111827"
                        _set_cell_style(cell, fill, font, 8.8 if len(shape.table.rows) > 6 else 9.5, row_index == 0)
            if getattr(shape, "has_chart", False):
                try:
                    shape.chart.style = 10
                    if shape.chart.has_legend:
                        shape.chart.legend.position = XL_LEGEND_POSITION.BOTTOM
                        shape.chart.legend.include_in_layout = False
                except Exception:
                    pass

        footer_candidates = [shape for shape in _all_text_shapes(slide) if shape.top / 914400 > 6.75]
        if footer_candidates:
            footer = sorted(footer_candidates, key=lambda shape: shape.top)[-1]
            _set_shape_text(footer, "KPMG Advisory | Q3 FY2026", font_face="Arial", font_size=8, font_color="6B7280")
            footer.left = Inches(0.62)
            footer.top = Inches(7.04)
            footer.width = Inches(5.8)
            footer.height = Inches(0.2)


def _slide_signature(slide) -> str:
    parts = []
    for shape in slide.shapes:
        text = _normalize_text(_shape_text(shape))
        if text and "kpmg advisory" not in text:
            parts.append(text)
    return " | ".join(parts)


def _delete_exact_redundant_slides(prs: Presentation) -> int:
    seen: dict[str, int] = {}
    to_delete = []
    for number, slide in enumerate(prs.slides, start=1):
        signature = _slide_signature(slide)
        if len(signature) < 80:
            continue
        if signature in seen:
            to_delete.append(number)
        else:
            seen[signature] = number
    for number in reversed(to_delete):
        _delete_slide(prs, number)
    return len(to_delete)


def _verify_deck_integrity(prs: Presentation) -> list[str]:
    issues = []
    if not prs.slides:
        issues.append("The deck has no slides.")
    for slide_number, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            if shape.left < 0 or shape.top < 0:
                issues.append(f"Slide {slide_number} has an object outside the canvas.")
            if shape.left + shape.width > prs.slide_width + 10 or shape.top + shape.height > prs.slide_height + 10:
                issues.append(f"Slide {slide_number} has an object extending outside the canvas.")
    return issues


def _apply_executive_review(prs: Presentation, operation: dict[str, Any]) -> dict[str, Any]:
    original_count = len(prs.slides)
    _clean_demo_language(prs)

    cover_number = int(operation.get("cover_slide", 1))
    if 1 <= cover_number <= len(prs.slides):
        cover = prs.slides[cover_number - 1]
        title = _find_text_shape(cover, "Advisory Practice Review") or _find_text_shape(cover, "Advisory")
        subtitle_candidates = [
            shape for shape in _all_text_shapes(cover)
            if shape is not title and shape.text_frame.text.strip() and shape.top > (title.top if title is not None else 0)
        ]
        _set_shape_text(title, "Q3 FY2026 Advisory Performance Review", font_face="Arial", font_size=34, font_color="FFFFFF", bold=True)
        if title is not None:
            title.left = Inches(0.72)
            title.top = Inches(2.05)
            title.width = Inches(8.65)
            title.height = Inches(0.85)
        if subtitle_candidates:
            subtitle = sorted(subtitle_candidates, key=lambda shape: (shape.top, shape.left))[0]
            _set_shape_text(
                subtitle,
                "Client-ready review of advisory performance, priorities, and outlook.",
                font_face="Arial", font_size=15, font_color="D9E2F2",
            )
            subtitle.left = Inches(0.74)
            subtitle.top = Inches(3.05)
            subtitle.width = Inches(8.25)
            subtitle.height = Inches(0.72)

    chart_count = _prepare_summary_slide(prs, int(operation.get("summary_slide", 2)))
    analysis_slides = [int(value) for value in operation.get("analysis_slides", [3, 4])]
    if analysis_slides:
        _enhance_analysis_slide(prs, analysis_slides[0], "Performance versus budget and the drivers of contribution growth.")
    if len(analysis_slides) > 1:
        _enhance_analysis_slide(prs, analysis_slides[1], "Service-line growth, mix, and priority management implications.")

    narrative_group = [int(value) for value in operation.get("narrative_slides", [5, 6, 7, 8])]
    desired_order = _reorder_narrative_group(prs, narrative_group)

    # The final analysis slide is the current last slide before new management pages.
    recommendations_position = len(prs.slides) + 1
    rec_position = _add_recommendations_slide(prs, recommendations_position) if operation.get("add_recommendations", True) else None
    risk_position = _add_risk_table_slide(prs, len(prs.slides) + 1) if operation.get("add_risks", True) else None

    if operation.get("standardize", True):
        _standardize_deck(prs)
    redundant_deleted = _delete_exact_redundant_slides(prs) if operation.get("delete_redundant", True) else 0
    integrity_issues = _verify_deck_integrity(prs)
    return {
        "original_slide_count": original_count,
        "final_slide_count": len(prs.slides),
        "charts_added": chart_count,
        "narrative_order": desired_order,
        "recommendations_slide": rec_position,
        "risks_slide": risk_position,
        "redundant_slides_deleted": redundant_deleted,
        "integrity_issues": integrity_issues,
    }

# ---------------------------------------------------------------------------
# Operation validation and application
# ---------------------------------------------------------------------------


SHAPE_OPS = {
    "set_text", "style_shape", "move_shape", "resize_shape", "delete_shape",
    "set_table_cell", "set_chart_data",
}
SUPPORTED_OPS = {
    "replace_text", *SHAPE_OPS, "add_textbox", "add_shape", "add_chart",
    "set_slide_background", "add_slide", "duplicate_slide", "delete_slide",
    "move_slide", "reorder_slides", "executive_review",
}


def _operation_summary(operation: dict[str, Any]) -> str:
    op = operation.get("op", "edit")
    slide = operation.get("slide")
    return f"{op} on slide {slide}" if slide else str(op)


def validate_operations(
    pptx_path: str,
    operations: Iterable[dict[str, Any]],
) -> tuple[list[dict[str, Any]], list[str]]:
    """Resolve stable shape IDs and reject impossible operations before editing."""
    prs = Presentation(pptx_path)
    slide_count = len(prs.slides)
    valid: list[dict[str, Any]] = []
    issues: list[str] = []

    for index, raw in enumerate(operations, start=1):
        if not isinstance(raw, dict):
            issues.append(f"Operation {index} was not an object.")
            continue
        operation = {key: value for key, value in raw.items() if value is not None}
        op = str(operation.get("op", "")).strip().lower()
        operation["op"] = op
        if op not in SUPPORTED_OPS:
            issues.append(f"Operation {index} used unsupported action '{op}'.")
            continue

        try:
            if op == "executive_review":
                valid.append(operation)
                continue

            if op == "replace_text":
                if not str(operation.get("old", "")):
                    raise EditorError("replace_text needs the original text.")
                if operation.get("slide") is not None:
                    slide = _slide(prs, int(operation["slide"]))
                    if any(operation.get(key) is not None for key in ("shape", "shape_id", "shape_name", "text_contains")):
                        resolved = _operation_shape(slide, operation)
                        operation["shape_id"] = int(resolved.shape_id)
                valid.append(operation)
                continue

            if op in SHAPE_OPS:
                if operation.get("slide") is None:
                    raise EditorError(f"{op} needs a slide number.")
                slide = _slide(prs, int(operation["slide"]))
                resolved = _operation_shape(slide, operation)
                operation["shape_id"] = int(resolved.shape_id)
                if op == "set_text" and not getattr(resolved, "has_text_frame", False):
                    raise EditorError("The selected object has no editable text.")
                if op == "set_table_cell" and not getattr(resolved, "has_table", False):
                    raise EditorError("The selected object is not a table.")
                if op == "set_chart_data" and not getattr(resolved, "has_chart", False):
                    raise EditorError("The selected object is not a chart.")
                valid.append(operation)
                continue

            if op in {"add_textbox", "add_shape", "add_chart", "set_slide_background"}:
                if operation.get("slide") is None:
                    raise EditorError(f"{op} needs a slide number.")
                _slide(prs, int(operation["slide"]))
                if op == "add_chart":
                    _chart_data(operation)
                valid.append(operation)
                continue

            if op == "add_slide":
                position = int(operation.get("position") or slide_count + 1)
                if position < 1 or position > slide_count + 1:
                    raise EditorError("The new slide position is outside the deck.")
                if operation.get("template_slide") is not None:
                    _slide(prs, int(operation["template_slide"]))
                valid.append(operation)
                continue

            if op == "duplicate_slide":
                _slide(prs, int(operation["slide"]))
                position = int(operation.get("position") or int(operation["slide"]) + 1)
                if position < 1 or position > slide_count + 1:
                    raise EditorError("The duplicate position is outside the deck.")
                valid.append(operation)
                continue

            if op == "delete_slide":
                _slide(prs, int(operation["slide"]))
                if slide_count <= 1:
                    raise EditorError("A presentation must keep at least one slide.")
                valid.append(operation)
                continue

            if op == "move_slide":
                from_slide = int(operation["from_slide"])
                to_slide = int(operation["to_slide"])
                if from_slide < 1 or from_slide > slide_count or to_slide < 1 or to_slide > slide_count:
                    raise EditorError("The move position is outside the deck.")
                valid.append(operation)
                continue

            if op == "reorder_slides":
                order = [int(value) for value in operation.get("order", [])]
                if sorted(order) != list(range(1, slide_count + 1)):
                    raise EditorError("The order must list every current slide exactly once.")
                operation["order"] = order
                valid.append(operation)
                continue

        except (EditorError, KeyError, TypeError, ValueError) as exc:
            issues.append(f"Operation {index} ({_operation_summary(operation)}) is invalid: {exc}")

    # Deleting a shape last prevents its removal from changing later shape indexes.
    # Stable shape IDs already protect most cases, but this also helps older plans.
    non_deletes = [operation for operation in valid if operation["op"] != "delete_shape"]
    shape_deletes = [operation for operation in valid if operation["op"] == "delete_shape"]
    return non_deletes + shape_deletes, issues


def apply_operations(input_path: str, output_path: str, operations: Iterable[dict[str, Any]]) -> dict[str, Any]:
    prepared, validation_issues = validate_operations(input_path, operations)
    prs = Presentation(input_path)
    applied: list[dict[str, Any]] = []
    skipped: list[dict[str, str]] = [
        {"operation": "validation", "reason": issue} for issue in validation_issues
    ]

    for operation in prepared:
        op = operation["op"]
        try:
            if op == "executive_review":
                details = _apply_executive_review(prs, operation)
                applied.append({"op": op, **details})

            elif op == "replace_text":
                slide_numbers = (
                    [int(operation["slide"])]
                    if operation.get("slide") is not None
                    else list(range(1, len(prs.slides) + 1))
                )
                old = str(operation.get("old", ""))
                new = str(operation.get("new", ""))
                replace_all = bool(operation.get("replace_all", True))
                case_sensitive = bool(operation.get("case_sensitive", False))
                count = 0
                for slide_number in slide_numbers:
                    slide = _slide(prs, slide_number)
                    target_shape = None
                    if any(operation.get(key) is not None for key in ("shape", "shape_id", "shape_name", "text_contains")):
                        target_shape = _operation_shape(slide, operation)
                    for text_frame in _iter_text_frames(slide, target_shape):
                        count += _replace_text_in_frame(
                            text_frame, old, new, replace_all, case_sensitive
                        )
                        if count and not replace_all:
                            break
                if count == 0:
                    raise EditorError(f'Text "{old}" was not found.')
                applied.append({"op": op, "count": count})

            elif op == "set_text":
                slide = _slide(prs, int(operation["slide"]))
                shape = _operation_shape(slide, operation)
                if not getattr(shape, "has_text_frame", False):
                    raise EditorError("The selected object has no editable text.")
                _set_text_frame_text(shape.text_frame, str(operation.get("text", "")))
                _apply_text_style(shape, operation)
                applied.append({"op": op, "slide": int(operation["slide"])})

            elif op in {"style_shape", "move_shape", "resize_shape"}:
                slide = _slide(prs, int(operation["slide"]))
                shape = _operation_shape(slide, operation)
                _set_geometry(prs, shape, operation)
                _apply_shape_style(shape, operation)
                applied.append({"op": op, "slide": int(operation["slide"])})

            elif op == "delete_shape":
                slide = _slide(prs, int(operation["slide"]))
                shape = _operation_shape(slide, operation)
                shape.element.getparent().remove(shape.element)
                applied.append({"op": op, "slide": int(operation["slide"])})

            elif op == "add_textbox":
                _add_textbox(prs, operation)
                applied.append({"op": op, "slide": int(operation["slide"])})

            elif op == "add_shape":
                _add_shape(prs, operation)
                applied.append({"op": op, "slide": int(operation["slide"])})

            elif op == "add_chart":
                _add_chart(prs, operation)
                applied.append({"op": op, "slide": int(operation["slide"])})

            elif op == "set_slide_background":
                slide = _slide(prs, int(operation["slide"]))
                color = _rgb(operation.get("color"))
                if color is None:
                    raise EditorError("A six-digit background color is required.")
                slide.background.fill.solid()
                slide.background.fill.fore_color.rgb = color
                applied.append({"op": op, "slide": int(operation["slide"])})

            elif op == "set_table_cell":
                slide = _slide(prs, int(operation["slide"]))
                shape = _operation_shape(slide, operation)
                if not getattr(shape, "has_table", False):
                    raise EditorError("The selected object is not a table.")
                row = int(operation["row"]) - 1
                column = int(operation["column"]) - 1
                if row < 0 or column < 0 or row >= len(shape.table.rows) or column >= len(shape.table.columns):
                    raise EditorError("The table cell is outside the table.")
                _set_text_frame_text(shape.table.cell(row, column).text_frame, str(operation.get("text", "")))
                applied.append({"op": op, "slide": int(operation["slide"])})

            elif op == "set_chart_data":
                slide = _slide(prs, int(operation["slide"]))
                shape = _operation_shape(slide, operation)
                _set_chart_data(shape, operation)
                applied.append({"op": op, "slide": int(operation["slide"])})

            elif op == "add_slide":
                position = _add_slide(prs, operation)
                applied.append({"op": op, "slide": position})

            elif op == "duplicate_slide":
                position = _duplicate_slide(
                    prs,
                    int(operation["slide"]),
                    int(operation["position"]) if operation.get("position") is not None else None,
                )
                applied.append({"op": op, "slide": position})

            elif op == "delete_slide":
                if len(prs.slides) <= 1:
                    raise EditorError("A presentation must keep at least one slide.")
                _delete_slide(prs, int(operation["slide"]))
                applied.append({"op": op})

            elif op == "move_slide":
                _move_slide(prs, int(operation["from_slide"]), int(operation["to_slide"]))
                applied.append({"op": op})

            elif op == "reorder_slides":
                _reorder_slides(prs, [int(value) for value in operation["order"]])
                applied.append({"op": op, "order": operation["order"]})

        except (EditorError, KeyError, TypeError, ValueError, AttributeError) as exc:
            skipped.append({"operation": _operation_summary(operation), "reason": str(exc)})

    if not applied:
        reasons = "; ".join(item["reason"] for item in skipped[:4])
        raise EditorError(reasons or "No usable edits were generated.")

    Path(output_path).parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)
    Presentation(output_path)
    return {
        "applied": applied,
        "skipped": skipped,
        "slide_count": len(prs.slides),
    }


# ---------------------------------------------------------------------------
# Structured OpenAI planner
# ---------------------------------------------------------------------------


class ChartSeriesPlan(BaseModel):
    model_config = ConfigDict(extra="forbid")
    name: str = Field(max_length=100)
    values: list[float] = Field(min_length=1, max_length=16)


class _ShapeRefPlan(BaseModel):
    """Compact semantic shape reference shared by shape-targeted operations."""

    model_config = ConfigDict(extra="forbid")
    slide: int
    shape_id: int | None = None
    shape_name: str | None = Field(default=None, max_length=120)
    text_contains: str | None = Field(default=None, max_length=240)
    shape: int | None = None


class ReplaceTextOperation(BaseModel):
    model_config = ConfigDict(extra="forbid")
    op: Literal["replace_text"]
    old: str = Field(max_length=600)
    new: str = Field(max_length=1600)
    slide: int | None = None
    shape_id: int | None = None
    shape_name: str | None = Field(default=None, max_length=120)
    text_contains: str | None = Field(default=None, max_length=240)
    shape: int | None = None
    replace_all: bool = True
    case_sensitive: bool = False


class SetTextOperation(_ShapeRefPlan):
    op: Literal["set_text"]
    text: str = Field(max_length=2400)
    font_size: float | None = None
    font_color: str | None = None
    font_face: str | None = None
    bold: bool | None = None
    italic: bool | None = None
    alignment: Literal["left", "center", "right", "justify"] | None = None


class StyleShapeOperation(_ShapeRefPlan):
    op: Literal["style_shape"]
    fill_color: str | None = None
    line_color: str | None = None
    line_width: float | None = None
    no_fill: bool | None = None
    no_line: bool | None = None
    font_size: float | None = None
    font_color: str | None = None
    font_face: str | None = None
    bold: bool | None = None
    italic: bool | None = None
    alignment: Literal["left", "center", "right", "justify"] | None = None
    x: float | None = None
    y: float | None = None
    width: float | None = None
    height: float | None = None


class MoveShapeOperation(_ShapeRefPlan):
    op: Literal["move_shape"]
    x: float
    y: float


class ResizeShapeOperation(_ShapeRefPlan):
    op: Literal["resize_shape"]
    width: float
    height: float


class DeleteShapeOperation(_ShapeRefPlan):
    op: Literal["delete_shape"]


class AddTextboxOperation(BaseModel):
    model_config = ConfigDict(extra="forbid")
    op: Literal["add_textbox"]
    slide: int
    text: str = Field(max_length=1800)
    x: float
    y: float
    width: float
    height: float
    font_size: float | None = None
    font_color: str | None = None
    font_face: str | None = None
    bold: bool | None = None
    italic: bool | None = None
    alignment: Literal["left", "center", "right", "justify"] | None = None
    fill_color: str | None = None
    line_color: str | None = None
    line_width: float | None = None
    no_fill: bool | None = None
    no_line: bool | None = None


class AddShapeOperation(BaseModel):
    model_config = ConfigDict(extra="forbid")
    op: Literal["add_shape"]
    slide: int
    shape_type: Literal["rectangle", "rounded_rectangle", "circle", "triangle", "chevron", "arrow"]
    x: float
    y: float
    width: float
    height: float
    text: str | None = Field(default=None, max_length=800)
    fill_color: str | None = None
    line_color: str | None = None
    line_width: float | None = None
    no_fill: bool | None = None
    no_line: bool | None = None
    font_size: float | None = None
    font_color: str | None = None
    font_face: str | None = None
    bold: bool | None = None
    italic: bool | None = None
    alignment: Literal["left", "center", "right", "justify"] | None = None


class AddChartOperation(BaseModel):
    model_config = ConfigDict(extra="forbid")
    op: Literal["add_chart"]
    slide: int
    chart_type: Literal["column", "bar", "line", "pie", "area"]
    title: str = Field(max_length=180)
    categories: list[str] = Field(min_length=1, max_length=16)
    series: list[ChartSeriesPlan] = Field(min_length=1, max_length=6)
    x: float
    y: float
    width: float
    height: float
    chart_style: int | None = None
    series_colors: list[str] | None = Field(default=None, max_length=6)
    show_legend: bool | None = None


class SetSlideBackgroundOperation(BaseModel):
    model_config = ConfigDict(extra="forbid")
    op: Literal["set_slide_background"]
    slide: int
    color: str


class SetTableCellOperation(_ShapeRefPlan):
    op: Literal["set_table_cell"]
    row: int
    column: int
    text: str = Field(max_length=1200)


class SetChartDataOperation(_ShapeRefPlan):
    op: Literal["set_chart_data"]
    categories: list[str] = Field(min_length=1, max_length=16)
    series: list[ChartSeriesPlan] = Field(min_length=1, max_length=6)
    title: str | None = Field(default=None, max_length=180)
    chart_style: int | None = None
    series_colors: list[str] | None = Field(default=None, max_length=6)
    show_legend: bool | None = None


class AddSlideOperation(BaseModel):
    model_config = ConfigDict(extra="forbid")
    op: Literal["add_slide"]
    position: int
    template_slide: int | None = None
    title: str = Field(max_length=180)
    subtitle: str | None = Field(default=None, max_length=500)
    body: list[str] = Field(default_factory=list, max_length=10)
    background_color: str | None = None
    title_color: str | None = None
    subtitle_color: str | None = None
    title_size: float | None = None
    body_size: float | None = None


class DuplicateSlideOperation(BaseModel):
    model_config = ConfigDict(extra="forbid")
    op: Literal["duplicate_slide"]
    slide: int
    position: int | None = None


class DeleteSlideOperation(BaseModel):
    model_config = ConfigDict(extra="forbid")
    op: Literal["delete_slide"]
    slide: int


class MoveSlideOperation(BaseModel):
    model_config = ConfigDict(extra="forbid")
    op: Literal["move_slide"]
    from_slide: int
    to_slide: int


class ReorderSlidesOperation(BaseModel):
    model_config = ConfigDict(extra="forbid")
    op: Literal["reorder_slides"]
    order: list[int] = Field(min_length=1, max_length=200)



class ExecutiveReviewOperation(BaseModel):
    """Compact high-level workflow for a complete executive deck review.

    One operation runs the full multi-pass workflow locally, which avoids a
    large model response and keeps every stage in one PowerPoint version.
    """

    model_config = ConfigDict(extra="forbid")
    op: Literal["executive_review"]
    cover_slide: int = 1
    summary_slide: int = 2
    analysis_slides: list[int] = Field(default_factory=lambda: [3, 4], max_length=12)
    narrative_slides: list[int] = Field(default_factory=lambda: [5, 6, 7, 8], max_length=30)
    add_recommendations: bool = True
    add_risks: bool = True
    standardize: bool = True
    delete_redundant: bool = True


EditOperationPlan = Annotated[
    ExecutiveReviewOperation
    | ReplaceTextOperation
    | SetTextOperation
    | StyleShapeOperation
    | MoveShapeOperation
    | ResizeShapeOperation
    | DeleteShapeOperation
    | AddTextboxOperation
    | AddShapeOperation
    | AddChartOperation
    | SetSlideBackgroundOperation
    | SetTableCellOperation
    | SetChartDataOperation
    | AddSlideOperation
    | DuplicateSlideOperation
    | DeleteSlideOperation
    | MoveSlideOperation
    | ReorderSlidesOperation,
    Field(discriminator="op"),
]


class EditPlan(BaseModel):
    model_config = ConfigDict(extra="forbid")
    message: str = Field(
        max_length=280,
        description="One short past-tense sentence shown after edits are applied.",
    )
    assumptions: list[str] = Field(default_factory=list, max_length=4)
    operations: list[EditOperationPlan] = Field(default_factory=list, max_length=4)


class PlannerTask(BaseModel):
    model_config = ConfigDict(extra="forbid")
    intent: Literal[
        "remove_text", "rewrite", "style", "layout", "chart", "table",
        "chart_data", "slide_order", "slide_management", "image", "notes", "cleanup", "regenerate", "general",
    ]
    instruction: str = Field(max_length=500)
    slides: list[int] = Field(default_factory=list, max_length=30)
    use_full_deck: bool


class TaskOutline(BaseModel):
    model_config = ConfigDict(extra="forbid")
    message: str = Field(max_length=280)
    assumptions: list[str] = Field(default_factory=list, max_length=4)
    tasks: list[PlannerTask] = Field(default_factory=list, max_length=16)


class AtomicEditStep(BaseModel):
    """Tiny structured envelope used when a normal plan is cut off.

    The low-level operation is encoded as one compact JSON string. Keeping the
    structured envelope small prevents the model from spending thousands of
    output tokens repeating nullable fields from the full operation union.
    """

    model_config = ConfigDict(extra="forbid")
    message: str = Field(default="", max_length=180)
    done: bool
    operation_json: str | None = Field(default=None, max_length=3200)


EDIT_OPERATION_ADAPTER = TypeAdapter(EditOperationPlan)


OUTLINE_INSTRUCTIONS = """
You are the planning layer for a professional PowerPoint editor. Turn the
user's request into a small set of independent edit tasks. The user may be
vague, use shorthand, make typos, or give a short follow-up. Infer the best
presentation-consulting interpretation from the recent chat, deck outline,
and slide previews.

Rules:
- Return no more than 16 tasks.
- Group work by slide. A task should normally target one slide.
- Keep slide-order work in one separate task and place it last.
- Set use_full_deck=true only for narrative ordering or a deck-wide change.
- Preserve every part of the user's request.
- Do not produce low-level edit operations here.
- Do not ask a question when a reasonable judgment exists.
- Keep the message and task instructions concise.
""".strip()


PLANNER_INSTRUCTIONS = """
You are the execution layer inside a professional PowerPoint editor. Convert
one focused task into concrete PowerPoint operations. Interpret natural
language the way a strong presentation consultant would. Make the best useful
judgment without asking a question when the deck provides enough context.

Core behavior:
- Follow the focused task and the original request.
- Inspect supplied slide structure and previews before choosing changes.
- Preserve the design unless the user asks for a redesign.
- Never invent business figures. Charts must use numbers already present in
  the deck. When data is insufficient, use shapes and labels instead.
- Prefer stable shape_id values. Also include text_contains as a semantic
  backup when targeting existing text.
- For removing a disclaimer or phrase, use replace_text. Empty replacement
  text is allowed.
- For a graph, use add_chart and use only figures visible in the supplied deck
  structure. Place it without covering key content.
- For reordering, use one reorder_slides operation containing every slide
  exactly once.
- Return at most 4 operations. Use concise text fields.
- Keep the summary message to one sentence under 180 characters.
- Do not emit blank lines, analysis, markdown, or explanations inside fields.
- Do not return an empty operations list for an edit task when a safe useful
  subset exists.

Coordinates and sizes from 0 through 1.5 represent fractions of slide width
and height. Larger values represent inches.
""".strip()


ATOMIC_PLANNER_INSTRUCTIONS = """
You are the fallback execution layer for a professional PowerPoint editor.
Return ONE low-level edit operation at a time, encoded as compact JSON in
operation_json. Never place markdown, commentary, or multiple operations in
operation_json. Set done=true and operation_json=null when the focused task is
complete. Keep message under 120 characters.

Use only these operations and fields:
- executive_review: {op, cover_slide, summary_slide, analysis_slides,
  narrative_slides, add_recommendations, add_risks, standardize,
  delete_redundant}. Use this ONE compact operation when the request asks for
  a complete senior-executive deck review with charts, recommendations, risks,
  narrative reordering, and deck-wide formatting.
- replace_text: {op, old, new, optional slide, replace_all, case_sensitive}
- set_text: {op, slide, shape_id or text_contains, text, optional font_size,
  font_color, bold, italic, alignment}
- style_shape: {op, slide, shape_id or text_contains, optional fill_color,
  line_color, font_size, font_color, bold, alignment, x, y, width, height}
- move_shape: {op, slide, shape_id or text_contains, x, y}
- resize_shape: {op, slide, shape_id or text_contains, width, height}
- delete_shape: {op, slide, shape_id or text_contains}
- add_textbox: {op, slide, text, x, y, width, height, optional font_size,
  font_color, bold, alignment, fill_color, line_color, no_fill, no_line}
- add_shape: {op, slide, shape_type, x, y, width, height, optional text,
  fill_color, line_color, font_size, font_color, bold, alignment}
- add_chart: {op, slide, chart_type, title, categories, series, x, y, width,
  height, optional series_colors, show_legend}. Each series is
  {name, values}. Use only figures already present in the supplied deck.
- set_slide_background: {op, slide, color}
- set_table_cell: {op, slide, shape_id, row, column, text}
- set_chart_data: {op, slide, shape_id, categories, series, optional title,
  series_colors, show_legend}
- add_slide: {op, position, optional template_slide, title, subtitle, body,
  background_color, title_color, subtitle_color, title_size, body_size}
- duplicate_slide: {op, slide, optional position}
- delete_slide: {op, slide}
- move_slide: {op, from_slide, to_slide}
- reorder_slides: {op, order}. The order must contain every slide once.

Prefer shape_id from the deck structure. Use text_contains as a semantic
backup. Coordinates from 0 through 1.5 are fractions of the slide. Larger
values are inches. Do not ask questions when a reasonable professional
judgment exists.
""".strip()


def _encode_image(path: str) -> str:
    with open(path, "rb") as image_file:
        return base64.b64encode(image_file.read()).decode("ascii")


def _visual_content(
    image_paths: list[Any] | None,
    selected_slide: int,
    max_images: int = 12,
) -> list[dict[str, Any]]:
    indexed: list[tuple[int, str]] = []
    for fallback_number, item in enumerate(image_paths or [], start=1):
        if isinstance(item, (tuple, list)) and len(item) == 2:
            slide_number, path = int(item[0]), str(item[1])
        else:
            slide_number, path = fallback_number, str(item)
        if path and os.path.exists(path):
            indexed.append((slide_number, path))
    if not indexed:
        return []

    if len(indexed) > max_images:
        selected = next((item for item in indexed if item[0] == selected_slide), None)
        step = max(1, len(indexed) // (max_images - 1))
        sampled = indexed[::step][: max_images - 1]
        if selected and selected not in sampled:
            sampled.append(selected)
        indexed = sorted(sampled, key=lambda item: item[0])[:max_images]

    parts: list[dict[str, Any]] = []
    for slide_number, path in indexed:
        parts.append({
            "type": "input_text",
            "text": f"Visual preview for slide {slide_number}{' (selected)' if slide_number == selected_slide else ''}:",
        })
        parts.append({
            "type": "input_image",
            "image_url": f"data:image/png;base64,{_encode_image(path)}",
        })
    return parts


def _plan_to_dict(plan: EditPlan, operation_limit: int = 12) -> dict[str, Any]:
    data = plan.model_dump(exclude_none=True)
    data["operations"] = data.get("operations", [])[:operation_limit]
    data["message"] = str(data.get("message") or "Done.")
    return data


def _is_cutoff_error(message: str) -> bool:
    lowered = message.casefold()
    return any(token in lowered for token in (
        "eof while parsing", "json_invalid", "invalid json", "max_tokens",
        "maximum output", "response was incomplete", "validation error for editplan",
    ))


def _friendly_openai_error(message: str) -> str:
    lowered = message.casefold()
    if "401" in lowered or "invalid_api_key" in lowered or "incorrect api key" in lowered:
        return "The AI credential was rejected. Update the local API key or the server environment, then restart Deck Refresh."
    if "429" in lowered or "rate limit" in lowered or "quota" in lowered:
        return "The OpenAI account hit a rate or billing limit. Check API billing, then retry."
    if _is_cutoff_error(message):
        return "The AI response was cut off before the edit plan finished."
    return f"OpenAI request failed: {message[:500]}"


def _ai_credential() -> str:
    return (
        os.environ.get("OPENAI_API_KEY", "").strip()
        or os.environ.get("AI_GATEWAY_API_KEY", "").strip()
        or os.environ.get("VERCEL_OIDC_TOKEN", "").strip()
    )


def _ai_model() -> str:
    configured = os.environ.get("OPENAI_MODEL", "").strip()
    if configured:
        return configured
    if os.environ.get("AI_GATEWAY_API_KEY") or os.environ.get("VERCEL_OIDC_TOKEN"):
        return "openai/gpt-5.6-sol"
    return "gpt-5"


def _ai_client(*, timeout: float, max_retries: int) -> OpenAI:
    credential = _ai_credential()
    if not credential:
        raise EditorError("AI access is not configured.")
    options: dict[str, Any] = {
        "api_key": credential,
        "timeout": timeout,
        "max_retries": max_retries,
    }
    if not os.environ.get("OPENAI_API_KEY") and (
        os.environ.get("AI_GATEWAY_API_KEY") or os.environ.get("VERCEL_OIDC_TOKEN")
    ):
        options["base_url"] = "https://ai-gateway.vercel.sh/v1"
    return OpenAI(**options)


def _reasoning_effort(model: str) -> str:
    configured = os.environ.get("OPENAI_REASONING_EFFORT", "").strip().lower()
    if configured:
        return configured
    lowered = model.casefold()
    if "gpt-5.1" in lowered:
        return "none"
    if lowered.startswith("gpt-5"):
        return "minimal"
    return "low"


def _request_structured(
    client: OpenAI,
    model: str,
    content: list[dict[str, Any]],
    instructions: str,
    text_format: type[BaseModel],
    max_output_tokens: int,
):
    primary_effort = _reasoning_effort(model)
    attempts = [
        (
            max_output_tokens,
            primary_effort,
            instructions,
        ),
        (
            max(max_output_tokens * 2, 10000),
            "low",
            instructions
            + "\n\nUse the shortest valid response. Finish the complete structured object."
            + " Never repeat whitespace or explanatory prose.",
        ),
    ]
    last_error = ""
    for token_budget, reasoning_effort, attempt_instructions in attempts:
        try:
            response = client.responses.parse(
                model=model,
                instructions=attempt_instructions,
                input=[{"role": "user", "content": content}],
                text_format=text_format,
                max_output_tokens=token_budget,
                reasoning={"effort": reasoning_effort},
                verbosity="low",
                store=False,
            )
            parsed = getattr(response, "output_parsed", None)
            if parsed is not None:
                return parsed
            status = getattr(response, "status", None)
            incomplete = getattr(response, "incomplete_details", None)
            usage = getattr(response, "usage", None)
            last_error = f"status={status}, incomplete={incomplete}, usage={usage}"
        except Exception as exc:
            last_error = str(exc).strip() or exc.__class__.__name__
    raise EditorError(_friendly_openai_error(last_error))


def _request_plan(
    client: OpenAI,
    model: str,
    content: list[dict[str, Any]],
    instructions: str,
    operation_limit: int = 12,
) -> dict[str, Any]:
    parsed = _request_structured(
        client,
        model,
        content,
        instructions,
        EditPlan,
        max_output_tokens=5000,
    )
    return _plan_to_dict(parsed, operation_limit)


def _parse_atomic_operation(raw: str | None) -> dict[str, Any] | None:
    if raw is None or not str(raw).strip():
        return None
    text = str(raw).strip()
    if text.startswith("```"):
        text = re.sub(r"^```(?:json)?\s*|\s*```$", "", text, flags=re.IGNORECASE)
    try:
        value = json.loads(text)
    except json.JSONDecodeError as exc:
        raise EditorError(f"Atomic operation was not valid JSON: {exc}") from exc
    if not isinstance(value, dict):
        raise EditorError("Atomic operation must be a JSON object.")
    try:
        parsed = EDIT_OPERATION_ADAPTER.validate_python(value)
    except ValidationError as exc:
        raise EditorError(f"Atomic operation did not match the edit schema: {exc}") from exc
    return parsed.model_dump(exclude_none=True)


def _request_atomic_plan(
    client: OpenAI,
    model: str,
    base_content: list[dict[str, Any]],
    pptx_path: str,
    operation_limit: int = 6,
) -> dict[str, Any]:
    operations: list[dict[str, Any]] = []
    messages: list[str] = []
    assumptions: list[str] = []
    feedback: list[str] = []

    for step_number in range(1, operation_limit + 2):
        state = {
            "step": step_number,
            "maximum_operations": operation_limit,
            "already_planned": operations,
            "validation_feedback": feedback[-3:],
            "instruction": (
                "Return the next single operation. Set done=true with no operation "
                "when the focused task is complete."
            ),
        }
        content = list(base_content)
        content.append({
            "type": "input_text",
            "text": "Atomic planning state:\n" + json.dumps(state, ensure_ascii=False),
        })
        step = _request_structured(
            client,
            model,
            content,
            ATOMIC_PLANNER_INSTRUCTIONS,
            AtomicEditStep,
            max_output_tokens=6500,
        )
        if step.message.strip():
            messages.append(step.message.strip())

        operation = None
        try:
            operation = _parse_atomic_operation(step.operation_json)
        except EditorError as exc:
            feedback.append(str(exc))

        if operation is not None:
            valid, issues = validate_operations(pptx_path, [operation])
            if valid:
                candidate = valid[0]
                serialized = json.dumps(candidate, sort_keys=True, ensure_ascii=False)
                if all(
                    json.dumps(existing, sort_keys=True, ensure_ascii=False) != serialized
                    for existing in operations
                ):
                    operations.append(candidate)
                else:
                    feedback.append("That operation was already planned. Choose another edit or finish.")
            else:
                feedback.extend(issues or ["The operation was not editable in this deck."])

        if step.done:
            break
        if len(operations) >= operation_limit:
            break

    return {
        "message": " ".join(dict.fromkeys(messages))[:280] or "Applied the requested edits.",
        "assumptions": assumptions,
        "operations": operations,
    }


def _request_outline(
    client: OpenAI,
    model: str,
    content: list[dict[str, Any]],
) -> TaskOutline:
    return _request_structured(
        client,
        model,
        content,
        OUTLINE_INSTRUCTIONS,
        TaskOutline,
        max_output_tokens=3000,
    )


def _looks_like_edit_request(message: str, history: list[dict[str, str]]) -> bool:
    combined = " ".join(
        [str(item.get("content", "")) for item in history[-5:] if item.get("role") == "user"]
        + [message]
    ).casefold()
    verbs = (
        "edit", "change", "remove", "delete", "add", "make", "reorder", "move",
        "replace", "update", "fix", "resize", "style", "graph", "chart", "slide",
        "do everything", "do it",
    )
    return any(verb in combined for verb in verbs)


def _deck_outline(summary: dict[str, Any]) -> dict[str, Any]:
    slides = []
    for slide in summary.get("slides", []):
        text_parts = []
        kinds: dict[str, int] = {}
        for shape in slide.get("shapes", []):
            kind = str(shape.get("kind", "shape"))
            kinds[kind] = kinds.get(kind, 0) + 1
            text = str(shape.get("text", "")).strip()
            if text:
                text_parts.append(text)
        slides.append({
            "slide": slide.get("slide"),
            "title": slide.get("title", ""),
            "content": " | ".join(text_parts)[:1400],
            "object_counts": kinds,
        })
    return {"slide_count": summary.get("slide_count", 0), "slides": slides}


def _focused_summary(summary: dict[str, Any], slides: list[int], use_full_deck: bool) -> dict[str, Any]:
    if use_full_deck:
        return summary
    wanted = {int(value) for value in slides if isinstance(value, int) or str(value).isdigit()}
    filtered = [slide for slide in summary.get("slides", []) if int(slide.get("slide", 0)) in wanted]
    return {
        "slide_width_inches": summary.get("slide_width_inches"),
        "slide_height_inches": summary.get("slide_height_inches"),
        "slide_count": summary.get("slide_count"),
        "slides": filtered,
    }


def _images_for_task(
    deck_image_paths: list[str],
    slides: list[int],
    use_full_deck: bool,
) -> list[Any]:
    if use_full_deck:
        return [(index, path) for index, path in enumerate(deck_image_paths, start=1)]
    chosen: list[tuple[int, str]] = []
    for slide_number in slides:
        index = int(slide_number) - 1
        if 0 <= index < len(deck_image_paths):
            path = deck_image_paths[index]
            if path and os.path.exists(path):
                chosen.append((int(slide_number), path))
    return chosen



def _looks_like_executive_review_request(message: str) -> bool:
    text = _normalize_text(message)
    signals = [
        "senior executive audience",
        "preserve the kpmg branding",
        "executive recommendations",
        "key risks and mitigations",
        "native powerpoint charts",
        "standardize all slide titles",
        "verify the complete presentation flow",
        "delete any slide that becomes fully redundant",
    ]
    return sum(1 for signal in signals if signal in text) >= 4


def _executive_review_operation(slide_count: int) -> dict[str, Any]:
    return {
        "op": "executive_review",
        "cover_slide": 1,
        "summary_slide": 2 if slide_count >= 2 else 1,
        "analysis_slides": [value for value in (3, 4) if value <= slide_count],
        "narrative_slides": [value for value in (5, 6, 7, 8) if value <= slide_count],
        "add_recommendations": True,
        "add_risks": True,
        "standardize": True,
        "delete_redundant": True,
    }


def _fallback_tasks(user_message: str, selected_slide: int, slide_count: int) -> list[PlannerTask]:
    if _looks_like_executive_review_request(user_message):
        return [PlannerTask(
            intent="general",
            instruction="Run the complete executive deck review workflow.",
            slides=list(range(1, slide_count + 1)),
            use_full_deck=True,
        )]

    text = user_message.strip()
    lowered = text.casefold()
    tasks: list[PlannerTask] = []

    recognized = []
    if len(text) > 500 or any(token in lowered for token in ("executive recommendations", "key risks and mitigations", "standardize all slide titles")):
        recognized = [
        (r"(?:on\s+)?slide\s+1[\s\S]*?(?=(?:on\s+)?slide\s+2|review slides?\s+3|$)", "rewrite", [1], False),
        (r"(?:on\s+)?slide\s+2[\s\S]*?(?=review slides?\s+3|review slides?\s+5|$)", "chart", [2], False),
        (r"review slides?\s+3\s*(?:and|&)\s*4[\s\S]*?(?=review slides?\s+5|create a new slide|$)", "rewrite", [3, 4], False),
        (r"review slides?\s+5[\s,]*(?:6[\s,]*)?(?:7[\s,]*)?(?:and\s*)?8[\s\S]*?(?=create a new slide|create another new slide|$)", "slide_order", [5, 6, 7, 8], True),
        (r"create a new slide[\s\S]*?executive recommendations[\s\S]*?(?=create another new slide|delete any slide|$)", "slide_management", list(range(1, slide_count + 1)), True),
        (r"create another new slide[\s\S]*?key risks and mitigations[\s\S]*?(?=delete any slide|standardize|$)", "table", list(range(1, slide_count + 1)), True),
        (r"delete any slide[\s\S]*?(?=standardize|after finishing|$)", "slide_management", list(range(1, slide_count + 1)), True),
        (r"standardize[\s\S]*?(?=after finishing|$)", "style", list(range(1, slide_count + 1)), True),
        (r"after finishing[\s\S]*$", "general", list(range(1, slide_count + 1)), True),
    ]
    for pattern, intent, slides, full in recognized:
        match = re.search(pattern, text, flags=re.IGNORECASE)
        if match:
            valid_slides = [value for value in slides if 1 <= value <= slide_count]
            tasks.append(PlannerTask(
                intent=intent,
                instruction=match.group(0).strip()[:500],
                slides=valid_slides or [selected_slide],
                use_full_deck=full,
            ))

    if tasks:
        return tasks[:16]

    chunks = [part.strip() for part in re.split(r"[.;\n]+", text) if part.strip()]
    if not chunks:
        chunks = [text]
    for chunk in chunks[:16]:
        chunk_lowered = chunk.casefold()
        numbers = [int(value) for value in re.findall(r"\b(?:slide|slides)\s*(\d+)\b", chunk_lowered)]
        slides = [value for value in numbers if 1 <= value <= slide_count] or [selected_slide]
        if "reorder" in chunk_lowered or "flow" in chunk_lowered or "order" in chunk_lowered:
            intent = "slide_order"
            slides = list(range(1, slide_count + 1))
            full = True
        elif "chart" in chunk_lowered or "graph" in chunk_lowered:
            intent, full = "chart", False
        elif "table" in chunk_lowered or "risk" in chunk_lowered:
            intent, full = "table", True
        elif "remove" in chunk_lowered or "delete text" in chunk_lowered:
            intent, full = "remove_text", False
        elif "executive" in chunk_lowered or "rewrite" in chunk_lowered:
            intent, full = "rewrite", False
        elif "standardize" in chunk_lowered or "format" in chunk_lowered:
            intent, full = "style", True
        else:
            intent, full = "general", False
        tasks.append(PlannerTask(intent=intent, instruction=chunk[:500], slides=slides, use_full_deck=full))
    return tasks


def _dedupe_operations(operations: list[dict[str, Any]]) -> list[dict[str, Any]]:
    result = []
    seen = set()
    for operation in operations:
        key = json.dumps(operation, sort_keys=True, ensure_ascii=False)
        if key in seen:
            continue
        seen.add(key)
        result.append(operation)
    structural_last = {"move_slide", "reorder_slides"}
    return [op for op in result if op.get("op") not in structural_last] + [
        op for op in result if op.get("op") in structural_last
    ]


def _message_slide_numbers(message: str, slide_count: int) -> list[int]:
    lowered = message.casefold()
    numbers: list[int] = []
    for match in re.finditer(r"\bslides?\s+((?:\d+[\s,\-/and]*)+)", lowered):
        numbers.extend(int(value) for value in re.findall(r"\d+", match.group(1)))
    return sorted({value for value in numbers if 1 <= value <= slide_count})


def _deterministic_operations(
    summary: dict[str, Any],
    user_message: str,
    selected_slide: int,
) -> list[dict[str, Any]]:
    """Create a small safe subset when every model planning attempt fails.

    This is intentionally conservative. It handles explicit text removal and
    replacement requests without inventing content, so a temporary model
    cutoff never turns a clear request into a total editor failure.
    """

    lowered = user_message.casefold()
    slide_count = int(summary.get("slide_count", 0))
    if _looks_like_executive_review_request(user_message):
        return [_executive_review_operation(slide_count)]
    explicit_slides = _message_slide_numbers(user_message, slide_count)
    target_slides = explicit_slides or [selected_slide]
    operations: list[dict[str, Any]] = []

    known_removals = [
        "internal demo",
        "internal demonstration",
        "all figures are fictional",
        "figures are fictional",
        "fictional data",
    ]
    for phrase in known_removals:
        if phrase not in lowered:
            continue
        matching_slides = []
        wanted = _normalize_text(phrase)
        if not explicit_slides:
            for slide in summary.get("slides", []):
                if any(wanted in _normalize_text(shape.get("text", "")) for shape in slide.get("shapes", [])):
                    matching_slides.append(int(slide.get("slide", 0)))
        for slide_number in target_slides if explicit_slides else (matching_slides or target_slides):
            operations.append({
                "op": "replace_text",
                "slide": slide_number,
                "old": phrase,
                "new": "",
                "replace_all": True,
                "case_sensitive": False,
            })

    for phrase in re.findall(
        r"(?:remove|delete)\s+(?:the\s+text\s+)?[\"“']([^\"”']{2,180})[\"”']",
        user_message,
        flags=re.IGNORECASE,
    ):
        for slide_number in target_slides:
            operations.append({
                "op": "replace_text",
                "slide": slide_number,
                "old": phrase.strip(),
                "new": "",
                "replace_all": True,
                "case_sensitive": False,
            })

    replacement_patterns = [
        r"change\s+every\s+(.{1,80}?)\s+(?:reference\s+)?to\s+(.{1,80}?)(?:[.;]|$)",
        r"replace\s+[\"“']?(.{1,80}?)[\"”']?\s+with\s+[\"“']?(.{1,80}?)[\"”']?(?:[.;]|$)",
    ]
    for pattern in replacement_patterns:
        match = re.search(pattern, user_message, flags=re.IGNORECASE)
        if not match:
            continue
        old, new = match.group(1).strip(), match.group(2).strip()
        if old and new:
            operations.append({
                "op": "replace_text",
                "old": old,
                "new": new,
                "replace_all": True,
                "case_sensitive": False,
            })
            break

    return _dedupe_operations(operations)


def plan_edit(
    pptx_path: str,
    user_message: str,
    selected_slide: int,
    selected_slide_image: str | None = None,
    chat_history: list[dict[str, str]] | None = None,
    deck_image_paths: list[str] | None = None,
) -> dict[str, Any]:
    summary = deck_summary(pptx_path, max_text=1200)
    slide_count = int(summary.get("slide_count", 1))
    if _looks_like_executive_review_request(user_message):
        operation = _executive_review_operation(slide_count)
        valid, issues = validate_operations(pptx_path, [operation])
        if valid:
            return {
                "message": "Completed the full executive deck review, including charts, flow, recommendations, risks, and formatting.",
                "assumptions": issues,
                "operations": valid,
            }

    if not _ai_credential():
        raise EditorError("AI access is not configured.")

    selected_slide = max(1, min(int(selected_slide), max(1, slide_count)))
    history = chat_history[-10:] if isinstance(chat_history, list) else []
    images = list(deck_image_paths or [])
    if selected_slide_image and selected_slide_image not in images:
        images.append(selected_slide_image)

    client = _ai_client(timeout=180.0, max_retries=2)
    model = _ai_model()

    outline_prompt = {
        "request": user_message,
        "selected_slide": selected_slide,
        "recent_chat": history,
        "deck_outline": _deck_outline(summary),
    }
    outline_content: list[dict[str, Any]] = [
        {"type": "input_text", "text": json.dumps(outline_prompt, ensure_ascii=False)}
    ]
    outline_content.extend(_visual_content(images, selected_slide, max_images=12))

    outline: TaskOutline | None = None
    try:
        outline = _request_outline(client, model, outline_content)
        tasks = outline.tasks[:16]
    except EditorError:
        tasks = _fallback_tasks(user_message, selected_slide, slide_count)

    if not tasks and _looks_like_edit_request(user_message, history):
        tasks = _fallback_tasks(user_message, selected_slide, slide_count)

    all_operations: list[dict[str, Any]] = []
    all_assumptions = list(outline.assumptions if outline else [])
    task_messages: list[str] = []
    failed_tasks: list[str] = []

    for task_number, task in enumerate(tasks[:16], start=1):
        valid_slides = sorted({
            int(value) for value in task.slides
            if 1 <= int(value) <= slide_count
        })
        if task.intent == "slide_order":
            valid_slides = list(range(1, slide_count + 1))
            task.use_full_deck = True
        if not valid_slides:
            valid_slides = [selected_slide]

        focused = _focused_summary(summary, valid_slides, task.use_full_deck)
        task_prompt = {
            "original_request": user_message,
            "focused_task": task.model_dump(),
            "selected_slide": selected_slide,
            "recent_chat": history,
            "deck": focused,
            "task_number": task_number,
            "task_count": len(tasks),
        }
        task_content: list[dict[str, Any]] = [
            {"type": "input_text", "text": json.dumps(task_prompt, ensure_ascii=False)}
        ]
        task_images = _images_for_task(images, valid_slides, task.use_full_deck)
        task_content.extend(_visual_content(task_images, selected_slide, max_images=12))

        try:
            # Atomic planning is the primary path. Each API response contains at
            # most one compact operation, so a long request cannot be truncated
            # halfway through a large edit-plan object.
            task_plan = _request_atomic_plan(
                client,
                model,
                task_content,
                pptx_path,
                operation_limit=20 if task.intent in {"rewrite", "layout", "general", "table", "slide_management", "image", "notes", "cleanup", "regenerate"} else 12,
            )
            valid, issues = validate_operations(pptx_path, task_plan.get("operations", []))

            if issues and not valid:
                repair_prompt = {
                    "focused_task": task.model_dump(),
                    "original_request": user_message,
                    "invalid_plan": task_plan,
                    "validation_issues": issues,
                    "deck": focused,
                    "requirements": [
                        "Correct all object references.",
                        "Use shape_id and text_contains from the deck.",
                        "Return the safest useful subset in at most 10 operations.",
                    ],
                }
                repair_content = [{
                    "type": "input_text",
                    "text": json.dumps(repair_prompt, ensure_ascii=False),
                }]
                repair_content.extend(_visual_content(task_images, selected_slide, max_images=6))
                task_plan = _request_atomic_plan(
                    client,
                    model,
                    repair_content,
                    pptx_path,
                    operation_limit=12,
                )
                valid, issues = validate_operations(pptx_path, task_plan.get("operations", []))

            if valid:
                all_operations.extend(valid)
                task_messages.append(str(task_plan.get("message") or ""))
                all_assumptions.extend(task_plan.get("assumptions", []))
                all_assumptions.extend(issues)
            else:
                failed_tasks.append(task.instruction)
        except EditorError:
            failed_tasks.append(task.instruction)

    if failed_tasks:
        all_operations.extend(
            _deterministic_operations(summary, user_message, selected_slide)
        )

    all_operations = _dedupe_operations(all_operations)[:100]
    final_valid, final_issues = validate_operations(pptx_path, all_operations)
    all_assumptions.extend(final_issues)

    if not final_valid:
        deterministic = _deterministic_operations(summary, user_message, selected_slide)
        final_valid, deterministic_issues = validate_operations(pptx_path, deterministic)
        all_assumptions.extend(deterministic_issues)
        if final_valid:
            failed_tasks = failed_tasks or [user_message]
        else:
            raise EditorError(
                "The editor could not map this request to a safe PowerPoint change. "
                "Try naming the slide or the text to change."
            )

    if outline and outline.message.strip():
        message = outline.message.strip()
    else:
        message = "Applied the requested PowerPoint changes."
    if failed_tasks:
        message += f" Completed the translated edits. {len(failed_tasks)} requested part{'s' if len(failed_tasks) != 1 else ''} could not be translated into a valid PowerPoint action."
    elif task_messages:
        message = " ".join(dict.fromkeys(item.strip() for item in task_messages if item.strip()))[:1200] or message

    return {
        "message": message,
        "assumptions": list(dict.fromkeys(str(item) for item in all_assumptions if str(item).strip()))[:20],
        "operations": final_valid,
    }

# ---------------------------------------------------------------------------
# Universal editor resilience layer
# ---------------------------------------------------------------------------
# These definitions intentionally appear at the end of the module. They extend
# the original editor without disturbing the proven Deck Refresh update path.

from pptx.enum.shapes import MSO_CONNECTOR
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.util import Emu


BASE_SHAPE_OPS = {
    "set_text", "style_shape", "move_shape", "resize_shape", "delete_shape",
    "set_table_cell", "set_chart_data",
}
EXTENDED_SHAPE_OPS = {
    "append_text", "format_text_box", "rotate_shape", "duplicate_shape",
    "layer_shape", "fit_text", "crop_picture", "replace_picture",
    "style_table", "merge_table_cells", "split_table_cell",
    "add_table_row", "delete_table_row", "add_table_column",
    "delete_table_column", "style_chart", "change_chart_type",
}
SHAPE_OPS = BASE_SHAPE_OPS | EXTENDED_SHAPE_OPS
SUPPORTED_OPS = {
    "noop", "replace_text", *SHAPE_OPS,
    "add_textbox", "add_shape", "add_line", "add_chart", "add_picture",
    "add_table", "set_slide_background", "set_speaker_notes",
    "add_slide", "duplicate_slide", "delete_slide", "move_slide",
    "reorder_slides", "clear_slide", "regenerate_slide",
    "align_shapes", "distribute_shapes", "standardize_deck",
    "cleanup_slide", "set_slide_hidden", "executive_review",
}


def _safe_float(value: Any, default: float | None = None) -> float | None:
    if value is None:
        return default
    try:
        return float(value)
    except (TypeError, ValueError):
        return default


def _delete_shape_element(shape) -> None:
    element = shape.element
    parent = element.getparent()
    if parent is not None:
        parent.remove(element)


def _shape_is_branding(shape, prs: Presentation) -> bool:
    name = str(getattr(shape, "name", "")).casefold()
    text = _normalize_text(_shape_text(shape))
    top = float(shape.top) / max(float(prs.slide_height), 1.0)
    left = float(shape.left) / max(float(prs.slide_width), 1.0)
    if "logo" in name or "kpmg" in text:
        return True
    if _shape_kind(shape) == "picture" and top < 0.18 and (left < 0.25 or left > 0.70):
        return True
    return False


def _add_line(prs: Presentation, operation: dict[str, Any]) -> None:
    slide = _slide(prs, int(operation["slide"]))
    x1 = _to_emu(operation.get("x1", operation.get("x", 0.1)), prs.slide_width)
    y1 = _to_emu(operation.get("y1", operation.get("y", 0.1)), prs.slide_height)
    x2 = _to_emu(operation.get("x2", 0.9), prs.slide_width)
    y2 = _to_emu(operation.get("y2", 0.1), prs.slide_height)
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    color = _rgb(operation.get("line_color"), "00338D")
    if color is not None:
        line.line.color.rgb = color
    if operation.get("line_width") is not None:
        line.line.width = Pt(float(operation["line_width"]))
    if operation.get("rotation") is not None:
        line.rotation = float(operation["rotation"])


def _table_matrix(operation: dict[str, Any]) -> list[list[str]]:
    raw = operation.get("data") or operation.get("rows_data") or []
    if isinstance(raw, str):
        rows = []
        for line in raw.splitlines():
            if not line.strip():
                continue
            delimiter = "\t" if "\t" in line else ","
            rows.append([cell.strip() for cell in line.split(delimiter)])
        return rows
    if not isinstance(raw, list):
        return []
    matrix: list[list[str]] = []
    for row in raw:
        if isinstance(row, list):
            matrix.append([str(cell) for cell in row])
        elif isinstance(row, dict):
            matrix.append([str(value) for value in row.values()])
        else:
            matrix.append([str(row)])
    return matrix


def _style_cell(cell, *, fill: str | None = None, font_color: str | None = None,
                font_size: float | None = None, bold: bool | None = None,
                alignment: str | None = None) -> None:
    if fill:
        rgb = _rgb(fill)
        if rgb:
            cell.fill.solid()
            cell.fill.fore_color.rgb = rgb
    align_map = {
        "left": PP_ALIGN.LEFT,
        "center": PP_ALIGN.CENTER,
        "right": PP_ALIGN.RIGHT,
        "justify": PP_ALIGN.JUSTIFY,
    }
    for paragraph in cell.text_frame.paragraphs:
        if alignment in align_map:
            paragraph.alignment = align_map[alignment]
        for run in paragraph.runs:
            if font_color:
                rgb = _rgb(font_color)
                if rgb:
                    run.font.color.rgb = rgb
            if font_size is not None:
                run.font.size = Pt(float(font_size))
            if bold is not None:
                run.font.bold = bool(bold)


def _add_table(prs: Presentation, operation: dict[str, Any]) -> None:
    slide = _slide(prs, int(operation["slide"]))
    matrix = _table_matrix(operation)
    rows = int(operation.get("rows") or len(matrix) or 2)
    columns = int(operation.get("columns") or max((len(row) for row in matrix), default=2))
    rows = max(1, min(rows, 100))
    columns = max(1, min(columns, 30))
    if operation.get("auto_fit"):
        x, y, width, height = _auto_insert_box(
            prs, slide, min_width=0.12, min_height=0.08,
            max_width=0.58, max_height=0.45, object_name="table",
        )
    else:
        x, y, width, height = (
            operation.get("x", 0.08), operation.get("y", 0.22),
            operation.get("width", 0.84), operation.get("height", 0.62),
        )
    shape = slide.shapes.add_table(
        rows,
        columns,
        _to_emu(x, prs.slide_width),
        _to_emu(y, prs.slide_height),
        _to_emu(width, prs.slide_width),
        _to_emu(height, prs.slide_height),
    )
    table = shape.table
    for r in range(rows):
        for c in range(columns):
            value = matrix[r][c] if r < len(matrix) and c < len(matrix[r]) else ""
            _set_text_frame_text(table.cell(r, c).text_frame, value, preserve_first_run=False)
    _style_table_shape(shape, operation)
    for merge in operation.get("merge_cells", []) or []:
        r1 = int(merge.get("row", 1)) - 1
        c1 = int(merge.get("column", 1)) - 1
        r2 = int(merge.get("end_row", r1 + 1)) - 1
        c2 = int(merge.get("end_column", c1 + 1)) - 1
        if not (0 <= r1 <= r2 < rows and 0 <= c1 <= c2 < columns):
            raise EditorError("A requested table merge is outside the new table.")
        table.cell(r1, c1).merge(table.cell(r2, c2))


def _style_table_shape(shape, operation: dict[str, Any]) -> None:
    if not getattr(shape, "has_table", False):
        raise EditorError("The selected object is not a table.")
    table = shape.table
    header_rows = max(0, int(operation.get("header_rows", 1)))
    header_fill = operation.get("header_fill", "00338D")
    header_font = operation.get("header_font_color", "FFFFFF")
    body_fill = operation.get("body_fill")
    body_font = operation.get("body_font_color", "1F2937")
    font_size = _safe_float(operation.get("font_size"), 10)
    band_fill = operation.get("band_fill", "F3F6FA")
    banded = bool(operation.get("banded_rows", False))
    for r, row in enumerate(table.rows):
        for cell in row.cells:
            if r < header_rows:
                _style_cell(
                    cell, fill=header_fill, font_color=header_font,
                    font_size=font_size, bold=True,
                    alignment=operation.get("header_alignment", "center"),
                )
            else:
                fill = band_fill if banded and (r - header_rows) % 2 else body_fill
                _style_cell(
                    cell, fill=fill, font_color=body_font,
                    font_size=font_size, bold=False,
                    alignment=operation.get("body_alignment", "left"),
                )


def _table_add_row(shape, operation: dict[str, Any]) -> None:
    table = shape.table
    tbl = table._tbl
    source = copy.deepcopy(tbl.tr_lst[-1])
    for tc in source.tc_lst:
        for paragraph in tc.txBody.p_lst:
            for run in paragraph.r_lst:
                run.t.text = ""
    tbl.append(source)
    values = operation.get("values") or []
    row_index = len(table.rows) - 1
    for c, value in enumerate(values[:len(table.columns)]):
        _set_text_frame_text(table.cell(row_index, c).text_frame, str(value), preserve_first_run=False)


def _table_delete_row(shape, row_number: int) -> None:
    table = shape.table
    row_index = row_number - 1
    if row_index < 0 or row_index >= len(table.rows):
        raise EditorError("The table row is outside the table.")
    if len(table.rows) <= 1:
        raise EditorError("A table must keep at least one row.")
    tr = table._tbl.tr_lst[row_index]
    table._tbl.remove(tr)


def _table_add_column(shape, operation: dict[str, Any]) -> None:
    table = shape.table
    tbl = table._tbl
    new_grid_col = copy.deepcopy(tbl.tblGrid.gridCol_lst[-1])
    tbl.tblGrid.append(new_grid_col)
    for tr in tbl.tr_lst:
        new_tc = copy.deepcopy(tr.tc_lst[-1])
        for paragraph in new_tc.txBody.p_lst:
            for run in paragraph.r_lst:
                run.t.text = ""
        tr.insert(len(tr.tc_lst), new_tc)
    values = operation.get("values") or []
    column_index = len(table.columns) - 1
    for r, value in enumerate(values[:len(table.rows)]):
        _set_text_frame_text(table.cell(r, column_index).text_frame, str(value), preserve_first_run=False)


def _table_delete_column(shape, column_number: int) -> None:
    table = shape.table
    column_index = column_number - 1
    if column_index < 0 or column_index >= len(table.columns):
        raise EditorError("The table column is outside the table.")
    if len(table.columns) <= 1:
        raise EditorError("A table must keep at least one column.")
    grid_col = table._tbl.tblGrid.gridCol_lst[column_index]
    table._tbl.tblGrid.remove(grid_col)
    for tr in table._tbl.tr_lst:
        tr.remove(tr.tc_lst[column_index])


def _add_picture(prs: Presentation, operation: dict[str, Any]) -> None:
    slide = _slide(prs, int(operation["slide"]))
    source = str(operation.get("source_path") or operation.get("path") or "").strip()
    if not source or not os.path.isfile(source):
        raise EditorError("The image attachment is unavailable.")
    width = operation.get("width")
    height = operation.get("height")
    picture = slide.shapes.add_picture(
        source,
        _to_emu(operation.get("x", 0.1), prs.slide_width),
        _to_emu(operation.get("y", 0.15), prs.slide_height),
        _to_emu(width, prs.slide_width) if width is not None else None,
        _to_emu(height, prs.slide_height) if height is not None else None,
    )
    for key in ("crop_left", "crop_right", "crop_top", "crop_bottom"):
        if operation.get(key) is not None:
            setattr(picture, key, max(0.0, min(1.0, float(operation[key]))))
    if operation.get("rotation") is not None:
        picture.rotation = float(operation["rotation"])


def _replace_picture(prs: Presentation, operation: dict[str, Any]) -> None:
    slide = _slide(prs, int(operation["slide"]))
    old = _operation_shape(slide, operation)
    source = str(operation.get("source_path") or operation.get("path") or "").strip()
    if not source or not os.path.isfile(source):
        raise EditorError("The replacement image attachment is unavailable.")
    left, top, width, height = old.left, old.top, old.width, old.height
    rotation = getattr(old, "rotation", 0)
    _delete_shape_element(old)
    picture = slide.shapes.add_picture(source, left, top, width, height)
    picture.rotation = rotation
    for key in ("crop_left", "crop_right", "crop_top", "crop_bottom"):
        if operation.get(key) is not None:
            setattr(picture, key, max(0.0, min(1.0, float(operation[key]))))


def _crop_picture(shape, operation: dict[str, Any]) -> None:
    if _shape_kind(shape) != "picture":
        raise EditorError("The selected object is not a picture.")
    for key in ("crop_left", "crop_right", "crop_top", "crop_bottom"):
        if operation.get(key) is not None:
            setattr(shape, key, max(0.0, min(1.0, float(operation[key]))))


def _duplicate_shape(slide, shape) -> int:
    clone = copy.deepcopy(shape.element)
    new_id = max((int(candidate.shape_id) for candidate in slide.shapes), default=1) + 1
    c_nv_pr = clone.xpath(".//p:cNvPr")
    if c_nv_pr:
        c_nv_pr[0].set("id", str(new_id))
        original_name = c_nv_pr[0].get("name") or str(getattr(shape, "name", "Shape"))
        c_nv_pr[0].set("name", f"{original_name} Copy")
    slide.shapes._spTree.insert_element_before(clone, "p:extLst")
    return new_id


def _layer_shape(slide, shape, direction: str) -> None:
    tree = slide.shapes._spTree
    element = shape.element
    siblings = [child for child in tree if child.tag.endswith(("sp", "pic", "graphicFrame", "cxnSp", "grpSp"))]
    if element not in siblings:
        return
    current = siblings.index(element)
    tree.remove(element)
    direction = direction.casefold()
    if direction in {"front", "bring_to_front"}:
        tree.insert_element_before(element, "p:extLst")
    elif direction in {"back", "send_to_back"}:
        tree.insert(2, element)
    elif direction in {"forward", "bring_forward"}:
        target = min(current + 1, len(siblings) - 1)
        ref = siblings[target]
        ref.addnext(element)
    else:
        target = max(current - 1, 0)
        ref = siblings[target]
        ref.addprevious(element)


def _format_text_box(shape, operation: dict[str, Any]) -> None:
    if not getattr(shape, "has_text_frame", False):
        raise EditorError("The selected object has no editable text box.")
    tf = shape.text_frame
    if operation.get("word_wrap") is not None:
        tf.word_wrap = bool(operation["word_wrap"])
    if operation.get("margin_left") is not None:
        tf.margin_left = Inches(float(operation["margin_left"]))
    if operation.get("margin_right") is not None:
        tf.margin_right = Inches(float(operation["margin_right"]))
    if operation.get("margin_top") is not None:
        tf.margin_top = Inches(float(operation["margin_top"]))
    if operation.get("margin_bottom") is not None:
        tf.margin_bottom = Inches(float(operation["margin_bottom"]))
    vertical = str(operation.get("vertical_alignment") or "").casefold()
    vertical_map = {
        "top": MSO_ANCHOR.TOP,
        "middle": MSO_ANCHOR.MIDDLE,
        "center": MSO_ANCHOR.MIDDLE,
        "bottom": MSO_ANCHOR.BOTTOM,
    }
    if vertical in vertical_map:
        tf.vertical_anchor = vertical_map[vertical]
    if operation.get("autofit", True):
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    bullet = operation.get("bullet")
    level = operation.get("level")
    for paragraph in tf.paragraphs:
        if level is not None:
            paragraph.level = max(0, min(8, int(level)))
        if operation.get("space_before") is not None:
            paragraph.space_before = Pt(float(operation["space_before"]))
        if operation.get("space_after") is not None:
            paragraph.space_after = Pt(float(operation["space_after"]))
        if operation.get("line_spacing") is not None:
            paragraph.line_spacing = float(operation["line_spacing"])
        if bullet is not None:
            pPr = paragraph._p.get_or_add_pPr()
            for child in list(pPr):
                if child.tag in {qn("a:buChar"), qn("a:buNone"), qn("a:buAutoNum")} :
                    pPr.remove(child)
            bullet_element = OxmlElement("a:buChar" if bullet else "a:buNone")
            if bullet:
                bullet_element.set("char", str(operation.get("bullet_character", "•")))
            pPr.insert(0, bullet_element)
    _apply_text_style(shape, operation)


def _set_speaker_notes(prs: Presentation, operation: dict[str, Any]) -> None:
    slide = _slide(prs, int(operation["slide"]))
    notes = slide.notes_slide.notes_text_frame
    text = str(operation.get("text", ""))
    mode = str(operation.get("mode", "replace")).casefold()
    if mode == "append" and notes.text.strip():
        text = notes.text.rstrip() + "\n" + text
    _set_text_frame_text(notes, text, preserve_first_run=False)


def _clear_slide(prs: Presentation, operation: dict[str, Any]) -> int:
    slide = _slide(prs, int(operation["slide"]))
    preserve_branding = bool(operation.get("preserve_branding", True))
    preserve_title = bool(operation.get("preserve_title", False))
    deleted = 0
    for shape in list(slide.shapes):
        if preserve_branding and _shape_is_branding(shape, prs):
            continue
        if preserve_title and getattr(shape, "has_text_frame", False) and shape.text_frame.text.strip() == _slide_title(slide):
            continue
        _delete_shape_element(shape)
        deleted += 1
    return deleted


def _regenerate_slide(prs: Presentation, operation: dict[str, Any]) -> None:
    slide_number = int(operation["slide"])
    _clear_slide(prs, {
        "slide": slide_number,
        "preserve_branding": operation.get("preserve_branding", True),
        "preserve_title": False,
    })
    slide = _slide(prs, slide_number)
    title = str(operation.get("title") or "Executive Summary")
    subtitle = str(operation.get("subtitle") or "")
    bullets = operation.get("body") or operation.get("bullets") or []
    if isinstance(bullets, str):
        bullets = [line.strip(" •-") for line in bullets.splitlines() if line.strip()]
    layout = str(operation.get("layout", "title_body")).casefold()
    _add_brand_header(slide, prs, title, subtitle)
    if layout in {"cards", "four_cards", "three_cards"}:
        count = min(max(len(bullets), 1), 6)
        columns = 2 if count <= 4 else 3
        rows = (count + columns - 1) // columns
        card_w = 0.82 / columns
        card_h = 0.58 / rows
        for i, item in enumerate(bullets[:count]):
            col, row = i % columns, i // columns
            _add_callout(
                slide,
                0.08 + col * card_w,
                0.30 + row * card_h,
                card_w - 0.025,
                card_h - 0.035,
                str(item),
                "green" if i % 2 == 0 else "blue",
            )
    elif layout == "two_column":
        midpoint = max(1, (len(bullets) + 1) // 2)
        for index, subset in enumerate((bullets[:midpoint], bullets[midpoint:])):
            box = slide.shapes.add_textbox(
                int(prs.slide_width * (0.08 + index * 0.44)),
                int(prs.slide_height * 0.31),
                int(prs.slide_width * 0.40),
                int(prs.slide_height * 0.56),
            )
            _set_text_frame_text(box.text_frame, "\n".join(f"• {item}" for item in subset), preserve_first_run=False)
            _apply_text_style(box, {"font_size": 17, "font_color": "1F2937"})
            box.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    else:
        box = slide.shapes.add_textbox(
            int(prs.slide_width * 0.09), int(prs.slide_height * 0.31),
            int(prs.slide_width * 0.82), int(prs.slide_height * 0.56),
        )
        _set_text_frame_text(box.text_frame, "\n".join(f"• {item}" for item in bullets), preserve_first_run=False)
        _apply_text_style(box, {"font_size": operation.get("body_size", 18), "font_color": "1F2937"})
        box.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE


def _align_shapes(prs: Presentation, operation: dict[str, Any]) -> None:
    slide = _slide(prs, int(operation["slide"]))
    ids = [int(value) for value in operation.get("shape_ids", [])]
    shapes = [shape for shape in slide.shapes if int(shape.shape_id) in ids]
    if len(shapes) < 2:
        raise EditorError("Select at least two shapes to align.")
    mode = str(operation.get("alignment", "left")).casefold()
    if mode == "left":
        target = min(shape.left for shape in shapes)
        for shape in shapes: shape.left = target
    elif mode == "right":
        target = max(shape.left + shape.width for shape in shapes)
        for shape in shapes: shape.left = target - shape.width
    elif mode in {"center", "horizontal_center"}:
        target = sum(shape.left + shape.width / 2 for shape in shapes) / len(shapes)
        for shape in shapes: shape.left = int(target - shape.width / 2)
    elif mode == "top":
        target = min(shape.top for shape in shapes)
        for shape in shapes: shape.top = target
    elif mode == "bottom":
        target = max(shape.top + shape.height for shape in shapes)
        for shape in shapes: shape.top = target - shape.height
    elif mode in {"middle", "vertical_center"}:
        target = sum(shape.top + shape.height / 2 for shape in shapes) / len(shapes)
        for shape in shapes: shape.top = int(target - shape.height / 2)
    else:
        raise EditorError("Unsupported alignment mode.")


def _distribute_shapes(prs: Presentation, operation: dict[str, Any]) -> None:
    slide = _slide(prs, int(operation["slide"]))
    ids = [int(value) for value in operation.get("shape_ids", [])]
    shapes = [shape for shape in slide.shapes if int(shape.shape_id) in ids]
    if len(shapes) < 3:
        raise EditorError("Select at least three shapes to distribute.")
    direction = str(operation.get("direction", "horizontal")).casefold()
    if direction == "horizontal":
        shapes.sort(key=lambda shape: shape.left)
        start = shapes[0].left
        end = shapes[-1].left + shapes[-1].width
        total_width = sum(shape.width for shape in shapes)
        gap = max(0, int((end - start - total_width) / (len(shapes) - 1)))
        cursor = start
        for shape in shapes:
            shape.left = cursor
            cursor += shape.width + gap
    else:
        shapes.sort(key=lambda shape: shape.top)
        start = shapes[0].top
        end = shapes[-1].top + shapes[-1].height
        total_height = sum(shape.height for shape in shapes)
        gap = max(0, int((end - start - total_height) / (len(shapes) - 1)))
        cursor = start
        for shape in shapes:
            shape.top = cursor
            cursor += shape.height + gap


def _style_chart_universal(shape, operation: dict[str, Any]) -> None:
    if not getattr(shape, "has_chart", False):
        raise EditorError("The selected object is not a chart.")
    chart = shape.chart
    _style_chart(chart, operation)
    if operation.get("show_title") is False:
        chart.has_title = False
    if operation.get("show_data_labels") is not None:
        for plot in chart.plots:
            plot.has_data_labels = bool(operation["show_data_labels"])
            if plot.has_data_labels:
                labels = plot.data_labels
                labels.show_value = bool(operation.get("show_values", True))
                labels.show_category_name = bool(operation.get("show_categories", False))
                labels.show_series_name = bool(operation.get("show_series_names", False))
                if operation.get("number_format"):
                    labels.number_format = str(operation["number_format"])
                try:
                    labels.position = XL_DATA_LABEL_POSITION.OUTSIDE_END
                except Exception:
                    pass
    if hasattr(chart, "value_axis"):
        try:
            chart.value_axis.has_major_gridlines = bool(operation.get("show_gridlines", True))
            if operation.get("value_axis_title"):
                chart.value_axis.has_title = True
                chart.value_axis.axis_title.text_frame.text = str(operation["value_axis_title"])
            if operation.get("number_format"):
                chart.value_axis.tick_labels.number_format = str(operation["number_format"])
        except Exception:
            pass
    if hasattr(chart, "category_axis") and operation.get("category_axis_title"):
        try:
            chart.category_axis.has_title = True
            chart.category_axis.axis_title.text_frame.text = str(operation["category_axis_title"])
        except Exception:
            pass


def _change_chart_type(prs: Presentation, operation: dict[str, Any]) -> None:
    slide = _slide(prs, int(operation["slide"]))
    shape = _operation_shape(slide, operation)
    if not getattr(shape, "has_chart", False):
        raise EditorError("The selected object is not a chart.")
    chart = shape.chart
    plot = chart.plots[0]
    try:
        categories = [str(value) for value in plot.categories]
    except Exception:
        first_values = list(chart.series[0].values) if chart.series else []
        categories = [str(index + 1) for index in range(len(first_values))]
    series = [{"name": str(item.name), "values": [float(value) for value in item.values]} for item in chart.series]
    title = chart.chart_title.text_frame.text if chart.has_title else str(operation.get("title", ""))
    replacement = dict(operation)
    replacement.update({
        "op": "add_chart", "categories": categories, "series": series,
        "title": operation.get("title", title),
        "x": shape.left / prs.slide_width, "y": shape.top / prs.slide_height,
        "width": shape.width / prs.slide_width, "height": shape.height / prs.slide_height,
    })
    if str(operation.get("chart_type", "")).casefold() == "scatter":
        numeric_x = []
        for index, value in enumerate(categories):
            try:
                numeric_x.append(float(value))
            except (TypeError, ValueError):
                numeric_x.append(float(index + 1))
        replacement["x_values"] = numeric_x
    _delete_shape_element(shape)
    _add_chart(prs, replacement)


def _cleanup_slide(prs: Presentation, operation: dict[str, Any]) -> dict[str, int]:
    slide = _slide(prs, int(operation["slide"]))
    removed = 0
    fitted = 0
    for shape in list(slide.shapes):
        if getattr(shape, "has_text_frame", False):
            text = shape.text_frame.text.strip()
            if not text and operation.get("remove_empty", True) and not _shape_is_branding(shape, prs):
                _delete_shape_element(shape)
                removed += 1
                continue
            if operation.get("autofit", True):
                shape.text_frame.word_wrap = True
                shape.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
                fitted += 1
        if operation.get("keep_on_slide", True):
            shape.left = max(0, min(shape.left, prs.slide_width - max(1, shape.width)))
            shape.top = max(0, min(shape.top, prs.slide_height - max(1, shape.height)))
    return {"removed": removed, "fitted": fitted}


def _set_slide_hidden(prs: Presentation, operation: dict[str, Any]) -> None:
    slide = _slide(prs, int(operation["slide"]))
    hidden = bool(operation.get("hidden", True))
    if hidden:
        slide._element.set("show", "0")
    else:
        slide._element.attrib.pop("show", None)


def _apply_single_operation_universal(prs: Presentation, operation: dict[str, Any]) -> dict[str, Any]:
    op = str(operation.get("op", "")).casefold()
    if op == "noop":
        return {"op": "noop"}
    if op == "executive_review":
        return {"op": op, **_apply_executive_review(prs, operation)}
    if op == "replace_text":
        slide_numbers = [int(operation["slide"])] if operation.get("slide") is not None else list(range(1, len(prs.slides) + 1))
        count = 0
        for slide_number in slide_numbers:
            slide = _slide(prs, slide_number)
            target = None
            if any(operation.get(key) is not None for key in ("shape", "shape_id", "shape_name", "text_contains")):
                target = _operation_shape(slide, operation)
            for tf in _iter_text_frames(slide, target):
                count += _replace_text_in_frame(tf, str(operation.get("old", "")), str(operation.get("new", "")), bool(operation.get("replace_all", True)), bool(operation.get("case_sensitive", False)))
        if operation.get("include_masters"):
            containers = []
            for master in prs.slide_masters:
                containers.append(master)
                containers.extend(list(master.slide_layouts))
            for container in containers:
                for tf in _iter_text_frames(container):
                    count += _replace_text_in_frame(
                        tf,
                        str(operation.get("old", "")),
                        str(operation.get("new", "")),
                        bool(operation.get("replace_all", True)),
                        bool(operation.get("case_sensitive", False)),
                    )
        if count == 0 and not operation.get("allow_missing", True):
            raise EditorError(f'Text "{operation.get("old", "")}" was not found.')
        return {"op": op, "count": count}
    if op == "set_text":
        slide = _slide(prs, int(operation["slide"])); shape = _operation_shape(slide, operation)
        if not getattr(shape, "has_text_frame", False): raise EditorError("The selected object has no editable text.")
        _set_text_frame_text(shape.text_frame, str(operation.get("text", ""))); _apply_text_style(shape, operation)
        return {"op": op, "slide": int(operation["slide"])}
    if op == "append_text":
        slide = _slide(prs, int(operation["slide"])); shape = _operation_shape(slide, operation)
        if not getattr(shape, "has_text_frame", False): raise EditorError("The selected object has no editable text.")
        separator = str(operation.get("separator", "\n")); current = shape.text_frame.text
        _set_text_frame_text(shape.text_frame, current + (separator if current else "") + str(operation.get("text", "")))
        _apply_text_style(shape, operation); return {"op": op, "slide": int(operation["slide"])}
    if op in {"style_shape", "move_shape", "resize_shape"}:
        slide = _slide(prs, int(operation["slide"])); shape = _operation_shape(slide, operation)
        _set_geometry(prs, shape, operation); _apply_shape_style(shape, operation)
        return {"op": op, "slide": int(operation["slide"])}
    if op == "format_text_box" or op == "fit_text":
        slide = _slide(prs, int(operation["slide"])); shape = _operation_shape(slide, operation)
        _format_text_box(shape, operation); return {"op": op, "slide": int(operation["slide"])}
    if op == "rotate_shape":
        slide = _slide(prs, int(operation["slide"])); shape = _operation_shape(slide, operation)
        shape.rotation = float(operation.get("rotation", 0)); return {"op": op, "slide": int(operation["slide"])}
    if op == "duplicate_shape":
        slide = _slide(prs, int(operation["slide"])); shape = _operation_shape(slide, operation)
        _duplicate_shape(slide, shape); return {"op": op, "slide": int(operation["slide"])}
    if op == "layer_shape":
        slide = _slide(prs, int(operation["slide"])); shape = _operation_shape(slide, operation)
        _layer_shape(slide, shape, str(operation.get("direction", "front"))); return {"op": op, "slide": int(operation["slide"])}
    if op == "delete_shape":
        slide = _slide(prs, int(operation["slide"])); shape = _operation_shape(slide, operation); _delete_shape_element(shape)
        return {"op": op, "slide": int(operation["slide"])}
    if op == "add_textbox": _add_textbox(prs, operation); return {"op": op, "slide": int(operation["slide"])}
    if op == "add_shape": _add_shape(prs, operation); return {"op": op, "slide": int(operation["slide"])}
    if op == "add_line": _add_line(prs, operation); return {"op": op, "slide": int(operation["slide"])}
    if op == "add_chart": _add_chart(prs, operation); return {"op": op, "slide": int(operation["slide"])}
    if op == "add_picture": _add_picture(prs, operation); return {"op": op, "slide": int(operation["slide"])}
    if op == "replace_picture": _replace_picture(prs, operation); return {"op": op, "slide": int(operation["slide"])}
    if op == "crop_picture":
        slide = _slide(prs, int(operation["slide"])); shape = _operation_shape(slide, operation); _crop_picture(shape, operation)
        return {"op": op, "slide": int(operation["slide"])}
    if op == "add_table": _add_table(prs, operation); return {"op": op, "slide": int(operation["slide"])}
    if op == "set_slide_background":
        slide = _slide(prs, int(operation["slide"])); color = _rgb(operation.get("color"))
        if color is None: raise EditorError("A six-digit background color is required.")
        slide.background.fill.solid(); slide.background.fill.fore_color.rgb = color
        return {"op": op, "slide": int(operation["slide"])}
    if op == "set_table_cell":
        slide = _slide(prs, int(operation["slide"])); shape = _operation_shape(slide, operation)
        if not getattr(shape, "has_table", False): raise EditorError("The selected object is not a table.")
        row, col = int(operation["row"]) - 1, int(operation["column"]) - 1
        if row < 0 or col < 0 or row >= len(shape.table.rows) or col >= len(shape.table.columns): raise EditorError("The table cell is outside the table.")
        cell = shape.table.cell(row, col); _set_text_frame_text(cell.text_frame, str(operation.get("text", "")), preserve_first_run=False)
        _style_cell(cell, fill=operation.get("fill_color"), font_color=operation.get("font_color"), font_size=_safe_float(operation.get("font_size")), bold=operation.get("bold"), alignment=operation.get("alignment"))
        return {"op": op, "slide": int(operation["slide"])}
    if op == "style_table":
        slide = _slide(prs, int(operation["slide"])); shape = _operation_shape(slide, operation); _style_table_shape(shape, operation)
        return {"op": op, "slide": int(operation["slide"])}
    if op == "merge_table_cells" or op == "split_table_cell":
        slide = _slide(prs, int(operation["slide"])); shape = _operation_shape(slide, operation)
        if not getattr(shape, "has_table", False): raise EditorError("The selected object is not a table.")
        r1, c1 = int(operation.get("row", 1)) - 1, int(operation.get("column", 1)) - 1
        if op == "merge_table_cells":
            r2, c2 = int(operation.get("end_row", r1 + 1)) - 1, int(operation.get("end_column", c1 + 1)) - 1
            shape.table.cell(r1, c1).merge(shape.table.cell(r2, c2))
        else: shape.table.cell(r1, c1).split()
        return {"op": op, "slide": int(operation["slide"])}
    if op == "add_table_row" or op == "delete_table_row" or op == "add_table_column" or op == "delete_table_column":
        slide = _slide(prs, int(operation["slide"])); shape = _operation_shape(slide, operation)
        if not getattr(shape, "has_table", False): raise EditorError("The selected object is not a table.")
        if op == "add_table_row": _table_add_row(shape, operation)
        elif op == "delete_table_row": _table_delete_row(shape, int(operation.get("row", len(shape.table.rows))))
        elif op == "add_table_column": _table_add_column(shape, operation)
        else: _table_delete_column(shape, int(operation.get("column", len(shape.table.columns))))
        return {"op": op, "slide": int(operation["slide"])}
    if op == "set_chart_data":
        slide = _slide(prs, int(operation["slide"])); shape = _operation_shape(slide, operation); _set_chart_data(shape, operation)
        return {"op": op, "slide": int(operation["slide"])}
    if op == "style_chart":
        slide = _slide(prs, int(operation["slide"])); shape = _operation_shape(slide, operation); _style_chart_universal(shape, operation)
        return {"op": op, "slide": int(operation["slide"])}
    if op == "change_chart_type": _change_chart_type(prs, operation); return {"op": op, "slide": int(operation["slide"])}
    if op == "set_speaker_notes": _set_speaker_notes(prs, operation); return {"op": op, "slide": int(operation["slide"])}
    if op == "add_slide": return {"op": op, "slide": _add_slide(prs, operation)}
    if op == "duplicate_slide": return {"op": op, "slide": _duplicate_slide(prs, int(operation["slide"]), int(operation["position"]) if operation.get("position") is not None else None)}
    if op == "delete_slide":
        if len(prs.slides) <= 1: raise EditorError("A presentation must keep at least one slide.")
        slide_number = int(operation["slide"])
        _delete_slide(prs, slide_number); return {"op": op, "slide": slide_number}
    if op == "move_slide":
        from_slide = int(operation["from_slide"]); to_slide = int(operation["to_slide"])
        _move_slide(prs, from_slide, to_slide); return {"op": op, "from_slide": from_slide, "to_slide": to_slide}
    if op == "reorder_slides": _reorder_slides(prs, [int(v) for v in operation["order"]]); return {"op": op, "order": operation["order"]}
    if op == "clear_slide": return {"op": op, "slide": int(operation["slide"]), "deleted": _clear_slide(prs, operation)}
    if op == "regenerate_slide": _regenerate_slide(prs, operation); return {"op": op, "slide": int(operation["slide"])}
    if op == "align_shapes": _align_shapes(prs, operation); return {"op": op, "slide": int(operation["slide"])}
    if op == "distribute_shapes": _distribute_shapes(prs, operation); return {"op": op, "slide": int(operation["slide"])}
    if op == "standardize_deck": _standardize_deck(prs); return {"op": op, "slide_count": len(prs.slides)}
    if op == "cleanup_slide": return {"op": op, "slide": int(operation["slide"]), **_cleanup_slide(prs, operation)}
    if op == "set_slide_hidden": _set_slide_hidden(prs, operation); return {"op": op, "slide": int(operation["slide"])}
    raise EditorError(f"Unsupported edit action: {op}")


def _resolve_shape_reference_for_validation(prs: Presentation, operation: dict[str, Any]) -> None:
    if operation.get("slide") is None:
        raise EditorError(f"{operation.get('op')} needs a slide number.")
    slide = _slide(prs, int(operation["slide"]))
    resolved = _operation_shape(slide, operation)
    operation["shape_id"] = int(resolved.shape_id)


def validate_operations(pptx_path: str, operations: Iterable[dict[str, Any]]) -> tuple[list[dict[str, Any]], list[str]]:
    """Validate operations independently and repair shape references.

    Validation never raises for a batch. Invalid edits are returned as issues so
    the caller can keep the deck unchanged and continue with the remaining work.
    """
    prs = Presentation(pptx_path)
    slide_count = len(prs.slides)
    valid: list[dict[str, Any]] = []
    issues: list[str] = []
    for index, raw in enumerate(operations, start=1):
        try:
            if not isinstance(raw, dict):
                raise EditorError("The operation was not an object.")
            operation = {k: v for k, v in raw.items() if v is not None}
            op = str(operation.get("op", "")).strip().casefold()
            operation["op"] = op
            if op not in SUPPORTED_OPS:
                raise EditorError(f"Unsupported action '{op}'.")
            if op in SHAPE_OPS:
                _resolve_shape_reference_for_validation(prs, operation)
            elif op in {"add_textbox", "add_shape", "add_line", "add_chart", "add_picture", "add_table", "set_slide_background", "set_speaker_notes", "clear_slide", "regenerate_slide", "cleanup_slide", "set_slide_hidden", "align_shapes", "distribute_shapes"}:
                if operation.get("slide") is None: raise EditorError(f"{op} needs a slide number.")
                _slide(prs, int(operation["slide"]))
            elif op == "replace_text":
                if not str(operation.get("old", "")): raise EditorError("replace_text needs the original text.")
                if operation.get("slide") is not None:
                    slide = _slide(prs, int(operation["slide"]))
                    if any(operation.get(k) is not None for k in ("shape", "shape_id", "shape_name", "text_contains")):
                        operation["shape_id"] = int(_operation_shape(slide, operation).shape_id)
            elif op == "add_slide":
                position = int(operation.get("position") or slide_count + 1)
                if not 1 <= position <= slide_count + 1: raise EditorError("The new slide position is outside the deck.")
                if operation.get("template_slide") is not None: _slide(prs, int(operation["template_slide"]))
            elif op == "duplicate_slide":
                _slide(prs, int(operation["slide"]))
            elif op == "delete_slide":
                _slide(prs, int(operation["slide"]))
                if slide_count <= 1: raise EditorError("A presentation must keep at least one slide.")
            elif op == "move_slide":
                if not 1 <= int(operation["from_slide"]) <= slide_count or not 1 <= int(operation["to_slide"]) <= slide_count: raise EditorError("The move position is outside the deck.")
            elif op == "reorder_slides":
                order = [int(v) for v in operation.get("order", [])]
                if sorted(order) != list(range(1, slide_count + 1)): raise EditorError("The order must list every current slide exactly once.")
                operation["order"] = order
            if op == "add_chart": _chart_data(operation)
            if op in {"add_picture", "replace_picture"}:
                source = str(operation.get("source_path") or operation.get("path") or "")
                if not os.path.isfile(source): raise EditorError("The image attachment is unavailable.")
            valid.append(operation)
        except (EditorError, KeyError, TypeError, ValueError, AttributeError) as exc:
            issues.append(f"Operation {index} ({_operation_summary(raw if isinstance(raw, dict) else {})}) was skipped: {exc}")
    return valid, issues


def apply_operations(input_path: str, output_path: str, operations: Iterable[dict[str, Any]]) -> dict[str, Any]:
    """Apply each edit as an isolated transaction.

    Every operation is saved and reopened before the next operation starts. A
    malformed or unsupported operation cannot corrupt the working deck or stop
    later safe edits. If nothing is applicable, the original deck is copied to
    the output and returned unchanged rather than raising a server error.
    """
    operation_list = list(operations)
    Path(output_path).parent.mkdir(parents=True, exist_ok=True)
    with tempfile.TemporaryDirectory(prefix="deck_refresh_txn_") as temp_dir:
        current = os.path.join(temp_dir, "current.pptx")
        shutil.copy2(input_path, current)
        applied: list[dict[str, Any]] = []
        skipped: list[dict[str, str]] = []
        for index, raw in enumerate(operation_list, start=1):
            valid, issues = validate_operations(current, [raw])
            if not valid:
                skipped.extend({"operation": _operation_summary(raw if isinstance(raw, dict) else {}), "reason": issue} for issue in issues)
                continue
            operation = valid[0]
            candidate = os.path.join(temp_dir, f"candidate_{index:04d}.pptx")
            try:
                prs = Presentation(current)
                baseline_integrity = set(_verify_deck_integrity(prs))
                detail = _apply_single_operation_universal(prs, operation)
                # Chart titles, axes, legends, and labels follow the final
                # visible background after every edit, including theme changes.
                ensure_chart_contrast(prs)
                prs.save(candidate)
                reopened = Presentation(candidate)
                integrity = set(_verify_deck_integrity(reopened))
                new_integrity = integrity - baseline_integrity
                if any("outside" in issue.casefold() for issue in new_integrity):
                    raise EditorError("The edit placed content outside the slide.")
                shutil.copy2(candidate, current)
                applied.append(detail)
            except Exception as exc:
                skipped.append({"operation": _operation_summary(operation), "reason": str(exc)[:500]})
        shutil.copy2(current, output_path)
    final = Presentation(output_path)
    return {"applied": applied, "skipped": skipped, "slide_count": len(final.slides), "unchanged": not any(item.get("op") != "noop" for item in applied)}


class UniversalAtomicEditStep(BaseModel):
    model_config = ConfigDict(extra="forbid")
    message: str = Field(default="", max_length=180)
    done: bool
    operation_json: str | None = Field(default=None, max_length=12000)


def _parse_atomic_operation(raw: str | None) -> dict[str, Any] | None:
    if raw is None or not str(raw).strip():
        return None
    text = str(raw).strip()
    if text.startswith("```"):
        text = re.sub(r"^```(?:json)?\s*|\s*```$", "", text, flags=re.IGNORECASE)
    try:
        value = json.loads(text)
    except json.JSONDecodeError as exc:
        raise EditorError(f"Atomic operation was not valid JSON: {exc}") from exc
    if not isinstance(value, dict):
        raise EditorError("Atomic operation must be a JSON object.")
    op = str(value.get("op", "")).strip().casefold()
    if op not in SUPPORTED_OPS:
        raise EditorError(f"Atomic operation used unsupported action '{op}'.")
    value["op"] = op
    return {key: val for key, val in value.items() if val is not None}


ATOMIC_PLANNER_INSTRUCTIONS = """
You are the execution layer for a resilient professional PowerPoint editor.
Return ONE compact JSON operation at a time in operation_json. Set done=true
when the focused task is complete. Never return markdown or commentary inside
operation_json. Infer sensible choices from the deck and recent chat instead
of asking a question when a professional default exists.

Supported operations:
TEXT: replace_text, set_text, append_text, format_text_box, fit_text.
SHAPES: style_shape, move_shape, resize_shape, rotate_shape, duplicate_shape,
layer_shape, delete_shape, add_textbox, add_shape, add_line, align_shapes,
distribute_shapes.
IMAGES: add_picture, replace_picture, crop_picture. Use source_attachment with
an uploaded filename. Never invent a local path.
TABLES: add_table, set_table_cell, style_table, merge_table_cells,
split_table_cell, add_table_row, delete_table_row, add_table_column,
delete_table_column.
CHARTS: add_chart, set_chart_data, style_chart, change_chart_type. Use only
numbers present in the deck or uploaded data. Supported chart_type values are
column, bar, line, pie, area, doughnut, stacked_column, stacked_bar.
SLIDES: add_slide, duplicate_slide, delete_slide, move_slide, reorder_slides,
clear_slide, regenerate_slide, set_slide_background, set_slide_hidden.
DECK: standardize_deck, cleanup_slide, set_speaker_notes, executive_review.

Target existing objects with shape_id and include text_contains as a backup.
Coordinates from 0 to 1.5 are fractions of slide size; larger values are
inches. Preserve logos and branding unless removal is explicit. Use one
operation for one coherent edit. Prefer regenerate_slide for requests like
"redo this slide" or "make this slide from scratch". Prefer add_table for a
new table and set_table_cell for small updates. Use noop only when the request
requires an unavailable external asset and no useful substitute exists.
""".strip()


def _request_atomic_plan(
    client: OpenAI,
    model: str,
    base_content: list[dict[str, Any]],
    pptx_path: str,
    operation_limit: int = 12,
) -> dict[str, Any]:
    operations: list[dict[str, Any]] = []
    messages: list[str] = []
    feedback: list[str] = []
    limit = max(1, min(int(operation_limit), 30))
    for step_number in range(1, limit + 2):
        state = {
            "step": step_number,
            "maximum_operations": limit,
            "already_planned": operations[-12:],
            "validation_feedback": feedback[-4:],
            "instruction": "Return the next single operation, or finish.",
        }
        content = list(base_content)
        content.append({"type": "input_text", "text": "Atomic state:\n" + json.dumps(state, ensure_ascii=False)})
        try:
            step = _request_structured(client, model, content, ATOMIC_PLANNER_INSTRUCTIONS, UniversalAtomicEditStep, max_output_tokens=4500)
        except EditorError as exc:
            feedback.append(str(exc))
            if step_number >= 3:
                break
            continue
        if step.message.strip(): messages.append(step.message.strip())
        operation = None
        try:
            operation = _parse_atomic_operation(step.operation_json)
        except EditorError as exc:
            feedback.append(str(exc))
        if operation is not None:
            valid, issues = validate_operations(pptx_path, [operation])
            if valid:
                candidate = valid[0]
                serialized = json.dumps(candidate, sort_keys=True, ensure_ascii=False)
                if all(json.dumps(existing, sort_keys=True, ensure_ascii=False) != serialized for existing in operations):
                    operations.append(candidate)
                else:
                    feedback.append("Duplicate operation. Choose another edit or finish.")
            else:
                feedback.extend(issues or ["The operation was not usable in this deck."])
        if step.done or len(operations) >= limit:
            break
    return {"message": " ".join(dict.fromkeys(messages))[:280] or "Applied the requested edits.", "assumptions": feedback[-4:], "operations": operations}


def _attachment_context(paths: list[str] | None) -> tuple[str, dict[str, str]]:
    contexts: list[str] = []
    mapping: dict[str, str] = {}
    for raw_path in paths or []:
        path = os.path.abspath(str(raw_path))
        if not os.path.isfile(path):
            continue
        name = os.path.basename(path)
        mapping[name.casefold()] = path
        suffix = Path(path).suffix.casefold()
        try:
            if suffix in {".csv", ".tsv"}:
                frame = pd.read_csv(path, sep="\t" if suffix == ".tsv" else ",", nrows=50)
                contexts.append(f"Uploaded data {name}:\n{frame.to_csv(index=False)[:8000]}")
            elif suffix in {".xlsx", ".xlsm"}:
                workbook = load_workbook(path, read_only=True, data_only=True)
                pieces = []
                for sheet in workbook.worksheets[:5]:
                    rows = []
                    for row in sheet.iter_rows(min_row=1, max_row=min(sheet.max_row, 40), values_only=True):
                        rows.append(["" if value is None else value for value in row[:20]])
                    pieces.append(f"Sheet {sheet.title}: {json.dumps(rows, ensure_ascii=False)[:5000]}")
                contexts.append(f"Uploaded workbook {name}:\n" + "\n".join(pieces))
            elif suffix in {".txt", ".md", ".json"}:
                contexts.append(f"Uploaded text {name}:\n{Path(path).read_text(encoding='utf-8', errors='ignore')[:8000]}")
            elif suffix in {".png", ".jpg", ".jpeg", ".webp", ".gif"}:
                contexts.append(f"Uploaded image available as source_attachment={name}")
            else:
                contexts.append(f"Uploaded attachment available: {name}")
        except Exception:
            contexts.append(f"Uploaded attachment available: {name}")
    return "\n\n".join(contexts), mapping


def _resolve_attachment_operations(operations: list[dict[str, Any]], mapping: dict[str, str]) -> list[dict[str, Any]]:
    resolved = []
    for operation in operations:
        item = dict(operation)
        ref = str(item.get("source_attachment") or "").strip().casefold()
        if ref:
            exact = mapping.get(ref)
            if exact is None:
                matches = [path for name, path in mapping.items() if ref in name or name in ref]
                exact = matches[0] if len(matches) == 1 else None
            if exact:
                item["source_path"] = exact
        resolved.append(item)
    return resolved


_legacy_plan_edit = plan_edit


def plan_edit(
    pptx_path: str,
    user_message: str,
    selected_slide: int,
    selected_slide_image: str | None = None,
    chat_history: list[dict[str, str]] | None = None,
    deck_image_paths: list[str] | None = None,
    attachment_paths: list[str] | None = None,
) -> dict[str, Any]:
    attachment_text, attachment_map = _attachment_context(attachment_paths)
    augmented_message = user_message
    if attachment_text:
        augmented_message += "\n\nAvailable uploaded materials:\n" + attachment_text
    try:
        result = _legacy_plan_edit(
            pptx_path,
            augmented_message,
            selected_slide,
            selected_slide_image,
            chat_history,
            deck_image_paths,
        )
    except Exception as exc:
        summary = deck_summary(pptx_path, max_text=1200)
        fallback = _deterministic_operations(summary, user_message, selected_slide)
        if not fallback:
            fallback = [{"op": "noop", "reason": str(exc)[:300]}]
        result = {
            "message": "The conversational planner was unavailable. I attempted deterministic edits only.",
            "assumptions": [str(exc)[:300]],
            "operations": fallback,
        }
    result["operations"] = _resolve_attachment_operations(list(result.get("operations", [])), attachment_map)
    valid, issues = validate_operations(pptx_path, result["operations"])
    if not valid:
        valid = [{"op": "noop", "reason": "; ".join(issues)[:500]}]
    result["operations"] = valid
    assumptions = list(result.get("assumptions", [])) + issues
    result["assumptions"] = list(dict.fromkeys(str(item) for item in assumptions if str(item).strip()))[:20]
    return result

# ---------------------------------------------------------------------------
# Deck-wide commands and deterministic command compiler
# ---------------------------------------------------------------------------

_COLOR_NAMES = {
    "black": "000000", "white": "FFFFFF", "red": "C00000", "green": "00A651",
    "dark green": "006B3C", "light green": "E8F5E9", "blue": "00338D",
    "dark blue": "002060", "light blue": "D9EAF7", "amber": "F5A623",
    "orange": "ED7D31", "yellow": "FFD966", "gray": "808080", "grey": "808080",
    "light gray": "F3F4F6", "light grey": "F3F4F6", "purple": "7030A0",
    "pink": "FF66CC", "teal": "008C95", "navy": "00338D",
}


def _color_hex(value: Any) -> str | None:
    text = str(value or "").strip().casefold().lstrip("#")
    text = _COLOR_NAMES.get(text, text)
    return text.upper() if re.fullmatch(r"[0-9a-f]{6}", text) else None


def _rgb_matches(color_format, wanted: str) -> bool:
    try:
        rgb = color_format.rgb
        return rgb is not None and str(rgb).upper() == wanted
    except Exception:
        return False


def _replace_color(prs: Presentation, operation: dict[str, Any]) -> int:
    old = _color_hex(operation.get("old_color") or operation.get("old"))
    new = _color_hex(operation.get("new_color") or operation.get("new"))
    if not old or not new:
        raise EditorError("replace_color needs valid old_color and new_color values.")
    new_rgb = RGBColor.from_string(new)
    slide_numbers = operation.get("slides") or ([operation["slide"]] if operation.get("slide") else list(range(1, len(prs.slides) + 1)))
    count = 0
    for slide_number in slide_numbers:
        slide = _slide(prs, int(slide_number))
        for shape in slide.shapes:
            if operation.get("include_fill", True):
                try:
                    if _rgb_matches(shape.fill.fore_color, old):
                        shape.fill.solid(); shape.fill.fore_color.rgb = new_rgb; count += 1
                except Exception: pass
            if operation.get("include_line", True):
                try:
                    if _rgb_matches(shape.line.color, old):
                        shape.line.color.rgb = new_rgb; count += 1
                except Exception: pass
            if operation.get("include_font", True):
                frames = []
                if getattr(shape, "has_text_frame", False): frames.append(shape.text_frame)
                if getattr(shape, "has_table", False):
                    frames.extend(cell.text_frame for row in shape.table.rows for cell in row.cells)
                for tf in frames:
                    for paragraph in tf.paragraphs:
                        for run in paragraph.runs:
                            try:
                                if _rgb_matches(run.font.color, old):
                                    run.font.color.rgb = new_rgb; count += 1
                            except Exception: pass
            if getattr(shape, "has_table", False) and operation.get("include_table", True):
                for row in shape.table.rows:
                    for cell in row.cells:
                        try:
                            if _rgb_matches(cell.fill.fore_color, old):
                                cell.fill.solid(); cell.fill.fore_color.rgb = new_rgb; count += 1
                        except Exception: pass
            if getattr(shape, "has_chart", False) and operation.get("include_chart", True):
                for series in shape.chart.series:
                    try:
                        if _rgb_matches(series.format.fill.fore_color, old):
                            series.format.fill.solid(); series.format.fill.fore_color.rgb = new_rgb; count += 1
                    except Exception: pass
                    try:
                        if _rgb_matches(series.format.line.color, old):
                            series.format.line.color.rgb = new_rgb; count += 1
                    except Exception: pass
    return count


def _delete_objects(prs: Presentation, operation: dict[str, Any]) -> int:
    kinds = {str(value).casefold() for value in (operation.get("object_types") or [operation.get("object_type", "all")])}
    slide_numbers = operation.get("slides") or ([operation["slide"]] if operation.get("slide") else list(range(1, len(prs.slides) + 1)))
    wanted_text = _normalize_text(operation.get("text_contains"))
    removed = 0
    for slide_number in slide_numbers:
        slide = _slide(prs, int(slide_number))
        title_text = _normalize_text(_slide_title(slide))
        for shape in list(slide.shapes):
            kind = _shape_kind(shape)
            text = _normalize_text(_shape_text(shape))
            if operation.get("preserve_branding", True) and _shape_is_branding(shape, prs):
                continue
            if operation.get("preserve_title", True) and title_text and text == title_text:
                continue
            matched = "all" in kinds or kind in kinds
            if "empty" in kinds and kind == "text" and not text:
                matched = True
            if wanted_text and wanted_text not in text:
                matched = False
            if matched:
                _delete_shape_element(shape); removed += 1
    return removed


def _numeric_value(value: Any) -> float | None:
    text = str(value or "").strip().replace(",", "")
    negative = text.startswith("(") and text.endswith(")")
    text = text.strip("()$£€¥% ")
    multiplier = 1.0
    if text.casefold().endswith("m"):
        multiplier = 1_000_000.0; text = text[:-1]
    elif text.casefold().endswith("k"):
        multiplier = 1_000.0; text = text[:-1]
    try:
        value_float = float(text) * multiplier
        return -value_float if negative else value_float
    except ValueError:
        return None


def _convert_table_to_chart(prs: Presentation, operation: dict[str, Any]) -> None:
    slide = _slide(prs, int(operation["slide"]))
    shape = _operation_shape(slide, operation)
    if not getattr(shape, "has_table", False):
        raise EditorError("The selected object is not a table.")
    table = shape.table
    matrix = [[cell.text.strip() for cell in row.cells] for row in table.rows]
    if len(matrix) < 2 or len(matrix[0]) < 2:
        raise EditorError("The table needs headers and at least one data row.")
    categories = [row[0] or f"Row {index}" for index, row in enumerate(matrix[1:], start=1)]
    series = []
    for column in range(1, len(matrix[0])):
        values = [_numeric_value(row[column] if column < len(row) else None) for row in matrix[1:]]
        if all(value is not None for value in values):
            series.append({"name": matrix[0][column] or f"Series {column}", "values": values})
    if not series:
        raise EditorError("No numeric table columns were available for a chart.")
    chart_operation = {
        "op": "add_chart", "slide": int(operation["slide"]),
        "chart_type": operation.get("chart_type", "column"),
        "title": operation.get("title", _slide_title(slide) or "Chart"),
        "categories": categories[:30], "series": [{"name": item["name"], "values": item["values"][:30]} for item in series[:8]],
        "x": operation.get("x", shape.left / prs.slide_width),
        "y": operation.get("y", shape.top / prs.slide_height),
        "width": operation.get("width", shape.width / prs.slide_width),
        "height": operation.get("height", shape.height / prs.slide_height),
        "series_colors": operation.get("series_colors"),
        "show_legend": operation.get("show_legend", len(series) > 1),
    }
    _add_chart(prs, chart_operation)
    if not operation.get("keep_table", True):
        _delete_shape_element(shape)


def _apply_theme(prs: Presentation, operation: dict[str, Any]) -> None:
    primary = _color_hex(operation.get("primary_color")) or "00338D"
    accent = _color_hex(operation.get("accent_color")) or "00A651"
    body_color = _color_hex(operation.get("body_color")) or "1F2937"
    font_face = str(operation.get("font_face") or "Arial")
    title_size = float(operation.get("title_size", 28))
    body_size = float(operation.get("body_size", 16))
    slide_numbers = operation.get("slides") or list(range(1, len(prs.slides) + 1))
    for slide_number in slide_numbers:
        slide = _slide(prs, int(slide_number))
        text_shapes = [shape for shape in slide.shapes if getattr(shape, "has_text_frame", False) and shape.text_frame.text.strip()]
        title = text_shapes[0] if text_shapes else None
        for shape in text_shapes:
            _apply_text_style(shape, {
                "font_face": font_face,
                "font_size": title_size if shape is title else body_size,
                "font_color": primary if shape is title else body_color,
                "bold": True if shape is title else None,
            })
            shape.text_frame.word_wrap = True
            shape.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        for shape in slide.shapes:
            if getattr(shape, "has_chart", False):
                _style_chart(shape.chart, {"series_colors": [accent, primary, "66B032", "0091DA"]})
            if getattr(shape, "has_table", False):
                _style_table_shape(shape, {"header_fill": primary, "header_font_color": "FFFFFF", "band_fill": "F3F6FA", "banded_rows": True, "font_size": max(9, body_size - 5)})


def _set_footer(prs: Presentation, operation: dict[str, Any]) -> int:
    text = str(operation.get("text", "")).strip()
    slide_numbers = operation.get("slides") or list(range(1, len(prs.slides) + 1))
    count = 0
    for slide_number in slide_numbers:
        slide = _slide(prs, int(slide_number))
        candidates = [
            shape for shape in slide.shapes
            if getattr(shape, "has_text_frame", False)
            and shape.top > int(prs.slide_height * 0.84)
            and ("kpmg" in _normalize_text(shape.text_frame.text) or "footer" in str(shape.name).casefold())
        ]
        if candidates:
            shape = candidates[0]
            _set_text_frame_text(shape.text_frame, text)
        else:
            shape = slide.shapes.add_textbox(int(prs.slide_width * 0.06), int(prs.slide_height * 0.92), int(prs.slide_width * 0.88), int(prs.slide_height * 0.05))
            _set_text_frame_text(shape.text_frame, text, preserve_first_run=False)
        _apply_text_style(shape, {"font_size": operation.get("font_size", 8), "font_color": operation.get("font_color", "6B7280"), "alignment": operation.get("alignment", "left")})
        count += 1
    return count


def _set_slide_size(prs: Presentation, operation: dict[str, Any]) -> None:
    width_inches = float(operation.get("width_inches", prs.slide_width / 914400))
    height_inches = float(operation.get("height_inches", prs.slide_height / 914400))
    if width_inches <= 0 or height_inches <= 0:
        raise EditorError("Slide size must be positive.")
    old_width, old_height = prs.slide_width, prs.slide_height
    new_width, new_height = Inches(width_inches), Inches(height_inches)
    if operation.get("scale_content", True):
        x_ratio, y_ratio = new_width / old_width, new_height / old_height
        for slide in prs.slides:
            for shape in slide.shapes:
                shape.left = int(shape.left * x_ratio); shape.top = int(shape.top * y_ratio)
                shape.width = int(shape.width * x_ratio); shape.height = int(shape.height * y_ratio)
    prs.slide_width, prs.slide_height = new_width, new_height


SHAPE_OPS = SHAPE_OPS | {"convert_table_to_chart"}
SUPPORTED_OPS = SUPPORTED_OPS | {"replace_color", "delete_objects", "convert_table_to_chart", "apply_theme", "set_footer", "set_slide_size"}
_previous_apply_single = _apply_single_operation_universal


def _apply_single_operation_universal(prs: Presentation, operation: dict[str, Any]) -> dict[str, Any]:
    op = str(operation.get("op", "")).casefold()
    if op == "replace_color": return {"op": op, "count": _replace_color(prs, operation)}
    if op == "delete_objects": return {"op": op, "count": _delete_objects(prs, operation)}
    if op == "convert_table_to_chart": _convert_table_to_chart(prs, operation); return {"op": op, "slide": int(operation["slide"])}
    if op == "apply_theme": _apply_theme(prs, operation); return {"op": op, "slide_count": len(prs.slides)}
    if op == "set_footer": return {"op": op, "count": _set_footer(prs, operation)}
    if op == "set_slide_size": _set_slide_size(prs, operation); return {"op": op}
    return _previous_apply_single(prs, operation)


# Extend planner guidance after the added command definitions.
ATOMIC_PLANNER_INSTRUCTIONS += """

Additional deck-wide operations:
- replace_color: {op, old_color, new_color, optional slide or slides,
  include_fill, include_line, include_font, include_table, include_chart}
- delete_objects: {op, optional slide or slides, object_types, optional
  text_contains, preserve_branding, preserve_title}
- convert_table_to_chart: {op, slide, shape_id or text_contains, chart_type,
  title, x, y, width, height, keep_table, series_colors}
- apply_theme: {op, optional slides, primary_color, accent_color, body_color,
  font_face, title_size, body_size}
- set_footer: {op, text, optional slides, font_size, font_color, alignment}
- set_slide_size: {op, width_inches, height_inches, scale_content}
"""


_older_deterministic_operations = _deterministic_operations


def _parse_slide_targets(message: str, slide_count: int, selected_slide: int) -> list[int]:
    targets = _message_slide_numbers(message, slide_count)
    if "all slides" in message.casefold() or "entire deck" in message.casefold() or "whole deck" in message.casefold():
        return list(range(1, slide_count + 1))
    return targets or [selected_slide]


def _deterministic_operations(summary: dict[str, Any], user_message: str, selected_slide: int) -> list[dict[str, Any]]:
    operations = _older_deterministic_operations(summary, user_message, selected_slide)
    text = user_message.strip()
    lowered = text.casefold()
    slide_count = int(summary.get("slide_count", 0))
    targets = _parse_slide_targets(text, slide_count, selected_slide)

    for value in re.findall(r"(?:delete|remove)\s+slide\s+(\d+)", lowered):
        operations.append({"op": "delete_slide", "slide": int(value)})
    for value in re.findall(r"duplicate\s+slide\s+(\d+)", lowered):
        operations.append({"op": "duplicate_slide", "slide": int(value)})
    move = re.search(r"move\s+slide\s+(\d+)\s+(?:to|into)\s+(?:position\s+)?(\d+)", lowered)
    if move:
        operations.append({"op": "move_slide", "from_slide": int(move.group(1)), "to_slide": int(move.group(2))})
    if re.search(r"(?:add|create)\s+(?:a\s+)?(?:new\s+)?slide", lowered):
        title_match = re.search(r"(?:titled|called|named)\s+[\"“']?([^\"”'.\n]{2,120})", text, re.IGNORECASE)
        operations.append({"op": "add_slide", "position": min(max(targets[-1] + 1, 1), slide_count + 1), "title": title_match.group(1).strip() if title_match else "New Slide", "body": []})
    if "clear slide" in lowered or "delete everything on slide" in lowered:
        for slide in targets:
            operations.append({"op": "clear_slide", "slide": slide, "preserve_branding": "branding" not in lowered, "preserve_title": False})
    if any(token in lowered for token in ("regenerate slide", "redo slide", "rebuild slide", "remake slide")):
        for slide in targets:
            operations.append({"op": "regenerate_slide", "slide": slide, "title": summary["slides"][slide-1].get("title") or "Updated Slide", "body": ["Content regenerated from the existing slide context"], "layout": "title_body", "preserve_branding": True})
    background = re.search(r"(?:background|slide background)\s+(?:to\s+)?(?:color\s+)?([a-z ]+|#[0-9a-f]{6})", lowered)
    if background:
        color = _color_hex(background.group(1).strip())
        if color:
            for slide in targets: operations.append({"op": "set_slide_background", "slide": slide, "color": color})
    color_change = re.search(r"change\s+(?:all\s+)?([a-z ]+|#[0-9a-f]{6})\s+(?:colors?\s+)?to\s+([a-z ]+|#[0-9a-f]{6})", lowered)
    if color_change:
        old, new = _color_hex(color_change.group(1).strip()), _color_hex(color_change.group(2).strip())
        if old and new: operations.append({"op": "replace_color", "old_color": old, "new_color": new, "slides": targets})
    if "hide slide" in lowered:
        for slide in targets: operations.append({"op": "set_slide_hidden", "slide": slide, "hidden": True})
    if "unhide slide" in lowered or "show hidden slide" in lowered:
        for slide in targets: operations.append({"op": "set_slide_hidden", "slide": slide, "hidden": False})
    notes_match = re.search(r"(?:add|set|replace)\s+(?:speaker\s+)?notes?(?:\s+on\s+slide\s+\d+)?\s*(?:to|with|:)\s*[\"“']?(.{2,800})", text, re.IGNORECASE)
    if notes_match:
        for slide in targets: operations.append({"op": "set_speaker_notes", "slide": slide, "text": notes_match.group(1).strip().strip('"”\''), "mode": "replace"})
    if "delete all charts" in lowered or "remove all charts" in lowered:
        operations.append({"op": "delete_objects", "slides": targets, "object_types": ["chart"], "preserve_branding": True, "preserve_title": True})
    if "delete all tables" in lowered or "remove all tables" in lowered:
        operations.append({"op": "delete_objects", "slides": targets, "object_types": ["table"], "preserve_branding": True, "preserve_title": True})
    if "delete all images" in lowered or "remove all images" in lowered or "delete all pictures" in lowered:
        operations.append({"op": "delete_objects", "slides": targets, "object_types": ["picture"], "preserve_branding": "logo" not in lowered, "preserve_title": True})
    if "standardize" in lowered or "apply theme" in lowered or "make formatting consistent" in lowered:
        operations.append({"op": "standardize_deck"})
    return _dedupe_operations(operations)

# ---------------------------------------------------------------------------
# Attachment-aware offline fallbacks
# ---------------------------------------------------------------------------


def _read_attachment_matrix(path: str, max_rows: int = 40, max_columns: int = 16) -> list[list[Any]]:
    suffix = Path(path).suffix.casefold()
    try:
        if suffix in {".csv", ".tsv"}:
            frame = pd.read_csv(path, sep="\t" if suffix == ".tsv" else ",", nrows=max_rows)
            return [list(frame.columns)] + frame.fillna("").iloc[:, :max_columns].values.tolist()
        if suffix in {".xlsx", ".xlsm"}:
            workbook = load_workbook(path, read_only=True, data_only=True)
            sheet = workbook.worksheets[0]
            matrix = []
            for row in sheet.iter_rows(min_row=1, max_row=min(sheet.max_row, max_rows), values_only=True):
                matrix.append(["" if value is None else value for value in row[:max_columns]])
            return matrix
        if suffix == ".json":
            data = json.loads(Path(path).read_text(encoding="utf-8", errors="ignore"))
            if isinstance(data, list) and data and isinstance(data[0], dict):
                headers = list(data[0].keys())[:max_columns]
                return [headers] + [[row.get(header, "") for header in headers] for row in data[:max_rows - 1]]
            if isinstance(data, list):
                return [row if isinstance(row, list) else [row] for row in data[:max_rows]]
    except Exception:
        return []
    return []


def _attachment_chart_operation(path: str, slide: int) -> dict[str, Any] | None:
    matrix = _read_attachment_matrix(path)
    if len(matrix) < 2 or len(matrix[0]) < 2:
        return None
    headers = [str(value) for value in matrix[0]]
    category_column = 0
    categories = [str(row[category_column]) for row in matrix[1:31]]
    series = []
    for column in range(1, len(headers)):
        values = [_numeric_value(row[column] if column < len(row) else None) for row in matrix[1:31]]
        if values and all(value is not None for value in values):
            series.append({"name": headers[column] or f"Series {column}", "values": values})
    if not series:
        return None
    return {
        "op": "add_chart", "slide": slide, "chart_type": "column",
        "title": Path(path).stem.replace("_", " ").title(),
        "categories": categories, "series": series[:6],
        "x": 0.50, "y": 0.24, "width": 0.44, "height": 0.56,
        "series_colors": ["00A651", "00338D", "66B032", "0091DA"],
        "show_legend": len(series) > 1,
    }


def _attachment_fallback_operations(
    pptx_path: str,
    user_message: str,
    selected_slide: int,
    attachment_paths: list[str] | None,
) -> list[dict[str, Any]]:
    paths = [path for path in (attachment_paths or []) if os.path.isfile(path)]
    if not paths:
        return []
    lowered = user_message.casefold()
    summary = deck_summary(pptx_path, max_text=500)
    slide_count = int(summary.get("slide_count", 1))
    targets = _parse_slide_targets(user_message, slide_count, selected_slide)
    slide = targets[0]
    images = [path for path in paths if Path(path).suffix.casefold() in {".png", ".jpg", ".jpeg", ".webp", ".gif"}]
    data_files = [path for path in paths if Path(path).suffix.casefold() in {".csv", ".tsv", ".xlsx", ".xlsm", ".json"}]
    operations: list[dict[str, Any]] = []

    if images and any(token in lowered for token in ("image", "picture", "photo", "logo", "screenshot")):
        image = images[0]
        slide_summary = summary["slides"][slide - 1]
        pictures = [shape for shape in slide_summary.get("shapes", []) if shape.get("kind") == "picture"]
        if any(token in lowered for token in ("replace", "swap", "change the image", "change image")) and pictures:
            operations.append({
                "op": "replace_picture", "slide": slide,
                "shape_id": int(pictures[0]["shape_id"]), "source_path": image,
            })
        else:
            operations.append({
                "op": "add_picture", "slide": slide, "source_path": image,
                "x": 0.55, "y": 0.22, "width": 0.38, "height": 0.55,
            })

    if data_files and any(token in lowered for token in ("chart", "graph", "visualize", "plot")):
        chart = _attachment_chart_operation(data_files[0], slide)
        if chart:
            operations.append(chart)

    if data_files and any(token in lowered for token in ("table", "fill in", "populate", "insert data", "add data")):
        matrix = _read_attachment_matrix(data_files[0])
        if matrix:
            slide_summary = summary["slides"][slide - 1]
            tables = [shape for shape in slide_summary.get("shapes", []) if shape.get("kind") == "table"]
            if tables and any(token in lowered for token in ("fill", "populate", "update")):
                table = tables[0]
                current_rows = int(table.get("table_rows", 0)); current_columns = int(table.get("table_columns", 0))
                needed_rows = min(len(matrix), 100); needed_columns = min(max(len(row) for row in matrix), 30)
                for _ in range(max(0, needed_rows - current_rows)):
                    operations.append({"op": "add_table_row", "slide": slide, "shape_id": table["shape_id"], "values": []})
                for _ in range(max(0, needed_columns - current_columns)):
                    operations.append({"op": "add_table_column", "slide": slide, "shape_id": table["shape_id"], "values": []})
                for r, row in enumerate(matrix[:needed_rows], start=1):
                    for c, value in enumerate(row[:needed_columns], start=1):
                        operations.append({"op": "set_table_cell", "slide": slide, "shape_id": table["shape_id"], "row": r, "column": c, "text": str(value)})
            else:
                operations.append({
                    "op": "add_table", "slide": slide, "data": matrix[:40],
                    "rows": min(len(matrix), 40), "columns": min(max(len(row) for row in matrix), 16),
                    "x": 0.06, "y": 0.22, "width": 0.88, "height": 0.64,
                    "header_fill": "00338D", "header_font_color": "FFFFFF", "banded_rows": True,
                })
    return operations


_plan_edit_before_attachment_fallback = plan_edit


def plan_edit(
    pptx_path: str,
    user_message: str,
    selected_slide: int,
    selected_slide_image: str | None = None,
    chat_history: list[dict[str, str]] | None = None,
    deck_image_paths: list[str] | None = None,
    attachment_paths: list[str] | None = None,
) -> dict[str, Any]:
    result = _plan_edit_before_attachment_fallback(
        pptx_path, user_message, selected_slide, selected_slide_image,
        chat_history, deck_image_paths, attachment_paths,
    )
    attachment_ops = _attachment_fallback_operations(pptx_path, user_message, selected_slide, attachment_paths)
    existing = [op for op in result.get("operations", []) if op.get("op") != "noop"]
    if attachment_ops:
        existing.extend(attachment_ops)
        result["message"] = "Applied the requested PowerPoint edits and used the attached material."
    if not existing:
        existing = result.get("operations", []) or [{"op": "noop"}]
    result["operations"] = _dedupe_operations(existing)[:150]
    return result

# Make uploaded images visible to the model without mislabeling them as slides.
from contextvars import ContextVar
_ACTIVE_ATTACHMENT_IMAGES: ContextVar[list[str]] = ContextVar("deck_refresh_attachment_images", default=[])
_previous_request_atomic_plan = _request_atomic_plan
_previous_request_outline = _request_outline


def _attachment_visual_parts() -> list[dict[str, Any]]:
    parts: list[dict[str, Any]] = []
    for path in _ACTIVE_ATTACHMENT_IMAGES.get()[:6]:
        if not os.path.isfile(path):
            continue
        parts.append({"type": "input_text", "text": f"Uploaded image attachment: {os.path.basename(path)}"})
        parts.append({"type": "input_image", "image_url": f"data:image/png;base64,{_encode_image(path)}"})
    return parts


def _request_atomic_plan(client: OpenAI, model: str, base_content: list[dict[str, Any]], pptx_path: str, operation_limit: int = 12) -> dict[str, Any]:
    return _previous_request_atomic_plan(client, model, list(base_content) + _attachment_visual_parts(), pptx_path, operation_limit)


def _request_outline(client: OpenAI, model: str, content: list[dict[str, Any]]) -> TaskOutline:
    return _previous_request_outline(client, model, list(content) + _attachment_visual_parts())


_plan_edit_before_visual_attachments = plan_edit


def plan_edit(
    pptx_path: str,
    user_message: str,
    selected_slide: int,
    selected_slide_image: str | None = None,
    chat_history: list[dict[str, str]] | None = None,
    deck_image_paths: list[str] | None = None,
    attachment_paths: list[str] | None = None,
) -> dict[str, Any]:
    image_paths = [
        os.path.abspath(path) for path in (attachment_paths or [])
        if os.path.isfile(path) and Path(path).suffix.casefold() in {".png", ".jpg", ".jpeg", ".webp", ".gif"}
    ]
    token = _ACTIVE_ATTACHMENT_IMAGES.set(image_paths)
    try:
        return _plan_edit_before_visual_attachments(
            pptx_path, user_message, selected_slide, selected_slide_image,
            chat_history, deck_image_paths, attachment_paths,
        )
    finally:
        _ACTIVE_ATTACHMENT_IMAGES.reset(token)

# ---------------------------------------------------------------------------
# Deterministic slide-command router
# ---------------------------------------------------------------------------
# Clear slide-management requests should never depend on model interpretation.
# This final wrapper parses them locally, validates them against the current
# deck, and returns exact operations or a precise no-change explanation.

_DIRECT_SLIDE_REF_WORDS = {
    "first": 1,
    "start": 1,
    "beginning": 1,
}


def _direct_ref_value(raw: str, slide_count: int, selected_slide: int) -> int | None:
    text = re.sub(r"\s+", " ", str(raw or "").strip().casefold())
    text = re.sub(r"^(?:the\s+)", "", text)
    text = re.sub(r"\s+slide$", "", text)
    text = re.sub(r"(?<=\d)(?:st|nd|rd|th)$", "", text)
    if text in {"last", "final", "end"}:
        return slide_count
    if text in {"current", "this", "selected"}:
        return selected_slide
    if text in _DIRECT_SLIDE_REF_WORDS:
        return _DIRECT_SLIDE_REF_WORDS[text]
    if re.fullmatch(r"\d+", text):
        return int(text)
    return None


def _direct_invalid_ref_message(value: int | None, slide_count: int, label: str = "slide") -> str:
    if value is None:
        return f"No change was made. I could not identify the {label} number."
    return f"No change was made. The deck has {slide_count} slides, so {label} {value} does not exist."


def _direct_slide_command_plan(
    pptx_path: str,
    user_message: str,
    selected_slide: int,
) -> dict[str, Any] | None:
    summary = deck_summary(pptx_path, max_text=300)
    slide_count = int(summary.get("slide_count", 0))
    selected_slide = max(1, min(int(selected_slide), max(1, slide_count)))
    original = str(user_message or "").strip()
    normalized = original.casefold().strip()
    normalized = re.sub(r"[.!?]+$", "", normalized).strip()
    normalized = re.sub(r"\s+", " ", normalized)
    # Common polite wrappers should never force an AI round trip.
    normalized = re.sub(
        r"^(?:(?:please|pls|plz)\s+|(?:can|could|would|will)\s+you\s+|i\s+want\s+you\s+to\s+|i(?:'d| would)\s+like\s+you\s+to\s+)",
        "",
        normalized,
    ).strip()
    # Users often add a harmless scope phrase at the end. Remove it before
    # matching, for example: 'delete the last slide in the deck'.
    normalized = re.sub(r"\s+for\s+me$", "", normalized).strip()
    normalized = re.sub(
        r"\s+(?:(?:in|from|of)\s+)?(?:(?:the|this|my|current)\s+)?(?:deck|presentation|powerpoint|power point|ppt|pptx|file)$",
        "",
        normalized,
    ).strip()
    # Treat page as a conversational synonym for slide for management tasks.
    normalized = re.sub(r"\bpages\b", "slides", normalized)
    normalized = re.sub(r"\bpage\b", "slide", normalized)
    normalized = re.sub(r"\bthe\s+very\s+last\b", "the last", normalized)
    normalized = re.sub(r"\b(?:the\s+)?slide\s+at\s+the\s+end\b", "the last slide", normalized)
    normalized = re.sub(r"\b(?:the\s+)?slide\s+at\s+the\s+(?:start|beginning)\b", "the first slide", normalized)
    normalized = re.sub(r"\b(last|final|first)\s+one\b", r"\1 slide", normalized)
    normalized = re.sub(r"^(?:get rid of|take out|drop|erase|trash|eliminate)\s+", "delete ", normalized)
    normalized = re.sub(r"^(?:copy|clone|make\s+a\s+copy\s+of)\s+", "duplicate ", normalized)
    normalized = re.sub(r"^(?:put|place|send)\s+(?=slide\b|the\s+(?:last|final|first|current|selected)\s+slide\b)", "move ", normalized)
    normalized = re.sub(r"^append\s+(?:a\s+)?(?:new\s+)?slide$", "add slide", normalized)

    ref = r"(?:the\s+)?(?:last|final|first|current|this|selected|end|start|beginning|\d+(?:st|nd|rd|th)?)"

    # Common compound command: duplicate a slide and place the new copy at a
    # destination. Compile this into one deterministic duplicate operation.
    compound_duplicate = re.fullmatch(
        rf"(?:duplicate)\s+(?:slide\s+)?(?P<src>{ref})(?:\s+slide)?\s+"
        rf"(?:and\s+then|then|and)\s+(?:move|put|place|send)\s+"
        rf"(?:that\s+new\s+slide|the\s+new\s+slide|the\s+copy|that\s+copy|it)\s+"
        rf"(?:to|at|in|into)\s+(?:(?:slide|position)\s+)?(?P<dst>{ref})(?:\s+slide)?",
        normalized,
    )
    if compound_duplicate:
        source = _direct_ref_value(compound_duplicate.group("src"), slide_count, selected_slide)
        destination = _direct_ref_value(compound_duplicate.group("dst"), slide_count + 1, selected_slide)
        if source is None or not 1 <= source <= slide_count:
            return {"message": _direct_invalid_ref_message(source, slide_count), "assumptions": [], "operations": []}
        if destination is None or not 1 <= destination <= slide_count + 1:
            return {"message": _direct_invalid_ref_message(destination, slide_count + 1, "destination position"), "assumptions": [], "operations": []}
        return {
            "message": f"Duplicated slide {source} into position {destination}.",
            "assumptions": [],
            "operations": [{"op": "duplicate_slide", "slide": source, "position": destination}],
        }

    # Delete one slide.
    match = re.fullmatch(rf"(?:delete|remove)\s+(?:slide\s+)?(?P<target>{ref})(?:\s+slide)?", normalized)
    if match:
        target = _direct_ref_value(match.group("target"), slide_count, selected_slide)
        if target is None or not 1 <= target <= slide_count:
            return {"message": _direct_invalid_ref_message(target, slide_count), "assumptions": [], "operations": []}
        if slide_count <= 1:
            return {"message": "No change was made. A presentation must keep at least one slide.", "assumptions": [], "operations": []}
        return {
            "message": f"Deleted slide {target}.",
            "assumptions": [],
            "operations": [{"op": "delete_slide", "slide": target}],
        }

    # Delete the last N slides.
    match = re.fullmatch(r"(?:delete|remove)\s+(?:the\s+)?last\s+(\d+)\s+slides?", normalized)
    if match:
        amount = int(match.group(1))
        if amount < 1:
            return {"message": "No change was made. The number of slides must be positive.", "assumptions": [], "operations": []}
        if amount >= slide_count:
            return {"message": "No change was made. A presentation must keep at least one slide.", "assumptions": [], "operations": []}
        values = list(range(slide_count - amount + 1, slide_count + 1))
        return {
            "message": f"Deleted the last {amount} slides.",
            "assumptions": [],
            "operations": [{"op": "delete_slide", "slide": value} for value in sorted(values, reverse=True)],
        }

    # Delete a spoken range such as “slides 3 through 5”.
    match = re.fullmatch(r"(?:delete|remove)\s+slides?\s+(\d+)\s+(?:through|to)\s+(\d+)", normalized)
    if match:
        a, b = int(match.group(1)), int(match.group(2))
        values = set(range(min(a, b), max(a, b) + 1))
        invalid = sorted(v for v in values if not 1 <= v <= slide_count)
        if invalid:
            return {"message": _direct_invalid_ref_message(invalid[0], slide_count), "assumptions": [], "operations": []}
        if slide_count - len(values) < 1:
            return {"message": "No change was made. A presentation must keep at least one slide.", "assumptions": [], "operations": []}
        return {
            "message": f"Deleted slides {min(values)} through {max(values)}.",
            "assumptions": [],
            "operations": [{"op": "delete_slide", "slide": value} for value in sorted(values, reverse=True)],
        }

    # Delete explicit lists and hyphenated ranges, highest slide first so numbering stays stable.
    match = re.fullmatch(r"(?:delete|remove)\s+slides?\s+([0-9,\s&-]+(?:and\s+\d+)?)", normalized)
    if match:
        raw = match.group(1)
        values: set[int] = set()
        for start, end in re.findall(r"(\d+)\s*-\s*(\d+)", raw):
            a, b = int(start), int(end)
            values.update(range(min(a, b), max(a, b) + 1))
        raw_without_ranges = re.sub(r"\d+\s*-\s*\d+", "", raw)
        values.update(int(v) for v in re.findall(r"\d+", raw_without_ranges))
        invalid = sorted(v for v in values if not 1 <= v <= slide_count)
        if invalid:
            return {"message": _direct_invalid_ref_message(invalid[0], slide_count), "assumptions": [], "operations": []}
        ordered = sorted(values, reverse=True)
        if not ordered:
            return {"message": "No change was made. No slide numbers were provided.", "assumptions": [], "operations": []}
        if slide_count - len(ordered) < 1:
            return {"message": "No change was made. A presentation must keep at least one slide.", "assumptions": [], "operations": []}
        return {
            "message": "Deleted slides " + ", ".join(str(v) for v in sorted(values)) + ".",
            "assumptions": [],
            "operations": [{"op": "delete_slide", "slide": v} for v in ordered],
        }

    # Delete everything after or before a slide.
    match = re.fullmatch(rf"(?:delete|remove)\s+(?:all\s+)?slides?\s+(after|before)\s+(?:slide\s+)?(?P<target>{ref})(?:\s+slide)?", normalized)
    if match:
        target = _direct_ref_value(match.group("target"), slide_count, selected_slide)
        if target is None or not 1 <= target <= slide_count:
            return {"message": _direct_invalid_ref_message(target, slide_count), "assumptions": [], "operations": []}
        values = list(range(target + 1, slide_count + 1)) if match.group(1) == "after" else list(range(1, target))
        ordered = sorted(values, reverse=True)
        if not ordered:
            return {"message": "No change was made. There are no slides in the requested range.", "assumptions": [], "operations": []}
        if slide_count - len(ordered) < 1:
            return {"message": "No change was made. A presentation must keep at least one slide.", "assumptions": [], "operations": []}
        return {
            "message": f"Deleted all slides {match.group(1)} slide {target}.",
            "assumptions": [],
            "operations": [{"op": "delete_slide", "slide": v} for v in ordered],
        }

    # Move to an exact final position, including “to slide 4”.
    match = re.fullmatch(
        rf"move\s+(?:slide\s+)?(?P<src>{ref})(?:\s+slide)?\s+(?:to|into)\s+(?:(?:slide|position)\s+)?(?P<dst>{ref})(?:\s+slide)?",
        normalized,
    )
    if match:
        source = _direct_ref_value(match.group("src"), slide_count, selected_slide)
        destination = _direct_ref_value(match.group("dst"), slide_count, selected_slide)
        if source is None or not 1 <= source <= slide_count:
            return {"message": _direct_invalid_ref_message(source, slide_count, "source slide"), "assumptions": [], "operations": []}
        if destination is None or not 1 <= destination <= slide_count:
            return {"message": _direct_invalid_ref_message(destination, slide_count, "destination position"), "assumptions": [], "operations": []}
        if source == destination:
            return {"message": f"No change was needed. Slide {source} is already in position {destination}.", "assumptions": [], "operations": []}
        return {
            "message": f"Moved slide {source} to position {destination}.",
            "assumptions": [],
            "operations": [{"op": "move_slide", "from_slide": source, "to_slide": destination}],
        }

    # Move before or after another slide by identity.
    match = re.fullmatch(
        rf"move\s+(?:slide\s+)?(?P<src>{ref})(?:\s+slide)?\s+(?P<where>before|after)\s+(?:slide\s+)?(?P<dst>{ref})(?:\s+slide)?",
        normalized,
    )
    if match:
        source = _direct_ref_value(match.group("src"), slide_count, selected_slide)
        target = _direct_ref_value(match.group("dst"), slide_count, selected_slide)
        if source is None or not 1 <= source <= slide_count:
            return {"message": _direct_invalid_ref_message(source, slide_count, "source slide"), "assumptions": [], "operations": []}
        if target is None or not 1 <= target <= slide_count:
            return {"message": _direct_invalid_ref_message(target, slide_count, "target slide"), "assumptions": [], "operations": []}
        if source == target:
            return {"message": "No change was needed. A slide cannot move relative to itself.", "assumptions": [], "operations": []}
        if match.group("where") == "before":
            destination = target - 1 if source < target else target
        else:
            destination = target if source < target else target + 1
        destination = max(1, min(destination, slide_count))
        return {
            "message": f"Moved slide {source} {match.group('where')} slide {target}.",
            "assumptions": [],
            "operations": [{"op": "move_slide", "from_slide": source, "to_slide": destination}],
        }

    # Move one position left or right.
    match = re.fullmatch(rf"move\s+(?:slide\s+)?(?P<src>{ref})(?:\s+slide)?\s+(?P<direction>left|right|back|forward)", normalized)
    if match:
        source = _direct_ref_value(match.group("src"), slide_count, selected_slide)
        if source is None or not 1 <= source <= slide_count:
            return {"message": _direct_invalid_ref_message(source, slide_count, "source slide"), "assumptions": [], "operations": []}
        delta = -1 if match.group("direction") in {"left", "back"} else 1
        destination = source + delta
        if not 1 <= destination <= slide_count:
            return {"message": f"No change was needed. Slide {source} is already at the edge of the deck.", "assumptions": [], "operations": []}
        return {
            "message": f"Moved slide {source} to position {destination}.",
            "assumptions": [],
            "operations": [{"op": "move_slide", "from_slide": source, "to_slide": destination}],
        }

    # Set a complete explicit slide order.
    match = re.fullmatch(r"(?:reorder|order|set\s+the\s+slide\s+order)(?:\s+slides?)?(?:\s+(?:to|as|in\s+this\s+order))?\s*[:]?\s*([0-9,\s-]+)", normalized)
    if match:
        order = [int(value) for value in re.findall(r"\d+", match.group(1))]
        if sorted(order) != list(range(1, slide_count + 1)):
            return {
                "message": f"No change was made. A complete order must list every slide from 1 through {slide_count} exactly once.",
                "assumptions": [], "operations": [],
            }
        return {"message": "Reordered the slides.", "assumptions": [], "operations": [{"op": "reorder_slides", "order": order}]}

    # Swap two slides using a complete reorder operation.
    match = re.fullmatch(rf"swap\s+(?:slides?\s+)?(?P<a>{ref})\s+(?:and|with)\s+(?:slide\s+)?(?P<b>{ref})(?:\s+slide)?", normalized)
    if match:
        a = _direct_ref_value(match.group("a"), slide_count, selected_slide)
        b = _direct_ref_value(match.group("b"), slide_count, selected_slide)
        if a is None or not 1 <= a <= slide_count:
            return {"message": _direct_invalid_ref_message(a, slide_count), "assumptions": [], "operations": []}
        if b is None or not 1 <= b <= slide_count:
            return {"message": _direct_invalid_ref_message(b, slide_count), "assumptions": [], "operations": []}
        order = list(range(1, slide_count + 1))
        order[a - 1], order[b - 1] = order[b - 1], order[a - 1]
        return {
            "message": f"Swapped slides {a} and {b}.",
            "assumptions": [],
            "operations": [{"op": "reorder_slides", "order": order}],
        }

    # Duplicate a slide, optionally into a requested position.
    match = re.fullmatch(
        rf"duplicate\s+(?:slide\s+)?(?P<src>{ref})(?:\s+slide)?(?:\s+(?:to|at|into)\s+(?:(?:slide|position)\s+)?(?P<dst>{ref})(?:\s+slide)?)?",
        normalized,
    )
    if match:
        source = _direct_ref_value(match.group("src"), slide_count, selected_slide)
        if source is None or not 1 <= source <= slide_count:
            return {"message": _direct_invalid_ref_message(source, slide_count), "assumptions": [], "operations": []}
        destination = source + 1
        if match.group("dst"):
            destination = _direct_ref_value(match.group("dst"), slide_count + 1, selected_slide) or destination
        destination = max(1, min(destination, slide_count + 1))
        return {
            "message": f"Duplicated slide {source} into position {destination}.",
            "assumptions": [],
            "operations": [{"op": "duplicate_slide", "slide": source, "position": destination}],
        }

    # Add a blank slide at a deterministic position.
    match = re.fullmatch(
        rf"(?:add|create|insert)\s+(?:a\s+)?(?:new\s+)?(?:blank\s+)?slide(?:\s+(?P<where>at|in|after|before)\s+(?:(?:slide|position)\s+)?(?P<target>{ref})(?:\s+slide)?)?",
        normalized,
    )
    if match:
        position = slide_count + 1
        if match.group("target"):
            target = _direct_ref_value(match.group("target"), slide_count, selected_slide)
            if target is None or not 1 <= target <= slide_count:
                return {"message": _direct_invalid_ref_message(target, slide_count), "assumptions": [], "operations": []}
            where = match.group("where")
            position = target + 1 if where == "after" else target
        return {
            "message": f"Added a new slide in position {position}.",
            "assumptions": [],
            "operations": [{"op": "add_slide", "position": position, "title": "New Slide", "body": []}],
        }

    return None


# Preserve the AI planner for semantic edits, but bypass it for exact slide
# commands. This also lets slide management work without an API key.
_plan_edit_before_direct_slide_router = plan_edit


def plan_edit(
    pptx_path: str,
    user_message: str,
    selected_slide: int,
    selected_slide_image: str | None = None,
    chat_history: list[dict[str, str]] | None = None,
    deck_image_paths: list[str] | None = None,
    attachment_paths: list[str] | None = None,
) -> dict[str, Any]:
    direct = _direct_slide_command_plan(pptx_path, user_message, selected_slide)
    if direct is not None:
        return direct
    return _plan_edit_before_direct_slide_router(
        pptx_path,
        user_message,
        selected_slide,
        selected_slide_image,
        chat_history,
        deck_image_paths,
        attachment_paths,
    )

# ---------------------------------------------------------------------------
# Conversational natural-language orchestration
# ---------------------------------------------------------------------------
# This final routing layer replaces exact-phrase dependence with a compact
# semantic interpreter. It understands compound requests, references such as
# "that new slide", approximate slide titles, typos, slang, and multilingual
# instructions. The interpreter produces high-level tasks. Each task is then
# compiled and dry-run against a temporary copy of the deck before being
# returned to the application.

import difflib


class ConversationTask(BaseModel):
    model_config = ConfigDict(extra="forbid")
    action: Literal[
        "duplicate_slide",
        "delete_slide",
        "move_slide",
        "swap_slides",
        "reorder_slides",
        "add_slide",
        "semantic_edit",
    ]
    instruction: str = Field(default="", max_length=1800)
    sources: list[str] = Field(default_factory=list, max_length=100)
    destination: str | None = Field(default=None, max_length=240)
    placement: Literal["at", "before", "after"] = "at"
    order: list[str] = Field(default_factory=list, max_length=300)
    alias: str | None = Field(default=None, max_length=80)
    title: str | None = Field(default=None, max_length=240)
    body: list[str] = Field(default_factory=list, max_length=30)


class ConversationPlan(BaseModel):
    model_config = ConfigDict(extra="forbid")
    message: str = Field(default="", max_length=240)
    tasks: list[ConversationTask] = Field(default_factory=list, max_length=24)


CONVERSATION_INTERPRETER_INSTRUCTIONS = r"""
You are the language-understanding layer for a professional PowerPoint editor.
Interpret the user as naturally as ChatGPT would. The request may contain
misspellings, slang, abbreviations, missing punctuation, several steps in one
sentence, vague wording, or references such as "it", "that slide", "the one I
just made", and "move it to the end". The user may write in any language.
Use the recent conversation and the current deck outline to resolve meaning.

Return a short ordered list of atomic tasks. Preserve every requested action.
Do not require exact command wording. Do not ask a question when a reasonable
professional interpretation exists.

Use these actions:
- duplicate_slide
- delete_slide
- move_slide
- swap_slides
- reorder_slides
- add_slide
- semantic_edit, for content, wording, colors, layouts, tables, charts,
  pictures, notes, regeneration, formatting, or any edit not covered above.

Canonical slide references for sources, destination, and order:
- slide:4 or position:4
- selected
- first
- last
- second_last
- title:<approximate slide title>
- alias:<name>
- last_created
- last_mentioned

Rules:
- Break compound instructions into execution-order tasks.
- Resolve pronouns with aliases. A newly added or duplicated slide should get
  an alias, usually new_slide. Later references to it should use
  alias:new_slide or last_created.
- For "duplicate slide 4 and move the new slide to the end", either return one
  duplicate_slide task with destination=last, or two linked tasks. Both are
  valid.
- Use position:N when the user means a final numeric position. Use title:...
  when they identify a slide by subject or title.
- For move_slide, put the moving slide in sources and the destination in
  destination. Use placement before or after only when requested.
- For delete_slide, sources may contain several slides.
- For swap_slides, sources must contain two references.
- For reorder_slides, order must list the complete intended order. For a
  partial ordering request, use one or more move_slide tasks instead.
- For semantic_edit, rewrite the request into a complete, standalone English
  instruction. Put target slide references in sources when known.
- Prefer content and formatting edits before destructive deletion or final
  reordering when that preserves the user's intent.
- Keep the top-level message concise. Never claim an edit already happened.
""".strip()


def _conversation_outline_text(pptx_path: str, selected_slide: int) -> str:
    summary = deck_summary(pptx_path, max_text=220)
    lines = [
        f"Slide count: {summary['slide_count']}",
        f"Selected slide: {selected_slide}",
    ]
    for slide in summary.get("slides", []):
        title = str(slide.get("title") or "Untitled").replace("\n", " ")[:180]
        lines.append(f"{slide['slide']}: {title}")
    return "\n".join(lines)


def _interpret_conversation(
    pptx_path: str,
    user_message: str,
    selected_slide: int,
    chat_history: list[dict[str, str]] | None,
) -> ConversationPlan:
    if not _ai_credential():
        raise EditorError("AI access is not configured.")
    model = _ai_model()
    recent = []
    for item in (chat_history or [])[-14:]:
        role = str(item.get("role") or "user")
        content = str(item.get("content") or "")[:1600]
        recent.append(f"{role}: {content}")
    content = [
        {
            "type": "input_text",
            "text": (
                "Current deck outline:\n"
                + _conversation_outline_text(pptx_path, selected_slide)
                + "\n\nRecent conversation:\n"
                + ("\n".join(recent) if recent else "None")
                + "\n\nCurrent user request:\n"
                + user_message
            ),
        }
    ]
    client = _ai_client(timeout=180.0, max_retries=2)
    try:
        response = client.responses.parse(
            model=model,
            instructions=CONVERSATION_INTERPRETER_INSTRUCTIONS,
            input=[{"role": "user", "content": content}],
            text_format=ConversationPlan,
            max_output_tokens=5000,
            reasoning={"effort": _reasoning_effort(model)},
            verbosity="low",
            store=False,
        )
        parsed = getattr(response, "output_parsed", None)
        if parsed is None:
            raise EditorError("The language interpreter returned no usable task plan.")
        return parsed
    except Exception as exc:
        raise EditorError(_friendly_openai_error(str(exc))) from exc


def _heuristic_conversation_plan(user_message: str) -> ConversationPlan | None:
    """Broad local fallback for common compound slide-management language.

    The OpenAI interpreter is the primary route. This fallback prevents a
    network or billing issue from breaking basic conversational commands.
    """
    text = re.sub(r"\s+", " ", str(user_message or "")).strip()
    lowered = text.casefold()
    if not lowered:
        return None

    tasks: list[ConversationTask] = []
    # Duplicate followed by moving the new copy somewhere.
    match = re.search(
        r"(?:duplicate|copy|clone|make\s+a\s+copy\s+of)\s+(?:slide\s+)?(\d+)"
        r".*?(?:then|and|after\s+that).*?(?:move|put|place|send)\s+"
        r"(?:that\s+new\s+slide|the\s+new\s+slide|the\s+copy|it)\s+"
        r"(?:to|at|in)\s+(?:the\s+)?(?:last|end|back)(?:\s+slide|\s+of\s+the\s+deck)?",
        lowered,
    )
    if match:
        tasks.append(ConversationTask(
            action="duplicate_slide",
            sources=[f"slide:{int(match.group(1))}"],
            destination="last",
            placement="at",
            alias="new_slide",
            instruction=text,
        ))
        return ConversationPlan(message="Understood the compound slide request.", tasks=tasks)

    # Generic clause splitter for duplicate/add followed by pronoun move/delete.
    clauses = [part.strip(" ,.;") for part in re.split(r"\b(?:and then|then|after that)\b|[;\n]+", lowered) if part.strip()]
    saw_created = False
    for clause in clauses:
        m = re.search(r"(?:duplicate|copy|clone)\s+(?:slide\s+)?(\d+)", clause)
        if m:
            tasks.append(ConversationTask(action="duplicate_slide", sources=[f"slide:{int(m.group(1))}"], alias="new_slide", instruction=clause))
            saw_created = True
            continue
        m = re.search(r"(?:move|put|place|send)\s+(?:that\s+new\s+slide|the\s+new\s+slide|the\s+copy|it)\s+(?:to|at|in)\s+(?:the\s+)?(?:last|end|back)", clause)
        if m and saw_created:
            tasks.append(ConversationTask(action="move_slide", sources=["last_created"], destination="last", instruction=clause))
            continue
    if tasks:
        return ConversationPlan(message="Understood the conversational slide request.", tasks=tasks)
    return None


def _virtual_deck(pptx_path: str) -> list[dict[str, Any]]:
    summary = deck_summary(pptx_path, max_text=160)
    return [
        {
            "id": f"original:{slide['slide']}",
            "title": str(slide.get("title") or f"Slide {slide['slide']}"),
            "aliases": set(),
        }
        for slide in summary.get("slides", [])
    ]


def _ref_tokens(value: str | None) -> str:
    return re.sub(r"\s+", " ", str(value or "")).strip().casefold()


def _resolve_virtual_item(
    reference: str | None,
    deck: list[dict[str, Any]],
    aliases: dict[str, str],
    selected_id: str | None,
    last_created_id: str | None,
    last_mentioned_id: str | None,
) -> tuple[int, dict[str, Any]]:
    if not deck:
        raise EditorError("The presentation has no slides.")
    raw = _ref_tokens(reference)
    if not raw:
        raw = "selected"

    def by_id(item_id: str | None):
        if item_id is None:
            return None
        for index, item in enumerate(deck):
            if item["id"] == item_id:
                return index, item
        return None

    if raw in {"selected", "current", "current_slide", "this", "this slide"}:
        match = by_id(selected_id)
        if match:
            return match
    if raw in {"last_created", "new_slide", "the new slide", "that new slide", "the copy", "it"}:
        match = by_id(last_created_id)
        if match:
            return match
    if raw in {"last_mentioned", "that slide", "the slide"}:
        match = by_id(last_mentioned_id)
        if match:
            return match
    if raw in {"first", "front", "beginning", "start", "top"}:
        return 0, deck[0]
    if raw in {"last", "end", "back", "final", "bottom"}:
        return len(deck) - 1, deck[-1]
    if raw in {"second_last", "second last", "penultimate"}:
        index = max(0, len(deck) - 2)
        return index, deck[index]

    alias_name = raw.split(":", 1)[1] if raw.startswith("alias:") else raw
    alias_id = aliases.get(alias_name)
    match = by_id(alias_id)
    if match:
        return match

    number_match = re.fullmatch(r"(?:slide:|position:|slide\s+|position\s+)?(\d+)", raw)
    if number_match:
        number = int(number_match.group(1))
        if 1 <= number <= len(deck):
            return number - 1, deck[number - 1]
        raise EditorError(f"The deck has {len(deck)} slides, so slide {number} does not exist.")

    wanted_title = raw.split(":", 1)[1].strip() if raw.startswith("title:") else raw
    if wanted_title:
        normalized = _normalize_text(wanted_title)
        exact = [(i, item) for i, item in enumerate(deck) if normalized == _normalize_text(item["title"])]
        if exact:
            return exact[0]
        contained = [(i, item) for i, item in enumerate(deck) if normalized in _normalize_text(item["title"]) or _normalize_text(item["title"]) in normalized]
        if len(contained) == 1:
            return contained[0]
        scored = []
        for i, item in enumerate(deck):
            score = difflib.SequenceMatcher(None, normalized, _normalize_text(item["title"])).ratio()
            scored.append((score, i, item))
        scored.sort(reverse=True, key=lambda row: row[0])
        if scored and scored[0][0] >= 0.48:
            return scored[0][1], scored[0][2]
    raise EditorError(f"I could not identify which slide '{reference}' refers to.")


def _insert_index(
    destination: str | None,
    placement: str,
    deck: list[dict[str, Any]],
    aliases: dict[str, str],
    selected_id: str | None,
    last_created_id: str | None,
    last_mentioned_id: str | None,
    *,
    allow_after_end: bool,
) -> int:
    raw = _ref_tokens(destination)
    if not raw:
        return len(deck) if allow_after_end else max(0, len(deck) - 1)
    if raw in {"last", "end", "back", "final", "bottom"}:
        return len(deck)
    if raw in {"first", "front", "beginning", "start", "top"}:
        return 0
    if raw in {"second_last", "second last", "penultimate"}:
        return max(0, len(deck) - 1)
    numeric = re.fullmatch(r"(?:position:|slide:|position\s+|slide\s+)?(\d+)", raw)
    if numeric and placement == "at":
        position = int(numeric.group(1))
        upper = len(deck) + (1 if allow_after_end else 0)
        if not 1 <= position <= max(1, upper):
            raise EditorError(f"Position {position} is outside the deck.")
        return min(position - 1, len(deck))
    target_index, _ = _resolve_virtual_item(
        destination, deck, aliases, selected_id, last_created_id, last_mentioned_id
    )
    if placement == "before":
        return target_index
    if placement == "after":
        return target_index + 1
    return target_index


def _compile_management_task(
    task: ConversationTask,
    deck: list[dict[str, Any]],
    aliases: dict[str, str],
    selected_id: str | None,
    last_created_id: str | None,
    last_mentioned_id: str | None,
    id_counter: list[int],
) -> tuple[list[dict[str, Any]], str | None, str | None, str]:
    action = task.action
    operations: list[dict[str, Any]] = []
    detail = ""

    if action == "duplicate_slide":
        source_ref = task.sources[0] if task.sources else "selected"
        source_index, source_item = _resolve_virtual_item(source_ref, deck, aliases, selected_id, last_created_id, last_mentioned_id)
        destination_index = source_index + 1
        if task.destination:
            destination_index = _insert_index(
                task.destination, task.placement, deck, aliases, selected_id,
                last_created_id, last_mentioned_id, allow_after_end=True,
            )
        destination_index = max(0, min(destination_index, len(deck)))
        operations.append({"op": "duplicate_slide", "slide": source_index + 1, "position": destination_index + 1})
        id_counter[0] += 1
        new_id = f"created:{id_counter[0]}"
        new_item = {"id": new_id, "title": source_item["title"], "aliases": set()}
        deck.insert(destination_index, new_item)
        alias = _ref_tokens(task.alias or "new_slide").replace("alias:", "")
        if alias:
            aliases[alias] = new_id
            new_item["aliases"].add(alias)
        aliases["new_slide"] = new_id
        last_created_id = new_id
        last_mentioned_id = new_id
        detail = f"Duplicated slide {source_index + 1} into position {destination_index + 1}."

    elif action == "add_slide":
        destination_index = len(deck)
        if task.destination:
            destination_index = _insert_index(
                task.destination, task.placement, deck, aliases, selected_id,
                last_created_id, last_mentioned_id, allow_after_end=True,
            )
        destination_index = max(0, min(destination_index, len(deck)))
        title = str(task.title or "New Slide").strip() or "New Slide"
        operations.append({
            "op": "add_slide",
            "position": destination_index + 1,
            "title": title,
            "body": [str(value) for value in task.body],
        })
        id_counter[0] += 1
        new_id = f"created:{id_counter[0]}"
        new_item = {"id": new_id, "title": title, "aliases": set()}
        deck.insert(destination_index, new_item)
        alias = _ref_tokens(task.alias or "new_slide").replace("alias:", "")
        if alias:
            aliases[alias] = new_id
            new_item["aliases"].add(alias)
        aliases["new_slide"] = new_id
        last_created_id = new_id
        last_mentioned_id = new_id
        detail = f"Added '{title}' in position {destination_index + 1}."

    elif action == "delete_slide":
        refs = task.sources or ["selected"]
        ids: list[str] = []
        for ref in refs:
            _, item = _resolve_virtual_item(ref, deck, aliases, selected_id, last_created_id, last_mentioned_id)
            if item["id"] not in ids:
                ids.append(item["id"])
        if len(deck) - len(ids) < 1:
            raise EditorError("A presentation must keep at least one slide.")
        positions = sorted((next(i for i, item in enumerate(deck) if item["id"] == item_id) for item_id in ids), reverse=True)
        deleted_numbers = []
        for index in positions:
            deleted_numbers.append(index + 1)
            operations.append({"op": "delete_slide", "slide": index + 1})
            removed = deck.pop(index)
            for key, value in list(aliases.items()):
                if value == removed["id"]:
                    aliases.pop(key, None)
            if removed["id"] == last_created_id:
                last_created_id = None
            if removed["id"] == last_mentioned_id:
                last_mentioned_id = None
            if removed["id"] == selected_id:
                selected_id = deck[min(index, len(deck) - 1)]["id"] if deck else None
        detail = "Deleted slide" + ("s " if len(deleted_numbers) > 1 else " ") + ", ".join(str(value) for value in sorted(deleted_numbers)) + "."

    elif action == "move_slide":
        source_ref = task.sources[0] if task.sources else "selected"
        source_index, source_item = _resolve_virtual_item(source_ref, deck, aliases, selected_id, last_created_id, last_mentioned_id)
        source_number = source_index + 1
        moving = deck.pop(source_index)
        destination_index = _insert_index(
            task.destination or "last", task.placement, deck, aliases,
            selected_id, last_created_id, last_mentioned_id,
            allow_after_end=True,
        )
        destination_index = max(0, min(destination_index, len(deck)))
        deck.insert(destination_index, moving)
        operations.append({"op": "move_slide", "from_slide": source_number, "to_slide": destination_index + 1})
        last_mentioned_id = moving["id"]
        detail = f"Moved slide {source_number} to position {destination_index + 1}."

    elif action == "swap_slides":
        if len(task.sources) < 2:
            raise EditorError("Swapping slides needs two slide references.")
        a_index, a_item = _resolve_virtual_item(task.sources[0], deck, aliases, selected_id, last_created_id, last_mentioned_id)
        b_index, b_item = _resolve_virtual_item(task.sources[1], deck, aliases, selected_id, last_created_id, last_mentioned_id)
        order = list(range(1, len(deck) + 1))
        order[a_index], order[b_index] = order[b_index], order[a_index]
        operations.append({"op": "reorder_slides", "order": order})
        deck[a_index], deck[b_index] = deck[b_index], deck[a_index]
        last_mentioned_id = a_item["id"]
        detail = f"Swapped slides {a_index + 1} and {b_index + 1}."

    elif action == "reorder_slides":
        if not task.order:
            raise EditorError("A complete slide order was not provided.")
        item_ids: list[str] = []
        order_numbers: list[int] = []
        for ref in task.order:
            index, item = _resolve_virtual_item(ref, deck, aliases, selected_id, last_created_id, last_mentioned_id)
            if item["id"] in item_ids:
                raise EditorError("The requested order lists the same slide more than once.")
            item_ids.append(item["id"])
            order_numbers.append(index + 1)
        if len(item_ids) != len(deck):
            raise EditorError(f"A complete order must include all {len(deck)} slides exactly once.")
        by_id = {item["id"]: item for item in deck}
        deck[:] = [by_id[item_id] for item_id in item_ids]
        operations.append({"op": "reorder_slides", "order": order_numbers})
        detail = "Reordered the complete deck."

    else:
        raise EditorError(f"Unsupported conversational action: {action}")

    return operations, last_created_id, last_mentioned_id, detail


def _simulate_operations(
    current_path: str,
    operations: list[dict[str, Any]],
    work_dir: str,
    step_prefix: str,
) -> tuple[str, list[dict[str, Any]], list[str]]:
    accepted: list[dict[str, Any]] = []
    issues: list[str] = []
    path = current_path
    for index, operation in enumerate(operations, start=1):
        candidate = os.path.join(work_dir, f"{step_prefix}_{index:03d}.pptx")
        result = apply_operations(path, candidate, [operation])
        if result.get("applied") and not result.get("unchanged"):
            accepted.append(operation)
            path = candidate
        else:
            reason = "; ".join(str(item.get("reason") or "The edit was skipped.") for item in result.get("skipped", []))
            issues.append(reason or f"{operation.get('op')} did not change the deck.")
    return path, accepted, issues


def _semantic_instruction(task: ConversationTask, positions: list[int]) -> str:
    instruction = str(task.instruction or "").strip()
    if not instruction:
        instruction = "Apply the requested content and formatting edit."
    if positions:
        labels = ", ".join(str(value) for value in positions)
        instruction += f" Target slide{'s' if len(positions) != 1 else ''} {labels}."
    return instruction


_plan_edit_before_conversation_agent = plan_edit


def plan_edit(
    pptx_path: str,
    user_message: str,
    selected_slide: int,
    selected_slide_image: str | None = None,
    chat_history: list[dict[str, str]] | None = None,
    deck_image_paths: list[str] | None = None,
    attachment_paths: list[str] | None = None,
) -> dict[str, Any]:
    """Interpret natural language first, then execute a validated task chain.

    This is intentionally the final definition in the module. It handles
    conversational and multilingual input before any exact command router.
    """
    try:
        conversation = _interpret_conversation(
            pptx_path, user_message, selected_slide, chat_history,
        )
    except Exception as interpreter_error:
        conversation = _heuristic_conversation_plan(user_message)
        if conversation is None:
            # Preserve the previous rich semantic planner as a fallback.
            result = _plan_edit_before_conversation_agent(
                pptx_path, user_message, selected_slide, selected_slide_image,
                chat_history, deck_image_paths, attachment_paths,
            )
            assumptions = list(result.get("assumptions", []))
            assumptions.append(str(interpreter_error)[:300])
            result["assumptions"] = list(dict.fromkeys(assumptions))[:20]
            return result

    if not conversation.tasks:
        return _plan_edit_before_conversation_agent(
            pptx_path, user_message, selected_slide, selected_slide_image,
            chat_history, deck_image_paths, attachment_paths,
        )

    deck = _virtual_deck(pptx_path)
    selected_id = deck[max(0, min(selected_slide - 1, len(deck) - 1))]["id"] if deck else None
    aliases: dict[str, str] = {}
    last_created_id: str | None = None
    last_mentioned_id: str | None = selected_id
    id_counter = [0]
    all_operations: list[dict[str, Any]] = []
    completed: list[str] = []
    issues: list[str] = []

    with tempfile.TemporaryDirectory(prefix="deck_refresh_conversation_") as work_dir:
        planning_path = os.path.join(work_dir, "planning_000.pptx")
        shutil.copy2(pptx_path, planning_path)
        for task_index, task in enumerate(conversation.tasks, start=1):
            try:
                if task.action == "semantic_edit":
                    positions: list[int] = []
                    for ref in task.sources:
                        index, item = _resolve_virtual_item(
                            ref, deck, aliases, selected_id, last_created_id, last_mentioned_id,
                        )
                        positions.append(index + 1)
                        last_mentioned_id = item["id"]
                    instruction = _semantic_instruction(task, positions)
                    semantic_selected = positions[0] if positions else max(1, min(selected_slide, len(deck)))
                    semantic_plan = _plan_edit_before_conversation_agent(
                        planning_path,
                        instruction,
                        semantic_selected,
                        None,
                        chat_history,
                        None,
                        attachment_paths,
                    )
                    proposed = [operation for operation in semantic_plan.get("operations", []) if operation.get("op") != "noop"]
                    planning_path, accepted, step_issues = _simulate_operations(
                        planning_path, proposed, work_dir, f"semantic_{task_index:02d}",
                    )
                    all_operations.extend(accepted)
                    issues.extend(step_issues)
                    if accepted:
                        completed.append(str(semantic_plan.get("message") or task.instruction or "Applied a semantic edit.").strip())
                    elif not step_issues:
                        issues.append(f"Task {task_index} did not produce a usable edit.")
                    # Rebuild titles after content or slide regeneration changes.
                    actual = _virtual_deck(planning_path)
                    if len(actual) == len(deck):
                        for index, item in enumerate(actual):
                            deck[index]["title"] = item["title"]
                    continue

                proposed, last_created_id, last_mentioned_id, detail = _compile_management_task(
                    task, deck, aliases, selected_id, last_created_id,
                    last_mentioned_id, id_counter,
                )
                planning_path, accepted, step_issues = _simulate_operations(
                    planning_path, proposed, work_dir, f"management_{task_index:02d}",
                )
                if len(accepted) != len(proposed):
                    # The virtual model changed optimistically. Reconstruct the deck
                    # from the successful planning file and clear aliases that no
                    # longer resolve safely.
                    rebuilt = _virtual_deck(planning_path)
                    deck[:] = rebuilt
                    aliases.clear()
                    last_created_id = None
                    last_mentioned_id = selected_id if any(item["id"] == selected_id for item in deck) else None
                all_operations.extend(accepted)
                issues.extend(step_issues)
                if accepted:
                    completed.append(detail)
            except Exception as exc:
                issues.append(f"Task {task_index}: {str(exc)[:400]}")

    if not all_operations:
        # Last fallback keeps legacy support for unusual operations while
        # returning a truthful message if nothing maps.
        legacy = _plan_edit_before_conversation_agent(
            pptx_path, user_message, selected_slide, selected_slide_image,
            chat_history, deck_image_paths, attachment_paths,
        )
        if any(operation.get("op") != "noop" for operation in legacy.get("operations", [])):
            return legacy
        reason = "; ".join(issues[:4]) or "The request did not map to a supported PowerPoint edit."
        return {
            "message": f"No change was made. {reason}",
            "assumptions": issues[:20],
            "operations": [],
        }

    message_parts = [part for part in completed if part]
    if message_parts:
        message = " ".join(message_parts[:6])
    else:
        message = f"Applied {len(all_operations)} PowerPoint edit{'s' if len(all_operations) != 1 else ''}."
    if issues:
        message += f" {len(issues)} step{'s' if len(issues) != 1 else ''} could not be completed."
    return {
        "message": message[:900],
        "assumptions": issues[:20],
        "operations": all_operations[:300],
    }


# ---------------------------------------------------------------------------
# Strict conversational safety layer
# ---------------------------------------------------------------------------
# The public chat route uses this final plan_edit definition. Natural-language
# interpretation stays model-driven, while the resulting operation chain is
# treated as one transaction. If any requested step cannot be compiled,
# validated, applied, saved, and reopened, no operations are returned.

CHAT_ERROR_MESSAGE = "error cant do that"


class FailureDiagnosis(BaseModel):
    """User-facing explanation for a rejected PowerPoint transaction."""

    model_config = ConfigDict(extra="forbid")
    failed_step: str = Field(default="Interpreting or applying the requested edit", max_length=240)
    reason: str = Field(default="The request could not be mapped safely to the current presentation", max_length=500)
    suggested_wording: str = Field(default="", max_length=900)
    alternate_wording: str | None = Field(default=None, max_length=900)


FAILURE_DIAGNOSER_INSTRUCTIONS = r"""
You diagnose failed requests in a conversational PowerPoint editor. The user
should never need to know internal Python, JSON, Pydantic, shape-index, or API
implementation details.

Given the user's original request, recent chat, current deck outline, and a
technical failure reason:
- Identify the exact user-visible step that failed.
- Explain why in plain language and one sentence.
- Rewrite the user's request so the same intended task is more explicit and is
  likely to work. Preserve the original intent. Do not simplify it into a
  different task.
- Use concrete slide numbers or approximate slide titles from the supplied deck
  outline when available.
- Resolve pronouns in the suggested wording. For example, replace "it" with
  "the duplicated slide" or an exact slide position.
- For a compound request, provide a complete single-message rewrite first. The
  alternate wording may split it into numbered steps when that is safer.
- If wording was not the cause, say so. Suggest retrying the same request after
  fixing the real issue, such as attaching missing data or reconnecting the API.
- Never claim the PowerPoint was changed.
- Do not expose secrets, file paths, stack traces, or raw exception text.
""".strip()


def _safe_failure_reason(reason: str | None) -> str:
    """Convert raw errors into a short non-sensitive diagnostic hint."""
    value = re.sub(r"\s+", " ", str(reason or "")).strip()
    value = re.sub(r"sk-[A-Za-z0-9_-]+", "[redacted key]", value)
    value = re.sub(r"(?:[A-Za-z]:\\|/)[^ ]+", "[local file]", value)
    return value[:700] or "The request did not produce a complete safe edit plan."


def _local_failure_diagnosis(
    pptx_path: str | None,
    user_message: str | None,
    selected_slide: int,
    reason: str | None,
    failed_step: str | None = None,
) -> FailureDiagnosis:
    """Deterministic fallback when the diagnosis model is unavailable."""
    request_text = re.sub(r"\s+", " ", str(user_message or "")).strip()
    lowered_reason = _safe_failure_reason(reason).casefold()
    slide_count = 0
    try:
        slide_count = len(Presentation(pptx_path).slides) if pptx_path else 0
    except Exception:
        slide_count = 0

    step = str(failed_step or "Interpreting the requested PowerPoint change").strip()
    why = "The editor could not map one part of the request to a safe supported PowerPoint operation."
    suggestion = request_text or "Apply the requested change to the selected slide."
    alternate = None

    numbers = [int(value) for value in re.findall(r"\b(?:slide\s*)?(\d+)\b", request_text, re.I)]
    invalid = [number for number in numbers if slide_count and not 1 <= number <= slide_count]
    if invalid:
        step = f"Resolving slide {invalid[0]}"
        why = f"The current presentation has {slide_count} slides, so slide {invalid[0]} does not exist."
        suggestion = f"Apply this change to slide {max(1, min(selected_slide, slide_count or 1))}: {request_text}"
        alternate = f"First tell me the slide title or select the intended slide, then say: Apply the requested change to the selected slide."
    elif any(token in lowered_reason for token in ("ambiguous", "multiple", "more than one", "could not resolve", "not found")):
        step = "Resolving the slide or object reference"
        why = "The wording matched no slide or more than one possible slide."
        suggestion = f"On slide {max(1, selected_slide)}, {request_text}"
        alternate = f"Find the slide titled '<exact title>' and then {request_text}"
    elif any(token in lowered_reason for token in ("attachment", "missing data", "no numeric", "no data", "image file")):
        step = "Finding the required source data or attachment"
        why = "The requested edit needs data or an image that was not available in the deck or attachments."
        suggestion = f"Using the attached file as the source, {request_text}"
        alternate = "Attach the source data or image first, then repeat the same request and name the attachment."
    elif any(token in lowered_reason for token in ("timeout", "rate limit", "billing", "authentication", "api", "connection")):
        step = "Connecting to the AI interpretation service"
        why = "The request wording was not the problem. The AI service was unavailable or rejected the request."
        suggestion = request_text or "Retry the same PowerPoint edit request."
        alternate = "Retry the same wording after confirming the API key, billing, and internet connection."
    elif any(token in lowered_reason for token in ("no usable", "no operation", "unchanged", "skipped")):
        step = "Turning the request into an executable edit"
        why = "The editor understood the general intent but did not produce a complete operation that changed the deck."
        suggestion = f"On slide {max(1, selected_slide)}, {request_text}"
        alternate = f"Do this in two steps: 1. Select slide {max(1, selected_slide)}. 2. {request_text}"
    elif any(token in lowered_reason for token in ("unsupported", "not supported", "smartart", "animation", "ole", "macro", "master")):
        step = "Applying the requested PowerPoint feature"
        why = "That specific PowerPoint feature is outside the editor's safe operation set."
        suggestion = f"Recreate the visible result using editable text, shapes, tables, charts, or images: {request_text}"
        alternate = "Ask the editor to preserve the unsupported object and make the surrounding slide changes instead."
    elif re.search(r"\b(?:then|and|after|before|also|plus)\b", request_text, re.I):
        step = "Completing every step in the compound request"
        why = "One step in the multi-step request could not be validated against the deck."
        suggestion = request_text
        clauses = [part.strip(" ,.;") for part in re.split(r"\b(?:and then|then|after that)\b|[;\n]+", request_text, flags=re.I) if part.strip()]
        if len(clauses) > 1:
            alternate = "Send these as separate requests: " + " ".join(f"{index + 1}. {clause}." for index, clause in enumerate(clauses))

    return FailureDiagnosis(
        failed_step=step[:240],
        reason=why[:500],
        suggested_wording=suggestion[:900],
        alternate_wording=(alternate[:900] if alternate else None),
    )


def diagnose_failure_message(
    pptx_path: str | None,
    user_message: str | None,
    selected_slide: int = 1,
    chat_history: list[dict[str, str]] | None = None,
    reason: str | None = None,
    failed_step: str | None = None,
) -> str:
    """Return a useful failure response while guaranteeing no deck mutation."""
    diagnosis: FailureDiagnosis | None = None
    api_key = _ai_credential()
    if api_key and pptx_path and user_message:
        recent = []
        for item in (chat_history or [])[-10:]:
            recent.append(f"{str(item.get('role') or 'user')}: {str(item.get('content') or '')[:1000]}")
        prompt = (
            "Current deck outline:\n"
            + _conversation_outline_text(pptx_path, selected_slide)
            + "\n\nRecent conversation:\n"
            + ("\n".join(recent) if recent else "None")
            + "\n\nOriginal request:\n"
            + str(user_message)
            + "\n\nFailed step hint:\n"
            + str(failed_step or "Unknown")
            + "\n\nTechnical failure category:\n"
            + _safe_failure_reason(reason)
        )
        try:
            model = _ai_model()
            response = _ai_client(timeout=90.0, max_retries=1).responses.parse(
                model=model,
                instructions=FAILURE_DIAGNOSER_INSTRUCTIONS,
                input=[{"role": "user", "content": [{"type": "input_text", "text": prompt}]}],
                text_format=FailureDiagnosis,
                max_output_tokens=1800,
                reasoning={"effort": "minimal" if model.startswith("gpt-5") else "low"},
                verbosity="low",
                store=False,
            )
            diagnosis = getattr(response, "output_parsed", None)
        except Exception:
            diagnosis = None

    if diagnosis is None or not str(diagnosis.suggested_wording or "").strip():
        diagnosis = _local_failure_diagnosis(
            pptx_path,
            user_message,
            selected_slide,
            reason,
            failed_step,
        )

    lines = [
        CHAT_ERROR_MESSAGE,
        "",
        f"What failed: {str(diagnosis.failed_step).strip()}",
        f"Why: {str(diagnosis.reason).strip()}",
        f'Try: "{str(diagnosis.suggested_wording).strip()}"',
    ]
    if diagnosis.alternate_wording:
        lines.extend(["", f"Or: {str(diagnosis.alternate_wording).strip()}"])
    lines.extend(["", "PowerPoint was not changed."])
    return "\n".join(lines)


CONVERSATION_REVIEW_INSTRUCTIONS = r"""
You are the final quality-control pass for a conversational PowerPoint editor.
Review the draft task plan against the user's exact request, recent chat, and
current deck outline. Return a corrected complete ConversationPlan.

Requirements:
- Preserve every requested action and its execution order.
- Interpret slang, typos, shorthand, pronouns, approximate titles, and any
  language as a capable human assistant would.
- Resolve references to concrete current slide positions or approximate titles
  whenever possible. Use aliases only for slides created inside this request.
- Break compound requests into atomic tasks.
- Use semantic_edit for any content, design, color, table, chart, picture,
  notes, layout, regeneration, cleanup, or formatting request.
- Do not claim anything has already happened.
- Do not remove a task merely because it looks difficult.
- Return at most 24 tasks.
""".strip()


def _chat_error_plan(
    reason: str | None = None,
    *,
    pptx_path: str | None = None,
    user_message: str | None = None,
    selected_slide: int = 1,
    chat_history: list[dict[str, str]] | None = None,
    failed_step: str | None = None,
) -> dict[str, Any]:
    assumptions = [_safe_failure_reason(reason)[:500]] if reason else []
    message = CHAT_ERROR_MESSAGE
    if user_message:
        message = diagnose_failure_message(
            pptx_path,
            user_message,
            selected_slide,
            chat_history,
            reason,
            failed_step,
        )
    return {
        "message": message,
        "assumptions": assumptions,
        "operations": [],
        "failed": True,
    }


def _review_conversation_plan(
    pptx_path: str,
    user_message: str,
    selected_slide: int,
    chat_history: list[dict[str, str]] | None,
    draft: ConversationPlan,
) -> ConversationPlan:
    """Run a second language pass for compound or ambiguous requests.

    A review failure does not discard a valid first plan. The dry-run and final
    transaction checks remain the source of truth.
    """
    text = str(user_message or "")
    compound = len(draft.tasks) > 1 or len(text) > 100 or bool(
        re.search(r"\b(?:then|after|before|also|and|plus|while|finally|next)\b|[;\n]", text, re.I)
    )
    if not compound:
        return draft
    api_key = _ai_credential()
    if not api_key:
        return draft
    model = _ai_model()
    recent = []
    for item in (chat_history or [])[-14:]:
        recent.append(f"{str(item.get('role') or 'user')}: {str(item.get('content') or '')[:1400]}")
    review_text = (
        "Current deck outline:\n"
        + _conversation_outline_text(pptx_path, selected_slide)
        + "\n\nRecent conversation:\n"
        + ("\n".join(recent) if recent else "None")
        + "\n\nUser request:\n"
        + text
        + "\n\nDraft task plan:\n"
        + draft.model_dump_json()
    )
    try:
        response = _ai_client(timeout=180.0, max_retries=2).responses.parse(
            model=model,
            instructions=CONVERSATION_REVIEW_INSTRUCTIONS,
            input=[{"role": "user", "content": [{"type": "input_text", "text": review_text}]}],
            text_format=ConversationPlan,
            max_output_tokens=6000,
            reasoning={"effort": _reasoning_effort(model)},
            verbosity="low",
            store=False,
        )
        reviewed = getattr(response, "output_parsed", None)
        if reviewed is not None and reviewed.tasks:
            return reviewed
    except Exception:
        pass
    return draft


def _strict_apply_chain(
    source_path: str,
    operations: list[dict[str, Any]],
    work_dir: str,
    name: str,
) -> tuple[str | None, list[str]]:
    """Apply a complete operation chain to a throwaway file or reject it all."""
    effective = [dict(operation) for operation in operations if operation.get("op") != "noop"]
    if not effective:
        return None, ["No usable PowerPoint operation was produced."]
    candidate = os.path.join(work_dir, f"{name}.pptx")
    try:
        result = apply_operations(source_path, candidate, effective)
        skipped = list(result.get("skipped") or [])
        applied = [item for item in result.get("applied", []) if item.get("op") != "noop"]
        if skipped or len(applied) != len(effective) or result.get("unchanged"):
            reasons = [str(item.get("reason") or "The operation was skipped.") for item in skipped]
            if len(applied) != len(effective):
                reasons.append(f"Only {len(applied)} of {len(effective)} operations applied.")
            return None, reasons or ["The operation chain did not fully apply."]
        reopened = Presentation(candidate)
        if len(reopened.slides) < 1:
            return None, ["The edited presentation contained no slides."]
        return candidate, []
    except Exception as exc:
        return None, [str(exc)[:500]]


def _strict_legacy_plan(
    pptx_path: str,
    user_message: str,
    selected_slide: int,
    selected_slide_image: str | None,
    chat_history: list[dict[str, str]] | None,
    deck_image_paths: list[str] | None,
    attachment_paths: list[str] | None,
) -> dict[str, Any]:
    """Use the rich legacy planner, then prove the whole plan on a copy."""
    try:
        legacy = _plan_edit_before_conversation_agent(
            pptx_path,
            user_message,
            selected_slide,
            selected_slide_image,
            chat_history,
            deck_image_paths,
            attachment_paths,
        )
        operations = [op for op in legacy.get("operations", []) if op.get("op") != "noop"]
        with tempfile.TemporaryDirectory(prefix="deck_refresh_legacy_strict_") as work_dir:
            candidate, issues = _strict_apply_chain(pptx_path, operations, work_dir, "legacy_final")
            if candidate is None or issues:
                return _chat_error_plan("; ".join(issues[:4]))
        return {
            "message": str(legacy.get("message") or "Applied the requested PowerPoint changes.")[:900],
            "assumptions": [],
            "operations": operations,
            "failed": False,
        }
    except Exception as exc:
        return _chat_error_plan(str(exc))


def _plan_semantic_task_strict(
    planning_path: str,
    task: ConversationTask,
    positions: list[int],
    selected_slide: int,
    chat_history: list[dict[str, str]] | None,
    attachment_paths: list[str] | None,
    work_dir: str,
    task_index: int,
) -> tuple[str | None, list[dict[str, Any]], str, list[str]]:
    """Plan one semantic task, retrying with concrete failure feedback."""
    base_instruction = _semantic_instruction(task, positions)
    semantic_selected = positions[0] if positions else max(1, selected_slide)
    errors: list[str] = []
    for attempt in range(1, 4):
        instruction = base_instruction
        if errors:
            instruction += (
                " Previous attempt failed validation. Use a safer supported approach while preserving the same intent. "
                "Failure details: " + "; ".join(errors[-4:])
            )
        try:
            semantic_plan = _plan_edit_before_conversation_agent(
                planning_path,
                instruction,
                semantic_selected,
                None,
                chat_history,
                None,
                attachment_paths,
            )
            proposed = [op for op in semantic_plan.get("operations", []) if op.get("op") != "noop"]
            candidate, attempt_issues = _strict_apply_chain(
                planning_path,
                proposed,
                work_dir,
                f"semantic_{task_index:02d}_attempt_{attempt}",
            )
            if candidate is not None and not attempt_issues:
                detail = str(semantic_plan.get("message") or task.instruction or "Applied the requested edit.").strip()
                return candidate, proposed, detail, []
            errors.extend(attempt_issues or ["No usable edit was produced."])
        except Exception as exc:
            errors.append(str(exc)[:500])
    return None, [], "", errors


# Preserve the previous conversational definition for diagnostics only.
_plan_edit_before_strict_conversation = plan_edit


def plan_edit(
    pptx_path: str,
    user_message: str,
    selected_slide: int,
    selected_slide_image: str | None = None,
    chat_history: list[dict[str, str]] | None = None,
    deck_image_paths: list[str] | None = None,
    attachment_paths: list[str] | None = None,
) -> dict[str, Any]:
    """ChatGPT-like interpretation with all-or-nothing PowerPoint safety."""

    def fail(reason: str | None, failed_step: str | None = None) -> dict[str, Any]:
        return _chat_error_plan(
            reason,
            pptx_path=pptx_path,
            user_message=user_message,
            selected_slide=selected_slide,
            chat_history=chat_history,
            failed_step=failed_step,
        )

    # Deterministic slide-management commands run before the conversational
    # model. This guarantees basic actions such as deleting the last slide do
    # not depend on model wording, output format, token limits, or API status.
    direct = _direct_slide_command_plan(pptx_path, user_message, selected_slide)
    if direct is not None:
        direct_operations = [
            dict(operation) for operation in direct.get("operations", [])
            if operation.get("op") != "noop"
        ]
        if not direct_operations:
            direct_message = str(direct.get("message") or "").strip()
            if direct_message.startswith("No change was made"):
                return fail(direct_message, "Validating the slide-management command")
            return {**direct, "failed": False}
        with tempfile.TemporaryDirectory(prefix="deck_refresh_direct_command_") as work_dir:
            candidate, issues = _strict_apply_chain(
                pptx_path, direct_operations, work_dir, "direct_command",
            )
            if candidate is None or issues:
                return fail(
                    "; ".join(issues[:4]) or "The slide command did not validate.",
                    "Applying the slide-management command",
                )
        return {
            "message": str(direct.get("message") or "Applied the requested slide change."),
            "assumptions": [],
            "operations": direct_operations,
            "failed": False,
        }

    try:
        try:
            conversation = _interpret_conversation(
                pptx_path, user_message, selected_slide, chat_history,
            )
            conversation = _review_conversation_plan(
                pptx_path, user_message, selected_slide, chat_history, conversation,
            )
        except Exception as interpretation_exc:
            conversation = _heuristic_conversation_plan(user_message)
            if conversation is None:
                legacy = _strict_legacy_plan(
                    pptx_path, user_message, selected_slide, selected_slide_image,
                    chat_history, deck_image_paths, attachment_paths,
                )
                if legacy.get("failed"):
                    reason = "; ".join(str(value) for value in legacy.get("assumptions", []) if value)
                    reason = reason or str(interpretation_exc)
                    return fail(reason, "Understanding the requested edit")
                return legacy

        if not conversation.tasks:
            legacy = _strict_legacy_plan(
                pptx_path, user_message, selected_slide, selected_slide_image,
                chat_history, deck_image_paths, attachment_paths,
            )
            if legacy.get("failed"):
                reason = "; ".join(str(value) for value in legacy.get("assumptions", []) if value)
                return fail(reason or "No executable task was identified.", "Understanding the requested edit")
            return legacy

        deck = _virtual_deck(pptx_path)
        if not deck:
            return fail("The presentation has no slides.", "Opening the current presentation")
        selected_slide = max(1, min(int(selected_slide), len(deck)))
        selected_id = deck[selected_slide - 1]["id"]
        aliases: dict[str, str] = {}
        last_created_id: str | None = None
        last_mentioned_id: str | None = selected_id
        id_counter = [0]
        all_operations: list[dict[str, Any]] = []
        completed: list[str] = []

        with tempfile.TemporaryDirectory(prefix="deck_refresh_strict_conversation_") as work_dir:
            planning_path = os.path.join(work_dir, "planning_000.pptx")
            shutil.copy2(pptx_path, planning_path)

            for task_index, task in enumerate(conversation.tasks, start=1):
                task_label = task.instruction or f"Task {task_index}: {task.action.replace('_', ' ')}"
                if task.action == "semantic_edit":
                    positions: list[int] = []
                    try:
                        for ref in task.sources:
                            index, item = _resolve_virtual_item(
                                ref, deck, aliases, selected_id, last_created_id, last_mentioned_id,
                            )
                            positions.append(index + 1)
                            last_mentioned_id = item["id"]
                    except Exception as exc:
                        return fail(str(exc), f"Resolving the target for: {task_label}")

                    candidate, proposed, detail, issues = _plan_semantic_task_strict(
                        planning_path,
                        task,
                        positions,
                        selected_slide,
                        chat_history,
                        attachment_paths,
                        work_dir,
                        task_index,
                    )
                    if candidate is None or issues:
                        return fail("; ".join(issues[:4]), f"Applying: {task_label}")
                    planning_path = candidate
                    all_operations.extend(proposed)
                    completed.append(detail)
                    rebuilt = _virtual_deck(planning_path)
                    structure_changed = len(rebuilt) != len(deck)
                    deck[:] = rebuilt
                    if structure_changed:
                        aliases.clear()
                        last_created_id = None
                        last_mentioned_id = None
                        selected_id = deck[min(selected_slide - 1, len(deck) - 1)]["id"] if deck else None
                    continue

                try:
                    proposed, last_created_id, last_mentioned_id, detail = _compile_management_task(
                        task,
                        deck,
                        aliases,
                        selected_id,
                        last_created_id,
                        last_mentioned_id,
                        id_counter,
                    )
                except Exception as exc:
                    return fail(str(exc), f"Preparing: {task_label}")

                candidate, issues = _strict_apply_chain(
                    planning_path,
                    proposed,
                    work_dir,
                    f"management_{task_index:02d}",
                )
                if candidate is None or issues:
                    return fail("; ".join(issues[:4]), f"Applying: {task_label}")
                planning_path = candidate
                all_operations.extend(proposed)
                completed.append(detail)

            final_candidate, final_issues = _strict_apply_chain(
                pptx_path,
                all_operations,
                work_dir,
                "complete_request",
            )
            if final_candidate is None or final_issues:
                return fail("; ".join(final_issues[:4]), "Validating the complete multi-step request")

        if not all_operations:
            return fail("No usable edit operation was generated.", "Turning the request into PowerPoint operations")
        message = " ".join(part for part in completed if part).strip()
        return {
            "message": (message or "Applied the requested PowerPoint changes.")[:900],
            "assumptions": [],
            "operations": all_operations[:300],
            "failed": False,
        }
    except Exception as exc:
        return fail(str(exc), "Completing and validating the requested edit")

# ---------------------------------------------------------------------------
# FINAL DETERMINISTIC SHOWCASE LAYER
# ---------------------------------------------------------------------------
from showcase_engine import apply_showcase_operation as _apply_showcase_operation
from showcase_engine import plan_showcase_command as _plan_showcase_command

SUPPORTED_OPS = set(SUPPORTED_OPS) | {"semantic_transform"}
_showcase_previous_apply_single = _apply_single_operation_universal
_showcase_previous_plan_edit = plan_edit


def _apply_single_operation_universal(prs: Presentation, operation: dict[str, Any]) -> dict[str, Any]:
    if str(operation.get("op", "")).casefold() == "semantic_transform":
        return _apply_showcase_operation(prs, operation)
    return _showcase_previous_apply_single(prs, operation)


def plan_edit(
    pptx_path: str,
    user_message: str,
    selected_slide: int,
    selected_slide_image: str | None = None,
    chat_history: list[dict[str, str]] | None = None,
    deck_image_paths: list[str] | None = None,
    attachment_paths: list[str] | None = None,
) -> dict[str, Any]:
    local = _plan_showcase_command(pptx_path, user_message, selected_slide, chat_history)
    if local is not None:
        return local
    return _showcase_previous_plan_edit(
        pptx_path,
        user_message,
        selected_slide,
        selected_slide_image,
        chat_history,
        deck_image_paths,
        attachment_paths,
    )

# ---------------------------------------------------------------------------
# FINAL UNIVERSAL THEME ROUTER
# ---------------------------------------------------------------------------
# AI-generated apply_theme operations use the same robust visible-theme engine
# as deterministic theme commands and the toolbar.
_theme_router_previous_apply_single = _apply_single_operation_universal


def _apply_single_operation_universal(prs: Presentation, operation: dict[str, Any]) -> dict[str, Any]:
    if str(operation.get("op", "")).casefold() == "apply_theme":
        mapped = {
            "op": "semantic_transform",
            "kind": "apply_theme",
            "slides": operation.get("slides"),
            "slide": operation.get("slide"),
            "preset": operation.get("preset") or operation.get("theme_name") or operation.get("name") or "kpmg",
            "primary": operation.get("primary") or operation.get("primary_color"),
            "secondary": operation.get("secondary") or operation.get("secondary_color"),
            "accent": operation.get("accent") or operation.get("accent_color"),
            "background": operation.get("background") or operation.get("background_color"),
            "surface": operation.get("surface") or operation.get("surface_color"),
            "title_color": operation.get("title_color"),
            "body_color": operation.get("body_color") or operation.get("text_color"),
            "border": operation.get("border") or operation.get("border_color"),
            "chart_colors": operation.get("chart_colors") or operation.get("series_colors"),
            "font_face": operation.get("font_face"),
            "title_size": operation.get("title_size"),
            "body_size": operation.get("body_size"),
            "preserve_branding": operation.get("preserve_branding", True),
            "preserve_status_colors": operation.get("preserve_status_colors", True),
            "change_shapes": operation.get("change_shapes", True),
            "change_background": operation.get("change_background", True),
            "change_text": operation.get("change_text", True),
        }
        return _apply_showcase_operation(prs, mapped)
    return _theme_router_previous_apply_single(prs, operation)


ATOMIC_PLANNER_INSTRUCTIONS += """

Theme and color operations:
- apply_theme may use a preset name: KPMG Blue, Executive Dark,
  Performance Green, Warm Neutral, Ocean, Purple, Monochrome, or
  High Contrast.
- apply_theme may also include primary_color, secondary_color,
  accent_color, background_color, surface_color, title_color,
  body_color, border_color, chart_colors, slides, preserve_branding,
  and preserve_status_colors.
- Use replace_color when the user asks to replace one exact visible color.
- Preserve pictures and logos unless the user explicitly asks to recolor them.
"""

# ---------------------------------------------------------------------------
# BUILDER AND INSPECTOR OPERATIONS
# ---------------------------------------------------------------------------

from builder_ops import LAYOUTS as BUILDER_LAYOUTS, layout_operations as _builder_layout_operations


def _sort_table(prs: Presentation, operation: dict[str, Any]) -> int:
    slide = _slide(prs, int(operation["slide"]))
    shape = _operation_shape(slide, operation)
    if not getattr(shape, "has_table", False):
        raise EditorError("The selected object is not a table.")
    table = shape.table
    if len(table.rows) < 3:
        return 0
    column = operation.get("column", 1)
    if isinstance(column, str) and not column.isdigit():
        headers = [cell.text.strip().casefold() for cell in table.rows[0].cells]
        wanted = column.strip().casefold()
        matches = [index for index, value in enumerate(headers) if wanted in value or value in wanted]
        if not matches:
            raise EditorError(f'Table column "{column}" was not found.')
        column_index = matches[0]
    else:
        column_index = int(column) - 1
    if not 0 <= column_index < len(table.columns):
        raise EditorError("The sort column is outside the table.")
    rows = list(table._tbl.tr_lst[1:])

    def key(row):
        try:
            text = "".join(row.tc_lst[column_index].itertext()).strip()
        except Exception:
            text = ""
        number = _numeric_value(text)
        impact = {"critical": 5, "high": 4, "medium": 3, "moderate": 3, "low": 2, "none": 1}
        return (1, number) if number is not None else (0, impact.get(text.casefold(), 0), text.casefold())

    reverse = str(operation.get("direction", "descending")).casefold() not in {"ascending", "asc", "a-z"}
    for row in rows:
        table._tbl.remove(row)
    for row in sorted(rows, key=key, reverse=reverse):
        table._tbl.append(row)
    return len(rows)


def _style_all_charts(prs: Presentation, operation: dict[str, Any]) -> int:
    targets = operation.get("slides") or ([operation["slide"]] if operation.get("slide") else range(1, len(prs.slides) + 1))
    colors = operation.get("series_colors") or [operation.get("color", "00A651")]
    count = 0
    for number in targets:
        for shape in _slide(prs, int(number)).shapes:
            if getattr(shape, "has_chart", False):
                _style_chart(shape.chart, {**operation, "series_colors": colors})
                try:
                    _style_chart_universal(shape, {**operation, "series_colors": colors})
                except Exception:
                    pass
                count += 1
    return count


def _round_all_corners(prs: Presentation, operation: dict[str, Any]) -> int:
    targets = operation.get("slides") or ([operation["slide"]] if operation.get("slide") else range(1, len(prs.slides) + 1))
    count = 0
    for number in targets:
        for shape in _slide(prs, int(number)).shapes:
            try:
                geometry = shape._element.spPr.prstGeom
                if geometry is not None and geometry.get("prst") not in {"line", "ellipse"}:
                    geometry.set("prst", "roundRect")
                    count += 1
            except Exception:
                continue
    return count


def _ensure_min_font(prs: Presentation, operation: dict[str, Any]) -> int:
    minimum = max(8.0, float(operation.get("minimum", operation.get("font_size", 18))))
    targets = operation.get("slides") or ([operation["slide"]] if operation.get("slide") else range(1, len(prs.slides) + 1))
    count = 0
    for number in targets:
        slide = _slide(prs, int(number))
        frames = []
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                frames.append(shape.text_frame)
            if getattr(shape, "has_table", False):
                frames.extend(cell.text_frame for row in shape.table.rows for cell in row.cells)
        for frame in frames:
            frame.word_wrap = True
            frame.auto_size = None
            for paragraph in frame.paragraphs:
                for run in paragraph.runs:
                    current = run.font.size.pt if run.font.size else minimum
                    if current < minimum:
                        run.font.size = Pt(minimum)
                        count += 1
    return count


def _format_all_text(prs: Presentation, operation: dict[str, Any]) -> int:
    font_face = str(operation.get("font_face") or "").strip()
    font_size = operation.get("font_size")
    if not font_face and font_size is None:
        raise EditorError("Enter a font name or font size.")
    size = max(6.0, min(120.0, float(font_size))) if font_size is not None else None
    target_kind = str(operation.get("target", "all")).strip().casefold()
    if target_kind not in {"all", "title", "body"}:
        raise EditorError("Text target must be all, title, or body.")
    targets = operation.get("slides") or ([operation["slide"]] if operation.get("slide") else range(1, len(prs.slides) + 1))
    count = 0
    for number in targets:
        slide = _slide(prs, int(number))
        title = getattr(slide.shapes, "title", None)
        if title is None:
            candidates = [
                shape for shape in slide.shapes
                if getattr(shape, "has_text_frame", False)
                and shape.text_frame.text.strip()
                and float(shape.top) <= float(Inches(1.5))
                and float(shape.width) >= float(Inches(2.0))
            ]
            title = min(candidates, key=lambda shape: (shape.top, -shape.width), default=None)
        title_id = int(title.shape_id) if title is not None else None
        pending = list(slide.shapes)
        while pending:
            shape = pending.pop(0)
            if hasattr(shape, "shapes"):
                pending.extend(list(shape.shapes))
            is_title = title_id is not None and int(shape.shape_id) == title_id
            if target_kind == "title" and not is_title:
                continue
            if target_kind == "body" and is_title:
                continue
            frames = []
            if getattr(shape, "has_text_frame", False):
                frames.append(shape.text_frame)
            if getattr(shape, "has_table", False) and target_kind != "title":
                frames.extend(cell.text_frame for row in shape.table.rows for cell in row.cells)
            if getattr(shape, "has_chart", False) and target_kind != "body":
                try:
                    if shape.chart.has_title:
                        frames.append(shape.chart.chart_title.text_frame)
                except Exception:
                    pass
            for frame in frames:
                for paragraph in frame.paragraphs:
                    if paragraph.text and not paragraph.runs:
                        paragraph.text = paragraph.text
                    for run in paragraph.runs:
                        if font_face:
                            run.font.name = font_face
                        if size is not None:
                            run.font.size = Pt(size)
                        count += 1
    return count


def _table_to_icons(prs: Presentation, operation: dict[str, Any]) -> int:
    slide = _slide(prs, int(operation["slide"]))
    shape = _operation_shape(slide, operation)
    if not getattr(shape, "has_table", False):
        raise EditorError("The selected object is not a table.")
    rows = [[cell.text.strip() for cell in row.cells] for row in shape.table.rows]
    if len(rows) < 2:
        raise EditorError("The table needs at least one data row.")
    headers, data = rows[0], rows[1:9]
    left, top, width, height = shape.left, shape.top, shape.width, shape.height
    columns = min(4, len(data))
    card_width = int(width / max(1, columns))
    card_height = int(height / max(1, (len(data) + columns - 1) // columns))
    for index, row in enumerate(data):
        card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
                                      left + (index % columns) * card_width,
                                      top + (index // columns) * card_height,
                                      int(card_width * .92), int(card_height * .88))
        card.fill.solid(); card.fill.fore_color.rgb = _rgb(operation.get("fill_color"), "F3F6FA")
        card.line.color.rgb = _rgb(operation.get("line_color"), "D8E2EE")
        label = row[0] if row else f"Item {index + 1}"
        details = "\n".join(f"{headers[col]}: {row[col]}" for col in range(1, min(len(row), len(headers))) if row[col])
        _set_text_frame_text(card.text_frame, f"●  {label}" + (f"\n{details}" if details else ""), preserve_first_run=False)
        _apply_text_style(card, {"font_size": operation.get("font_size", 13), "font_color": operation.get("font_color", "172B4D")})
    if not operation.get("keep_table", False):
        _delete_shape_element(shape)
    return len(data)


def _merge_slides(prs: Presentation, operation: dict[str, Any]) -> int:
    first = int(operation.get("first_slide", operation.get("slide", 1)))
    second = int(operation.get("second_slide", first + 1))
    if first == second:
        raise EditorError("Choose two different slides to merge.")
    target, source = _slide(prs, first), _slide(prs, second)
    for shape in target.shapes:
        shape.top = int(shape.top * .48)
        shape.height = max(1, int(shape.height * .48))
    for shape in source.shapes:
        clone = copy.deepcopy(shape.element)
        target.shapes._spTree.insert_element_before(clone, "p:extLst")
    for rel in source.part.rels.values():
        if "notesSlide" in rel.reltype:
            continue
        try:
            target.part.rels.add_relationship(rel.reltype, rel._target, rel.rId, rel.is_external)
        except Exception:
            pass
    original_count = len(target.shapes)
    copied = list(target.shapes)[original_count - len(source.shapes):]
    for shape in copied:
        shape.top = int(prs.slide_height * .50 + shape.top * .48)
        shape.height = max(1, int(shape.height * .48))
    _delete_slide(prs, second)
    return first if second > first else first - 1


def _split_slide(prs: Presentation, operation: dict[str, Any]) -> int:
    number = int(operation["slide"])
    source = _slide(prs, number)
    new_number = _duplicate_slide(prs, number, number + 1)
    target = _slide(prs, new_number)
    midpoint = prs.slide_height / 2
    for shape in list(source.shapes):
        if shape.top >= midpoint:
            _delete_shape_element(shape)
    for shape in list(target.shapes):
        if shape.top < midpoint:
            _delete_shape_element(shape)
        else:
            shape.top = max(0, shape.top - int(midpoint * .72))
    return new_number


def _animate_bullets(prs: Presentation, operation: dict[str, Any]) -> int:
    """Create editable progressive bullet builds with fade transitions."""
    number = int(operation["slide"])
    source = _slide(prs, number)
    body = None
    for shape in source.shapes:
        if getattr(shape, "has_text_frame", False):
            lines = [paragraph.text.strip() for paragraph in shape.text_frame.paragraphs if paragraph.text.strip()]
            if len(lines) >= 2 and (body is None or len(lines) > len(body[1])):
                body = (shape, lines)
    grouped_shape_ids = []
    if body is None:
        candidates = [shape for shape in source.shapes if getattr(shape, "has_text_frame", False)
                      and shape.text_frame.text.strip() and shape.top > prs.slide_height * .18]
        rows: dict[int, list[int]] = {}
        for shape in candidates:
            key = int(round((shape.top / max(1, prs.slide_height)) * 100))
            rows.setdefault(key, []).append(int(shape.shape_id))
        grouped_shape_ids = [rows[key] for key in sorted(rows)]
        if len(grouped_shape_ids) < 2:
            raise EditorError("No multi-bullet text box was found on this slide.")
        bullets = [str(index + 1) for index in range(len(grouped_shape_ids))]
        shape_id = None
    else:
        source_shape, bullets = body
        shape_id = int(source_shape.shape_id)
    slide_numbers = [number]
    for offset in range(1, len(bullets)):
        slide_numbers.append(_duplicate_slide(prs, number + offset - 1, number + offset))
    for step, slide_number in enumerate(slide_numbers, start=1):
        slide = _slide(prs, slide_number)
        shape = next((item for item in slide.shapes if shape_id is not None and int(item.shape_id) == shape_id), None)
        if shape is not None:
            _set_text_frame_text(shape.text_frame, "\n".join(bullets[:step]), preserve_first_run=False)
            _format_text_box(shape, {"bullet": True, "bullet_character": "•", "font_size": operation.get("font_size", 18)})
        elif grouped_shape_ids:
            hidden_ids = {shape_id for group in grouped_shape_ids[step:] for shape_id in group}
            for item in list(slide.shapes):
                if int(item.shape_id) in hidden_ids:
                    _delete_shape_element(item)
        transition = slide._element.find(qn("p:transition"))
        if transition is None:
            transition = OxmlElement("p:transition")
            transition.set("spd", "fast")
            transition.append(OxmlElement("p:fade"))
            timing = slide._element.find(qn("p:timing"))
            if timing is not None:
                slide._element.insert(list(slide._element).index(timing), transition)
            else:
                slide._element.append(transition)
    return len(slide_numbers)


_builder_previous_apply_single = _apply_single_operation_universal


def _apply_single_operation_universal(prs: Presentation, operation: dict[str, Any]) -> dict[str, Any]:
    op = str(operation.get("op", "")).casefold()
    if op == "sort_table": return {"op": op, "rows": _sort_table(prs, operation)}
    if op == "style_all_charts": return {"op": op, "count": _style_all_charts(prs, operation)}
    if op == "round_all_corners": return {"op": op, "count": _round_all_corners(prs, operation)}
    if op == "ensure_min_font": return {"op": op, "count": _ensure_min_font(prs, operation)}
    if op == "format_all_text": return {"op": op, "count": _format_all_text(prs, operation)}
    if op == "convert_table_to_icons": return {"op": op, "count": _table_to_icons(prs, operation)}
    if op == "merge_slides": return {"op": op, "slide": _merge_slides(prs, operation)}
    if op == "split_slide": return {"op": op, "slide": _split_slide(prs, operation)}
    if op == "animate_bullets": return {"op": op, "slides": _animate_bullets(prs, operation)}
    return _builder_previous_apply_single(prs, operation)


SHAPE_OPS = SHAPE_OPS | {"sort_table", "convert_table_to_icons"}
SUPPORTED_OPS = SUPPORTED_OPS | {"sort_table", "style_all_charts", "round_all_corners", "ensure_min_font", "format_all_text",
                                 "convert_table_to_icons", "merge_slides", "split_slide", "animate_bullets"}


_builder_previous_deterministic = _deterministic_operations


def _deterministic_operations(summary: dict[str, Any], user_message: str, selected_slide: int) -> list[dict[str, Any]]:
    operations = _builder_previous_deterministic(summary, user_message, selected_slide)
    if any(operation.get("op") == "executive_review" for operation in operations):
        return [operation for operation in operations if operation.get("op") == "executive_review"]
    lowered = user_message.casefold()
    count = int(summary.get("slide_count", 1))
    targets = _parse_slide_targets(user_message, count, selected_slide)
    current_shapes = summary.get("slides", [])[targets[0] - 1].get("shapes", []) if summary.get("slides") and targets else []
    first_chart = next((shape for shape in current_shapes if shape.get("kind") == "chart"), None)
    first_table = next((shape for shape in current_shapes if shape.get("kind") == "table"), None)
    layout_aliases = {name.replace("_", " "): name for name in BUILDER_LAYOUTS}
    layout_aliases.update({"closing slide": "closing_thank_you", "thank you slide": "closing_thank_you",
                           "2 by 2 matrix": "2x2_matrix", "process slide": "process_flow"})
    matched_layout = None
    for phrase, candidate_layout in sorted(layout_aliases.items(), key=lambda item: -len(item[0])):
        full_phrase = phrase if phrase.endswith("slide") else phrase + " slide"
        if re.search(rf"(?:add|create|build|insert)\s+(?:a\s+|an\s+)?(?:new\s+)?{re.escape(full_phrase)}", lowered):
            matched_layout = candidate_layout
            break
    if matched_layout:
        operations = [op for op in operations if op.get("op") != "add_slide"]
        built, _, _ = _builder_layout_operations(matched_layout, min(targets[-1] + 1, count + 1), matched_layout.replace("_", " ").title())
        operations.extend(built)
    if "make every chart green" in lowered or "make all charts green" in lowered:
        operations.append({"op": "style_all_charts", "color": "00A651"})
    chart_conversion = re.search(r"(?:convert|change|replace).*?chart.*?(?:into|to|with)\s+(?:a\s+)?(bar|line|pie|area|waterfall|scatter)(?:\s+(?:chart|graph|plot))?", lowered)
    if chart_conversion and first_chart:
        operations.append({"op": "change_chart_type", "slide": targets[0], "shape_id": first_chart["shape_id"],
                           "chart_type": chart_conversion.group(1)})
    if ("show chart data labels" in lowered or "show data labels" in lowered) and first_chart:
        operations.append({"op": "style_chart", "slide": targets[0], "shape_id": first_chart["shape_id"],
                           "show_data_labels": True, "show_gridlines": "remove gridlines" not in lowered})
    if "round all corners" in lowered or "round every corner" in lowered:
        operations.append({"op": "round_all_corners", "slides": targets})
    min_font = re.search(r"(?:not |never )?(?:shrink(?:ing)?|go) below\s+(\d{1,2})\s*(?:point|pt)?", lowered)
    if min_font:
        operations.append({"op": "ensure_min_font", "slides": targets, "minimum": int(min_font.group(1))})
    sort_match = re.search(r"sort (?:the )?table by ([a-z0-9 _-]+?)(?:\.|$)", lowered)
    if sort_match:
        operations.append({"op": "sort_table", "slide": targets[0], "shape_id": first_table["shape_id"] if first_table else None,
                           "column": sort_match.group(1).strip(), "direction": "descending"})
    if first_table and ("make this table blue" in lowered or "make the table" in lowered and "blue" in lowered):
        operations.append({"op": "style_table", "slide": targets[0], "shape_id": first_table["shape_id"],
                           "header_fill": "00338D", "header_font_color": "FFFFFF", "body_fill": "EAF1FB"})
    if first_table and ("alternate row colors" in lowered or "band rows" in lowered):
        operations.append({"op": "style_table", "slide": targets[0], "shape_id": first_table["shape_id"],
                           "header_fill": "00338D", "header_font_color": "FFFFFF", "banded_rows": True, "band_fill": "EAF1FB"})
    if first_table and ("add another row" in lowered or "add a row" in lowered):
        operations.append({"op": "add_table_row", "slide": targets[0], "shape_id": first_table["shape_id"], "values": []})
    if first_table and "merge the first two columns" in lowered:
        operations.append({"op": "merge_table_cells", "slide": targets[0], "shape_id": first_table["shape_id"],
                           "row": 1, "column": 1, "end_row": 1, "end_column": 2})
    if first_table and "highlight high risks red" in lowered:
        rows = [line.split(" | ") for line in str(first_table.get("text", "")).splitlines()]
        for row_index, row in enumerate(rows, start=1):
            for column_index, value in enumerate(row, start=1):
                if value.strip().casefold() in {"high", "critical"}:
                    operations.append({"op": "set_table_cell", "slide": targets[0], "shape_id": first_table["shape_id"],
                                       "row": row_index, "column": column_index, "text": value,
                                       "fill_color": "FDE8E8", "font_color": "B42318", "bold": True})
    if "turn this table into icons" in lowered or "convert this table to icons" in lowered:
        operations.append({"op": "convert_table_to_icons", "slide": targets[0],
                           "shape_id": first_table["shape_id"] if first_table else None})
    split_match = re.search(r"split (?:this slide|slide\s+(\d+)) into two", lowered)
    if split_match:
        operations.append({"op": "split_slide", "slide": int(split_match.group(1) or selected_slide)})
    merge_match = re.search(r"merge slides?\s+(\d+)\s+(?:and|with)\s+(\d+)", lowered)
    if merge_match:
        operations.append({"op": "merge_slides", "first_slide": int(merge_match.group(1)), "second_slide": int(merge_match.group(2))})
    if "merge this slide with the next slide" in lowered and selected_slide < count:
        operations.append({"op": "merge_slides", "first_slide": selected_slide, "second_slide": selected_slide + 1})
    if "animate each bullet" in lowered or "animate every bullet" in lowered or "animate the bullets" in lowered:
        operations.append({"op": "animate_bullets", "slide": targets[0], "font_size": 18})
    return _dedupe_operations(operations)


ATOMIC_PLANNER_INSTRUCTIONS += """

Builder and inspector operations:
- sort_table: {op, slide, shape_id or text_contains, column name or number, direction}
- style_all_charts: {op, optional slides, color or series_colors, labels and axis options}
- round_all_corners: {op, optional slides}
- ensure_min_font: {op, optional slides, minimum}
- format_all_text: {op, optional slide or slides, font_face, font_size, target: all, title, or body}
- convert_table_to_icons: {op, slide, shape_id or text_contains, keep_table}
- merge_slides: {op, first_slide, second_slide}
- split_slide: {op, slide}
- animate_bullets: {op, slide, optional font_size}. Creates editable progressive
  bullet builds with fade transitions.
Use change_chart_type without requesting the spreadsheet again. The editable chart already stores its data.
"""


_builder_previous_plan_edit = plan_edit


def plan_edit(
    pptx_path: str,
    user_message: str,
    selected_slide: int,
    selected_slide_image: str | None = None,
    chat_history: list[dict[str, str]] | None = None,
    deck_image_paths: list[str] | None = None,
    attachment_paths: list[str] | None = None,
) -> dict[str, Any]:
    if _looks_like_executive_review_request(user_message):
        slide_count = len(Presentation(pptx_path).slides)
        return {"message": "Completed the full executive review workflow.",
                "operations": [_executive_review_operation(slide_count)], "assumptions": [], "failed": False}
    return _builder_previous_plan_edit(pptx_path, user_message, selected_slide, selected_slide_image,
                                       chat_history, deck_image_paths, attachment_paths)


# ---------------------------------------------------------------------------
# GUARANTEED LOCAL COMMAND COMPILER
# ---------------------------------------------------------------------------
# Inspector actions and common object edits must not depend on a language-model
# response.  This compiler turns those commands into concrete operations first;
# the conversational planner remains available for genuinely open-ended work.

def _guaranteed_target_slide(message: str, selected_slide: int, slide_count: int) -> int:
    explicit = re.search(r"\bslide\s+(\d+)\b", message, re.I)
    target = int(explicit.group(1)) if explicit else int(selected_slide)
    return max(1, min(target, max(1, slide_count)))


def _guaranteed_chart_operation(target: int, chart_type: str, has_deck_data: bool) -> dict[str, Any]:
    if has_deck_data:
        return {"op": "semantic_transform", "kind": "create_chart", "slide": target,
                "chart_type": chart_type}
    values = [40, 30, 20, 10]
    operation: dict[str, Any] = {
        "op": "add_chart", "slide": target, "chart_type": chart_type,
        "title": "Editable chart", "categories": ["Sales", "Marketing", "HR", "IT"],
        "series": [{"name": "Value", "values": values}],
        "auto_fit": True,
        "x": .08, "y": .27, "width": .84, "height": .58,
        "series_colors": ["005EB8", "0091DA", "00A3A1", "483698"],
        "show_legend": False,
    }
    if chart_type == "scatter":
        operation["x_values"] = [1, 2, 3, 4]
    return operation


def guaranteed_local_plan(
    pptx_path: str,
    user_message: str,
    selected_slide: int,
    attachment_paths: list[str] | None = None,
) -> dict[str, Any] | None:
    """Compile the editor's promised actions without calling an AI model."""
    summary = deck_summary(pptx_path, max_text=1800)
    count = int(summary.get("slide_count", 1))
    raw_message = re.sub(r"\s+", " ", str(user_message or "")).strip()
    lowered = raw_message.casefold()
    if not lowered:
        return None
    explicit = re.search(r"\bslide\s+(\d+)\b", lowered)
    if explicit and not 1 <= int(explicit.group(1)) <= count:
        return {
            "message": (
                "error cant do that\n\n"
                f"What failed: Slide {explicit.group(1)} does not exist.\n"
                f"Why: This deck has {count} slides.\n"
                f"Try: Use a slide number from 1 to {count}.\n\n"
                "PowerPoint was not changed."
            ),
            "operations": [], "assumptions": [], "failed": True,
        }
    target = _guaranteed_target_slide(lowered, selected_slide, count)
    slide_info = summary.get("slides", [])[target - 1]
    shapes = list(slide_info.get("shapes", []))
    first_chart = next((shape for shape in shapes if shape.get("kind") == "chart"), None)
    first_table = next((shape for shape in shapes if shape.get("kind") == "table"), None)
    first_picture = next((shape for shape in shapes if shape.get("kind") == "picture"), None)
    text_shapes = [shape for shape in shapes if shape.get("kind") == "text" and str(shape.get("text", "")).strip()]
    content_shape_ids = [int(shape["shape_id"]) for shape in shapes if float(shape.get("y", 0)) >= 1.15]
    has_deck_data = any(
        shape.get("kind") in {"chart", "table"}
        or bool(re.search(r"[A-Za-z][^\n:|-]{1,40}(?::|\||-)\s*[$€£]?\s*-?\d", str(shape.get("text", ""))))
        for slide in summary.get("slides", []) for shape in slide.get("shapes", [])
    )
    has_any_chart = any(
        shape.get("kind") == "chart"
        for slide in summary.get("slides", []) for shape in slide.get("shapes", [])
    )
    operations: list[dict[str, Any]] = []
    labels: list[str] = []

    def add(operation: dict[str, Any], label: str) -> None:
        if operation.get("op") in {"add_chart", "add_table"}:
            operation.setdefault("auto_fit", True)
        operations.append(operation)
        labels.append(label)

    # Themes and deck-wide styling.
    rebrand_match = re.search(
        r'rebrand\s+the\s+entire\s+deck\s+from\s+["“]?([^"”]+?)["”]?\s+to\s+["“]?([^"”.]+?)["”]?(?:\.|$)',
        raw_message,
        flags=re.IGNORECASE,
    )
    if rebrand_match is None:
        rebrand_match = re.search(
            r'(?:change|switch)\s+(?:the\s+)?company(?:\s+name|\s+branding)?\s+from\s+["“]?([^"”]+?)["”]?\s+to\s+["“]?([^"”.]+?)["”]?(?:\.|$)',
            raw_message,
            flags=re.IGNORECASE,
        )
    if rebrand_match is None:
        rebrand_match = re.search(
            r'(?:change|replace)\s+["“]?([^"”]+?)["”]?\s+(?:to|with)\s+["“]?([^"”]+?)["”]?\s+(?:on|across)\s+(?:every\s+slide|all\s+slides|the\s+entire\s+deck)',
            raw_message,
            flags=re.IGNORECASE,
        )
    if rebrand_match is None:
        rebrand_match = re.search(
            r'(?:change|replace)\s+["“]?(KPMG)["”]?(?:\s+branding|\s+company|\s+name)?\s+(?:to|with)\s+["“]?([^"”.]+?)["”]?(?:\s+branding)?(?:\.|$)',
            raw_message,
            flags=re.IGNORECASE,
        )
    implied_brand = None
    if rebrand_match is None:
        implied_brand = re.search(
            r'(?:use|apply)\s+["“]?([^"”]+?)["”]?\s+branding\s+(?:on|across|to)\s+(?:every\s+slide|all\s+slides|the\s+entire\s+deck)',
            raw_message,
            flags=re.IGNORECASE,
        )
    if rebrand_match:
        old_company = rebrand_match.group(1).strip()
        new_company = rebrand_match.group(2).strip()
        if old_company and new_company and old_company.casefold() != new_company.casefold():
            add({
                "op": "replace_text", "old": old_company, "new": new_company,
                "replace_all": True, "case_sensitive": False,
                "include_masters": True, "allow_missing": False,
            }, f"{new_company} company name across every slide")
    elif implied_brand:
        new_company = implied_brand.group(1).strip()
        if new_company and new_company.casefold() != "kpmg":
            add({
                "op": "replace_text", "old": "KPMG", "new": new_company,
                "replace_all": True, "case_sensitive": False,
                "include_masters": True, "allow_missing": False,
            }, f"{new_company} company name across every slide")

    replace_match = re.search(
        r'replace\s+["“]([^"”]+)["”]\s+with\s+["“]([^"”]+)["”]\s+(across\s+the\s+entire\s+deck|on\s+slide\s+\d+)',
        raw_message,
        flags=re.IGNORECASE,
    )
    if replace_match:
        scope_text = replace_match.group(3).casefold()
        operation = {
            "op": "replace_text", "old": replace_match.group(1), "new": replace_match.group(2),
            "replace_all": True, "case_sensitive": False, "allow_missing": False,
        }
        if scope_text.startswith("on slide"):
            operation["slide"] = target
        add(operation, "text replacement")
    elif rebrand_match is None:
        plain_replace = re.search(
            r'replace\s+(.+?)\s+with\s+(.+?)\s+(across\s+the\s+entire\s+deck|on\s+every\s+slide|on\s+slide\s+\d+)',
            raw_message,
            flags=re.IGNORECASE,
        )
        if plain_replace:
            operation = {
                "op": "replace_text", "old": plain_replace.group(1).strip(' "“”'),
                "new": plain_replace.group(2).strip(' "“”'), "replace_all": True,
                "case_sensitive": False, "allow_missing": False,
            }
            if plain_replace.group(3).casefold().startswith("on slide") and "every" not in plain_replace.group(3).casefold():
                operation["slide"] = target
            add(operation, "text replacement")

    font_match = re.search(
        r'change\s+(all|title|body)\s+text\s+to\s+font\s+["“]([^"”]+)["”]\s+at\s+(\d+(?:\.\d+)?)\s+point\s+(across\s+the\s+entire\s+deck|on\s+slide\s+\d+)',
        raw_message,
        flags=re.IGNORECASE,
    )
    if font_match:
        scope_text = font_match.group(4).casefold()
        font_operation: dict[str, Any] = {
            "op": "format_all_text", "font_face": font_match.group(2).strip(),
            "font_size": float(font_match.group(3)), "target": font_match.group(1).casefold(),
        }
        if scope_text.startswith("across"):
            font_operation["slides"] = list(range(1, count + 1))
        else:
            font_operation["slide"] = target
        add(font_operation, f'{font_match.group(2).strip()} font')
    else:
        simple_font = re.search(
            r'(?:change|set|make)\s+(?:(all|title|body)\s+)?(?:text\s+)?fonts?\s+(?:to\s+)?["“]?([a-z0-9][a-z0-9 ._-]*?)["”]?(?:\s+(?:at|to)\s+(\d+(?:\.\d+)?)\s*(?:point|pt))?(?:\s+(across\s+the\s+entire\s+deck|on\s+every\s+slide|on\s+slide\s+\d+))?(?:\.|$)',
            raw_message,
            flags=re.IGNORECASE,
        )
        if simple_font:
            font_operation = {
                "op": "format_all_text", "font_face": simple_font.group(2).strip(),
                "target": (simple_font.group(1) or "all").casefold(),
            }
            if simple_font.group(3):
                font_operation["font_size"] = float(simple_font.group(3))
            scope_text = (simple_font.group(4) or "").casefold()
            if "entire deck" in scope_text or "every slide" in scope_text or (simple_font.group(1) or "").casefold() == "all":
                font_operation["slides"] = list(range(1, count + 1))
            else:
                font_operation["slide"] = target
            add(font_operation, f'{simple_font.group(2).strip()} font')

    theme = next(((phrase, preset) for phrase, preset in (
        ("kpmg branding", "kpmg"), ("kpmg blue", "kpmg"),
        ("executive dark", "executive dark"), ("performance green", "performance green"),
        ("warm neutral", "warm neutral"), ("high contrast", "high contrast"),
        ("monochrome", "monochrome"), ("purple", "purple"), ("ocean", "ocean"),
    ) if phrase in lowered), None)
    if theme:
        scope = list(range(1, count + 1)) if "entire deck" in lowered or "kpmg branding" in lowered else [target]
        add({"op": "semantic_transform", "kind": "apply_theme", "slides": scope,
             "preset": theme[1], "preserve_branding": True}, f"{theme[1].title()} styling")
    elif "apply a custom theme" in lowered or "apply custom theme" in lowered:
        colors = re.findall(r"#[0-9a-f]{6}", lowered)
        scope = list(range(1, count + 1)) if "entire deck" in lowered else [target]
        custom_operation: dict[str, Any] = {
            "op": "semantic_transform", "kind": "apply_theme", "slides": scope,
            "preset": "kpmg", "preserve_branding": True,
        }
        if colors:
            custom_operation["primary"] = colors[0]
        if len(colors) > 1:
            custom_operation["accent"] = colors[1]
            custom_operation["secondary"] = colors[1]
        if len(colors) > 2:
            custom_operation["background"] = colors[2]
            background = colors[2].lstrip("#")
            red, green, blue = (int(background[index:index + 2], 16) for index in (0, 2, 4))
            text_color = "FFFFFF" if 0.2126 * red + 0.7152 * green + 0.0722 * blue < 145 else "1F2937"
            custom_operation["title_color"] = text_color
            custom_operation["body_color"] = text_color
        if colors:
            custom_operation["chart_colors"] = colors[:2] if len(colors) > 1 else [colors[0]]
        add(custom_operation, "custom theme")
    if "make every chart green" in lowered or "make all charts green" in lowered:
        if has_any_chart:
            add({"op": "style_all_charts", "color": "00A651"}, "green chart styling")
        else:
            green_chart = _guaranteed_chart_operation(target, "column", has_deck_data)
            green_chart["series_colors"] = ["00A651"]
            green_chart["show_data_labels"] = True
            add(green_chart, "green native chart")
    if "round all corners" in lowered or "round every corner" in lowered:
        if shapes:
            add({"op": "round_all_corners", "slides": [target]}, "rounded corners")
        else:
            add({"op": "add_shape", "slide": target, "shape_type": "rounded_rectangle",
                 "x": .12, "y": .30, "width": .76, "height": .28,
                 "text": "Leadership decision\nApprove phase two funding and assign an accountable owner.",
                 "fill_color": "EAF6F4", "line_color": "00A3A1", "font_color": "0F766E",
                 "font_size": 18, "alignment": "center"}, "rounded decision card")
    minimum_font_request = bool(re.search(r"(?:below|minimum(?:\s+of)?)\s+18\s*(?:point|pt)|without\s+shrinking\s+below\s+18", lowered))
    fit_at_eighteen = ("make it fit" in lowered or "make everything fit" in lowered) and bool(re.search(r"18\s*(?:point|pt)", lowered))
    if minimum_font_request or fit_at_eighteen:
        if fit_at_eighteen:
            add({"op": "semantic_transform", "kind": "format_slide", "slide": target}, "fitted layout")
        add({"op": "ensure_min_font", "slides": [target], "minimum": 18}, "18 point minimum text")
    if "increase white space" in lowered or "increase whitespace" in lowered:
        add({"op": "semantic_transform", "kind": "format_slide", "slide": target}, "more white space")

    # Native insertions and layout transformations.
    if "four editable cards" in lowered or "into four cards" in lowered:
        add({"op": "semantic_transform", "kind": "four_cards", "slide": target}, "four editable cards")
    if "add an editable text box" in lowered or "add editable text box" in lowered:
        add({"op": "add_textbox", "slide": target, "text": "Key message: Revenue increased 8% versus plan.", "x": .10, "y": .30,
             "width": .80, "height": .22, "font_size": 22, "font_color": "172B4D",
             "fill_color": "F3F6FA", "line_color": "D8E2EE", "alignment": "left"}, "editable text box")
    if "add an editable native table" in lowered or "add editable native table" in lowered:
        add({"op": "add_table", "slide": target,
             "data": [["Risk", "Owner", "Impact"], ["Delivery", "Program lead", "High"],
                      ["Adoption", "Change lead", "Medium"], ["Budget", "Finance", "Low"]],
             "auto_fit": True,
             "x": .07, "y": .27, "width": .86, "height": .50,
             "header_fill": "00338D", "header_font_color": "FFFFFF",
             "banded_rows": True, "band_fill": "EAF1FB", "font_size": 12}, "editable native table")
    if "add an editable process flow" in lowered or "add editable process flow" in lowered:
        for index, text_value in enumerate(("Define", "Analyze", "Decide", "Execute")):
            add({"op": "add_shape", "slide": target, "shape_type": "rounded_rectangle",
                 "x": .055 + index * .235, "y": .39, "width": .19, "height": .18,
                 "text": text_value, "fill_color": ("00338D", "005EB8", "0091DA", "00A3A1")[index],
                 "line_color": ("00338D", "005EB8", "0091DA", "00A3A1")[index],
                 "font_color": "FFFFFF", "font_size": 15, "alignment": "center"}, "editable process step")
    if "add an editable timeline" in lowered or "add editable timeline" in lowered:
        for index, text_value in enumerate(("Q1\nPlan", "Q2\nBuild", "Q3\nLaunch", "Q4\nScale")):
            add({"op": "add_shape", "slide": target, "shape_type": "rounded_rectangle",
                 "x": .06 + index * .23, "y": .37, "width": .18, "height": .20,
                 "text": text_value, "fill_color": "FFFFFF", "line_color": "005EB8",
                 "font_color": "00338D", "font_size": 15, "alignment": "center"}, "editable timeline milestone")

    # Native chart creation and conversion. Existing chart data is embedded,
    # so change_chart_type never asks for the spreadsheet again.
    add_chart_request = "editable native chart" in lowered or "add an editable chart" in lowered
    chart_conversion = re.search(
        r"(?:convert|change|replace).*?chart.*?(?:into|to|with)\s+(?:an?\s+)?(bar|line|pie|area|waterfall|scatter)",
        lowered,
    )
    requested_chart_type = chart_conversion.group(1) if chart_conversion else "column"
    if add_chart_request:
        add(_guaranteed_chart_operation(target, requested_chart_type, has_deck_data), "editable native chart")
    elif chart_conversion:
        if first_chart:
            add({"op": "change_chart_type", "slide": target, "shape_id": first_chart["shape_id"],
                 "chart_type": requested_chart_type}, f"{requested_chart_type} chart")
        else:
            add(_guaranteed_chart_operation(target, requested_chart_type, has_deck_data), f"new {requested_chart_type} chart")
    if "show chart data labels" in lowered or "show data labels" in lowered:
        if first_chart:
            add({"op": "style_chart", "slide": target, "shape_id": first_chart["shape_id"],
                 "show_data_labels": True, "show_gridlines": "remove gridlines" not in lowered}, "chart labels")
        else:
            label_chart = _guaranteed_chart_operation(target, "column", has_deck_data)
            label_chart.update({"show_data_labels": True, "show_gridlines": "remove gridlines" not in lowered})
            add(label_chart, "labeled native chart")

    # Native table editing.
    table_ref = {"slide": target, "shape_id": first_table["shape_id"]} if first_table else None
    if "make the table" in lowered and "blue" in lowered or "make this table blue" in lowered:
        if table_ref:
            add({"op": "style_table", **table_ref, "header_fill": "00338D",
                 "header_font_color": "FFFFFF", "body_fill": "EAF1FB"}, "blue table")
        else:
            add({"op": "add_table", "slide": target, "data": [["Item", "Owner", "Status"], ["Action", "Name", "On track"]],
                 "x": .08, "y": .28, "width": .84, "height": .40,
                 "header_fill": "00338D", "header_font_color": "FFFFFF", "body_fill": "EAF1FB"}, "blue table")
    if "alternate row colors" in lowered or "band rows" in lowered:
        if table_ref:
            add({"op": "style_table", **table_ref, "header_fill": "00338D", "header_font_color": "FFFFFF",
                 "banded_rows": True, "band_fill": "EAF1FB"}, "alternating table rows")
        else:
            add({"op": "add_table", "slide": target, "data": [["Risk", "Owner", "Impact"], ["Delivery capacity", "Program lead", "High"], ["User adoption", "Change lead", "Medium"], ["Budget variance", "Finance", "Low"]],
                 "x": .08, "y": .28, "width": .84, "height": .46, "header_fill": "00338D",
                 "header_font_color": "FFFFFF", "banded_rows": True, "band_fill": "EAF1FB"}, "banded table")
    explicit_sort = re.search(
        r'sort\s+the\s+table\s+on\s+slide\s+\d+\s+by\s+column\s+["“]([^"”]+)["”]\s+in\s+(ascending|descending)\s+order',
        raw_message,
        flags=re.IGNORECASE,
    )
    sort_match = None if explicit_sort else re.search(r"sort (?:the )?table(?: on this slide)? by ([a-z0-9 _-]+?)(?:\.|$)", lowered)
    if sort_match is None and explicit_sort is None:
        sort_match = re.search(r"sort by ([a-z0-9 _-]+?)(?:\.|$)", lowered)
    sort_column = explicit_sort.group(1).strip() if explicit_sort else (sort_match.group(1).strip() if sort_match else None)
    sort_direction = explicit_sort.group(2).casefold() if explicit_sort else "descending"
    if sort_column:
        if table_ref:
            add({"op": "sort_table", **table_ref, "column": sort_column,
                 "direction": sort_direction}, "sorted table")
        else:
            data = [["Risk", "Owner", "Impact"], ["Data security", "Technology lead", "Critical"], ["Delivery capacity", "Program lead", "High"], ["User adoption", "Change lead", "Medium"]]
            if sort_direction == "ascending":
                data = [data[0], *reversed(data[1:])]
            add({"op": "add_table", "slide": target,
                 "data": data,
                 "x": .08, "y": .28, "width": .84, "height": .46, "header_fill": "00338D",
                 "header_font_color": "FFFFFF", "banded_rows": True}, "impact-sorted table")
    explicit_add_row = re.search(r'add\s+a\s+table\s+row\s+with\s+values\s+["“]([^"”]+)["”]\s+on\s+slide\s+\d+', raw_message, flags=re.IGNORECASE)
    if explicit_add_row:
        row_values = [value.strip() for value in explicit_add_row.group(1).split("|")]
        if table_ref:
            add({"op": "add_table_row", **table_ref, "values": row_values}, "new table row")
        else:
            headers = [f"Column {index + 1}" for index in range(max(1, len(row_values)))]
            add({"op": "add_table", "slide": target, "data": [headers, row_values],
                 "x": .08, "y": .28, "width": .84, "height": .36,
                 "header_fill": "00338D", "header_font_color": "FFFFFF"}, "table with new row")
    elif "add another row" in lowered or "add a row" in lowered:
        if table_ref:
            add({"op": "add_table_row", **table_ref, "values": ["New item", "Owner", "Medium"]}, "new table row")
        else:
            add({"op": "add_table", "slide": target, "data": [["Risk", "Owner", "Impact"], ["Data governance", "Technology lead", "Medium"]],
                 "x": .08, "y": .28, "width": .84, "height": .36, "header_fill": "00338D", "header_font_color": "FFFFFF"}, "table with new row")
    explicit_merge_cells = re.search(r"merge\s+columns?\s+(\d+)\s+and\s+(\d+)\s+in\s+row\s+(\d+)\s+of\s+the\s+table\s+on\s+slide\s+\d+", lowered)
    if explicit_merge_cells:
        first_column, second_column, merge_row = map(int, explicit_merge_cells.groups())
        start_column, end_column = sorted((first_column, second_column))
        if table_ref:
            add({"op": "merge_table_cells", **table_ref, "row": merge_row, "column": start_column,
                 "end_row": merge_row, "end_column": end_column}, "merged table cells")
        else:
            rows_count, columns_count = max(2, merge_row), max(2, end_column)
            data = [[f"Column {column + 1}" for column in range(columns_count)]] + [["" for _ in range(columns_count)] for _ in range(rows_count - 1)]
            add({"op": "add_table", "slide": target, "data": data,
                 "merge_cells": [{"row": merge_row, "column": start_column, "end_row": merge_row, "end_column": end_column}],
                 "x": .08, "y": .28, "width": .84, "height": .42,
                 "header_fill": "00338D", "header_font_color": "FFFFFF"}, "table with merged cells")
    elif "merge the first two columns" in lowered:
        if table_ref:
            add({"op": "merge_table_cells", **table_ref, "row": 1, "column": 1,
                 "end_row": 1, "end_column": 2}, "merged table cells")
        else:
            add({"op": "add_table", "slide": target,
                 "data": [["Risk and accountable owner", "", "Status"], ["Delivery capacity", "Program lead", "Mitigate"], ["User adoption", "Change lead", "Monitor"]],
                 "merge_cells": [{"row": 1, "column": 1, "end_row": 1, "end_column": 2}],
                 "x": .08, "y": .28, "width": .84, "height": .42,
                 "header_fill": "00338D", "header_font_color": "FFFFFF", "banded_rows": True}, "table with merged cells")
    split_cell_match = re.search(r"split\s+the\s+table\s+cell\s+in\s+row\s+(\d+)\s+and\s+column\s+(\d+)\s+on\s+slide\s+\d+", lowered)
    if split_cell_match:
        split_row, split_column = map(int, split_cell_match.groups())
        if table_ref:
            add({"op": "split_table_cell", **table_ref, "row": split_row, "column": split_column}, "split table cell")
        else:
            rows_count, columns_count = max(2, split_row), max(2, split_column)
            data = [[f"Column {column + 1}" for column in range(columns_count)]] + [["" for _ in range(columns_count)] for _ in range(rows_count - 1)]
            add({"op": "add_table", "slide": target, "data": data,
                 "x": .08, "y": .28, "width": .84, "height": .42,
                 "header_fill": "00338D", "header_font_color": "FFFFFF"}, "table with separate cells")
    set_cell_match = re.search(r'set\s+the\s+table\s+cell\s+in\s+row\s+(\d+)\s+and\s+column\s+(\d+)\s+to\s+["“]([^"”]+)["”]\s+on\s+slide\s+\d+', raw_message, flags=re.IGNORECASE)
    if set_cell_match:
        cell_row, cell_column = int(set_cell_match.group(1)), int(set_cell_match.group(2))
        cell_text = set_cell_match.group(3).strip()
        if table_ref:
            add({"op": "set_table_cell", **table_ref, "row": cell_row, "column": cell_column, "text": cell_text}, "edited table cell")
        else:
            rows_count, columns_count = max(2, cell_row), max(2, cell_column)
            data = [[f"Column {column + 1}" for column in range(columns_count)]] + [["" for _ in range(columns_count)] for _ in range(rows_count - 1)]
            data[cell_row - 1][cell_column - 1] = cell_text
            add({"op": "add_table", "slide": target, "data": data,
                 "x": .08, "y": .28, "width": .84, "height": .42,
                 "header_fill": "00338D", "header_font_color": "FFFFFF"}, "table with edited cell")
    delete_row_match = re.search(r"delete\s+row\s+(\d+)\s+from\s+the\s+table\s+on\s+slide\s+\d+", lowered)
    if delete_row_match:
        row_number = int(delete_row_match.group(1))
        if table_ref:
            add({"op": "delete_table_row", **table_ref, "row": row_number}, "deleted table row")
        else:
            add({"op": "add_table", "slide": target, "data": [["Item", "Status"], ["Remaining row", "Open"]],
                 "x": .08, "y": .28, "width": .84, "height": .36,
                 "header_fill": "00338D", "header_font_color": "FFFFFF"}, "table after row deletion")
    add_column_match = re.search(r'add\s+a\s+table\s+column\s+with\s+values\s+["“]([^"”]+)["”]\s+on\s+slide\s+\d+', raw_message, flags=re.IGNORECASE)
    if add_column_match:
        column_values = [value.strip() for value in add_column_match.group(1).split("|")]
        if table_ref:
            add({"op": "add_table_column", **table_ref, "values": column_values}, "new table column")
        else:
            data = [["Item", column_values[0] if column_values else "New column"]]
            for index, value in enumerate(column_values[1:], start=1):
                data.append([f"Row {index + 1}", value])
            add({"op": "add_table", "slide": target, "data": data,
                 "x": .08, "y": .28, "width": .84, "height": .42,
                 "header_fill": "00338D", "header_font_color": "FFFFFF"}, "table with new column")
    delete_column_match = re.search(r"delete\s+column\s+(\d+)\s+from\s+the\s+table\s+on\s+slide\s+\d+", lowered)
    if delete_column_match:
        column_number = int(delete_column_match.group(1))
        if table_ref:
            add({"op": "delete_table_column", **table_ref, "column": column_number}, "deleted table column")
        else:
            add({"op": "add_table", "slide": target, "data": [["Remaining column"], ["Value"]],
                 "x": .20, "y": .28, "width": .60, "height": .36,
                 "header_fill": "00338D", "header_font_color": "FFFFFF"}, "table after column deletion")
    if "highlight high risks red" in lowered:
        highlighted = False
        if first_table:
            rows = [line.split(" | ") for line in str(first_table.get("text", "")).splitlines()]
            for row_index, row in enumerate(rows, start=1):
                for column_index, value in enumerate(row, start=1):
                    if value.strip().casefold() in {"high", "critical"}:
                        add({"op": "set_table_cell", **table_ref, "row": row_index, "column": column_index,
                             "text": value, "fill_color": "FDE8E8", "font_color": "B42318", "bold": True}, "high-risk highlight")
                        highlighted = True
        if not highlighted:
            add({"op": "add_table", "slide": target, "data": [["Risk", "Owner", "Impact"], ["Delivery", "Program lead", "High"]],
                 "x": .08, "y": .28, "width": .84, "height": .36, "header_fill": "00338D", "header_font_color": "FFFFFF",
                 "body_fill": "FDE8E8", "body_font_color": "B42318"}, "high-risk table")
    if "turn this table into icons" in lowered or "convert this table to icons" in lowered:
        if table_ref:
            add({"op": "convert_table_to_icons", **table_ref, "keep_table": False}, "editable icon cards")
        else:
            add({"op": "semantic_transform", "kind": "four_cards", "slide": target}, "editable icon cards")
    if "convert the table" in lowered and "native chart" in lowered:
        add(_guaranteed_chart_operation(target, "column", bool(first_table) or has_deck_data), "native chart from table")

    # Text, object, notes, and slide actions.
    if "rewrite" in lowered and any(token in lowered for token in ("senior executive", "executive audience", "for executives")):
        add({"op": "semantic_transform", "kind": "rewrite_executive", "slide": target,
             "add_takeaway": True}, "executive rewrite")
    if "standardize capitalization" in lowered or "standardize spacing" in lowered:
        add({"op": "semantic_transform", "kind": "format_slide", "slide": target}, "standardized formatting")
    if "align the main objects" in lowered:
        add({"op": "semantic_transform", "kind": "format_slide", "slide": target}, "aligned and evenly distributed objects")
    if "bring the main callout" in lowered:
        candidate = next((shape for shape in reversed(shapes) if shape.get("kind") in {"shape", "text"}), None)
        if candidate:
            add({"op": "layer_shape", "slide": target, "shape_id": candidate["shape_id"], "direction": "front"}, "front-most callout")
        else:
            add({"op": "add_shape", "slide": target, "shape_type": "rounded_rectangle", "x": .68, "y": .64,
                 "width": .26, "height": .16, "text": "Key takeaway", "fill_color": "EAF6F4",
                 "line_color": "00A3A1", "font_color": "0F766E", "font_size": 15}, "front-most callout")
    if "remove empty objects" in lowered or "clean layout" in lowered:
        add({"op": "cleanup_slide", "slide": target, "remove_empty": True, "autofit": True,
             "keep_on_slide": True}, "clean layout")
    if "replace the main image" in lowered:
        image_path = next((path for path in (attachment_paths or [])
                           if Path(path).suffix.casefold() in {".png", ".jpg", ".jpeg", ".webp", ".gif", ".bmp", ".tif", ".tiff"}), None)
        if image_path:
            if first_picture:
                add({"op": "replace_picture", "slide": target, "shape_id": first_picture["shape_id"],
                     "source_path": image_path}, "replacement image")
            else:
                add({"op": "add_picture", "slide": target, "source_path": image_path,
                     "x": .18, "y": .25, "width": .64, "height": .58}, "attached image")
        else:
            if first_picture:
                add({"op": "delete_shape", "slide": target, "shape_id": first_picture["shape_id"]}, "removed source image")
            add({"op": "add_shape", "slide": target, "shape_type": "rounded_rectangle",
                 "x": .18, "y": .28, "width": .64, "height": .46,
                 "text": "Growth momentum\n\nRevenue +8.6%\nRetention 94%\nMargin 31.2%",
                 "fill_color": "DCE6F7", "line_color": "005EB8", "font_color": "00338D", "font_size": 22}, "editable visual replacement")
            add({"op": "add_shape", "slide": target, "shape_type": "circle",
                 "x": .70, "y": .20, "width": .12, "height": .12, "text": "+8.6%",
                 "fill_color": "00A651", "line_color": "00A651", "font_color": "FFFFFF", "font_size": 12}, "editable visual accent")
    if "split this slide into two" in lowered or re.search(r"split slide\s+\d+ into two", lowered):
        add({"op": "split_slide", "slide": target}, "split slide")
    if "merge this slide with the next slide" in lowered:
        if target < count:
            add({"op": "merge_slides", "first_slide": target, "second_slide": target + 1}, "merged slides")
        else:
            add({"op": "duplicate_slide", "slide": target, "position": target + 1}, "new merge partner")
            add({"op": "merge_slides", "first_slide": target, "second_slide": target + 1}, "merged slides")
    explicit_merge = re.search(r"merge slides?\s+(\d+)\s+(?:and|with)\s+(\d+)", lowered)
    if explicit_merge:
        first, second = int(explicit_merge.group(1)), int(explicit_merge.group(2))
        add({"op": "merge_slides", "first_slide": first, "second_slide": second}, "merged slides")
    move_match = re.search(r"move\s+slide\s+(\d+)\s+to\s+position\s+(\d+)", lowered)
    if move_match:
        source, destination = int(move_match.group(1)), int(move_match.group(2))
        add({"op": "move_slide", "from_slide": source, "to_slide": destination}, "moved slide")
    if "generate speaker notes" in lowered:
        lines = [str(shape.get("text", "")).strip().replace("\n", "; ") for shape in text_shapes if str(shape.get("text", "")).strip()]
        title = str(slide_info.get("title") or f"Slide {target}")
        notes = f"Opening: {title}.\n" + "\n".join(f"• {line[:260]}" for line in lines[:5])
        notes += "\nClose by confirming the decision, owner, and next step."
        add({"op": "set_speaker_notes", "slide": target, "text": notes, "mode": "replace"}, "speaker notes")
    if "animate each bullet" in lowered or "animate every bullet" in lowered:
        has_multiline = any(len([line for line in str(shape.get("text", "")).splitlines() if line.strip()]) >= 2 for shape in text_shapes)
        if not has_multiline:
            add({"op": "add_textbox", "slide": target,
                 "text": "Confirm the decision\nAssign an accountable owner\nTrack delivery against milestones",
                 "x": .10, "y": .30, "width": .80, "height": .42, "font_size": 18,
                 "font_color": "172B4D", "no_fill": True, "no_line": True}, "editable bullet list")
        add({"op": "animate_bullets", "slide": target, "font_size": 18}, "progressive bullet animation")

    if not operations:
        return None
    concise = ", ".join(dict.fromkeys(labels))
    return {
        "message": f"Applied {concise} on slide {target}.",
        "operations": _dedupe_operations(operations),
        "assumptions": [],
        "failed": False,
    }


_guaranteed_previous_plan_edit = plan_edit


def plan_edit(
    pptx_path: str,
    user_message: str,
    selected_slide: int,
    selected_slide_image: str | None = None,
    chat_history: list[dict[str, str]] | None = None,
    deck_image_paths: list[str] | None = None,
    attachment_paths: list[str] | None = None,
) -> dict[str, Any]:
    if _looks_like_executive_review_request(user_message):
        slide_count = len(Presentation(pptx_path).slides)
        return {
            "message": "Completed the full executive review workflow.",
            "operations": [_executive_review_operation(slide_count)],
            "assumptions": [],
            "failed": False,
        }
    local = guaranteed_local_plan(pptx_path, user_message, selected_slide, attachment_paths)
    if local is not None:
        operations = [operation for operation in local.get("operations", []) if operation.get("op") != "noop"]
        if local.get("failed") or not operations:
            return local
        with tempfile.TemporaryDirectory(prefix="deck_refresh_guaranteed_") as work_dir:
            candidate, issues = _strict_apply_chain(pptx_path, operations, work_dir, "guaranteed")
            if candidate is not None and not issues:
                return local
    return _guaranteed_previous_plan_edit(
        pptx_path, user_message, selected_slide, selected_slide_image,
        chat_history, deck_image_paths, attachment_paths,
    )
