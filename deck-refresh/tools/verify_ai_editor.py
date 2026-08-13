"""Offline verification for the Deck Refresh AI editing engine."""

from __future__ import annotations

import json
import os
import sys
import tempfile
from pathlib import Path

ROOT = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(ROOT))

from openai.lib._pydantic import to_strict_json_schema
from pptx import Presentation

from pydantic import ValidationError

import pptx_editor as editor_module
from pptx_editor import (
    AtomicEditStep,
    EditPlan,
    EditorError,
    PlannerTask,
    TaskOutline,
    _deterministic_operations,
    _fallback_tasks,
    _parse_atomic_operation,
    apply_operations,
    deck_summary,
    validate_operations,
)


def check(condition: bool, message: str) -> None:
    if not condition:
        raise AssertionError(message)
    print(f"PASS: {message}")


def main() -> None:
    sample = ROOT / "sample_files" / "kpmg_advisory_q3_original.pptx"
    check(sample.exists(), "sample deck exists")

    schema = to_strict_json_schema(EditPlan)
    check(schema.get("additionalProperties") is False, "AI plan uses a strict schema")
    check(len(json.dumps(schema)) > 1000, "AI plan schema includes operation fields")
    replace_schema = schema["$defs"]["ReplaceTextOperation"]
    check(len(replace_schema.get("required", [])) == 10, "operation schema is compact instead of emitting every possible field")
    check(schema["properties"]["message"].get("maxLength") == 280, "AI summary length is capped to prevent runaway output")
    check(schema["properties"]["operations"].get("maxItems") == 4, "each structured response is capped to four operations")
    broad = EditPlan.model_validate({
        "message": "Updated the requested slides.",
        "assumptions": [],
        "operations": [
            {
                "op": "replace_text", "old": "INTERNAL DEMO", "new": "",
                "slide": 1, "shape_id": None, "shape_name": None,
                "text_contains": None, "shape": None,
                "replace_all": True, "case_sensitive": False,
            },
            {
                "op": "reorder_slides",
                "order": [1, 2, 3, 4, 5, 7, 6, 8, 9],
            },
        ],
    })
    check(len(broad.operations) == 2, "multi-part structured plans parse without oversized null-filled operations")
    try:
        EditPlan.model_validate({"message": "x" * 281, "assumptions": [], "operations": []})
        raise AssertionError("oversized message was accepted")
    except ValidationError:
        print("PASS: oversized model summaries are rejected before they consume the response budget")

    atomic = AtomicEditStep.model_validate({
        "message": "Removed the disclaimer.",
        "done": False,
        "operation_json": '{"op":"replace_text","slide":1,"old":"INTERNAL DEMO","new":""}',
    })
    parsed_atomic = _parse_atomic_operation(atomic.operation_json)
    check(parsed_atomic and parsed_atomic["op"] == "replace_text", "cutoff fallback parses one compact atomic operation")
    fallback = _fallback_tasks(
        "Remove the disclaimer on slide 1. Add a graph on slide 2. Reorder slides for flow.",
        1,
        9,
    )
    check(len(fallback) == 3 and fallback[-1].intent == "slide_order", "large requests split into focused fallback batches")

    summary = deck_summary(str(sample))
    check(summary["slide_count"] == 9, "deck inspection reads all slides")
    check("shape_id" in summary["slides"][0]["shapes"][0], "deck inspection includes stable shape IDs")
    deterministic = _deterministic_operations(
        summary,
        "Remove internal demo and all figures are fictional from slide 1.",
        1,
    )
    check(
        len(deterministic) >= 2 and all(op.get("slide") == 1 for op in deterministic),
        "clear text removals still work when every AI planning request is unavailable",
    )

    original_openai = editor_module.OpenAI
    original_outline = editor_module._request_outline
    original_plan = editor_module._request_plan
    original_atomic = editor_module._request_atomic_plan
    previous_key = os.environ.get("OPENAI_API_KEY")
    try:
        os.environ["OPENAI_API_KEY"] = "test-key"
        editor_module.OpenAI = lambda **_: object()
        editor_module._request_outline = lambda *_args, **_kwargs: TaskOutline(
            message="Updated slide 1.",
            assumptions=[],
            tasks=[PlannerTask(
                intent="remove_text",
                instruction="Remove the internal demo disclaimer from slide 1.",
                slides=[1],
                use_full_deck=False,
            )],
        )

        def _raise_cutoff(*_args, **_kwargs):
            raise EditorError("The AI response was cut off before the edit plan finished.")

        editor_module._request_plan = _raise_cutoff
        editor_module._request_atomic_plan = lambda *_args, **_kwargs: {
            "message": "Removed the disclaimer.",
            "assumptions": [],
            "operations": [{
                "op": "replace_text",
                "slide": 1,
                "old": "INTERNAL DEMO",
                "new": "",
                "replace_all": True,
                "case_sensitive": False,
            }],
        }
        cutoff_plan = editor_module.plan_edit(
            str(sample),
            "Remove the internal demo disclaimer from slide 1.",
            1,
        )
        check(
            cutoff_plan["operations"] and cutoff_plan["operations"][0]["op"] == "replace_text",
            "a max-output cutoff automatically switches to atomic planning",
        )
    finally:
        editor_module.OpenAI = original_openai
        editor_module._request_outline = original_outline
        editor_module._request_plan = original_plan
        editor_module._request_atomic_plan = original_atomic
        if previous_key is None:
            os.environ.pop("OPENAI_API_KEY", None)
        else:
            os.environ["OPENAI_API_KEY"] = previous_key


    exact_request = """Review the entire presentation and improve it for a senior executive audience. Preserve the KPMG branding, theme, fonts, logos, slide size, and overall visual identity.

On slide 1, remove every reference to internal demo, sample, and all figures are fictional. Rewrite the subtitle so it presents the deck as a client-ready Q3 advisory performance review. Keep the title concise and executive-focused.

On slide 2, analyze all existing numbers and labels. Create two native PowerPoint charts using only figures already present in the presentation. Use a green color palette. Add one chart showing the strongest positive trend and one chart showing the most important comparison between business areas. Add short insight headlines above each chart. Do not invent data. Reduce or reposition existing content to prevent overlap.

Review slides 3 and 4. Rewrite long paragraphs into concise executive bullets. Keep every important fact. Highlight the three most important findings using bold text or callout boxes. Remove repeated information across both slides.

Review slides 5, 6, 7, and 8. Determine their logical presentation order based on their titles, data, and narrative purpose. Reorder them so the flow moves from current performance, to drivers, to risks, to recommended actions. Do not rely only on the current slide numbers.

Create a new slide after the final analysis slide titled Executive Recommendations. Match the visual style of the existing deck. Add four recommendations based only on the presentation's existing findings. Each recommendation should include an action, a business reason, an owner type, and a suggested timing. Use a clean four-column or four-card layout.

Create another new slide titled Key Risks and Mitigations. Use a table with five rows and these columns: Risk, Impact, Likelihood, Mitigation, Owner. Derive the risks from the deck's existing content. Use red for high-risk indicators, amber for medium-risk indicators, and green for low-risk indicators.

Delete any slide that becomes fully redundant after the edits. Keep unique data and insights. Do not delete slides only because their writing is weak.

Standardize all slide titles, spacing, capitalization, font sizes, chart labels, and footer placement. Fix text overflow, overlapping objects, inconsistent alignment, and excessive empty space. Keep charts editable as native PowerPoint charts.

After finishing, verify the complete presentation flow. Make any additional layout or wording improvements needed for a polished client-ready deck."""

    previous_key_for_full = os.environ.pop("OPENAI_API_KEY", None)
    try:
        full_plan = editor_module.plan_edit(str(sample), exact_request, 1)
    finally:
        if previous_key_for_full is not None:
            os.environ["OPENAI_API_KEY"] = previous_key_for_full
    check(
        len(full_plan["operations"]) == 1 and full_plan["operations"][0]["op"] == "executive_review",
        "the exact complex executive request compiles to one local multi-pass workflow without model output risk",
    )

    with tempfile.TemporaryDirectory(prefix="deck_refresh_full_request_") as temp_dir:
        output = os.path.join(temp_dir, "full_request.pptx")
        full_result = apply_operations(str(sample), output, full_plan["operations"])
        check(not full_result["skipped"], "the exact complex request completes without skipped operations")
        prs = Presentation(output)
        check(len(prs.slides) == 11, "the workflow adds recommendations and risk slides")
        titles = []
        for slide in prs.slides:
            title = ""
            for shape in slide.shapes:
                if getattr(shape, "has_text_frame", False) and shape.text_frame.text.strip() and shape.text_frame.text.strip() != "KPMG":
                    title = shape.text_frame.text.strip().splitlines()[0]
                    break
            titles.append(title)
        check(titles[9] == "Executive Recommendations", "recommendations slide is created in the requested format")
        check(titles[10] == "Key Risks and Mitigations", "risk and mitigation slide is created")
        check(sum(1 for shape in prs.slides[1].shapes if getattr(shape, "has_chart", False)) == 2, "slide 2 receives two native charts")
        check(sum(1 for shape in prs.slides[10].shapes if getattr(shape, "has_table", False)) == 1, "risk slide contains an editable native table")
        check(
            [titles[index] for index in range(4, 8)] == [
                "Q3 Service Line Profitability", "Q3 Sector Performance", "Q3 Practice Costs", "Q3 Client Pipeline"
            ],
            "slides 5 through 8 are ordered by narrative role instead of original number",
        )
        all_text = "\n".join(
            shape.text_frame.text
            for slide in prs.slides
            for shape in slide.shapes
            if getattr(shape, "has_text_frame", False)
        ).casefold()
        check("internal demonstration" not in all_text and "all figures are fictional" not in all_text, "demo and fictional disclaimers are removed deck-wide")

    # Deliberately provide an invalid shape index with a semantic fallback.
    prepared, issues = validate_operations(
        str(sample),
        [{
            "op": "set_text",
            "slide": 1,
            "shape": 999,
            "text_contains": "All figures are fictional",
            "text": "",
        }],
    )
    check(not issues and prepared[0].get("shape_id") == 7, "semantic shape fallback repairs stale indexes")

    operations = [
        {"op": "replace_text", "slide": 1, "old": "INTERNAL DEMO", "new": ""},
        {"op": "replace_text", "slide": 1, "old": "All figures are fictional", "new": ""},
        {
            "op": "add_chart",
            "slide": 2,
            "chart_type": "bar",
            "title": "Revenue and Contribution Profit",
            "categories": ["Advisory Revenue", "Contribution Profit"],
            "series": [{"name": "Q3 FY2026 ($M)", "values": [486.2, 112.8]}],
            "x": 0.54,
            "y": 0.23,
            "width": 0.40,
            "height": 0.31,
            "series_colors": ["00A651"],
            "show_legend": False,
        },
        {
            "op": "add_chart",
            "slide": 2,
            "chart_type": "column",
            "title": "Operating KPIs",
            "categories": ["Margin", "Repeat Client Rate", "Utilization"],
            "series": [{"name": "Percent", "values": [23.2, 82.4, 74.8]}],
            "x": 0.54,
            "y": 0.57,
            "width": 0.40,
            "height": 0.31,
            "series_colors": ["66B032"],
            "show_legend": False,
        },
        {"op": "reorder_slides", "order": [1, 2, 3, 4, 5, 7, 6, 8, 9]},
    ]

    with tempfile.TemporaryDirectory(prefix="deck_refresh_ai_test_") as temp_dir:
        output = os.path.join(temp_dir, "edited.pptx")
        result = apply_operations(str(sample), output, operations)
        check(len(result["applied"]) == 5, "multi-part request applies every safe operation")
        check(not result["skipped"], "valid request has no skipped operations")

        prs = Presentation(output)
        check(len(prs.slides) == 9, "edited PowerPoint reopens successfully")
        slide1_text = "\n".join(
            shape.text for shape in prs.slides[0].shapes if getattr(shape, "has_text_frame", False)
        )
        check("INTERNAL" not in slide1_text and "All figures" not in slide1_text, "whitespace-tolerant disclaimer removal works")
        chart_count = sum(1 for shape in prs.slides[1].shapes if getattr(shape, "has_chart", False))
        check(chart_count == 2, "new green charts are embedded as native PowerPoint charts")
        check("Sector Performance" in prs.slides[5].shapes[0].text, "full-deck narrative reorder works")

    print("\nAll AI editor checks passed.")


if __name__ == "__main__":
    main()
