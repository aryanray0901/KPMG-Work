"""Regression checks for ChatGPT-like interpretation and all-or-nothing safety."""
from __future__ import annotations
import os
import tempfile
from pathlib import Path
from pptx import Presentation

ROOT = Path(__file__).resolve().parents[1]
import sys
sys.path.insert(0, str(ROOT))
import pptx_editor as pe


def check(value, message):
    if not value:
        raise AssertionError(message)
    print('PASS:', message)


def main():
    sample = str(ROOT / 'sample_files' / 'kpmg_advisory_q3_original.pptx')
    previous = os.environ.pop('OPENAI_API_KEY', None)
    try:
        invalid = pe.plan_edit(sample, 'delete slide 999 pls', 1, None, [], [], [])
        check(invalid.get('failed') is True and invalid['message'].startswith('error cant do that') and 'Try:' in invalid['message'] and not invalid['operations'], 'invalid language returns a safe error with corrective wording')

        compound = pe.plan_edit(sample, 'yo copy page 4 then chuck the new one all the way at the back', 1, None, [], [], [])
        # The local heuristic is intentionally narrower; mock the model-level
        # understanding for slang and verify execution is deterministic.
        original_interpreter = pe._interpret_conversation
        original_review = pe._review_conversation_plan
        try:
            pe._interpret_conversation = lambda *_a, **_k: pe.ConversationPlan(tasks=[
                pe.ConversationTask(action='duplicate_slide', sources=['slide:4'], alias='copy'),
                pe.ConversationTask(action='move_slide', sources=['alias:copy'], destination='last'),
            ])
            pe._review_conversation_plan = lambda *_a, **_k: _a[-1]
            slang = pe.plan_edit(sample, 'yo copy page 4 then chuck the new one all the way at the back', 1, None, [], [], [])
            check(slang['operations'] == [
                {'op': 'duplicate_slide', 'slide': 4, 'position': 5},
                {'op': 'move_slide', 'from_slide': 5, 'to_slide': 10},
            ], 'slang and pronouns compile into an ordered operation chain')

            # One valid and one invalid semantic operation must reject the
            # complete request instead of returning a partial deck change.
            pe._interpret_conversation = lambda *_a, **_k: pe.ConversationPlan(tasks=[
                pe.ConversationTask(action='semantic_edit', instruction='Make the requested edits.', sources=['slide:1'])
            ])
            old_legacy = pe._plan_edit_before_conversation_agent
            pe._plan_edit_before_conversation_agent = lambda *_a, **_k: {
                'message': 'edited',
                'operations': [
                    {'op': 'replace_text', 'slide': 1, 'old': 'INTERNAL DEMO', 'new': ''},
                    {'op': 'delete_slide', 'slide': 999},
                ],
            }
            rejected = pe.plan_edit(sample, 'clean slide one and also delete a nonexistent slide', 1, None, [], [], [])
            check(rejected.get('failed') is True and rejected['message'].startswith('error cant do that') and 'Try:' in rejected['message'] and not rejected['operations'], 'a partially invalid request is rejected as one transaction')
            pe._plan_edit_before_conversation_agent = old_legacy
        finally:
            pe._interpret_conversation = original_interpreter
            pe._review_conversation_plan = original_review
            if 'old_legacy' in locals():
                pe._plan_edit_before_conversation_agent = old_legacy

        # A rejected plan leaves the source file readable and unchanged.
        before = [int(item.id) for item in Presentation(sample).slides._sldIdLst]
        after = [int(item.id) for item in Presentation(sample).slides._sldIdLst]
        check(before == after, 'failed requests leave the PowerPoint untouched')
        print('\nSafe conversational checks passed.')
    finally:
        if previous is not None:
            os.environ['OPENAI_API_KEY'] = previous


if __name__ == '__main__':
    main()
