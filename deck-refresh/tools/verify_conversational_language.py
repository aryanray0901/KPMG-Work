"""Regression checks for conversational, compound, and referential slide language."""
from __future__ import annotations
import os, sys, tempfile
from pathlib import Path
ROOT = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(ROOT))
from pptx import Presentation
import pptx_editor as pe


def check(value, message):
    if not value:
        raise AssertionError(message)
    print('PASS:', message)


def title(path, number):
    return pe._slide_title(Presentation(path).slides[number - 1])


def main():
    sample = str(ROOT / 'sample_files' / 'kpmg_advisory_q3_original.pptx')
    previous = os.environ.pop('OPENAI_API_KEY', None)
    try:
        # Exact user phrasing must work even when the API is unavailable.
        request = 'duplicate slide 4 and then move that new slide to the last slide in the deck'
        plan = pe.plan_edit(sample, request, 1, None, [], [], [])
        check(plan['operations'] == [{'op': 'duplicate_slide', 'slide': 4, 'position': 10}], 'compound duplicate and pronoun move compiles into one exact operation')
        with tempfile.TemporaryDirectory(prefix='conversation_language_') as td:
            out = os.path.join(td, 'compound.pptx')
            result = pe.apply_operations(sample, out, plan['operations'])
            check(not result['skipped'] and len(Presentation(out).slides) == 10, 'compound request executes without skipped edits')
            check(title(out, 10) == title(sample, 4), 'the new duplicate reaches the final slide position')

        # Mock the semantic interpreter to test aliases and approximate titles
        # without making a network request.
        original_interpreter = pe._interpret_conversation
        try:
            def linked_plan(*args, **kwargs):
                return pe.ConversationPlan(tasks=[
                    pe.ConversationTask(action='duplicate_slide', sources=['slide:4'], alias='copy'),
                    pe.ConversationTask(action='move_slide', sources=['alias:copy'], destination='last'),
                ])
            pe._interpret_conversation = linked_plan
            linked = pe.plan_edit(sample, 'make a copy of the fourth page and stick that one at the back', 1, None, [], [], [])
            check(linked['operations'] == [
                {'op': 'duplicate_slide', 'slide': 4, 'position': 5},
                {'op': 'move_slide', 'from_slide': 5, 'to_slide': 10},
            ], 'alias references link separate conversational steps')

            def title_plan(*args, **kwargs):
                return pe.ConversationPlan(tasks=[
                    pe.ConversationTask(
                        action='move_slide',
                        sources=['title:client pipeline'],
                        destination='title:financial overview',
                        placement='before',
                    )
                ])
            pe._interpret_conversation = title_plan
            title_based = pe.plan_edit(sample, 'put the pipeline page before the financial page', 1, None, [], [], [])
            check(title_based['operations'] == [{'op': 'move_slide', 'from_slide': 6, 'to_slide': 3}], 'approximate slide-title references resolve semantically')

            def multilingual_plan(*args, **kwargs):
                return pe.ConversationPlan(tasks=[
                    pe.ConversationTask(action='delete_slide', sources=['last'])
                ])
            pe._interpret_conversation = multilingual_plan
            multilingual = pe.plan_edit(sample, 'elimina la última diapositiva por favor', 1, None, [], [], [])
            check(multilingual['operations'] == [{'op': 'delete_slide', 'slide': 9}], 'multilingual intent uses the same deterministic executor')
        finally:
            pe._interpret_conversation = original_interpreter

        print('\nConversational language checks passed.')
    finally:
        if previous is not None:
            os.environ['OPENAI_API_KEY'] = previous


if __name__ == '__main__':
    main()
