"""Regression suite for the 30-slide conversational PowerPoint editor."""
from __future__ import annotations

import sys
import tempfile
import types
from pathlib import Path

ROOT = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(ROOT))

try:
    import openai  # noqa: F401
except Exception:
    module = types.ModuleType("openai")
    module.OpenAI = type("OpenAI", (), {"__init__": lambda self, *args, **kwargs: None})
    sys.modules["openai"] = module

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from pptx_editor import apply_operations, plan_edit
from builtin_preview import render_pptx

SOURCE = ROOT / "sample_files" / "kpmg_advisory_30_slide_original.pptx"


def title(slide) -> str:
    candidates = [s for s in slide.shapes if getattr(s, "has_text_frame", False) and s.text_frame.text.strip()]
    candidates = [s for s in candidates if s.text_frame.text.strip().casefold() != "kpmg"]
    return max(candidates, key=lambda s: s.width).text_frame.text.strip().splitlines()[0] if candidates else ""


def run(command: str, selected=1):
    with tempfile.TemporaryDirectory(prefix="deck_refresh_test_") as td:
        output = Path(td) / "output.pptx"
        plan = plan_edit(str(SOURCE), command, selected)
        if plan.get("failed"):
            return plan, None, None
        result = apply_operations(str(SOURCE), str(output), plan.get("operations", []))
        assert not result.get("skipped"), (command, result)
        assert result.get("applied"), (command, result)
        prs = Presentation(output)
        preview_dir = Path(td) / "previews"
        preview_dir.mkdir()
        images = render_pptx(str(output), str(preview_dir), prefix="slide", canvas_width=480)
        assert len(images) == len(prs.slides), (command, len(images), len(prs.slides))
        assert all(Path(path).exists() for path in images)
        return plan, result, prs


def all_text(slide) -> str:
    parts = []
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False):
            parts.append(shape.text_frame.text)
        if getattr(shape, "has_table", False):
            for row in shape.table.rows:
                parts.extend(cell.text for cell in row.cells)
    return "\n".join(parts)


def test_all():
    assert len(Presentation(SOURCE).slides) == 30

    plan, result, prs = run("Delete the last slide in the deck.")
    assert len(prs.slides) == 29

    plan, result, prs = run("Duplicate slide 4, then move the new copy to the end of the presentation.")
    assert len(prs.slides) == 31
    assert title(prs.slides[-1]) == title(Presentation(SOURCE).slides[3])

    plan, result, prs = run("Move slide 5 to position 2.")
    assert len(prs.slides) == 30

    plan, result, prs = run("Swap slide 3 and slide 7.")
    assert len(prs.slides) == 30

    plan, result, prs = run("Rewrite slide 2 for a senior executive audience. Shorten the text, preserve every important fact, and add a one-sentence takeaway at the top.")
    assert "Takeaway:" in all_text(prs.slides[1])

    plan, result, prs = run("On slide 3, make the title 28 point, bold, and dark blue. Change the body text to dark gray. Align all objects evenly and fix any overlap.")
    assert "Strategic Context" in all_text(prs.slides[2])

    plan, result, prs = run("Turn the three most important findings on slide 4 into three editable callout boxes. Use green for positive performance, amber for concerns, and red for major risks.")
    callouts = [s for s in prs.slides[3].shapes if s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE and getattr(s, "has_text_frame", False)]
    assert len(callouts) >= 3

    plan, result, prs = run("Convert slide 5 into an editable table with the columns Finding, Impact, Priority, Owner, and Next Step. Use red for High, amber for Medium, and green for Low.")
    assert any(getattr(s, "has_table", False) and len(s.table.columns) == 5 for s in prs.slides[4].shapes)

    plan, result, prs = run("Review the data on slide 2. Create an editable clustered column chart showing the most important comparison. Use only existing numbers and do not invent data.")
    assert any(getattr(s, "has_chart", False) for s in prs.slides[1].shapes)

    plan, result, prs = run("Replace the table on slide 6 with an editable line chart using the same data. Keep the title and add data labels.")
    assert any(getattr(s, "has_chart", False) for s in prs.slides[5].shapes)
    assert not any(getattr(s, "has_table", False) for s in prs.slides[5].shapes)

    plan, result, prs = run("Create a new slide after slide 6 titled Executive Recommendations. Add four recommendation cards based on the deck’s existing findings. Each card should include an action, reason, owner, and timing.")
    assert len(prs.slides) == 31
    assert "Executive Recommendations" in all_text(prs.slides[6])

    plan, result, prs = run("Regenerate slide 4 as an executive summary with three sections titled What Changed, Why It Matters, and Recommended Action. Preserve all important facts.")
    summary_text = all_text(prs.slides[3])
    assert all(x in summary_text for x in ("What Changed", "Why It Matters", "Recommended Action"))

    plan, result, prs = run("Review the entire presentation. Standardize titles, fonts, capitalization, spacing, footers, chart labels, table formatting, and object alignment. Fix text overflow, overlaps, tiny text, and excessive empty space. Preserve the current branding.")
    assert len(prs.slides) == 30

    plan, result, prs = run("Delete the final slide, duplicate slide 3, move the duplicate to position 6, rewrite it as an executive summary, and add a green takeaway box at the top.")
    assert len(prs.slides) == 30
    assert "Executive Summary" in all_text(prs.slides[5])
    assert "Takeaway:" in all_text(prs.slides[5])

    plan, result, prs = run("The slide about risks looks messy. Clean it up, shorten the writing, make the major risks stand out, and keep the current KPMG style.")
    assert "Key Risks" in all_text(prs.slides[6])

    plan, result, prs = run("Duplicate slide 4 and move it to the end. Then change its title to Executive Summary and make it more concise.")
    assert len(prs.slides) == 31
    assert "Executive Summary" in all_text(prs.slides[-1])

    plan, result, prs = run("Delete the last slide. Then duplicate slide 3. Then move slide 4 to the end.")
    assert len(prs.slides) == 30

    failed = plan_edit(str(SOURCE), "Move slide 50 to position 2.", 1)
    assert failed.get("failed") and failed.get("operations") == []
    assert "error cant do that" in failed.get("message", "")

    with tempfile.TemporaryDirectory(prefix="deck_refresh_preview_") as td:
        images = render_pptx(str(SOURCE), td, prefix="slide", canvas_width=480)
        assert len(images) == 30
        assert all(Path(path).exists() for path in images)

    print("18/18 conversational editor tests passed")
    print("Every successful command rendered a complete refreshed preview set")


if __name__ == "__main__":
    test_all()
