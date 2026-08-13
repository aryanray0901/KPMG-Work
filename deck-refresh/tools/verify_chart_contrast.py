"""Regression test for automatic chart text contrast on dark and light slides."""

from __future__ import annotations

import shutil
import sys
import tempfile
from pathlib import Path

ROOT = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(ROOT))

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.util import Inches

from builtin_preview import render_pptx
from chart_contrast import _contrast_ratio, ensure_chart_contrast


def _rgb(font) -> tuple[int, int, int]:
    return tuple(bytes(font.color.rgb))


def _add_chart(slide, chart_type, title):
    data = CategoryChartData()
    data.categories = ["Revenue", "Margin", "Engagements", "Milestones", "Risks"]
    data.add_series("Q3 Actual", [486.2, 22.9, 1214, 84, 14])
    frame = slide.shapes.add_chart(chart_type, Inches(1), Inches(0.7), Inches(8), Inches(5), data)
    chart = frame.chart
    chart.has_title = True
    chart.chart_title.text_frame.text = title
    chart.has_legend = True
    plot = chart.plots[0]
    plot.has_data_labels = True
    plot.data_labels.show_value = True
    return frame


def main():
    temp_dir = tempfile.mkdtemp(prefix="deck_refresh_contrast_")
    try:
        prs = Presentation()
        dark_slide = prs.slides.add_slide(prs.slide_layouts[6])
        dark_slide.background.fill.solid()
        dark_slide.background.fill.fore_color.rgb = RGBColor(0x1E, 0x3A, 0x8A)
        _add_chart(dark_slide, XL_CHART_TYPE.PIE, "Dark background")

        light_slide = prs.slides.add_slide(prs.slide_layouts[6])
        light_slide.background.fill.solid()
        light_slide.background.fill.fore_color.rgb = RGBColor(0xF8, 0xFA, 0xFC)
        _add_chart(light_slide, XL_CHART_TYPE.COLUMN_CLUSTERED, "Light background")

        shape_slide = prs.slides.add_slide(prs.slide_layouts[6])
        rectangle = shape_slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
        )
        rectangle.fill.solid()
        rectangle.fill.fore_color.rgb = RGBColor(0x00, 0x4A, 0x98)
        rectangle.line.fill.background()
        _add_chart(shape_slide, XL_CHART_TYPE.LINE_MARKERS, "Background shape")

        assert ensure_chart_contrast(prs) == 3
        path = str(Path(temp_dir) / "chart_contrast.pptx")
        prs.save(path)
        reopened = Presentation(path)
        expected_backgrounds = [(30, 58, 138), (248, 250, 252), (0, 74, 152)]
        for slide, background in zip(reopened.slides, expected_backgrounds):
            chart_shape = next(shape for shape in slide.shapes if getattr(shape, "has_chart", False))
            chart = chart_shape.chart
            title_color = _rgb(chart.chart_title.text_frame.paragraphs[0].font)
            legend_color = _rgb(chart.legend.font)
            label_color = _rgb(chart.plots[0].data_labels.font)
            for color in (title_color, legend_color, label_color):
                assert _contrast_ratio(background, color) >= 4.5, (background, color)

        previews = render_pptx(path, temp_dir, "contrast", 1000)
        assert previews and len(previews) == 3
        assert all(Path(preview).stat().st_size > 1000 for preview in previews)
        print("PASS: dark, light, and full-slide background shapes produce readable native chart text.")
    finally:
        shutil.rmtree(temp_dir, ignore_errors=True)


if __name__ == "__main__":
    main()
