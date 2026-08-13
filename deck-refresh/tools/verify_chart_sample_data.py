#!/usr/bin/env python3
"""Verify every chart button imports the packaged Excel sample into a native chart."""

from __future__ import annotations

import sys
import tempfile
from pathlib import Path

from pptx import Presentation
from pptx.enum.chart import XL_LEGEND_POSITION

ROOT = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(ROOT))

from builder_ops import create_blank_deck, layout_operations, profile, read_frame
from pptx_editor import apply_operations


SAMPLE = ROOT / "sample_files" / "Deck-Refresh-Chart-Test-Data.xlsx"
EXPECTED_MONTHS = [
    "Jan 2026", "Feb 2026", "Mar 2026", "Apr 2026", "May 2026", "Jun 2026",
    "Jul 2026", "Aug 2026", "Sep 2026", "Oct 2026", "Nov 2026", "Dec 2026",
]
EXPECTED_REVENUE = [
    120000.0, 128000.0, 135000.0, 142000.0, 151000.0, 160000.0,
    172000.0, 181000.0, 193000.0, 207000.0, 222000.0, 238000.0,
]
CHART_LAYOUTS = {
    "column": "bar_chart",
    "bar": "bar_chart",
    "line": "line_chart",
    "pie": "pie_chart",
    "area": "area_chart",
    "waterfall": "waterfall_chart",
    "scatter": "scatter_plot",
}


def main() -> None:
    assert SAMPLE.exists(), SAMPLE
    frame = read_frame(str(SAMPLE))
    info = profile(frame)
    assert frame.shape == (13, 5)
    assert info["category_column"] == "Month"
    assert info["date_columns"] == ["Month"]
    assert info["numeric_columns"] == ["Revenue", "Operating Cost", "Operating Profit", "Budget Revenue"]
    assert info["currency_columns"] == info["numeric_columns"]
    assert info["total_rows"] == [12]
    assert info["data_kind"] == "time_series"
    assert info["insight"] == "Revenue increased 98.3% from the first to latest observation."

    with tempfile.TemporaryDirectory() as temp:
        root = Path(temp)
        for chart_type, layout in CHART_LAYOUTS.items():
            source = root / f"{chart_type}_source.pptx"
            output = root / f"{chart_type}_output.pptx"
            create_blank_deck(str(source))
            operations, detected, selected_layout = layout_operations(
                layout,
                2,
                data_path=str(SAMPLE),
                chart_type_override=chart_type,
            )
            assert detected == info
            assert selected_layout == layout
            assert operations[1]["number_format"] == '"$"#,##0'
            result = apply_operations(str(source), str(output), operations)
            assert not result["skipped"], (chart_type, result["skipped"])

            presentation = Presentation(str(output))
            chart = next(shape.chart for shape in presentation.slides[1].shapes if shape.has_chart)
            assert all(len(list(series.values)) == 12 for series in chart.series), chart_type
            assert chart.has_title
            expected_title = "Operating Cost vs Revenue" if chart_type == "scatter" else "Revenue"
            assert chart.chart_title.text_frame.text == expected_title

            if chart_type == "scatter":
                x_nodes = chart.series[0]._ser.xpath("./c:xVal/c:numRef/c:numCache/c:pt/c:v")
                assert [float(node.text) for node in x_nodes] == EXPECTED_REVENUE
                assert [series.name for series in chart.series] == ["Operating Cost", "Operating Profit", "Budget Revenue"]
                assert chart.has_legend
                assert all(series.marker.size == 7 for series in chart.series)
            else:
                assert [category.label for category in chart.plots[0].categories] == EXPECTED_MONTHS
                assert "Grand Total" not in [category.label for category in chart.plots[0].categories]
                if chart_type == "waterfall":
                    assert [series.name for series in chart.series] == ["Base", "Increase", "Decrease"]
                    base, increase, decrease = [list(series.values) for series in chart.series]
                    assert [a + b - c for a, b, c in zip(base, increase, decrease)] == EXPECTED_REVENUE
                    assert chart.series[0]._ser.xpath("./c:spPr/a:noFill")
                else:
                    assert list(chart.series[0].values) == EXPECTED_REVENUE, (chart_type, list(chart.series[0].values))
                    expected_series = 1 if chart_type == "pie" else 4
                    assert len(chart.series) == expected_series, chart_type
                if chart_type == "pie":
                    colors = [str(point.format.fill.fore_color.rgb) for point in chart.series[0].points]
                    assert len(set(colors)) == 12
                    assert chart.plots[0].has_data_labels
                    assert chart.plots[0].data_labels.show_percentage
                    assert chart.has_legend and chart.legend.position == XL_LEGEND_POSITION.RIGHT
                elif chart_type == "bar":
                    assert chart.has_legend and chart.legend.position == XL_LEGEND_POSITION.RIGHT
                elif chart_type == "line":
                    assert all(series.marker.size == 6 for series in chart.series)
                elif chart_type == "area":
                    assert all(series._ser.xpath("./c:spPr/a:solidFill/a:srgbClr/a:alpha") for series in chart.series)

        titled_ops, _, _ = layout_operations(
            "pie_chart", 2, title="Monthly Revenue", data_path=str(SAMPLE), chart_type_override="pie"
        )
        assert titled_ops[0]["title"] == "Monthly Revenue"
        assert titled_ops[1]["title"] == ""

    print("PASS: all 7 native chart types received 12 exact points, excluded totals, and retained chart-specific formatting.")


if __name__ == "__main__":
    main()
