#!/usr/bin/env python3
"""Regression check for every inspector command and new-slide template."""

from __future__ import annotations

import html
import os
import re
import sys
import tempfile
from pathlib import Path

from pptx import Presentation

ROOT = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(ROOT))

from builder_ops import CHART_LAYOUTS, LAYOUTS, create_blank_deck, layout_operations
from pptx_editor import _insertion_obstacle_box, apply_operations, guaranteed_local_plan, plan_edit
from showcase_engine import _is_branding, _title_shape


SAMPLE = ROOT / "sample_files" / "kpmg_advisory_30_slide_original.pptx"
TEMPLATE = ROOT / "templates" / "editor.html"


def apply_strict(source: str, operations: list[dict]) -> str:
    handle, output = tempfile.mkstemp(suffix=".pptx")
    os.close(handle)
    result = apply_operations(source, output, operations)
    if result.get("skipped") or len(result.get("applied", [])) != len(operations):
        os.remove(output)
        raise AssertionError(result)
    for item in result.get("applied", []):
        effect_values = [item[key] for key in ("count", "rows", "slides", "removed", "fitted") if key in item]
        if effect_values and not any(bool(value) for value in effect_values):
            os.remove(output)
            raise AssertionError(f"Operation reported no effect: {item}")
    return output


def verify_reported_chat_failures() -> None:
    commands = [
        "On slide 30, Add an editable native chart to this slide using data already in the deck.",
        "On slide 30, Rewrite this slide for a senior executive audience.",
    ]
    for command in commands:
        plan = plan_edit(str(SAMPLE), command, 30)
        assert plan and not plan.get("failed") and plan.get("operations"), plan
        output = apply_strict(str(SAMPLE), plan["operations"])
        try:
            if "chart" in command.casefold():
                slide = Presentation(output).slides[29]
                assert any(getattr(shape, "has_chart", False) for shape in slide.shapes)
        finally:
            os.remove(output)


def verify_inspector_commands() -> int:
    source = TEMPLATE.read_text(encoding="utf-8")
    commands = [html.unescape(value) for value in re.findall(r'data-prompt="([^"]+)"', source)]
    for command in commands:
        plan = guaranteed_local_plan(str(SAMPLE), command, 30)
        assert plan and not plan.get("failed") and plan.get("operations"), (command, plan)
        output = apply_strict(str(SAMPLE), plan["operations"])
        os.remove(output)
    return len(commands)


def verify_auto_fit() -> int:
    commands = [
        ("Add an editable native chart to this slide using data already in the deck.", "chart"),
        ("Add an editable native table to this slide.", "table"),
    ]
    for command, kind in commands:
        plan = guaranteed_local_plan(str(SAMPLE), command, 30)
        handle, output = tempfile.mkstemp(suffix=".pptx")
        os.close(handle)
        result = apply_operations(str(SAMPLE), output, plan["operations"])
        if result.get("skipped"):
            assert kind == "table", (kind, result)
            assert "no blank area" in str(result["skipped"][0].get("reason", "")).casefold(), result
            os.remove(output)
            continue
        assert len(result.get("applied", [])) == len(plan["operations"]), result
        try:
            prs = Presentation(output)
            slide = prs.slides[29]
            inserted = [
                shape for shape in slide.shapes
                if (kind == "chart" and getattr(shape, "has_chart", False))
                or (kind == "table" and getattr(shape, "has_table", False))
            ][-1]
            title = _title_shape(slide)

            def box(shape):
                return (
                    shape.left / prs.slide_width, shape.top / prs.slide_height,
                    shape.width / prs.slide_width, shape.height / prs.slide_height,
                )

            def intersection(first, second):
                width = max(0, min(first[0] + first[2], second[0] + second[2]) - max(first[0], second[0]))
                height = max(0, min(first[1] + first[3], second[1] + second[3]) - max(first[1], second[1]))
                return width * height

            inserted_box = box(inserted)
            for shape in slide.shapes:
                if int(shape.shape_id) == int(inserted.shape_id):
                    continue
                if title is not None and int(shape.shape_id) == int(title.shape_id):
                    continue
                if _is_branding(shape, prs):
                    continue
                shape_box = _insertion_obstacle_box(prs, shape)
                if shape_box is None:
                    continue
                area = shape_box[2] * shape_box[3]
                if 0.008 <= area <= 0.72 and shape_box[1] <= 0.91:
                    assert intersection(inserted_box, shape_box) <= 0.002, (kind, inserted_box, shape.name, shape_box)
        finally:
            os.remove(output)
    return len(commands)


def verify_layouts() -> int:
    handle, blank = tempfile.mkstemp(suffix=".pptx")
    os.close(handle)
    create_blank_deck(blank)
    try:
        for layout in LAYOUTS:
            operations, _, _ = layout_operations(layout, 2)
            output = apply_strict(blank, operations)
            try:
                created = Presentation(output).slides[1]
                visible = [
                    shape for shape in created.shapes
                    if getattr(shape, "has_chart", False)
                    or getattr(shape, "has_table", False)
                    or (getattr(shape, "has_text_frame", False) and shape.text_frame.text.strip())
                ]
                if layout != "blank":
                    assert visible, f"{layout} created an empty slide"
                if layout in CHART_LAYOUTS:
                    assert any(getattr(shape, "has_chart", False) for shape in created.shapes), layout
                text = "\n".join(
                    shape.text_frame.text for shape in created.shapes
                    if getattr(shape, "has_text_frame", False)
                ).casefold()
                forbidden = [
                    "tbd", "drop image", "import excel to populate", "add subtitle",
                    "add key message", "add supporting evidence", "add action",
                    "map canvas", "function a", "level 1", "primary analysis",
                    "what changed and why", "what to do next", "placeholder",
                ]
                assert not [phrase for phrase in forbidden if phrase in text], (layout, text)
                if layout in {"comparison", "kpi_dashboard"}:
                    assert any(getattr(shape, "has_chart", False) for shape in created.shapes), layout
                if layout in {"table", "financial_statement"}:
                    assert any(getattr(shape, "has_table", False) for shape in created.shapes), layout
            finally:
                os.remove(output)
    finally:
        os.remove(blank)
    return len(LAYOUTS)


def main() -> None:
    verify_reported_chat_failures()
    command_count = verify_inspector_commands()
    auto_fit_count = verify_auto_fit()
    layout_count = verify_layouts()
    print(f"PASS: 2 previously failing chat commands, {command_count} inspector commands, {auto_fit_count} collision-free inserts, and {layout_count} layouts.")


if __name__ == "__main__":
    main()
