"""Broad offline capability and failure-isolation test for Deck Refresh."""
from __future__ import annotations
import os, sys, tempfile
from pathlib import Path
ROOT=Path(__file__).resolve().parents[1]
sys.path.insert(0,str(ROOT))
from pptx import Presentation
from pptx_editor import apply_operations, deck_summary, SUPPORTED_OPS

def check(x,msg):
    if not x: raise AssertionError(msg)
    print('PASS:',msg)

def main():
    sample=ROOT/'sample_files'/'kpmg_advisory_q3_original.pptx'
    image=ROOT/'screenshots'/'data-slide.png'
    check(sample.exists() and image.exists(),'sample files exist')
    required={
        'append_text','format_text_box','rotate_shape','duplicate_shape','layer_shape','add_line',
        'add_table','style_table','merge_table_cells','split_table_cell','add_table_row',
        'delete_table_row','add_table_column','delete_table_column','add_picture','replace_picture',
        'crop_picture','style_chart','change_chart_type','set_speaker_notes','clear_slide',
        'regenerate_slide','align_shapes','distribute_shapes','standardize_deck','cleanup_slide',
        'set_slide_hidden','replace_color','delete_objects','convert_table_to_chart','apply_theme','set_footer','set_slide_size'
    }
    check(required <= SUPPORTED_OPS,'universal operation registry is complete')

    with tempfile.TemporaryDirectory(prefix='universal_editor_') as td:
        out=os.path.join(td,'edited.pptx')
        ops=[
            {'op':'replace_text','slide':1,'old':'INTERNAL DEMO','new':'CLIENT READY','allow_missing':True},
            {'op':'append_text','slide':1,'shape_id':3,'text':'Executive Edition','separator':' | ','font_size':25,'bold':True},
            {'op':'format_text_box','slide':1,'shape_id':3,'autofit':True,'word_wrap':True,'vertical_alignment':'middle'},
            {'op':'style_shape','slide':1,'shape_id':5,'fill_color':'00A651','line_color':'00338D','line_width':1.5},
            {'op':'rotate_shape','slide':1,'shape_id':5,'rotation':5},
            {'op':'duplicate_shape','slide':1,'shape_id':5},
            {'op':'layer_shape','slide':1,'shape_id':5,'direction':'front'},
            {'op':'add_textbox','slide':1,'text':'Universal editor test','x':0.1,'y':0.82,'width':0.35,'height':0.08,'font_size':12,'font_color':'00338D','no_fill':True,'no_line':True},
            {'op':'add_shape','slide':1,'shape_type':'rounded_rectangle','x':0.50,'y':0.82,'width':0.12,'height':0.08,'text':'A','fill_color':'E8F5E9'},
            {'op':'add_shape','slide':1,'shape_type':'rounded_rectangle','x':0.65,'y':0.82,'width':0.12,'height':0.08,'text':'B','fill_color':'E8F5E9'},
            {'op':'add_shape','slide':1,'shape_type':'rounded_rectangle','x':0.80,'y':0.82,'width':0.12,'height':0.08,'text':'C','fill_color':'E8F5E9'},
            {'op':'add_line','slide':1,'x1':0.08,'y1':0.78,'x2':0.92,'y2':0.78,'line_color':'00A651','line_width':2},
            {'op':'align_shapes','slide':1,'shape_ids':[9,10,11],'alignment':'top'},
            {'op':'distribute_shapes','slide':1,'shape_ids':[9,10,11],'direction':'horizontal'},
            {'op':'add_table','slide':2,'data':[['Metric','Value'],['Revenue','486.2'],['Profit','112.8']], 'x':0.05,'y':0.70,'width':0.40,'height':0.22,'header_fill':'00338D','banded_rows':True},
            {'op':'convert_table_to_chart','slide':2,'text_contains':'Metric | Value','chart_type':'line','title':'Table-derived chart','x':0.52,'y':0.42,'width':0.43,'height':0.22,'keep_table':True},
            {'op':'style_table','slide':2,'text_contains':'Metric | Value','header_fill':'00A651','header_font_color':'FFFFFF','banded_rows':True},
            {'op':'set_table_cell','slide':2,'text_contains':'Metric | Value','row':2,'column':2,'text':'$486.2M','font_color':'00338D','bold':True},
            {'op':'add_table_row','slide':2,'text_contains':'Metric | Value','values':['Margin','23.2%']},
            {'op':'add_table_column','slide':2,'text_contains':'Metric | Value','values':['Status','Good','Strong','On plan']},
            {'op':'delete_table_row','slide':2,'text_contains':'Metric | Value','row':4},
            {'op':'delete_table_column','slide':2,'text_contains':'Metric | Value','column':3},
            {'op':'merge_table_cells','slide':2,'text_contains':'Metric | Value','row':1,'column':1,'end_row':1,'end_column':2},
            {'op':'split_table_cell','slide':2,'text_contains':'Metric','row':1,'column':1},
            {'op':'add_chart','slide':2,'chart_type':'column','title':'KPI Comparison','categories':['Revenue','Profit'],'series':[{'name':'Q3','values':[486.2,112.8]}],'x':0.52,'y':0.68,'width':0.43,'height':0.25,'series_colors':['00A651'],'show_legend':False},
            {'op':'style_chart','slide':2,'text_contains':'KPI Comparison','show_data_labels':True,'show_gridlines':False,'series_colors':['66B032'],'number_format':'0.0'},
            {'op':'change_chart_type','slide':2,'text_contains':'KPI Comparison','chart_type':'bar','series_colors':['00A651']},
            {'op':'add_picture','slide':3,'source_path':str(image),'x':0.70,'y':0.72,'width':0.22,'height':0.20,'crop_left':0.02},
            {'op':'set_speaker_notes','slide':1,'text':'Present the executive headline first.','mode':'replace'},
            {'op':'add_slide','position':10,'title':'Generated Slide','body':['Initial content']},
            {'op':'regenerate_slide','slide':10,'title':'Regenerated Executive Slide','subtitle':'Built from a natural-language request','body':['Action: improve margin','Owner: Finance','Timing: Q4'],'layout':'cards'},
            {'op':'cleanup_slide','slide':10,'remove_empty':True,'autofit':True,'keep_on_slide':True},
            {'op':'set_slide_hidden','slide':9,'hidden':True},
            {'op':'replace_color','old_color':'00338D','new_color':'002060','slides':[10]},
            {'op':'set_footer','text':'KPMG | Confidential','slides':[10]},
            {'op':'apply_theme','slides':[10],'primary_color':'00338D','accent_color':'00A651','font_face':'Arial'},
            {'op':'delete_objects','slide':10,'object_types':['empty'],'preserve_branding':True,'preserve_title':True},
            {'op':'set_slide_size','width_inches':13.333,'height_inches':7.5,'scale_content':True},
            {'op':'standardize_deck'},
            # Deliberately bad operation must be isolated, not crash the deck.
            {'op':'set_text','slide':99,'shape_id':999,'text':'invalid'},
        ]
        result=apply_operations(str(sample),out,ops)
        check(len(result['applied']) >= 37,'broad operation set applies transactionally')
        check(len(result['skipped']) == 1 and 'Slide 99' in result['skipped'][0]['reason'],'invalid operation is isolated and reported')
        prs=Presentation(out)
        check(len(prs.slides)==10,'slide creation survives all operations')
        check(any(getattr(s,'has_table',False) for s in prs.slides[1].shapes),'editable table remains')
        check(any(getattr(s,'has_chart',False) for s in prs.slides[1].shapes),'editable chart remains')
        check(any(str(s.shape_type)=='PICTURE (13)' for s in prs.slides[2].shapes),'image insertion remains')
        check('Present the executive headline' in prs.slides[0].notes_slide.notes_text_frame.text,'speaker notes are editable')
        check('Regenerated Executive Slide' in '\n'.join(getattr(s,'text','') for s in prs.slides[9].shapes),'slide regeneration works')
        check(prs.slides[8]._element.get('show')=='0','slide hiding works')
        Presentation(out)
        print('\nUniversal editor checks passed.')

if __name__=='__main__': main()
