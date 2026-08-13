"""Verify Deck Refresh theme and color editing against the 30-slide sample."""
from __future__ import annotations

import os
import shutil
import tempfile
from pathlib import Path
import sys

ROOT = Path(__file__).resolve().parents[1]
if str(ROOT) not in sys.path:
    sys.path.insert(0, str(ROOT))

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.util import Inches

from showcase_engine import apply_showcase_operation, plan_showcase_command, _title_shape
from pptx_editor import guaranteed_local_plan


SOURCE = ROOT / "sample_files" / "kpmg_advisory_30_slide_original.pptx"


def _apply_prompt(prompt: str, selected_slide: int = 1) -> tuple[dict, Presentation]:
    plan = plan_showcase_command(str(SOURCE), prompt, selected_slide, [])
    assert plan is not None, f"No deterministic plan for: {prompt}"
    assert not plan.get("failed"), plan
    prs = Presentation(str(SOURCE))
    for operation in plan.get("operations", []):
        if operation.get("op") != "semantic_transform":
            raise AssertionError(f"Unexpected non-theme operation in theme test: {operation}")
        apply_showcase_operation(prs, operation)
    return plan, prs


def _background(prs: Presentation, slide_number: int) -> str | None:
    try:
        return str(prs.slides[slide_number - 1].background.fill.fore_color.rgb)
    except Exception:
        return None


def _title_color(prs: Presentation, slide_number: int) -> str | None:
    slide = prs.slides[slide_number - 1]
    title = _title_shape(slide)
    if title is None:
        return None
    for paragraph in title.text_frame.paragraphs:
        for run in paragraph.runs:
            if run.text.strip() and run.font.color.rgb is not None:
                return str(run.font.color.rgb)
    return None


def _first_chart_colors(prs: Presentation) -> list[str]:
    for slide in prs.slides:
        for shape in slide.shapes:
            if getattr(shape, "has_chart", False):
                result = []
                for series in shape.chart.series:
                    try:
                        result.append(str(series.format.fill.fore_color.rgb))
                    except Exception:
                        try:
                            result.append(str(series.format.line.color.rgb))
                        except Exception:
                            pass
                return result
    return []


def _first_table_header(prs: Presentation) -> list[str]:
    for slide in prs.slides:
        for shape in slide.shapes:
            if getattr(shape, "has_table", False):
                return [str(cell.fill.fore_color.rgb) for cell in shape.table.rows[0].cells]
    return []


def main() -> None:
    assert SOURCE.exists(), SOURCE

    dark_plan, dark = _apply_prompt("Apply the Executive Dark theme to the entire deck.")
    assert len(dark_plan["operations"][0]["slides"]) == 30
    assert _background(dark, 1) == "101827"
    assert _background(dark, 30) == "101827"
    assert _title_color(dark, 1) == "F8FAFC"
    assert _first_chart_colors(dark)[0] == "60A5FA"
    assert set(_first_table_header(dark)) == {"2563EB"}

    green_plan, green = _apply_prompt("Apply the Performance Green theme to slide 3.")
    assert green_plan["operations"][0]["slides"] == [3]
    assert _background(green, 3) == "F7FBF7"
    assert _background(green, 2) != "F7FBF7"

    custom_plan, custom = _apply_prompt(
        "Use a navy, teal, and light gray palette on slide 4 with a white background."
    )
    custom_op = custom_plan["operations"][0]
    assert custom_op["primary"] == "0B1F3A"
    assert custom_op["secondary"] == "0F766E"
    assert custom_op["accent"] == "E5E7EB"
    assert _background(custom, 4) == "FFFFFF"

    source_prs = Presentation(str(SOURCE))
    picker_slide_number = next(
        index for index, slide in enumerate(source_prs.slides, start=1)
        if any(getattr(shape, "has_chart", False) for shape in slide.shapes)
    )
    picker_plan = guaranteed_local_plan(
        str(SOURCE),
        f"Apply a custom theme to slide {picker_slide_number} with primary #112233, accent #44AA88, and background #F4F5F6. Preserve logos and status colors.",
        picker_slide_number,
    )
    picker_op = picker_plan["operations"][0]
    assert picker_op["slides"] == [picker_slide_number]
    assert picker_op["primary"] == "#112233"
    assert picker_op["accent"] == "#44aa88"
    assert picker_op["chart_colors"] == ["#112233", "#44aa88"]
    picker_prs = Presentation(str(SOURCE))
    apply_showcase_operation(picker_prs, picker_op)
    assert _background(picker_prs, picker_slide_number) == "F4F5F6"
    picker_chart = next(shape for shape in picker_prs.slides[picker_slide_number - 1].shapes if getattr(shape, "has_chart", False))
    assert str(picker_chart.chart.series[0].format.fill.fore_color.rgb) == "112233"

    body_plan, body = _apply_prompt(
        "Change the background to black and body text to white on slide 2."
    )
    assert body_plan["operations"][0]["background"] == "000000"
    assert body_plan["operations"][0]["body_color"] == "FFFFFF"
    assert _background(body, 2) == "000000"

    chart_plan, chart_prs = _apply_prompt(
        "Change all chart colors to green, dark blue, and light blue."
    )
    assert len(chart_plan["operations"][0]["slides"]) == 30
    assert _first_chart_colors(chart_prs)[0] == "2E7D32"

    invalid = plan_showcase_command(
        str(SOURCE), "Apply the Ocean theme to slide 50.", 1, []
    )
    assert invalid and invalid.get("failed")
    assert not invalid.get("operations")

    with tempfile.TemporaryDirectory(prefix="deck-theme-test-") as tmp:
        output = Path(tmp) / "themed.pptx"
        dark.save(str(output))
        reopened = Presentation(str(output))
        assert len(reopened.slides) == 30
        assert _background(reopened, 1) == "101827"

    compact = Presentation()
    slide = compact.slides.add_slide(compact.slide_layouts[6])
    transparent = slide.shapes.add_textbox(Inches(0.7), Inches(0.7), Inches(10), Inches(0.5))
    transparent.text = "Wide unfilled text container"
    transparent_fill_type = transparent.fill.type
    label = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(0.7), Inches(1.5), Inches(7.5), Inches(0.55),
    )
    label.text = "Short label"
    label.fill.solid()
    label.fill.fore_color.rgb = RGBColor(0, 51, 141)
    original_width = label.width
    apply_showcase_operation(compact, {
        "op": "semantic_transform", "kind": "apply_theme", "preset": "performance green", "slides": [1]
    })
    assert transparent.fill.type == transparent_fill_type
    assert label.width < original_width * 0.6

    print("Theme editor verification: 8/8 passed.")
    print("Presets, compact label fills, unfilled text protection, chart recoloring, slide scope, validation, and reopen checks passed.")


if __name__ == "__main__":
    main()
