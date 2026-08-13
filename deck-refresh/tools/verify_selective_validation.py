"""Verify checkbox-level approval for native text, chart, and table changes."""
from __future__ import annotations

import sys
import tempfile
from pathlib import Path

from pptx import Presentation

ROOT = Path(__file__).resolve().parents[1]
if str(ROOT) not in sys.path:
    sys.path.insert(0, str(ROOT))

from replacement_engine import (
    apply_selected_deck_replacements,
    compare_deck_replacements,
    replace_deck_1to1,
    validate_replacement_structure,
)


DECK = ROOT / "sample_files" / "Goldman Sachs Q3 Analysis.pptx"
DATA = ROOT / "sample_files" / "new goldman sachs Q4.xlsx"


def _located_value(prs, locator):
    shape = prs.slides[locator["slide_index"]].shapes[locator["shape_index"]]
    if locator["kind"] == "text_shape":
        return shape.text
    if locator["kind"] == "table_cell":
        return shape.table.cell(locator["row_index"], locator["column_index"]).text
    if locator["kind"] == "chart_point":
        return list(shape.chart.series[locator["series_index"]].values)[locator["point_index"]]
    raise AssertionError(locator)


def main():
    assert DECK.exists() and DATA.exists()
    with tempfile.TemporaryDirectory(prefix="deck-refresh-selective-") as temp_dir:
        pending = Path(temp_dir) / "pending.pptx"
        selective = Path(temp_dir) / "selective.pptx"
        replace_deck_1to1(DECK, DATA, pending, DATA.name)
        summary = compare_deck_replacements(DECK, pending)
        entries = summary["entries"]
        chosen = [
            next(entry for entry in entries if entry["object_type"] == "Text" and entry["slide"] == 1),
            next(entry for entry in entries if entry["object_type"] == "Chart"),
            next(entry for entry in entries if entry["object_type"] == "Table"),
        ]
        selected_ids = {entry["id"] for entry in chosen}
        applied = apply_selected_deck_replacements(DECK, pending, entries, selected_ids, selective)
        assert applied == 3
        assert validate_replacement_structure(DECK, selective)["passed"]

        original_prs = Presentation(DECK)
        pending_prs = Presentation(pending)
        selective_prs = Presentation(selective)
        for entry in chosen:
            locator = entry["locator"]
            assert _located_value(selective_prs, locator) == _located_value(pending_prs, locator)

        unchecked_table = next(
            entry for entry in entries
            if entry["object_type"] == "Table" and entry["id"] not in selected_ids
        )
        assert _located_value(selective_prs, unchecked_table["locator"]) == _located_value(original_prs, unchecked_table["locator"])
        slide_10_text = "\n".join(
            shape.text for shape in selective_prs.slides[9].shapes if getattr(shape, "has_text_frame", False)
        )
        assert "Q3 1:1 replacement ready" in slide_10_text
        assert "Q4 1:1 replacement ready" not in slide_10_text

        template = (ROOT / "templates" / "replacement_review.html").read_text(encoding="utf-8")
        assert 'name="mapping_id"' in template
        assert 'value="approve_selected"' in template
        assert "select-all-mappings" in template and "clear-all-mappings" in template

    print("Selective validation verification passed: text, chart, and table rows apply independently.")


if __name__ == "__main__":
    main()
