from pathlib import Path

from openpyxl import Workbook
from openpyxl.chart import BarChart, Reference
from openpyxl.styles import Alignment, Font, PatternFill
from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "sample_files"
OUT.mkdir(parents=True, exist_ok=True)

NAVY = RGBColor(0, 51, 141)
BLUE = RGBColor(0, 94, 184)
TEAL = RGBColor(0, 145, 218)
ORANGE = RGBColor(71, 10, 104)
GREEN = RGBColor(0, 163, 161)
RED = RGBColor(198, 33, 39)
LIGHT = RGBColor(243, 246, 250)
MID = RGBColor(216, 225, 238)
TEXT = RGBColor(31, 41, 55)
MUTED = RGBColor(100, 116, 139)
WHITE = RGBColor(255, 255, 255)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def set_cell_text(cell, text, size=12, bold=False, color=TEXT, align=PP_ALIGN.LEFT):
    cell.text = str(text)
    p = cell.text_frame.paragraphs[0]
    p.alignment = align
    p.font.name = "Aptos"
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE


def add_text(slide, text, x, y, w, h, size=18, bold=False, color=TEXT,
             align=PP_ALIGN.LEFT, font="Aptos", fill=None, margin=0.06):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is not None:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
        shape.line.fill.background()
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    p.font.name = font
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    return shape


def add_title(slide, title, subtitle=None):
    add_text(slide, title, 0.6, 0.35, 8.9, 0.5, size=24, bold=True, color=NAVY)
    if subtitle:
        add_text(slide, subtitle, 0.62, 0.86, 9.8, 0.32, size=10.5, color=MUTED)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(1.24), Inches(12.1), Inches(0.03))
    line.fill.solid(); line.fill.fore_color.rgb = BLUE; line.line.fill.background()


def add_footer(slide):
    add_text(slide, "KPMG | Internal demonstration | Fictional data", 0.62, 7.08, 6.0, 0.2, size=8.5, color=MUTED)


def add_kpi(slide, x, y, w, h, value, label, accent=BLUE):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    card.fill.solid(); card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = MID
    card.shadow.inherit = False
    stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(0.08), Inches(h))
    stripe.fill.solid(); stripe.fill.fore_color.rgb = accent; stripe.line.fill.background()
    tf = card.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.margin_left = Inches(0.18); tf.margin_right = Inches(0.12)
    tf.margin_top = Inches(0.10); tf.margin_bottom = Inches(0.08)
    p1 = tf.paragraphs[0]
    p1.text = value
    p1.font.name = "Aptos Display"; p1.font.size = Pt(24); p1.font.bold = True; p1.font.color.rgb = NAVY
    p2 = tf.add_paragraph()
    p2.text = label
    p2.font.name = "Aptos"; p2.font.size = Pt(10); p2.font.color.rgb = MUTED
    return card


def add_table(slide, x, y, w, h, headers, rows, widths=None):
    shape = slide.shapes.add_table(len(rows) + 1, len(headers), Inches(x), Inches(y), Inches(w), Inches(h))
    table = shape.table
    if widths:
        for i, width in enumerate(widths):
            table.columns[i].width = Inches(width)
    for c, header in enumerate(headers):
        cell = table.cell(0, c)
        cell.fill.solid(); cell.fill.fore_color.rgb = NAVY
        set_cell_text(cell, header, size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER if c else PP_ALIGN.LEFT)
    for r, row in enumerate(rows, start=1):
        for c, value in enumerate(row):
            cell = table.cell(r, c)
            cell.fill.solid(); cell.fill.fore_color.rgb = WHITE if r % 2 else LIGHT
            set_cell_text(cell, value, size=10.5, bold=(c == 0), color=TEXT,
                          align=PP_ALIGN.LEFT if c == 0 else PP_ALIGN.RIGHT)
    return shape


def add_chart(slide, x, y, w, h, title, categories, series, chart_type=XL_CHART_TYPE.COLUMN_CLUSTERED):
    data = ChartData()
    data.categories = categories
    for name, values in series:
        data.add_series(name, values)
    chart = slide.shapes.add_chart(chart_type, Inches(x), Inches(y), Inches(w), Inches(h), data).chart
    chart.has_title = True
    chart.chart_title.text_frame.text = title
    chart.chart_title.text_frame.paragraphs[0].font.name = "Aptos"
    chart.chart_title.text_frame.paragraphs[0].font.size = Pt(13)
    chart.chart_title.text_frame.paragraphs[0].font.bold = True
    chart.chart_title.text_frame.paragraphs[0].font.color.rgb = NAVY
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    chart.legend.font.name = "Aptos"
    chart.legend.font.size = Pt(9)
    chart.value_axis.has_major_gridlines = True
    chart.value_axis.major_gridlines.format.line.color.rgb = MID
    chart.value_axis.tick_labels.font.size = Pt(8)
    chart.category_axis.tick_labels.font.size = Pt(8)
    colors = [BLUE, TEAL, ORANGE, GREEN]
    for idx, ser in enumerate(chart.series):
        ser.format.fill.solid(); ser.format.fill.fore_color.rgb = colors[idx % len(colors)]
        ser.format.line.color.rgb = colors[idx % len(colors)]
    return chart


def build_deck(path: Path):
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]

    # Slide 1
    slide = prs.slides.add_slide(blank)
    bg = slide.background.fill; bg.solid(); bg.fore_color.rgb = NAVY
    add_text(slide, "KPMG", 0.72, 0.55, 3.0, 0.45, size=18, bold=True, color=WHITE, font="Arial")
    add_text(slide, "Q3 FY2026 Advisory Practice Review", 0.72, 2.05, 9.5, 1.0, size=32, bold=True, color=WHITE)
    add_text(slide, "Fictional KPMG data for testing presentation updates, native tables, charts, formatting, and synchronized comparison.",
             0.75, 3.10, 8.9, 1.0, size=15, color=RGBColor(215, 228, 245))
    block = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(10.55), Inches(0), Inches(2.78), Inches(7.5))
    block.fill.solid(); block.fill.fore_color.rgb = BLUE; block.line.fill.background()
    add_text(slide, "INTERNAL\nDEMO", 10.87, 2.65, 2.1, 1.3, size=25, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, "All figures are fictional", 10.82, 4.25, 2.2, 0.55, size=10, color=WHITE, align=PP_ALIGN.CENTER)

    # Slide 2
    slide = prs.slides.add_slide(blank)
    add_title(slide, "Q3 Executive Summary", "Fictional advisory practice KPI snapshot")
    kpis = [
        ("$486.2M", "Advisory Revenue", BLUE),
        ("$112.8M", "Contribution Profit", TEAL),
        ("23.2%", "Contribution Margin", ORANGE),
        ("1,148", "Active Engagements", BLUE),
        ("82.4%", "Repeat Client Rate", TEAL),
        ("74.8%", "Billable Utilization", ORANGE),
    ]
    for idx, (value, label, accent) in enumerate(kpis):
        col = idx % 3; row = idx // 3
        add_kpi(slide, 0.7 + col * 4.15, 1.58 + row * 2.28, 3.72, 1.75, value, label, accent)
    add_footer(slide)

    # Slide 3
    slide = prs.slides.add_slide(blank)
    add_title(slide, "Q3 Financial Overview", "Advisory practice performance versus budget")
    fin_rows = [
        ["Advisory Revenue", "$486.2M", "$472.0M"],
        ["Client Delivery Costs", "$238.7M", "$232.0M"],
        ["Gross Contribution", "$247.5M", "$240.0M"],
        ["Practice Expenses", "$134.7M", "$133.5M"],
        ["Contribution Profit", "$112.8M", "$106.5M"],
    ]
    add_table(slide, 0.65, 1.55, 5.25, 4.72, ["Metric", "Q3 Actual", "Q3 Budget"], fin_rows, [2.25, 1.5, 1.5])
    add_chart(slide, 6.18, 1.55, 6.48, 4.72, "Q3 Actual vs Q3 Budget",
              ["Advisory Revenue", "Gross Contribution", "Contribution Profit"],
              [("Q3 Actual", [486.2, 247.5, 112.8]), ("Q3 Budget", [472.0, 240.0, 106.5])])
    add_footer(slide)

    # Slide 4
    slide = prs.slides.add_slide(blank)
    add_title(slide, "Q3 Revenue by Service Line", "Fictional KPMG advisory service-line performance")
    service_lines = ["Deal Advisory", "Risk Consulting", "Technology Consulting", "Strategy & Operations"]
    service_actual = [142.6, 128.4, 121.7, 93.5]
    service_budget = [136.0, 126.5, 119.0, 90.5]
    service_rows = [[name, f"${actual:.1f}M", f"${budget:.1f}M"] for name, actual, budget in zip(service_lines, service_actual, service_budget)]
    add_table(slide, 0.65, 1.55, 5.3, 4.55, ["Service Line", "Q3 Actual Revenue", "Q3 Budget Revenue"], service_rows, [2.1, 1.6, 1.6])
    add_chart(slide, 6.18, 1.55, 6.48, 4.55, "Q3 Service Line Revenue", service_lines,
              [("Q3 Actual Revenue", service_actual), ("Q3 Budget Revenue", service_budget)])
    add_footer(slide)

    # Slide 5
    slide = prs.slides.add_slide(blank)
    add_title(slide, "Q3 Service Line Profitability", "Delivery and contribution margins by advisory service line")
    delivery_margin = [56.4, 52.8, 49.6, 45.9]
    contribution_margin = [27.2, 23.1, 20.4, 18.3]
    margin_rows = [[name, f"{delivery:.1f}%", f"{contribution:.1f}%"] for name, delivery, contribution in zip(service_lines, delivery_margin, contribution_margin)]
    add_table(slide, 0.65, 1.55, 5.3, 4.55, ["Service Line", "Q3 Delivery Margin", "Q3 Contribution Margin"], margin_rows, [2.1, 1.6, 1.6])
    add_chart(slide, 6.18, 1.55, 6.48, 4.55, "Q3 Service Line Margins", service_lines,
              [("Q3 Delivery Margin", delivery_margin), ("Q3 Contribution Margin", contribution_margin)], chart_type=XL_CHART_TYPE.LINE_MARKERS)
    add_footer(slide)

    # Slide 6
    slide = prs.slides.add_slide(blank)
    add_title(slide, "Q3 Client Pipeline", "Engagement growth and pursuit efficiency")
    add_kpi(slide, 0.7, 1.48, 2.75, 1.42, "286", "New Engagements", BLUE)
    add_kpi(slide, 3.65, 1.48, 2.75, 1.42, "38.6%", "Proposal Win Rate", TEAL)
    add_kpi(slide, 6.6, 1.48, 2.75, 1.42, "$18,450", "Average Pursuit Cost", ORANGE)
    add_kpi(slide, 9.55, 1.48, 2.75, 1.42, "$42.8M", "Expansion Revenue", GREEN)
    channels = ["Existing Clients", "Strategic Partners", "Digital Campaigns", "Executive Events"]
    new_engagements = [118, 72, 61, 35]
    opportunities = [284, 208, 196, 102]
    pursuit_cost = [10400, 16800, 22100, 27600]
    channel_rows = [[channel, str(new), str(opps), f"${cost:,}"] for channel, new, opps, cost in zip(channels, new_engagements, opportunities, pursuit_cost)]
    add_table(slide, 0.65, 3.22, 6.0, 3.0,
              ["Channel", "Q3 New Engagements", "Q3 Qualified Opportunities", "Q3 Average Pursuit Cost"],
              channel_rows, [1.55, 1.4, 1.4, 1.65])
    add_chart(slide, 6.92, 3.22, 5.74, 3.0, "Q3 Engagement Funnel", channels,
              [("Q3 New Engagements", new_engagements), ("Q3 Qualified Opportunities", opportunities)])
    add_footer(slide)

    # Slide 7
    slide = prs.slides.add_slide(blank)
    add_title(slide, "Q3 Sector Performance", "Revenue and contribution by client sector")
    sectors = ["Financial Services", "Consumer & Retail", "Healthcare", "Technology"]
    sector_revenue = [156.4, 118.7, 109.6, 101.5]
    sector_contribution = [39.8, 26.1, 24.2, 22.7]
    sector_margin = [25.4, 22.0, 22.1, 22.4]
    sector_rows = [[sector, f"${revenue:.1f}M", f"${contribution:.1f}M", f"{margin:.1f}%"] for sector, revenue, contribution, margin in zip(sectors, sector_revenue, sector_contribution, sector_margin)]
    add_table(slide, 0.65, 1.55, 6.0, 4.65, ["Sector", "Q3 Revenue", "Q3 Contribution", "Q3 Margin"], sector_rows, [2.1, 1.3, 1.3, 1.3])
    add_chart(slide, 6.92, 1.55, 5.74, 4.65, "Q3 Revenue and Contribution", sectors,
              [("Q3 Revenue", sector_revenue), ("Q3 Contribution", sector_contribution)])
    add_footer(slide)

    # Slide 8
    slide = prs.slides.add_slide(blank)
    add_title(slide, "Q3 Practice Costs", "Fictional advisory practice spend versus plan")
    functions = ["Business Development", "Technology Enablement", "Talent & Learning", "Practice Operations"]
    actual_costs = [39.8, 32.6, 29.4, 32.9]
    budget_costs = [38.5, 31.8, 30.2, 33.0]
    cost_rows = [[function, f"${actual:.1f}M", f"${budget:.1f}M"] for function, actual, budget in zip(functions, actual_costs, budget_costs)]
    add_table(slide, 0.65, 1.55, 5.3, 4.55, ["Function", "Q3 Actual Costs", "Q3 Budget Costs"], cost_rows, [2.1, 1.6, 1.6])
    add_chart(slide, 6.18, 1.55, 6.48, 4.55, "Q3 Practice Costs", functions,
              [("Q3 Actual Costs", actual_costs), ("Q3 Budget Costs", budget_costs)])
    add_footer(slide)

    # Slide 9
    slide = prs.slides.add_slide(blank)
    add_title(slide, "Q3 Outlook", "Forward indicators and fictional management forecast")
    outlook = [
        ("$522.0M", "Forecast Revenue", BLUE),
        ("$126.5M", "Forecast Contribution", TEAL),
        ("24.2%", "Forecast Contribution Margin", ORANGE),
        ("3.1", "Pipeline Coverage", BLUE),
        ("$318.0M", "Contracted Backlog", TEAL),
        ("78.5%", "Forecast Utilization", ORANGE),
    ]
    for idx, (value, label, accent) in enumerate(outlook):
        col = idx % 3; row = idx // 3
        add_kpi(slide, 0.7 + col * 4.15, 1.58 + row * 2.28, 3.72, 1.75, value, label, accent)
    add_footer(slide)

    prs.save(path)

def style_sheet(ws):
    for cell in ws[1]:
        cell.font = Font(bold=True, color="FFFFFF")
        cell.fill = PatternFill("solid", fgColor="00338D")
        cell.alignment = Alignment(horizontal="center")
    for col in ws.columns:
        width = max(len(str(c.value)) if c.value is not None else 0 for c in col) + 3
        ws.column_dimensions[col[0].column_letter].width = min(max(width, 14), 34)
    ws.freeze_panes = "A2"


def add_sheet(wb, title, headers, rows):
    ws = wb.create_sheet(title)
    ws.append(headers)
    for row in rows:
        ws.append(row)
    style_sheet(ws)
    return ws


def build_data(path: Path):
    wb = Workbook()
    wb.remove(wb.active)

    add_sheet(wb, "Executive Summary", ["Metric", "Q4 Actual"], [
        ["Advisory Revenue", 522.8],
        ["Contribution Profit", 129.6],
        ["Contribution Margin", 24.8],
        ["Active Engagements", 1286],
        ["Repeat Client Rate", 84.9],
        ["Billable Utilization", 78.6],
    ])

    add_sheet(wb, "Financial Overview", ["Metric", "Q4 Actual", "Q4 Budget"], [
        ["Advisory Revenue", 522.8, 510.0],
        ["Client Delivery Costs", 254.4, 249.0],
        ["Gross Contribution", 268.4, 261.0],
        ["Practice Expenses", 138.8, 137.0],
        ["Contribution Profit", 129.6, 124.0],
    ])

    add_sheet(wb, "Service Line Revenue", ["Service Line", "Q4 Actual Revenue", "Q4 Budget Revenue"], [
        ["Deal Advisory", 154.7, 149.0],
        ["Risk Consulting", 139.2, 136.0],
        ["Technology Consulting", 131.8, 128.0],
        ["Strategy & Operations", 97.1, 97.0],
    ])

    add_sheet(wb, "Service Line Profitability", ["Service Line", "Q4 Delivery Margin", "Q4 Contribution Margin"], [
        ["Deal Advisory", 58.2, 29.1],
        ["Risk Consulting", 54.1, 25.0],
        ["Technology Consulting", 51.3, 22.2],
        ["Strategy & Operations", 47.4, 19.7],
    ])

    pipeline_ws = add_sheet(wb, "Client Pipeline", ["Channel", "Q4 New Engagements", "Q4 Qualified Opportunities", "Q4 Average Pursuit Cost"], [
        ["Existing Clients", 132, 318, 9600],
        ["Strategic Partners", 83, 226, 15900],
        ["Digital Campaigns", 74, 224, 20700],
        ["Executive Events", 39, 112, 25300],
    ])
    chart = BarChart()
    chart.type = "col"
    chart.title = "Q4 Engagement Funnel"
    chart.y_axis.title = "Count"
    chart.x_axis.title = "Channel"
    data = Reference(pipeline_ws, min_col=2, max_col=3, min_row=1, max_row=5)
    cats = Reference(pipeline_ws, min_col=1, min_row=2, max_row=5)
    chart.add_data(data, titles_from_data=True)
    chart.set_categories(cats)
    chart.height = 7
    chart.width = 13
    pipeline_ws.add_chart(chart, "F2")

    add_sheet(wb, "Client KPIs", ["Metric", "Q4 Actual"], [
        ["New Engagements", 328],
        ["Proposal Win Rate", 41.3],
        ["Average Pursuit Cost", 17600],
        ["Expansion Revenue", 49.6],
    ])

    add_sheet(wb, "Sector Performance", ["Sector", "Q4 Revenue", "Q4 Contribution", "Q4 Margin"], [
        ["Financial Services", 168.2, 44.9, 26.7],
        ["Consumer & Retail", 127.5, 30.6, 24.0],
        ["Healthcare", 118.4, 28.1, 23.7],
        ["Technology", 108.7, 26.0, 23.9],
    ])

    add_sheet(wb, "Practice Costs", ["Function", "Q4 Actual Costs", "Q4 Budget Costs"], [
        ["Business Development", 41.2, 40.5],
        ["Technology Enablement", 34.7, 34.0],
        ["Talent & Learning", 30.8, 31.0],
        ["Practice Operations", 32.1, 31.5],
    ])

    add_sheet(wb, "Outlook", ["Metric", "Q4 Forecast"], [
        ["Forecast Revenue", 558.0],
        ["Forecast Contribution", 143.5],
        ["Forecast Contribution Margin", 25.7],
        ["Pipeline Coverage", 3.5],
        ["Contracted Backlog", 352.4],
        ["Forecast Utilization", 81.2],
    ])

    wb.save(path)


if __name__ == "__main__":
    original = OUT / "kpmg_advisory_q3_original.pptx"
    data = OUT / "kpmg_advisory_q4_data.xlsx"
    build_deck(original)
    build_data(data)
    print(original)
    print(data)
