"""Checks technical chart setup, the broader layout engine, and native objects."""

import csv
import sys
import tempfile
from pathlib import Path

ROOT = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(ROOT))

from pptx import Presentation
from builder_ops import LAYOUTS, create_blank_deck, layout_operations
from pptx_editor import apply_operations

XLSX_SOURCE = ROOT / "sample_files" / "kpmg_advisory_30_slide_data.xlsx"


def main():
    with tempfile.TemporaryDirectory() as temp:
        root = Path(temp)
        data = root / "monthly_revenue.csv"
        with data.open("w", newline="", encoding="utf-8") as handle:
            writer = csv.writer(handle)
            writer.writerow(["Month", "Revenue", "Margin"])
            writer.writerows([["Jan", 100, "20%"], ["Feb", 112, "22%"], ["Mar", 125, "24%"]])

        current = root / "start.pptx"
        create_blank_deck(str(current))

        blank_chart_ops, blank_profile, _ = layout_operations(
            "bar_chart", 2, chart_type_override="column"
        )
        assert blank_profile is None
        assert [operation["op"] for operation in blank_chart_ops] == ["add_slide", "add_chart"]
        assert blank_chart_ops[0]["title"] == ""
        assert blank_chart_ops[1]["chart_type"] == "column"
        assert blank_chart_ops[1]["series"][0]["values"] == [0, 0, 0]

        data_chart_ops, data_profile, _ = layout_operations(
            "pie_chart", 2, data_path=str(data), chart_type_override="pie"
        )
        assert data_profile and data_profile["category_column"] == "Month"
        assert data_chart_ops[1]["chart_type"] == "pie"
        assert len(data_chart_ops[1]["series"]) == 1

        xlsx_chart_ops, xlsx_profile, _ = layout_operations(
            "line_chart", 2, data_path=str(XLSX_SOURCE), chart_type_override="line"
        )
        assert xlsx_profile and xlsx_profile["category_column"] == "Metric"
        assert xlsx_chart_ops[1]["chart_type"] == "line"
        assert len(xlsx_chart_ops[1]["series"]) == 2
        for index, layout in enumerate(LAYOUTS, start=1):
            target = root / f"layout_{index}.pptx"
            data_path = str(data) if layout in {"table", "bar_chart", "line_chart", "pie_chart", "area_chart",
                                                        "waterfall_chart", "scatter_plot", "kpi_dashboard",
                                                        "financial_statement", "gantt_chart", "comparison"} else None
            operations, _, _ = layout_operations(layout, index + 1, layout.replace("_", " ").title(), "Alpha\nBeta\nGamma", data_path)
            result = apply_operations(str(current), str(target), operations)
            assert not result["skipped"], (layout, result["skipped"])
            current = target

        prs = Presentation(str(current))
        assert len(prs.slides) == len(LAYOUTS) + 1
        assert any(shape.has_chart for slide in prs.slides for shape in slide.shapes)
        assert any(shape.has_table for slide in prs.slides for shape in slide.shapes)

        chart_slide = next(index for index, slide in enumerate(prs.slides, 1) if any(shape.has_chart for shape in slide.shapes))
        chart_shape = next(shape for shape in prs.slides[chart_slide - 1].shapes if shape.has_chart)
        converted = root / "converted.pptx"
        result = apply_operations(str(current), str(converted), [
            {"op": "change_chart_type", "slide": chart_slide, "shape_id": chart_shape.shape_id, "chart_type": "line"},
            {"op": "style_all_charts", "color": "00A651"},
            {"op": "round_all_corners", "slides": [chart_slide]},
            {"op": "ensure_min_font", "slides": [chart_slide], "minimum": 18},
        ])
        assert not result["skipped"], result["skipped"]
        Presentation(str(converted))
        print(f"PASS: blank or Excel chart setup, {len(LAYOUTS)} internal layouts, chart reuse, and inspector operations")


if __name__ == "__main__":
    main()
