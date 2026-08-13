from pathlib import Path
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt
from openpyxl import Workbook

ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "sample_files"
OUT.mkdir(exist_ok=True)

BLUE = RGBColor(0x00, 0x33, 0x8D)
DARK_BLUE = RGBColor(0x00, 0x1E, 0x5A)
GREEN = RGBColor(0x2E, 0x7D, 0x32)
AMBER = RGBColor(0xF9, 0xA8, 0x25)
RED = RGBColor(0xC6, 0x28, 0x28)
DARK = RGBColor(0x37, 0x41, 0x51)
LIGHT = RGBColor(0xF3, 0xF4, 0xF6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]


def add_brand(slide, number):
    bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.12))
    bar.fill.solid(); bar.fill.fore_color.rgb = BLUE; bar.line.fill.background()
    logo = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.25), Inches(0.24), Inches(0.74), Inches(0.34))
    logo.fill.solid(); logo.fill.fore_color.rgb = BLUE; logo.line.fill.background()
    logo.text_frame.text = "KPMG"
    p = logo.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    for r in p.runs: r.font.size = Pt(15); r.font.bold = True; r.font.color.rgb = WHITE
    footer = slide.shapes.add_textbox(Inches(0.35), Inches(7.12), Inches(12.55), Inches(0.22))
    footer.text_frame.text = f"Advisory Transformation Steering Committee | Confidential | {number}"
    p = footer.text_frame.paragraphs[0]
    for r in p.runs: r.font.size = Pt(8.5); r.font.color.rgb = RGBColor(0x6B,0x72,0x80)


def add_title(slide, title, subtitle=None):
    box = slide.shapes.add_textbox(Inches(1.15), Inches(0.28), Inches(11.5), Inches(0.55))
    box.text_frame.text = title
    for r in box.text_frame.paragraphs[0].runs: r.font.size = Pt(26); r.font.bold = True; r.font.color.rgb = DARK_BLUE
    if subtitle:
        sub = slide.shapes.add_textbox(Inches(1.15), Inches(0.86), Inches(11.2), Inches(0.38))
        sub.text_frame.text = subtitle
        for r in sub.text_frame.paragraphs[0].runs: r.font.size = Pt(12.5); r.font.color.rgb = RGBColor(0x6B,0x72,0x80)


def add_bullets(slide, bullets, x=0.9, y=1.45, w=11.6, h=5.2, size=17):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame; tf.clear(); tf.word_wrap = True
    for i, text in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {text}"
        for r in p.runs: r.font.size = Pt(size); r.font.color.rgb = DARK
        p.space_after = Pt(9)
    return box


def add_table(slide, headers, rows, x=0.7, y=1.55, w=12.0, h=4.9):
    shape = slide.shapes.add_table(len(rows)+1, len(headers), Inches(x), Inches(y), Inches(w), Inches(h))
    table = shape.table
    for c, header in enumerate(headers):
        cell = table.cell(0,c); cell.text = header; cell.fill.solid(); cell.fill.fore_color.rgb = DARK_BLUE
        for r in cell.text_frame.paragraphs[0].runs: r.font.size = Pt(10.5); r.font.bold = True; r.font.color.rgb = WHITE
    for rr, row in enumerate(rows, start=1):
        for c, value in enumerate(row):
            cell = table.cell(rr,c); cell.text = str(value); cell.fill.solid(); cell.fill.fore_color.rgb = WHITE if rr%2 else LIGHT
            for r in cell.text_frame.paragraphs[0].runs: r.font.size = Pt(10); r.font.color.rgb = DARK
    return shape


def add_chart(slide, categories, values, title, x=1.0, y=1.55, w=11.2, h=4.9, line=False):
    data = CategoryChartData(); data.categories = categories; data.add_series("Actual", values)
    ctype = XL_CHART_TYPE.LINE_MARKERS if line else XL_CHART_TYPE.COLUMN_CLUSTERED
    shape = slide.shapes.add_chart(ctype, Inches(x), Inches(y), Inches(w), Inches(h), data)
    chart = shape.chart; chart.has_title = True; chart.chart_title.text_frame.text = title; chart.has_legend = False
    chart.plots[0].has_data_labels = True; chart.plots[0].data_labels.show_value = True
    try:
        chart.series[0].format.fill.solid(); chart.series[0].format.fill.fore_color.rgb = GREEN
        chart.series[0].format.line.color.rgb = GREEN
    except Exception:
        pass
    return shape


def new_slide(title, subtitle=None):
    slide = prs.slides.add_slide(blank)
    add_brand(slide, len(prs.slides))
    add_title(slide, title, subtitle)
    return slide

# 1
s = new_slide("Q3 Advisory Transformation Review", "Steering committee decision pack | September 2026")
hero = s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(1.2), Inches(1.8), Inches(10.9), Inches(3.8))
hero.fill.solid(); hero.fill.fore_color.rgb = RGBColor(0xE8,0xEE,0xFA); hero.line.color.rgb = BLUE
hero.text_frame.text = "Enterprise performance, delivery risks, and actions required for Q4"
p=hero.text_frame.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
for r in p.runs: r.font.size=Pt(27); r.font.bold=True; r.font.color.rgb=DARK_BLUE

# 2 numeric table
s = new_slide("Executive Performance Snapshot", "Q3 actuals versus plan")
add_table(s, ["Metric","Q2 Actual","Q3 Plan","Q3 Actual","Variance"], [
    ["Advisory revenue","$486.2M","$505.0M","$522.8M","+3.5%"],
    ["Contribution margin","22.9%","23.8%","24.6%","+0.8 pts"],
    ["Active engagements","1,214","1,250","1,286","+2.9%"],
    ["On-time milestones","84%","90%","92%","+2.0 pts"],
    ["Critical risks","14","10","8","-20.0%"],
])

# 3 dense prose
s = new_slide("Strategic Context", "What changed during the quarter")
add_bullets(s, [
    "Demand accelerated across cloud modernization, cyber resilience, and operating-model redesign, with the strongest pipeline conversion in regulated industries.",
    "Delivery capacity remained constrained in data engineering and program management, increasing reliance on specialist contractors and offshore teams.",
    "Client buying decisions shifted toward smaller phased commitments, creating more opportunities but increasing governance and coordination requirements.",
    "Leadership aligned on a Q4 focus of margin protection, milestone reliability, and selective investment in the highest-growth capabilities.",
], size=16)

# 4 findings
s = new_slide("Key Findings", "Three issues require steering committee attention")
add_bullets(s, [
    "Revenue exceeded plan by $17.8M, led by cloud transformation and cyber resilience programs.",
    "Contribution margin improved 80 basis points, but contractor spend remains 12% above target.",
    "Eight critical delivery risks remain open, concentrated in data migration, client dependencies, and resource availability.",
    "Milestone performance reached 92%, the highest level in six quarters.",
    "Cross-sell conversion improved from 18% to 23% after account planning changes.",
], size=17)

# 5 table-target content
s = new_slide("Operating Performance", "Delivery and commercial observations")
add_bullets(s, [
    "Cloud modernization delivered $184M revenue and 27% margin.",
    "Cyber resilience delivered $126M revenue and 25% margin.",
    "Data and AI delivered $98M revenue and 21% margin.",
    "Operating model delivered $74M revenue and 22% margin.",
    "Managed services delivered $41M revenue and 19% margin.",
], size=17)

# 6 table to chart
s = new_slide("Quarterly Revenue Trend", "Revenue progression by quarter")
add_table(s, ["Quarter","Revenue ($M)"], [["Q4 2025","448.0"],["Q1 2026","462.5"],["Q2 2026","486.2"],["Q3 2026","522.8"]], x=2.2,y=1.6,w=8.9,h=4.5)

# 7 risks
s = new_slide("Key Risks", "Current exposure and mitigation status")
add_table(s, ["Risk","Impact","Likelihood","Mitigation","Status"], [
    ["Data migration delay","High","Medium","Daily cutover governance","Amber"],
    ["Specialist capacity","High","High","Accelerated hiring and vendor bench","Red"],
    ["Client dependency","Medium","High","Executive escalation path","Amber"],
    ["Margin leakage","Medium","Medium","Weekly commercial review","Green"],
    ["Regulatory change","High","Low","Policy monitoring and scenario plan","Green"],
])

# 8 recommendations
s = new_slide("Recommended Actions", "Actions proposed for steering committee approval")
add_bullets(s, [
    "Approve targeted specialist hiring for data engineering and program management.",
    "Require recovery plans for every milestone with red or amber delivery status.",
    "Expand contractor rate-card governance to protect Q4 margin.",
    "Prioritize cloud and cyber pipeline conversion while limiting low-margin custom work.",
], size=17)

slide_specs = [
("Revenue by Capability","Q3 revenue mix and growth", "chart"),
("Margin Bridge","Primary drivers of margin improvement", "bullets"),
("Pipeline Health","Qualified pipeline and conversion", "table"),
("Client Portfolio","Concentration and account momentum", "bullets"),
("Delivery Reliability","Milestone performance by month", "chart_line"),
("Resource Capacity","Demand versus available specialist capacity", "table"),
("Contractor Economics","Spend, rates, and utilization", "bullets"),
("Cloud Modernization","Performance and outlook", "chart"),
("Cyber Resilience","Performance and outlook", "chart"),
("Data and AI","Performance and outlook", "chart"),
("Operating Model","Performance and outlook", "bullets"),
("Managed Services","Performance and outlook", "bullets"),
("Cross-Sell Performance","Conversion by priority account", "table"),
("Regional Performance","Revenue and margin by region", "chart"),
("Client Satisfaction","NPS and executive feedback", "chart_line"),
("Transformation Milestones","Status of major initiatives", "table"),
("Q4 Outlook","Base, upside, and downside scenarios", "table"),
("Decision Requests","Approvals required today", "bullets"),
("90-Day Action Plan","Actions, owners, and timing", "table"),
("Governance Cadence","Meeting and escalation structure", "bullets"),
("Appendix: Metric Definitions","Definitions and calculation notes", "bullets"),
("Closing and Next Steps","Confirmed actions and follow-up", "bullets"),
]

for idx, (title, subtitle, kind) in enumerate(slide_specs, start=9):
    s = new_slide(title, subtitle)
    if kind == "chart":
        add_chart(s, ["Q4","Q1","Q2","Q3"], [72+idx, 78+idx, 84+idx, 92+idx], title)
    elif kind == "chart_line":
        add_chart(s, ["Apr","May","Jun","Jul","Aug","Sep"], [78,82,84,87,90,92], title, line=True)
    elif kind == "table":
        add_table(s, ["Area","Current","Target","Status"], [
            ["Priority 1", f"{82+idx}%", "90%", "Amber"],
            ["Priority 2", f"{75+idx/2:.0f}%", "85%", "Green"],
            ["Priority 3", f"{68+idx/3:.0f}%", "80%", "Red"],
            ["Priority 4", f"{88+idx/4:.0f}%", "90%", "Green"],
        ])
    else:
        add_bullets(s, [
            f"Performance improved during Q3, with the strongest contribution from priority initiative {idx-8}.",
            "The current trajectory supports the base plan, subject to timely resolution of identified dependencies.",
            "Leadership action is required to confirm ownership, timing, and resource allocation.",
            "The team will report progress through the weekly transformation governance cadence.",
        ], size=16)

assert len(prs.slides) == 30, len(prs.slides)
path = OUT / "kpmg_advisory_30_slide_original.pptx"
prs.save(path)

wb = Workbook(); ws = wb.active; ws.title = "Q4 Update"
ws.append(["Metric","Q3 Actual","Q4 Forecast"])
for row in [
    ["Advisory revenue",522.8,548.4],
    ["Contribution margin",24.6,25.1],
    ["Active engagements",1286,1324],
    ["On-time milestones",92,94],
    ["Critical risks",8,5],
]: ws.append(row)
wb.save(OUT / "kpmg_advisory_30_slide_data.xlsx")
print(path)
