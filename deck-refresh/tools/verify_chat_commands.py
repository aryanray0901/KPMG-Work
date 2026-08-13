"""Regression tests for deterministic chat slide-management commands."""
from __future__ import annotations

import os
import shutil
import sys
import tempfile
from pathlib import Path

ROOT = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(ROOT))

from pptx import Presentation
from pptx_editor import apply_operations, plan_edit


def check(condition: bool, message: str) -> None:
    if not condition:
        raise AssertionError(message)
    print(f"PASS: {message}")


def slide_ids(path: str | Path) -> list[int]:
    prs = Presentation(str(path))
    return [int(item.id) for item in prs.slides._sldIdLst]


def apply_chat(source: str, destination: str, message: str, selected_slide: int = 1):
    plan = plan_edit(source, message, selected_slide, None, [], [], [])
    operations = [op for op in plan.get("operations", []) if op.get("op") != "noop"]
    if not operations:
        shutil.copy2(source, destination)
        return plan, None
    result = apply_operations(source, destination, operations)
    return plan, result


def main() -> None:
    previous_key = os.environ.pop("OPENAI_API_KEY", None)
    try:
        sample = ROOT / "sample_files" / "kpmg_advisory_q3_original.pptx"
        check(len(Presentation(sample).slides) == 9, "sample starts with nine slides")

        with tempfile.TemporaryDirectory(prefix="chat_slide_commands_") as temp_dir:
            # Exact wording reported from the browser must stay deterministic.
            exact = plan_edit(str(sample), "Delete the last slide in the deck.", 1, None, [], [], [])
            check(exact["operations"] == [{"op": "delete_slide", "slide": 9}], "browser wording deletes the last slide without AI")

            variants = [
                "Remove the final slide from this presentation",
                "Please delete the last page in my PowerPoint.",
                "Can you remove slide 9 from the deck?",
                "I want you to delete the slide at the end of the file",
                "get rid of the very last slide",
                "Take out the final one in the presentation",
                "erase slide 9 from the pptx",
            ]
            for wording in variants:
                variant_plan = plan_edit(str(sample), wording, 1, None, [], [], [])
                check(variant_plan["operations"] == [{"op": "delete_slide", "slide": 9}], f"understands delete wording: {wording}")

            compound = plan_edit(
                str(sample),
                "duplicate slide 4 and then move that new slide to the last slide in the deck",
                1, None, [], [], [],
            )
            check(compound["operations"] == [{"op": "duplicate_slide", "slide": 4, "position": 10}], "compound duplicate-and-move runs locally")

            step1 = os.path.join(temp_dir, "step1.pptx")
            plan1, result1 = apply_chat(str(sample), step1, "delete last slide")
            check(plan1["operations"] == [{"op": "delete_slide", "slide": 9}], "delete last slide bypasses AI and targets slide 9")
            check(result1 and not result1["skipped"] and len(Presentation(step1).slides) == 8, "delete last slide executes")

            step2 = os.path.join(temp_dir, "step2.pptx")
            plan2, result2 = apply_chat(step1, step2, "delete slide 9")
            check(not plan2["operations"] and plan2["message"].startswith("error cant do that") and "Try:" in plan2["message"], "invalid slide returns the safe error message")
            check(result2 is None and slide_ids(step2) == slide_ids(step1), "invalid delete leaves the deck unchanged")

            before_move = slide_ids(step2)
            moved_identity = before_move[7]
            step3 = os.path.join(temp_dir, "step3.pptx")
            plan3, result3 = apply_chat(step2, step3, "move slide 8 to slide 4")
            after_move = slide_ids(step3)
            check(plan3["operations"] == [{"op": "move_slide", "from_slide": 8, "to_slide": 4}], "move slide X to slide Y compiles locally")
            check(result3 and not result3["skipped"] and after_move[3] == moved_identity, "move command places the requested slide in position 4")

            plan4 = plan_edit(step3, "delete slides 2 and 4", 1, None, [], [], [])
            check([op["slide"] for op in plan4["operations"]] == [4, 2], "multiple deletes run from highest position to lowest")

            plan5 = plan_edit(step3, "move this slide right", 3, None, [], [], [])
            check(plan5["operations"] == [{"op": "move_slide", "from_slide": 3, "to_slide": 4}], "current-slide directional movement works")

            plan6 = plan_edit(step3, "duplicate last slide", 1, None, [], [], [])
            check(plan6["operations"] == [{"op": "duplicate_slide", "slide": 8, "position": 9}], "duplicate last slide works without AI")

            plan7 = plan_edit(step3, "reorder slides to 1, 2, 3, 4, 5, 6, 8, 7", 1, None, [], [], [])
            check(plan7["operations"] and plan7["operations"][0]["op"] == "reorder_slides", "complete explicit slide order works")

        print("\nChat slide-command checks passed.")
    finally:
        if previous_key is not None:
            os.environ["OPENAI_API_KEY"] = previous_key


if __name__ == "__main__":
    main()
