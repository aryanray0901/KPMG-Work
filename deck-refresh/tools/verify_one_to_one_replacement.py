"""Regression test for the home-page 1:1 replacement workflow."""

from __future__ import annotations

import io
import os
import re
import shutil
import sys
import tempfile
from pathlib import Path

from openpyxl import Workbook
from pptx import Presentation

ROOT = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(ROOT))

import app as webapp
from replacement_engine import (
    compare_deck_replacements,
    inspect_deck_structure,
    replace_deck_1to1,
    validate_replacement_structure,
)


CASES = [
    {
        "name": "screening",
        "deck": ROOT / "sample_files" / "Company Screening Analysis.pptx",
        "data": ROOT / "sample_files" / "new data screening.xlsx",
        "profile": "Company screening",
        "rows": 1810,
        "charts": 4,
        "tables": 2,
        "checks": {
            "slide2": "86.7%",
            "slide3": "Europe is 40% of the universe",
            "slide8": "Atlas Test Company",
        },
    },
    {
        "name": "goldman",
        "deck": ROOT / "sample_files" / "Goldman Sachs Q3 Analysis.pptx",
        "data": ROOT / "sample_files" / "new goldman sachs Q4.xlsx",
        "profile": "Goldman Sachs 424B2",
        "rows": 32,
        "charts": 3,
        "tables": 2,
        "checks": {
            "slide1": "Q4 client presentation",
            "slide10": "Q4 1:1 replacement ready",
            "slide2": "53.1%",
            "slide3": "$53.34",
            "slide6": "$128.00",
        },
    },
]


def slide_text(prs, number):
    slide = prs.slides[number - 1]
    values = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            values.append(shape.text)
        if shape.has_table:
            values.extend(cell.text for row in shape.table.rows for cell in row.cells)
    return "\n".join(values)


def all_deck_text(prs):
    values = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                values.append(shape.text)
            if getattr(shape, "has_table", False):
                values.extend(cell.text for row in shape.table.rows for cell in row.cells)
            if getattr(shape, "has_chart", False):
                values.extend(str(series.name or "") for series in shape.chart.series)
                values.extend(str(value) for value in shape.chart.plots[0].categories)
                if shape.chart.has_title:
                    values.append(shape.chart.chart_title.text_frame.text)
        if slide.has_notes_slide:
            values.extend(
                shape.text for shape in slide.notes_slide.shapes if getattr(shape, "has_text_frame", False)
            )
    return "\n".join(values)


def direct_engine_test(case, temp_dir):
    output = Path(temp_dir) / f"{case['name']}_updated.pptx"
    report = replace_deck_1to1(case["deck"], case["data"], output, case["data"].name)
    assert report["profile"] == case["profile"]
    assert report["source_rows"] == case["rows"]
    assert report["slide_count"] == 10
    assert report["charts"] == case["charts"]
    assert report["tables"] == case["tables"]
    assert report["geometry_preserved"] is True
    assert report["theme_preserved"] is True
    assert report["total_updates"] > 50
    prs = Presentation(output)
    for key, expected in case["checks"].items():
        number = int(key.replace("slide", ""))
        assert expected in slide_text(prs, number), f"{case['name']} missing {expected} on slide {number}"
    if case["name"] == "goldman":
        assert not re.search(r"\bQ3\b", all_deck_text(prs), re.IGNORECASE)
        assert re.search(r"\bQ4\b", all_deck_text(prs), re.IGNORECASE)
    assert len(prs.slides) == len(Presentation(case["deck"]).slides)
    mappings = compare_deck_replacements(case["deck"], output)
    assert mappings["changed_fields"] > 25
    assert mappings["chart_changes"] > 0
    assert mappings["table_changes"] > 0
    structure = inspect_deck_structure(case["deck"], mappings["entries"])
    assert structure["slide_count"] == 10
    assert structure["charts"] == case["charts"]
    assert structure["tables"] == case["tables"]
    assert structure["shapes"] > 0
    assert any(slide["objects"] for slide in structure["slides"])
    validation = validate_replacement_structure(case["deck"], output)
    assert validation["passed"] is True
    assert len(validation["checks"]) == 6
    assert all(check["passed"] for check in validation["checks"])
    tampered = Path(temp_dir) / f"{case['name']}_tampered.pptx"
    tampered_prs = Presentation(output)
    tampered_prs.slides[0].shapes[0].left += 914400
    tampered_prs.save(tampered)
    rejected = validate_replacement_structure(case["deck"], tampered)
    assert rejected["passed"] is False
    assert any(check["name"] == "Object positions and sizes" and not check["passed"] for check in rejected["checks"])
    return report


def home_page_test(case):
    original_renderer = webapp.render_pptx_to_images
    webapp.render_pptx_to_images = lambda *_args, **_kwargs: (None, None)
    try:
        client = webapp.app.test_client()
        home = client.get("/")
        assert home.status_code == 200
        assert b'action="/replace1to1"' in home.data
        assert b'Analyze and match data' in home.data
        with open(case["deck"], "rb") as deck_handle, open(case["data"], "rb") as data_handle:
            response = client.post(
                "/replace1to1",
                data={
                    "primary_file": (io.BytesIO(deck_handle.read()), case["deck"].name),
                    "data_file": (io.BytesIO(data_handle.read()), case["data"].name),
                },
                content_type="multipart/form-data",
                follow_redirects=False,
            )
        assert response.status_code == 302
        location = response.headers["Location"]
        assert "/replace1to1/review/" in location, location
        review = client.get(location)
        assert review.status_code == 200
        assert b"Review the presentation structure" in review.data
        assert b"Structure validation passed" in review.data
        assert b"Presentation structure" in review.data
        assert b"Proposed changes" in review.data
        assert b"Native and editable" in review.data
        assert b"Approve all" in review.data
        assert b"Apply selected changes" in review.data
        assert b"mapping-checkbox" in review.data
        assert b"Select all" in review.data
        assert b"Old value" in review.data and b"Replacement" in review.data
        sid = location.rstrip("/").split("/")[-1]
        blocked = client.post(f"/replace1to1/apply/{sid}", follow_redirects=False)
        assert blocked.status_code == 302
        assert f"/replace1to1/review/{sid}" in blocked.headers["Location"]
        assert not (Path(webapp.SESSIONS_DIR) / sid / "updated.pptx").exists()
        apply_response = client.post(
            f"/replace1to1/apply/{sid}",
            data={"approval": "approve_all"},
            follow_redirects=False,
        )
        assert apply_response.status_code == 302
        result_location = apply_response.headers["Location"]
        assert "/result/" in result_location, result_location
        result = client.get(result_location)
        assert result.status_code == 200
        assert b"1:1 replacement complete" in result.data
        shutil.rmtree(Path(webapp.SESSIONS_DIR) / sid, ignore_errors=True)
    finally:
        webapp.render_pptx_to_images = original_renderer


def markerless_profile_test(case, temp_dir):
    marker = "screening sheet unchanged" if case["name"] == "screening" else "same sheets and row anchors"
    source = Presentation(case["deck"])
    replacements = 0
    for slide in source.slides:
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    if marker in run.text.casefold():
                        run.text = run.text.lower().replace(marker, "workbook structure verified")
                        replacements += 1
    assert replacements > 0
    markerless_path = Path(temp_dir) / f"{case['name']}_markerless.pptx"
    output = Path(temp_dir) / f"{case['name']}_markerless_updated.pptx"
    source.save(markerless_path)
    report = replace_deck_1to1(markerless_path, case["data"], output, case["data"].name)
    assert report["profile"] == case["profile"]
    assert report["total_updates"] > 50


def generic_fallback_test():
    deck = ROOT / "sample_files" / "kpmg_advisory_q3_original.pptx"
    data = ROOT / "sample_files" / "kpmg_advisory_q4_data.xlsx"
    original_renderer = webapp.render_pptx_to_images
    webapp.render_pptx_to_images = lambda *_args, **_kwargs: (None, None)
    try:
        client = webapp.app.test_client()
        with open(deck, "rb") as deck_handle, open(data, "rb") as data_handle:
            response = client.post(
                "/replace1to1",
                data={
                    "primary_file": (io.BytesIO(deck_handle.read()), deck.name),
                    "data_file": (io.BytesIO(data_handle.read()), data.name),
                },
                content_type="multipart/form-data",
                follow_redirects=False,
            )
        assert response.status_code == 302
        location = response.headers["Location"]
        assert "/replace1to1/review/" in location
        review = client.get(location)
        assert review.status_code == 200
        assert b"Structure validation passed" in review.data
        assert b"Approve all" in review.data
        sid = location.rstrip("/").split("/")[-1]
        meta = webapp._load_session_meta(Path(webapp.SESSIONS_DIR) / sid)
        assert meta["replacement_report"]["profile"] == "Generic label-matched deck"
        assert meta["replacement_report"]["total_updates"] == 120
        shutil.rmtree(Path(webapp.SESSIONS_DIR) / sid, ignore_errors=True)
    finally:
        webapp.render_pptx_to_images = original_renderer


def arbitrary_relative_matching_test(temp_dir):
    deck = ROOT / "sample_files" / "kpmg_advisory_q3_original.pptx"
    workbook_path = Path(temp_dir) / "unrelated_layout.xlsx"
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "Operating Inputs"
    sheet.append(["Operational Data Extract"])
    sheet.append(["Indicator", "Current", "Next"])
    for index in range(1, 31):
        sheet.append([f"Indicator {index}", index * 11, index * 17])
    second = workbook.create_sheet("Regional Measures")
    second.append(["Territory", "Measure A", "Measure B"])
    for index, region in enumerate(("North", "South", "East", "West"), start=1):
        second.append([region, index * 125, index / 10])
    workbook.save(workbook_path)

    original_renderer = webapp.render_pptx_to_images
    webapp.render_pptx_to_images = lambda *_args, **_kwargs: (None, None)
    try:
        client = webapp.app.test_client()
        with open(deck, "rb") as deck_handle, open(workbook_path, "rb") as data_handle:
            response = client.post(
                "/replace1to1",
                data={
                    "primary_file": (io.BytesIO(deck_handle.read()), "arbitrary_presentation.pptx"),
                    "data_file": (io.BytesIO(data_handle.read()), "arbitrary_workbook.xlsx"),
                },
                content_type="multipart/form-data",
                follow_redirects=False,
            )
        assert response.status_code == 302
        location = response.headers["Location"]
        assert "/replace1to1/review/" in location
        review = client.get(location)
        assert review.status_code == 200
        assert b"Relative matching active" in review.data
        assert b"Matched from" in review.data
        assert b"Relative position" in review.data or b"Similar variable" in review.data
        sid = location.rstrip("/").split("/")[-1]
        meta = webapp._load_session_meta(Path(webapp.SESSIONS_DIR) / sid)
        report = meta["replacement_report"]
        assert report["profile"] == "Generic label-matched deck"
        assert report["total_updates"] > 0
        assert report["relative_matches"]
        apply_response = client.post(
            f"/replace1to1/apply/{sid}",
            data={"approval": "approve_all"},
            follow_redirects=False,
        )
        assert apply_response.status_code == 302
        assert "/result/" in apply_response.headers["Location"]
        updated_path = Path(webapp.SESSIONS_DIR) / sid / "updated.pptx"
        assert updated_path.exists()
        assert len(Presentation(updated_path).slides) == len(Presentation(deck).slides)
        shutil.rmtree(Path(webapp.SESSIONS_DIR) / sid, ignore_errors=True)
    finally:
        webapp.render_pptx_to_images = original_renderer


def no_numeric_data_still_validates_test(temp_dir):
    deck = ROOT / "sample_files" / "kpmg_advisory_q3_original.pptx"
    workbook_path = Path(temp_dir) / "text_only.xlsx"
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "Notes"
    sheet.append(["Topic", "Comment"])
    sheet.append(["Status", "Pending management review"])
    workbook.save(workbook_path)

    original_renderer = webapp.render_pptx_to_images
    webapp.render_pptx_to_images = lambda *_args, **_kwargs: (None, None)
    try:
        client = webapp.app.test_client()
        with open(deck, "rb") as deck_handle, open(workbook_path, "rb") as data_handle:
            response = client.post(
                "/replace1to1",
                data={
                    "primary_file": (io.BytesIO(deck_handle.read()), "presentation.pptx"),
                    "data_file": (io.BytesIO(data_handle.read()), "text_only.xlsx"),
                },
                content_type="multipart/form-data",
                follow_redirects=False,
            )
        assert response.status_code == 302
        location = response.headers["Location"]
        assert "/replace1to1/review/" in location
        review = client.get(location)
        assert review.status_code == 200
        assert b"Structure validation passed" in review.data
        assert b"No comparable numeric fields found" in review.data
        sid = location.rstrip("/").split("/")[-1]
        shutil.rmtree(Path(webapp.SESSIONS_DIR) / sid, ignore_errors=True)
    finally:
        webapp.render_pptx_to_images = original_renderer


def main():
    for case in CASES:
        for path in (case["deck"], case["data"]):
            assert path.exists(), f"Missing sample file: {path}"
    with tempfile.TemporaryDirectory(prefix="deck_refresh_1to1_") as temp_dir:
        reports = []
        for case in CASES:
            reports.append(direct_engine_test(case, temp_dir))
            markerless_profile_test(case, temp_dir)
            home_page_test(case)
        generic_fallback_test()
        arbitrary_relative_matching_test(temp_dir)
        no_numeric_data_still_validates_test(temp_dir)
    print("1:1 replacement verification passed")
    print("Markerless, generic, arbitrary relative, and no-numeric validation passed")
    for report in reports:
        print(
            f"{report['profile']}: {report['slide_count']} slides, {report['charts']} charts, "
            f"{report['tables']} tables, {report['total_updates']} mapped fields"
        )


if __name__ == "__main__":
    main()
