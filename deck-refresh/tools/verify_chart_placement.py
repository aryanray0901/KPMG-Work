#!/usr/bin/env python3
"""Verify all chart buttons create blank native charts in either placement."""

from __future__ import annotations

import os
import shutil
import sys
import tempfile
import time
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE

ROOT = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(ROOT))

import app as deck_app
from builder_ops import create_blank_deck
from pptx_editor import _insertion_obstacle_box, _shape_is_branding


SOURCE = ROOT / "sample_files" / "kpmg_advisory_30_slide_original.pptx"
CHART_LAYOUTS = {
    "column": "bar_chart",
    "bar": "bar_chart",
    "line": "line_chart",
    "pie": "pie_chart",
    "area": "area_chart",
    "waterfall": "waterfall_chart",
    "scatter": "scatter_plot",
}


def fake_render(sess_dir: str, version: int) -> dict:
    count = len(Presentation(deck_app._editor_version_path(sess_dir, version)).slides)
    return {
        "slide_count": count,
        "rendering_ok": False,
        "render_engine": None,
        "render_help": "Test renderer disabled.",
        "preview_revision": time.time_ns(),
    }


def box(prs: Presentation, shape) -> tuple[float, float, float, float]:
    return (
        shape.left / prs.slide_width,
        shape.top / prs.slide_height,
        shape.width / prs.slide_width,
        shape.height / prs.slide_height,
    )


def intersection(first, second) -> float:
    width = max(0, min(first[0] + first[2], second[0] + second[2]) - max(first[0], second[0]))
    height = max(0, min(first[1] + first[3], second[1] + second[3]) - max(first[1], second[1]))
    return width * height


def assert_blank_chart(chart, chart_type: str) -> None:
    assert chart.series, chart_type
    assert all(all(abs(float(value or 0)) < 1e-9 for value in series.values) for series in chart.series), chart_type
    assert not chart.has_title, chart_type


def assert_current_slide_geometry(prs: Presentation, shapes_before: dict[int, tuple[float, float, float, float]],
                                  slide_number: int = 1, require_right_side: bool = False) -> None:
    slide = prs.slides[slide_number - 1]
    inserted = [shape for shape in slide.shapes if getattr(shape, "has_chart", False)][-1]
    inserted_box = box(prs, inserted)
    assert inserted_box[0] >= 0 and inserted_box[1] >= 0
    assert inserted_box[0] + inserted_box[2] <= 1.001
    assert inserted_box[1] + inserted_box[3] <= 1.001
    assert 0.18 <= inserted_box[2] <= 0.581
    assert 0.13 <= inserted_box[3] <= 0.501
    if require_right_side:
        assert inserted_box[0] >= 0.60, inserted_box
    for shape in slide.shapes:
        if int(shape.shape_id) == int(inserted.shape_id):
            continue
        other = _insertion_obstacle_box(prs, shape)
        if other is None:
            continue
        area = other[2] * other[3]
        if area < 0.0005 or area > 0.78:
            continue
        assert intersection(inserted_box, other) <= 0.00001, (inserted_box, shape.name, other)

    shapes_after = {int(shape.shape_id): box(prs, shape) for shape in slide.shapes}
    for shape_id, original_box in shapes_before.items():
        assert shape_id in shapes_after, shape_id
        assert all(abs(first - second) <= 0.0002 for first, second in zip(original_box, shapes_after[shape_id])), (
            shape_id, original_box, shapes_after[shape_id]
        )

    branding = [shape for shape in slide.shapes if _shape_is_branding(shape, prs)]
    branding_ids = {int(shape.shape_id) for shape in branding}
    top_text = [
        shape for shape in slide.shapes
        if getattr(shape, "has_text_frame", False)
        and str(shape.text_frame.text or "").strip()
        and int(shape.shape_id) not in branding_ids
        and box(prs, shape)[1] < 0.30
    ]
    for brand in branding:
        for text_shape in top_text:
            assert intersection(box(prs, brand), box(prs, text_shape)) <= 0.0002, (
                brand.name, text_shape.name, box(prs, brand), box(prs, text_shape)
            )


def main() -> None:
    original_renderer = deck_app._render_editor_version
    deck_app._render_editor_version = fake_render
    deck_app.app.testing = True
    client = deck_app.app.test_client()
    session_ids: list[str] = []
    try:
        for chart_type, layout in CHART_LAYOUTS.items():
            sid = deck_app._create_editor_session(str(SOURCE), SOURCE.name)
            session_ids.append(sid)
            source_deck = Presentation(str(SOURCE))
            shapes_before = {
                int(shape.shape_id): box(source_deck, shape)
                for shape in source_deck.slides[4].shapes
            }
            current_response = client.post(f"/editor/build/{sid}", data={
                "layout": layout,
                "chart_type": chart_type,
                "data_mode": "blank",
                "placement": "current",
                "selected_slide": "5",
                "smart": "false",
            })
            current_payload = current_response.get_json(silent=True) or {}
            assert current_response.status_code == 200 and current_payload.get("ok"), (chart_type, current_payload)
            assert current_payload["state"]["slide_count"] == 30 and current_payload["state"]["selected_slide"] == 5
            current_path = deck_app._editor_version_path(
                deck_app._session_dir(sid), current_payload["state"]["version"]
            )
            current_deck = Presentation(current_path)
            current_chart = [shape.chart for shape in current_deck.slides[4].shapes if shape.has_chart][-1]
            assert_blank_chart(current_chart, chart_type)
            assert_current_slide_geometry(current_deck, shapes_before, slide_number=5, require_right_side=True)

            new_response = client.post(f"/editor/build/{sid}", data={
                "layout": layout,
                "chart_type": chart_type,
                "data_mode": "blank",
                "placement": "new",
                "selected_slide": "5",
                "smart": "false",
            })
            new_payload = new_response.get_json(silent=True) or {}
            assert new_response.status_code == 200 and new_payload.get("ok"), (chart_type, new_payload)
            assert new_payload["state"]["slide_count"] == 31 and new_payload["state"]["selected_slide"] == 6
            new_path = deck_app._editor_version_path(deck_app._session_dir(sid), new_payload["state"]["version"])
            new_deck = Presentation(new_path)
            new_charts = [shape.chart for shape in new_deck.slides[5].shapes if shape.has_chart]
            assert len(new_charts) == 1, chart_type
            assert_blank_chart(new_charts[0], chart_type)

        handle, dense_path = tempfile.mkstemp(suffix=".pptx")
        os.close(handle)
        create_blank_deck(dense_path)
        dense_deck = Presentation(dense_path)
        dense_slide = dense_deck.slides[0]
        dense_slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            int(dense_deck.slide_width * 0.025),
            int(dense_deck.slide_height * 0.025),
            int(dense_deck.slide_width * 0.95),
            int(dense_deck.slide_height * 0.90),
        )
        dense_deck.save(dense_path)
        dense_sid = deck_app._create_editor_session(dense_path, "dense-slide.pptx")
        session_ids.append(dense_sid)
        dense_response = client.post(f"/editor/build/{dense_sid}", data={
            "layout": "line_chart",
            "chart_type": "line",
            "data_mode": "blank",
            "placement": "current",
            "selected_slide": "1",
            "smart": "false",
        })
        dense_payload = dense_response.get_json(silent=True) or {}
        assert dense_response.status_code == 200 and dense_payload.get("ok"), dense_payload
        assert dense_payload["state"]["slide_count"] == 1 and dense_payload["state"]["version"] == 1
        dense_output = deck_app._editor_version_path(
            deck_app._session_dir(dense_sid), dense_payload["state"]["version"]
        )
        dense_result = Presentation(dense_output)
        dense_charts = [shape for shape in dense_result.slides[0].shapes if getattr(shape, "has_chart", False)]
        assert len(dense_charts) == 1, dense_payload
        dense_box = box(dense_result, dense_charts[0])
        assert dense_box[2] > 0 and dense_box[3] > 0, dense_box
        assert dense_box[0] + dense_box[2] <= 1.001 and dense_box[1] + dense_box[3] <= 1.001, dense_box
        os.remove(dense_path)
    finally:
        deck_app._render_editor_version = original_renderer
        for sid in session_ids:
            shutil.rmtree(os.path.join(deck_app.SESSIONS_DIR, sid), ignore_errors=True)

    print("PASS: all 7 chart buttons place on the current slide, shrink into corners when needed, preserve every object, and avoid charts and metrics.")


if __name__ == "__main__":
    main()
