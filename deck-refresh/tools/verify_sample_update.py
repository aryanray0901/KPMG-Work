from pathlib import Path
import sys

from pptx import Presentation

ROOT = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(ROOT))

from app import (  # noqa: E402
    apply_edits_pptx,
    collect_all_pptx_text,
    compute_period_replacements,
    extract_from_spreadsheet,
    extract_number_value,
    extract_targets_pptx,
    match_targets_to_source,
)

SAMPLES = ROOT / "sample_files"
ORIGINAL = SAMPLES / "kpmg_advisory_q3_original.pptx"
DATA = SAMPLES / "kpmg_advisory_q4_data.xlsx"
UPDATED = SAMPLES / "kpmg_advisory_q4_expected.pptx"
REPORT = SAMPLES / "sample_verification.txt"


def chart_period_text(prs):
    values = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_chart:
                continue
            chart = shape.chart
            try:
                values.extend(str(c) for c in chart.plots[0].categories)
                values.extend((series.name or "") for series in chart.series)
            except Exception:
                pass
    return values


def structure_signature(path):
    prs = Presentation(path)
    slides = []
    for slide in prs.slides:
        shapes = []
        for shape in slide.shapes:
            kind = "chart" if shape.has_chart else "table" if shape.has_table else "text" if shape.has_text_frame else "shape"
            extra = None
            if shape.has_chart:
                extra = (str(shape.chart.chart_type), len(shape.chart.series), len(list(shape.chart.plots[0].categories)))
            elif shape.has_table:
                extra = (len(shape.table.rows), len(shape.table.columns))
            shapes.append((kind, shape.left, shape.top, shape.width, shape.height, extra))
        slides.append(shapes)
    return slides


def main():
    source_pairs = extract_from_spreadsheet(str(DATA))
    prs, targets = extract_targets_pptx(str(ORIGINAL))
    matches = match_targets_to_source(targets, source_pairs)
    unmatched = [m for m in matches if not m.get("matched_label")]
    if unmatched:
        raise AssertionError(f"Unmatched sample targets: {[m['id'] for m in unmatched]}")

    period_replacements = compute_period_replacements(collect_all_pptx_text(prs), source_pairs, "")
    if period_replacements != {"Q3": "Q4"}:
        raise AssertionError(f"Unexpected period replacements: {period_replacements}")

    confirmed = []
    for match in matches:
        entry = dict(match)
        entry["new_text"] = match["new_text_preview"]
        confirmed.append(entry)

    apply_edits_pptx(str(ORIGINAL), str(UPDATED), confirmed, period_replacements)

    _, updated_targets = extract_targets_pptx(str(UPDATED))
    updated_by_id = {item["id"]: item for item in updated_targets}
    failures = []
    for match in confirmed:
        changed = updated_by_id.get(match["id"])
        if not changed:
            failures.append(f"Missing target after update: {match['id']}")
            continue
        if match["kind"] == "chart_point":
            actual = extract_number_value(changed["original"])
            expected = float(match["new_value"])
            if actual is None or abs(actual - expected) > 1e-8:
                failures.append(f"Chart mismatch {match['id']}: {actual} != {expected}")
        elif changed["original"] != match["new_text"]:
            failures.append(f"Text mismatch {match['id']}: {changed['original']} != {match['new_text']}")

    updated_prs = Presentation(str(UPDATED))
    period_text = collect_all_pptx_text(updated_prs) + chart_period_text(updated_prs)
    q3_left = [text for text in period_text if "Q3" in str(text).upper()]
    if q3_left:
        failures.append(f"Q3 text remained: {q3_left[:5]}")

    if structure_signature(ORIGINAL) != structure_signature(UPDATED):
        failures.append("Slide, shape, table, or chart structure changed during update")

    if failures:
        raise AssertionError("\n".join(failures))

    lines = [
        "KPMG Deck Refresh fictional sample verification: PASS",
        f"Slides checked: {len(prs.slides)}",
        f"Numeric targets checked: {len(targets)}",
        f"Matched targets: {len(matches)}",
        "Unmatched targets: 0",
        "Period update checked: Q3 to Q4",
        "Slide and shape structure preserved: yes",
        "Table values checked: yes",
        "Chart values checked: yes",
        f"Expected updated file: {UPDATED.name}",
    ]
    REPORT.write_text("\n".join(lines) + "\n", encoding="utf-8")
    print("\n".join(lines))


if __name__ == "__main__":
    main()
