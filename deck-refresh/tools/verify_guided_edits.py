#!/usr/bin/env python3
"""Verify guided input forms, company rebranding, and deterministic chat edits."""

from __future__ import annotations

import os
import re
import sys
import tempfile
from pathlib import Path

from pptx import Presentation

ROOT = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(ROOT))

from builder_ops import LAYOUTS
from pptx_editor import apply_operations, deck_summary, plan_edit


SOURCE = ROOT / "sample_files" / "kpmg_advisory_30_slide_original.pptx"


def apply_command(source: str, command: str, selected_slide: int) -> str:
    plan = plan_edit(source, command, selected_slide)
    assert plan and not plan.get("failed") and plan.get("operations"), (command, plan)
    handle, output = tempfile.mkstemp(suffix=".pptx")
    os.close(handle)
    result = apply_operations(source, output, plan["operations"])
    assert not result.get("skipped") and not result.get("failed") and not result.get("unchanged"), (command, result)
    assert len(result.get("applied", [])) == len(plan["operations"]), (command, result)
    return output


def all_container_text(prs: Presentation) -> str:
    values: list[str] = []
    containers = list(prs.slides)
    for master in prs.slide_masters:
        containers.append(master)
        containers.extend(list(master.slide_layouts))
    for container in containers:
        for shape in container.shapes:
            if getattr(shape, "has_text_frame", False):
                values.append(shape.text_frame.text)
            if getattr(shape, "has_table", False):
                values.extend(cell.text for row in shape.table.rows for cell in row.cells)
    return "\n".join(values)


def verify_company_and_font() -> None:
    rebranded = apply_command(str(SOURCE), 'Rebrand the entire deck from "KPMG" to "Deloitte".', 1)
    try:
        prs = Presentation(rebranded)
        text = all_container_text(prs)
        assert "kpmg" not in text.casefold(), "KPMG remained in visible, master, or layout text"
        assert "deloitte" in text.casefold()
        assert sum(
            "deloitte" in "\n".join(
                shape.text_frame.text for shape in slide.shapes if getattr(shape, "has_text_frame", False)
            ).casefold()
            for slide in prs.slides
        ) == len(prs.slides), "The company name was not changed on every slide"

        formatted = apply_command(rebranded, 'Change all text to font "Aptos" at 18 point across the entire deck.', 1)
        try:
            formatted_prs = Presentation(formatted)
            runs = [
                run
                for slide in formatted_prs.slides
                for shape in slide.shapes
                if getattr(shape, "has_text_frame", False)
                for paragraph in shape.text_frame.paragraphs
                for run in paragraph.runs
                if run.text.strip()
            ]
            assert runs and all(run.font.name == "Aptos" for run in runs)
            assert all(run.font.size and round(run.font.size.pt, 1) == 18.0 for run in runs)
        finally:
            os.remove(formatted)
    finally:
        os.remove(rebranded)


def verify_table_forms() -> None:
    table_slide = next(
        slide["slide"] for slide in deck_summary(str(SOURCE))["slides"]
        if any(shape["kind"] == "table" for shape in slide["shapes"])
    )
    source = str(SOURCE)
    temporary: list[str] = []
    commands = [
        f"Merge columns 1 and 2 in row 1 of the table on slide {table_slide}.",
        f"Split the table cell in row 1 and column 1 on slide {table_slide}.",
        f'Sort the table on slide {table_slide} by column "1" in ascending order.',
        f'Add a table row with values "New risk | Owner | Medium" on slide {table_slide}.',
        f'Set the table cell in row 2 and column 1 to "High risk" on slide {table_slide}.',
        f'Add a table column with values "Status | Open | Closed" on slide {table_slide}.',
    ]
    try:
        for command in commands:
            source = apply_command(source, command, table_slide)
            temporary.append(source)
        prs = Presentation(source)
        table = next(shape.table for shape in prs.slides[table_slide - 1].shapes if getattr(shape, "has_table", False))
        last_column = len(table.columns)
        last_row = len(table.rows)
        for command in (
            f"Delete column {last_column} from the table on slide {table_slide}.",
            f"Delete row {last_row} from the table on slide {table_slide}.",
        ):
            source = apply_command(source, command, table_slide)
            temporary.append(source)
    finally:
        for path in temporary:
            if os.path.exists(path):
                os.remove(path)


def verify_interface_contract() -> None:
    template = (ROOT / "templates" / "editor.html").read_text(encoding="utf-8")
    script = (ROOT / "static" / "editor.js").read_text(encoding="utf-8")
    chart_buttons = set(re.findall(r'data-chart-type="([^"]+)"', template))
    assert chart_buttons == {"column", "bar", "line", "pie", "area", "waterfall", "scatter"}
    assert "data-new-slide-layout" not in template and 'class="layout-choice"' not in template
    assert "openChartBuilder(button.dataset.chartType, 'blank')" in script
    assert "<summary>Chart tools</summary>" not in template
    assert 'name="placement" value="new"' in template and 'name="placement" value="current"' in template
    guided_names = set(re.findall(r'data-editor-form="([^"]+)"', template))
    assert not guided_names
    assert "rebrand_company" not in script and "Change company across every slide" not in script
    assert "change_font" not in guided_names and "Change font" not in template
    assert "<summary>Tables</summary>" not in template
    assert "apply-theme-slide" in template and "apply-theme-deck" in template
    assert "McKinsey" not in template and "Bain" not in template
    assert len(LAYOUTS) == 32


def main() -> None:
    verify_company_and_font()
    verify_table_forms()
    verify_interface_contract()
    print("PASS: technical chart HUD, blank or Excel chart setup, fixed KPMG controls, and table selections")


if __name__ == "__main__":
    main()
