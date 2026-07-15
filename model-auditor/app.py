"""
Model Auditor
-------------
Reads an Excel workbook's actual formulas (not just its values) and finds
the kinds of errors that show up in real financial model reviews: hardcoded
assumptions mixed into formulas, formulas that break the pattern of their
row or column, circular references, and cached error values. Produces a
model health report slide and an annotated copy of the workbook with
flagged cells highlighted.

This is a different problem from every other tool in this repository:
those all read cell *values*. This one reads formula *structure*, which
means parsing Excel formula syntax, building a cell dependency graph, and
comparing formula "shapes" across a range, none of which existed anywhere
else in this codebase.

Single self-contained Flask application. No external services.
"""

import os
import re
import io
import json
import time
import uuid
import shutil
import subprocess
import tempfile
import threading
from collections import defaultdict, Counter

from flask import Flask, request, render_template, send_file, redirect, url_for, flash, abort
from openpyxl import load_workbook
from openpyxl.utils import column_index_from_string, get_column_letter
from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE

BASE_DIR = os.path.dirname(os.path.abspath(__file__))
SESSIONS_DIR = os.path.join(BASE_DIR, "sessions")
os.makedirs(SESSIONS_DIR, exist_ok=True)

app = Flask(__name__)
app.secret_key = "model-auditor-local-secret"
app.config["MAX_CONTENT_LENGTH"] = 50 * 1024 * 1024

KPMG_BLUE = RGBColor(0x00, 0x33, 0x8D)
KPMG_COBALT = RGBColor(0x00, 0x91, 0xDA)
KPMG_LIGHT = RGBColor(0xF2, 0xF5, 0xFA)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x1A, 0x1A, 0x1A)
GREY = RGBColor(0x77, 0x77, 0x77)
GOOD = RGBColor(0x1E, 0x84, 0x49)
BAD = RGBColor(0xC0, 0x39, 0x2B)
WARN = RGBColor(0xA9, 0x70, 0x0D)
GOOD_HEX, BAD_HEX, WARN_HEX = "1E8449", "C0392B", "A9700D"
PALETTE = [KPMG_BLUE, KPMG_COBALT, RGBColor(0x00, 0xA3, 0x93), RGBColor(0x5A, 0x28, 0x8C), RGBColor(0x8A, 0x8D, 0x8F)]


@app.after_request
def _no_cache(response):
    response.headers["Cache-Control"] = "no-store, no-cache, must-revalidate, max-age=0"
    return response


# ---------------------------------------------------------------------------
# FORMULA PARSING ENGINE
# ---------------------------------------------------------------------------

CELL_REF_RE = re.compile(
    r"(?:'[^']+'!|[A-Za-z_][A-Za-z0-9_ ]*!)?(\$?)([A-Z]{1,3})(\$?)([0-9]{1,7})"
    r"(?::(\$?)([A-Z]{1,3})(\$?)([0-9]{1,7}))?"
)
SHEET_PREFIX_RE = re.compile(r"(?:'([^']+)'|([A-Za-z_][A-Za-z0-9_ ]*))!")
NUMBER_RE = re.compile(r"(?<![A-Za-z$])\d+\.?\d*")
ERROR_VALUES = {"#REF!", "#DIV/0!", "#N/A", "#VALUE!", "#NAME?", "#NULL!", "#NUM!"}
SUPPRESSED_LITERALS = {1.0, -1.0}


def extract_cell_refs_with_sheets(formula, current_sheet):
    """Returns a list of (sheet_name, cell_address) tuples referenced by a
    formula, resolving unqualified refs to the current sheet. Range refs
    (A1:B5) are expanded into every cell in the range."""
    refs = []
    for m in CELL_REF_RE.finditer(formula):
        full_match = m.group(0)
        sheet_m = SHEET_PREFIX_RE.match(full_match)
        sheet = sheet_m.group(1) or sheet_m.group(2) if sheet_m else current_sheet
        col1, row1 = m.group(2), m.group(4)
        col2, row2 = m.group(6), m.group(8)
        if col2 and row2:
            c1, c2 = column_index_from_string(col1), column_index_from_string(col2)
            r1, r2 = int(row1), int(row2)
            if (c2 - c1 + 1) * (r2 - r1 + 1) <= 500:  # cap range expansion for sanity
                for c in range(min(c1, c2), max(c1, c2) + 1):
                    for r in range(min(r1, r2), max(r1, r2) + 1):
                        refs.append((sheet, f"{get_column_letter(c)}{r}"))
        else:
            refs.append((sheet, f"{col1}{row1}"))
    return refs


def extract_hardcoded_literals(formula):
    no_strings = re.sub(r'"[^"]*"', '""', formula)
    stripped = CELL_REF_RE.sub(" ", no_strings)
    literals = [float(l) for l in NUMBER_RE.findall(stripped) if l]
    return [l for l in literals if l not in SUPPRESSED_LITERALS]


def normalize_formula_relative(formula, cur_col_idx, cur_row):
    """Replace each same-sheet cell ref with its offset from the current
    cell, so two formulas pointing at different cells but following the
    same structural pattern normalize to the same string."""
    def repl(m):
        full_match = m.group(0)
        if SHEET_PREFIX_RE.match(full_match):
            return "XSHEET"  # cross-sheet refs aren't comparable this way; treat as opaque
        col_abs, col_letters, row_abs, row_digits, col_abs2, col_letters2, row_abs2, row_digits2 = m.groups()
        col_idx = column_index_from_string(col_letters)
        row_idx = int(row_digits)
        col_delta = "A" if col_abs else str(col_idx - cur_col_idx)
        row_delta = "A" if row_abs else str(row_idx - cur_row)
        result = f"R[{row_delta}]C[{col_delta}]"
        if col_letters2 and row_digits2:
            col_idx2 = column_index_from_string(col_letters2)
            row_idx2 = int(row_digits2)
            col_delta2 = "A" if col_abs2 else str(col_idx2 - cur_col_idx)
            row_delta2 = "A" if row_abs2 else str(row_idx2 - cur_row)
            result += f":R[{row_delta2}]C[{col_delta2}]"
        return result
    return CELL_REF_RE.sub(repl, formula)


# ---------------------------------------------------------------------------
# WORKBOOK SCANNER
# ---------------------------------------------------------------------------

SEVERITY_ORDER = {"High": 0, "Medium": 1, "Low": 2}


def scan_workbook(path):
    """Runs every check against the workbook and returns (issues, stats).
    Each issue: {sheet, cell, type, severity, description}."""
    wb_formulas = load_workbook(path, data_only=False)
    wb_values = load_workbook(path, data_only=True)

    issues = []
    formula_cells = []  # (sheet, cell_addr, col_idx, row_idx, formula)
    total_formula_count = 0

    for ws in wb_formulas.worksheets:
        ws_values = wb_values[ws.title]
        for row in ws.iter_rows():
            for cell in row:
                if isinstance(cell.value, str) and cell.value.startswith("="):
                    total_formula_count += 1
                    formula_cells.append((ws.title, cell.coordinate, cell.column, cell.row, cell.value))

                    # --- Check: hardcoded literal in formula ---
                    literals = extract_hardcoded_literals(cell.value)
                    if literals:
                        lit_str = ", ".join(f"{l:g}" for l in literals[:3])
                        issues.append({
                            "sheet": ws.title, "cell": cell.coordinate, "type": "Hardcoded input",
                            "severity": "Medium",
                            "description": f"Formula mixes a hardcoded value ({lit_str}) with cell references. "
                                           f"Consider moving this to a labeled input cell.",
                        })

                    # --- Check: cached error value ---
                    cached = ws_values[cell.coordinate].value
                    if isinstance(cached, str) and cached in ERROR_VALUES:
                        issues.append({
                            "sheet": ws.title, "cell": cell.coordinate, "type": "Error value",
                            "severity": "High",
                            "description": f"Formula evaluates to {cached}.",
                        })

    # --- Check: row/column consistency ---
    issues.extend(check_consistency(formula_cells, axis="row"))
    issues.extend(check_consistency(formula_cells, axis="col"))

    # --- Check: circular references ---
    issues.extend(check_circular_refs(wb_formulas, formula_cells))

    issues.sort(key=lambda i: SEVERITY_ORDER.get(i["severity"], 3))

    stats = {
        "total_formulas": total_formula_count,
        "high_count": sum(1 for i in issues if i["severity"] == "High"),
        "medium_count": sum(1 for i in issues if i["severity"] == "Medium"),
        "low_count": sum(1 for i in issues if i["severity"] == "Low"),
        "sheets_scanned": len(wb_formulas.worksheets),
    }
    stats["health_score"] = compute_health_score(stats)
    return issues, stats


def compute_health_score(stats):
    """0-100 score. Errors and circular refs hurt the most; hardcodes and
    inconsistencies are penalized more lightly since they're not always
    wrong, just worth a second look."""
    if stats["total_formulas"] == 0:
        return 100
    penalty = (stats["high_count"] * 8) + (stats["medium_count"] * 2) + (stats["low_count"] * 1)
    score = 100 - penalty
    return max(0, min(100, round(score)))


def health_grade(score):
    if score >= 90:
        return "A"
    if score >= 80:
        return "B"
    if score >= 65:
        return "C"
    if score >= 50:
        return "D"
    return "F"


def check_consistency(formula_cells, axis):
    """Groups formula cells into contiguous runs along a row or column,
    normalizes each formula relative to its own position, and flags any
    cell whose normalized pattern doesn't match the majority pattern in
    its run (minimum run length 4, so isolated formulas aren't flagged)."""
    issues = []
    groups = defaultdict(list)
    for sheet, addr, col, row, formula in formula_cells:
        key = (sheet, row) if axis == "row" else (sheet, col)
        groups[key].append((addr, col, row, formula))

    for key, cells in groups.items():
        if len(cells) < 4:
            continue
        cells.sort(key=lambda c: c[1] if axis == "row" else c[2])
        # only compare cells that are contiguous (no gaps) to avoid comparing
        # unrelated formulas that happen to share a row/column
        runs = []
        current_run = [cells[0]]
        for prev, cur in zip(cells, cells[1:]):
            prev_pos = prev[1] if axis == "row" else prev[2]
            cur_pos = cur[1] if axis == "row" else cur[2]
            if cur_pos - prev_pos == 1:
                current_run.append(cur)
            else:
                runs.append(current_run)
                current_run = [cur]
        runs.append(current_run)

        for run in runs:
            if len(run) < 4:
                continue
            normalized = [normalize_formula_relative(f, col, row) for _, col, row, f in run]
            pattern_counts = Counter(normalized)
            majority_pattern, majority_count = pattern_counts.most_common(1)[0]
            if majority_count < len(run) * 0.6:
                continue  # no clear dominant pattern, too heterogeneous to judge
            for (addr, col, row, formula), norm in zip(run, normalized):
                if norm != majority_pattern:
                    issues.append({
                        "sheet": key[0], "cell": addr, "type": "Inconsistent formula",
                        "severity": "High",
                        "description": f"This {'row' if axis=='row' else 'column'} follows the pattern "
                                       f"'{majority_pattern}' in {majority_count} of {len(run)} adjacent cells; "
                                       f"this cell's formula ('{formula}') breaks that pattern.",
                    })
    return issues


def check_circular_refs(wb_formulas, formula_cells):
    """Builds a directed graph (cell -> cells it depends on) and detects
    cycles via DFS. Reports the first cell of each distinct cycle found."""
    graph = {}
    for sheet, addr, col, row, formula in formula_cells:
        refs = extract_cell_refs_with_sheets(formula, sheet)
        graph[(sheet, addr)] = set(refs)

    issues = []
    visited = set()
    in_stack = set()
    reported_cycles = set()

    def dfs(node, path):
        if node in in_stack:
            cycle_start = path.index(node)
            cycle = tuple(sorted(path[cycle_start:]))
            if cycle not in reported_cycles:
                reported_cycles.add(cycle)
                sheet, addr = node
                issues.append({
                    "sheet": sheet, "cell": addr, "type": "Circular reference",
                    "severity": "High",
                    "description": f"This cell is part of a circular reference chain "
                                    f"({len(path) - cycle_start} cells involved).",
                })
            return
        if node in visited or node not in graph:
            return
        visited.add(node)
        in_stack.add(node)
        path.append(node)
        for ref in graph.get(node, ()):
            dfs(ref, path)
        path.pop()
        in_stack.discard(node)

    for node in list(graph.keys()):
        if node not in visited:
            dfs(node, [])
    return issues


# ---------------------------------------------------------------------------
# PPTX HELPERS (word_wrap + auto_size set explicitly on every text box)
# ---------------------------------------------------------------------------

def add_textbox(slide, x, y, w, h, text, size=14, bold=False, italic=False,
                 color=DARK, align=PP_ALIGN.LEFT, font="Arial", wrap=True):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = wrap
    tf.auto_size = MSO_AUTO_SIZE.NONE
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font
    return box


def add_header_bar(slide, prs, title_text, subtitle_text=None):
    bar = slide.shapes.add_shape(1, 0, 0, prs.slide_width, Inches(0.9))
    bar.fill.solid()
    bar.fill.fore_color.rgb = KPMG_BLUE
    bar.line.fill.background()
    tf = bar.text_frame
    tf.word_wrap = False
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.margin_left = Inches(0.3)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = "KPMG"
    r.font.size = Pt(22)
    r.font.bold = True
    r.font.color.rgb = WHITE
    r.font.name = "Arial"
    add_textbox(slide, 3.4, 0.12, 9.6, 0.65, title_text, size=16, bold=True, color=WHITE, align=PP_ALIGN.RIGHT)
    if subtitle_text:
        add_textbox(slide, 0.4, 1.0, 12.5, 0.35, subtitle_text, size=11, italic=True, color=GREY)


def add_kpi_tile(slide, x, y, w, h, value, label, value_color=KPMG_BLUE):
    tile = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    tile.fill.solid()
    tile.fill.fore_color.rgb = KPMG_LIGHT
    tile.line.color.rgb = value_color
    tile.line.width = Pt(1)
    tf = tile.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.margin_left = Inches(0.15)
    tf.margin_right = Inches(0.15)
    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    r1 = p1.add_run()
    r1.text = value
    r1.font.size = Pt(26)
    r1.font.bold = True
    r1.font.color.rgb = value_color
    r1.font.name = "Arial"
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = label
    r2.font.size = Pt(11)
    r2.font.color.rgb = DARK
    r2.font.name = "Arial"


def add_table(slide, x, y, w, h, headers, rows, col_widths=None, row_colorizer=None, font_size=10.5):
    n_rows = 1 + len(rows)
    shape = slide.shapes.add_table(n_rows, len(headers), Inches(x), Inches(y), Inches(w), Inches(h))
    table = shape.table
    if col_widths:
        for i, cw in enumerate(col_widths):
            table.columns[i].width = Inches(cw)
    for c, htext in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = str(htext)
        run = cell.text_frame.paragraphs[0].runs[0]
        run.font.size = Pt(font_size + 1)
        run.font.bold = True
        run.font.color.rgb = WHITE
        run.font.name = "Arial"
        cell.fill.solid()
        cell.fill.fore_color.rgb = KPMG_BLUE
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT
    for r_idx, row in enumerate(rows, start=1):
        for c_idx, val in enumerate(row):
            cell = table.cell(r_idx, c_idx)
            cell.text = str(val)
            run = cell.text_frame.paragraphs[0].runs[0]
            run.font.size = Pt(font_size)
            run.font.name = "Arial"
            run.font.color.rgb = DARK
            cell.fill.solid()
            cell.fill.fore_color.rgb = KPMG_LIGHT if (r_idx % 2 == 0) else WHITE
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER if c_idx > 0 else PP_ALIGN.LEFT
            if row_colorizer:
                override = row_colorizer(r_idx - 1, c_idx, val)
                if override:
                    run.font.color.rgb = override
                    run.font.bold = True
    return table


def new_slide(prs=None):
    if prs is None:
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
    return prs, prs.slides.add_slide(prs.slide_layouts[6])


# ---------------------------------------------------------------------------
# REPORT GENERATION
# ---------------------------------------------------------------------------

def build_report_pptx(issues, stats, title, subtitle, filename, out_path):
    prs, slide = new_slide()
    add_header_bar(slide, prs, title, subtitle)

    score = stats["health_score"]
    grade = health_grade(score)
    score_color = GOOD if score >= 80 else (WARN if score >= 50 else BAD)

    add_kpi_tile(slide, 0.4, 1.5, 2.6, 1.6, f"{score}", f"Health score (Grade {grade})", value_color=score_color)
    add_kpi_tile(slide, 3.15, 1.5, 2.6, 1.6, str(stats["total_formulas"]), "Formulas scanned")
    add_kpi_tile(slide, 5.9, 1.5, 2.6, 1.6, str(stats["high_count"]), "High severity", value_color=BAD)
    add_kpi_tile(slide, 8.65, 1.5, 2.6, 1.6, str(stats["medium_count"]), "Medium severity", value_color=WARN)
    add_kpi_tile(slide, 11.4, 1.5, 1.5, 1.6, str(stats["sheets_scanned"]), "Sheets")

    add_textbox(slide, 0.4, 3.35, 6, 0.35, "Findings by Type", size=14, bold=True, color=KPMG_BLUE)
    type_counts = Counter(i["type"] for i in issues)
    if type_counts:
        chart_data = CategoryChartData()
        chart_data.categories = list(type_counts.keys())
        chart_data.add_series("Count", tuple(type_counts.values()))
        gf = slide.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, Inches(0.4), Inches(3.75),
                                     Inches(5.8), Inches(3.3), chart_data)
        chart = gf.chart
        chart.has_legend = False
        for i, pt in enumerate(chart.series[0].points):
            pt.format.fill.solid()
            pt.format.fill.fore_color.rgb = PALETTE[i % len(PALETTE)]
        plot = chart.plots[0]
        plot.has_data_labels = True
        plot.data_labels.font.size = Pt(10)

    add_textbox(slide, 6.5, 3.35, 6, 0.35, "Top Findings", size=14, bold=True, color=KPMG_BLUE)
    top_issues = issues[:8]
    top_rows = [[f"{i['sheet']}!{i['cell']}", i["type"], i["severity"]] for i in top_issues]

    def colorize(row_idx, col_idx, val):
        if col_idx == 2:
            if val == "High":
                return BAD
            if val == "Medium":
                return WARN
        return None

    add_table(slide, 6.5, 3.75, 6.4, min(3.3, 0.35 * (len(top_rows) + 1)),
              ["Location", "Type", "Severity"], top_rows,
              col_widths=[2.2, 2.9, 1.3], row_colorizer=colorize, font_size=9.5)

    # --- Slide 2: full findings table ---
    if issues:
        prs, detail_slide = new_slide(prs)
        add_header_bar(detail_slide, prs, "Detailed Findings", filename)
        detail_rows = [[f"{i['sheet']}!{i['cell']}", i["type"], i["severity"], i["description"]]
                        for i in issues[:18]]

        def colorize2(row_idx, col_idx, val):
            if col_idx == 2:
                if val == "High":
                    return BAD
                if val == "Medium":
                    return WARN
            return None

        table_h = min(5.6, 0.4 * (len(detail_rows) + 1))
        add_table(detail_slide, 0.4, 1.5, 12.5, table_h,
                  ["Location", "Type", "Severity", "Description"], detail_rows,
                  col_widths=[1.6, 2.2, 1.2, 7.5], row_colorizer=colorize2, font_size=9.5)
        if len(issues) > 18:
            add_textbox(detail_slide, 0.4, 1.5 + table_h + 0.1, 10, 0.3,
                        f"+ {len(issues) - 18} additional finding(s) in the annotated workbook.",
                        size=10, italic=True, color=GREY)

    prs.save(out_path)


SEVERITY_FILL = {"High": "FBE4E1", "Medium": "FFF1D6", "Low": "E9ECEF"}
SEVERITY_FONT = {"High": BAD_HEX, "Medium": WARN_HEX, "Low": "555555"}


def build_annotated_xlsx(input_path, issues, out_path):
    """Copies the workbook and highlights every flagged cell in place, then
    adds an Audit Findings sheet listing every issue."""
    wb = load_workbook(input_path, data_only=False)

    issues_by_cell = defaultdict(list)
    for issue in issues:
        if issue["sheet"] in wb.sheetnames:
            issues_by_cell[(issue["sheet"], issue["cell"])].append(issue)

    from openpyxl.comments import Comment
    for (sheet, addr), cell_issues in issues_by_cell.items():
        ws = wb[sheet]
        try:
            cell = ws[addr]
        except Exception:
            continue
        cell_issues.sort(key=lambda i: SEVERITY_ORDER.get(i["severity"], 3))
        worst_severity = cell_issues[0]["severity"]
        cell.fill = PatternFill("solid", fgColor=SEVERITY_FILL.get(worst_severity, "E9ECEF"))
        note = "\n".join(f"[{i['severity']}] {i['type']}: {i['description']}" for i in cell_issues)
        cell.comment = Comment(note, "Model Auditor")

    ws2 = wb.create_sheet("Audit Findings")
    ws2["A1"] = "Model Audit Findings"
    ws2["A1"].font = Font(bold=True, size=14, color="00338D")
    headers = ["Sheet", "Cell", "Type", "Severity", "Description"]
    for c, h in enumerate(headers, start=1):
        cell = ws2.cell(row=3, column=c, value=h)
        cell.font = Font(bold=True, color="FFFFFF")
        cell.fill = PatternFill("solid", fgColor="00338D")
    border = Border(*(Side(style="thin", color="D7DEEA"),) * 4)
    for i, issue in enumerate(issues, start=1):
        r = 3 + i
        vals = [issue["sheet"], issue["cell"], issue["type"], issue["severity"], issue["description"]]
        for c, v in enumerate(vals, start=1):
            cell = ws2.cell(row=r, column=c, value=v)
            cell.border = border
            cell.alignment = Alignment(wrap_text=(c == 5), vertical="top")
            if c == 4:
                cell.font = Font(color=SEVERITY_FONT.get(v, "000000"), bold=True)
        if i % 2 == 0:
            for c in range(1, 6):
                ws2.cell(row=r, column=c).fill = PatternFill("solid", fgColor="F2F5FA")
    ws2.column_dimensions["A"].width = 16
    ws2.column_dimensions["B"].width = 10
    ws2.column_dimensions["C"].width = 22
    ws2.column_dimensions["D"].width = 12
    ws2.column_dimensions["E"].width = 80

    wb.save(out_path)


# ---------------------------------------------------------------------------
# REAL SLIDE PREVIEW (LibreOffice + PyMuPDF)
# ---------------------------------------------------------------------------

SOFFICE_PATH = shutil.which("soffice") or shutil.which("libreoffice")
_soffice_lock = threading.Lock()
try:
    import fitz
    HAVE_FITZ = True
except ImportError:
    HAVE_FITZ = False


def render_pptx_preview(pptx_path, out_png_path):
    if not (SOFFICE_PATH and HAVE_FITZ):
        return False
    out_dir = os.path.dirname(out_png_path)
    profile_dir = tempfile.mkdtemp(prefix="lo_profile_")
    try:
        with _soffice_lock:
            result = subprocess.run(
                [SOFFICE_PATH, "--headless", "--norestore",
                 f"-env:UserInstallation=file://{profile_dir}",
                 "--convert-to", "pdf", "--outdir", out_dir, pptx_path],
                capture_output=True, timeout=90,
            )
        pdf_path = os.path.join(out_dir, os.path.splitext(os.path.basename(pptx_path))[0] + ".pdf")
        if result.returncode != 0 or not os.path.exists(pdf_path):
            return False
        doc = fitz.open(pdf_path)
        page = doc[0]
        pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
        pix.save(out_png_path)
        doc.close()
        os.remove(pdf_path)
        return True
    except Exception:
        return False
    finally:
        shutil.rmtree(profile_dir, ignore_errors=True)


def _session_dir(sid):
    d = os.path.join(SESSIONS_DIR, sid)
    if not os.path.isdir(d):
        abort(404)
    return d


# ---------------------------------------------------------------------------
# ROUTES
# ---------------------------------------------------------------------------

@app.route("/", methods=["GET"])
def index():
    return render_template("index.html")


@app.route("/process", methods=["POST"])
def process():
    model_file = request.files.get("model_file")
    if not model_file or model_file.filename == "":
        flash("Please upload an Excel workbook (.xlsx) to audit.")
        return redirect(url_for("index"))

    title = request.form.get("title", "").strip() or "Model Health Report"
    subtitle = request.form.get("subtitle", "").strip()

    sid = uuid.uuid4().hex[:12]
    sess_dir = os.path.join(SESSIONS_DIR, sid)
    os.makedirs(sess_dir, exist_ok=True)
    model_path = os.path.join(sess_dir, "model.xlsx")
    model_file.save(model_path)

    try:
        issues, stats = scan_workbook(model_path)
    except Exception as e:
        flash(f"Couldn't read that workbook: {e}")
        shutil.rmtree(sess_dir, ignore_errors=True)
        return redirect(url_for("index"))

    with open(os.path.join(sess_dir, "issues.json"), "w") as f:
        json.dump(issues, f)
    with open(os.path.join(sess_dir, "stats.json"), "w") as f:
        json.dump(stats, f)
    with open(os.path.join(sess_dir, "config.json"), "w") as f:
        json.dump({"title": title, "subtitle": subtitle, "filename": model_file.filename}, f)

    return redirect(url_for("review", sid=sid))


@app.route("/review/<sid>", methods=["GET"])
def review(sid):
    sess_dir = _session_dir(sid)
    issues_path = os.path.join(sess_dir, "issues.json")
    if not os.path.exists(issues_path):
        flash("That session has expired. Please upload again.")
        return redirect(url_for("index"))
    with open(issues_path) as f:
        issues = json.load(f)
    with open(os.path.join(sess_dir, "stats.json")) as f:
        stats = json.load(f)
    with open(os.path.join(sess_dir, "config.json")) as f:
        config = json.load(f)
    grade = health_grade(stats["health_score"])
    return render_template("review.html", sid=sid, issues=issues, stats=stats, config=config, grade=grade)


@app.route("/generate/<sid>", methods=["POST"])
def generate(sid):
    sess_dir = _session_dir(sid)
    issues_path = os.path.join(sess_dir, "issues.json")
    if not os.path.exists(issues_path):
        flash("That session has expired. Please upload again.")
        return redirect(url_for("index"))
    with open(issues_path) as f:
        issues = json.load(f)
    with open(os.path.join(sess_dir, "stats.json")) as f:
        stats = json.load(f)
    with open(os.path.join(sess_dir, "config.json")) as f:
        config = json.load(f)

    model_path = os.path.join(sess_dir, "model.xlsx")
    report_path = os.path.join(sess_dir, "model_health_report.pptx")
    annotated_path = os.path.join(sess_dir, "annotated_model.xlsx")

    build_report_pptx(issues, stats, config["title"], config["subtitle"], config["filename"], report_path)
    build_annotated_xlsx(model_path, issues, annotated_path)

    preview_path = os.path.join(sess_dir, "preview.png")
    rendering_ok = render_pptx_preview(report_path, preview_path)

    with open(os.path.join(sess_dir, "meta.json"), "w") as f:
        json.dump({"rendering_ok": rendering_ok}, f)
    return redirect(url_for("result", sid=sid))


@app.route("/result/<sid>", methods=["GET"])
def result(sid):
    sess_dir = _session_dir(sid)
    report_path = os.path.join(sess_dir, "model_health_report.pptx")
    if not os.path.exists(report_path):
        flash("That session has expired. Please upload again.")
        return redirect(url_for("index"))
    with open(os.path.join(sess_dir, "stats.json")) as f:
        stats = json.load(f)
    meta_path = os.path.join(sess_dir, "meta.json")
    rendering_ok = False
    if os.path.exists(meta_path):
        with open(meta_path) as f:
            rendering_ok = json.load(f).get("rendering_ok", False)
    grade = health_grade(stats["health_score"])
    return render_template("result.html", sid=sid, stats=stats, grade=grade, rendering_ok=rendering_ok,
                            soffice_missing=not (SOFFICE_PATH and HAVE_FITZ))


@app.route("/preview_image/<sid>", methods=["GET"])
def preview_image(sid):
    sess_dir = _session_dir(sid)
    path = os.path.join(sess_dir, "preview.png")
    if not os.path.exists(path):
        abort(404)
    return send_file(path, mimetype="image/png")


@app.route("/download/<sid>/<which>", methods=["GET"])
def download(sid, which):
    sess_dir = _session_dir(sid)
    files = {
        "report": ("model_health_report.pptx",
                   "application/vnd.openxmlformats-officedocument.presentationml.presentation"),
        "annotated": ("annotated_model.xlsx",
                      "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"),
    }
    if which not in files:
        abort(404)
    fname, mimetype = files[which]
    path = os.path.join(sess_dir, fname)
    if not os.path.exists(path):
        abort(404)
    return send_file(path, as_attachment=True, download_name=fname, mimetype=mimetype)


def _open_browser():
    import webbrowser
    import sys as _sys
    time.sleep(1.2)
    url = "http://127.0.0.1:5090"
    chrome_candidates = []
    if _sys.platform == "darwin":
        chrome_candidates.append("open -a 'Google Chrome' %s")
    elif _sys.platform.startswith("win"):
        for p in (r"C:\Program Files\Google\Chrome\Application\chrome.exe",
                  r"C:\Program Files (x86)\Google\Chrome\Application\chrome.exe"):
            if os.path.exists(p):
                chrome_candidates.append(p.replace("\\", "\\\\") + " %s")
    else:
        chrome_candidates += ["google-chrome %s", "chromium-browser %s", "chromium %s"]
    for template in chrome_candidates:
        try:
            webbrowser.get(template).open(url)
            return
        except webbrowser.Error:
            continue
    webbrowser.open(url)


if __name__ == "__main__":
    import sys
    if os.environ.get("WERKZEUG_RUN_MAIN") != "true":
        threading.Thread(target=_open_browser, daemon=True).start()
    print("Starting Model Auditor at http://127.0.0.1:5090")
    app.run(host="127.0.0.1", port=5090, debug=("--debug" in sys.argv))
