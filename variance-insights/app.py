"""
Variance Insights
------------------
Upload a financial data table (Actual vs. Budget/Prior Period/Forecast),
and this app computes dollar and percent variances, flags which ones are
material, classifies each as favorable/unfavorable based on the metric
type (revenue-like vs. cost-like), and drafts consulting-style commentary
bullets for the material ones -- a first-pass draft an analyst can review
and finish, not a black box. Exports a client-ready PowerPoint slide and
a formatted Excel workbook, both KPMG-styled.

Single self-contained file, no external services, nothing leaves this
machine.
"""

import os
import re
import io
import json
import time
import uuid
import shutil
import subprocess
import tempfile
import threading

from flask import Flask, request, render_template, send_file, redirect, url_for, flash, abort
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import openpyxl
from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
from openpyxl.utils import get_column_letter

BASE_DIR = os.path.dirname(os.path.abspath(__file__))
SESSIONS_DIR = os.path.join(BASE_DIR, "sessions")
os.makedirs(SESSIONS_DIR, exist_ok=True)

app = Flask(__name__)
app.secret_key = "variance-insights-local-secret"
app.config["MAX_CONTENT_LENGTH"] = 50 * 1024 * 1024  # 50MB

KPMG_BLUE = RGBColor(0x00, 0x33, 0x8D)
KPMG_COBALT = RGBColor(0x00, 0x91, 0xDA)
KPMG_LIGHT = RGBColor(0xF2, 0xF5, 0xFA)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x1A, 0x1A, 0x1A)
GOOD_HEX = "1E8449"
BAD_HEX = "C0392B"


@app.after_request
def _no_cache(response):
    response.headers["Cache-Control"] = "no-store, no-cache, must-revalidate, max-age=0"
    return response


# ---------------------------------------------------------------------------
# COLUMN AUTO-DETECTION
# ---------------------------------------------------------------------------

ACTUAL_HINTS = ["actual", "current", "this year", "ty"]
COMPARISON_HINTS = ["budget", "plan", "forecast", "prior", "py", "last year", "ly", "target"]


def detect_columns(df):
    """Given a dataframe, guess which column is the line-item label, which
    is 'Actual', and which is the comparison (Budget/Prior/Forecast)."""
    cols = list(df.columns)
    numeric_cols = [c for c in cols if pd.api.types.is_numeric_dtype(df[c])]
    non_numeric_cols = [c for c in cols if c not in numeric_cols]

    label_col = non_numeric_cols[0] if non_numeric_cols else cols[0]

    actual_col = None
    comparison_col = None
    for c in numeric_cols:
        cl = str(c).lower()
        if actual_col is None and any(h in cl for h in ACTUAL_HINTS):
            actual_col = c
        elif comparison_col is None and any(h in cl for h in COMPARISON_HINTS):
            comparison_col = c

    remaining = [c for c in numeric_cols if c not in (actual_col, comparison_col)]
    if actual_col is None and remaining:
        actual_col = remaining.pop(0)
    if comparison_col is None and remaining:
        comparison_col = remaining.pop(0)

    return label_col, actual_col, comparison_col


# ---------------------------------------------------------------------------
# METRIC CLASSIFICATION (unit type + favorable direction)
# ---------------------------------------------------------------------------

PERCENT_WORDS = ["%", "margin", "rate", "utilization", "percentage", "ratio"]
COUNT_WORDS = ["count", "engagement", "headcount", "clients", "employees", "fte", "days", "staff"]
COST_LIKE_WORDS = ["cost", "expense", "tax", "interest", "loss", "spend", "churn", "attrition", "days sales"]
REVENUE_LIKE_WORDS = ["revenue", "income", "profit", "margin", "contribution", "gross", "ebitda",
                       "ebit", "earnings", "utilization", "rate"]


def classify_unit(label):
    l = label.lower()
    if any(w in l for w in PERCENT_WORDS):
        return "percent"
    if any(w in l for w in COUNT_WORDS):
        return "count"
    return "currency"


def classify_direction(label):
    """Returns 'higher_is_good', 'lower_is_good', or 'neutral'.
    Revenue-like words are checked first so a line like "Tax Revenue"
    (which contains "tax", a cost-like word) is correctly read as revenue
    rather than as a cost, unlike "Tax Expense"."""
    l = label.lower()
    if any(w in l for w in REVENUE_LIKE_WORDS):
        return "higher_is_good"
    if any(w in l for w in COST_LIKE_WORDS):
        return "lower_is_good"
    return "neutral"


def format_value(val, unit, currency_suffix):
    if val is None:
        return "-"
    if unit == "percent":
        return f"{val:.1f}%"
    if unit == "count":
        return f"{val:,.0f}" if val == int(val) else f"{val:,.1f}"
    return f"${val:,.1f}{currency_suffix}"


def format_variance(var, unit, currency_suffix):
    sign = "+" if var >= 0 else "-"
    absval = abs(var)
    if unit == "percent":
        return f"{sign}{absval:.1f} pts"
    if unit == "count":
        return f"{sign}{absval:,.0f}" if absval == int(absval) else f"{sign}{absval:,.1f}"
    return f"{sign}${absval:,.1f}{currency_suffix}"


# ---------------------------------------------------------------------------
# CORE ANALYSIS
# ---------------------------------------------------------------------------

def analyze(df, label_col, actual_col, comparison_col, comparison_label,
            materiality_pct, currency_suffix):
    rows = []
    for _, row in df.iterrows():
        label = str(row[label_col]).strip() if pd.notna(row[label_col]) else ""
        actual = row[actual_col]
        comparison = row[comparison_col]
        if not label or pd.isna(actual) or pd.isna(comparison):
            continue
        try:
            actual = float(actual)
            comparison = float(comparison)
        except (ValueError, TypeError):
            continue

        unit = classify_unit(label)
        direction = classify_direction(label)

        if unit == "percent":
            variance = actual - comparison  # percentage-point difference
            variance_pct = None
        else:
            variance = actual - comparison
            variance_pct = (variance / comparison * 100) if comparison != 0 else None

        pct_for_materiality = variance_pct if variance_pct is not None else (
            abs(variance) if unit == "percent" else None
        )
        is_material = pct_for_materiality is not None and abs(pct_for_materiality) >= materiality_pct

        if direction == "higher_is_good":
            favorable = variance > 0
        elif direction == "lower_is_good":
            favorable = variance < 0
        else:
            favorable = None

        if favorable is True:
            tag = "Favorable"
        elif favorable is False:
            tag = "Unfavorable"
        else:
            tag = "Notable" if is_material else "Neutral"

        direction_word = "up" if variance > 0 else ("down" if variance < 0 else "flat")
        commentary = draft_commentary(label, variance, variance_pct, unit, tag,
                                       comparison_label, currency_suffix, direction_word)

        rows.append({
            "label": label,
            "actual": actual,
            "comparison": comparison,
            "actual_display": format_value(actual, unit, currency_suffix),
            "comparison_display": format_value(comparison, unit, currency_suffix),
            "variance": variance,
            "variance_display": format_variance(variance, unit, currency_suffix),
            "variance_pct": variance_pct,
            "variance_pct_display": f"{variance_pct:+.1f}%" if variance_pct is not None else "n/a",
            "unit": unit,
            "tag": tag,
            "is_material": is_material,
            "commentary": commentary,
        })
    return rows


def draft_commentary(label, variance, variance_pct, unit, tag, comparison_label,
                      currency_suffix, direction_word):
    if unit == "percent":
        magnitude = f"{abs(variance):.1f} percentage points"
    else:
        magnitude = format_variance(variance, unit, currency_suffix).lstrip("+-")
        if variance_pct is not None:
            magnitude += f" ({abs(variance_pct):.1f}%)"

    base = f"{label} came in {direction_word} {magnitude} versus {comparison_label}"
    if tag == "Favorable":
        return base + ", a favorable variance. [Add driver commentary]"
    if tag == "Unfavorable":
        return base + ", an unfavorable variance that warrants explanation. [Add driver commentary]"
    if tag == "Notable":
        return base + ". [Add context]"
    return base + "."


# ---------------------------------------------------------------------------
# EXPORT: PPTX SLIDE
# ---------------------------------------------------------------------------

def build_pptx(rows, title, subtitle, comparison_label, out_path):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bar = slide.shapes.add_shape(1, 0, 0, prs.slide_width, Inches(0.9))
    bar.fill.solid()
    bar.fill.fore_color.rgb = KPMG_BLUE
    bar.line.fill.background()
    tf = bar.text_frame
    tf.margin_left = Inches(0.3)
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "KPMG"
    r.font.size = Pt(24)
    r.font.bold = True
    r.font.color.rgb = WHITE
    r.font.name = "Arial"

    title_box = slide.shapes.add_textbox(Inches(3.4), Inches(0.12), Inches(9.6), Inches(0.65))
    tp = title_box.text_frame.paragraphs[0]
    tp.alignment = PP_ALIGN.RIGHT
    tr = tp.add_run()
    tr.text = title
    tr.font.size = Pt(17)
    tr.font.bold = True
    tr.font.color.rgb = WHITE
    tr.font.name = "Arial"

    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.4), Inches(1.0), Inches(9), Inches(0.35))
        sp = sub_box.text_frame.paragraphs[0]
        sr = sp.add_run()
        sr.text = subtitle
        sr.font.size = Pt(11)
        sr.font.italic = True
        sr.font.color.rgb = RGBColor(0x77, 0x77, 0x77)

    # table: Line Item | Actual | Comparison | Variance | Variance %
    headers = ["Line Item", "Actual", comparison_label, "Variance", "Variance %"]
    table_rows = 1 + len(rows)
    table_w = Inches(12.5)
    tbl_shape = slide.shapes.add_table(table_rows, 5, Inches(0.4), Inches(1.5), table_w, Inches(0.35 * table_rows))
    table = tbl_shape.table
    widths = [Inches(4.3), Inches(2.0), Inches(2.0), Inches(2.1), Inches(2.1)]
    for i, w in enumerate(widths):
        table.columns[i].width = w

    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = h
        run = cell.text_frame.paragraphs[0].runs[0]
        run.font.size = Pt(12)
        run.font.bold = True
        run.font.color.rgb = WHITE
        run.font.name = "Arial"
        cell.fill.solid()
        cell.fill.fore_color.rgb = KPMG_BLUE
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT

    for r_idx, row in enumerate(rows, start=1):
        values = [row["label"], row["actual_display"], row["comparison_display"],
                  row["variance_display"], row["variance_pct_display"]]
        for c_idx, val in enumerate(values):
            cell = table.cell(r_idx, c_idx)
            cell.text = str(val)
            run = cell.text_frame.paragraphs[0].runs[0]
            run.font.size = Pt(11)
            run.font.name = "Arial"
            cell.fill.solid()
            cell.fill.fore_color.rgb = KPMG_LIGHT if (r_idx % 2 == 0) else WHITE
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER if c_idx > 0 else PP_ALIGN.LEFT
            if c_idx in (3, 4):
                if row["tag"] == "Favorable":
                    run.font.color.rgb = RGBColor(0x1E, 0x84, 0x49)
                    run.font.bold = True
                elif row["tag"] == "Unfavorable":
                    run.font.color.rgb = RGBColor(0xC0, 0x39, 0x2B)
                    run.font.bold = True
                else:
                    run.font.color.rgb = DARK
            else:
                run.font.color.rgb = DARK

    # commentary box below the table
    commentary_rows = [r for r in rows if r["is_material"]]
    if commentary_rows:
        cy = Inches(1.5) + Inches(0.35 * table_rows) + Inches(0.25)
        max_h = prs.slide_height - cy - Inches(0.3)
        label_box = slide.shapes.add_textbox(Inches(0.4), cy, Inches(4), Inches(0.3))
        lp = label_box.text_frame.paragraphs[0]
        lr = lp.add_run()
        lr.text = "Key Variance Commentary"
        lr.font.size = Pt(13)
        lr.font.bold = True
        lr.font.color.rgb = KPMG_BLUE

        box = slide.shapes.add_textbox(Inches(0.4), cy + Inches(0.35), Inches(12.5), max_h - Inches(0.35))
        btf = box.text_frame
        btf.word_wrap = True
        for i, row in enumerate(commentary_rows):
            para = btf.paragraphs[0] if i == 0 else btf.add_paragraph()
            run = para.add_run()
            run.text = "•  " + row["commentary"]
            run.font.size = Pt(11)
            run.font.color.rgb = DARK

    prs.save(out_path)


# ---------------------------------------------------------------------------
# EXPORT: EXCEL WORKBOOK
# ---------------------------------------------------------------------------

def build_xlsx(rows, title, comparison_label, out_path):
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Variance Analysis"

    ws.merge_cells("A1:E1")
    ws["A1"] = title
    ws["A1"].font = Font(bold=True, size=14, color="00338D")

    headers = ["Line Item", "Actual", comparison_label, "Variance", "Variance %"]
    header_row = 3
    for c, h in enumerate(headers, start=1):
        cell = ws.cell(row=header_row, column=c, value=h)
        cell.font = Font(bold=True, color="FFFFFF")
        cell.fill = PatternFill("solid", fgColor="00338D")
        cell.alignment = Alignment(horizontal="center" if c > 1 else "left")

    thin = Side(style="thin", color="D7DEEA")
    border = Border(left=thin, right=thin, top=thin, bottom=thin)

    for i, row in enumerate(rows, start=1):
        r = header_row + i
        ws.cell(row=r, column=1, value=row["label"]).border = border
        ws.cell(row=r, column=2, value=row["actual_display"]).border = border
        ws.cell(row=r, column=3, value=row["comparison_display"]).border = border
        var_cell = ws.cell(row=r, column=4, value=row["variance_display"])
        pct_cell = ws.cell(row=r, column=5, value=row["variance_pct_display"])
        var_cell.border = border
        pct_cell.border = border
        for cell in (var_cell, pct_cell):
            cell.alignment = Alignment(horizontal="center")
            if row["tag"] == "Favorable":
                cell.font = Font(color=GOOD_HEX, bold=True)
            elif row["tag"] == "Unfavorable":
                cell.font = Font(color=BAD_HEX, bold=True)
        if i % 2 == 0:
            for c in range(1, 6):
                ws.cell(row=r, column=c).fill = PatternFill("solid", fgColor="F2F5FA")

    ws.column_dimensions["A"].width = 32
    for col in "BCDE":
        ws.column_dimensions[col].width = 16

    # commentary sheet
    ws2 = wb.create_sheet("Commentary")
    ws2["A1"] = "Key Variance Commentary"
    ws2["A1"].font = Font(bold=True, size=13, color="00338D")
    r = 3
    for row in rows:
        if row["is_material"]:
            ws2.cell(row=r, column=1, value=row["label"]).font = Font(bold=True)
            ws2.cell(row=r, column=1).alignment = Alignment(wrap_text=True)
            ws2.merge_cells(start_row=r+1, start_column=1, end_row=r+1, end_column=4)
            c = ws2.cell(row=r+1, column=1, value=row["commentary"])
            c.alignment = Alignment(wrap_text=True)
            r += 3
    ws2.column_dimensions["A"].width = 100

    wb.save(out_path)


# ---------------------------------------------------------------------------
# REAL SLIDE PREVIEW (reuses the LibreOffice + PyMuPDF approach)
# ---------------------------------------------------------------------------

SOFFICE_PATH = shutil.which("soffice") or shutil.which("libreoffice")
_soffice_lock = threading.Lock()
try:
    import fitz
    HAVE_FITZ = True
except ImportError:
    HAVE_FITZ = False


def render_pptx_preview(pptx_path, out_png_path):
    if not (SOFFICE_PATH and HAVE_FITZ):
        return False
    out_dir = os.path.dirname(out_png_path)
    profile_dir = tempfile.mkdtemp(prefix="lo_profile_")
    try:
        with _soffice_lock:
            result = subprocess.run(
                [SOFFICE_PATH, "--headless", "--norestore",
                 f"-env:UserInstallation=file://{profile_dir}",
                 "--convert-to", "pdf", "--outdir", out_dir, pptx_path],
                capture_output=True, timeout=90,
            )
        pdf_path = os.path.join(out_dir, os.path.splitext(os.path.basename(pptx_path))[0] + ".pdf")
        if result.returncode != 0 or not os.path.exists(pdf_path):
            return False
        doc = fitz.open(pdf_path)
        page = doc[0]
        pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
        pix.save(out_png_path)
        doc.close()
        os.remove(pdf_path)
        return True
    except Exception:
        return False
    finally:
        shutil.rmtree(profile_dir, ignore_errors=True)


# ---------------------------------------------------------------------------
# ROUTES
# ---------------------------------------------------------------------------

def _session_dir(sid):
    d = os.path.join(SESSIONS_DIR, sid)
    if not os.path.isdir(d):
        abort(404)
    return d


@app.route("/", methods=["GET"])
def index():
    return render_template("index.html")


@app.route("/process", methods=["POST"])
def process():
    data_file = request.files.get("data_file")
    if not data_file or data_file.filename == "":
        flash("Please upload an Excel or CSV file with your variance data.")
        return redirect(url_for("index"))

    title = request.form.get("title", "").strip() or "Variance Analysis"
    subtitle = request.form.get("subtitle", "").strip()
    comparison_label = request.form.get("comparison_label", "").strip() or "Budget"
    currency_suffix = request.form.get("currency_suffix", "").strip() or "M"
    try:
        materiality_pct = float(request.form.get("materiality_pct", "10") or 10)
    except ValueError:
        materiality_pct = 10.0

    sid = uuid.uuid4().hex[:12]
    sess_dir = os.path.join(SESSIONS_DIR, sid)
    os.makedirs(sess_dir, exist_ok=True)
    data_path = os.path.join(sess_dir, "input_" + data_file.filename)
    data_file.save(data_path)

    try:
        if data_path.lower().endswith(".csv"):
            df = pd.read_csv(data_path)
        else:
            df = pd.read_excel(data_path)
    except Exception:
        flash("Couldn't read that file. Please upload a valid .xlsx or .csv.")
        shutil.rmtree(sess_dir, ignore_errors=True)
        return redirect(url_for("index"))

    df = df.dropna(how="all").dropna(axis=1, how="all")
    label_col, actual_col, comparison_col = detect_columns(df)
    if actual_col is None or comparison_col is None:
        flash("Couldn't find two numeric columns to compare (e.g. Actual and Budget). "
              "Make sure your file has a label column plus at least two numeric columns.")
        shutil.rmtree(sess_dir, ignore_errors=True)
        return redirect(url_for("index"))

    rows = analyze(df, label_col, actual_col, comparison_col, comparison_label,
                   materiality_pct, currency_suffix)
    if not rows:
        flash("No usable rows found. Check that your label column and numeric columns line up.")
        shutil.rmtree(sess_dir, ignore_errors=True)
        return redirect(url_for("index"))

    config = {
        "title": title, "subtitle": subtitle, "comparison_label": comparison_label,
        "currency_suffix": currency_suffix, "materiality_pct": materiality_pct,
        "label_col": str(label_col), "actual_col": str(actual_col), "comparison_col": str(comparison_col),
    }
    with open(os.path.join(sess_dir, "config.json"), "w") as f:
        json.dump(config, f)
    with open(os.path.join(sess_dir, "rows.json"), "w") as f:
        json.dump(rows, f)

    return redirect(url_for("review", sid=sid))


@app.route("/review/<sid>", methods=["GET"])
def review(sid):
    sess_dir = _session_dir(sid)
    rows_path = os.path.join(sess_dir, "rows.json")
    config_path = os.path.join(sess_dir, "config.json")
    if not os.path.exists(rows_path):
        flash("That session has expired. Please upload again.")
        return redirect(url_for("index"))
    with open(rows_path) as f:
        rows = json.load(f)
    with open(config_path) as f:
        config = json.load(f)

    material_count = sum(1 for r in rows if r["is_material"])
    favorable_count = sum(1 for r in rows if r["tag"] == "Favorable")
    unfavorable_count = sum(1 for r in rows if r["tag"] == "Unfavorable")

    return render_template(
        "review.html", sid=sid, rows=rows, config=config,
        material_count=material_count, favorable_count=favorable_count,
        unfavorable_count=unfavorable_count, total_count=len(rows),
    )


@app.route("/export/<sid>", methods=["POST"])
def export(sid):
    sess_dir = _session_dir(sid)
    rows_path = os.path.join(sess_dir, "rows.json")
    config_path = os.path.join(sess_dir, "config.json")
    if not os.path.exists(rows_path):
        flash("That session has expired. Please upload again.")
        return redirect(url_for("index"))
    with open(rows_path) as f:
        rows = json.load(f)
    with open(config_path) as f:
        config = json.load(f)

    for i, row in enumerate(rows):
        row["is_material"] = request.form.get(f"include_{i}") == "on"
        edited = request.form.get(f"commentary_{i}", "").strip()
        if edited:
            row["commentary"] = edited

    with open(rows_path, "w") as f:
        json.dump(rows, f)

    pptx_path = os.path.join(sess_dir, "variance_analysis.pptx")
    xlsx_path = os.path.join(sess_dir, "variance_analysis.xlsx")
    build_pptx(rows, config["title"], config["subtitle"], config["comparison_label"], pptx_path)
    build_xlsx(rows, config["title"], config["comparison_label"], xlsx_path)

    preview_path = os.path.join(sess_dir, "preview.png")
    rendering_ok = render_pptx_preview(pptx_path, preview_path)

    with open(os.path.join(sess_dir, "meta.json"), "w") as f:
        json.dump({"rendering_ok": rendering_ok}, f)

    return redirect(url_for("result", sid=sid))


@app.route("/result/<sid>", methods=["GET"])
def result(sid):
    sess_dir = _session_dir(sid)
    pptx_path = os.path.join(sess_dir, "variance_analysis.pptx")
    if not os.path.exists(pptx_path):
        flash("That session has expired. Please upload again.")
        return redirect(url_for("index"))
    meta_path = os.path.join(sess_dir, "meta.json")
    rendering_ok = False
    if os.path.exists(meta_path):
        with open(meta_path) as f:
            rendering_ok = json.load(f).get("rendering_ok", False)
    return render_template("result.html", sid=sid, rendering_ok=rendering_ok,
                            soffice_missing=not (SOFFICE_PATH and HAVE_FITZ))


@app.route("/preview_image/<sid>", methods=["GET"])
def preview_image(sid):
    sess_dir = _session_dir(sid)
    path = os.path.join(sess_dir, "preview.png")
    if not os.path.exists(path):
        abort(404)
    return send_file(path, mimetype="image/png")


@app.route("/download/<sid>/<which>", methods=["GET"])
def download(sid, which):
    sess_dir = _session_dir(sid)
    if which == "pptx":
        path = os.path.join(sess_dir, "variance_analysis.pptx")
        name = "variance_analysis.pptx"
        mimetype = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    elif which == "xlsx":
        path = os.path.join(sess_dir, "variance_analysis.xlsx")
        name = "variance_analysis.xlsx"
        mimetype = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
    else:
        abort(404)
    if not os.path.exists(path):
        abort(404)
    return send_file(path, as_attachment=True, download_name=name, mimetype=mimetype)


def _open_browser():
    import webbrowser
    import sys as _sys
    time.sleep(1.2)
    url = "http://127.0.0.1:5060"
    chrome_candidates = []
    if _sys.platform == "darwin":
        chrome_candidates.append("open -a 'Google Chrome' %s")
    elif _sys.platform.startswith("win"):
        for p in (r"C:\Program Files\Google\Chrome\Application\chrome.exe",
                  r"C:\Program Files (x86)\Google\Chrome\Application\chrome.exe"):
            if os.path.exists(p):
                chrome_candidates.append(p.replace("\\", "\\\\") + " %s")
    else:
        chrome_candidates += ["google-chrome %s", "chromium-browser %s", "chromium %s"]
    for template in chrome_candidates:
        try:
            webbrowser.get(template).open(url)
            return
        except webbrowser.Error:
            continue
    webbrowser.open(url)


if __name__ == "__main__":
    import sys
    if os.environ.get("WERKZEUG_RUN_MAIN") != "true":
        threading.Thread(target=_open_browser, daemon=True).start()
    print("Starting Variance Insights at http://127.0.0.1:5060")
    app.run(host="127.0.0.1", port=5060, debug=("--debug" in sys.argv))
